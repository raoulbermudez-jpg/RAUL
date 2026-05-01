"""
pendrive_pipeline.py
Pipeline nocturno: inventario D: + comparación KB + extracción de texto
Genteca Knowledge Base — Sistema /RAUL/
Ejecutar: python pendrive_pipeline.py
"""

import os
import sys
import json
import hashlib
import shutil
from pathlib import Path
from datetime import datetime
import logging

# ── Configuración ──────────────────────────────────────────────────────────────
PENDRIVE         = Path("D:/")
KB_GENTECA       = Path("C:/RAUL/02-knowledge-base/02-domains/01-genteca")
STAGING          = Path("C:/RAUL/01-inbox/03-raw-sources/genteca/pendrive-D")
ASSETS_BASE      = Path("C:/RAUL/02-knowledge-base/02-domains/01-genteca/assets")
REPORT_DIR       = Path("C:/RAUL/04-system/05-indexes")
LOG_FILE         = REPORT_DIR / f"pendrive_pipeline_{datetime.now().strftime('%Y%m%d_%H%M')}.log"

# Líneas de producto
LINEAS = {
    "01. Exceline":            "exceline",
    "02. Exceline-Profesional": "exceline-profesional",
    "03. Genius":              "genius",
}

# Tipos procesables para extracción de texto
TEXT_EXTRACTABLE = {".pdf", ".docx", ".doc", ".pptx", ".ppt", ".txt", ".xlsx"}
# Tipos de imagen/asset (catalogar, no extraer texto)
ASSET_TYPES      = {".jpg", ".jpeg", ".png", ".tif", ".tiff", ".gif", ".bmp", ".svg", ".eps"}
# Tipos de diseño (catalogar solo, no procesar)
DESIGN_TYPES     = {".ai", ".psd", ".indd", ".fh11", ".fh9", ".fh10", ".dxf", ".dwg", ".vsd"}

# ── Logging ────────────────────────────────────────────────────────────────────
REPORT_DIR.mkdir(parents=True, exist_ok=True)
STAGING.mkdir(parents=True, exist_ok=True)

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    handlers=[
        logging.FileHandler(LOG_FILE, encoding="utf-8"),
        logging.StreamHandler(sys.stdout),
    ],
)
log = logging.getLogger(__name__)


# ── Utilidades ─────────────────────────────────────────────────────────────────
def slugify(name: str) -> str:
    import re
    name = name.lower()
    name = re.sub(r"[áàä]", "a", name)
    name = re.sub(r"[éèë]", "e", name)
    name = re.sub(r"[íìï]", "i", name)
    name = re.sub(r"[óòö]", "o", name)
    name = re.sub(r"[úùü]", "u", name)
    name = re.sub(r"ñ", "n", name)
    name = re.sub(r"[^a-z0-9]+", "-", name)
    return name.strip("-")


def classify_asset(filename: str) -> str:
    """Classify asset type based on filename patterns."""
    fname = filename.upper()
    if "GLA-T" in fname or "GLAT" in fname or "_T_" in fname and "GLA" in fname:
        return "packaging/gla-t"
    if "GLA-R" in fname or "GLAR" in fname or "_R_" in fname and "GLA" in fname:
        return "packaging/gla-r"
    if any(x in fname for x in ["EMPAQU", "PACKAG", "BOX", "CAJA", "BURBUJA"]):
        return "packaging"
    if any(x in fname for x in ["DIAGRAMA", "DIAGRAM", "WIRING", "CONEXION", "ESQUEMA"]):
        return "diagrams"
    if any(x in fname for x in ["FOTO", "PHOTO", "PRODUCT", "PROTECTOR"]):
        return "products"
    return "uncoded"


def pdf_has_text(path: Path) -> bool:
    """Check if PDF has extractable text (not image-only)."""
    try:
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages[:3]:  # check first 3 pages
                if page.extract_text():
                    return True
        return False
    except Exception:
        return False


def extract_pdf(path: Path) -> str:
    """Extract text from PDF using pdfplumber, fallback to PyMuPDF."""
    text_parts = []
    try:
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                t = page.extract_text()
                if t:
                    text_parts.append(t)
        if text_parts:
            return "\n\n".join(text_parts)
    except Exception:
        pass

    # Fallback: PyMuPDF
    try:
        import fitz
        doc = fitz.open(str(path))
        for page in doc:
            text_parts.append(page.get_text())
        doc.close()
        return "\n\n".join(text_parts)
    except Exception as e:
        return f"[ERROR EXTRACCION: {e}]"


def extract_docx(path: Path) -> str:
    try:
        from docx import Document
        doc = Document(str(path))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except Exception as e:
        return f"[ERROR EXTRACCION: {e}]"


def extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text)
        return "\n\n".join(parts)
    except Exception as e:
        return f"[ERROR EXTRACCION: {e}]"


def extract_xlsx(path: Path) -> str:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        parts = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                line = "\t".join(str(c) for c in row if c is not None)
                if line.strip():
                    parts.append(line)
        return "\n".join(parts)
    except Exception as e:
        return f"[ERROR EXTRACCION: {e}]"


# ── Fase 1: Inventario + Comparación ──────────────────────────────────────────
def phase1_inventory():
    log.info("=" * 60)
    log.info("FASE 1 — Inventario D: + Comparación KB")
    log.info("=" * 60)

    # Leer KB existente
    kb_files = set()
    for f in KB_GENTECA.rglob("*.md"):
        kb_files.add(slugify(f.stem))
    log.info(f"KB actual: {len(kb_files)} archivos MD")

    inventory = []
    delta_text = []   # procesables que no están en KB
    delta_assets = [] # imágenes nuevas
    already_have = []
    design_catalog = []

    for folder_name, linea in LINEAS.items():
        folder = PENDRIVE / folder_name
        if not folder.exists():
            log.warning(f"Carpeta no encontrada: {folder}")
            continue

        for f in folder.rglob("*"):
            if not f.is_file():
                continue
            ext = f.suffix.lower()
            slug = slugify(f.stem)
            size_kb = round(f.stat().st_size / 1024, 1)

            entry = {
                "linea": linea,
                "path": str(f),
                "name": f.name,
                "slug": slug,
                "ext": ext,
                "size_kb": size_kb,
                "in_kb": slug in kb_files,
            }
            inventory.append(entry)

            if slug in kb_files:
                already_have.append(entry)
            elif ext in TEXT_EXTRACTABLE:
                delta_text.append(entry)
            elif ext in ASSET_TYPES:
                delta_assets.append(entry)
            elif ext in {e.lower() for e in DESIGN_TYPES}:
                design_catalog.append(entry)

    # Guardar inventario JSON
    inv_file = REPORT_DIR / "pendrive_D_inventory.json"
    with open(inv_file, "w", encoding="utf-8") as fh:
        json.dump(inventory, fh, ensure_ascii=False, indent=2)

    # Reporte MD
    report_lines = [
        f"# Inventario Pendrive D: — {datetime.now().strftime('%Y-%m-%d %H:%M')}",
        "",
        "## Resumen",
        f"- **Total archivos en D:** {len(inventory)}",
        f"- **Ya en KB:** {len(already_have)}",
        f"- **Procesables nuevos (texto):** {len(delta_text)}",
        f"- **Assets nuevos (imágenes):** {len(delta_assets)}",
        f"- **Diseño solo (catalogar):** {len(design_catalog)}",
        "",
        "## Delta — Procesables por línea",
    ]
    by_linea = {}
    for e in delta_text:
        by_linea.setdefault(e["linea"], []).append(e)
    for linea, entries in sorted(by_linea.items()):
        by_ext = {}
        for e in entries:
            by_ext.setdefault(e["ext"], []).append(e)
        report_lines.append(f"\n### {linea} ({len(entries)} archivos)")
        for ext, items in sorted(by_ext.items(), key=lambda x: -len(x[1])):
            report_lines.append(f"- `{ext}`: {len(items)} archivos")

    report_lines += [
        "",
        "## Assets por clasificar",
        f"- GLA-T (frente empaque): {sum(1 for e in delta_assets if 'gla-t' in classify_asset(e['name']))}",
        f"- GLA-R (reverso empaque): {sum(1 for e in delta_assets if 'gla-r' in classify_asset(e['name']))}",
        f"- Productos/fotos: {sum(1 for e in delta_assets if classify_asset(e['name']) == 'products')}",
        f"- Diagramas: {sum(1 for e in delta_assets if classify_asset(e['name']) == 'diagrams')}",
        f"- Sin clasificar: {sum(1 for e in delta_assets if classify_asset(e['name']) == 'uncoded')}",
    ]

    report_file = REPORT_DIR / "pendrive_D_report.md"
    with open(report_file, "w", encoding="utf-8") as fh:
        fh.write("\n".join(report_lines))

    log.info(f"Reporte guardado: {report_file}")
    log.info(f"Total: {len(inventory)} | Ya en KB: {len(already_have)} | Delta texto: {len(delta_text)} | Delta assets: {len(delta_assets)}")

    return delta_text, delta_assets, design_catalog


# ── Fase 2: Extracción de texto ────────────────────────────────────────────────
def phase2_extract(delta_text: list):
    log.info("=" * 60)
    log.info(f"FASE 2 — Extracción de texto ({len(delta_text)} archivos)")
    log.info("=" * 60)

    done = 0
    skipped = 0
    errors = 0

    for i, entry in enumerate(delta_text, 1):
        path = Path(entry["path"])
        linea = entry["linea"]
        ext = entry["ext"]
        slug = entry["slug"]

        # Destino en staging
        staging_linea = STAGING / linea
        staging_linea.mkdir(parents=True, exist_ok=True)
        out_path = staging_linea / f"{slug}.txt"

        # Skip si ya extraído
        if out_path.exists():
            skipped += 1
            continue

        try:
            if ext == ".pdf":
                if not pdf_has_text(path):
                    # PDF de solo imagen — catalogar como asset visual
                    with open(out_path, "w", encoding="utf-8") as fh:
                        fh.write(f"[PDF_IMAGEN_SOLO]\nArchivo: {path.name}\nLinea: {linea}\nTamaño: {entry['size_kb']} KB\n")
                else:
                    text = extract_pdf(path)
                    with open(out_path, "w", encoding="utf-8") as fh:
                        fh.write(f"[FUENTE: {path.name}]\n[LINEA: {linea}]\n\n{text}")
            elif ext in {".docx", ".doc"}:
                text = extract_docx(path)
                with open(out_path, "w", encoding="utf-8") as fh:
                    fh.write(f"[FUENTE: {path.name}]\n[LINEA: {linea}]\n\n{text}")
            elif ext in {".pptx", ".ppt"}:
                text = extract_pptx(path)
                with open(out_path, "w", encoding="utf-8") as fh:
                    fh.write(f"[FUENTE: {path.name}]\n[LINEA: {linea}]\n\n{text}")
            elif ext in {".xlsx"}:
                text = extract_xlsx(path)
                with open(out_path, "w", encoding="utf-8") as fh:
                    fh.write(f"[FUENTE: {path.name}]\n[LINEA: {linea}]\n\n{text}")
            elif ext == ".txt":
                shutil.copy2(path, out_path)
            done += 1

            if i % 50 == 0:
                log.info(f"  Progreso: {i}/{len(delta_text)} | OK: {done} | Skip: {skipped} | Error: {errors}")

        except Exception as e:
            errors += 1
            log.error(f"  ERROR {path.name}: {e}")

    log.info(f"Fase 2 completa — Extraídos: {done} | Saltados: {skipped} | Errores: {errors}")


# ── Fase 3: Catálogo de assets ─────────────────────────────────────────────────
def phase3_catalog_assets(delta_assets: list, design_catalog: list):
    log.info("=" * 60)
    log.info(f"FASE 3 — Catálogo de assets ({len(delta_assets)} imágenes + {len(design_catalog)} diseño)")
    log.info("=" * 60)

    catalog = []

    for entry in delta_assets + design_catalog:
        asset_type = classify_asset(entry["name"])
        catalog.append({
            "name": entry["name"],
            "linea": entry["linea"],
            "type": asset_type,
            "ext": entry["ext"],
            "size_kb": entry["size_kb"],
            "source_path": entry["path"],
            "cite": True,  # todos son Genteca
        })

    # Guardar catálogo JSON
    cat_file = REPORT_DIR / "pendrive_D_assets_catalog.json"
    with open(cat_file, "w", encoding="utf-8") as fh:
        json.dump(catalog, fh, ensure_ascii=False, indent=2)

    # Resumen MD para assets
    lines = [
        f"# Catálogo de Assets — Pendrive D: — {datetime.now().strftime('%Y-%m-%d')}",
        "",
        "| Archivo | Línea | Tipo | Ext | KB |",
        "|---|---|---|---|---|",
    ]
    for a in sorted(catalog, key=lambda x: (x["linea"], x["type"], x["name"])):
        lines.append(f"| {a['name']} | {a['linea']} | {a['type']} | {a['ext']} | {a['size_kb']} |")

    cat_md = REPORT_DIR / "pendrive_D_assets_catalog.md"
    with open(cat_md, "w", encoding="utf-8") as fh:
        fh.write("\n".join(lines))

    log.info(f"Catálogo guardado: {cat_md} ({len(catalog)} assets)")


# ── Main ───────────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    log.info("PENDRIVE PIPELINE — Genteca KB — Sistema /RAUL/")
    log.info(f"Inicio: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
    log.info(f"Fuente: {PENDRIVE}")
    log.info(f"Destino staging: {STAGING}")

    delta_text, delta_assets, design_catalog = phase1_inventory()
    phase2_extract(delta_text)
    phase3_catalog_assets(delta_assets, design_catalog)

    log.info("=" * 60)
    log.info(f"PIPELINE COMPLETO — {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
    log.info(f"Log completo: {LOG_FILE}")
    log.info(f"Reporte: {REPORT_DIR}/pendrive_D_report.md")
    log.info(f"Staging: {STAGING}")
    log.info("Mañana: ejecuta Fase 4 (formateo KB con Claude) sobre el staging.")
    log.info("=" * 60)
