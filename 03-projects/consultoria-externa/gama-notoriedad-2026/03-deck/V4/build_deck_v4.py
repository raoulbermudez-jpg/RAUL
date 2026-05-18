"""
build_deck_v4.py — Vivienne render script
Notoriedad Gama 2026 — Deck V4 (VI-2)
Genera: 2026-05-17_Notoriedad-Gama-2026_V4.pptx

Basado en:
  VI-1a_OUTLINE-DECK-V4.md (SSOT)
  BR-2 gate V4 (2026-05-17_genteca-V4_claim-approval-log_v1.md)

Aplicaciones BR-2 codificadas:
  DW-1 → versión [A] con nombre Azahara Betancourt + C-LIT-DW1 (D3, R5) y C-LIT-DW1-bis (C7)
  DW-2 → lenguaje condicional + C-LIT-DW2-a (D4, R6, E3) + C-LIT-DW2-b (E3 long-term)
  DW-3 → wording aprobado C-LIT-DW3-a (D6) + C-LIT-DW3-b (no renderizado en deck V4 — es R6 del Resumen Ejecutivo)
  DW-4 → atribución por segmento/edad C-LIT-DW4-c + enmarcado C-LIT-DW4-a ANTES + caveat C-LIT-DW4-b DESPUÉS (D5)
  DW-5 → etiqueta [Hipótesis V4 — evidencia convergente] sin explicar cambio de status (A2, A5, E2)
  CO-1 → EXCLUIR bullet Rio-proteínas de C3
  CO-2 → mezcla equilibrada: wording literal de BR-2 §3 CO-2 en C1, C2, C8, E2
  CO-3 → Opción C híbrida: hallazgos digitales con nota "ampliación del brief original" en E3, E5
  Placeholders [PENDIENTE DECISIÓN DW-X BRUNA:...] eliminados

Caveats literales obligatorios (12 strings):
  C-LIT-DW1, C-LIT-DW1-bis, C-LIT-DW2-a, C-LIT-DW2-b
  C-LIT-DW3-a, C-LIT-DW4-a, C-LIT-DW4-b, C-LIT-DW4-c
  C-LIT-CV-WOW-001, C-LIT-CV-WOW-002, C-LIT-CV-WOW-005
  C-LIT-IN7-1.4

Notas técnicas:
  - Sección B (~25 slides V3) renderizada como 2 slides placeholder con instrucción al presenter
  - Total deck: ~32 slides reales (A x6 + B-placeholder x2 + C x8 + D x6 + E x5 + portada + cierre)
  - Paleta heredada de V3: rojo #7A1212, neutro #333333, gris #808080, fondo blanco #FFFFFF, sección header dark red #4A0A0A, warning #FFF3E0/#D97306
  - Dimensiones V3: 12191695 x 6858000 EMU (widescreen 16:9 aprox 13.4 x 7.55 in)
  - Sin doble codificación font-name porque V3 usa default theme font; se aplica Calibri como fallback.
"""

import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Emu, Pt, Inches, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu as EMU
import copy

# ── Output path ──────────────────────────────────────────────────────────────
OUTPUT_PATH = r"C:\RAUL\03-projects\consultoria-externa\gama-notoriedad-2026\03-deck\V4\2026-05-17_Notoriedad-Gama-2026_V4.pptx"

# ── Palette (inherited from V3) ───────────────────────────────────────────────
C_RED       = RGBColor(0x7A, 0x12, 0x12)   # primary brand red
C_DARKRED   = RGBColor(0x4A, 0x0A, 0x0A)   # dark red (highlights/starred items)
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK      = RGBColor(0x33, 0x33, 0x33)   # body text
C_GRAY      = RGBColor(0x80, 0x80, 0x80)   # footer text
C_WARN_BG   = RGBColor(0xFF, 0xF3, 0xE0)   # warning box bg
C_WARN_FG   = RGBColor(0xD9, 0x73, 0x06)   # warning text/border
C_LIGHT_BG  = RGBColor(0xF8, 0xF8, 0xF8)   # light box bg
C_GREEN     = RGBColor(0x1B, 0x7A, 0x3C)   # approved/positive indicator
C_BLUE      = RGBColor(0x0D, 0x47, 0x8A)   # Tipo B label
C_AMBER     = RGBColor(0xD9, 0x73, 0x06)   # Tipo C label

# ── Slide dimensions (matching V3) ────────────────────────────────────────────
SLIDE_W = 12191695
SLIDE_H = 6858000

# ── Caveat literal strings (exact from BR-2 §5) ──────────────────────────────
C_LIT_DW1 = (
    "Cita textual de participante de focus group (investigacion cualitativa, mayo 2026). "
    "El termino 'sifrinaje' es vocabulario espontaneo del informante — no una caracterizacion "
    "de Gama hecha por el equipo analitico. Refleja como segmentos de no-usuarios codifican "
    "la barrera de entrada percibida."
)

C_LIT_DW1_BIS = C_LIT_DW1  # mismo texto, adaptado para verbatim de David

C_LIT_DW2_A = (
    "Este hallazgo describe el imaginario de marca que el shopper ya proyecta espontaneamente "
    "— no propone una campana especifica de personificacion. La materializacion creativa "
    "(figura concreta, voz, presencia visual) requiere prueba de concepto cuantitativa "
    "con el target antes de cualquier ejecucion."
)

C_LIT_DW2_B = (
    "Cualquier desarrollo creativo en esta direccion requiere verificacion previa de Gama "
    "sobre compatibilidad con (a) contratos de imagen vigentes, (b) plataforma de marca "
    "existente, (c) posicionamiento global del grupo (si aplica). El hallazgo cuali no "
    "sustituye esta verificacion."
)

C_LIT_DW3_A = (
    "El cuali V4 enriquece el perfil de los Pragmaticos Convertibles (33% del mercado) con "
    "una dimension que el cuestionario cuantitativo no podia capturar: la barrera de "
    "conversion no es solo el precio como calculo economico, tambien es el precio como senal "
    "de pertenencia social. Esta capa adicional refina la estrategia de activacion: el anzuelo "
    "correcto es una razon especifica para venir, no una oferta de descuento."
)

C_LIT_DW4_A = (
    "El siguiente hallazgo del canal digital identifica una oportunidad operativa de alto "
    "retorno con baja inversion. La friccion documentada es resoluble con cambio de producto "
    "digital (sincronizacion de inventario en tiempo real), no requiere reposicionamiento "
    "ni nueva investigacion."
)

C_LIT_DW4_B = (
    "El verbatim ilustra el patron documentado en 3 de los 4 documentos de uso digital "
    "(Corpus B). La inconsistencia de inventario es la friccion mas reportada del canal "
    "y la de mayor costo de confianza."
)

C_LIT_DW4_C = "Mujer, 42 anos, segmento Frecuentes 31-50 (verbatim de sesion online cualitativa)"

C_LIT_CV_WOW_001 = (
    "Los resultados 2025 se presentan sin ponderacion muestral (factor de ponderacion no "
    "disponible en BBDD). Los estimados 2025 pueden tener sesgo de diseno si la muestra "
    "original era estratificada."
)

C_LIT_CV_WOW_002 = (
    "Los cambios WoW deben interpretarse con precaucion: la composicion geografica de las "
    "muestras difiere significativamente entre olas (Libertador sobrerepresentado en 2025). "
    "Los cambios en variables con dispersion geografica pueden reflejar composicion, no cambio real."
)

C_LIT_CV_WOW_005 = (
    "Los analisis por NSE no fueron pre-registrados en el diseno 2025. Deben interpretarse "
    "como hipotesis a confirmar en ola 2027."
)

C_LIT_IN7_14 = (
    "Los hallazgos cualitativos derivan de analisis tematico de unico analista (sin doble "
    "codificacion inter-codificador). Es la practica estandar en investigacion cualitativa "
    "aplicada cuando recursos no permiten Kappa; limita replicabilidad pero no validez "
    "de la interpretacion."
)

FOOTER_TEXT = "Notoriedad Gama 2026 · V4 · n=402 (2026) n=785 (2025) · WoW + Cuali · Confidencial NDA"

# CO-2 approved wording (literal from BR-2 §3 CO-2)
CO2_MAIN = (
    "Gama es la unica gran cadena del mercado que mantiene posicion estable entre 2025 y 2026 "
    "— fortaleza defensiva en mercado turbulento. La estabilidad es resultado de los activos "
    "diferenciadores (24h, atencion relacional) que los competidores no replican. Al mismo tiempo, "
    "en el segmento C+/C — natural de Gama — Rio y Paramo crecen +25.8pp y +22.3pp TOM. La senal "
    "direccional es que el espacio competitivo se esta estrechando, aunque la posicion actual de Gama es solida."
)

CO2_SYNTHESIS = (
    "La estrategia defensiva (retener Nucleo Leal + Mayoria Exigente) preserva la posicion actual. "
    "La estrategia ofensiva (capturar el espacio disponible de CM + responder al avance de Rio/Paramo "
    "en C+/C) es necesaria para que la estabilidad de hoy no se erosione en 2027."
)

# ── Helpers ───────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = Emu(SLIDE_W)
    prs.slide_height = Emu(SLIDE_H)
    return prs

def blank_layout(prs):
    return prs.slide_layouts[6]  # Blank

def add_slide(prs):
    return prs.slides.add_slide(blank_layout(prs))

def set_bg(slide, color=None):
    """Set slide background to solid color or leave transparent."""
    bg = slide.background
    fill = bg.fill
    if color:
        fill.solid()
        fill.fore_color.rgb = color
    else:
        fill.background()

def add_rect(slide, l, t, w, h, fill_color=None, line_color=None, line_width=Pt(0)):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Emu(l), Emu(t), Emu(w), Emu(h)
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, l, t, w, h, text, font_size=Pt(11), color=C_DARK,
                bold=False, italic=False, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = "Calibri"
    return txb

def add_para(tf, text, font_size=Pt(10), color=C_DARK, bold=False, italic=False,
             align=PP_ALIGN.LEFT, indent_level=0):
    from pptx.oxml.ns import qn
    import lxml.etree as etree
    p = tf.add_paragraph()
    p.alignment = align
    p.level = indent_level
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = "Calibri"
    return p

def add_footer(slide, text=FOOTER_TEXT):
    add_textbox(slide,
                l=200000, t=SLIDE_H - 220000,
                w=SLIDE_W - 400000, h=200000,
                text=text,
                font_size=Pt(7.5),
                color=C_GRAY,
                align=PP_ALIGN.CENTER)

def header_bar(slide, section_label, title_text, section_bg=C_RED):
    """Dark red top header bar with section label and title."""
    # Top bar (dark)
    bar_h = 950000
    add_rect(slide, 0, 0, SLIDE_W, bar_h, fill_color=section_bg)
    # Section label (small, white)
    add_textbox(slide, 300000, 60000, 2000000, 350000,
                text=section_label, font_size=Pt(9), color=C_WHITE, bold=False)
    # Slide title (larger, white)
    add_textbox(slide, 300000, 350000, SLIDE_W - 600000, 550000,
                text=title_text, font_size=Pt(19), color=C_WHITE, bold=True)

def label_box(slide, label_text, l, t, w=2800000, h=240000,
              bg=C_LIGHT_BG, fg=C_DARKRED):
    """Claim-type label badge."""
    add_rect(slide, l, t, w, h, fill_color=bg)
    add_textbox(slide, l + 60000, t + 20000, w - 120000, h - 40000,
                text=label_text, font_size=Pt(7.5), color=fg, bold=True)

def body_block(slide, l, t, w, h, title=None, bullets=None,
               font_title=Pt(12), font_body=Pt(10)):
    """
    Adds a white rounded-rect-style block with optional sub-title and bullet list.
    Returns the text frame for further manipulation.
    """
    # background box
    box = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = title
        r.font.size = font_title
        r.font.color.rgb = C_RED
        r.font.bold = True
        r.font.name = "Calibri"
    if bullets:
        for b in bullets:
            p = tf.add_paragraph() if not first else tf.paragraphs[0]
            first = False
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = b
            r.font.size = font_body
            r.font.color.rgb = C_DARK
            r.font.name = "Calibri"
    return tf

def caveat_box(slide, text, l, t, w, h, bg=C_WARN_BG, fg=C_WARN_FG, size=Pt(8)):
    """Warning/caveat box."""
    add_rect(slide, l, t, w, h, fill_color=bg, line_color=fg, line_width=Pt(1))
    add_textbox(slide, l + 100000, t + 80000, w - 200000, h - 160000,
                text=text, font_size=size, color=fg, wrap=True)

def info_box(slide, text, l, t, w, h, bg=C_LIGHT_BG, fg=C_DARKRED, size=Pt(8.5)):
    """Info/highlight box (light gray)."""
    add_rect(slide, l, t, w, h, fill_color=bg)
    add_textbox(slide, l + 100000, t + 80000, w - 200000, h - 160000,
                text=text, font_size=size, color=fg, wrap=True)

def quote_box(slide, quote, attribution, l, t, w, h):
    """Verbatim quote block."""
    add_rect(slide, l, t, w, h, fill_color=C_LIGHT_BG)
    # quote mark strip
    add_rect(slide, l, t, 120000, h, fill_color=C_RED)
    txb = slide.shapes.add_textbox(Emu(l + 160000), Emu(t + 80000),
                                   Emu(w - 280000), Emu(h - 200000))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = f'"{quote}"'
    r.font.size = Pt(9.5)
    r.font.color.rgb = C_DARKRED
    r.font.italic = True
    r.font.name = "Calibri"
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = attribution
    r2.font.size = Pt(8)
    r2.font.color.rgb = C_GRAY
    r2.font.name = "Calibri"

# ── CONTENT DEFINITIONS ────────────────────────────────────────────────────────

def make_cover(prs):
    slide = add_slide(prs)
    set_bg(slide, C_WHITE)
    # left dark panel
    add_rect(slide, 0, 0, 3400000, SLIDE_H, fill_color=C_RED)
    add_textbox(slide, 200000, 500000, 3000000, 900000,
                "GAMA", Pt(54), C_WHITE, bold=True, align=PP_ALIGN.LEFT)
    add_textbox(slide, 200000, 1300000, 3000000, 600000,
                "2026", Pt(40), C_WHITE, bold=False)
    # Version badge
    badge = slide.shapes.add_shape(1, Emu(200000), Emu(1950000), Emu(700000), Emu(320000))
    badge.fill.solid(); badge.fill.fore_color.rgb = C_WHITE
    badge.line.fill.background()
    add_textbox(slide, 230000, 1970000, 640000, 280000,
                "V4", Pt(13), C_RED, bold=True, align=PP_ALIGN.CENTER)
    # Right panel text
    add_textbox(slide, 3700000, 700000, SLIDE_W - 4000000, 700000,
                "Notoriedad y Preferencia de Marca — Gama 2026",
                Pt(26), C_RED, bold=True)
    add_textbox(slide, 3700000, 1450000, SLIDE_W - 4000000, 450000,
                "Capa Analitica Complementaria V4 · Wave-over-Wave 2025<>2026 + Cuali Profundo + Triangulacion",
                Pt(12), C_DARK)
    # Metadata block
    meta_lines = [
        "Fecha: 2026-05-17",
        "Equipo analitico: Cora Urrea + Raoul Bermudez",
        "Version: V4 — complementa V3 (2026-05-17)",
        "Gate Bruna: BR-2 V4 aprobado con ajustes | Confidencial / NDA",
    ]
    for idx, line in enumerate(meta_lines):
        add_textbox(slide, 3700000, 2200000 + idx * 280000,
                    SLIDE_W - 4200000, 260000,
                    line, Pt(9), C_GRAY)
    add_footer(slide)
    return slide

# ── SECCIÓN A ──────────────────────────────────────────────────────────────────

def make_A1(prs):
    slide = add_slide(prs)
    set_bg(slide, C_WHITE)
    header_bar(slide, "Seccion A — Recap V3 y Encuadre V4",
               "A1 — V4 amplia la foto sin cambiar el veredicto de V3")
    body_y = 1050000
    info_box(slide,
             "TAKEAWAY: V4 suma evidencia longitudinal (WoW 2025<>2026) y mecanismo cualitativo "
             "profundo para enriquecer — no contradecir — las conclusiones V3.",
             l=300000, t=body_y, w=SLIDE_W - 600000, h=360000)
    bullets = [
        "V3 (2026-05-17) respondia: que mueve la preferencia de Gama hoy? "
        "-> cuatro argumentos MECE con convergencia de 4 metodos estadisticos.",
        "V4 responde dos preguntas nuevas: la posicion de Gama cambio entre 2025 y 2026? "
        "por que operan los mecanismos que mide el cuanti?",
        "Las tres capas V4: (1) wave-over-wave cuanti 2025<>2026; "
        "(2) analisis cualitativo profundo de 12 documentos FG reales; "
        "(3) triangulacion sistematica de ambas capas contra los claims V3.",
        "Resultado de la triangulacion: 0 contradicciones directas entre V4 y V3. "
        "4 matices/reformulaciones de mecanismo. 6 hallazgos genuinamente nuevos.",
        "Regla de lectura: donde V4 reformula, el claim V3 es mas preciso — no incorrecto.",
    ]
    body_block(slide, 300000, body_y + 420000,
               SLIDE_W - 600000, SLIDE_H - body_y - 720000,
               bullets=["  * " + b for b in bullets])
    # Diagram description (text placeholder for visual)
    add_textbox(slide, 300000, SLIDE_H - 650000, SLIDE_W - 600000, 200000,
                "[Visual: Diagrama capas — V3 base + V4 WoW + V4 Cuali. Flecha: 'V3->V4: enriquecimiento sin ruptura']",
                Pt(7.5), C_GRAY, italic=True)
    add_footer(slide)

def make_A2(prs):
    slide = add_slide(prs)
    set_bg(slide, C_WHITE)
    header_bar(slide, "Seccion A — Recap V3 y Encuadre V4",
               "A2 — Los 4 claims V3 de alta certeza siguen en pie")
    info_box(slide,
             "TAKEAWAY: Los cuatro argumentos MECE de V3 estan confirmados o enriquecidos por V4 "
             "— ninguno fue refutado.",
             l=300000, t=1050000, w=SLIDE_W - 600000, h=320000)
    # Table of 4 claims
    claims = [
        ("V3-1 — Atencion driver #1",
         "OR=5.73*** (IC95: 1.6-20.4), triple convergencia.",
         "[Triangulado: cuanti V3 + cuali V4]",
         "Mecanismo reformulado (acompanamiento-guia), driver intacto."),
        ("V3-2 — Precio no driver",
         "OR=1.03 (p=0.966), SHAP #10.",
         "[Triangulado]",
         "El precio SI es barrera de entrada identitaria para no-usuarios; prediccion de preferencia no cambia."),
        ("V3-3 — Gap comunicacional",
         "Recall 0/17 espontaneo, 65% interpreta PTL como mensaje de precio.",
         "[Hipotesis V4 — evidencia convergente]",
         # DW-5 wording obligatorio (BR-2 §2 DW-5)
         "El corpus FG no menciona la campana espontaneamente; WoW muestra debilitamiento "
         "'promociones' -5.1pp (tendencia). El claim sigue siendo la hipotesis mas plausible."),
        ("V3-4 — 3 segmentos k-means",
         "Mayoria Exigente (59%), Pragmaticos Convertibles (33%), Nucleo Leal (8%).",
         "[Triangulado]",
         "Exigente y Leal confirmados por FG. Pragmaticos: mecanismo reformulado (barrera identitaria, no solo precio)."),
    ]
    row_h = 700000
    row_y = 1450000
    for (title, v3, label, v4) in claims:
        # row bg alternating
        add_rect(slide, 300000, row_y, SLIDE_W - 600000, row_h - 40000,
                 fill_color=C_LIGHT_BG)
        add_textbox(slide, 340000, row_y + 40000, 2600000, 280000,
                    title, Pt(9), C_RED, bold=True)
        add_textbox(slide, 340000, row_y + 300000, 2600000, 360000,
                    v3, Pt(8), C_DARK)
        add_textbox(slide, 3000000, row_y + 40000, 2000000, 280000,
                    label, Pt(8), C_BLUE, bold=True)
        add_textbox(slide, 5100000, row_y + 40000, SLIDE_W - 5400000, row_h - 80000,
                    v4, Pt(8.5), C_DARK)
        row_y += row_h

    add_footer(slide)

def make_A3(prs):
    slide = add_slide(prs)
    set_bg(slide, C_WHITE)
    header_bar(slide, "Seccion A — Recap V3 y Encuadre V4",
               "A3 — Los 4 claims V3 vistos con ojos de 2026 vs 2025")
    info_box(slide,
             "TAKEAWAY: La vision longitudinal confirma que V3 fotografiaba una posicion real "
             "y estable — no un artefacto de un solo corte.",
             l=300000, t=1050000, w=SLIDE_W - 600000, h=320000)
    bullets = [
        "Atencion driver #1: el activo relacional de Gama no tiene equivalente en los competidores "
        "que crecen (Rio, Paramo). La estabilidad del embudo de Gama mientras otros crecen "
        "indica que el activo existe y es defensivo — aunque no esta convirtiendo crecimiento aun.",
        "Precio como no-driver: el debilitamiento tendencial de 'menor precio' (-3.9pp, tendencia p=0.091) "
        "y 'promociones' (-5.1pp, p=0.053) es coherente con la lectura V3.",
        "3 segmentos: el NSE C+/C (dominio del Nucleo Leal y Mayoria Exigente) es estable sin captura "
        "del crecimiento disponible — coherente con los segmentos ya-convertidos vs. no-convertidos.",
        "Gap comunicacional: los FG del cuali no mencionan la campana espontaneamente en ninguno "
        "de los 12 documentos / 42,094 chars. Silencio consistente con el gap medido cuantitativamente.",
    ]
    y = 1450000
    for b in bullets:
        add_textbox(slide, 350000, y, SLIDE_W - 700000, 500000,
                    "  " + b, Pt(10), C_DARK, wrap=True)
        y += 520000
    caveat_box(slide, "Caveat: " + C_LIT_CV_WOW_002,
               l=300000, t=SLIDE_H - 550000, w=SLIDE_W - 600000, h=280000)
    add_footer(slide)

def make_A4(prs):
    slide = add_slide(prs)
    set_bg(slide, C_WHITE)
    header_bar(slide, "Seccion A — Recap V3 y Encuadre V4",
               "A4 — Los 6 hallazgos que V4 agrega y V3 no tenia")
    info_box(slide,
             "TAKEAWAY: Seis insights genuinamente nuevos emergen de la combinacion de evidencia "
             "longitudinal y cualitativa — ninguno contradice V3; todos lo hacen mas accionable.",
             l=300000, t=1050000, w=SLIDE_W - 600000, h=320000)
    news = [
        ("V4-NEW-1", "[Triangulado]", "El mecanismo del driver de atencion es acompanamiento-guia, "
         "no reconocimiento nominal. El claim mas comunicable del estudio."),
        ("V4-NEW-2", "[Triangulado]", "La barrera de precio es identitaria-simbolica ('sifrinaje'), "
         "no un calculo economico. Las promociones de precio bajo no la mueven."),
        ("V4-NEW-3", "[Cuali V4]", "Gama Club = activo subexplotado con solucion de bajo costo "
         "(solo visibilidad digital del saldo). La accion de mayor accionabilidad a corto plazo."),
        ("V4-NEW-4", "[Cuali V4 FG]", "Personificacion femenina espontanea ('guia experta de 35-45 anos') "
         "en 6/7 grupos FG — territorio de identidad de marca disponible, sin competidor que lo ocupe."),
        ("V4-NEW-5", "[WoW Confirmado]", "Dinamica competitiva 2025->2026: Rio y Paramo crecen "
         "explosivamente en el segmento natural de Gama (C+/C) mientras Gama permanece estatica."),
        ("V4-NEW-6", "[Cuali V4]", "Canal digital con dos barreras distintas (inconsistencia de "
         "inventario + desconfianza en frescos) que requieren soluciones distintas."),
    ]
    row_h = 550000
    y = 1450000
    for (tag, label, text) in news:
        add_rect(slide, 300000, y, 1100000, row_h - 60000, fill_color=C_RED)
        add_textbox(slide, 320000, y + 80000, 1060000, 380000,
                    tag, Pt(9), C_WHITE, bold=True)
        add_rect(slide, 1460000, y, 1500000, row_h - 60000, fill_color=C_LIGHT_BG)
        add_textbox(slide, 1500000, y + 60000, 1420000, 420000,
                    label, Pt(8), C_BLUE, bold=True)
        add_textbox(slide, 3020000, y + 60000, SLIDE_W - 3320000, row_h - 120000,
                    text, Pt(9.5), C_DARK, wrap=True)
        y += row_h
    add_footer(slide)

def make_A5(prs):
    slide = add_slide(prs)
    set_bg(slide, C_WHITE)
    header_bar(slide, "Seccion A — Recap V3 y Encuadre V4",
               "A5 — 4 reformulaciones de mecanismo: que cambia exactamente en V4 vs V3")
    info_box(slide,
             "TAKEAWAY: V4 reforma el 'como' de cuatro mecanismos V3 sin tocar el 'que' "
             "— el hallazgo se vuelve mas preciso y mas accionable.",
             l=300000, t=1050000, w=SLIDE_W - 600000, h=320000)
    # Table header
    cols = ["Wording V3", "Wording V4", "Implicacion accionable"]
    col_x = [300000, 3900000, 7700000]
    col_w = [3500000, 3700000, 4000000]
    y_hdr = 1450000
    for ci, (cx, cw, ch) in enumerate(zip(col_x, col_w, cols)):
        add_rect(slide, cx, y_hdr, cw - 80000, 340000, fill_color=C_RED)
        add_textbox(slide, cx + 60000, y_hdr + 60000, cw - 140000, 240000,
                    ch, Pt(9), C_WHITE, bold=True)
    rows = [
        ("'Reconocimiento personal por nombre' [V3]",
         "'Acompanamiento-guia — en Gama siempre hay alguien que te ayuda' [V4]",
         "El OR=5.73 es el mismo; el mensaje correcto cambia."),
        ("'54% percibe Gama como cara' — precio como calculo [V3]",
         "'El precio es senal de posicion social (sifrinaje)' [V4]",
         "La solucion no es bajar precios — es la primera experiencia."),
        ("'Sensibles al precio' [V3]",
         "'Aspiracionales con barrera identitaria' [V4]",
         "La activacion correcta es una razon especifica para venir, no precio bajo."),
        # DW-5 gap comunicacional wording aprobado
        ("Gap comunicacional: Tipo B certeza media [V3]",
         "[Hipotesis V4 — evidencia convergente]",
         "La hipotesis sigue siendo la mas plausible; confirmacion directa requiere "
         "modulo comunicacion en ola 2027."),
    ]
    row_h = 620000
    y = 1850000
    for (v3, v4, impl) in rows:
        for ci, (cx, cw, txt) in enumerate(zip(col_x, col_w, [v3, v4, impl])):
            clr = C_AMBER if "Hipotesis" in txt else C_DARK
            add_textbox(slide, cx + 40000, y + 40000, cw - 120000, row_h - 80000,
                        txt, Pt(8.5), clr, wrap=True)
        add_rect(slide, 300000, y + row_h - 20000, SLIDE_W - 600000, 20000,
                 fill_color=RGBColor(0xDD, 0xDD, 0xDD))
        y += row_h
    add_footer(slide)

def make_A6(prs):
    slide = add_slide(prs)
    set_bg(slide, C_WHITE)
    header_bar(slide, "Seccion A — Recap V3 y Encuadre V4",
               "A6 — Metodologia V4: como se anado la capa analitica")
    info_box(slide,
             "TAKEAWAY: V4 usa tres metodos complementarios al V3 sin reabrir el estudio base "
             "— cada metodo responde una pregunta que V3 no podia.",
             l=300000, t=1050000, w=SLIDE_W - 600000, h=320000)
    methods = [
        ("Wave-over-wave (CU-7)",
         "Pregunta: la posicion de Gama cambio 2025<>2026?\n"
         "57 items. Tests Newcombe-Wilson + correccion BH-FDR. n_2025=785, n_2026=402.\n"
         "Limitacion: ponderacion 2025 no disponible (D-CU7-001)."),
        ("Analisis cualitativo profundo (IN-7)",
         "Pregunta: por que operan los mecanismos del cuanti?\n"
         "Analisis tematico Braun & Clarke (2006) sobre 12 docs FG reales (42,094 chars).\n"
         "9 cat a priori + 9 emergentes. 52 unidades codificadas.\n"
         "Limitacion: unico analista sin doble codificacion (IN-7 §1.4)."),
        ("Triangulacion (IN-8)",
         "Pregunta: que cambia, que se confirma, que es nuevo?\n"
         "20 hallazgos triangulados. Cuanti-como-ancla segun ME-5 §4.1.\n"
         "0 contradicciones directas. 4 matices. 6 nuevos hallazgos V4."),
    ]
    col_w = (SLIDE_W - 800000) // 3
    for mi, (mname, mdesc) in enumerate(methods):
        cx = 300000 + mi * (col_w + 100000)
        add_rect(slide, cx, 1450000, col_w, SLIDE_H - 1950000, fill_color=C_LIGHT_BG)
        add_textbox(slide, cx + 80000, 1500000, col_w - 160000, 380000,
                    mname, Pt(10), C_RED, bold=True)
        add_textbox(slide, cx + 80000, 1900000, col_w - 160000, SLIDE_H - 2500000,
                    mdesc, Pt(8.5), C_DARK, wrap=True)

    # Caveats WOW + IN7
    caveat_box(slide,
               "Caveats metodologicos V4:\n"
               + C_LIT_CV_WOW_001 + "\n\n"
               + C_LIT_CV_WOW_002 + "\n\n"
               + C_LIT_IN7_14,
               l=300000, t=SLIDE_H - 620000,
               w=SLIDE_W - 600000, h=380000, size=Pt(7))
    add_footer(slide)

# ── SECCIÓN B (PLACEHOLDER) ────────────────────────────────────────────────────

def make_B_placeholder(prs):
    # Placeholder slide 1: section intro
    slide = add_slide(prs)
    set_bg(slide, C_RED)
    add_textbox(slide, 500000, 800000, SLIDE_W - 1000000, 700000,
                "SECCION B", Pt(48), C_WHITE, bold=True)
    add_textbox(slide, 500000, 1600000, SLIDE_W - 1000000, 500000,
                "Nucleo V3 Preservado — ~25 slides", Pt(28), C_WHITE)
    add_textbox(slide, 500000, 2300000, SLIDE_W - 1000000, 600000,
                "El deck V3 (2026-05-17, 46 slides) constituye el espinazo analitico de este estudio. "
                "Las slides de V3 se presentan en el archivo: 2026-05-17_Notoriedad-Gama-2026_V3.pptx",
                Pt(14), C_WHITE)

    # Placeholder slide 2: index of V3 blocks
    slide2 = add_slide(prs)
    set_bg(slide2, C_WHITE)
    header_bar(slide2, "Seccion B — Nucleo V3 Preservado (referencia)",
               "B — Contenido V3: Bloques 0-9 + Apendice (slides 1-46 del deck V3)")
    blocks = [
        "Slides 1-5:   Bloque 0 — Apertura (portada, advertencia metodologica, ficha tecnica, agenda)",
        "Slides 6-9:   Bloque 1 — Tesis central V3",
        "Slides 10-14: Bloque 2 — El activo experiencial (driver atencion OR=5.73)",
        "Slides 15-18: Bloque 3 — Las promociones como anzuelo",
        "Slides 19-23: Bloque 4 — Dos misiones (urgencia vs mercado grande)",
        "Slides 24-28: Bloque 5 — Los tres segmentos k-means",
        "Slides 29-32: Bloque 6 — La voz del consumidor (cuali V3; actualizable con IN-7 v2)",
        "Slides 33-34: Bloque 7 — Diagnostico comunicacion actual",
        "Slides 35-39: Bloque 8 — Recomendaciones priorizadas V3",
        "Slides 40-42: Bloque 9 — Roadmap 2026-2027",
        "Slides 43-46: Apendice — Metodologia ampliada + comparativa V2 vs V3",
    ]
    y = 1150000
    for bl in blocks:
        add_textbox(slide2, 400000, y, SLIDE_W - 800000, 360000,
                    bl, Pt(10), C_DARK)
        y += 380000
    caveat_box(slide2,
               "NOTA AL PRESENTER: insertar aqui las 46 slides del deck V3 "
               "(2026-05-17_Notoriedad-Gama-2026_V3.pptx) para la presentacion completa. "
               "El deck V4 agrega las secciones A, C, D, E como capas complementarias.",
               l=300000, t=SLIDE_H - 560000, w=SLIDE_W - 600000, h=320000)
    add_footer(slide2)

# ── SECCIÓN C ──────────────────────────────────────────────────────────────────

def make_C1(prs):
    slide = add_slide(prs)
    set_bg(slide, C_WHITE)
    header_bar(slide, "Seccion C — Capa Wave-over-Wave 2025<>2026",
               "C1 — 2025->2026: el mercado se mueve, Gama se mantiene")
    # CO-2 wording obligatorio
    info_box(slide, CO2_MAIN,
             l=300000, t=1050000, w=SLIDE_W - 600000, h=520000)
    bullets = [
        "57 comparaciones WoW ejecutadas (embudo de 9 marcas x 5 etapas + asociaciones). Correccion BH-FDR.",
        "10 items estadisticamente significativos (p_adj<0.05). Los 9 positivos son de competidores; "
        "el unico cambio negativo es Central Madeirense.",
        "Gama: 0/8 items de embudo significativos. Ninguna variacion en awareness, consideracion, "
        "compra ni preferencia. Esto es evidencia positiva en un mercado dinamico.",
    ]
    y = 1650000
    for b in bullets:
        add_textbox(slide, 350000, y, SLIDE_W - 700000, 400000,
                    "  * " + b, Pt(10), C_DARK, wrap=True)
        y += 420000
    label_box(slide, "[WoW Confirmado — n_2025=785, n_2026=402, BH-FDR]",
              l=300000, t=y + 80000, w=3600000, fg=C_GREEN)
    caveat_box(slide,
               "Caveats: " + C_LIT_CV_WOW_001 + "  |  " + C_LIT_CV_WOW_002,
               l=300000, t=SLIDE_H - 560000, w=SLIDE_W - 600000, h=320000, size=Pt(7.5))
    add_footer(slide)

def make_C2(prs):
    slide = add_slide(prs)
    set_bg(slide, C_WHITE)
    header_bar(slide, "Seccion C — Capa Wave-over-Wave 2025<>2026",
               "C2 — Gama estable en los ocho eslabones del embudo — evidencia positiva en mercado turbulento")
    info_box(slide,
             "TAKEAWAY: Ninguna metrica de Gama — desde TOM hasta compra habitual — vario significativamente "
             "entre 2025 y 2026, con alta certeza estadistica.",
             l=300000, t=1050000, w=SLIDE_W - 600000, h=320000)
    # Embudo table
    items = [
        ("TOM",          "+2.2pp", "p_adj>0.10"),
        ("Asistida",     "+2.2pp", "p_adj>0.10"),
        ("Consideracion","+4.3pp", "p_adj>0.10"),
        ("Compra 3m",    "-0.2pp", "p_adj>0.10"),
        ("Preferida",    "-1.7pp", "p_adj>0.10"),
        ("Ultima compra","-1.2pp", "p_adj>0.10"),
        ("Habitual",     "+0.8pp", "p_adj>0.10"),
        ("Misiones",     "+4.9pp", "p_adj>0.10"),
    ]
    col_titles = ["Metrica", "Delta WoW", "Significancia"]
    col_x = [300000, 3500000, 6500000]
    col_w = [3100000, 2900000, 5300000]
    y = 1450000
    for ci, (cx, cw, ch) in enumerate(zip(col_x, col_w, col_titles)):
        add_rect(slide, cx, y, cw - 60000, 320000, fill_color=C_RED)
        add_textbox(slide, cx + 60000, y + 60000, cw - 120000, 220000,
                    ch, Pt(9), C_WHITE, bold=True)
    y += 360000
    for (metric, delta, sig) in items:
        add_textbox(slide, 340000, y, 3050000, 340000, metric, Pt(9.5), C_DARK)
        clr = C_GREEN if delta.startswith("+") else C_DARK
        add_textbox(slide, 3540000, y, 2840000, 340000, delta, Pt(9.5), clr, bold=True)
        add_textbox(slide, 6540000, y, 5200000, 340000,
                    sig + " — sin cambio significativo", Pt(9), C_GRAY)
        y += 360000
    label_box(slide, "[WoW Confirmado — H-1: ninguna variacion sig p_adj>0.10]",
              l=300000, t=y + 80000, w=4200000, fg=C_GREEN)
    label_box(slide, "[Hipotesis V4 — estabilidad como alerta estrategica pendiente ola 2027]",
              l=4700000, t=y + 80000, w=4800000, fg=C_AMBER)
    caveat_box(slide,
               C_LIT_CV_WOW_001 + "  |  " + C_LIT_CV_WOW_002,
               l=300000, t=SLIDE_H - 500000, w=SLIDE_W - 600000, h=280000, size=Pt(7))
    add_footer(slide)

def make_C3(prs):
    """CO-1: Rio bullet de proteinas EXCLUIDO. Solo datos cuantitativos sig_99."""
    slide = add_slide(prs)
    set_bg(slide, C_WHITE)
    header_bar(slide, "Seccion C — Capa Wave-over-Wave 2025<>2026",
               "C3 — Rio: crecimiento explosivo en awareness y conversion, especialmente en C+/C")
    info_box(slide,
             "TAKEAWAY: Rio es el mayor ganador competitivo 2025->2026: TOM +17pp, Consideracion +19.6pp, "
             "Compra +12.5pp, con efectos especialmente marcados en el segmento C+/C de Gama.",
             l=300000, t=1050000, w=SLIDE_W - 600000, h=320000)
    bullets = [
        "TOM: 28.0% -> 45.0% (+17.0pp, p_adj<0.0001, sig_99)",
        "Consideracion: 19.2% -> 38.8% (+19.6pp, p_adj<0.0001, sig_99)",
        "Compra 3m: 14.4% -> 26.9% (+12.5pp, p_adj<0.0001, sig_99)",
        "Preferida: 6.2% -> 10.2% (+4.0pp, tendencia p_adj=0.059)",
        "En NSE C+/C (el segmento natural de Gama): TOM Rio 30.9% -> 56.7% (+25.8pp, p_adj=0.0002). "
        "El efecto en el segmento premium es aun mas marcado.",
        # CO-1 applied: Rio-proteinas bullet ELIMINATED per Bruna recommendation accepted by main Claude
    ]
    y = 1450000
    for b in bullets:
        add_textbox(slide, 350000, y, SLIDE_W - 700000, 400000,
                    "  * " + b, Pt(10), C_DARK, wrap=True)
        y += 420000
    label_box(slide, "[WoW Confirmado — H-2: TOM, Consideracion, Compra sig_99]",
              l=300000, t=y + 60000, w=4200000, fg=C_GREEN)
    caveat_box(slide,
               C_LIT_CV_WOW_002 + "  |  " + C_LIT_CV_WOW_005,
               l=300000, t=SLIDE_H - 500000, w=SLIDE_W - 600000, h=280000, size=Pt(7))
    add_footer(slide)

def make_C4(prs):
    slide = add_slide(prs)
    set_bg(slide, C_WHITE)
    header_bar(slide, "Seccion C — Capa Wave-over-Wave 2025<>2026",
               "C4 — Paramo consolida liderazgo en toda la cadena del embudo")
    info_box(slide,
             "TAKEAWAY: Paramo crece en los tres eslabones superiores del embudo con alta certeza, "
             "consolidando su posicion como la opcion aspiracional para el mercado grande.",
             l=300000, t=1050000, w=SLIDE_W - 600000, h=320000)
    bullets = [
        "TOM: 27.0% -> 39.1% (+12.1pp, p_adj=0.0003, sig_99)",
        "Asistida: 40.9% -> 51.7% (+10.8pp, p_adj=0.003, sig_99)",
        "Consideracion: 31.7% -> 49.0% (+17.3pp, p_adj<0.0001, sig_99)",
        "En NSE C+/C: TOM Paramo 18.0% -> 40.4% (+22.3pp, p_adj=0.0002, sig_99). "
        "Paramo crece en el mismo segmento natural de Gama.",
        "El corpus FG confirma a Paramo como destino del mercado grande: multiples participantes "
        "de todos los segmentos mencionan Paramo para la compra semanal/quincenal de alto volumen.",
        "Lectura estrategica: el espacio competitivo en C+/C se esta reduciendo. "
        "Gama tiene su posicion hoy; la velocidad de Rio y Paramo en ese segmento es una senal de presion futura.",
    ]
    y = 1450000
    for b in bullets:
        add_textbox(slide, 350000, y, SLIDE_W - 700000, 380000,
                    "  * " + b, Pt(9.5), C_DARK, wrap=True)
        y += 400000
    label_box(slide, "[WoW Confirmado — H-3: TOM, Asistida, Consideracion sig_99]",
              l=300000, t=y + 40000, w=4200000, fg=C_GREEN)
    caveat_box(slide,
               C_LIT_CV_WOW_002 + "  |  " + C_LIT_CV_WOW_005,
               l=300000, t=SLIDE_H - 500000, w=SLIDE_W - 600000, h=280000, size=Pt(7))
    add_footer(slide)

def make_C5(prs):
    slide = add_slide(prs)
    set_bg(slide, C_WHITE)
    header_bar(slide, "Seccion C — Capa Wave-over-Wave 2025<>2026",
               "C5 — Central Madeirense pierde compra efectiva: la oportunidad no capturada de Gama")
    info_box(slide,
             "TAKEAWAY: CM es la unica cadena con perdida significativa de conversion (-7.7pp compra "
             "efectiva) — el espacio disponible que Gama aun no ha capturado.",
             l=300000, t=1050000, w=SLIDE_W - 600000, h=320000)
    bullets = [
        "Compra 3m CM: 30.3% -> 22.6% (-7.7pp, p_adj=0.033, sig_95). Unica perdida significativa en el panel.",
        "Preferida CM: 16.7% -> 11.2% (-5.5pp, tendencia p_adj=0.053). La perdida de compra anticipa perdida de preferencia.",
        "Gama mantiene su posicion mientras CM pierde — la aritmetica del mercado es favorable, "
        "pero Gama no ha capturado el espacio que CM libera.",
        "El corpus FG provee el mecanismo: el shopper que deja CM va a Paramo para el mercado grande "
        "o a Rio para proteinas — no a Gama, porque Gama compite en la mision de urgencia.",
        # C5 approved wording (BR-2 §4 C5 adjustment)
        "Para capturar el espacio de CM, Gama puede ampliar las misiones en que es elegida "
        "— la propuesta actual ancla en urgencia, no en mercado grande.",
    ]
    y = 1450000
    for b in bullets:
        add_textbox(slide, 350000, y, SLIDE_W - 700000, 380000,
                    "  * " + b, Pt(9.5), C_DARK, wrap=True)
        y += 400000
    label_box(slide, "[WoW Confirmado — H-4: CM -7.7pp compra, sig_95]",
              l=300000, t=y + 40000, w=3800000, fg=C_GREEN)
    add_textbox(slide, 4200000, y + 40000, 3800000, 240000,
                "Nivel sig_95 — certeza MEDIA. Causalidad no inferible.", Pt(8), C_AMBER)
    add_footer(slide)

def make_C6(prs):
    slide = add_slide(prs)
    set_bg(slide, C_WHITE)
    header_bar(slide, "Seccion C — Capa Wave-over-Wave 2025<>2026",
               "C6 — La imagen experiencial de Gama gana terreno — tendencia, no certeza")
    info_box(slide,
             "TAKEAWAY: Los datos WoW sugieren — sin confirmar — que Gama se aleja de la imagen "
             "de precio y se acerca a la imagen de experiencia.",
             l=300000, t=1050000, w=SLIDE_W - 600000, h=320000)
    add_textbox(slide, 300000, 1450000, SLIDE_W - 600000, 340000,
                "Tendencias positivas en atributos de experiencia:", Pt(10), C_GREEN, bold=True)
    pos = [
        ("'Tienda atractiva'", "+6.3pp", "p_adj=0.058 — tendencia"),
        ("'Seguro'",           "+5.8pp", "p_adj=0.089 — tendencia"),
    ]
    y = 1840000
    for (attr, delta, sig) in pos:
        add_textbox(slide, 400000, y, 3000000, 320000, attr, Pt(9.5), C_DARK)
        add_textbox(slide, 3500000, y, 1200000, 320000, delta, Pt(9.5), C_GREEN, bold=True)
        add_textbox(slide, 4800000, y, 3000000, 320000, sig, Pt(9), C_AMBER)
        y += 340000
    add_textbox(slide, 300000, y + 80000, SLIDE_W - 600000, 340000,
                "Tendencias negativas en atributos economicos:", Pt(10), C_WARN_FG, bold=True)
    y += 480000
    negs = [
        ("'Promociones'",   "-5.1pp", "p_adj=0.053 — tendencia"),
        ("'Menor precio'",  "-3.9pp", "p_adj=0.091 — tendencia"),
    ]
    for (attr, delta, sig) in negs:
        add_textbox(slide, 400000, y, 3000000, 320000, attr, Pt(9.5), C_DARK)
        add_textbox(slide, 3500000, y, 1200000, 320000, delta, Pt(9.5), C_WARN_FG, bold=True)
        add_textbox(slide, 4800000, y, 3000000, 320000, sig, Pt(9), C_AMBER)
        y += 340000
    caveat_box(slide,
               "ADVERTENCIA VISIBLE: NINGUNO de estos items pasa la correccion BH-FDR (umbral q<0.05). "
               "Son informativos, no confirmados. Presentar siempre con esta advertencia explícita.",
               l=300000, t=y + 80000, w=SLIDE_W - 600000, h=280000)
    label_box(slide, "[Hipotesis V4 — H-5 y H-8: solo tendencia, p_adj 0.053-0.091, no pasa BH-FDR]",
              l=300000, t=SLIDE_H - 560000, w=5400000, fg=C_AMBER)
    add_footer(slide)

def make_C7(prs):
    slide = add_slide(prs)
    set_bg(slide, C_WHITE)
    header_bar(slide, "Seccion C — Capa Wave-over-Wave 2025<>2026",
               "C7 — La imagen de precio bajo de Gama se debilita: preparacion para la barrera identitaria")
    info_box(slide,
             "TAKEAWAY: La erosion tendencial de 'menor precio' con Gama es coherente con el hallazgo cuali "
             "del 'sifrinaje' — el precio de Gama opera como senal de posicion, y esa senal se aleja de lo economico.",
             l=300000, t=1050000, w=SLIDE_W - 600000, h=320000)
    bullets_top = [
        "'Menor precio' Gama: 11.1% -> 7.2% (-3.9pp, p_adj=0.091, tendencia). No confirmado estadisticamente.",
        "'Promociones' Gama: 14.0% -> 9.0% (-5.1pp, p_adj=0.053, tendencia). No confirmado estadisticamente.",
        "El corpus cualitativo V4 provee el mecanismo: la percepcion de precio de Gama no es un calculo "
        "comparativo — es una categorizacion de posicion social.",
    ]
    y = 1450000
    for b in bullets_top:
        add_textbox(slide, 350000, y, SLIDE_W - 700000, 380000,
                    "  * " + b, Pt(9.5), C_DARK, wrap=True)
        y += 400000
    # Verbatim David — DW-1-bis applies (same caveat as DW1)
    quote_box(slide,
              "una nina muy cool del este, muy catira, muy cuchi... una veinteañera de esas cool. "
              "Sifrina, de Las Mercedes.",
              "David, Frecuentes 31-50+",
              l=300000, t=y + 60000, w=SLIDE_W * 2 // 3, h=500000)
    # C-LIT-DW1-bis obligatorio
    caveat_box(slide, C_LIT_DW1_BIS,
               l=300000, t=y + 620000, w=SLIDE_W - 600000, h=280000, size=Pt(7.5))
    label_box(slide, "[Hipotesis V4 — H-8: tendencia p_adj=0.091, no pasa BH-FDR] "
              "+ [Triangulado: cuali V4 — mecanismo identitario]",
              l=300000, t=SLIDE_H - 540000, w=6000000, fg=C_AMBER)
    add_footer(slide)

def make_C8(prs):
    slide = add_slide(prs)
    set_bg(slide, C_WHITE)
    header_bar(slide, "Seccion C — Capa Wave-over-Wave 2025<>2026",
               "C8 — Sintesis WoW: el mercado se recompone, Gama necesita estrategia ofensiva y defensiva")
    # CO-2 synthesis wording obligatorio
    info_box(slide, CO2_SYNTHESIS,
             l=300000, t=1050000, w=SLIDE_W - 600000, h=360000)
    bullets = [
        "Ganadores 2025->2026: Rio (TOM +17pp, Consideracion +19.6pp, sig_99) + Paramo "
        "(TOM +12.1pp, Consideracion +17.3pp, sig_99). Crecimiento de awareness masiva y penetracion de compra.",
        "Perdedor: Central Madeirense (-7.7pp compra, sig_95). El espacio disponible existe.",
        "Las tres implicaciones competitivas para Gama: (1) defender el activo relacional que Rio y Paramo "
        "no tienen (24h + acompanamiento); (2) activar el Gama Club como mecanismo de conversion del espacio "
        "disponible de CM; (3) disenar una estrategia ofensiva en el segmento C+/C antes de que Rio y Paramo "
        "completen su penetracion.",
    ]
    y = 1500000
    for b in bullets:
        add_textbox(slide, 350000, y, SLIDE_W - 700000, 380000,
                    "  * " + b, Pt(9.5), C_DARK, wrap=True)
        y += 420000
    add_textbox(slide, 300000, y + 40000, SLIDE_W - 600000, 280000,
                "[Visual: Mapa competidores 2x2 — eje X=cambio consideracion WoW, eje Y=cambio compra WoW. "
                "Rio y Paramo cuadrante sup-der; CM inf-izq; Gama centro sin movimiento.]",
                Pt(7.5), C_GRAY, italic=True)
    label_box(slide, "[WoW Confirmado para Rio H-2, Paramo H-3, CM H-4] + [Hipotesis V4 para implicaciones estrategicas]",
              l=300000, t=SLIDE_H - 560000, w=7200000, fg=C_GREEN)
    add_footer(slide)

# ── SECCIÓN D ──────────────────────────────────────────────────────────────────

def make_D1(prs):
    slide = add_slide(prs)
    set_bg(slide, C_WHITE)
    header_bar(slide, "Seccion D — Capa Cualitativa + Triangulacion",
               "D1 — Los focus groups revelan el mecanismo detras de los numeros")
    info_box(slide,
             "TAKEAWAY: El analisis de 12 documentos FG (42,094 chars, 6 segmentos) provee el 'por que' "
             "de los drivers cuantitativos y documenta barreras que el cuestionario no podia capturar.",
             l=300000, t=1050000, w=SLIDE_W - 600000, h=320000)
    bullets = [
        "Corpus analizado: 8 sesiones presenciales + 4 documentos de verbatims online. "
        "Total: 6 segmentos, ~42,094 chars, 52 unidades codificadas.",
        "Metodologia: analisis tematico Braun & Clarke (2006), codificacion deductiva-inductiva en dos pasadas. "
        "9 categorias a priori + 9 categorias emergentes (inductivas).",
        "6 temas principales: (Q1) Drivers de experiencia fisica; (Q2) Precio como identidad; (Q3) 24h y conveniencia; "
        "(Q4) Gama Club; (Q5) Canal digital; (Q6) Personificacion de marca.",
        "Los hallazgos cualitativos no prueban prevalencia (eso es el cuanti); prueban mecanismo, "
        "barreras y territorio de marca. Las dos fuentes se complementan, no compiten.",
    ]
    y = 1450000
    for b in bullets:
        add_textbox(slide, 350000, y, SLIDE_W - 700000, 400000,
                    "  * " + b, Pt(9.5), C_DARK, wrap=True)
        y += 420000
    # C-LIT-IN7-1.4 obligatorio
    caveat_box(slide, C_LIT_IN7_14,
               l=300000, t=SLIDE_H - 560000, w=SLIDE_W - 600000, h=320000, size=Pt(7.5))
    add_footer(slide)

def make_D2(prs):
    slide = add_slide(prs)
    set_bg(slide, C_WHITE)
    header_bar(slide, "Seccion D — Capa Cualitativa + Triangulacion",
               "D2 — La atencion de Gama genera acompanamiento, no solo servicio — el mecanismo reformulado")
    info_box(slide,
             "TAKEAWAY: El corpus FG documenta en 6 de 7 segmentos la misma imagen espontanea: Gama es una "
             "'guia experta de confianza' que acompana la compra — exactamente el territorio relacional que el OR=5.73 mide.",
             l=300000, t=1050000, w=SLIDE_W - 600000, h=320000)
    # Split layout
    # Left: cuanti
    add_rect(slide, 300000, 1450000, 5400000, SLIDE_H - 1950000, fill_color=C_LIGHT_BG)
    add_textbox(slide, 380000, 1520000, 5200000, 380000,
                "Cuanti V3 — el numero", Pt(11), C_RED, bold=True)
    add_textbox(slide, 380000, 1950000, 5200000, 680000,
                "OR = 5.73*** (IC95: 1.6-20.4)\nAtención driver #1\nTriple convergencia logit + RF + SHAP",
                Pt(28), C_DARKRED, bold=True)
    add_textbox(slide, 380000, 2700000, 5200000, 500000,
                "Quien experimenta ser guiado tiene 5.7 veces mas odds de preferir Gama.",
                Pt(10), C_DARK)
    # Right: verbatims
    add_rect(slide, 5900000, 1450000, SLIDE_W - 6200000, SLIDE_H - 1950000, fill_color=C_LIGHT_BG)
    add_textbox(slide, 5980000, 1520000, SLIDE_W - 6280000, 380000,
                "Cuali V4 — el mecanismo", Pt(11), C_RED, bold=True)
    quote_box(slide,
              "Una mujer joven, pero con experiencia. Que sea como una guia, que te diga: "
              "'Mira, hoy llego el pan arabe fresco, llévatelo'.",
              "Mary Francis Bossia, Ocasionales 18-30",
              l=5900000, t=1950000, w=SLIDE_W - 6200000, h=500000)
    quote_box(slide,
              "Alguien de Caracas que te oriente y sea como una guia de compras, "
              "que no se quede atras con la tecnologia.",
              "Moises Torrealba, Ocasionales 18-30",
              l=5900000, t=2520000, w=SLIDE_W - 6200000, h=500000)
    add_textbox(slide, 5980000, 3080000, SLIDE_W - 6280000, 500000,
                "Wording comunicacional correcto: 'en Gama siempre hay alguien que te ayuda'.\n"
                "El claim correcto V4: 'La atencion de Gama genera la experiencia de ser "
                "acompanado y guiado — quien experimenta ser guiado tiene 5.7x mas odds de preferir Gama.'",
                Pt(8.5), C_DARK, wrap=True)
    label_box(slide, "[Triangulado: cuanti V3 OR=5.73 + cuali V4 — mecanismo reformulado. Tipo B.]",
              l=300000, t=SLIDE_H - 540000, w=5400000, fg=C_BLUE)
    add_footer(slide)

def make_D3(prs):
    """DW-1: Version [A] con nombre Azahara Betancourt + C-LIT-DW1 obligatorio."""
    slide = add_slide(prs)
    set_bg(slide, C_WHITE)
    header_bar(slide, "Seccion D — Capa Cualitativa + Triangulacion",
               "D3 — 'Sifrinaje': la barrera de precio que los descuentos no mueven")
    info_box(slide,
             "TAKEAWAY: La barrera de precio de Gama no es un calculo economico — es una senal "
             "de posicion social que el shopper usa para autoexcluirse antes de experimentar el diferencial.",
             l=300000, t=1050000, w=SLIDE_W - 600000, h=320000)
    add_textbox(slide, 350000, 1450000, SLIDE_W - 700000, 320000,
                "Lo que el cuanti media: 54% percibe Gama como cara (V3). Precio no driver de preferencia "
                "(OR=1.03, SHAP #10). Paradoja documentada — V4 revela el mecanismo.",
                Pt(9.5), C_DARK, wrap=True)
    # Verbatim Azahara — DW-1 version [A] con nombre
    quote_box(slide,
              "A mi me parece que es un poco costoso... el sifrinaje y tal, pero para hacer un mercado "
              "grande... es mejor ir al Lux o al Paramo.",
              "Azahara Betancourt, Ocasionales 18-30",
              l=300000, t=1850000, w=SLIDE_W * 2 // 3, h=520000)
    # C-LIT-DW1 obligatorio inmediatamente despues del verbatim (BR-2 §5)
    caveat_box(slide, C_LIT_DW1,
               l=300000, t=2440000, w=SLIDE_W - 600000, h=300000, size=Pt(8))
    bullets = [
        "El mecanismo: el shopper se autoexcluye de Gama no porque haya comparado precios "
        "sino porque percibe que 'Gama es para otro tipo de persona'. La barrera es simbolica.",
        "La implicacion directa: los descuentos de precio no mueven esta barrera. "
        "La barrera se neutraliza con una primera experiencia que haga creible que el diferencial vale.",
        "Convergencia WoW: debilitamiento de 'menor precio' (-3.9pp) y 'promociones' (-5.1pp) — tendencia.",
    ]
    y = 2800000
    for b in bullets:
        add_textbox(slide, 350000, y, SLIDE_W - 700000, 360000,
                    "  * " + b, Pt(9), C_DARK, wrap=True)
        y += 380000
    label_box(slide, "[Triangulado: cuanti V3 + cuali V4. Tipo B.]",
              l=300000, t=SLIDE_H - 540000, w=3400000, fg=C_BLUE)
    add_footer(slide)

def make_D4(prs):
    """DW-2: lenguaje condicional + C-LIT-DW2-a obligatorio."""
    slide = add_slide(prs)
    set_bg(slide, C_WHITE)
    header_bar(slide, "Seccion D — Capa Cualitativa + Triangulacion",
               "D4 — La marca tiene un arquetipo ya construido — solo falta apropiarlo")
    info_box(slide,
             "TAKEAWAY: Sin ninguna induccion del moderador, 6 de 7 grupos FG convergen en el mismo "
             "arquetipo: Gama como mujer de 35-45 anos, caraqueña, profesional, 'guia experta de confianza' "
             "— territorio disponible sin ocupacion por ningun competidor.",
             l=300000, t=1050000, w=SLIDE_W - 600000, h=360000)
    bullets = [
        "La pregunta en los FG era sobre personificacion de la marca. Sin mascota, sin campana de personificacion activa.",
        "Respuestas convergentes en 6/7 segmentos: 'Yo me la imagino como una mujer, como de unos "
        "35 anos, super activa, lider, entusiasta, que sabe lo que hace.' (Moises Torrealba, Ocasionales 18-30)",
        "La mayoria coincide en 'una mujer entre 35 y 45 anos, caraqueña, profesional, activa, lider y amable.' (Compendio, transversal)",
        "Coherencia con los drivers cuantitativos: el arquetipo 'guia experta' es la personificacion exacta "
        "de lo que OR=5.73 mide (atencion como acompanamiento). Son el mismo activo desde dos angulos.",
        "Ningún competidor de Gama (Rio, Paramo, CM) ha apropiado este arquetipo. La ventana esta disponible.",
        # DW-2: lenguaje condicional — "se recomienda explorar", no "debe materializar"
        "El equipo analitico se recomienda explorar la materializacion de este imaginario en "
        "creatividad concreta — previa verificacion de restricciones de marca.",
    ]
    y = 1500000
    for b in bullets:
        add_textbox(slide, 350000, y, SLIDE_W - 700000, 380000,
                    "  * " + b, Pt(9.5), C_DARK, wrap=True)
        y += 400000
    # C-LIT-DW2-a obligatorio
    caveat_box(slide, C_LIT_DW2_A,
               l=300000, t=SLIDE_H - 580000, w=SLIDE_W - 600000, h=340000, size=Pt(8))
    label_box(slide, "[Cuali V4 FG — emergente no inducido. 6/7 segmentos. Tipo B.]",
              l=300000, t=SLIDE_H - 620000 - 260000, w=4200000, fg=C_BLUE)
    add_footer(slide)

def make_D5(prs):
    """DW-4 + CO-3 Opción C: hallazgos digitales con nota scope. Atribución por segmento."""
    slide = add_slide(prs)
    set_bg(slide, C_WHITE)
    header_bar(slide, "Seccion D — Capa Cualitativa + Triangulacion",
               "D5 — Gama Club y canal digital: los activos subexplotados mas reparables")
    info_box(slide,
             "TAKEAWAY: El Gama Club tiene interes genuino y transversal en todos los segmentos, "
             "pero su saldo es invisible digitalmente — la unica barrera es operacional y resoluble.",
             l=300000, t=1050000, w=SLIDE_W - 600000, h=320000)
    # Block 1: Gama Club
    add_rect(slide, 300000, 1450000, SLIDE_W - 600000, 800000, fill_color=C_LIGHT_BG)
    add_textbox(slide, 360000, 1490000, SLIDE_W - 720000, 280000,
                "Gama Club — El activo invisible", Pt(10), C_RED, bold=True)
    quote_box(slide,
              "Las millas... yo ni sabia que eso se podia usar para pagar hasta hace poco. "
              "Yo pensaba que era como los puntos del banco que uno nunca usa.",
              "Gregory, Frecuentes 18-30",
              l=360000, t=1790000, w=SLIDE_W - 720000, h=420000)

    # Block 2: inventario digital — DW-4 completo
    add_rect(slide, 300000, 2280000, SLIDE_W - 600000, 1100000, fill_color=C_LIGHT_BG)
    add_textbox(slide, 360000, 2320000, SLIDE_W - 720000, 280000,
                "Canal digital — Barrera 1: Inconsistencia de inventario", Pt(10), C_RED, bold=True)
    # C-LIT-DW4-a ANTES del verbatim (obligatorio)
    info_box(slide, C_LIT_DW4_A,
             l=360000, t=2640000, w=SLIDE_W - 720000, h=280000, size=Pt(8))
    # Verbatim con atribucion C-LIT-DW4-c (no nombre propio — por segmento/edad)
    quote_box(slide,
              "Yo deje de usar la pagina porque a veces pides algo que dice que hay y a la media "
              "hora te llaman para decirte que no tienen.",
              C_LIT_DW4_C,   # atribucion aprobada
              l=360000, t=2960000, w=SLIDE_W - 720000, h=420000)
    # C-LIT-DW4-b DESPUES del verbatim (obligatorio)
    caveat_box(slide, C_LIT_DW4_B,
               l=300000, t=3440000, w=SLIDE_W - 600000, h=260000, size=Pt(8))

    # Block 3: Frescos + CO-3 scope note
    add_textbox(slide, 300000, 3760000, SLIDE_W - 600000, 300000,
                "Canal digital — Barrera 2: Desconfianza en frescos (5/7 segmentos FG). "
                "Esta barrera no se mueve con promesas de calidad — requiere segmentacion de la propuesta.",
                Pt(9), C_DARK, wrap=True)

    # CO-3 Opcion C — nota scope visible (obligatoria per brief)
    caveat_box(slide,
               "NOTA DE ALCANCE (CO-3 — Opcion C hibrida): Los hallazgos de canal digital y Gama Club "
               "son ampliacion del brief original de notoriedad/brand health. Se incluyen por su alta "
               "accionabilidad. Las recomendaciones de implementacion detalladas sugieren una conversion "
               "de scope con Gama. El Owner puede ajustar este wording pre-envio a Cora.",
               l=300000, t=SLIDE_H - 520000, w=SLIDE_W - 600000, h=300000)
    label_box(slide, "[Cuali V4 FG + digital. Tipo B. Sin contraparte directa en cuanti V3 o CU-7.]",
              l=300000, t=SLIDE_H - 570000 - 260000, w=5400000, fg=C_BLUE)
    add_footer(slide)

def make_D6(prs):
    """DW-3: wording aprobado C-LIT-DW3-a. Lenguaje prohibido no aparece."""
    slide = add_slide(prs)
    set_bg(slide, C_WHITE)
    header_bar(slide, "Seccion D — Capa Cualitativa + Triangulacion",
               "D6 — K-means revisitados: los Pragmaticos no son solo sensibles al precio — son aspiracionales con barrera identitaria")
    info_box(slide,
             "TAKEAWAY: El corpus FG confirma los tres segmentos k-means y reforma el mecanismo "
             "del Segmento 2 (Pragmaticos Convertibles).",
             l=300000, t=1050000, w=SLIDE_W - 600000, h=280000)
    segs = [
        ("Segmento 3 — Nucleo Leal (8%)", "CONFIRMADO",
         "Reconocible en frecuentes de 18-30 y 50+ con lenguaje afectivo: "
         "'Voy al de 24 horas tipo 10 u 11 de la noche. Es mi momento de relax, literal.' (Gaitskel, Frecuentes 18-30). "
         "Alta frecuencia, uso fuera de horario habitual, experiencia emocional positiva.", C_GREEN),
        ("Segmento 1 — Mayoria Exigente (59%)", "CONFIRMADO",
         "El perfil mas reconocible en los FG. Los dos sub-perfiles (exigente sin preferencia exclusiva Gama "
         "+ tension de misiones) son visibles en el corpus.", C_GREEN),
        ("Segmento 2 — Pragmaticos Convertibles (33%)", "CONFIRMADO EN EXISTENCIA, MECANISMO REFORMULADO",
         "V3: 'sensibles al precio'. V4 cuali: el seg_2 no es simplemente precio-sensitivo. Es un shopper "
         "para quien Gama representa un territorio aspiracional con connotacion de clase. "
         "Verbatim: 'A mi me parece que es un poco costoso... el sifrinaje y tal...' "
         "(Azahara Betancourt, Ocasionales 18-30 — reconocible como seg_2).", C_AMBER),
    ]
    y = 1400000
    for (title, status, desc, clr) in segs:
        add_rect(slide, 300000, y, SLIDE_W - 600000, 700000, fill_color=C_LIGHT_BG)
        add_textbox(slide, 360000, y + 40000, 3800000, 280000, title, Pt(10), C_RED, bold=True)
        add_textbox(slide, 4300000, y + 40000, SLIDE_W - 4700000, 280000,
                    status, Pt(9), clr, bold=True)
        add_textbox(slide, 360000, y + 340000, SLIDE_W - 720000, 340000,
                    desc, Pt(8.5), C_DARK, wrap=True)
        y += 760000

    # C-LIT-DW3-a wording obligatorio (sustituye cualquier formulacion prohibida)
    info_box(slide, C_LIT_DW3_A,
             l=300000, t=y + 80000, w=SLIDE_W - 600000, h=340000, size=Pt(9))
    label_box(slide, "[Triangulado: cuanti V3 k-means + cuali V4. Posicion II — mecanismo reformulado. Tipo B parcial.]",
              l=300000, t=SLIDE_H - 560000, w=6600000, fg=C_BLUE)
    add_footer(slide)

# ── SECCIÓN E ──────────────────────────────────────────────────────────────────

def make_E1(prs):
    slide = add_slide(prs)
    set_bg(slide, C_WHITE)
    header_bar(slide, "Seccion E — Sintesis V4 + Recomendaciones Actualizadas",
               "E1 — Los 6 hallazgos V4 nuevos en una sola vista")
    info_box(slide,
             "TAKEAWAY: V4 agrega seis hallazgos que ninguna capa analitica anterior podia generar "
             "— todos accionables, todos con trazabilidad de evidencia explicita.",
             l=300000, t=1050000, w=SLIDE_W - 600000, h=280000)
    news = [
        ("V4-NEW-1", "[Triangulado]", C_BLUE,
         "El mecanismo del driver de atencion es acompanamiento-guia (no reconocimiento nominal). El claim mas comunicable."),
        ("V4-NEW-2", "[Triangulado]", C_BLUE,
         "La barrera de precio es identitaria-simbolica ('sifrinaje'). Promociones de precio bajo no la mueven."),
        ("V4-NEW-3", "[Cuali V4]",    C_BLUE,
         "Gama Club = activo subexplotado. La unica barrera es la opacidad digital del saldo — resoluble con bajo costo."),
        ("V4-NEW-4", "[Cuali V4 FG emergente]", C_BLUE,
         "Personificacion femenina espontanea ('guia experta 35-45') en 6/7 grupos FG. Territorio disponible sin competidor."),
        ("V4-NEW-5", "[WoW Confirmado Tipo A]",  C_GREEN,
         "Rio (+17pp TOM sig_99) y Paramo (+12pp TOM sig_99) crecen explosivamente en el segmento natural de Gama."),
        ("V4-NEW-6", "[Cuali V4 digital]", C_BLUE,
         "Canal digital con dos barreras distintas: inconsistencia de inventario (solucion tecnologica) + desconfianza en frescos (solucion de propuesta)."),
    ]
    row_h = 580000
    y = 1430000
    for (tag, label, lclr, text) in news:
        add_rect(slide, 300000, y, 1200000, row_h - 60000, fill_color=C_RED)
        add_textbox(slide, 320000, y + 80000, 1160000, 380000, tag, Pt(9), C_WHITE, bold=True)
        add_rect(slide, 1580000, y, 2400000, row_h - 60000, fill_color=C_LIGHT_BG)
        add_textbox(slide, 1620000, y + 60000, 2320000, row_h - 120000,
                    label, Pt(8), lclr, bold=True)
        add_textbox(slide, 4060000, y + 60000, SLIDE_W - 4360000, row_h - 120000,
                    text, Pt(9.5), C_DARK, wrap=True)
        y += row_h
    add_footer(slide)

def make_E2(prs):
    """DW-5 aplicado (etiqueta visible sin explicar cambio) + CO-2 en Argumento 3."""
    slide = add_slide(prs)
    set_bg(slide, C_WHITE)
    header_bar(slide, "Seccion E — Sintesis V4 + Recomendaciones Actualizadas",
               "E2 — Piramide Minto actualizada V4: los argumentos con su nueva evidencia")
    info_box(slide,
             "TAKEAWAY: Los cuatro argumentos MECE de V3 se mantienen intactos — V4 los enriquece "
             "con dos capas de evidencia que refuerzan y afinan cada uno.",
             l=300000, t=1050000, w=SLIDE_W - 600000, h=280000)
    # Apex
    add_rect(slide, SLIDE_W // 2 - 2500000, 1430000, 5000000, 480000, fill_color=C_RED)
    add_textbox(slide, SLIDE_W // 2 - 2480000, 1480000, 4960000, 380000,
                "TESIS V4: Gama tiene el activo relacional mas diferenciador del mercado y su mecanismo "
                "es el acompanamiento-guia — la comunicacion correcta materializa el arquetipo que el "
                "shopper ya proyecta espontaneamente.",
                Pt(9), C_WHITE, bold=True, wrap=True)

    args = [
        ("QUE — activo a comunicar",
         "V3: Atencion driver #1 OR=5.73. Limpieza driver secundario.\n"
         "V4: El mecanismo es acompanamiento-guia (6/7 FG), no reconocimiento nominal. "
         "El arquetipo femenino (D4) es la materializacion.",
         "[Triangulado]", C_BLUE),
        ("COMO — palanca tactica",
         "V3: Promociones OR=3.64 como anzuelo de primera prueba.\n"
         "V4: El anzuelo no puede comunicar precio bajo (refuerza la distancia simbolica del 'sifrinaje'). "
         "Debe comunicar 'razon especifica para venir'.",
         "[Tipo B]", C_BLUE),
        # CO-2 obligatorio en Argumento 3
        ("DONDE — territorio (CO-2 equilibrado)",
         "V3: Urgencia hoy (#2 mision), mercado grande como horizonte.\n"
         "V4: " + CO2_MAIN[:200] + "...\n[Wording CO-2 mezcla equilibrada — ver slide C1 y C8 para wording completo]",
         "[WoW Confirmado H-2+H-3]", C_GREEN),
        ("A QUIEN — target",
         "V3: Seg_2 Pragmaticos Convertibles (33%), 0% preferencia, menor resistencia al precio.\n"
         "V4: La barrera del seg_2 es identitaria, no de precio. La activacion correcta es la primera "
         "experiencia que resignifique pertenencia.",
         "[Triangulado]", C_BLUE),
    ]
    arg_w = (SLIDE_W - 700000) // 4
    for ai, (atitle, adesc, alabel, aclr) in enumerate(args):
        ax = 300000 + ai * (arg_w + 60000)
        add_rect(slide, ax, 2000000, arg_w, SLIDE_H - 2500000, fill_color=C_LIGHT_BG)
        add_textbox(slide, ax + 60000, 2060000, arg_w - 120000, 340000,
                    atitle, Pt(9), C_RED, bold=True)
        add_textbox(slide, ax + 60000, 2440000, arg_w - 120000, SLIDE_H - 3100000,
                    adesc, Pt(8), C_DARK, wrap=True)
        add_rect(slide, ax, SLIDE_H - 480000, arg_w, 200000, fill_color=C_LIGHT_BG)
        add_textbox(slide, ax + 40000, SLIDE_H - 460000, arg_w - 80000, 180000,
                    alabel, Pt(7.5), aclr, bold=True)
    add_footer(slide)

def make_E3(prs):
    """DW-2 + CO-3 Opcion C. Caveats DW2-a y DW2-b en lugares correctos. Nota scope en E3."""
    slide = add_slide(prs)
    set_bg(slide, C_WHITE)
    header_bar(slide, "Seccion E — Sintesis V4 + Recomendaciones Actualizadas",
               "E3 — Recomendaciones priorizadas V4: quick wins, mid-term, long-term")
    info_box(slide,
             "TAKEAWAY: Las tres recomendaciones top V3 se mantienen — V4 las reordena por accionabilidad "
             "considerando los hallazgos de canal digital y la dinamica WoW.",
             l=300000, t=1050000, w=SLIDE_W - 600000, h=280000)

    horizons = [
        ("QUICK WIN\n(inmediato, bajo costo)",
         [
             "Integracion digital basica del Gama Club: saldo visible en app + equivalente en dinero + canje en caja. "
             "La disposicion de lealtad ya existe en todos los segmentos — solo falta la visibilidad. [V4-NEW-3 Tipo B]",
             "Reactivar stickers con paraguas que no comunique precio bajo — 'razon especifica para venir esta semana'. "
             "El mecanismo del anzuelo funciona (V3 OR=3.64) pero el mensaje debe ser correcto per barrera identitaria V4. [Triangulado V3+V4]",
         ], C_GREEN),
        ("MID-TERM\n(Q3-Q4 2026)",
         [
             # CO-3 Opcion C: canal digital incluido con nota scope
             "Activar canal digital: resolver inconsistencia de inventario en tiempo real (Barrera 1) + "
             "disenar propuesta de delivery de no-perecederos primero (Barrera 2). "
             "[V4-NEW-6 — AMPLIACION DEL BRIEF ORIGINAL — sugiere conversacion de scope con Gama]",
             "Diseno de activacion para Seg_2: secuencia de 3 pasos (anzuelo -> experiencia guiada -> retencion por Gama Club). "
             "El mensaje debe hacer creible que el diferencial vale, no que el precio es bajo. [V3 Rec.3 + mecanismo V4]",
         ], C_BLUE),
        ("LONG-TERM\n(Q1 2027 en adelante)",
         [
             # DW-2 lenguaje condicional — "se recomienda explorar"
             "Se recomienda explorar la plataforma de comunicacion hacia el registro de acompanamiento-guia: "
             "materializar el arquetipo femenino que el shopper ya proyecta. [V4-NEW-4 — sujeto DW-2]",
             "Ampliar misiones en que Gama es elegida para capturar el espacio disponible de CM. "
             "[Triangulado WoW H-4 + cuali B-V4-6]",
         ], C_DARKRED),
    ]

    col_w = (SLIDE_W - 700000) // 3
    for ci, (hor_title, hor_bullets, hor_clr) in enumerate(horizons):
        cx = 300000 + ci * (col_w + 80000)
        add_rect(slide, cx, 1430000, col_w, 380000, fill_color=hor_clr)
        add_textbox(slide, cx + 60000, 1450000, col_w - 120000, 340000,
                    hor_title, Pt(9), C_WHITE, bold=True)
        add_rect(slide, cx, 1860000, col_w, SLIDE_H - 2900000, fill_color=C_LIGHT_BG)
        y_b = 1920000
        for b in hor_bullets:
            add_textbox(slide, cx + 60000, y_b, col_w - 120000, 480000,
                        "* " + b, Pt(8), C_DARK, wrap=True)
            y_b += 520000

    # C-LIT-DW2-a (Long-term arquetipo) + C-LIT-DW2-b
    caveat_box(slide, "Recomendacion Long-Term — " + C_LIT_DW2_A,
               l=300000, t=SLIDE_H - 720000, w=SLIDE_W - 600000, h=240000, size=Pt(7.5))
    caveat_box(slide, C_LIT_DW2_B,
               l=300000, t=SLIDE_H - 460000, w=SLIDE_W - 600000, h=240000, size=Pt(7.5),
               bg=C_WARN_BG, fg=C_WARN_FG)
    # CO-3 scope note visible
    add_textbox(slide, 300000, SLIDE_H - 200000, SLIDE_W - 600000, 180000,
                "NOTA: Las recomendaciones de canal digital (Mid-Term) son ampliacion del brief original — "
                "sugieren conversacion de scope con Gama. El Owner puede ajustar este wording pre-envio a Cora.",
                Pt(7.5), C_WARN_FG, italic=True, wrap=True)
    add_footer(slide)

def make_E4(prs):
    slide = add_slide(prs)
    set_bg(slide, C_WHITE)
    header_bar(slide, "Seccion E — Sintesis V4 + Recomendaciones Actualizadas",
               "E4 — Recomendaciones con priorizacion por evidencia")
    info_box(slide,
             "TAKEAWAY: La priorizacion no es arbitraria — se deriva directamente del nivel de certeza "
             "de la evidencia y la accionabilidad estimada de cada recomendacion.",
             l=300000, t=1050000, w=SLIDE_W - 600000, h=280000)
    headers = ["Recomendacion", "Evidencia", "Accionabilidad", "Prioridad"]
    col_x = [300000, 4600000, 7200000, 9800000]
    col_w = [4200000, 2500000, 2500000, 2100000]
    y = 1430000
    for ci, (cx, cw, ch) in enumerate(zip(col_x, col_w, headers)):
        add_rect(slide, cx, y, cw - 60000, 320000, fill_color=C_RED)
        add_textbox(slide, cx + 40000, y + 50000, cw - 80000, 220000,
                    ch, Pt(9), C_WHITE, bold=True)
    rows = [
        ("Gama Club: integracion digital", "Tipo B — 6/7 FG robusto", "Alta: cambio producto digital", "1 — Quick win"),
        ("Anzuelo con mensaje correcto (sin precio)", "Triangulado V3+V4", "Alta: reorientar campana existente", "2 — Quick win"),
        ("Activacion Seg_2 con barrera identitaria", "Triangulado + cuali", "Media: requiere diseno activacion", "3 — Mid-term"),
        ("Canal digital: inventario en tiempo real", "Tipo B digital", "Media: inversion tecnologica", "4 — Mid-term"),
        ("Comunicacion arquetipo-guia", "Tipo B emergente", "Media: requiere desarrollo creativo", "5 — Long-term"),
        ("Ampliar misiones (captura espacio CM)", "WoW H-4 + cuali", "Media-baja: requiere repositioning", "6 — Long-term"),
    ]
    row_h = 460000
    y += 360000
    for ri, (rec, evid, acc, pri) in enumerate(rows):
        bg = C_LIGHT_BG if ri % 2 == 0 else C_WHITE
        for ci, (cx, cw, txt) in enumerate(zip(col_x, col_w, [rec, evid, acc, pri])):
            add_rect(slide, cx, y, cw - 60000, row_h - 40000, fill_color=bg)
            add_textbox(slide, cx + 40000, y + 40000, cw - 80000, row_h - 80000,
                        txt, Pt(8.5), C_DARK, wrap=True)
        y += row_h
    caveat_box(slide,
               "La tabla es una recomendacion del equipo analitico — no una decision ejecutiva. "
               "La prioridad final la define Gama en funcion de recursos, timing y estrategia de marca.",
               l=300000, t=SLIDE_H - 500000, w=SLIDE_W - 600000, h=280000, size=Pt(8))
    add_footer(slide)

def make_E5(prs):
    """CO-3 scope note en E5."""
    slide = add_slide(prs)
    set_bg(slide, C_WHITE)
    header_bar(slide, "Seccion E — Sintesis V4 + Recomendaciones Actualizadas",
               "E5 — Que sigue para ola 2027 y decisiones abiertas para Owner/Gama")
    info_box(slide,
             "TAKEAWAY: V4 cierra la cadena analitica disponible con los datos actuales; tres conversaciones "
             "con Owner y cinco mejoras metodologicas para ola 2027 aseguran que la siguiente oleada resuelva "
             "las hipotesis que V4 no puede cerrar.",
             l=300000, t=1050000, w=SLIDE_W - 600000, h=320000)
    # Left column: hipotesis abiertas
    add_rect(slide, 300000, 1450000, 5700000, SLIDE_H - 1950000, fill_color=C_LIGHT_BG)
    add_textbox(slide, 360000, 1490000, 5600000, 320000,
                "Hipotesis que requieren ola 2027:", Pt(10), C_RED, bold=True)
    hyps = [
        "H-5 y H-8: tendencias imagen experiencial (+6.3pp) y debilitamiento imagen precio (-3.9pp) "
        "— solo tendencia, requieren sig_99 en ola 2027.",
        "Gap comunicacional (Tipo C): necesita modulo de comunicacion con estimulos de campana en 2027.",
        "Validacion cuantitativa del arquetipo femenino (HQ-7): prueba de concepto creativa con target.",
        "5 mejoras metodologicas: MaxDiff, NPS+switching, CEPs expandidos, penetracion 12m, booster Pref-Gama. "
        "V4 agrega: modulo comunicacion + Rio como benchmark en asociaciones.",
    ]
    yh = 1870000
    for h in hyps:
        add_textbox(slide, 380000, yh, 5520000, 400000,
                    "  * " + h, Pt(8.5), C_DARK, wrap=True)
        yh += 420000

    # Right column: 3 decisiones Owner (CO-1/CO-2/CO-3)
    add_rect(slide, 6200000, 1450000, SLIDE_W - 6500000, SLIDE_H - 1950000, fill_color=C_WARN_BG)
    add_textbox(slide, 6280000, 1490000, SLIDE_W - 6580000, 320000,
                "3 decisiones Owner para esta reunion:", Pt(10), C_WARN_FG, bold=True)
    cos = [
        ("CO-1", "Senal de alerta Rio en proteinas (baja certeza, un verbatim) — "
         "entra al deck para Cora/Gama o se reserva para el memo interno?"),
        ("CO-2", "Tono de la posicion competitiva de Gama en el deck — "
         "ya incorporado como mezcla equilibrada per BR-2 §3. Owner puede ajustar."),
        ("CO-3", "Recomendaciones del canal digital (Gama Club integracion, inventario, frescos) — "
         "incorporadas en Opcion C hibrida con nota de scope. "
         "Owner confirma o ajusta wording de nota antes del envio a Cora."),
    ]
    yco = 1870000
    for (tag, desc) in cos:
        add_rect(slide, 6220000, yco, 600000, 300000, fill_color=C_WARN_FG)
        add_textbox(slide, 6240000, yco + 50000, 560000, 220000,
                    tag, Pt(8), C_WHITE, bold=True)
        add_textbox(slide, 6900000, yco, SLIDE_W - 7200000, 400000,
                    desc, Pt(8), C_DARK, wrap=True)
        yco += 480000
    add_footer(slide)

def make_closing(prs):
    slide = add_slide(prs)
    set_bg(slide, C_RED)
    add_textbox(slide, 500000, 1200000, SLIDE_W - 1000000, 1000000,
                "Notoriedad Gama 2026 — V4", Pt(48), C_WHITE, bold=True)
    add_textbox(slide, 500000, 2300000, SLIDE_W - 1000000, 600000,
                "Deck Principal V4 · Gate Bruna BR-2 V4 aprobado",
                Pt(20), C_WHITE)
    add_textbox(slide, 500000, 3000000, SLIDE_W - 1000000, 400000,
                "Owner: revisar resoluciones CO-1 / CO-2 / CO-3 antes del envio a Cora.",
                Pt(14), C_WARN_BG)
    add_textbox(slide, 500000, 3500000, SLIDE_W - 1000000, 400000,
                "12 caveats literales BR-2 §5 incorporados en slides correspondientes.",
                Pt(12), C_WHITE)
    meta = [
        "Cora Urrea + Raoul Bermudez — 2026-05-17",
        "Confidencial / NDA",
        FOOTER_TEXT,
    ]
    for i, line in enumerate(meta):
        add_textbox(slide, 500000, 4200000 + i * 320000,
                    SLIDE_W - 1000000, 280000,
                    line, Pt(10), C_WHITE)

# ── MAIN BUILD ────────────────────────────────────────────────────────────────

def build():
    prs = new_prs()

    # Cover
    make_cover(prs)

    # Section A
    make_A1(prs)
    make_A2(prs)
    make_A3(prs)
    make_A4(prs)
    make_A5(prs)
    make_A6(prs)

    # Section B placeholder (2 slides)
    make_B_placeholder(prs)

    # Section C
    make_C1(prs)
    make_C2(prs)
    make_C3(prs)
    make_C4(prs)
    make_C5(prs)
    make_C6(prs)
    make_C7(prs)
    make_C8(prs)

    # Section D
    make_D1(prs)
    make_D2(prs)
    make_D3(prs)
    make_D4(prs)
    make_D5(prs)
    make_D6(prs)

    # Section E
    make_E1(prs)
    make_E2(prs)
    make_E3(prs)
    make_E4(prs)
    make_E5(prs)

    # Closing
    make_closing(prs)

    prs.save(OUTPUT_PATH)
    print(f"Saved: {OUTPUT_PATH}")
    print(f"Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    build()
