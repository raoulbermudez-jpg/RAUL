import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

# Paleta heredada de V3
RED_GAMA     = RGBColor(0x7A, 0x12, 0x12)
RED_DARK     = RGBColor(0x4A, 0x0A, 0x0A)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT    = RGBColor(0x33, 0x33, 0x33)
GRAY_LIGHT   = RGBColor(0xF8, 0xF8, 0xF8)
GRAY_MID     = RGBColor(0x80, 0x80, 0x80)
ORANGE_WARN  = RGBColor(0xD9, 0x73, 0x06)
ORANGE_BG    = RGBColor(0xFF, 0xF3, 0xE0)
BLUE_INFO    = RGBColor(0x1A, 0x56, 0x8A)
BLUE_LIGHT   = RGBColor(0xE8, 0xF4, 0xFD)

# Dimensiones 16:9 widescreen, heredadas de V3
SW = 12191695
SH = 6858000

HEADER_H     = 640080
HEADER_TOP   = 0
MARGIN_L     = 274320
MARGIN_R     = 274320
CONTENT_W    = SW - MARGIN_L - MARGIN_R
CONTENT_TOP  = HEADER_H + 228600
CONTENT_H    = 4206240
CALLOUT_TOP  = 868680
CALLOUT_H    = 640080
BODY_TOP     = 1645919
WARN_TOP     = 6080760
WARN_H       = 365760
FOOTER_TOP   = 6492240
FOOTER_H     = 274320

FONT_TITLE   = 20
FONT_SECTION = 10
FONT_BODY    = 11
FONT_CALLOUT = 11
FONT_FOOTER  = 8
FONT_CAVEAT  = 9
FONT_BADGE   = 9

OUTPUT_PATH = r"C:\rAUL\03-projects\consultoria-externa\gama-notoriedad-2026\03-deck\V4\2026-05-17_Notoriedad-Gama-2026_Resumen-Ejecutivo-V4.pptx"

FOOTER_TEXT = "Notoriedad Gama 2026 · V4 · WoW 2025↔2026 + Cuali Profundo · Confidencial NDA"


def new_prs():
    prs = Presentation()
    prs.slide_width  = Emu(SW)
    prs.slide_height = Emu(SH)
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]
    return prs.slides.add_slide(layout)


def add_rect(slide, left, top, width, height, fill_rgb, line_rgb=None):
    shape = slide.shapes.add_shape(
        1,
        Emu(left), Emu(top), Emu(width), Emu(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_rgb
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_rgb, line_rgb=None):
    shape = slide.shapes.add_shape(
        5,
        Emu(left), Emu(top), Emu(width), Emu(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_rgb
    return shape


def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))


def set_run(run, text, size_pt, bold=False, color=DARK_TEXT, italic=False):
    run.text = text
    run.font.size  = Pt(size_pt)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color


def add_para_with_runs(tf, runs_spec, space_before=0, space_after=0, align=PP_ALIGN.LEFT):
    para = tf.add_paragraph()
    para.alignment = align
    para.space_before = Pt(space_before)
    para.space_after  = Pt(space_after)
    for (text, size_pt, bold, color) in runs_spec:
        run = para.add_run()
        set_run(run, text, size_pt, bold, color)
    return para


def clear_tf(tf):
    from pptx.oxml.ns import qn
    from lxml import etree
    txBody = tf._txBody
    paras = txBody.findall(qn('a:p'))
    for p in paras[1:]:
        txBody.remove(p)
    first = paras[0]
    for child in list(first):
        first.remove(child)


def add_standard_header(slide, slide_num, slide_title):
    add_rect(slide, 0, 0, SW, HEADER_H, RED_GAMA)
    tb = add_textbox(slide, SW - 2133600, 137160, 2011680, 365760)
    tf = tb.text_frame
    tf.word_wrap = True
    clear_tf(tf)
    add_para_with_runs(tf, [("Resumen Ejecutivo", FONT_SECTION, True, WHITE)])
    tb2 = add_textbox(slide, MARGIN_L, 137160, SW - MARGIN_L - 2133600, 457200)
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    clear_tf(tf2)
    add_para_with_runs(tf2, [(slide_title, FONT_TITLE, True, WHITE)])


def add_callout_box(slide, text, fill=GRAY_LIGHT, text_color=RED_DARK):
    shape = add_rounded_rect(slide, MARGIN_L, CALLOUT_TOP, CONTENT_W, CALLOUT_H, fill)
    tf = shape.text_frame
    tf.word_wrap = True
    clear_tf(tf)
    add_para_with_runs(tf, [("★ " + text, FONT_CALLOUT, True, text_color)])


def add_warn_box(slide, text, fill=ORANGE_BG, text_color=ORANGE_WARN):
    shape = add_rounded_rect(slide, MARGIN_L, WARN_TOP, CONTENT_W, WARN_H, fill)
    tf = shape.text_frame
    tf.word_wrap = True
    clear_tf(tf)
    add_para_with_runs(tf, [(text, FONT_CAVEAT, True, text_color)])


def add_info_box(slide, left, top, width, height, text, text_color=BLUE_INFO, fill=BLUE_LIGHT):
    shape = add_rounded_rect(slide, left, top, width, height, fill)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    clear_tf(tf)
    add_para_with_runs(tf, [(text, FONT_CAVEAT, False, text_color)])


def add_footer(slide, slide_num):
    tb = add_textbox(slide, MARGIN_L, FOOTER_TOP, CONTENT_W, FOOTER_H)
    tf = tb.text_frame
    tf.word_wrap = True
    clear_tf(tf)
    add_para_with_runs(tf, [
        (f"V4 · slide {slide_num} · {FOOTER_TEXT}", FONT_FOOTER, False, GRAY_MID)
    ])


def body_textbox(slide):
    tb = add_textbox(slide, MARGIN_L, BODY_TOP, CONTENT_W, CONTENT_H)
    tf = tb.text_frame
    tf.word_wrap = True
    clear_tf(tf)
    return tf


def bullet(tf, label, label_bold=True, value="", value_bold=False,
           indent=False, size=FONT_BODY, before=2, after=1):
    indent_prefix = "    " if indent else ""
    para = tf.add_paragraph()
    para.space_before = Pt(before)
    para.space_after  = Pt(after)
    r1 = para.add_run()
    set_run(r1, indent_prefix + "▸ ", size, True, RED_GAMA)
    if label:
        r2 = para.add_run()
        set_run(r2, label, size, label_bold, DARK_TEXT)
    if value:
        r3 = para.add_run()
        set_run(r3, value, size, value_bold, DARK_TEXT)


def badge_text(tf, text, size=FONT_BADGE, before=3, after=1, color=BLUE_INFO):
    para = tf.add_paragraph()
    para.space_before = Pt(before)
    para.space_after  = Pt(after)
    run = para.add_run()
    set_run(run, text, size, False, color)


def plain_text(tf, text, size=FONT_BODY, bold=False, before=2, after=1, color=DARK_TEXT, italic=False):
    para = tf.add_paragraph()
    para.space_before = Pt(before)
    para.space_after  = Pt(after)
    run = para.add_run()
    set_run(run, text, size, bold, color, italic=italic)


def section_divider(tf, text, size=FONT_BODY, before=5, after=2):
    plain_text(tf, text, size=size, bold=True, before=before, after=after, color=RED_GAMA)


def build_r1(prs):
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, 2286000, SH, RED_GAMA)

    tb = add_textbox(slide, 365760, 457200, 1828800, 914400)
    tf = tb.text_frame; clear_tf(tf)
    add_para_with_runs(tf, [("GAMA", 36, True, WHITE)])

    tb2 = add_textbox(slide, 365760, 5760720, 1828800, 731520)
    tf2 = tb2.text_frame; clear_tf(tf2)
    add_para_with_runs(tf2, [("2026", 28, True, WHITE)])

    shape = add_rounded_rect(slide, 365760, 5120640, 1828800, 457200, WHITE)
    tf3 = shape.text_frame; clear_tf(tf3)
    add_para_with_runs(tf3, [("V4 · RE", 18, True, RED_GAMA)])

    add_rect(slide, 2286000, 0, SW - 2286000, 45720, BLUE_INFO)

    tb4 = add_textbox(slide, 2743200, 2103120, 8961120, 1371600)
    tf4 = tb4.text_frame; clear_tf(tf4)
    add_para_with_runs(tf4, [("Notoriedad y Preferencia de Marca — Gama 2026", 32, True, RED_GAMA)])

    tb5 = add_textbox(slide, 2743200, 3566160, 8961120, 914400)
    tf5 = tb5.text_frame; clear_tf(tf5)
    add_para_with_runs(tf5, [
        ("Resumen Ejecutivo V4 · Capa WoW + Análisis Cualitativo Profundo + Triangulación",
         18, False, DARK_TEXT)
    ])

    meta_lines = [
        "Fecha: 2026-05-17",
        "Equipo analítico: Cora Urrea + Raoul Bermudez",
        "Versión: V4 Resumen Ejecutivo — complementa V3 RE (2026-05-17) y V2 RE (2026-05-16)",
        "Confidencial / NDA",
        "Nota al lector: este RE V4 puede leerse de forma independiente o como secuencia del RE V3. "
        "Las 3 capas de V4 (WoW + Cuali + Triangulación) se sintetizan en las 7 slides siguientes. "
        "El deck principal V4 (~50 slides) contiene todo el análisis detallado con verbatims y caveats completos.",
    ]
    y_offset = 4663440
    for line in meta_lines:
        tb = add_textbox(slide, 2743200, y_offset, 8961120, 320040)
        tf = tb.text_frame; clear_tf(tf)
        add_para_with_runs(tf, [(line, 11, False, DARK_TEXT)])
        y_offset += 320040


def build_r2(prs):
    slide = blank_slide(prs)
    add_standard_header(slide, 2, "Tres preguntas que V3 no podía responder — V4 las responde")
    add_callout_box(slide,
        "V4 no contradice V3: suma la dimensión longitudinal (¿cambió el mercado?), "
        "el mecanismo cualitativo (¿por qué operan los drivers?) y seis hallazgos genuinamente nuevos.")

    tf = body_textbox(slide)

    section_divider(tf, "Pregunta 1 — ¿La posición de Gama cambió entre 2025 y 2026?", before=0)
    bullet(tf,
        "Respuesta: ",
        value="No. El embudo de Gama (TOM, consideración, compra, preferencia) es "
              "estadísticamente estable entre olas. 0/8 ítems con variación significativa. "
              "En un mercado donde Rio crece +17pp TOM y Paramo +12pp TOM, Gama no retrocedió.")
    badge_text(tf, "    [WoW Confirmado — H-1]")

    section_divider(tf, "Pregunta 2 — ¿Por qué la atención genera preferencia? ¿Cuál es el mecanismo real?")
    bullet(tf,
        "Respuesta: ",
        value="El corpus de 12 documentos FG (42.094 chars, 6 segmentos) documenta en 6 de 7 grupos el mismo "
              "imaginario espontáneo: Gama como una “guía experta de confianza” que acompaña la compra. "
              "No es el reconocimiento nominal — es el acompañamiento. El OR=5.73 cuanti mide exactamente ese territorio.")
    badge_text(tf, "    [Triangulado: cuanti V3 + cuali V4]")

    section_divider(tf, "Pregunta 3 — ¿Por qué el 92% que no prefiere Gama no convierte, si la atención es diferenciadora?")
    bullet(tf,
        "Respuesta: ",
        value="La barrera no es el precio como cálculo económico. Es el precio como señal de posición social. "
              "El shopper se autoexcluye antes de experimentar el diferencial. Los descuentos no mueven esta barrera — la primera experiencia sí.")
    badge_text(tf, "    [Triangulado: cuanti V3 + cuali V4]")

    section_divider(tf, "4 claims de alta certeza V3 siguen vigentes:")
    bullet(tf, "Atención driver #1 (OR=5.73) · Precio no driver (OR=1.03) · 3 segmentos k-means · Gap comunicacional ",
           value="[Hipótesis V4 — evidencia convergente]")
    bullet(tf, "4 reformulaciones de mecanismo (no de dirección): todas son Posición II — el hallazgo es más preciso, no contradictorio.")

    add_footer(slide, 2)


def build_r3(prs):
    slide = blank_slide(prs)
    add_standard_header(slide, 3,
        "La dinámica del mercado venezolano de supermercados: un año de movimiento")
    add_callout_box(slide,
        "10 ítems estadísticamente significativos en 57 comparaciones WoW: 9 son aumentos de competidores "
        "(Rio y Paramo), 1 es la caída de CM. Gama es la única gran cadena que no se movió — "
        "en ninguna dirección — fortaleza defensiva en mercado turbulento. Al mismo tiempo, el mercado "
        "se mueve a su alrededor.")

    tf = body_textbox(slide)

    section_divider(tf, "Rio — ganador absoluto:", before=0)
    bullet(tf, "TOM: +17.0pp (sig_99) · Consideración: +19.6pp (sig_99) · Compra: +12.5pp (sig_99)")
    bullet(tf, "NSE C+/C (segmento natural de Gama): TOM Rio +25.8pp (sig_99, exploratorio)",
           indent=True)
    badge_text(tf, "    [WoW Confirmado — H-2 Tipo A sig_99]  [Subgrupo NSE exploratorio: CV-WOW-005*]")

    section_divider(tf, "Paramo — consolidación en toda la cadena:")
    bullet(tf, "TOM: +12.1pp (sig_99) · Asistida: +10.8pp (sig_99) · Consideración: +17.3pp (sig_99)")
    bullet(tf, "NSE C+/C: TOM Paramo +22.3pp (sig_99, exploratorio)", indent=True)
    badge_text(tf, "    [WoW Confirmado — H-3 Tipo A sig_99]  [Subgrupo NSE exploratorio: CV-WOW-005*]")

    section_divider(tf, "Central Madeirense — única pérdida significativa:")
    bullet(tf, "Compra 3m: -7.7pp (sig_95). Preferida: -5.5pp (tendencia).")
    badge_text(tf, "    [WoW Confirmado — H-4 Tipo A sig_95]")

    section_divider(tf, "Gama — posición defensiva sólida en mercado turbulento:")
    bullet(tf, "0/8 ítems del embudo con variación significativa (p_adj>0.10 en todos).")
    bullet(tf,
        "Gama es la única gran cadena del mercado que mantiene posición estable entre 2025 y 2026 — "
        "fortaleza defensiva resultado de activos diferenciadores (24h, atención relacional). "
        "Al mismo tiempo, en el segmento C+/C, Rio y Paramo crecen +25.8pp y +22.3pp TOM. "
        "La señal direccional es que el espacio competitivo se está estrechando, aunque la posición actual de Gama es sólida.")
    badge_text(tf, "    [WoW Confirmado — H-1 Tipo A · 0/8 sig]")

    add_warn_box(slide,
        "CV-WOW-001* Los resultados 2025 se presentan sin ponderación muestral (factor no disponible en BBDD). "
        "Los estimados 2025 pueden tener sesgo de diseño si la muestra original era estratificada.  ·  "
        "CV-WOW-002* Cambios WoW con precaución: composición geográfica difiere entre olas (Libertador sobrerepresentado en 2025).  ·  "
        "CV-WOW-005* Análisis por NSE no pre-registrados en diseño 2025. Hipótesis a confirmar en ola 2027.  ·  "
        "n_2025=785, n_2026=402, BH-FDR. Causalidad no inferible.")

    add_footer(slide, 3)


def build_r4(prs):
    slide = blank_slide(prs)
    add_standard_header(slide, 4, "Gama no retrocedió — pero el mercado se mueve a su alrededor")
    add_callout_box(slide,
        "La estabilidad de Gama en un mercado dinámico es positiva hoy; el crecimiento de Rio y Paramo "
        "en el segmento C+/C es la señal de alerta para ola 2027. La posición actual es sólida.")

    tf = body_textbox(slide)

    section_divider(tf, "Embudo de Gama 2025↔2026 (8 indicadores):", before=0)
    items = [
        ("TOM:", "42.0% → 44.3% (+2.2pp, no significativo)"),
        ("Consideración:", "27.5% → 31.8% (+4.3pp, no significativo)"),
        ("Compra 3m:", "17.8% → 17.7% (-0.2pp, no significativo)"),
        ("Preferida:", "9.7% → 8.0% (-1.7pp, no significativo)"),
        ("Habitual:", "19.4% → 20.2% (+0.8pp, no significativo)"),
    ]
    for lbl, val in items:
        bullet(tf, lbl + " ", value=val, indent=True, before=1, after=0)
    badge_text(tf, "    [WoW Confirmado — H-1 Tipo A · 0/8 sig]")

    section_divider(tf, "Lectura correcta:")
    bullet(tf,
        "Gama es la única gran cadena del mercado que mantiene posición estable entre 2025 y 2026 — "
        "fortaleza defensiva en mercado turbulento. La estabilidad es resultado de los activos diferenciadores "
        "(24h, atención relacional) que los competidores no replican.")
    bullet(tf,
        "Al mismo tiempo, en el segmento C+/C — natural de Gama — Rio y Paramo crecen +25.8pp y +22.3pp TOM. "
        "La señal direccional es que el espacio competitivo se está estrechando, aunque la posición actual de Gama es sólida.")

    section_divider(tf, "Activo defensivo clave — el 24h:")
    bullet(tf,
        "Rio y Paramo no tienen equivalente funcional al 24 horas de Gama. "
        "El corpus FG lo confirma: \"El de 24 horas es un tiro al piso, de verdad que sí.\" (Frecuentes 18-30) + "
        "\"Eso es impelable.\" (Ocasionales 50+). Defender este activo es la primera línea estratégica.")
    badge_text(tf, "    [Triangulado: cuanti V3 + cuali V4]  [Señal alerta C+/C: Hipótesis V4 — pendiente confirmación ola 2027 · CV-WOW-005*]")

    add_warn_box(slide,
        "CV-WOW-001* Los resultados 2025 se presentan sin ponderación muestral (factor no disponible en BBDD). "
        "Los estimados 2025 pueden tener sesgo de diseño si la muestra original era estratificada.  ·  "
        "CV-WOW-005* Los análisis por NSE no fueron pre-registrados en el diseño 2025. "
        "Deben interpretarse como hipótesis a confirmar en ola 2027.  (n_2026 subgrupo C+/C=104)")

    add_footer(slide, 4)


def build_r5(prs):
    slide = blank_slide(prs)
    add_standard_header(slide, 5,
        "V4 revela el POR QUÉ del driver #1 — y por qué los descuentos no convierten al 92% restante")
    add_callout_box(slide,
        "El OR=5.73 mide en realidad el mecanismo de acompañamiento-guía; la barrera que impide la "
        "conversión del 92% no es precio como cálculo sino precio como identidad — eso cambia la estrategia de activación.")

    tf = body_textbox(slide)

    section_divider(tf, "El mecanismo del driver de atención (V4-NEW-1):", before=0)
    bullet(tf, "Cuanti V3: ",
           value="OR=5.73 (IC95: 1.6-20.4), SHAP #1, triple convergencia. El activo es real y robusto.")
    bullet(tf, "Cuali V4 (6 de 7 FGs, emergente no inducido): ")
    plain_text(tf,
        "    “Una mujer joven, pero con experiencia. Que sea como una guía, que te diga: "
        "'Mira, hoy llegó el pan árabe fresco, llévatelo'.”",
        size=FONT_BODY, italic=True, color=DARK_TEXT, before=1, after=1)
    plain_text(tf, "    Mary Francis Bossia, Ocasionales 18-30",
               size=FONT_BADGE, color=GRAY_MID, before=0, after=1)
    bullet(tf, "Implicación comunicacional: ",
           value="el mensaje correcto es 'en Gama siempre hay alguien que te ayuda', "
                 "no 'el personal te conoce'. El primero es creíble; el segundo era una hipérbole.")
    badge_text(tf, "    [Triangulado: cuanti V3 OR=5.73 + cuali V4 — Tipo B]")

    section_divider(tf, "La barrera del sifrinaje (V4-NEW-2):")
    bullet(tf, "El 92% que no prefiere a Gama no llegó a experimentar el activo de atención. "
               "La barrera es la percepción de precio — pero no como cálculo comparativo in-store.")
    bullet(tf, "Mecanismo: ",
           value="el precio de Gama opera como señal de posición social. "
                 "El shopper se autoexcluye sin haber verificado los precios.")
    plain_text(tf,
        "    “Gama es sifrina, es para otro tipo de persona.”",
        size=FONT_BODY, italic=True, color=DARK_TEXT, before=1, after=0)
    plain_text(tf, "    Azahara Betancourt, Ocasionales 18-30",
               size=FONT_BADGE, color=GRAY_MID, before=0, after=1)
    plain_text(tf,
        "    Cita textual de participante de focus group (investigación cualitativa, mayo 2026). "
        "El término ‘sifrinaje’ es vocabulario espontáneo del informante — no una caracterización de Gama "
        "hecha por el equipo analítico. Refleja cómo segmentos de no-usuarios codifican la barrera de entrada percibida.",
        size=FONT_CAVEAT, color=BLUE_INFO, before=1, after=1)
    bullet(tf, "Implicación directa: ",
           value="los descuentos y promociones de precio bajo refuerzan la distancia simbólica. "
                 "La activación correcta es una razón específica para la primera visita que no sea el precio.")
    badge_text(tf, "    [Triangulado: cuanti V3 + cuali V4 — Tipo B]  "
                   "[Hipótesis V4 para tendencia WoW — p_adj 0.053-0.091, no pasa BH-FDR]")

    add_footer(slide, 5)


def build_r6(prs):
    slide = blank_slide(prs)
    add_standard_header(slide, 6,
        "De los hallazgos V4 a la acción: 3 recomendaciones con badge de evidencia y timeline")
    add_callout_box(slide,
        "Las tres recomendaciones de mayor retorno a corto plazo derivan directamente de los hallazgos V4 "
        "más robustos — todas son accionables con activos existentes de Gama o con inversión de bajo costo.")

    tf = body_textbox(slide)

    section_divider(tf, "Rec. 1 — Quick win: Integrar digitalmente el Gama Club (saldo visible + canje en caja)", before=0)
    bullet(tf, "Evidencia base: ",
           value="V4-NEW-3. El Gama Club tiene interés genuino en todos los segmentos (I4, 6/7 FGs). "
                 "La única barrera es operacional: el saldo es invisible en el canal digital. "
                 "Un activo de lealtad que el usuario no puede ver es funcionalmente equivalente a no tenerlo.")
    bullet(tf, "Accionabilidad: alta. Cambio de producto digital sin nueva investigación ni inversión de capital mayor.")
    bullet(tf, "Oportunidad WoW: ",
           value="el espacio disponible por la caída de CM (-7.7pp compra, sig_95) puede capturarse con "
                 "un mecanismo de lealtad visible. El Gama Club bien ejecutado es ese mecanismo.")
    plain_text(tf,
        "    Los hallazgos de canal digital (inconsistencia de inventario, desconfianza en frescos) identifican "
        "una agenda accionable adicional — ampliación del brief original que sugiere conversación de scope con Gama.",
        size=FONT_CAVEAT, color=BLUE_INFO, before=1, after=1)
    badge_text(tf, "    [Cuali V4 FG + digital — Tipo B]  [Validado WoW — oportunidad CM H-4]")

    section_divider(tf, "Rec. 2 — Quick win: Reorientar el anzuelo promocional — 'razón para venir esta semana', no 'precio bajo'")
    bullet(tf, "Evidencia base: ",
           value="V3 Rec. 2 (OR=3.64 promociones) + V4 mecanismo sifrinaje.")
    bullet(tf, "El mensaje correcto no puede comunicar precio bajo (refuerza la distancia simbólica) — "
               "debe comunicar acceso a una experiencia específica ('esta semana hay una razón para venir a Gama').")
    badge_text(tf, "    [Triangulado: cuanti V3 OR=3.64 + cuali V4 mecanismo V4-NEW-2]")

    section_divider(tf, "Rec. 3 — Mid-term: Diseñar la activación para Segmento 2 con el mecanismo identitario correcto")
    plain_text(tf,
        "    La activación del Seg_2 se construye con tres pasos: "
        "(1) anzuelo con razón específica accesible — no precio bajo; "
        "(2) primera experiencia guiada que active el driver de atención; "
        "(3) retención por Gama Club visible. "
        "El cuali V4 confirma que el motor de conversión es resignificar la pertenencia, no neutralizar la diferencia de precio.",
        size=FONT_BODY, color=DARK_TEXT, before=2, after=1)
    bullet(tf, "El seg_2 ya tiene menor resistencia de precio cuanti (diferencia 0.22 puntos en escala 5) y 0% preferencia — alta convertibilidad potencial.")
    bullet(tf, "Se recomienda explorar: ",
           value="el imaginario del shopper sugiere un arquetipo femenino disponible como plataforma de comunicación "
                 "(territorio sin apropiar, 6/7 segmentos FG, emergente no inducido).")
    plain_text(tf,
        "    Este hallazgo describe el imaginario de marca que el shopper ya proyecta espontáneamente — "
        "no propone una campaña específica de personificación. La materialización creativa "
        "(figura concreta, voz, presencia visual) requiere prueba de concepto cuantitativa con el target antes de cualquier ejecución.",
        size=FONT_CAVEAT, color=BLUE_INFO, before=1, after=1)
    badge_text(tf, "    [Triangulado: cuanti V3 k-means + cuali V4 HQ-8]  [Rec.3 Seg_2: silhouette moderado ~0.20 — perfiles como tendencias, no categorías discretas]")

    add_warn_box(slide,
        "Las recomendaciones son valoración del equipo analítico, no un dato empírico. "
        "La accionabilidad estimada asume continuidad de condiciones de mercado. "
        "Rec. 3 (arquetipo) requiere verificación previa de Gama sobre compatibilidad con contratos de imagen vigentes y plataforma de marca existente.")

    add_footer(slide, 6)


def build_r7(prs):
    slide = blank_slide(prs)
    add_standard_header(slide, 7,
        "Antes de Cora — 3 decisiones del Owner y 5 decisiones de wording para Bruna")
    add_callout_box(slide,
        "V4 tiene 3 decisiones estratégicas que solo el Owner puede tomar y 5 decisiones de wording "
        "que solo Bruna puede gatear — el deck no sale hasta que estén resueltas.")

    tf = body_textbox(slide)

    section_divider(tf, "Decisiones Owner (CO-1 / CO-2 / CO-3):", before=0)
    bullet(tf, "CO-1 — Señal de alerta Rio en proteínas: ",
           value="un participante FG (Ocasionales 18-30) reporta comprar proteínas en Rio por calidad. "
                 "Certeza baja (un verbatim). ¿Entra al deck para Cora o al memo interno? "
                 "Owner define el umbral de evidencia para señales de alerta en presentación ejecutiva.")
    bullet(tf, "CO-2 — Tono competitivo del deck: ",
           value="¿La posición de Gama se presenta como 'liderazgo defensivo sólido' "
                 "o como 'señal de alerta competitiva'? Los dos son verdaderos — "
                 "el tono afecta cómo Gama recibe el reporte.")
    bullet(tf, "CO-3 — Alcance de recomendaciones digitales: ",
           value="los hallazgos de canal digital implican recomendaciones sobre el producto tecnológico de Gama. "
                 "El brief original era notoriedad y brand health. ¿Estas recomendaciones entran al deck V4 "
                 "o se reservan para una propuesta de scope adicional?")

    section_divider(tf, "Decisiones wording Bruna (DW-1..DW-5) — resumen:")
    dw_items = [
        ("DW-1:", "¿Verbatim 'sifrinaje' de Azahara Betancourt con nombre o anonimizado en deck Cora/Gama? "
                  "→ Bruna aprobó Opción [A] con nombre + caveat literal. APLICADO en este render."),
        ("DW-2:", "¿Personificación femenina como recomendación estratégica? Verificar restricciones de marca Gama. "
                  "→ Bruna aprobó: lenguaje condicional + caveat DW-2 obligatorio. APLICADO."),
        ("DW-3:", "¿Nivel de detalle del matiz en Pragmáticos Convertibles en deck vs memo técnico? "
                  "→ Bruna aprobó con wording suavizado DW-3. APLICADO."),
        ("DW-4:", "¿Verbatim inventario digital enmarcado como oportunidad o queja? "
                  "→ Bruna aprobó con enmarcado de oportunidad. Aplica a Slide R6."),
        ("DW-5:", "¿El downgrade del gap comunicacional de Tipo B a Tipo C se hace visible al cliente? "
                  "→ Bruna aprobó: visible solo como etiqueta [Hipótesis V4 — evidencia convergente]. APLICADO."),
    ]
    for lbl, val in dw_items:
        bullet(tf, lbl + " ", value=val, before=1, after=0)

    plain_text(tf,
        "Estado del gate Bruna: BR-2 V4 APROBADO con ajustes. "
        "La circulación a Cora/Gama no puede ocurrir hasta que CO-1/CO-2/CO-3 estén resueltos por Owner.",
        size=FONT_CAVEAT, bold=True, color=RED_GAMA, before=4, after=1)

    add_footer(slide, 7)


def build_r8(prs):
    slide = blank_slide(prs)
    add_standard_header(slide, 8,
        "Lo que V4 no puede cerrar — agenda de ola 2027 y decisiones de investigación")
    add_callout_box(slide,
        "V4 cierra la cadena analítica disponible con los datos actuales; cuatro hipótesis quedan abiertas "
        "para ola 2027, y cinco mejoras metodológicas de bajo costo aseguran que la siguiente oleada las resuelva.")

    tf = body_textbox(slide)

    section_divider(tf, "4 hipótesis que requieren ola 2027 para confirmación:", before=0)
    hypo_items = [
        "Tendencia de imagen experiencial (+6.3pp 'atractiva', +5.8pp 'seguro') y debilitamiento de imagen "
        "de precio (-3.9pp) de Gama — actualmente solo tendencia (p_adj 0.053-0.091). Ola 2027 confirma si el patrón alcanza sig_99.",
        "Gap comunicacional: corpus FG no menciona la campaña espontáneamente en ninguno de los 12 documentos. "
        "Hipótesis robusta pero sin confirmación directa — requiere módulo de comunicación con estímulos 'PRECIOS DE TU LADO' en ola 2027.",
        "Señal de alerta Rio en proteínas (pendiente CO-1): si Owner la incluye, ola 2027 debe incluir a Rio "
        "como benchmark en el módulo de asociaciones de calidad/proteínas.",
        "Validación cuantitativa del arquetipo femenino: la convergencia espontánea en FG necesita prueba de "
        "concepto creativa cuantitativa (concept testing con el target) para confirmar accionabilidad como plataforma de comunicación.",
    ]
    for item in hypo_items:
        bullet(tf, item, label_bold=False, before=1, after=0)

    section_divider(tf, "5 mejoras metodológicas obligatorias para ola 2027:")
    meth_items = [
        "MaxDiff reemplaza Likert saturado (P22) — elimina la saturación de T2B >90% en todos los atributos.",
        "NPS + switching explícito (3 preguntas) — mide la fidelidad activa, no solo la preferencia declarada.",
        "CEPs expandidos (15-20 ocasiones vs 5 misiones genéricas) — captura la fragmentación de misiones documentada en FG.",
        "Penetración 12 meses + frecuencia + ticket — construye el modelo de share of wallet que V3/V4 no pueden calcular.",
        "Booster campo Pref-Gama: n=80 real + módulo de comunicación con estímulos de campaña + Rio como benchmark en asociaciones.",
    ]
    for item in meth_items:
        bullet(tf, item, label_bold=False, indent=True, before=1, after=0)

    section_divider(tf, "Horizonte 2028 (opcional, inversión mayor):")
    bullet(tf, "Conjoint/DCE para medir la barrera de precio de forma experimental (USD 10-20K).")
    bullet(tf, "DBA battery básica (colores, logo, slogan sin marca) para medir activos de identidad visual.")

    add_warn_box(slide,
        "ME-4 §8 — el diseño de ola 2027 depende de confirmación de Cora sobre metodología de fieldwork 2026.  ·  "
        "ME-5 §2.4 — los análisis WoW se mejoran significativamente si la ola 2025 puede ser re-ponderada en la próxima iteración.")

    add_footer(slide, 8)


def main():
    prs = new_prs()

    build_r1(prs)
    build_r2(prs)
    build_r3(prs)
    build_r4(prs)
    build_r5(prs)
    build_r6(prs)
    build_r7(prs)
    build_r8(prs)

    prs.save(OUTPUT_PATH)
    print(f"Guardado: {OUTPUT_PATH}")


if __name__ == "__main__":
    main()
