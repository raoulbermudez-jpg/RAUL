import sys
sys.stdout.reconfigure(encoding='utf-8')

import os
import io
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np

from pptx import Presentation
from pptx.util import Emu, Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUTPUT_PATH = r"C:\rAUL\03-projects\consultoria-externa\gama-notoriedad-2026\03-deck\V5\2026-05-18_Notoriedad-Gama-2026_Resumen-Ejecutivo-V5.pptx"
CHARTS_DIR  = r"C:\rAUL\03-projects\consultoria-externa\gama-notoriedad-2026\03-deck\V5\charts"
os.makedirs(CHARTS_DIR, exist_ok=True)

RED_GAMA    = RGBColor(0x7A, 0x12, 0x12)
RED_DARK    = RGBColor(0x4A, 0x0A, 0x0A)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT   = RGBColor(0x33, 0x33, 0x33)
GRAY_LIGHT  = RGBColor(0xF8, 0xF8, 0xF8)
GRAY_MID    = RGBColor(0x80, 0x80, 0x80)
ORANGE_WARN = RGBColor(0xD9, 0x73, 0x06)
ORANGE_BG   = RGBColor(0xFF, 0xF3, 0xE0)
BLUE_INFO   = RGBColor(0x1A, 0x56, 0x8A)
BLUE_LIGHT  = RGBColor(0xE8, 0xF4, 0xFD)
GREEN_SIG   = RGBColor(0x1A, 0x70, 0x2A)
GREEN_LIGHT = RGBColor(0xE6, 0xF4, 0xEA)

HEX_RED     = '#7A1212'
HEX_GRAY    = '#808080'
HEX_BLUE    = '#1A568A'
HEX_GREEN   = '#1A702A'
HEX_ORANGE  = '#D97306'

SW           = 12191695
SH           = 6858000
HEADER_H     = 640080
MARGIN_L     = 274320
MARGIN_R     = 274320
CONTENT_W    = SW - MARGIN_L - MARGIN_R
CONTENT_TOP  = HEADER_H + 228600
CONTENT_H    = 4206240
CALLOUT_TOP  = 868680
CALLOUT_H    = 640080
BODY_TOP     = 1645919
WARN_TOP     = 6080760
WARN_H       = 365760
FOOTER_TOP   = 6492240
FOOTER_H     = 274320

FONT_TITLE   = 20
FONT_SECTION = 10
FONT_BODY    = 11
FONT_CALLOUT = 11
FONT_FOOTER  = 8
FONT_CAVEAT  = 9
FONT_BADGE   = 9
FONT_SMALL   = 9

FOOTER_TEXT = "Notoriedad Gama 2026 · V5 · Analisis independiente · Eje 2026 · Confidencial NDA"


def new_prs():
    prs = Presentation()
    prs.slide_width  = Emu(SW)
    prs.slide_height = Emu(SH)
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_rect(slide, left, top, width, height, fill_rgb, line_rgb=None):
    shape = slide.shapes.add_shape(1, Emu(left), Emu(top), Emu(width), Emu(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_rgb
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_rgb, line_rgb=None):
    shape = slide.shapes.add_shape(5, Emu(left), Emu(top), Emu(width), Emu(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_rgb
    return shape


def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))


def set_run(run, text, size_pt, bold=False, color=DARK_TEXT, italic=False):
    run.text = text
    run.font.size   = Pt(size_pt)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color


def clear_tf(tf):
    from pptx.oxml.ns import qn
    txBody = tf._txBody
    paras  = txBody.findall(qn('a:p'))
    for p in paras[1:]:
        txBody.remove(p)
    first = paras[0]
    for child in list(first):
        first.remove(child)


def add_para_with_runs(tf, runs_spec, space_before=0, space_after=0, align=PP_ALIGN.LEFT):
    para = tf.add_paragraph()
    para.alignment    = align
    para.space_before = Pt(space_before)
    para.space_after  = Pt(space_after)
    for (text, size_pt, bold, color) in runs_spec:
        run = para.add_run()
        set_run(run, text, size_pt, bold, color)
    return para


def add_standard_header(slide, slide_num, slide_title):
    add_rect(slide, 0, 0, SW, HEADER_H, RED_GAMA)
    tb = add_textbox(slide, SW - 2133600, 137160, 2011680, 365760)
    tf = tb.text_frame; tf.word_wrap = True; clear_tf(tf)
    add_para_with_runs(tf, [("V5 · RE", FONT_SECTION, True, WHITE)])
    tb2 = add_textbox(slide, MARGIN_L, 137160, SW - MARGIN_L - 2133600, 457200)
    tf2 = tb2.text_frame; tf2.word_wrap = True; clear_tf(tf2)
    add_para_with_runs(tf2, [(slide_title, FONT_TITLE, True, WHITE)])


def add_callout_box(slide, text, fill=GRAY_LIGHT, text_color=RED_DARK):
    shape = add_rounded_rect(slide, MARGIN_L, CALLOUT_TOP, CONTENT_W, CALLOUT_H, fill)
    tf = shape.text_frame; tf.word_wrap = True; clear_tf(tf)
    add_para_with_runs(tf, [("★ " + text, FONT_CALLOUT, True, text_color)])


def add_warn_box(slide, text, fill=ORANGE_BG, text_color=ORANGE_WARN):
    shape = add_rounded_rect(slide, MARGIN_L, WARN_TOP, CONTENT_W, WARN_H, fill)
    tf = shape.text_frame; tf.word_wrap = True; clear_tf(tf)
    add_para_with_runs(tf, [(text, FONT_CAVEAT, True, text_color)])


def add_info_box(slide, left, top, width, height, text, text_color=BLUE_INFO, fill=BLUE_LIGHT):
    shape = add_rounded_rect(slide, left, top, width, height, fill)
    tf = shape.text_frame; tf.word_wrap = True; tf.auto_size = None; clear_tf(tf)
    add_para_with_runs(tf, [(text, FONT_CAVEAT, False, text_color)])


def add_footer(slide, slide_num):
    tb = add_textbox(slide, MARGIN_L, FOOTER_TOP, CONTENT_W, FOOTER_H)
    tf = tb.text_frame; tf.word_wrap = True; clear_tf(tf)
    add_para_with_runs(tf, [(f"V5 · RE · slide {slide_num} · {FOOTER_TEXT}", FONT_FOOTER, False, GRAY_MID)])


def body_textbox(slide, top=None, height=None):
    t = BODY_TOP if top is None else top
    h = CONTENT_H if height is None else height
    tb = add_textbox(slide, MARGIN_L, t, CONTENT_W, h)
    tf = tb.text_frame; tf.word_wrap = True; clear_tf(tf)
    return tf


def bullet(tf, label, label_bold=True, value="", value_bold=False,
           indent=False, size=FONT_BODY, before=2, after=1):
    prefix = "    " if indent else ""
    para = tf.add_paragraph()
    para.space_before = Pt(before)
    para.space_after  = Pt(after)
    r1 = para.add_run(); set_run(r1, prefix + "▸ ", size, True, RED_GAMA)
    if label:
        r2 = para.add_run(); set_run(r2, label, size, label_bold, DARK_TEXT)
    if value:
        r3 = para.add_run(); set_run(r3, value, size, value_bold, DARK_TEXT)


def badge_text(tf, text, size=FONT_BADGE, before=3, after=1, color=BLUE_INFO):
    para = tf.add_paragraph()
    para.space_before = Pt(before)
    para.space_after  = Pt(after)
    run = para.add_run(); set_run(run, text, size, False, color)


def plain_text(tf, text, size=FONT_BODY, bold=False, before=2, after=1,
               color=DARK_TEXT, italic=False):
    para = tf.add_paragraph()
    para.space_before = Pt(before)
    para.space_after  = Pt(after)
    run = para.add_run(); set_run(run, text, size, bold, color, italic=italic)


def section_divider(tf, text, size=FONT_BODY, before=5, after=2):
    plain_text(tf, text, size=size, bold=True, before=before, after=after, color=RED_GAMA)


def embed_png(slide, png_path, left, top, width, height):
    slide.shapes.add_picture(png_path, Emu(left), Emu(top), Emu(width), Emu(height))


def make_chart_r1_cpc_vs_total():
    path = os.path.join(CHARTS_DIR, "RE_cpc_vs_total.png")
    metrics   = ["TOM", "Preferida"]
    cpc_vals  = [60.6, 13.5]
    tot_vals  = [44.3, 8.0]
    x = np.arange(len(metrics))
    w = 0.32

    fig, ax = plt.subplots(figsize=(6.2, 3.4), dpi=300)
    bars_cpc = ax.bar(x - w/2, cpc_vals, w, label="NSE C+/C (n=104)", color=HEX_RED,    zorder=3)
    bars_tot = ax.bar(x + w/2, tot_vals, w, label="Total (n=402)",     color=HEX_GRAY,   zorder=3)

    for bar, val in zip(bars_cpc, cpc_vals):
        ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.8,
                f"{val:.1f}%", ha='center', va='bottom', fontsize=10, fontweight='bold', color=HEX_RED)
    for bar, val in zip(bars_tot, tot_vals):
        ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.8,
                f"{val:.1f}%", ha='center', va='bottom', fontsize=10, fontweight='bold', color='#444444')

    ax.set_xticks(x)
    ax.set_xticklabels(metrics, fontsize=12)
    ax.set_ylabel("% respondientes", fontsize=10)
    ax.set_ylim(0, 80)
    ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda v, _: f"{v:.0f}%"))
    ax.legend(fontsize=9, loc='upper right')
    ax.set_title("Gama: C+/C vs Total", fontsize=11, color=HEX_RED, fontweight='bold', pad=8)
    ax.grid(axis='y', linestyle='--', alpha=0.4, zorder=0)
    ax.spines[['top','right']].set_visible(False)
    ax.text(0.5, -0.18, "Caveat: n=104 C+/C implica m.e. ±9.8% al 95%",
            ha='center', va='bottom', fontsize=7.5, color='#888888',
            transform=ax.transAxes, style='italic')
    fig.tight_layout(rect=[0, 0.06, 1, 1])
    fig.savefig(path, dpi=300, bbox_inches='tight', facecolor='white')
    plt.close(fig)
    return path


def make_chart_r2_wow_movement():
    path = os.path.join(CHARTS_DIR, "RE_wow_movement.png")

    cadenas = ["Rio", "Paramo", "Gama", "CM"]
    metrics = ["TOM", "Consideracion", "Compra 3m"]

    data = {
        "Rio":    [17.0, 19.6, 12.5],
        "Paramo": [12.1, 17.3, None],
        "Gama":   [0.0,  0.0,  0.0],
        "CM":     [None, None, -7.7],
    }
    sig_map = {
        "Rio":    ['sig_99', 'sig_99', 'sig_99'],
        "Paramo": ['sig_99', 'sig_99', None],
        "Gama":   ['ns',     'ns',     'ns'],
        "CM":     [None,     None,     'sig_95'],
    }

    n_cadenas = len(cadenas)
    n_metrics = len(metrics)
    y_base    = np.arange(n_cadenas) * (n_metrics + 1.2)
    h         = 0.65

    fig, ax = plt.subplots(figsize=(7.5, 4.2), dpi=300)

    for mi, met in enumerate(metrics):
        for ci, cad in enumerate(cadenas):
            val = data[cad][mi]
            sig = sig_map[cad][mi]
            y   = y_base[ci] + mi * h
            if val is None:
                continue
            if sig in ('sig_99', 'sig_95'):
                col = HEX_GREEN if val > 0 else '#C0392B'
                alpha = 1.0
            else:
                col   = HEX_GRAY
                alpha = 0.55
            bar = ax.barh(y, val, h*0.85, color=col, alpha=alpha, zorder=3)
            label = f"{val:+.1f}pp"
            if sig in ('sig_99', 'sig_95'):
                marker = "***" if sig == 'sig_99' else "**"
                label  = f"{val:+.1f}pp {marker}"
            x_pos = val + (0.3 if val >= 0 else -0.3)
            ha    = 'left' if val >= 0 else 'right'
            ax.text(x_pos, y, label, va='center', ha=ha, fontsize=7.5,
                    color=col if sig in ('sig_99','sig_95') else '#666666',
                    fontweight='bold' if sig in ('sig_99','sig_95') else 'normal')

    tick_y = [y_base[ci] + (n_metrics-1)*h/2 for ci in range(n_cadenas)]
    ax.set_yticks(tick_y)
    ax.set_yticklabels(cadenas, fontsize=11, fontweight='bold')
    ax.set_xlabel("Cambio WoW (pp)", fontsize=9)
    ax.axvline(0, color='black', linewidth=0.8)
    ax.set_xlim(-14, 28)
    ax.grid(axis='x', linestyle='--', alpha=0.3, zorder=0)
    ax.spines[['top','right']].set_visible(False)
    ax.set_title("Movimiento competitivo WoW 2025→2026 (Total)", fontsize=10,
                 color=HEX_RED, fontweight='bold', pad=8)

    legend_patches = [
        mpatches.Patch(color=HEX_GREEN, label='Significativo (≥sig_95)'),
        mpatches.Patch(color=HEX_GRAY,  alpha=0.55, label='No significativo / nd'),
    ]
    met_legend = [
        mpatches.Patch(color='none', label=f"Metricas: {' | '.join(metrics)}")
    ]
    ax.legend(handles=legend_patches + met_legend, fontsize=7.5,
              loc='lower right', framealpha=0.85)

    ax.text(0.5, -0.16,
            "n_2025=785, n_2026=402, BH-FDR. CV-WOW-001/002: caveats comparabilidad WoW aplican.",
            ha='center', va='bottom', fontsize=7, color='#888888',
            transform=ax.transAxes, style='italic')
    fig.tight_layout(rect=[0, 0.08, 1, 1])
    fig.savefig(path, dpi=300, bbox_inches='tight', facecolor='white')
    plt.close(fig)
    return path


def make_chart_r3_forest_plot():
    path = os.path.join(CHARTS_DIR, "RE_forest_plot.png")

    drivers = [
        ("Atencion al cliente", 5.73, 1.6,  20.4,  "sig_99"),
        ("Limpieza",            3.99, 0.94, 16.91, "sig_90"),
        ("Promociones",         3.64, 1.1,  11.8,  "sig_95"),
        ("Precio (menor)",      1.03, None, None,   "ns"),
    ]
    drivers = drivers[::-1]

    fig, ax = plt.subplots(figsize=(6.8, 3.2), dpi=300)

    y_positions = list(range(len(drivers)))

    for yi, (name, orval, lo, hi, sig) in enumerate(drivers):
        if sig == 'ns':
            col    = HEX_GRAY
            marker = 'D'
            ms     = 7
        elif sig == 'sig_99':
            col    = HEX_RED
            marker = 'o'
            ms     = 9
        elif sig == 'sig_95':
            col    = HEX_BLUE
            marker = 'o'
            ms     = 8
        else:
            col    = HEX_ORANGE
            marker = 'o'
            ms     = 7

        ax.plot(orval, yi, marker=marker, color=col, markersize=ms, zorder=5)

        if lo is not None and hi is not None:
            ax.hlines(yi, lo, hi, colors=col, linewidth=2.2, zorder=4)
            ax.vlines(lo, yi - 0.08, yi + 0.08, colors=col, linewidth=1.4, zorder=4)
            ax.vlines(hi, yi - 0.08, yi + 0.08, colors=col, linewidth=1.4, zorder=4)
            ic_str = f"IC95 [{lo}, {hi}]"
        else:
            ic_str = "NS"

        label_x = (hi if hi else orval) + 0.25
        sig_label = {"sig_99": "*** p<0.01", "sig_95": "** p<0.05",
                     "sig_90": "* p<0.10",   "ns": "NS"}.get(sig, "")
        ax.text(label_x, yi,
                f"OR={orval:.2f}  {ic_str}  {sig_label}",
                va='center', ha='left', fontsize=7.5, color=col,
                fontweight='bold' if sig != 'ns' else 'normal')

    ax.axvline(1, color='black', linewidth=1, linestyle='--', zorder=3)
    ax.set_yticks(y_positions)
    ax.set_yticklabels([d[0] for d in drivers], fontsize=10)
    ax.set_xlabel("Odds Ratio (log scale)", fontsize=9)
    ax.set_xscale('log')
    ax.set_xlim(0.4, 45)
    ax.set_xticks([0.5, 1, 2, 5, 10, 20])
    ax.get_xaxis().set_major_formatter(matplotlib.ticker.ScalarFormatter())
    ax.grid(axis='x', linestyle='--', alpha=0.35, zorder=0)
    ax.spines[['top','right']].set_visible(False)
    ax.set_title("Drivers de preferencia — Forest plot (simplificado RE)", fontsize=10,
                 color=HEX_RED, fontweight='bold', pad=8)
    ax.text(0.5, -0.20,
            "n pref-Gama=32 (referencial). IC95 amplios; robustez por convergencia 4 metodos.",
            ha='center', va='bottom', fontsize=7, color='#888888',
            transform=ax.transAxes, style='italic')
    fig.tight_layout(rect=[0, 0.09, 1, 1])
    fig.savefig(path, dpi=300, bbox_inches='tight', facecolor='white')
    plt.close(fig)
    return path


def build_r1_portada(prs):
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, 2286000, SH, RED_GAMA)
    tb = add_textbox(slide, 365760, 457200, 1828800, 914400)
    tf = tb.text_frame; clear_tf(tf)
    add_para_with_runs(tf, [("GAMA", 36, True, WHITE)])
    tb2 = add_textbox(slide, 365760, 5760720, 1828800, 731520)
    tf2 = tb2.text_frame; clear_tf(tf2)
    add_para_with_runs(tf2, [("2026", 28, True, WHITE)])
    shape = add_rounded_rect(slide, 365760, 5120640, 1828800, 457200, WHITE)
    tf3 = shape.text_frame; clear_tf(tf3)
    add_para_with_runs(tf3, [("V5 · RE", 18, True, RED_GAMA)])
    add_rect(slide, 2286000, 0, SW - 2286000, 45720, BLUE_INFO)
    tb4 = add_textbox(slide, 2743200, 1828800, 8961120, 1371600)
    tf4 = tb4.text_frame; clear_tf(tf4)
    add_para_with_runs(tf4, [("Notoriedad y Preferencia de Marca — Gama 2026", 30, True, RED_GAMA)])
    tb5 = add_textbox(slide, 2743200, 3292000, 8961120, 914400)
    tf5 = tb5.text_frame; clear_tf(tf5)
    add_para_with_runs(tf5, [("Resumen Ejecutivo V5 · Analisis independiente · Eje 2026", 18, False, DARK_TEXT)])
    meta_lines = [
        "Fecha: 2026-05-18",
        "Equipo analitico: Cora Urrea + Raoul Bermudez",
        "Version: V5 Resumen Ejecutivo — analisis fresh start desde datos 2026",
        "Confidencial / NDA",
        "Nota al lector: este RE V5 es autocontenido y puede leerse sin V3 ni V4. "
        "Eje de analisis: datos 2026. Los datos 2025 son referencia comparativa. "
        "El cualitativo es hipotesis e insight (no conclusiones co-iguales al cuanti). "
        "El deck principal V5 (~40-45 slides) contiene el analisis completo.",
    ]
    y_off = 4389120
    for line in meta_lines:
        tb_m = add_textbox(slide, 2743200, y_off, 8961120, 320040)
        tf_m = tb_m.text_frame; clear_tf(tf_m)
        add_para_with_runs(tf_m, [(line, 10, False, DARK_TEXT)])
        y_off += 320040


def build_r2_tesis(prs):
    slide = blank_slide(prs)
    add_standard_header(slide, 2,
        "Tesis V5: activo real, posicion solida — pero comunicacion y activacion mal dirigidas")

    add_callout_box(slide,
        "Gama tiene el activo relacional mas diferenciador del mercado y su posicion en 2026 "
        "es solida — pero comunica en el territorio equivocado, no activa al tercio de "
        "shoppers convertibles, y su segmento natural (C+/C) esta siendo penetrado por Rio y "
        "Paramo a velocidad que exige respuesta antes de ola 2027.")

    tf = body_textbox(slide)

    section_divider(tf, "Los 3 datos que definen el estado de Gama en 2026:", before=0)

    bullet(tf, "Dato 1 — El activo existe y es diferenciador: ",
           value="Atencion driver #1, OR=5.73*** (convergencia 4 metodos independientes). "
                 "53% de preferentes lo citan espontaneamente. Rio y Paramo no tienen equivalente.")
    badge_text(tf, "    ✅ Conclusion — alta certeza")

    bullet(tf, "Dato 2 — Recall campana = 0%: ",
           value='Recall espontaneo "PRECIOS DE TU LADO" = 0/17 = 0%. '
                 "65% interpreta el mensaje como precio. Precio OR=1.03 NS — atributo menos "
                 "predictivo de preferencia (SHAP #10).")
    badge_text(tf, "    ✅ Conclusion — alta certeza en gap · certeza media en causalidad")

    bullet(tf, "Dato 3 — Rio +17pp TOM (sig_99). Paramo +12pp TOM (sig_99). ",
           value="En C+/C: posiblemente +25.8pp y +22.3pp TOM. "
                 "Gama +2.2pp (no significativo). El mercado se mueve — Gama no.")
    badge_text(tf, "    ✅ Total: sig_99 · ⚠ C+/C: hipotesis WoW — CV-WOW-005")

    section_divider(tf, "Iconografia de certeza V5 (aplica a todo el RE):")
    plain_text(tf,
        "✅ Conclusion cuanti (>=95%, base >=30)  |  "
        "⚠ Hipotesis apoyada (tendencia o soporte cuali robusto)  |  "
        "💡 Insight cuali (FG, no proyectable)  |  "
        "📊 Referencia evolutiva (dato 2025 con caveats)",
        size=FONT_SMALL, color=DARK_TEXT, before=2, after=1)

    add_footer(slide, 2)


def build_r3_cpc(prs, chart_r1_path, chart_r2_path):
    slide = blank_slide(prs)
    add_standard_header(slide, 3,
        "NSE C+/C: donde Gama gana — y donde Rio y Paramo avanzan mas rapido")

    add_callout_box(slide,
        "TOM Gama en C+/C = 60.6% (vs 44.3% Total). El segmento natural de Gama es su "
        "posicion mas fuerte — y posiblemente el segmento donde la presion competitiva "
        "se intensifica mas: Rio +25.8pp y Paramo +22.3pp TOM (exploratorio sig_99, CV-WOW-005).")

    half_w = CONTENT_W // 2 - 91440
    tb = add_textbox(slide, MARGIN_L, BODY_TOP, half_w, CONTENT_H)
    tf = tb.text_frame; tf.word_wrap = True; clear_tf(tf)

    section_divider(tf, "Posicion de Gama en C+/C 2026:", before=0)
    bullet(tf, "TOM: 60.6% (vs 44.3% Total — +16.3pp)")
    bullet(tf, "Preferida: 13.5% (vs 8.0% Total — +5.5pp)")
    bullet(tf, "C+/C es donde el DNA de Gama (atencion, calidad, 24h) "
               "tiene mayor resonancia")
    badge_text(tf, "    ✅ Conclusion — sobrerepresentacion confirmada")

    section_divider(tf, "Amenaza competitiva en C+/C (WoW):")
    bullet(tf, "Rio TOM C+/C: ~30.9% → 56.7% (+25.8pp, sig_99, exploratorio)")
    bullet(tf, "Paramo TOM C+/C: ~18.0% → 40.4% (+22.3pp, sig_99, exploratorio)")
    bullet(tf, "Gama TOM C+/C: ~58.4% → 60.6% (+2.2pp, no sig)")
    badge_text(tf, "    ⚠ Hipotesis WoW C+/C — CV-WOW-005")

    plain_text(tf,
        "Lectura: la posicion actual de Gama en C+/C es solida. "
        "La velocidad de Rio y Paramo en el segmento propio es la "
        "senal de alerta mas importante del estudio.",
        size=FONT_SMALL, color=DARK_TEXT, italic=True, before=3, after=1)

    chart_left  = MARGIN_L + half_w + 182880
    chart_top   = BODY_TOP
    chart_w     = SW - chart_left - MARGIN_R
    chart_h_each = CONTENT_H // 2 - 91440
    embed_png(slide, chart_r1_path, chart_left, chart_top, chart_w, chart_h_each)
    embed_png(slide, chart_r2_path, chart_left, chart_top + chart_h_each + 91440, chart_w, chart_h_each)

    add_warn_box(slide,
        "CV-WOW-005*: analisis WoW por NSE no pre-registrado en diseno 2025. "
        "Hipotesis a confirmar en ola 2027. Composicion geografica 2025 puede sesgar deltas. "
        "n C+/C=104 — m.e. +-9.8% al 95%. Causalidad no inferible.")

    add_footer(slide, 3)


def build_r4_hallazgos(prs, chart_r3_path):
    slide = blank_slide(prs)
    add_standard_header(slide, 4,
        "Los 4 hallazgos centrales del estudio 2026")

    add_callout_box(slide,
        "Cuatro argumentos MECE construidos directamente desde los datos 2026. "
        "Los dos primeros son conclusiones de alta certeza cuanti. "
        "El tercero combina certeza cuanti (segmento) + hipotesis cuali (mecanismo). "
        "El cuarto es hipotesis WoW apoyada — pendiente confirmacion ola 2027.")

    tw = int(CONTENT_W * 0.56)
    tb = add_textbox(slide, MARGIN_L, BODY_TOP, tw, CONTENT_H)
    tf = tb.text_frame; tf.word_wrap = True; clear_tf(tf)

    section_divider(tf, "Arg. 1 — El activo es real, robusto y sin equivalente competitivo", before=0)
    bullet(tf, "OR=5.73*** (4 metodos convergentes). 53% preferentes citan atencion espontaneamente #1.")
    bullet(tf, "Limpieza: driver secundario OR=3.99* (IC95 [0.94, 16.91]).")
    bullet(tf, "Rio y Paramo no tienen equivalente relacional documentado.")
    badge_text(tf, "    ✅ Conclusion — alta certeza")

    section_divider(tf, "Arg. 2 — El activo no esta siendo comunicado")
    bullet(tf, 'Recall espontaneo "PRECIOS DE TU LADO" = 0% (0/17).')
    bullet(tf, "65% interpreta PTL como precio. Precio = OR=1.03 NS (SHAP #10).")
    bullet(tf, "Gama invierte recursos en el territorio menos predictivo de la preferencia.")
    badge_text(tf, "    ✅ Gap: alta certeza · causalidad: certeza media — [Hipotesis V4 — evidencia convergente]")

    section_divider(tf, "Arg. 3 — El segmento de conversion existe y tiene mecanismo conocido (hipotesis)")
    bullet(tf, "Seg 2 (Pragmaticos Convertibles, 33%): 0% preferencia actual, "
               "menor resistencia precio cuanti (3.44 vs 3.66/5).")
    bullet(tf, "Hipotesis cuali: barrera identitaria (sifrinaje). "
               "Activacion correcta: primera experiencia, no descuento.")
    plain_text(tf,
        "    La activacion del Seg_2 se construye con tres pasos: "
        "(1) anzuelo con razon especifica accesible — no precio bajo; "
        "(2) primera experiencia guiada que active el driver de atencion; "
        "(3) retencion por Gama Club visible.",
        size=FONT_SMALL, color=DARK_TEXT, italic=True, before=1, after=1)
    badge_text(tf, "    ✅ Segmento: k-means cuanti · 💡 Mecanismo: hipotesis cuali")

    section_divider(tf, "Arg. 4 — La presion en C+/C se intensifica (hipotesis WoW)")
    bullet(tf, "Rio +25.8pp TOM en C+/C (exploratorio sig_99). Paramo +22.3pp (exploratorio sig_99).")
    bullet(tf, "Posicion actual Gama: solida (TOM 60.6%). La ventana de respuesta ofensiva esta abierta hoy.")
    badge_text(tf, "    ⚠ Hipotesis WoW C+/C — pendiente confirmacion ola 2027 · CV-WOW-005")

    cw = int(CONTENT_W * 0.40)
    cl = MARGIN_L + tw + 182880
    embed_png(slide, chart_r3_path, cl, BODY_TOP, cw, int(CONTENT_H * 0.95))

    add_warn_box(slide,
        "CV-WOW-001*: resultados 2025 sin ponderacion muestral (factor no disponible). "
        "CV-WOW-002*: composicion geografica difiere entre olas. "
        "n_2025=785, n_2026=402, BH-FDR. Causalidad no inferible.")

    add_footer(slide, 4)


def build_r5_recomendaciones(prs):
    slide = blank_slide(prs)
    add_standard_header(slide, 5,
        "Recomendaciones priorizadas: 3 quick wins + mid-term con badge de evidencia")

    add_callout_box(slide,
        "Las tres recomendaciones de mayor retorno a corto plazo derivan de los hallazgos "
        "mas robustos del estudio. Todas son accionables con activos existentes de Gama "
        "o con inversion de bajo costo. La mid-term requiere diseno de activacion.")

    tf = body_textbox(slide)

    section_divider(tf, "Rec. 1 — Quick win: Activar digitalmente el Gama Club", before=0)
    bullet(tf, "Que: ",
           value="saldo visible en app + equivalente en dinero + opcion de canje en caja.")
    bullet(tf, "Evidencia: ",
           value="interes genuino y transversal en todos los segmentos. "
                 "La unica barrera documentada es operacional — opacidad digital del saldo.")
    bullet(tf, "Accionabilidad: alta. Cambio de producto digital, sin nueva investigacion ni capex mayor.")
    bullet(tf, "Potencial adicional: ",
           value="espacio disponible de CM (-7.7pp compra, sig_95) puede capturarse con "
                 "un mecanismo de lealtad visible bien ejecutado.")
    badge_text(tf, "    💡 Insight cuali 6/7 FG · Validar con cuanti ola 2027 · ✅ WoW oportunidad CM (sig_95)")

    section_divider(tf, "Rec. 2 — Quick win: Reorientar el anzuelo — 'razon para venir', no 'precio bajo'")
    bullet(tf, "Que: ",
           value="mantener stickers/Cashea/Club como palanca tactica (OR=3.64** confirmado), "
                 "pero cambiar el mensaje de 'precio bajo' a 'razon especifica para venir esta semana'.")
    bullet(tf, "Por que: ",
           value="OR=3.64** confirma que las promociones son el driver #3. "
                 "Pero el mensaje de precio bajo refuerza la barrera identitaria del Seg 2 "
                 "(hipotesis cuali sifrinaje). El mismo activo con el mensaje correcto no activa esa barrera.")
    badge_text(tf, "    ✅ OR=3.64** cuanti · 💡 Mecanismo identitario: hipotesis cuali")

    section_divider(tf, "Rec. 3 — Mid-term: Activacion Seg 2 (33% del mercado)")
    plain_text(tf,
        "    La activacion del Seg_2 se construye con tres pasos: "
        "(1) anzuelo con razon especifica accesible — no precio bajo; "
        "(2) primera experiencia guiada que active el driver de atencion (OR=5.73); "
        "(3) retencion por Gama Club visible + saldo acumulado comunicado. "
        "El cuali V4 confirma que el motor de conversion es resignificar la pertenencia, "
        "no neutralizar la diferencia de precio.",
        size=FONT_BODY, color=DARK_TEXT, before=2, after=1)
    bullet(tf, "Seg 2: 0% preferencia actual + menor resistencia precio cuanti (3.44 vs 3.66/5). "
               "Mayor retorno C/P de los tres segmentos.")
    badge_text(tf, "    ✅ k-means cuanti · 💡 Mecanismo cuali · [Hipotesis V4 — evidencia convergente]")

    plain_text(tf,
        "    Rec. 4 (mid-term, sujeta a CO-3): canal digital — inventario en tiempo real + propuesta "
        "inicial no perecederos. Rec. 5 (long-term): respuesta ofensiva en C+/C. "
        "Rec. 6 (long-term, sujeta a CO-3 + DW-2): explorar arquetipo 'acompanamiento-guia'.",
        size=FONT_SMALL, color=GRAY_MID, before=4, after=1)

    add_warn_box(slide,
        "Las recomendaciones son valoracion del equipo analitico — no un dato empirico. "
        "Accionabilidad estimada asume continuidad de condiciones de mercado. "
        "Rec. 3 (sifrinaje + activacion) requiere validacion cuantitativa del mecanismo de barrera. "
        "IN7-1.4: el cualitativo establece mecanismos e hipotesis — NO prevalencia estadistica.")

    add_footer(slide, 5)


def build_r6_tono_dual(prs):
    slide = blank_slide(prs)
    add_standard_header(slide, 6,
        "Dos lecturas validas — el Owner elige cual presentar a la Junta (CO-2)")

    add_callout_box(slide,
        "Ambas lecturas son estadisticamente correctas. La diferencia es estrategica, "
        "no estadistica. El Owner elige el tono que mejor sirve los objetivos de la relacion "
        "con Gama antes de la presentacion a la Junta.")

    col_w   = CONTENT_W // 2 - 137160
    col_gap = 274320
    col1_l  = MARGIN_L
    col2_l  = MARGIN_L + col_w + col_gap

    hdr1 = add_rounded_rect(slide, col1_l, BODY_TOP, col_w, 365760, GRAY_LIGHT)
    tf_h1 = hdr1.text_frame; clear_tf(tf_h1)
    add_para_with_runs(tf_h1, [("CO-2 Opcion A — Liderazgo defensivo solido",
                                FONT_BODY, True, RED_DARK)], align=PP_ALIGN.CENTER)

    col1_body_top = BODY_TOP + 365760 + 91440
    tb1 = add_textbox(slide, col1_l, col1_body_top, col_w, 3291840)
    tf1 = tb1.text_frame; tf1.word_wrap = True; clear_tf(tf1)
    plain_text(tf1,
        '"0/8 metricas del embudo de Gama variaron significativamente. '
        'Gama es la unica gran cadena del mercado que mantuvo posicion '
        'en un ano de movimiento."',
        size=FONT_BODY, italic=True, color=DARK_TEXT, before=0, after=4)
    bullet(tf1, "TOM Gama 44.3% — solido y estable en un mercado dinamico.")
    bullet(tf1, "Activo relacional (OR=5.73) existe y es robusto.")
    bullet(tf1, "Posicion en C+/C (TOM 60.6%) — la mas fuerte del mercado.")
    bullet(tf1, "El 24h: ventaja funcional sin equivalente en Rio ni Paramo.")
    badge_text(tf1, "    ✅ Todos los datos de soporte son sig_99 o alta certeza cuanti.")

    hdr2 = add_rounded_rect(slide, col2_l, BODY_TOP, col_w, 365760, ORANGE_BG)
    tf_h2 = hdr2.text_frame; clear_tf(tf_h2)
    add_para_with_runs(tf_h2, [("CO-2 Opcion B — Senal de alerta competitiva",
                                FONT_BODY, True, ORANGE_WARN)], align=PP_ALIGN.CENTER)

    col2_body_top = BODY_TOP + 365760 + 91440
    tb2 = add_textbox(slide, col2_l, col2_body_top, col_w, 3291840)
    tf2 = tb2.text_frame; tf2.word_wrap = True; clear_tf(tf2)
    plain_text(tf2,
        '"Rio +17pp TOM (sig_99). Paramo +12pp TOM (sig_99). '
        'Posiblemente +25pp y +22pp en C+/C — el segmento de Gama."',
        size=FONT_BODY, italic=True, color=DARK_TEXT, before=0, after=4)
    bullet(tf2, "Conversion TOM→Preferida = 44.3% → 8.0%.")
    bullet(tf2, "El 92% que no prefiere no esta siendo activado.")
    bullet(tf2, "El anzuelo comunicacional esta en el territorio equivocado (OR=1.03 NS).")
    bullet(tf2, "La ventana de respuesta ofensiva en C+/C se cierra si Rio/Paramo consolidan el segmento.")
    badge_text(tf2, "    ⚠ C+/C: hipotesis WoW pendiente confirmacion ola 2027 · CV-WOW-005")

    add_info_box(slide, MARGIN_L, 5669280, CONTENT_W, 320040,
        "CO-1: senal Rio-proteinas — excluir del deck (1 verbatim, certeza muy baja). "
        "CO-3: recomendaciones canal digital — Owner decide si van al deck V5 o a propuesta de scope adicional.")

    add_footer(slide, 6)


def build_r7_decisiones(prs):
    slide = blank_slide(prs)
    add_standard_header(slide, 7,
        "Decisiones Owner antes de presentar a Cora y a la Junta de Gama")

    add_callout_box(slide,
        "Tres decisiones estrategicas solo el Owner puede tomar antes de que el deck V5 "
        "salga a Cora y a la Junta. El gate Bruna BR-2 V4 sigue vigente en V5 "
        "(mismos datos, mismos caveats, mismos wordings DW-1..DW-5).")

    tf = body_textbox(slide)

    section_divider(tf, "Decisiones Owner (CO-1 / CO-2 / CO-3):", before=0)
    bullet(tf, "CO-1 — Senal de alerta Rio en proteinas: ",
           value="1 verbatim FG, certeza muy baja. "
                 "Bruna recomienda excluir del deck — trasladar a memo interno hasta ola 2027. "
                 "Owner define el umbral de evidencia para senales de alerta en presentacion ejecutiva.")
    bullet(tf, "CO-2 — Tono del deck: ",
           value="Opcion A (liderazgo defensivo solido) vs Opcion B (senal de alerta competitiva). "
                 "Las dos son verdaderas. El tono afecta como Gama recibe el reporte. "
                 "Bruna recomienda mezcla equilibrada — pero Owner decide.")
    bullet(tf, "CO-3 — Alcance recomendaciones digitales: ",
           value="Gama Club + canal digital van mas alla del brief original de notoriedad/brand health. "
                 "Owner decide si entran al deck V5 o se reservan para propuesta de scope adicional.")

    section_divider(tf, "Gate Bruna BR-2: vigente en V5")
    bullet(tf, "DW-1: verbatim Azahara Betancourt con nombre + caveat literal obligatorio. APLICADO.")
    bullet(tf, "DW-2: personificacion femenina como recomendacion condicional + caveat DW-2. APLICADO.")
    bullet(tf, "DW-3: wording suavizado para Seg_2 ('cuali V4 enriquece'). APLICADO.")
    bullet(tf, "DW-4: verbatim inventario digital con enmarcado de oportunidad. APLICADO.")
    bullet(tf, "DW-5: etiqueta [Hipotesis V4 — evidencia convergente] sin explicar cambio de status. APLICADO.")

    section_divider(tf, "Agenda ola 2027 (5 mejoras metodologicas obligatorias):")
    meth = [
        "MaxDiff reemplaza Likert saturado (P22) — eliminar saturacion T2B >90%.",
        "NPS + switching explicito (3 preguntas) — medir fidelidad activa.",
        "CEPs expandidos (15-20 ocasiones vs 5 misiones genericas).",
        "Penetracion 12m + frecuencia + ticket promedio — modelo share of wallet.",
        "Booster Pref-Gama n=80 + modulo comunicacion con estimulos + Rio como benchmark.",
    ]
    for m in meth:
        bullet(tf, m, label_bold=False, indent=True, before=1, after=0)

    plain_text(tf,
        "    Hipotesis abiertas para ola 2027: imagen experiencial en ascenso · "
        "gap comunicacional PTL · arquetipo femenino accionable · WoW en NSE C+/C.",
        size=FONT_SMALL, color=GRAY_MID, before=3, after=1)

    add_warn_box(slide,
        "El deck V5 no puede circular a Cora/Gama hasta que CO-1/CO-2/CO-3 esten "
        "resueltos por Owner. El render del .pptx se entrega en paralelo para agilizar "
        "— las CO son decisiones de contenido, no de estructura.")

    add_footer(slide, 7)


def build_r8_cierre(prs):
    slide = blank_slide(prs)
    add_standard_header(slide, 8,
        "Gama 2026: activo real, posicion solida, activacion pendiente")

    add_callout_box(slide,
        "El estudio entrega tres certezas de alta evidencia y un mapa de hipotesis accionables "
        "para la siguiente ola. La ventana de respuesta ofensiva en C+/C esta abierta hoy.")

    tf = body_textbox(slide)

    section_divider(tf, "Lo que el estudio V5 cierra con alta certeza:", before=0)
    bullet(tf, "El activo de atencion (OR=5.73***) es real, robusto y sin equivalente competitivo documentado.")
    bullet(tf, "El precio no predice preferencia (OR=1.03 NS). La campana PTL invierte en el territorio equivocado.")
    bullet(tf, "El Seg 2 (33% del mercado) tiene 0% preferencia actual y mayor potencial de conversion C/P.")
    bullet(tf, "Gama es la unica gran cadena con posicion estadisticamente estable en 2026 (0/8 sig).")
    badge_text(tf, "    ✅ Conclusiones de alta certeza cuanti · n=402 · AUC=0.929 · convergencia 4 metodos")

    section_divider(tf, "Lo que requiere ola 2027 para confirmar:")
    bullet(tf, "WoW en NSE C+/C (Rio +25.8pp, Paramo +22.3pp — exploratorio, CV-WOW-005).")
    bullet(tf, "Gap comunicacional PTL: modulo con estimulos de campana + benchmark Rio.")
    bullet(tf, "Arquetipo femenino: concept testing cuantitativo con target antes de cualquier ejecucion creativa.")
    bullet(tf, "Imagen experiencial en ascenso (+6.3pp 'atractiva', +5.8pp 'seguro' — solo tendencia p_adj 0.058-0.089).")
    badge_text(tf, "    ⚠ Hipotesis apoyadas · pendiente confirmacion · 5 mejoras metodologicas disenadas para ola 2027")

    section_divider(tf, "Las dos preguntas que el Owner responde antes de la Junta:")
    bullet(tf, "CO-2: liderazgo defensivo o senal de alerta? (Ambas verdaderas — el tono es estrategico.)")
    bullet(tf, "CO-3: Gama Club + canal digital en el deck o en propuesta de scope adicional?")

    plain_text(tf,
        "    El equipo analitico (Cora Urrea + Raoul Bermudez) queda disponible para "
        "preguntas, calibraciones de tono y diseno del brief de ola 2027.",
        size=FONT_BODY, italic=True, color=DARK_TEXT, before=5, after=1)

    plain_text(tf,
        "    Notoriedad y Preferencia de Marca — Gama 2026 · V5 · n=402 · "
        "m.e. ±4.89% · Confidencial NDA · Cora Urrea + Raoul Bermudez · 2026-05-18",
        size=FONT_SMALL, color=GRAY_MID, before=3, after=1)

    add_footer(slide, 8)


def main():
    print("Generando charts PNG...")
    chart_r1 = make_chart_r1_cpc_vs_total()
    print(f"  Chart R-1 guardado: {chart_r1}")
    chart_r2 = make_chart_r2_wow_movement()
    print(f"  Chart R-2 guardado: {chart_r2}")
    chart_r3 = make_chart_r3_forest_plot()
    print(f"  Chart R-3 guardado: {chart_r3}")

    print("Construyendo presentacion...")
    prs = new_prs()

    build_r1_portada(prs)
    print("  R1 Portada OK")
    build_r2_tesis(prs)
    print("  R2 Tesis OK")
    build_r3_cpc(prs, chart_r1, chart_r2)
    print("  R3 C+/C OK")
    build_r4_hallazgos(prs, chart_r3)
    print("  R4 Hallazgos OK")
    build_r5_recomendaciones(prs)
    print("  R5 Recomendaciones OK")
    build_r6_tono_dual(prs)
    print("  R6 Tono dual OK")
    build_r7_decisiones(prs)
    print("  R7 Decisiones OK")
    build_r8_cierre(prs)
    print("  R8 Cierre OK")

    prs.save(OUTPUT_PATH)
    print(f"\nGuardado: {OUTPUT_PATH}")
    print(f"Charts en: {CHARTS_DIR}")


if __name__ == "__main__":
    main()
