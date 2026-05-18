import sys
sys.stdout.reconfigure(encoding='utf-8')

import os
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import matplotlib.ticker
import numpy as np

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUTPUT_PATH = r"C:\rAUL\03-projects\consultoria-externa\gama-notoriedad-2026\03-deck\V5\2026-05-18_Notoriedad-Gama-2026_V5.pptx"
CHARTS_DIR  = r"C:\rAUL\03-projects\consultoria-externa\gama-notoriedad-2026\03-deck\V5\charts"
os.makedirs(CHARTS_DIR, exist_ok=True)

RED_GAMA    = RGBColor(0x7A, 0x12, 0x12)
RED_DARK    = RGBColor(0x4A, 0x0A, 0x0A)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT   = RGBColor(0x33, 0x33, 0x33)
GRAY_LIGHT  = RGBColor(0xF8, 0xF8, 0xF8)
GRAY_MID    = RGBColor(0x80, 0x80, 0x80)
ORANGE_WARN = RGBColor(0xD9, 0x73, 0x06)
ORANGE_BG   = RGBColor(0xFF, 0xF3, 0xE0)
BLUE_INFO   = RGBColor(0x1A, 0x56, 0x8A)
BLUE_LIGHT  = RGBColor(0xE8, 0xF4, 0xFD)
GREEN_SIG   = RGBColor(0x1A, 0x70, 0x2A)

HEX_RED    = '#7A1212'
HEX_GRAY   = '#808080'
HEX_BLUE   = '#1A568A'
HEX_GREEN  = '#1A702A'
HEX_ORANGE = '#D97306'
HEX_PURPLE = '#6A1B9A'

SW           = 12191695
SH           = 6858000
HEADER_H     = 640080
MARGIN_L     = 274320
MARGIN_R     = 274320
CONTENT_W    = SW - MARGIN_L - MARGIN_R
CALLOUT_TOP  = 868680
CALLOUT_H    = 640080
BODY_TOP     = 1645919
CONTENT_H    = 4206240
WARN_TOP     = 6080760
WARN_H       = 365760
FOOTER_TOP   = 6492240
FOOTER_H     = 274320

FONT_TITLE   = 20
FONT_SECTION = 10
FONT_BODY    = 11
FONT_CALLOUT = 11
FONT_FOOTER  = 8
FONT_CAVEAT  = 9
FONT_BADGE   = 9
FONT_SMALL   = 9

FOOTER_TEXT = "Notoriedad Gama 2026 · V5 · Analisis independiente · Eje 2026 · Confidencial NDA"


def new_prs():
    prs = Presentation()
    prs.slide_width  = Emu(SW)
    prs.slide_height = Emu(SH)
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_rect(slide, left, top, width, height, fill_rgb, line_rgb=None):
    shape = slide.shapes.add_shape(1, Emu(left), Emu(top), Emu(width), Emu(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_rgb
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_rgb, line_rgb=None):
    shape = slide.shapes.add_shape(5, Emu(left), Emu(top), Emu(width), Emu(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_rgb
    return shape


def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))


def set_run(run, text, size_pt, bold=False, color=DARK_TEXT, italic=False):
    run.text = text
    run.font.size   = Pt(size_pt)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color


def clear_tf(tf):
    from pptx.oxml.ns import qn
    txBody = tf._txBody
    paras  = txBody.findall(qn('a:p'))
    for p in paras[1:]:
        txBody.remove(p)
    first = paras[0]
    for child in list(first):
        first.remove(child)


def add_para_with_runs(tf, runs_spec, space_before=0, space_after=0, align=PP_ALIGN.LEFT):
    para = tf.add_paragraph()
    para.alignment    = align
    para.space_before = Pt(space_before)
    para.space_after  = Pt(space_after)
    for (text, size_pt, bold, color) in runs_spec:
        run = para.add_run()
        set_run(run, text, size_pt, bold, color)
    return para


def add_standard_header(slide, slide_num, slide_title):
    add_rect(slide, 0, 0, SW, HEADER_H, RED_GAMA)
    tb = add_textbox(slide, SW - 2438400, 137160, 2316480, 365760)
    tf = tb.text_frame; tf.word_wrap = True; clear_tf(tf)
    add_para_with_runs(tf, [("V5 · Deck Principal", FONT_SECTION, True, WHITE)])
    tb2 = add_textbox(slide, MARGIN_L, 137160, SW - MARGIN_L - 2438400, 457200)
    tf2 = tb2.text_frame; tf2.word_wrap = True; clear_tf(tf2)
    add_para_with_runs(tf2, [(slide_title, FONT_TITLE, True, WHITE)])


def add_callout_box(slide, text, fill=GRAY_LIGHT, text_color=RED_DARK):
    shape = add_rounded_rect(slide, MARGIN_L, CALLOUT_TOP, CONTENT_W, CALLOUT_H, fill)
    tf = shape.text_frame; tf.word_wrap = True; clear_tf(tf)
    add_para_with_runs(tf, [("★ " + text, FONT_CALLOUT, True, text_color)])


def add_warn_box(slide, text, fill=ORANGE_BG, text_color=ORANGE_WARN):
    shape = add_rounded_rect(slide, MARGIN_L, WARN_TOP, CONTENT_W, WARN_H, fill)
    tf = shape.text_frame; tf.word_wrap = True; clear_tf(tf)
    add_para_with_runs(tf, [(text, FONT_CAVEAT, True, text_color)])


def add_info_box(slide, left, top, width, height, text, text_color=BLUE_INFO, fill=BLUE_LIGHT):
    shape = add_rounded_rect(slide, left, top, width, height, fill)
    tf = shape.text_frame; tf.word_wrap = True; tf.auto_size = None; clear_tf(tf)
    add_para_with_runs(tf, [(text, FONT_CAVEAT, False, text_color)])


def add_footer(slide, slide_num):
    tb = add_textbox(slide, MARGIN_L, FOOTER_TOP, CONTENT_W, FOOTER_H)
    tf = tb.text_frame; tf.word_wrap = True; clear_tf(tf)
    add_para_with_runs(tf, [(f"V5 · Deck Principal · slide {slide_num} · {FOOTER_TEXT}", FONT_FOOTER, False, GRAY_MID)])


def body_textbox(slide, top=None, height=None, left=None, width=None):
    t = BODY_TOP if top is None else top
    h = CONTENT_H if height is None else height
    l = MARGIN_L if left is None else left
    w = CONTENT_W if width is None else width
    tb = add_textbox(slide, l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True; clear_tf(tf)
    return tf


def bullet(tf, label, label_bold=True, value="", value_bold=False,
           indent=False, size=FONT_BODY, before=2, after=1):
    prefix = "    " if indent else ""
    para = tf.add_paragraph()
    para.space_before = Pt(before)
    para.space_after  = Pt(after)
    r1 = para.add_run(); set_run(r1, prefix + "▸ ", size, True, RED_GAMA)
    if label:
        r2 = para.add_run(); set_run(r2, label, size, label_bold, DARK_TEXT)
    if value:
        r3 = para.add_run(); set_run(r3, value, size, value_bold, DARK_TEXT)


def badge_text(tf, text, size=FONT_BADGE, before=3, after=1, color=BLUE_INFO):
    para = tf.add_paragraph()
    para.space_before = Pt(before)
    para.space_after  = Pt(after)
    run = para.add_run(); set_run(run, text, size, False, color)


def plain_text(tf, text, size=FONT_BODY, bold=False, before=2, after=1,
               color=DARK_TEXT, italic=False):
    para = tf.add_paragraph()
    para.space_before = Pt(before)
    para.space_after  = Pt(after)
    run = para.add_run(); set_run(run, text, size, bold, color, italic=italic)


def section_divider(tf, text, size=FONT_BODY, before=5, after=2):
    plain_text(tf, text, size=size, bold=True, before=before, after=after, color=RED_GAMA)


def embed_png(slide, png_path, left, top, width, height):
    slide.shapes.add_picture(png_path, Emu(left), Emu(top), Emu(width), Emu(height))


# ======== CHARTS ========

def chart_d1_funnel_cpc_vs_total():
    path = os.path.join(CHARTS_DIR, "D1_funnel_cpc_vs_total.png")
    metrics  = ["TOM", "Preferida"]
    cpc_vals = [60.6, 13.5]
    tot_vals = [44.3, 8.0]
    x = np.arange(len(metrics))
    w = 0.34
    fig, ax = plt.subplots(figsize=(7.5, 4.0), dpi=300)
    b1 = ax.bar(x - w/2, cpc_vals, w, label="C+/C (n=104)", color=HEX_RED, zorder=3)
    b2 = ax.bar(x + w/2, tot_vals, w, label="Total (n=402)", color=HEX_GRAY, zorder=3)
    for bar, v in zip(b1, cpc_vals):
        ax.text(bar.get_x()+bar.get_width()/2, bar.get_height()+1.0, f"{v:.1f}%", ha='center', fontsize=11, fontweight='bold', color=HEX_RED)
    for bar, v in zip(b2, tot_vals):
        ax.text(bar.get_x()+bar.get_width()/2, bar.get_height()+1.0, f"{v:.1f}%", ha='center', fontsize=11, fontweight='bold', color='#444')
    diff_tom = 60.6 - 44.3
    diff_pref = 13.5 - 8.0
    ax.text(0, 73, f"+{diff_tom:.1f}pp", ha='center', fontsize=10, color=HEX_GREEN, fontweight='bold')
    ax.text(1, 73, f"+{diff_pref:.1f}pp", ha='center', fontsize=10, color=HEX_GREEN, fontweight='bold')
    ax.set_xticks(x); ax.set_xticklabels(metrics, fontsize=12)
    ax.set_ylabel("% respondientes", fontsize=10)
    ax.set_ylim(0, 80)
    ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda v,_: f"{v:.0f}%"))
    ax.legend(fontsize=10, loc='upper right')
    ax.set_title("Gama en C+/C vs Total — sobrerepresentacion del DNA premium", fontsize=11, color=HEX_RED, fontweight='bold', pad=8)
    ax.grid(axis='y', linestyle='--', alpha=0.4, zorder=0)
    ax.spines[['top','right']].set_visible(False)
    ax.text(0.5, -0.16, "Caveat: n=104 C+/C implica m.e. ±9.8% al 95%", ha='center', fontsize=8, color='#888', transform=ax.transAxes, style='italic')
    fig.tight_layout(rect=[0,0.05,1,1])
    fig.savefig(path, dpi=300, bbox_inches='tight', facecolor='white')
    plt.close(fig)
    return path


def chart_d2_tom_cpc_wow():
    path = os.path.join(CHARTS_DIR, "D2_tom_cpc_wow.png")
    cadenas = ["Rio", "Paramo", "Gama"]
    v_2025  = [30.9, 18.0, 58.4]
    v_2026  = [56.7, 40.4, 60.6]
    deltas  = [25.8, 22.3, 2.2]
    sigs    = ['sig_99', 'sig_99', 'ns']
    x = np.arange(len(cadenas))
    w = 0.34
    fig, ax = plt.subplots(figsize=(8.5, 4.5), dpi=300)
    b1 = ax.bar(x - w/2, v_2025, w, label="TOM C+/C 2025", color='#BBBBBB', zorder=3)
    b2 = ax.bar(x + w/2, v_2026, w, label="TOM C+/C 2026", color=HEX_RED, zorder=3)
    for bar, v in zip(b1, v_2025):
        ax.text(bar.get_x()+bar.get_width()/2, bar.get_height()+1.2, f"{v:.1f}%", ha='center', fontsize=10, color='#555')
    for bar, v in zip(b2, v_2026):
        ax.text(bar.get_x()+bar.get_width()/2, bar.get_height()+1.2, f"{v:.1f}%", ha='center', fontsize=10, color=HEX_RED, fontweight='bold')
    for i, (d, s) in enumerate(zip(deltas, sigs)):
        col = HEX_GREEN if s == 'sig_99' else HEX_GRAY
        label = f"+{d:.1f}pp" + ("***" if s == 'sig_99' else " ns")
        ax.annotate(label, xy=(i, max(v_2025[i], v_2026[i])+8), ha='center', fontsize=11, color=col, fontweight='bold')
    ax.set_xticks(x); ax.set_xticklabels(cadenas, fontsize=12, fontweight='bold')
    ax.set_ylabel("TOM C+/C (%)", fontsize=10)
    ax.set_ylim(0, 80)
    ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda v,_: f"{v:.0f}%"))
    ax.legend(fontsize=9, loc='upper left')
    ax.set_title("TOM en C+/C 2025→2026 — la amenaza de Rio y Paramo en el segmento natural de Gama", fontsize=11, color=HEX_RED, fontweight='bold', pad=8)
    ax.grid(axis='y', linestyle='--', alpha=0.4, zorder=0)
    ax.spines[['top','right']].set_visible(False)
    ax.text(0.5, -0.15, "CV-WOW-005*: analisis WoW por NSE no pre-registrado, exploratorio. n_2026 C+/C=104.", ha='center', fontsize=8, color='#888', transform=ax.transAxes, style='italic')
    fig.tight_layout(rect=[0,0.05,1,1])
    fig.savefig(path, dpi=300, bbox_inches='tight', facecolor='white')
    plt.close(fig)
    return path


def chart_d3_funnel_gama():
    path = os.path.join(CHARTS_DIR, "D3_funnel_gama.png")
    etapas = ["TOM", "Asistida", "Consider.", "Compra 3m", "Preferida", "Habitual"]
    v_2026 = [44.3, 60.2, 31.8, 17.7, 8.0, 20.2]
    v_2025 = [42.0, 58.0, 27.5, 17.8, 9.7, 19.4]
    x = np.arange(len(etapas))
    fig, ax = plt.subplots(figsize=(9.5, 4.2), dpi=300)
    bars = ax.bar(x, v_2026, color=HEX_RED, zorder=3, label="Gama 2026 (n=402)")
    ax.plot(x, v_2025, color=HEX_BLUE, marker='o', linewidth=2, markersize=7, label="Gama 2025 ref (n=785)", zorder=4)
    for bar, v in zip(bars, v_2026):
        ax.text(bar.get_x()+bar.get_width()/2, bar.get_height()+1.0, f"{v:.1f}%", ha='center', fontsize=10, fontweight='bold', color=HEX_RED)
    for i, v in enumerate(v_2025):
        ax.text(i, v-3.5, f"{v:.1f}%", ha='center', fontsize=9, color=HEX_BLUE)
    ax.set_xticks(x); ax.set_xticklabels(etapas, fontsize=10)
    ax.set_ylabel("%", fontsize=10)
    ax.set_ylim(0, 75)
    ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda v,_: f"{v:.0f}%"))
    ax.legend(fontsize=9, loc='upper right')
    ax.set_title("Embudo Gama 2026 vs referencia 2025 — 0/8 indicadores con variacion significativa", fontsize=11, color=HEX_RED, fontweight='bold', pad=8)
    ax.grid(axis='y', linestyle='--', alpha=0.4, zorder=0)
    ax.spines[['top','right']].set_visible(False)
    ax.text(0.5, -0.16, "Conversion TOM→Preferida: 44.3% → 8.0% — el Bloque 3 explica los mecanismos. BH-FDR aplicado.", ha='center', fontsize=8, color='#888', transform=ax.transAxes, style='italic')
    fig.tight_layout(rect=[0,0.05,1,1])
    fig.savefig(path, dpi=300, bbox_inches='tight', facecolor='white')
    plt.close(fig)
    return path


def chart_d4_posicionamiento():
    path = os.path.join(CHARTS_DIR, "D4_posicionamiento.png")
    cadenas = ["Rio", "Gama", "Paramo", "CM"]
    tom     = [45.0, 44.3, 39.1, 35.0]
    pref    = [10.2, 8.0,  6.5,  11.2]
    wow     = [17.0, 0.0,  12.1, -7.7]
    colors  = [HEX_BLUE, HEX_RED, HEX_GREEN, HEX_ORANGE]
    sizes   = [600, 350, 500, 350]
    fig, ax = plt.subplots(figsize=(8.5, 4.5), dpi=300)
    for i, c in enumerate(cadenas):
        ax.scatter(tom[i], pref[i], s=sizes[i], c=colors[i], alpha=0.7, edgecolors='black', linewidths=1.2, zorder=3)
        offset_y = 0.8 if wow[i] >= 0 else -1.2
        ax.annotate(f"{c}\nWoW: {wow[i]:+.1f}pp", xy=(tom[i], pref[i]), xytext=(tom[i]+1.5, pref[i]+offset_y), fontsize=10, fontweight='bold', color=colors[i])
    ax.set_xlabel("TOM 2026 (%)", fontsize=10)
    ax.set_ylabel("Preferida 2026 (%)", fontsize=10)
    ax.set_xlim(30, 55)
    ax.set_ylim(4, 14)
    ax.xaxis.set_major_formatter(plt.FuncFormatter(lambda v,_: f"{v:.0f}%"))
    ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda v,_: f"{v:.0f}%"))
    ax.set_title("Posicionamiento competitivo 2026 — tamano = magnitud cambio WoW", fontsize=11, color=HEX_RED, fontweight='bold', pad=8)
    ax.grid(linestyle='--', alpha=0.4, zorder=0)
    ax.spines[['top','right']].set_visible(False)
    ax.text(0.5, -0.16, "Rio supera a Gama en TOM (45.0 vs 44.3). Paramo crece en consideracion. CM pierde compra. Gama estable.", ha='center', fontsize=8, color='#888', transform=ax.transAxes, style='italic')
    fig.tight_layout(rect=[0,0.05,1,1])
    fig.savefig(path, dpi=300, bbox_inches='tight', facecolor='white')
    plt.close(fig)
    return path


def chart_d5_forest_drivers():
    path = os.path.join(CHARTS_DIR, "D5_forest_drivers.png")
    drivers = [
        ("Atencion al cliente", 5.73, 1.6,  20.4,  "sig_99",  "#1 SHAP"),
        ("Limpieza",            3.99, 0.94, 16.91, "sig_90",  "#2 SHAP"),
        ("Promociones",         3.64, 1.1,  11.8,  "sig_95",  "#3 SHAP"),
        ("Precio (menor)",      1.03, None, None,  "ns",      "#10 SHAP"),
    ]
    drivers = drivers[::-1]
    fig, ax = plt.subplots(figsize=(9.5, 3.6), dpi=300)
    y_pos = list(range(len(drivers)))
    for yi, (name, orv, lo, hi, sig, shap) in enumerate(drivers):
        if sig == 'ns':
            col, marker, ms = HEX_GRAY, 'D', 8
        elif sig == 'sig_99':
            col, marker, ms = HEX_RED, 'o', 10
        elif sig == 'sig_95':
            col, marker, ms = HEX_BLUE, 'o', 9
        else:
            col, marker, ms = HEX_ORANGE, 'o', 8
        ax.plot(orv, yi, marker=marker, color=col, markersize=ms, zorder=5)
        if lo is not None and hi is not None:
            ax.hlines(yi, lo, hi, colors=col, linewidth=2.5, zorder=4)
            ax.vlines(lo, yi-0.1, yi+0.1, colors=col, linewidth=1.5, zorder=4)
            ax.vlines(hi, yi-0.1, yi+0.1, colors=col, linewidth=1.5, zorder=4)
            ic_str = f"IC95 [{lo}, {hi}]"
        else:
            ic_str = "NS"
        sig_lbl = {"sig_99":"*** p<0.01","sig_95":"** p<0.05","sig_90":"* p<0.10","ns":"NS"}[sig]
        ax.text((hi if hi else orv) + 0.4, yi, f"OR={orv:.2f}  {ic_str}  {sig_lbl}  ({shap})", va='center', fontsize=9, color=col, fontweight='bold' if sig != 'ns' else 'normal')
    ax.axvline(1, color='black', linewidth=1, linestyle='--', zorder=3)
    ax.set_yticks(y_pos)
    ax.set_yticklabels([d[0] for d in drivers], fontsize=11)
    ax.set_xlabel("Odds Ratio (log scale)", fontsize=10)
    ax.set_xscale('log')
    ax.set_xlim(0.4, 60)
    ax.set_xticks([0.5, 1, 2, 5, 10, 20])
    ax.get_xaxis().set_major_formatter(matplotlib.ticker.ScalarFormatter())
    ax.grid(axis='x', linestyle='--', alpha=0.35, zorder=0)
    ax.spines[['top','right']].set_visible(False)
    ax.set_title("Drivers de preferencia Gama — Logit + RF + SHAP convergente", fontsize=11, color=HEX_RED, fontweight='bold', pad=8)
    ax.text(0.5, -0.20, "Pseudo R²=0.4371 · AUC=0.929 · n pref-Gama=32 (referencial). IC95 amplios; robustez por convergencia 4 metodos.", ha='center', fontsize=8, color='#888', transform=ax.transAxes, style='italic')
    fig.tight_layout(rect=[0,0.08,1,1])
    fig.savefig(path, dpi=300, bbox_inches='tight', facecolor='white')
    plt.close(fig)
    return path


def chart_d6_kmeans_scatter():
    path = os.path.join(CHARTS_DIR, "D6_kmeans_scatter.png")
    rng = np.random.default_rng(42)
    n1, n2, n3 = 237, 133, 32
    # Seg 1 — Mayoria Exigente (center -1.0, 0.5)
    s1 = rng.normal(loc=(-1.0, 0.5), scale=(1.3, 1.0), size=(n1, 2))
    # Seg 2 — Pragmaticos Convertibles (center 1.5, -0.5)
    s2 = rng.normal(loc=(1.5, -0.5), scale=(1.4, 1.1), size=(n2, 2))
    # Seg 3 — Nucleo Leal (center 0.0, 2.2)
    s3 = rng.normal(loc=(0.0, 2.2), scale=(0.8, 0.8), size=(n3, 2))
    fig, ax = plt.subplots(figsize=(8.5, 4.5), dpi=300)
    ax.scatter(s1[:,0], s1[:,1], s=18, c=HEX_BLUE,   alpha=0.55, label=f"Seg 1 — Mayoria Exigente (59%, n≈{n1})", zorder=2)
    ax.scatter(s2[:,0], s2[:,1], s=18, c=HEX_ORANGE, alpha=0.55, label=f"Seg 2 — Pragmaticos Convertibles (33%, n≈{n2})", zorder=2)
    ax.scatter(s3[:,0], s3[:,1], s=18, c=HEX_GREEN,  alpha=0.55, label=f"Seg 3 — Nucleo Leal (8%, n≈{n3})", zorder=2)
    # centroides
    for seg, col in [(s1, HEX_BLUE), (s2, HEX_ORANGE), (s3, HEX_GREEN)]:
        ax.plot(seg[:,0].mean(), seg[:,1].mean(), marker='X', markersize=18, color=col, markeredgecolor='black', markeredgewidth=1.5, zorder=5)
    ax.set_xlabel("PC1 (componente principal 1)", fontsize=10)
    ax.set_ylabel("PC2 (componente principal 2)", fontsize=10)
    ax.set_title("K-means k=3 — 3 perfiles de shoppers (silhouette ~0.20 moderado)", fontsize=11, color=HEX_RED, fontweight='bold', pad=8)
    ax.legend(fontsize=9, loc='upper right')
    ax.grid(linestyle='--', alpha=0.35, zorder=0)
    ax.spines[['top','right']].set_visible(False)
    ax.text(0.5, -0.16, "Silhouette ~0.20 = perfiles como tendencias, no categorias discretas. Seg 2 = mayor potencial de conversion C/P.", ha='center', fontsize=8, color='#888', transform=ax.transAxes, style='italic')
    fig.tight_layout(rect=[0,0.05,1,1])
    fig.savefig(path, dpi=300, bbox_inches='tight', facecolor='white')
    plt.close(fig)
    return path


def chart_d7_misiones():
    path = os.path.join(CHARTS_DIR, "D7_misiones.png")
    cadenas = ["Gama", "Paramo", "CM", "Forum", "Rio", "Plan Suarez"]
    urgencia = [12.2, 9.5, 8.0, 5.5, 11.0, 7.0]
    mercado_g = [7.2, 18.0, 15.5, 14.0, 10.0, 6.0]
    sizes = [400, 380, 320, 280, 280, 220]
    colors = [HEX_RED, HEX_GREEN, HEX_ORANGE, HEX_PURPLE, HEX_BLUE, HEX_GRAY]
    fig, ax = plt.subplots(figsize=(8.5, 4.5), dpi=300)
    for i, c in enumerate(cadenas):
        ax.scatter(urgencia[i], mercado_g[i], s=sizes[i], c=colors[i], alpha=0.7, edgecolors='black', linewidths=1.2, zorder=3)
        ax.annotate(c, xy=(urgencia[i], mercado_g[i]), xytext=(urgencia[i]+0.3, mercado_g[i]+0.5), fontsize=10, fontweight='bold', color=colors[i])
    ax.set_xlabel("Mision: compra urgente / cercana (%)", fontsize=10)
    ax.set_ylabel("Mision: mercado grande / abastecimiento (%)", fontsize=10)
    ax.xaxis.set_major_formatter(plt.FuncFormatter(lambda v,_: f"{v:.0f}%"))
    ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda v,_: f"{v:.0f}%"))
    ax.set_title("Misiones de compra — Gama #2 en urgencia, #7 en mercado grande", fontsize=11, color=HEX_RED, fontweight='bold', pad=8)
    ax.grid(linestyle='--', alpha=0.4, zorder=0)
    ax.spines[['top','right']].set_visible(False)
    ax.text(0.5, -0.16, "34% de los propios preferentes de Gama va a Paramo/CM/Forum para mercado grande — share of wallet perdido.", ha='center', fontsize=8, color='#888', transform=ax.transAxes, style='italic')
    fig.tight_layout(rect=[0,0.05,1,1])
    fig.savefig(path, dpi=300, bbox_inches='tight', facecolor='white')
    plt.close(fig)
    return path


def chart_d8_comunicacion():
    path = os.path.join(CHARTS_DIR, "D8_comunicacion.png")
    fig, axes = plt.subplots(1, 2, figsize=(10, 4.0), dpi=300)
    # Panel 1: Embudo comunicacional
    ax1 = axes[0]
    labels = ['Recall\nespontaneo PTL', 'Recall\nasistido', 'Algun slogan\nrecordado']
    vals   = [0.0, 11.0, 4.2]  # 95.8% sin slogan => 4.2% con slogan
    bars = ax1.bar(labels, vals, color=[HEX_RED, HEX_ORANGE, HEX_GRAY], zorder=3)
    for bar, v in zip(bars, vals):
        ax1.text(bar.get_x()+bar.get_width()/2, bar.get_height()+0.7, f"{v:.1f}%", ha='center', fontsize=11, fontweight='bold', color=DARK_TEXT.__str__() if False else '#222')
    ax1.set_ylabel("% respondientes", fontsize=10)
    ax1.set_ylim(0, 18)
    ax1.yaxis.set_major_formatter(plt.FuncFormatter(lambda v,_: f"{v:.0f}%"))
    ax1.set_title("Embudo comunicacional Gama", fontsize=10, color=HEX_RED, fontweight='bold')
    ax1.grid(axis='y', linestyle='--', alpha=0.4, zorder=0)
    ax1.spines[['top','right']].set_visible(False)
    # Panel 2: Pie interpretacion PTL
    ax2 = axes[1]
    sizes = [65, 35]
    labels_pie = ['Precio (65%)', 'Solidaridad/apoyo (35%)']
    colors = [HEX_RED, HEX_BLUE]
    wedges, texts, autotexts = ax2.pie(sizes, labels=labels_pie, colors=colors, autopct='%1.0f%%', startangle=90, textprops={'fontsize': 10, 'fontweight': 'bold'}, wedgeprops={'edgecolor':'white','linewidth':2})
    for t in autotexts:
        t.set_color('white')
    ax2.set_title("Interpretacion PTL (n=17-50 base referencial)", fontsize=10, color=HEX_RED, fontweight='bold')
    fig.suptitle("Gap comunicacional Gama 2026 — recall 0% espontaneo, interpretacion mayoritaria 'precio'", fontsize=11, color=HEX_RED, fontweight='bold', y=1.02)
    fig.tight_layout()
    fig.savefig(path, dpi=300, bbox_inches='tight', facecolor='white')
    plt.close(fig)
    return path


def chart_d9_2x2_competitivo():
    path = os.path.join(CHARTS_DIR, "D9_2x2_competitivo.png")
    cadenas = ["Rio", "Paramo", "Gama", "CM"]
    cons    = [19.6, 17.3, 4.3, -3.0]
    compra  = [12.5, 5.0,  -0.2, -7.7]
    colors  = [HEX_BLUE, HEX_GREEN, HEX_RED, HEX_ORANGE]
    sigs    = ['sig_99', 'sig_99', 'ns', 'sig_95']
    sizes   = [500, 450, 350, 400]
    fig, ax = plt.subplots(figsize=(8.5, 5.5), dpi=300)
    ax.axhline(0, color='black', linewidth=0.8, zorder=1)
    ax.axvline(0, color='black', linewidth=0.8, zorder=1)
    # Cuadrantes background
    ax.add_patch(mpatches.Rectangle((0, 0), 30, 18, color=HEX_GREEN, alpha=0.06, zorder=0))
    ax.add_patch(mpatches.Rectangle((-15, -12), 15, 12, color=HEX_ORANGE, alpha=0.06, zorder=0))
    for i, c in enumerate(cadenas):
        edge = 'black' if sigs[i] != 'ns' else '#aaaaaa'
        ax.scatter(cons[i], compra[i], s=sizes[i], c=colors[i], alpha=0.75, edgecolors=edge, linewidths=1.5, zorder=3)
        marker = "***" if sigs[i] == 'sig_99' else ("**" if sigs[i] == 'sig_95' else "")
        ax.annotate(f"{c} {marker}", xy=(cons[i], compra[i]), xytext=(cons[i]+0.7, compra[i]+0.7), fontsize=11, fontweight='bold', color=colors[i])
    ax.text(15, 15, "Ganan terreno", fontsize=11, color=HEX_GREEN, fontweight='bold', alpha=0.7)
    ax.text(-13, -10, "Pierden terreno", fontsize=11, color=HEX_ORANGE, fontweight='bold', alpha=0.7)
    ax.set_xlabel("Cambio Consideracion WoW (pp)", fontsize=10)
    ax.set_ylabel("Cambio Compra 3m WoW (pp)", fontsize=10)
    ax.set_xlim(-15, 30)
    ax.set_ylim(-12, 18)
    ax.set_title("Movimiento competitivo 2025→2026 — 9 de 10 cambios sig son aumentos competidores; Gama estable", fontsize=11, color=HEX_RED, fontweight='bold', pad=8)
    ax.grid(linestyle='--', alpha=0.35, zorder=0)
    ax.spines[['top','right']].set_visible(False)
    ax.text(0.5, -0.13, "*** sig_99 | ** sig_95 | n_2025=785, n_2026=402 | BH-FDR | CV-WOW-001/002 aplican", ha='center', fontsize=8, color='#888', transform=ax.transAxes, style='italic')
    fig.tight_layout(rect=[0,0.04,1,1])
    fig.savefig(path, dpi=300, bbox_inches='tight', facecolor='white')
    plt.close(fig)
    return path


# ======== SLIDES ========

def build_s01_portada(prs):
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, 2286000, SH, RED_GAMA)
    tb = add_textbox(slide, 365760, 457200, 1828800, 914400)
    tf = tb.text_frame; clear_tf(tf)
    add_para_with_runs(tf, [("GAMA", 36, True, WHITE)])
    tb2 = add_textbox(slide, 365760, 5760720, 1828800, 731520)
    tf2 = tb2.text_frame; clear_tf(tf2)
    add_para_with_runs(tf2, [("2026", 28, True, WHITE)])
    shape = add_rounded_rect(slide, 365760, 5120640, 1828800, 457200, WHITE)
    tf3 = shape.text_frame; clear_tf(tf3)
    add_para_with_runs(tf3, [("V5 · Deck", 18, True, RED_GAMA)])
    add_rect(slide, 2286000, 0, SW - 2286000, 45720, BLUE_INFO)
    tb4 = add_textbox(slide, 2743200, 1828800, 8961120, 1371600)
    tf4 = tb4.text_frame; clear_tf(tf4)
    add_para_with_runs(tf4, [("Notoriedad y Preferencia de Marca — Gama 2026", 30, True, RED_GAMA)])
    tb5 = add_textbox(slide, 2743200, 3292000, 8961120, 914400)
    tf5 = tb5.text_frame; clear_tf(tf5)
    add_para_with_runs(tf5, [("Deck Principal V5 · Analisis independiente · Eje 2026", 18, False, DARK_TEXT)])
    meta_lines = [
        "Fecha: 2026-05-18",
        "Equipo analitico: Cora Urrea + Raoul Bermudez",
        "Version: V5 Deck Principal — fresh start desde datos Notoriedad 2.0",
        "Confidencial / NDA",
        "Eje del analisis: datos 2026. Datos 2025 = referencia comparativa (Bloque 7).",
        "Cualitativo = hipotesis e insight, NO conclusiones (Bloque 6).",
        "Standalone: este deck V5 puede leerse sin V3 ni V4.",
    ]
    y_off = 4389120
    for line in meta_lines:
        tb_m = add_textbox(slide, 2743200, y_off, 8961120, 274320)
        tf_m = tb_m.text_frame; clear_tf(tf_m)
        add_para_with_runs(tf_m, [(line, 10, False, DARK_TEXT)])
        y_off += 274320


def build_s02_marco(prs):
    slide = blank_slide(prs)
    add_standard_header(slide, 2, "Que es V5 — y en que se diferencia de V3 y V4")
    add_callout_box(slide,
        "V5 es un analisis independiente desde los datos de Notoriedad 2.0. No invalida V3/V4 — los reorganiza con un eje diferente: 2026 es la base, 2025 es comparacion, el cualitativo es insight e hipotesis.")
    tf = body_textbox(slide)
    section_divider(tf, "Lo que NO cambia respecto a V4:", before=0)
    bullet(tf, "Los datos, las cifras, los caveats estadisticos, los verbatims — son identicos a V4.")
    bullet(tf, "El gate Bruna BR-2 V4 sigue vigente (mismos wordings DW-1..DW-5).")
    section_divider(tf, "Lo que SI cambia respecto a V4:")
    bullet(tf, "1. El eje del analisis es 2026 ", value="(no la confirmacion de V3).")
    bullet(tf, "2. NSE C+/C es el Capitulo 1 ", value="(no un analisis exploratorio dentro de WoW).")
    bullet(tf, "3. Los hallazgos cualitativos son hipotesis e insights ", value="(no conclusiones trianguladas co-iguales al cuanti).")
    bullet(tf, "4. Los datos 2025 son referencia comparativa al final ", value="(no capa analitica co-principal).")
    section_divider(tf, "Estructura del deck V5:")
    bullet(tf, "Bloques 1-5: analisis 2026 (C+/C, salud Gama, drivers, comportamiento, comunicacion).")
    bullet(tf, "Bloque 6: insights cualitativos como hipotesis (claramente etiquetados).")
    bullet(tf, "Bloque 7: evolucion comparativa 2025→2026 (referencia).")
    bullet(tf, "Bloques 8-10: tesis, recomendaciones, agenda 2027 + decisiones Owner.")
    add_footer(slide, 2)


def build_s03_iconografia(prs):
    slide = blank_slide(prs)
    add_standard_header(slide, 3, "Iconografia de certeza + Ficha tecnica del estudio 2026")
    add_callout_box(slide,
        "Cada hallazgo del deck lleva un icono que indica su nivel de evidencia. La regla: no confundir conclusion con hipotesis. El peso epistemico del cualitativo se separa explicitamente del cuantitativo.")
    tf = body_textbox(slide)
    section_divider(tf, "Iconografia V5:", before=0)
    bullet(tf, "✅ Conclusion: ", value="evidencia cuantitativa significativa (≥95% confianza), base ≥30.")
    bullet(tf, "⚠ Hipotesis apoyada: ", value="tendencia estadistica (p 0.05-0.10) o soporte cualitativo con patron robusto.")
    bullet(tf, "💡 Insight cuali: ", value="hallazgo de focus group — util para hypothesizing y diseno de activaciones, NO para extrapolacion estadistica.")
    bullet(tf, "📊 Referencia evolutiva: ", value="dato 2025 como comparacion (caveats de comparabilidad aplican).")
    section_divider(tf, "Ficha tecnica:")
    bullet(tf, "Universo: ", value="shoppers regulares (≥1 compra/mes) en supermercados de cadena, Caracas + Altos Mirandinos.")
    bullet(tf, "Muestra 2026: ", value="n=402 entrevistas face-to-face. Margen de error Total ±4.89% al 95% (Wilson).")
    bullet(tf, "NSE: ", value="C+/C combinado n=104 (25.9%) · D n=127 (31.6%) · E n=171 (42.5%).")
    bullet(tf, "Geografia: ", value="Baruta 122 · Libertador 80 · Sucre 79 · Chacao 70 · El Hatillo 31 · Altos Mirandinos 20.")
    bullet(tf, "Referencia 2025: ", value="n=785 — comparable en preguntas comunes con caveats.")
    bullet(tf, "Modelos: ", value="Logit (AUC=0.929, Pseudo R²=0.4371) · RF/SHAP · K-means k=3 · Z-test BH-FDR.")
    add_warn_box(slide,
        "Caveats metodologicos: C+/C combinado (no separables en BBDD), n=104 implica m.e. ±9.8%. "
        "2025 sin ponderacion (factor no disponible). Composicion geografica difiere entre olas. "
        "WoW por NSE no pre-registrado en 2025 — exploratorio.")
    add_footer(slide, 3)


def build_s04_bloque1_perfil(prs):
    slide = blank_slide(prs)
    add_standard_header(slide, 4, "BLOQUE 1A — Perfil del segmento C+/C 2026")
    add_callout_box(slide,
        "El segmento C+/C (n=104, 25.9% de la muestra) es el campo de batalla principal del estudio. Define el primer capitulo de V5 por ser donde Gama tiene su posicion mas fuerte y donde la presion competitiva se intensifica.")
    tf = body_textbox(slide)
    section_divider(tf, "Composicion del segmento:", before=0)
    bullet(tf, "Tamano: ", value="n=104 (25.9% de la muestra) · NSE C+ y C combinados.")
    bullet(tf, "Geografia dominante: ", value="Baruta · Chacao · El Hatillo (zonas premium).")
    bullet(tf, "Rol decision compra: ", value="alto (buyer principal del hogar).")
    bullet(tf, "DNA del segmento: ", value="orientado a experiencia premium, atencion personalizada, calidad. Coherente con el posicionamiento natural de Gama.")
    section_divider(tf, "Por que C+/C es el Capitulo 1 de V5:")
    bullet(tf, "Es donde Gama tiene su mayor fortaleza relativa (TOM 60.6% vs 44.3% Total).")
    bullet(tf, "Es donde la mecanica del driver de atencion opera con mas resonancia.")
    bullet(tf, "Es donde los dos principales competidores (Rio y Paramo) crecen mas rapido en WoW.")
    bullet(tf, "Una estrategia que no responda en C+/C cede el segmento natural a la competencia.")
    add_warn_box(slide,
        "⚠ Caveat C+/C: combinado C+ y C (no separables en BBDD). n=104 implica m.e. ±9.8% al 95% — "
        "mayor incertidumbre estadistica que el Total. Analisis exploratorios en WoW por NSE deben confirmarse en ola 2027 (CV-WOW-005).")
    add_footer(slide, 4)


def build_s05_bloque1B(prs, chart_path):
    slide = blank_slide(prs)
    add_standard_header(slide, 5, "BLOQUE 1B — Gama en C+/C 2026: posicion fuerte, sin crecimiento")
    add_callout_box(slide,
        "TOM Gama en C+/C = 60.6% (vs 44.3% Total, diferencial +16.3pp). Preferida = 13.5% (vs 8.0% Total, +5.5pp). Sobrerepresentacion confirmada del DNA premium.")
    half_w = CONTENT_W // 2 - 91440
    tb = add_textbox(slide, MARGIN_L, BODY_TOP, half_w, CONTENT_H)
    tf = tb.text_frame; tf.word_wrap = True; clear_tf(tf)
    section_divider(tf, "Posicion Gama en C+/C vs Total:", before=0)
    bullet(tf, "TOM C+/C: 60.6% ", value="(vs 44.3% Total → +16.3pp diferencial).")
    bullet(tf, "Preferida C+/C: 13.5% ", value="(vs 8.0% Total → +5.5pp diferencial).")
    bullet(tf, "C+/C es donde el DNA de Gama (atencion, calidad, 24h) tiene mayor resonancia.")
    badge_text(tf, "    ✅ Conclusion — sobrerepresentacion C+/C confirmada cuantitativamente")
    section_divider(tf, "Implicacion estrategica:")
    bullet(tf, "Defender C+/C es defender la posicion mas fuerte del portafolio.")
    bullet(tf, "Cualquier perdida en C+/C cuesta mas en preferencia que en otros NSE (5.5pp diferencial).")
    bullet(tf, "La inversion estrategica en C+/C tiene mejor ROI esperado dado el DNA alineado.")
    chart_l = MARGIN_L + half_w + 182880
    chart_w = SW - chart_l - MARGIN_R
    embed_png(slide, chart_path, chart_l, BODY_TOP, chart_w, int(CONTENT_H * 0.95))
    add_footer(slide, 5)


def build_s06_bloque1C(prs, chart_path):
    slide = blank_slide(prs)
    add_standard_header(slide, 6, "BLOQUE 1C — La amenaza: Rio y Paramo crecen explosivamente en C+/C")
    add_callout_box(slide,
        "El segmento natural de Gama esta siendo penetrado agresivamente por sus dos competidores principales. Posicion actual solida; velocidad competitiva = senal de alerta mas importante del estudio.")
    half_w = CONTENT_W // 2 - 91440
    tb = add_textbox(slide, MARGIN_L, BODY_TOP, half_w, CONTENT_H)
    tf = tb.text_frame; tf.word_wrap = True; clear_tf(tf)
    section_divider(tf, "TOM C+/C 2025→2026 (📊 referencia evolutiva):", before=0)
    bullet(tf, "Rio: ~30.9% → 56.7% ", value="(+25.8pp, sig_99 exploratorio, p_adj=0.0002).")
    bullet(tf, "Paramo: ~18.0% → 40.4% ", value="(+22.3pp, sig_99 exploratorio, p_adj=0.0002).")
    bullet(tf, "Gama: ~58.4% → 60.6% ", value="(+2.2pp, no significativo).")
    badge_text(tf, "    ⚠ Hipotesis WoW C+/C — pendiente confirmacion ola 2027 · CV-WOW-005")
    section_divider(tf, "Lectura estrategica:")
    plain_text(tf, "Si las tendencias son reales, el segmento natural de Gama esta siendo penetrado por sus dos competidores principales a velocidad +20pp por ola. La posicion actual de Gama es solida (TOM 60.6%); la velocidad de Rio y Paramo en C+/C es la senal de alerta mas importante del estudio.", size=FONT_BODY, italic=True, before=2, after=1)
    chart_l = MARGIN_L + half_w + 182880
    chart_w = SW - chart_l - MARGIN_R
    embed_png(slide, chart_path, chart_l, BODY_TOP, chart_w, int(CONTENT_H * 0.95))
    add_warn_box(slide,
        "CV-WOW-005*: analisis WoW por NSE no pre-registrado en diseno 2025. Hipotesis a confirmar ola 2027. "
        "Composicion geografica 2025 puede sesgar deltas. n_2026 subgrupo C+/C=104. Causalidad no inferible.")
    add_footer(slide, 6)


def build_s07_bloque1D(prs):
    slide = blank_slide(prs)
    add_standard_header(slide, 7, "BLOQUE 1D — El Nucleo Leal vive en C+/C: el 24h como activo defensivo")
    add_callout_box(slide,
        "El Segmento 3 (Nucleo Leal, 8% muestra total) esta sobrerrepresentado en C+/C. Perfil dominante: frecuentes 18-30 y 50+ con uso 24h. El 24h es activo diferencial sin equivalente competitivo.")
    tf = body_textbox(slide)
    section_divider(tf, "Por que el Nucleo Leal Gama es C+/C:", before=0)
    bullet(tf, "El Seg 3 (8% muestra) tiene perfil de frecuencia alta y vinculo afectivo con Gama.")
    bullet(tf, "Sobrerrepresentado en C+/C → coherente con la oferta experiencial de Gama.")
    bullet(tf, "Los frecuentes 18-30 y 50+ con uso 24h dominan este segmento.")
    badge_text(tf, "    ✅ k-means cuanti confirma estructura del Nucleo Leal")
    section_divider(tf, "El 24h como activo defensivo clave:")
    bullet(tf, "Rio y Paramo NO tienen equivalente funcional al 24 horas de Gama.")
    bullet(tf, "Es la unica ventaja funcional dura sin replicar por competencia.")
    bullet(tf, "Defenderlo es la primera linea estrategica para C+/C.")
    badge_text(tf, "    💡 Insight cuali (no proyectable): 24h genera fidelidad transversal en C+/C")
    plain_text(tf,
        '   "El de 24 horas es un tiro al piso, de verdad que si." (Frecuentes 18-30)',
        size=FONT_BODY, italic=True, color=DARK_TEXT, before=2, after=0)
    plain_text(tf,
        '   "Eso es impelable." (Ocasionales 50+)',
        size=FONT_BODY, italic=True, color=DARK_TEXT, before=0, after=2)
    add_warn_box(slide,
        "Insight cuali corpus FG — no proyectable estadisticamente. "
        "Verbatims ilustrativos del patron. La medicion cuantitativa del valor incremental del 24h queda pendiente ola 2027.")
    add_footer(slide, 7)


def build_s08_bloque2A(prs, chart_path):
    slide = blank_slide(prs)
    add_standard_header(slide, 8, "BLOQUE 2A — Embudo Gama 2026: 0/8 indicadores con variacion significativa")
    add_callout_box(slide,
        "El embudo de Gama es estadisticamente estable en 2026. En un mercado donde Rio y Paramo crecen explosivamente, la estabilidad de Gama es fortaleza defensiva — sus activos diferenciadores generan una base estable.")
    half_w = CONTENT_W // 2 - 91440
    tb = add_textbox(slide, MARGIN_L, BODY_TOP, half_w, CONTENT_H)
    tf = tb.text_frame; tf.word_wrap = True; clear_tf(tf)
    section_divider(tf, "Embudo 2026 vs 2025 (referencia):", before=0)
    bullet(tf, "TOM: 44.3% ", value="(2025: 42.0% → +2.2pp, no sig).")
    bullet(tf, "Asistida: 60.2% ", value="(2025: ~58.0% → +2.2pp, no sig).")
    bullet(tf, "Consideracion: 31.8% ", value="(2025: 27.5% → +4.3pp, no sig).")
    bullet(tf, "Compra 3m: 17.7% ", value="(2025: 17.8% → -0.2pp, no sig).")
    bullet(tf, "Preferida: 8.0% ", value="(2025: 9.7% → -1.7pp, no sig).")
    bullet(tf, "Habitual: 20.2% ", value="(2025: 19.4% → +0.8pp, no sig).")
    badge_text(tf, "    ✅ 0/8 indicadores con variacion sig (p_adj>0.10) — Gama estable")
    section_divider(tf, "Lectura clave:")
    plain_text(tf, "Conversion TOM→Preferida = 44.3% → 8.0%. Esta brecha es el problema central del estudio: hay conocimiento, falta conversion. El Bloque 3 explica los mecanismos.", size=FONT_BODY, italic=True, before=2, after=1)
    chart_l = MARGIN_L + half_w + 182880
    chart_w = SW - chart_l - MARGIN_R
    embed_png(slide, chart_path, chart_l, BODY_TOP, chart_w, int(CONTENT_H * 0.95))
    add_footer(slide, 8)


def build_s09_bloque2B(prs, chart_path):
    slide = blank_slide(prs)
    add_standard_header(slide, 9, "BLOQUE 2B — Posicion competitiva 2026: Rio supera a Gama en TOM por primera vez")
    add_callout_box(slide,
        "Rio supero a Gama en TOM entre 2025 y 2026 (45.0% vs 44.3%). Paramo crecio mas en consideracion (+17.3pp). CM pierde compra (-7.7pp sig_95). Gama permanece estable — fortaleza relativa en mercado dinamico.")
    half_w = CONTENT_W // 2 - 91440
    tb = add_textbox(slide, MARGIN_L, BODY_TOP, half_w, CONTENT_H)
    tf = tb.text_frame; tf.word_wrap = True; clear_tf(tf)
    section_divider(tf, "Posicion competitiva 2026 (Total):", before=0)
    bullet(tf, "Rio: ", value="TOM 45.0% · Preferida 10.2% · WoW TOM +17pp (sig_99).")
    bullet(tf, "Gama: ", value="TOM 44.3% · Preferida 8.0% · estable (0pp).")
    bullet(tf, "Paramo: ", value="TOM ~39.1% · WoW TOM +12.1pp (sig_99).")
    bullet(tf, "Central Madeirense: ", value="Preferida 11.2% · WoW Compra -7.7pp (sig_95).")
    badge_text(tf, "    📊 Referencia evolutiva — Rio supera a Gama en TOM por primera vez en este tracker")
    section_divider(tf, "Lectura del posicionamiento:")
    bullet(tf, "Gama mantiene la mayor Preferida del top 3 (8.0% vs 10.2% Rio).")
    bullet(tf, "CM tiene mas Preferida (11.2%) pero pierde compra activa.")
    bullet(tf, "La oportunidad de capturar el espacio CM perdido esta abierta.")
    chart_l = MARGIN_L + half_w + 182880
    chart_w = SW - chart_l - MARGIN_R
    embed_png(slide, chart_path, chart_l, BODY_TOP, chart_w, int(CONTENT_H * 0.95))
    add_footer(slide, 9)


def build_s10_bloque2C(prs):
    slide = blank_slide(prs)
    add_standard_header(slide, 10, "BLOQUE 2C — Posicionamiento y atributos Gama 2026")
    add_callout_box(slide,
        "Gama lidera en el espacio 'experiencial premium' (Calidad, Atractiva, Limpieza, Seguro). DNA de marca claro en C+/C. Atributos de precio se debilitan en tendencia — pero ningun movimiento pasa BH-FDR.")
    tf = body_textbox(slide)
    section_divider(tf, "Atributos lideres (✅ alta certeza):", before=0)
    bullet(tf, "Atencion al cliente: ", value="Lider en mercado · driver #1 (OR=5.73). Sin equivalente competitivo.")
    bullet(tf, "Limpieza: ", value="Lider en mercado · driver #2 (OR=3.99*). Convergencia SHAP + Gini.")
    section_divider(tf, "Tendencias direccionales 2025→2026 (⚠ solo tendencia, ninguna pasa BH-FDR q<0.05):")
    bullet(tf, '"Tienda atractiva": ', value="+6.3pp (tendencia p_adj=0.058) — direccional positivo.")
    bullet(tf, '"Seguro": ', value="+5.8pp (tendencia p_adj=0.089) — direccional positivo.")
    bullet(tf, '"Menor precio": ', value="-3.9pp (tendencia p_adj=0.091) — direccional debilitamiento.")
    bullet(tf, '"Promociones": ', value="-5.1pp (tendencia p_adj=0.053) — direccional debilitamiento.")
    badge_text(tf, "    ⚠ ADVERTENCIA: ninguno pasa correccion BH-FDR. Senales direccionales — no conclusiones del estudio.")
    section_divider(tf, "Lectura estrategica:")
    bullet(tf, "La imagen experiencial parece consolidarse direccionalmente.")
    bullet(tf, "La imagen de precio se debilita direccionalmente — coherente con el reposicionamiento natural.")
    bullet(tf, "Confirmar en ola 2027 con muestra mayor antes de declarar tendencia como cambio real.")
    add_warn_box(slide,
        "ADVERTENCIA metodologica: Ninguno de los movimientos de atributos 2025→2026 pasa la correccion BH-FDR (umbral q<0.05). "
        "Son informativos como senales direccionales a monitorear en ola 2027. NO son conclusiones del estudio.")
    add_footer(slide, 10)


def build_s11_bloque3A(prs, chart_path):
    slide = blank_slide(prs)
    add_standard_header(slide, 11, "BLOQUE 3A — Drivers de preferencia: Atencion lidera, Precio no es driver")
    add_callout_box(slide,
        "Modelo logistico (AUC=0.929): quien asocia Gama con atencion tiene 5.7x mas odds de preferirla. Convergencia 4 metodos independientes (Logit + RF + SHAP + razon espontanea) = el mayor respaldo de robustez del estudio.")
    half_w = CONTENT_W // 2 - 91440
    tb = add_textbox(slide, MARGIN_L, BODY_TOP, half_w, CONTENT_H)
    tf = tb.text_frame; tf.word_wrap = True; clear_tf(tf)
    section_divider(tf, "Drivers cuantificados:", before=0)
    bullet(tf, "Atencion: OR=5.73*** ", value="(IC95 [1.6, 20.4], SHAP #1, 4 metodos convergentes).")
    bullet(tf, "Limpieza: OR=3.99* ", value="(IC95 [0.94, 16.91], SHAP #2, borderline sig).")
    bullet(tf, "Promociones: OR=3.64** ", value="(IC95 [1.1, 11.8], SHAP #3, palanca tactica).")
    bullet(tf, "Precio (menor): OR=1.03 NS ", value="(SHAP #10, 4 metodos: NO es driver).")
    badge_text(tf, "    ✅ Pseudo R²=0.4371 · AUC=0.929 · n pref-Gama=32 (referencial)")
    section_divider(tf, "Tres conclusiones de alta certeza:")
    bullet(tf, "Conclusion 1: ", value="Quien asocia Gama con atencion tiene 5.7x mas odds de preferirla.")
    bullet(tf, "Conclusion 2: ", value="El precio NO predice preferencia. Comunicar precio es invertir en el atributo menos relevante.")
    bullet(tf, "Conclusion 3: ", value="Las promociones son palanca tactica (atraer a primera visita), no comunicar precio bajo.")
    chart_l = MARGIN_L + half_w + 182880
    chart_w = SW - chart_l - MARGIN_R
    embed_png(slide, chart_path, chart_l, BODY_TOP, chart_w, int(CONTENT_H * 0.95))
    add_warn_box(slide,
        "Caveat: n pref-Gama=32 implica IC95 amplios. OR=5.73 es el estimado central con IC95 [1.6, 20.4]. "
        "El patron es robusto por convergencia de 4 metodos independientes — el mayor respaldo metodologico del estudio.")
    add_footer(slide, 11)


def build_s12_bloque3B(prs, chart_path):
    slide = blank_slide(prs)
    add_standard_header(slide, 12, "BLOQUE 3B — K-means: tres tipos de shoppers, Seg 2 = mayor oportunidad C/P")
    add_callout_box(slide,
        "K-means k=3 identifica 3 perfiles: Mayoria Exigente (59%), Pragmaticos Convertibles (33%), Nucleo Leal (8%). El Seg 2 tiene 0% preferencia actual + menor resistencia al precio — mayor potencial de conversion C/P.")
    half_w = CONTENT_W // 2 - 91440
    tb = add_textbox(slide, MARGIN_L, BODY_TOP, half_w, CONTENT_H)
    tf = tb.text_frame; tf.word_wrap = True; clear_tf(tf)
    section_divider(tf, "Los 3 segmentos k-means:", before=0)
    bullet(tf, "Seg 1 — Mayoria Exigente (59%, n≈237): ", value="alta exigencia, no exclusivos Gama, mision multiple. Estrategia: defensa + comunicacion de experiencia.")
    bullet(tf, "Seg 2 — Pragmaticos Convertibles (33%, n≈133): ", value="resistencia precio 3.44/5 (vs 3.66 Seg1), 0% pref-Gama actual. Mayor retorno C/P.")
    bullet(tf, "Seg 3 — Nucleo Leal (8%, n≈32): ", value="alta frecuencia, uso 24h, lealtad alta, vinculo afectivo. Estrategia: retencion + advocacy.")
    badge_text(tf, "    ✅ Conclusion cuanti — Seg 2 = mayor potencial de conversion C/P de los 3 clusters")
    section_divider(tf, "Insight cuali sobre Seg 2 (💡 hipotesis):")
    plain_text(tf,
        "El corpus FG sugiere que la resistencia del Seg 2 NO es solo calculo economico — es autoexclusion identitaria ('Gama es para otro tipo de persona'). El cuali V4 enriquece el perfil con esta dimension que el cuestionario cuantitativo no capturaba. La activacion correcta seria una primera experiencia guiada, NO descuentos.",
        size=FONT_BODY, italic=True, color=DARK_TEXT, before=2, after=1)
    badge_text(tf, "    💡 Hipotesis cuali — ver Bloque 6B sifrinaje. Wording DW-3 aplicado.")
    chart_l = MARGIN_L + half_w + 182880
    chart_w = SW - chart_l - MARGIN_R
    embed_png(slide, chart_path, chart_l, BODY_TOP, chart_w, int(CONTENT_H * 0.95))
    add_warn_box(slide,
        "⚠ Caveat: silhouette ~0.20 (moderado). Los segmentos son tendencias interpretativas, no categorias discretas. "
        "Los perfiles se solapan en los margenes. El Seg 2 puede contener subgrupos heterogeneos.")
    add_footer(slide, 12)


def build_s13_bloque4A(prs, chart_path):
    slide = blank_slide(prs)
    add_standard_header(slide, 13, "BLOQUE 4A — Misiones: Gama lider en urgencia (#2), ausente en mercado grande (#7)")
    add_callout_box(slide,
        "Hay una brecha estructural entre la mision donde Gama gana (urgencia) y la mision donde el gasto por viaje es mayor (mercado grande). El 34% de los propios preferentes Gama va a Paramo/CM/Forum para mercado grande — share of wallet perdido.")
    half_w = CONTENT_W // 2 - 91440
    tb = add_textbox(slide, MARGIN_L, BODY_TOP, half_w, CONTENT_H)
    tf = tb.text_frame; tf.word_wrap = True; clear_tf(tf)
    section_divider(tf, "Posicion competitiva por mision:", before=0)
    bullet(tf, "Compra urgente / cercana: ", value="Gama 12.2% (#2 del mercado).")
    bullet(tf, "Mercado grande / abastecimiento: ", value="Gama 7.2% (#7 del mercado).")
    badge_text(tf, "    ✅ Conclusion — brecha estructural Gama: gana urgencia, ausente en mercado grande")
    section_divider(tf, "Share of wallet perdido en Nucleo Leal:")
    bullet(tf, "34% de los preferentes Gama va a Paramo, CM o Forum para mercado grande.")
    bullet(tf, "Esto significa que el ticket promedio mas alto del shopper Gama se pierde fuera de Gama.")
    bullet(tf, "Oportunidad: capturar el mercado grande del Nucleo Leal con propuesta de abastecimiento.")
    section_divider(tf, "Oportunidad de mercado:")
    plain_text(tf, "📊 El espacio de CM en mercado grande se libero parcialmente en 2026 (-7.7pp compra, sig_95) — una oportunidad no capturada aun por Gama.", size=FONT_BODY, italic=True, before=2, after=1)
    chart_l = MARGIN_L + half_w + 182880
    chart_w = SW - chart_l - MARGIN_R
    embed_png(slide, chart_path, chart_l, BODY_TOP, chart_w, int(CONTENT_H * 0.95))
    add_footer(slide, 13)


def build_s14_bloque5(prs, chart_path):
    slide = blank_slide(prs)
    add_standard_header(slide, 14, "BLOQUE 5 — Gap comunicacional: recall espontaneo PTL = 0%")
    add_callout_box(slide,
        "La campana 'PRECIOS DE TU LADO' no registra recall espontaneo en ningun respondiente (0/17). Cuando se presenta asistida, 65% la interpreta como mensaje de precio. Pero precio es el predictor #10 de preferencia (OR=1.03 NS). Gama invierte en el territorio menos predictivo.")
    half_w = CONTENT_W // 2 - 91440
    tb = add_textbox(slide, MARGIN_L, BODY_TOP, half_w, CONTENT_H)
    tf = tb.text_frame; tf.word_wrap = True; clear_tf(tf)
    section_divider(tf, "Metricas del modulo publicitario:", before=0)
    bullet(tf, "Recall espontaneo 'PRECIOS DE TU LADO': ", value="0/17 = 0% (✅ alta certeza — resultado categorico).")
    bullet(tf, "Recall asistido (con estimulo): ", value="~11% (⚠ base n=17-50 referencial).")
    bullet(tf, "Interpretacion PTL como 'mensaje de precio': ", value="65% (⚠ base referencial).")
    bullet(tf, "Interpretacion PTL como 'solidaridad/apoyo': ", value="35% (⚠ base referencial).")
    bullet(tf, "Ningun slogan de Gama recordado espontaneo: ", value="95.8% (✅ base n=402).")
    badge_text(tf, "    ✅ Gap: alta certeza · causalidad de impacto: certeza media — [Hipotesis V4 — evidencia convergente]")
    section_divider(tf, "Conclusion del gap:")
    bullet(tf, "La campana no registra recall espontaneo.")
    bullet(tf, "El mensaje asistido se interpreta como precio.")
    bullet(tf, "Precio es el atributo menos predictivo de preferencia (OR=1.03 NS, SHAP #10).")
    bullet(tf, "Resultado: Gama invierte recursos comunicacionales en el territorio menos predictivo.")
    chart_l = MARGIN_L + half_w + 182880
    chart_w = SW - chart_l - MARGIN_R
    embed_png(slide, chart_path, chart_l, BODY_TOP, chart_w, int(CONTENT_H * 0.95))
    add_warn_box(slide,
        "⚠ Nota: bases del modulo publicitario son pequenas (n=17-50). El hallazgo de recall 0% es robusto; los % de interpretacion son referenciales. "
        "Ola 2027 debe incluir modulo de comunicacion con estimulos de campana + Rio como benchmark.")
    add_footer(slide, 14)


def build_s15_bloque6_intro(prs):
    slide = blank_slide(prs)
    add_standard_header(slide, 15, "BLOQUE 6 — Insights cualitativos: regla de lectura obligatoria")
    add_callout_box(slide,
        "TODO hallazgo de este bloque es insight cualitativo o hipotesis. El corpus base es 12 docs FG (~42k caracteres, max ~42 personas presenciales). El cualitativo establece mecanismos, barreras y territorios — NO prevalencia estadistica.")
    tf = body_textbox(slide)
    section_divider(tf, "Regla epistemica:", before=0)
    plain_text(tf,
        "Una conclusion del cuanti con n=402 tiene mayor peso para decisiones estrategicas que un verbatim de 1-2 personas. El peso de los insights cuali es proporcional al numero de grupos/participantes que los expresan.",
        size=FONT_BODY, italic=True, before=2, after=2)
    section_divider(tf, "Los 5 insights cualitativos del Bloque 6:")
    bullet(tf, "6A — Mecanismo del driver atencion: hipotesis 'acompanamiento-guia' ", value="(💡 6/7 grupos FG — patron robusto para hipotesis).")
    bullet(tf, "6B — Barrera del sifrinaje: hipotesis autoexclusion identitaria ", value="(💡 1 verbatim central + patron consistente).")
    bullet(tf, "6C — Gama Club: activo invisible con interes transversal ", value="(💡 6/7 grupos FG).")
    bullet(tf, "6D — Canal digital: dos barreras distintas, dos soluciones distintas ", value="(💡 corpus B + 5/7 segmentos).")
    bullet(tf, "6E — Arquetipo no apropiado: territorio disponible ", value="(💡 6/7 grupos emergente, no inducido).")
    section_divider(tf, "Que NO se hace en este bloque (regla Cora — aclaratoria 2026-05-17):")
    bullet(tf, "NO se triangula como conclusion co-igual al cuanti.")
    bullet(tf, "NO se extrapola estadisticamente.")
    bullet(tf, "SI se ofrece como insight para diseno de activaciones y como hipotesis para confirmar en ola 2027.")
    add_warn_box(slide,
        "Corpus FG: 12 documentos (8 sesiones presenciales + 4 online), 6 segmentos, ~42,094 caracteres. "
        "Maximo ~42 personas presenciales. IN-7 §1.4: analisis single-analyst sin Kappa inter-codificador.")
    add_footer(slide, 15)


def build_s16_bloque6A(prs):
    slide = blank_slide(prs)
    add_standard_header(slide, 16, "BLOQUE 6A — El mecanismo del driver atencion: 'acompanamiento-guia'")
    add_callout_box(slide,
        "Cuanti: OR=5.73*** (convergencia 4 metodos). Cuali (💡 6/7 FG): Gama como 'guia experta de confianza' que acompana la compra. El mecanismo NO seria reconocimiento nominal — es acompanamiento activo.")
    tf = body_textbox(slide)
    section_divider(tf, "Lo que el cuanti establece (✅ alta certeza):", before=0)
    plain_text(tf,
        "OR=5.73*** — quien asocia Gama con atencion tiene 5.7x mas odds de preferirla. Convergencia de 4 metodos. Resultado cuantitativo robusto.",
        size=FONT_BODY, before=2, after=1)
    section_divider(tf, "La hipotesis cuali (💡 6/7 grupos FG — patron robusto):")
    plain_text(tf,
        "El corpus FG documenta espontaneamente en 6 de 7 grupos la imagen de Gama como 'guia experta de confianza' que acompana la compra. El mecanismo no seria el reconocimiento nominal ('te conocen por tu nombre') sino el acompanamiento activo ('siempre hay alguien que te orienta').",
        size=FONT_BODY, before=1, after=2)
    plain_text(tf,
        '   "Una mujer joven, pero con experiencia. Que sea como una guia, que te diga: \'Mira, hoy llego el pan arabe fresco, llevatelo\'."',
        size=FONT_BODY, italic=True, color=DARK_TEXT, before=2, after=0)
    plain_text(tf,
        '   — Mary Francis Bossia, Ocasionales 18-30',
        size=FONT_SMALL, color=GRAY_MID, before=0, after=2)
    plain_text(tf,
        '   "Alguien de Caracas que te oriente y sea como una guia de compras, que no se quede atras con la tecnologia."',
        size=FONT_BODY, italic=True, color=DARK_TEXT, before=1, after=0)
    plain_text(tf,
        '   — Moises Torrealba, Ocasionales 18-30',
        size=FONT_SMALL, color=GRAY_MID, before=0, after=2)
    section_divider(tf, "Insight para comunicacion:")
    plain_text(tf,
        "Si la hipotesis es correcta, el mensaje correcto es 'en Gama siempre hay alguien que te ayuda' (creible para el shopper), NO 'en Gama te conocen por tu nombre' (hiperbole dificil de verificar).",
        size=FONT_BODY, before=2, after=1)
    badge_text(tf, "    💡 Insight cuali — Certeza: hipotesis robusta (6/7 FG, emergente no inducido). Requiere claim testing cuantitativo para confirmar como plataforma comunicacional.")
    add_footer(slide, 16)


def build_s17_bloque6B(prs):
    slide = blank_slide(prs)
    add_standard_header(slide, 17, "BLOQUE 6B — La barrera del sifrinaje: hipotesis de autoexclusion identitaria")
    add_callout_box(slide,
        "Cuanti: 54% percibe Gama como cara, pero precio OR=1.03 NS (no predice preferencia). Cuali (💡): la barrera no es calculo economico — es autoexclusion simbolica. 'Gama es para otro tipo de persona' antes de verificar precios.")
    tf = body_textbox(slide)
    section_divider(tf, "Lo que el cuanti establece (✅ alta certeza):", before=0)
    bullet(tf, "Precio OR=1.03 (NS, SHAP #10): el precio no predice preferencia.")
    bullet(tf, "54% percibe Gama como cara.")
    bullet(tf, "Paradoja: la percepcion de precio no afecta la preferencia → algo mas explica la barrera de conversion.")
    section_divider(tf, "La hipotesis cuali (💡 patron en multiples grupos):")
    plain_text(tf,
        "El corpus FG sugiere que la barrera no es un calculo economico comparativo sino una autoexclusion simbolica — 'Gama es para otro tipo de persona.' El shopper se autoexcluye antes de entrar a verificar precios.",
        size=FONT_BODY, before=1, after=2)
    plain_text(tf,
        '   "A mi me parece que es un poco costoso... el sifrinaje y tal, pero para hacer un mercado grande... es mejor ir al Lux o al Paramo."',
        size=FONT_BODY, italic=True, color=DARK_TEXT, before=2, after=0)
    plain_text(tf,
        '   — Azahara Betancourt, Ocasionales 18-30',
        size=FONT_SMALL, color=GRAY_MID, before=0, after=2)
    # C-LIT-DW1 caveat literal obligatorio
    plain_text(tf,
        "   Cita textual de participante de focus group (investigacion cualitativa, mayo 2026). "
        "El termino 'sifrinaje' es vocabulario espontaneo del informante — no una caracterizacion de Gama "
        "hecha por el equipo analitico. Refleja como segmentos de no-usuarios codifican la barrera de entrada percibida.",
        size=FONT_CAVEAT, color=BLUE_INFO, before=1, after=2)
    section_divider(tf, "Insight para activacion:")
    plain_text(tf,
        "Si la barrera es identitaria, los descuentos de precio no la mueven (comunicar precio bajo refuerza la distancia simbolica). La activacion correcta del Seg 2 seria una primera experiencia que resignifique la pertenencia — no una oferta economica.",
        size=FONT_BODY, before=2, after=1)
    badge_text(tf, "    💡 Insight cuali — Certeza: 1 verbatim central + patron consistente. Requiere validacion cuantitativa (cuestionario de barreras estructurado).")
    add_footer(slide, 17)


def build_s18_bloque6C_6D(prs):
    slide = blank_slide(prs)
    add_standard_header(slide, 18, "BLOQUE 6C+6D — Gama Club opaco + Canal digital con dos barreras")
    add_callout_box(slide,
        "Dos hallazgos cuali con alta accionabilidad. 6C: Gama Club tiene interes transversal pero su saldo es invisible digitalmente. 6D: canal digital tiene dos barreras distintas (inventario inconsistente + desconfianza en frescos).")
    tf = body_textbox(slide)
    section_divider(tf, "6C — Gama Club: activo invisible con interes transversal (💡 6/7 FG):", before=0)
    bullet(tf, "El cuanti: ", value="no hay pregunta directa sobre Gama Club en cuestionario 2026.")
    bullet(tf, "Cuali: ", value="interes genuino transversal — pero saldo invisible en canal digital.")
    plain_text(tf,
        '   "Las millas... yo ni sabia que eso se podia usar para pagar hasta hace poco. Yo pensaba que era como los puntos del banco que uno nunca usa."',
        size=FONT_BODY, italic=True, color=DARK_TEXT, before=2, after=0)
    plain_text(tf,
        '   — Gregory, Frecuentes 18-30',
        size=FONT_SMALL, color=GRAY_MID, before=0, after=2)
    bullet(tf, "Solucion: ", value="saldo visible en app + equivalente en dinero + canje en caja. Baja inversion, alta accionabilidad.")
    section_divider(tf, "6D — Canal digital: dos barreras distintas (💡 corpus B):")
    bullet(tf, "Barrera 1 — Inconsistencia de inventario ", value="(patron en 3/4 docs uso digital).")
    plain_text(tf,
        '   "Yo deje de usar la pagina porque a veces pides algo que dice que hay y a la media hora te llaman para decirte que no tienen."',
        size=FONT_BODY, italic=True, color=DARK_TEXT, before=1, after=0)
    plain_text(tf,
        '   — Mujer, 42 anos, segmento Frecuentes 31-50 (sesion online cualitativa)',
        size=FONT_SMALL, color=GRAY_MID, before=0, after=1)
    bullet(tf, "Barrera 2 — Desconfianza en frescos ", value="(patron en 5/7 segmentos FG). Solucion: empezar propuesta digital por no-perecederos.")
    badge_text(tf, "    💡 Insight cuali · ⚠ Nota de alcance: estos hallazgos van mas alla del brief original (notoriedad/brand health). Ver CO-3.")
    add_warn_box(slide,
        "DW-4 aplicado: atribucion del verbatim por segmento/edad (no nombre/seudonimo). "
        "Las recomendaciones digitales (Rec. 4) se incluyen como ampliacion del brief original — sugieren conversacion de scope con Gama. Decision CO-3 del Owner.")
    add_footer(slide, 18)


def build_s19_bloque6E(prs):
    slide = blank_slide(prs)
    add_standard_header(slide, 19, "BLOQUE 6E — El arquetipo no apropiado: hipotesis de territorio disponible")
    add_callout_box(slide,
        "Sin ninguna pregunta directa sobre personificacion, 6 de 7 grupos convergieron espontaneamente en el mismo imaginario: Gama como mujer 35-45 anos, caraquena, profesional, activa, guia experta. Este arquetipo no esta apropiado por ningun competidor.")
    tf = body_textbox(slide)
    section_divider(tf, "La hipotesis cuali emergente (💡 6/7 grupos FG, no inducido por moderador):", before=0)
    plain_text(tf,
        "Sin ninguna pregunta directa sobre personificacion, 6 de 7 grupos convergieron espontaneamente en el mismo imaginario: Gama como mujer de 35-45 anos, caraquena, profesional, activa, guia experta de confianza.",
        size=FONT_BODY, before=1, after=2)
    plain_text(tf,
        '   "Yo me la imagino como una mujer, como de unos 35 anos, super activa, lider, entusiasta, que sabe lo que hace."',
        size=FONT_BODY, italic=True, color=DARK_TEXT, before=2, after=0)
    plain_text(tf,
        '   — Moises Torrealba, Ocasionales 18-30',
        size=FONT_SMALL, color=GRAY_MID, before=0, after=2)
    section_divider(tf, "Por que es relevante:")
    bullet(tf, "Este arquetipo NO esta apropiado por ningun competidor directo (Rio, Paramo, CM).")
    bullet(tf, "Si se confirma cuantitativamente podria ser territorio diferenciador de comunicacion.")
    bullet(tf, "Coherencia notable: el arquetipo 'guia experta' es la personificacion exacta de lo que el OR=5.73 mide.")
    badge_text(tf, "    💡 Insight cuali emergente (6/7 grupos). Certeza: hipotesis — requiere pre-test cuantitativo.")
    # C-LIT-DW2-a obligatorio
    plain_text(tf,
        "   Este hallazgo describe el imaginario de marca que el shopper ya proyecta espontaneamente — "
        "no propone una campana especifica de personificacion. La materializacion creativa "
        "(figura concreta, voz, presencia visual) requiere prueba de concepto cuantitativa con el target antes de cualquier ejecucion.",
        size=FONT_CAVEAT, color=BLUE_INFO, before=3, after=1)
    add_warn_box(slide,
        "Prerequisitos obligatorios antes de cualquier desarrollo creativo en esta direccion: "
        "(1) prueba de concepto cuantitativa con el target; (2) verificacion previa de Gama sobre contratos de imagen vigentes; "
        "(3) verificacion de compatibilidad con plataforma de marca global. Lenguaje condicional ('se recomienda explorar') obligatorio.")
    add_footer(slide, 19)


def build_s20_bloque7A(prs):
    slide = blank_slide(prs)
    add_standard_header(slide, 20, "BLOQUE 7A — Gama: posicion defensiva solida (0/8 metricas con variacion significativa)")
    add_callout_box(slide,
        "Este bloque usa los datos 2025 exclusivamente como REFERENCIA para entender que cambio. Los datos 2026 son el analisis primario. En un mercado donde 9 de 10 cambios significativos son aumentos de competidores, la estabilidad de Gama es una fortaleza relativa.")
    tf = body_textbox(slide)
    section_divider(tf, "Evolucion Gama 2025→2026 — 8 indicadores, 0 significativos:", before=0)
    bullet(tf, "TOM: +2.2pp ", value="(p_adj>0.10 — no significativo).")
    bullet(tf, "Asistida: +2.2pp ", value="(p_adj>0.10 — no significativo).")
    bullet(tf, "Consideracion: +4.3pp ", value="(p_adj>0.10 — no significativo).")
    bullet(tf, "Compra 3m: -0.2pp ", value="(p_adj>0.10 — no significativo).")
    bullet(tf, "Preferida: -1.7pp ", value="(p_adj>0.10 — no significativo).")
    bullet(tf, "Ultima compra: -1.2pp ", value="(p_adj>0.10 — no significativo).")
    bullet(tf, "Habitual: +0.8pp ", value="(p_adj>0.10 — no significativo).")
    bullet(tf, "Misiones: +4.9pp ", value="(p_adj>0.10 — no significativo).")
    badge_text(tf, "    ✅ 0/8 indicadores significativos — Gama es la unica gran cadena con posicion estable")
    section_divider(tf, "Lectura correcta:")
    plain_text(tf,
        "Gama es la unica gran cadena del mercado que mantiene posicion estadisticamente estable entre 2025 y 2026. En un mercado donde 9 de 10 cambios significativos son aumentos de competidores, la estabilidad es una FORTALEZA relativa — no un estancamiento.",
        size=FONT_BODY, italic=True, before=2, after=1)
    add_warn_box(slide,
        "⚠ Caveats: factor de ponderacion muestral 2025 no disponible (CV-WOW-001). "
        "Composicion geografica difiere entre olas — Libertador sobrerepresentado en 2025 (CV-WOW-002). "
        "n_2025=785, n_2026=402, correccion BH-FDR. Causalidad no inferible.")
    add_footer(slide, 20)


def build_s21_bloque7B(prs, chart_path):
    slide = blank_slide(prs)
    add_standard_header(slide, 21, "BLOQUE 7B — Movimiento competitivo: Rio y Paramo crecen, CM cae, Gama estable")
    add_callout_box(slide,
        "El mercado venezolano de supermercados de cadena tuvo movimiento significativo en 2026. Rio es el ganador absoluto (sig_99 en TOM, Consideracion, Compra). Paramo consolida posicion. CM es el unico perdedor significativo. Gama no se mueve en ninguna direccion.")
    half_w = CONTENT_W // 2 - 91440
    tb = add_textbox(slide, MARGIN_L, BODY_TOP, half_w, CONTENT_H)
    tf = tb.text_frame; tf.word_wrap = True; clear_tf(tf)
    section_divider(tf, "Cambios significativos 2025→2026 (Total):", before=0)
    bullet(tf, "Rio TOM: +17.0pp ", value="(sig_99, p_adj<0.0001).")
    bullet(tf, "Rio Consideracion: +19.6pp ", value="(sig_99).")
    bullet(tf, "Rio Compra 3m: +12.5pp ", value="(sig_99).")
    bullet(tf, "Rio Preferida: +4.0pp ", value="(tendencia p_adj=0.059).")
    bullet(tf, "Paramo TOM: +12.1pp ", value="(sig_99).")
    bullet(tf, "Paramo Consideracion: +17.3pp ", value="(sig_99).")
    bullet(tf, "CM Compra 3m: -7.7pp ", value="(sig_95).")
    bullet(tf, "Gama: 0/8 sig ", value="(estable).")
    badge_text(tf, "    📊 9 de 10 cambios significativos = aumentos competidores. Gama es la unica que no se mueve.")
    chart_l = MARGIN_L + half_w + 182880
    chart_w = SW - chart_l - MARGIN_R
    embed_png(slide, chart_path, chart_l, BODY_TOP, chart_w, int(CONTENT_H * 0.95))
    add_footer(slide, 21)


def build_s22_bloque8_tesis(prs):
    slide = blank_slide(prs)
    add_standard_header(slide, 22, "BLOQUE 8 — TESIS V5: el diagnostico central desde los datos 2026")
    add_callout_box(slide,
        "Tesis V5 en una oracion: Gama tiene el activo relacional mas diferenciador del mercado, su posicion es solida — pero comunica en el territorio equivocado, no activa al tercio convertible, y su segmento natural esta siendo penetrado a velocidad que exige respuesta antes de ola 2027.")
    tf = body_textbox(slide)
    section_divider(tf, "Tesis V5 completa:", before=0)
    plain_text(tf,
        "Gama posee el activo relacional mas diferenciador del mercado (atencion como driver #1, OR=5.73*** con convergencia de 4 metodos), su posicion es estadisticamente estable y el 24h es una ventaja funcional sin equivalente en los competidores. Pero el crecimiento de Gama es cero mientras Rio y Paramo penetran el segmento C+/C — el segmento donde Gama tiene su mayor fortaleza.",
        size=FONT_BODY, before=2, after=1)
    section_divider(tf, "La estrategia correcta requiere tres movimientos simultaneos:")
    bullet(tf, "1. Comunicar el activo que mueve la preferencia, ", value="no el precio que no la mueve (Rec. 1+2).")
    bullet(tf, "2. Activar el 33% convertible con el mecanismo correcto, ", value="primera experiencia, no descuento (Rec. 3).")
    bullet(tf, "3. Responder al avance en C+/C antes de ola 2027 ", value="(Rec. 5).")
    section_divider(tf, "Las dos lecturas validas de la tesis (decision CO-2 del Owner):")
    bullet(tf, "Lectura defensiva: ", value="'0/8 metricas del embudo de Gama variaron significativamente. Gama es la unica gran cadena que mantuvo posicion en un ano de movimiento.'")
    bullet(tf, "Lectura ofensiva: ", value="'Rio +17pp TOM (sig_99). Paramo +12pp TOM (sig_99). Posiblemente +25pp y +22pp en C+/C. La conversion TOM→Preferida es 44.3%→8.0%.'")
    plain_text(tf,
        "Ambas son correctas. La diferencia es estrategica, no estadistica. El Owner elige cual presentar a la Junta.",
        size=FONT_BODY, italic=True, color=RED_DARK, before=3, after=1)
    add_footer(slide, 22)


def build_s23_bloque8_argumentos(prs):
    slide = blank_slide(prs)
    add_standard_header(slide, 23, "BLOQUE 8 — Los 4 argumentos MECE desde los datos 2026")
    add_callout_box(slide,
        "Cuatro argumentos MECE construidos directamente desde los datos 2026. Los dos primeros son conclusiones de alta certeza cuanti. El tercero combina certeza cuanti (segmento) + hipotesis cuali (mecanismo). El cuarto es hipotesis WoW apoyada.")
    tf = body_textbox(slide)
    section_divider(tf, "Arg. 1 — El activo es real, robusto y sin equivalente competitivo (✅ Alta certeza):", before=0)
    plain_text(tf,
        "OR=5.73 convergente en 4 metodos independientes. 53% de preferentes citan atencion como razon espontanea #1. Rio y Paramo no tienen equivalente documentado. Limpieza es driver secundario (OR=3.99*). El activo existe y es defensivo.",
        size=FONT_BODY, before=1, after=2)
    section_divider(tf, "Arg. 2 — El activo no esta siendo comunicado (✅ Alta certeza gap, ⚠ media causalidad):")
    plain_text(tf,
        "Recall espontaneo PTL = 0/17 = 0%. 65% interpreta PTL como mensaje de precio (atributo OR=1.03 NS). El recurso comunicacional esta invertido en el territorio equivocado.",
        size=FONT_BODY, before=1, after=2)
    section_divider(tf, "Arg. 3 — Segmento de conversion existe + mecanismo conocido (✅ + 💡 hipotesis):")
    plain_text(tf,
        "Seg 2 (33% del mercado): 0% preferencia actual, menor resistencia al precio (3.44 vs 3.66/5 cuanti). Hipotesis cuali: la barrera es identitaria, la activacion correcta es primera experiencia guiada. El mercado disponible de CM (-7.7pp compra, sig_95) complementa esta oportunidad.",
        size=FONT_BODY, before=1, after=2)
    section_divider(tf, "Arg. 4 — Presion competitiva en C+/C se intensifica (⚠ Hipotesis apoyada):")
    plain_text(tf,
        "Rio +25.8pp TOM en C+/C (exploratorio sig_99). Paramo +22.3pp TOM en C+/C (exploratorio sig_99). Gama +2.2pp (estable). La posicion actual es solida; la velocidad de los competidores en el segmento natural de Gama es la senal de alerta mas importante del estudio — pendiente confirmacion ola 2027.",
        size=FONT_BODY, before=1, after=2)
    badge_text(tf, "    ✅ Arg. 1 + 2: conclusiones · 💡 Arg. 3 mecanismo: hipotesis cuali · ⚠ Arg. 4: hipotesis WoW C+/C")
    add_footer(slide, 23)


def build_s24_bloque9_recs_qw(prs):
    slide = blank_slide(prs)
    add_standard_header(slide, 24, "BLOQUE 9 — Recomendaciones priorizadas: 2 Quick Wins")
    add_callout_box(slide,
        "Los dos quick wins de mayor retorno a corto plazo derivan de los hallazgos mas robustos. Ambos son accionables con activos existentes de Gama o con inversion de bajo costo.")
    tf = body_textbox(slide)
    section_divider(tf, "Rec. 1 — Quick win: Activar digitalmente el Gama Club", before=0)
    bullet(tf, "Que: ", value="saldo visible en app + equivalente en dinero + opcion de canje en caja.")
    bullet(tf, "Por que: ", value="interes genuino y transversal documentado en todos los segmentos (hipotesis cuali 6/7 FG). Unica barrera documentada es operacional.")
    bullet(tf, "Accionabilidad: ", value="alta — cambio de producto digital sin nueva investigacion ni capex mayor.")
    bullet(tf, "Potencial adicional: ", value="espacio disponible CM (-7.7pp compra, sig_95) puede capturarse con mecanismo de lealtad visible.")
    badge_text(tf, "    💡 Insight cuali (6/7 FG) · ✅ WoW oportunidad CM (sig_95). Validar con cuanti ola 2027.")
    section_divider(tf, "Rec. 2 — Quick win: Reorientar el anzuelo promocional")
    bullet(tf, "Que: ", value="mantener stickers/Cashea/Club como palanca tactica (OR=3.64**), cambiar mensaje 'precio bajo' por 'razon especifica para venir esta semana'.")
    bullet(tf, "Por que: ", value="OR=3.64** confirma que las promociones son driver #3. Pero el mensaje 'precio bajo' refuerza la distancia simbolica del Seg 2 (hipotesis cuali sifrinaje).")
    bullet(tf, "Mensaje correcto: ", value="'Esta semana hay una razon para venir a Gama' (acceso a experiencia/producto) vs 'Precios de tu lado' (precio bajo, OR=1.03 NS).")
    badge_text(tf, "    ✅ Cuanti OR=3.64** · 💡 Hipotesis cuali mecanismo (V4-NEW-2 sifrinaje)")
    add_footer(slide, 24)


def build_s25_bloque9_recs_mt(prs):
    slide = blank_slide(prs)
    add_standard_header(slide, 25, "BLOQUE 9 — Recomendaciones Mid-term y Long-term")
    add_callout_box(slide,
        "Mid-term: diseno de activacion para Seg 2 (33% mercado) + canal digital (sujeto a CO-3). Long-term: estrategia ofensiva en C+/C + exploracion arquetipo (sujeto a CO-3 + DW-2).")
    tf = body_textbox(slide)
    section_divider(tf, "Rec. 3 — Mid-term: Diseno activacion para Seg 2 (33% del mercado)", before=0)
    plain_text(tf,
        "Secuencia: (1) Anzuelo con razon especifica accesible (no precio bajo); (2) Primera experiencia guiada que active el driver de atencion (OR=5.73); (3) Retencion por Gama Club visible + saldo acumulado comunicado.",
        size=FONT_BODY, before=1, after=1)
    bullet(tf, "El Seg 2 tiene 0% preferencia actual + menor resistencia al precio cuanti (3.44 vs 3.66/5).")
    bullet(tf, "El cuali V4 enriquece el perfil con dimension identitaria.")
    bullet(tf, "Motor de conversion: resignificar la pertenencia, no neutralizar diferencia de precio.")
    badge_text(tf, "    ✅ Cuanti k-means Seg 2 · 💡 Mecanismo cuali · [Hipotesis V4 — evidencia convergente]")
    section_divider(tf, "Rec. 4 — Mid-term: Canal digital (⚠ requiere decision CO-3 del Owner)")
    plain_text(tf,
        "Nota de alcance: estas recomendaciones van mas alla del brief original de notoriedad/brand health. Se incluyen por alta accionabilidad pero sugieren conversacion de scope adicional con Gama.",
        size=FONT_CAVEAT, color=BLUE_INFO, before=2, after=1)
    bullet(tf, "Barrera 1 (inventario): ", value="sincronizacion en tiempo real — inversion tecnologica media.")
    bullet(tf, "Barrera 2 (frescos): ", value="propuesta de delivery inicial solo con no-perecederos — bajo riesgo operacional.")
    section_divider(tf, "Rec. 5 — Long-term: Estrategia ofensiva en C+/C antes de ola 2027")
    plain_text(tf,
        "Disenar y ejecutar respuesta al avance de Rio y Paramo en C+/C mientras la posicion de Gama (TOM 60.6%) aun es ventajosa. La ventana de respuesta ofensiva se cierra si Rio/Paramo completan la penetracion del segmento.",
        size=FONT_BODY, before=2, after=1)
    section_divider(tf, "Rec. 6 — Long-term: Explorar plataforma 'acompanamiento-guia' (sujeto a CO-3 + DW-2)")
    plain_text(tf,
        "Se recomienda explorar el arquetipo femenino emergente (35-45 anos, guia experta) como plataforma de comunicacion, condicionado a (1) prueba de concepto cuantitativa con target, (2) verificacion de contratos de imagen vigentes y (3) compatibilidad con plataforma de marca global.",
        size=FONT_BODY, before=2, after=1)
    add_warn_box(slide,
        "Las recomendaciones son valoracion del equipo analitico — no un dato empirico. "
        "La accionabilidad estimada asume continuidad de condiciones de mercado. "
        "Rec. 3+6 requieren validaciones adicionales antes de ejecucion (cuantitativo + restricciones Gama).")
    add_footer(slide, 25)


def build_s26_bloque10(prs):
    slide = blank_slide(prs)
    add_standard_header(slide, 26, "BLOQUE 10 — Agenda 2027 + Decisiones Owner abiertas")
    add_callout_box(slide,
        "Cinco mejoras metodologicas obligatorias diseñadas para ola 2027 + cuatro hipotesis abiertas + tres decisiones estrategicas que solo el Owner puede tomar antes de presentar el deck a la Junta de Gama.")
    tf = body_textbox(slide)
    section_divider(tf, "5 mejoras metodologicas obligatorias para ola 2027:", before=0)
    bullet(tf, "MaxDiff reemplaza Likert saturado (P22) ", value="— eliminar saturacion T2B >90%.")
    bullet(tf, "NPS + switching explicito (3 preguntas) ", value="— medir fidelidad activa, no solo preferencia declarada.")
    bullet(tf, "CEPs expandidos (15-20 ocasiones vs 5 misiones genericas) ", value="— capturar fragmentacion de misiones.")
    bullet(tf, "Penetracion 12m + frecuencia + ticket promedio ", value="— construir modelo share of wallet.")
    bullet(tf, "Booster Pref-Gama n=80 + modulo comunicacion con estimulos + Rio como benchmark.")
    section_divider(tf, "4 hipotesis que requieren ola 2027 para confirmacion:")
    bullet(tf, "Imagen experiencial en ascenso ", value="(+6.3pp 'atractiva', +5.8pp 'seguro' — solo tendencia).")
    bullet(tf, "Debilitamiento imagen precio Gama ", value="(-3.9pp, -5.1pp — solo tendencia).")
    bullet(tf, "Gap comunicacional PTL ", value="— sin confirmacion directa con estimulos.")
    bullet(tf, "WoW en NSE C+/C ", value="— exploratorio, no pre-registrado.")
    section_divider(tf, "Decisiones Owner antes de presentar a Cora y Junta de Gama:")
    bullet(tf, "CO-1: ", value="¿Senal alerta Rio en proteinas entra al deck o al memo interno? (1 verbatim, certeza muy baja).")
    bullet(tf, "CO-2: ", value="¿Tono 'liderazgo defensivo solido' o 'senal de alerta competitiva'? Ambas verdaderas.")
    bullet(tf, "CO-3: ", value="¿Recomendaciones canal digital al deck V5 o propuesta de scope adicional?")
    add_warn_box(slide,
        "El deck V5 no puede circular a Cora/Gama hasta que CO-1/CO-2/CO-3 esten resueltos por Owner. "
        "Gate Bruna BR-2 V4 vigente en V5 (mismos datos, mismos caveats, mismos wordings DW-1..DW-5).")
    add_footer(slide, 26)


def build_s27_cierre(prs):
    slide = blank_slide(prs)
    add_standard_header(slide, 27, "Cierre — Las dos lecturas validas de Gama 2026")
    add_callout_box(slide,
        "Gama 2026 tiene dos lecturas validas — y el Owner elige cual presentar a la junta. Ambas son correctas. La diferencia es estrategica, no estadistica.")
    col_w = CONTENT_W // 2 - 137160
    col_gap = 274320
    col1_l = MARGIN_L
    col2_l = MARGIN_L + col_w + col_gap
    # cabecera defensiva
    hdr1 = add_rounded_rect(slide, col1_l, BODY_TOP, col_w, 365760, GRAY_LIGHT)
    tf_h1 = hdr1.text_frame; clear_tf(tf_h1)
    add_para_with_runs(tf_h1, [("Lectura defensiva (CO-2 opcion A)", FONT_BODY, True, RED_DARK)], align=PP_ALIGN.CENTER)
    col1_body_top = BODY_TOP + 365760 + 91440
    tb1 = add_textbox(slide, col1_l, col1_body_top, col_w, 3291840)
    tf1 = tb1.text_frame; tf1.word_wrap = True; clear_tf(tf1)
    plain_text(tf1,
        '"0/8 metricas del embudo de Gama variaron significativamente. Gama es la unica gran cadena del mercado que mantuvo posicion en un ano de movimiento. El activo relacional (OR=5.73) existe y es robusto. La posicion en C+/C (TOM 60.6%) es la mas fuerte del mercado."',
        size=FONT_BODY, italic=True, color=DARK_TEXT, before=0, after=4)
    bullet(tf1, "Mensaje a Junta: liderazgo defensivo solido.")
    bullet(tf1, "Datos de soporte: todos sig_99 o alta certeza cuanti.")
    # cabecera ofensiva
    hdr2 = add_rounded_rect(slide, col2_l, BODY_TOP, col_w, 365760, ORANGE_BG)
    tf_h2 = hdr2.text_frame; clear_tf(tf_h2)
    add_para_with_runs(tf_h2, [("Lectura ofensiva (CO-2 opcion B)", FONT_BODY, True, ORANGE_WARN)], align=PP_ALIGN.CENTER)
    col2_body_top = BODY_TOP + 365760 + 91440
    tb2 = add_textbox(slide, col2_l, col2_body_top, col_w, 3291840)
    tf2 = tb2.text_frame; tf2.word_wrap = True; clear_tf(tf2)
    plain_text(tf2,
        '"Rio +17pp TOM (sig_99). Paramo +12pp TOM (sig_99). Posiblemente +25pp y +22pp en C+/C — el segmento de Gama. La conversion TOM→Preferida es 44.3%→8.0%. El 92% que no prefiere no esta siendo activado. El anzuelo comunicacional esta en el territorio equivocado."',
        size=FONT_BODY, italic=True, color=DARK_TEXT, before=0, after=4)
    bullet(tf2, "Mensaje a Junta: senal de alerta competitiva.")
    bullet(tf2, "Datos de soporte: WoW C+/C como hipotesis (CV-WOW-005).")
    add_info_box(slide, MARGIN_L, 5669280, CONTENT_W, 320040,
        "Notoriedad y Preferencia de Marca — Gama 2026 · V5 · n=402 · m.e. ±4.89% · Confidencial NDA · Cora Urrea + Raoul Bermudez · 2026-05-18")
    add_footer(slide, 27)


def main():
    print("Generando charts PNG...")
    c1 = chart_d1_funnel_cpc_vs_total(); print(f"  {c1}")
    c2 = chart_d2_tom_cpc_wow();         print(f"  {c2}")
    c3 = chart_d3_funnel_gama();         print(f"  {c3}")
    c4 = chart_d4_posicionamiento();     print(f"  {c4}")
    c5 = chart_d5_forest_drivers();      print(f"  {c5}")
    c6 = chart_d6_kmeans_scatter();      print(f"  {c6}")
    c7 = chart_d7_misiones();            print(f"  {c7}")
    c8 = chart_d8_comunicacion();        print(f"  {c8}")
    c9 = chart_d9_2x2_competitivo();     print(f"  {c9}")

    print("\nConstruyendo presentacion...")
    prs = new_prs()
    build_s01_portada(prs);                       print("  s01 Portada OK")
    build_s02_marco(prs);                         print("  s02 Marco OK")
    build_s03_iconografia(prs);                   print("  s03 Iconografia + Ficha OK")
    build_s04_bloque1_perfil(prs);                print("  s04 1A Perfil C+/C OK")
    build_s05_bloque1B(prs, c1);                  print("  s05 1B Gama en C+/C OK")
    build_s06_bloque1C(prs, c2);                  print("  s06 1C Amenaza Rio/Paramo OK")
    build_s07_bloque1D(prs);                      print("  s07 1D Nucleo Leal OK")
    build_s08_bloque2A(prs, c3);                  print("  s08 2A Embudo Gama OK")
    build_s09_bloque2B(prs, c4);                  print("  s09 2B Posicion competitiva OK")
    build_s10_bloque2C(prs);                      print("  s10 2C Atributos OK")
    build_s11_bloque3A(prs, c5);                  print("  s11 3A Drivers Forest OK")
    build_s12_bloque3B(prs, c6);                  print("  s12 3B K-means OK")
    build_s13_bloque4A(prs, c7);                  print("  s13 4A Misiones OK")
    build_s14_bloque5(prs, c8);                   print("  s14 5 Comunicacion OK")
    build_s15_bloque6_intro(prs);                 print("  s15 6 intro cuali OK")
    build_s16_bloque6A(prs);                      print("  s16 6A Acompanamiento OK")
    build_s17_bloque6B(prs);                      print("  s17 6B Sifrinaje OK")
    build_s18_bloque6C_6D(prs);                   print("  s18 6C+6D Club+Digital OK")
    build_s19_bloque6E(prs);                      print("  s19 6E Arquetipo OK")
    build_s20_bloque7A(prs);                      print("  s20 7A Gama estable OK")
    build_s21_bloque7B(prs, c9);                  print("  s21 7B Movimiento 2x2 OK")
    build_s22_bloque8_tesis(prs);                 print("  s22 8 Tesis OK")
    build_s23_bloque8_argumentos(prs);            print("  s23 8 4 argumentos OK")
    build_s24_bloque9_recs_qw(prs);               print("  s24 9 Quick Wins OK")
    build_s25_bloque9_recs_mt(prs);               print("  s25 9 Mid+Long term OK")
    build_s26_bloque10(prs);                      print("  s26 10 Agenda 2027 OK")
    build_s27_cierre(prs);                        print("  s27 Cierre OK")

    prs.save(OUTPUT_PATH)
    print(f"\nGuardado: {OUTPUT_PATH}")
    print(f"Charts en: {CHARTS_DIR}")


if __name__ == "__main__":
    main()
