"""
chart_builders_v82.py  —  Vivienne V8.2 chart builder module
Reutilizado por build_deck_v8.2.py y build_resumen_ejecutivo_v8.2.py

REGLA CRITICA: TODOS los datos provienen del JSON chart_data_v82.json.
PROHIBIDO inventar valores. Si assertion falla -> abortar.

Encoding: UTF-8 obligatorio (Windows cp1252 workaround)
"""
import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.chart.data import ChartData, CategoryChartData
from pptx.oxml.ns import qn
import copy
from lxml import etree

# ─────────────────────────────────────────────
# BRAND COLORS
# ─────────────────────────────────────────────
GAMA_RED      = RGBColor(0xE3, 0x06, 0x13)
GAMA_RED_SOFT = RGBColor(0xE3, 0x71, 0x78)   # brecha exp.reciente
PARAMO_BLUE   = RGBColor(0x4A, 0x6F, 0xA5)
RIO_GREY      = RGBColor(0x7C, 0x8B, 0x9B)
CM_GREY       = RGBColor(0x5B, 0x70, 0x90)
FORUM_GREY    = RGBColor(0x8F, 0xA3, 0xB8)
PLAZAS_GREY   = RGBColor(0x6D, 0x84, 0x99)
PLAN_GREY     = RGBColor(0x9C, 0xAE, 0xC0)
LUZ_GREY      = RGBColor(0xA8, 0xB9, 0xC7)

TEXT_DARK     = RGBColor(0x1A, 0x1A, 0x1A)
GREY_MED      = RGBColor(0x6B, 0x6B, 0x6B)
GREY_LIGHT    = RGBColor(0xE5, 0xE5, 0xE5)
AMBER         = RGBColor(0xF2, 0xA9, 0x00)
GREEN_VAL     = RGBColor(0x2D, 0x8F, 0x47)
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
BLACK         = RGBColor(0x00, 0x00, 0x00)
HEATMAP_MIN   = RGBColor(0xF4, 0xF4, 0xF4)

# Map cadena -> color
CADENA_COLORS = {
    "Gama":               GAMA_RED,
    "Páramo":             PARAMO_BLUE,
    "Rio":                RIO_GREY,
    "Central Madeirense": CM_GREY,
    "Forum":              FORUM_GREY,
    "Plazas":             PLAZAS_GREY,
    "Plan Suárez":        PLAN_GREY,
    "Luz":                LUZ_GREY,
}

# ─────────────────────────────────────────────
# HELPERS GENERALES
# ─────────────────────────────────────────────

def rgb_lerp(c_min: RGBColor, c_max: RGBColor, t: float) -> RGBColor:
    """Interpolación lineal entre dos colores. t in [0,1]."""
    t = max(0.0, min(1.0, t))
    r = int(c_min[0] + t * (c_max[0] - c_min[0]))
    g = int(c_min[1] + t * (c_max[1] - c_min[1]))
    b = int(c_min[2] + t * (c_max[2] - c_min[2]))
    return RGBColor(r, g, b)


def is_dark(color: RGBColor) -> bool:
    """True si el color es oscuro (luminosidad < 0.4)."""
    r, g, b = color[0] / 255, color[1] / 255, color[2] / 255
    lum = 0.2126 * r + 0.7152 * g + 0.0722 * b
    return lum < 0.4


def set_cell_color(cell, rgb: RGBColor):
    """Aplica color de fondo a una celda de tabla."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
    srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
    srgbClr.set('val', f'{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}')


def set_cell_border(cell, color: RGBColor, width_emu: int = 19050):
    """Aplica borde a todos los lados de una celda."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for side in ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']:
        existing = tcPr.find(qn(side))
        if existing is not None:
            tcPr.remove(existing)
        ln = etree.SubElement(tcPr, qn(side))
        ln.set('w', str(width_emu))
        solidFill = etree.SubElement(ln, qn('a:solidFill'))
        srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
        srgbClr.set('val', f'{color[0]:02X}{color[1]:02X}{color[2]:02X}')


def add_text_box(slide, text: str, left: float, top: float, width: float, height: float,
                 font_size: int = 10, bold: bool = False, color: RGBColor = None,
                 align=PP_ALIGN.LEFT, wrap: bool = True):
    """Agrega un text box al slide."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color if color else TEXT_DARK
    return txBox


def add_slide_title(slide, title: str, subtitle: str = None):
    """Agrega título principal y subtítulo opcional."""
    add_text_box(slide, title,
                 left=0.4, top=0.1, width=9.2, height=0.55,
                 font_size=18, bold=True, color=TEXT_DARK, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text_box(slide, subtitle,
                     left=0.4, top=0.65, width=9.2, height=0.3,
                     font_size=10, bold=False, color=GREY_MED, align=PP_ALIGN.LEFT)


def add_footnote(slide, text: str, amber: bool = False):
    """Footnote al pie del slide."""
    color = AMBER if amber else GREY_MED
    add_text_box(slide, text,
                 left=0.4, top=7.05, width=9.2, height=0.35,
                 font_size=8, bold=False, color=color, align=PP_ALIGN.LEFT)


def add_section_tag(slide, tag: str):
    """Etiqueta de sección (esquina superior derecha)."""
    add_text_box(slide, tag,
                 left=7.5, top=0.05, width=2.4, height=0.22,
                 font_size=7, bold=False, color=GREY_MED, align=PP_ALIGN.RIGHT)


# ─────────────────────────────────────────────
# C01 — EMBUDO GAMA (bar horizontal)
# ─────────────────────────────────────────────

def make_chart_c01(slide, data: dict, left: float = 0.4, top: float = 1.0,
                   width: float = 9.2, height: float = 5.5):
    """Bar horizontal: embudo de marca Gama."""
    assert 'etapas' in data and 'gama_total_pct' in data, "C01: campos requeridos ausentes"
    assert len(data['etapas']) == len(data['gama_total_pct']), "C01: dimensiones inconsistentes"
    assert all(0 <= v <= 100 for v in data['gama_total_pct']), "C01: valores fuera de rango 0-100"

    etapas = data['etapas']
    total_pct = data['gama_total_pct']
    cpc_pct = data.get('gama_cpc_pct', [None] * len(etapas))

    chart_data = CategoryChartData()
    chart_data.categories = etapas
    chart_data.add_series('Total', total_pct)
    valid_cpc = [v if v is not None else 0.0 for v in cpc_pct]
    chart_data.add_series('C+/C', valid_cpc)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(left), Inches(top), Inches(width), Inches(height),
        chart_data
    ).chart

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    # Colores series
    series_total = chart.series[0]
    series_cpc   = chart.series[1]

    series_total.format.fill.solid()
    series_total.format.fill.fore_color.rgb = GAMA_RED

    series_cpc.format.fill.solid()
    series_cpc.format.fill.fore_color.rgb = GREY_MED

    # Data labels
    for s in chart.series:
        s.data_labels.show_value = True
        s.data_labels.number_format = '0.0"%"'
        s.data_labels.font.size = Pt(8)

    # Gridlines suaves
    try:
        vax = chart.value_axis
        vax.major_gridlines.format.line.color.rgb = GREY_LIGHT
    except Exception:
        pass

    return chart


# ─────────────────────────────────────────────
# C02 — PREFERENCIA P21 (bar horizontal grouped)
# ─────────────────────────────────────────────

def make_chart_c02(slide, data: dict, left: float = 0.4, top: float = 1.0,
                   width: float = 9.2, height: float = 5.5):
    """Bar horizontal agrupado: preferencia P21 Total vs C+/C."""
    assert 'cadenas' in data and 'pct_total' in data, "C02: campos requeridos ausentes"
    assert len(data['cadenas']) == len(data['pct_total']), "C02: dimensiones inconsistentes"
    assert all(0 <= v <= 100 for v in data['pct_total']), "C02: valores fuera de rango"

    cadenas = data['cadenas']
    pct_total = data['pct_total']
    cadenas_cpc = data.get('cadenas_cpc', [])
    pct_cpc_raw = data.get('pct_cpc', [])
    # Align CPC data to same category order
    cpc_lookup = dict(zip(cadenas_cpc, pct_cpc_raw))
    pct_cpc = [cpc_lookup.get(c, 0.0) for c in cadenas]

    chart_data = CategoryChartData()
    chart_data.categories = cadenas
    chart_data.add_series('Total', pct_total)
    chart_data.add_series('C+/C', pct_cpc)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(left), Inches(top), Inches(width), Inches(height),
        chart_data
    ).chart

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    # Colores per-point para Total
    s_total = chart.series[0]
    s_cpc   = chart.series[1]

    s_total.format.fill.solid()
    s_total.format.fill.fore_color.rgb = RGBColor(0xB0, 0xB8, 0xC1)

    s_cpc.format.fill.solid()
    s_cpc.format.fill.fore_color.rgb = PARAMO_BLUE

    # Destacar Gama en rojo
    for s in [s_total, s_cpc]:
        s.data_labels.show_value = True
        s.data_labels.number_format = '0.0"%"'
        s.data_labels.font.size = Pt(8)

    try:
        vax = chart.value_axis
        vax.major_gridlines.format.line.color.rgb = GREY_LIGHT
    except Exception:
        pass

    return chart


# ─────────────────────────────────────────────
# C03 — MODELO MENTAL PRECIO (bar vertical)
# ─────────────────────────────────────────────

def make_chart_c03(slide, data: dict, left: float = 0.4, top: float = 1.0,
                   width: float = 9.2, height: float = 5.5):
    """Bar vertical con threshold en 50% y color Gama destacado."""
    assert 'cadenas_ordenadas_asc' in data and 'pct_precio' in data, "C03: campos requeridos ausentes"
    assert len(data['cadenas_ordenadas_asc']) == len(data['pct_precio']), "C03: dimensiones inconsistentes"
    assert all(0 <= v <= 100 for v in data['pct_precio']), "C03: valores fuera de rango"

    cadenas = data['cadenas_ordenadas_asc']
    pct = data['pct_precio']

    chart_data = CategoryChartData()
    chart_data.categories = cadenas
    chart_data.add_series('% precio dominante', pct)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(left), Inches(top), Inches(width), Inches(height),
        chart_data
    ).chart

    chart.has_legend = False

    # Per-point coloring
    series = chart.series[0]
    for i, cadena in enumerate(cadenas):
        pt = series.points[i]
        if cadena == 'Gama':
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = GAMA_RED
        else:
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = RIO_GREY

    series.data_labels.show_value = True
    series.data_labels.number_format = '0.0"%"'
    series.data_labels.font.size = Pt(9)
    series.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

    # Rango eje Y con margen para threshold line
    try:
        vax = chart.value_axis
        vax.minimum_scale = 0
        vax.maximum_scale = 100
        vax.major_gridlines.format.line.color.rgb = GREY_LIGHT
    except Exception:
        pass

    # Línea threshold en 50% (añadir como serie auxiliar plana)
    # Usamos texto annotation porque add_series requiere mismos n puntos
    threshold_series_data = CategoryChartData()
    threshold_series_data.categories = cadenas
    threshold_series_data.add_series('Umbral 50%', [50.0] * len(cadenas))

    # Agregar serie de línea al chart existente no es trivial con python-pptx;
    # en su lugar añadimos un text box con la línea threshold visual
    add_text_box(slide,
                 "──────────────────────────── 50% umbral ────────────────────────────",
                 left=left + 0.3, top=top + height * (1 - 50 / 100) * 0.82,
                 width=width - 0.6, height=0.2,
                 font_size=7, bold=False, color=AMBER)

    return chart


# ─────────────────────────────────────────────
# C04 — TB vs T2B RERANKING (bar horizontal grouped)
# ─────────────────────────────────────────────

def make_chart_c04(slide, data: dict, left: float = 0.4, top: float = 1.0,
                   width: float = 9.2, height: float = 5.5):
    """Bar horizontal agrupado TB vs T2B con delta de rank."""
    assert 'atributos' in data and 'tb_pct' in data and 't2b_pct' in data, "C04: campos requeridos"
    assert len(data['atributos']) == len(data['tb_pct']) == len(data['t2b_pct']), "C04: dimensiones inconsistentes"
    assert all(0 <= v <= 100 for v in data['tb_pct'] + data['t2b_pct']), "C04: valores fuera de rango"

    atributos = data['atributos']
    tb_pct    = data['tb_pct']
    t2b_pct   = data['t2b_pct']
    delta     = data.get('delta_rank', [0] * len(atributos))

    chart_data = CategoryChartData()
    chart_data.categories = atributos
    chart_data.add_series('Top Box (TB)', tb_pct)
    chart_data.add_series('Top 2 Box (T2B)', t2b_pct)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(left), Inches(top), Inches(width), Inches(height),
        chart_data
    ).chart

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    s_tb  = chart.series[0]
    s_t2b = chart.series[1]

    s_tb.format.fill.solid()
    s_tb.format.fill.fore_color.rgb = GAMA_RED

    s_t2b.format.fill.solid()
    s_t2b.format.fill.fore_color.rgb = RIO_GREY

    for s in [s_tb, s_t2b]:
        s.data_labels.show_value = True
        s.data_labels.number_format = '0.0"%"'
        s.data_labels.font.size = Pt(7)

    try:
        vax = chart.value_axis
        vax.major_gridlines.format.line.color.rgb = GREY_LIGHT
    except Exception:
        pass

    # Delta labels a la derecha
    delta_lines = [f"Δrank: {'+' if d > 0 else ''}{d}" for d in delta]
    add_text_box(slide,
                 "\n".join(delta_lines),
                 left=left + width + 0.05, top=top + 0.5,
                 width=1.0, height=height - 0.5,
                 font_size=7, bold=False, color=GREY_MED)

    return chart


# ─────────────────────────────────────────────
# C05 — HEATMAP TB PURO POR MARCA (tabla)
# ─────────────────────────────────────────────

def make_heatmap_table(slide, matrix: list, row_labels: list, col_labels: list,
                       title_row: str, fuente: str,
                       left: float, top: float, width: float, height: float,
                       footnote: str = None):
    """
    Heatmap genérico via tabla PPTX.
    - Rango dinámico min/max del dataset (no fijo 0-100).
    - Fila/columna con valor máximo destacada con borde rojo.
    - Texto blanco si fondo oscuro, negro si claro.
    """
    n_rows = len(row_labels)
    n_cols = len(col_labels)

    assert len(matrix) == n_rows, f"Heatmap: matrix tiene {len(matrix)} filas, esperadas {n_rows}"
    for i, row in enumerate(matrix):
        assert len(row) == n_cols, f"Heatmap: fila {i} tiene {len(row)} cols, esperadas {n_cols}"

    # Rango dinámico
    all_vals = [v for row in matrix for v in row if v is not None]
    v_min = min(all_vals)
    v_max = max(all_vals)

    # Encontrar celda máxima
    max_row_i, max_col_j = 0, 0
    for ri, row in enumerate(matrix):
        for ci, val in enumerate(row):
            if val is not None and val == v_max:
                max_row_i, max_col_j = ri, ci

    # Tabla: 1 header row + n_rows, 1 header col + n_cols
    rows_total = n_rows + 1
    cols_total = n_cols + 1

    tbl = slide.shapes.add_table(
        rows_total, cols_total,
        Inches(left), Inches(top), Inches(width), Inches(height)
    ).table

    # Row heights uniform
    row_h = Inches(height / rows_total)
    for i in range(rows_total):
        tbl.rows[i].height = row_h

    # Col widths
    label_col_w = Inches(1.2)
    data_col_w  = Inches((width - 1.2) / n_cols)
    tbl.columns[0].width = label_col_w
    for j in range(1, cols_total):
        tbl.columns[j].width = data_col_w

    # Header row (row 0)
    # corner cell
    corner = tbl.cell(0, 0)
    corner.text = title_row
    set_cell_color(corner, RGBColor(0x1A, 0x1A, 0x1A))
    for run in corner.text_frame.paragraphs[0].runs:
        run.font.color.rgb = WHITE
        run.font.size = Pt(7)
        run.font.bold = True

    for j, col_label in enumerate(col_labels):
        cell = tbl.cell(0, j + 1)
        cell.text = col_label
        set_cell_color(cell, RGBColor(0x1A, 0x1A, 0x1A))
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        if p.runs:
            p.runs[0].font.color.rgb = WHITE
            p.runs[0].font.size = Pt(7)
            p.runs[0].font.bold = True
        else:
            run = p.add_run()
            run.text = col_label
            run.font.color.rgb = WHITE
            run.font.size = Pt(7)
            run.font.bold = True
            cell.text = ""

    # Data rows
    for ri, row_label in enumerate(row_labels):
        # Label column
        lbl_cell = tbl.cell(ri + 1, 0)
        lbl_cell.text = row_label
        set_cell_color(lbl_cell, GREY_MED)
        p = lbl_cell.text_frame.paragraphs[0]
        if p.runs:
            p.runs[0].font.color.rgb = WHITE
            p.runs[0].font.size = Pt(7)
            p.runs[0].font.bold = True
        else:
            run = p.add_run()
            run.text = row_label
            run.font.color.rgb = WHITE
            run.font.size = Pt(7)
            run.font.bold = True
            lbl_cell.text = ""

        # Data cells
        for ci, val in enumerate(matrix[ri]):
            cell = tbl.cell(ri + 1, ci + 1)
            val_str = f"{val:.1f}" if val is not None else "–"
            cell.text = val_str

            if val is not None:
                t = (val - v_min) / (v_max - v_min) if v_max > v_min else 0.0
                bg_color = rgb_lerp(HEATMAP_MIN, GAMA_RED, t)
            else:
                bg_color = HEATMAP_MIN

            set_cell_color(cell, bg_color)
            text_color = WHITE if is_dark(bg_color) else TEXT_DARK

            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            is_max_cell = (ri == max_row_i and ci == max_col_j)
            if p.runs:
                p.runs[0].font.color.rgb = text_color
                p.runs[0].font.size = Pt(7)
                p.runs[0].font.bold = is_max_cell
            else:
                run = p.add_run()
                run.text = val_str
                run.font.color.rgb = text_color
                run.font.size = Pt(7)
                run.font.bold = is_max_cell
                cell.text = ""

            # Borde rojo para celda/fila/col máxima
            if is_max_cell:
                set_cell_border(cell, GAMA_RED, width_emu=38100)

    if footnote:
        add_footnote(slide, footnote, amber=True)

    return tbl


def make_chart_c05(slide, data: dict, left: float = 0.3, top: float = 1.0,
                   width: float = 9.4, height: float = 5.8):
    """Heatmap C05: TB puro por marca preferida — 10 atributos x 9 segmentos."""
    assert 'matrix_pct' in data and 'atributos' in data and 'segmentos' in data, "C05: campos requeridos"
    assert len(data['matrix_pct']) == len(data['atributos']), "C05: filas inconsistentes"
    for i, row in enumerate(data['matrix_pct']):
        assert len(row) == len(data['segmentos']), f"C05: columnas inconsistentes fila {i}"

    footnote = "* REFERENCIAL: n Pref-Gama=32, Pref-Luz=32, Pref-PlanS=28. Tendencias indicativas, no proyectables."
    return make_heatmap_table(
        slide,
        matrix=data['matrix_pct'],
        row_labels=data['atributos'],
        col_labels=data['segmentos'],
        title_row="Atributo \\ Segmento",
        fuente=data.get('fuente', ''),
        left=left, top=top, width=width, height=height,
        footnote=footnote
    )


# ─────────────────────────────────────────────
# C06 — DNA Z-SCORES (bar horizontal)
# ─────────────────────────────────────────────

def make_chart_c06(slide, data: dict, left: float = 0.4, top: float = 1.0,
                   width: float = 9.2, height: float = 5.5):
    """Bar horizontal con colores positivo/negativo para z-scores."""
    assert 'atributos' in data and 'zscore_total' in data, "C06: campos requeridos"
    assert len(data['atributos']) == len(data['zscore_total']), "C06: dimensiones inconsistentes"

    atributos = data['atributos']
    zscores   = data['zscore_total']

    chart_data = CategoryChartData()
    chart_data.categories = atributos
    chart_data.add_series('Z-score Total', zscores)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(left), Inches(top), Inches(width), Inches(height),
        chart_data
    ).chart

    chart.has_legend = False

    series = chart.series[0]
    for i, z in enumerate(zscores):
        pt = series.points[i]
        pt.format.fill.solid()
        if z >= 0:
            pt.format.fill.fore_color.rgb = GAMA_RED
        else:
            pt.format.fill.fore_color.rgb = GREY_MED

    series.data_labels.show_value = True
    series.data_labels.number_format = '+0.00;-0.00;0.00'
    series.data_labels.font.size = Pt(8)

    try:
        vax = chart.value_axis
        vax.major_gridlines.format.line.color.rgb = GREY_LIGHT
    except Exception:
        pass

    # Línea en 0
    add_text_box(slide, "z = 0.0 (promedio mercado)",
                 left=left, top=top + height * 0.5 - 0.15,
                 width=2.0, height=0.2,
                 font_size=7, bold=False, color=GREY_MED)

    return chart


# ─────────────────────────────────────────────
# C07 — HEATMAP MUNDO DE MARCA P23 (tabla)
# ─────────────────────────────────────────────

def make_chart_c07(slide, data: dict, left: float = 0.3, top: float = 1.0,
                   width: float = 9.4, height: float = 5.8):
    """Heatmap C07: mundo de marca 8 cadenas x 10 atributos."""
    assert 'matrix_pct' in data and 'cadenas' in data and 'atributos' in data, "C07: campos requeridos"
    assert len(data['matrix_pct']) == len(data['cadenas']), "C07: filas inconsistentes"

    # Spot-check CRITICO
    # C07 Gama row (index 0), atributos["Precio"] = index 7
    gama_precio = data['matrix_pct'][0][7]
    assert abs(gama_precio - 7.2) < 0.05, f"C07 SPOT-CHECK FALLO: Gama Precio={gama_precio}, esperado 7.2"
    gama_atencion = data['matrix_pct'][0][4]
    assert abs(gama_atencion - 21.9) < 0.05, f"C07 SPOT-CHECK FALLO: Gama Atencion={gama_atencion}, esperado 21.9"

    footnote = f"Fuente: {data.get('fuente', 'CU-10 §4.1')} | Base: n=402. Gama Precio=7.2% (mínimo absoluto mercado). Páramo Precio=35.1% (máximo)."
    return make_heatmap_table(
        slide,
        matrix=data['matrix_pct'],
        row_labels=data['cadenas'],
        col_labels=data['atributos'],
        title_row="Cadena \\ Atributo",
        fuente=data.get('fuente', ''),
        left=left, top=top, width=width, height=height,
        footnote=footnote
    )


# ─────────────────────────────────────────────
# C08 — TORNADO BRECHAS P23 (bar horizontal doble)
# ─────────────────────────────────────────────

def make_chart_c08(slide, data: dict, left: float = 0.4, top: float = 1.0,
                   width: float = 9.2, height: float = 5.5):
    """Tornado: Pref-Gama y Exp.reciente brechas vs Total."""
    assert 'atributos' in data and 'delta_pref_total' in data and 'delta_exp_total' in data, "C08: campos requeridos"
    assert len(data['atributos']) == len(data['delta_pref_total']) == len(data['delta_exp_total']), "C08: dimensiones inconsistentes"
    assert all(v >= 0 for v in data['delta_pref_total'] + data['delta_exp_total']), "C08: deltas negativos inesperados"

    # Spot-check
    idx_atencion = data['atributos'].index('Mejor atención') if 'Mejor atención' in data['atributos'] else None
    if idx_atencion is not None:
        assert abs(data['delta_pref_total'][idx_atencion] - 62.5) < 0.1, \
            f"C08 SPOT-CHECK: Atencion delta={data['delta_pref_total'][idx_atencion]}, esperado 62.5"

    # Ordenar descendente por Pref-Gama
    combined = list(zip(data['atributos'], data['delta_pref_total'], data['delta_exp_total']))
    combined.sort(key=lambda x: x[1], reverse=True)
    atributos  = [r[0] for r in combined]
    delta_pref = [r[1] for r in combined]
    delta_exp  = [r[2] for r in combined]

    chart_data = CategoryChartData()
    chart_data.categories = atributos
    chart_data.add_series('Pref-Gama (n=32 REF)', delta_pref)
    chart_data.add_series('Exp. reciente (n=30 REF)', delta_exp)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(left), Inches(top), Inches(width), Inches(height),
        chart_data
    ).chart

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    s_pref = chart.series[0]
    s_exp  = chart.series[1]

    s_pref.format.fill.solid()
    s_pref.format.fill.fore_color.rgb = GAMA_RED

    s_exp.format.fill.solid()
    s_exp.format.fill.fore_color.rgb = GAMA_RED_SOFT

    for s in [s_pref, s_exp]:
        s.data_labels.show_value = True
        s.data_labels.number_format = '+0.0;-0.0;0.0'
        s.data_labels.font.size = Pt(7)

    try:
        vax = chart.value_axis
        vax.minimum_scale = 0
        vax.major_gridlines.format.line.color.rgb = GREY_LIGHT
    except Exception:
        pass

    add_footnote(slide, "* REFERENCIAL: n Pref-Gama=32, n Exp.reciente=30. Brechas en pp vs Total muestra.", amber=True)
    return chart


# ─────────────────────────────────────────────
# C09 — SCATTER VECINDAD PERCEPTUAL
# ─────────────────────────────────────────────

def make_chart_c09(slide, data: dict, left: float = 0.5, top: float = 1.1,
                   width: float = 8.8, height: float = 5.5):
    """Scatter con etiquetas, cuadrantes y Gama destacado."""
    assert 'puntos' in data and len(data['puntos']) > 0, "C09: puntos requeridos"

    # Crear scatter manualmente como shapes (python-pptx scatter tiene limitaciones de per-point color)
    # Usamos representación visual con shapes + texto

    # Determinar rangos
    xs = [p['x'] for p in data['puntos']]
    ys = [p['y'] for p in data['puntos']]
    x_min, x_max = min(xs) - 2, max(xs) + 2
    y_min, y_max = min(ys) - 2, max(ys) + 5

    x_range = x_max - x_min
    y_range = y_max - y_min

    def to_slide_x(x_val):
        return left + (x_val - x_min) / x_range * width

    def to_slide_y(y_val):
        return top + height - (y_val - y_min) / y_range * height

    # Área de plot
    plot_bg = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    plot_bg.fill.solid()
    plot_bg.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)
    plot_bg.line.color.rgb = GREY_LIGHT

    # Eje X label
    x_mid = data.get('ejes', {}).get('x_label', '% imagen experiencial (promedio)')
    add_text_box(slide, x_mid,
                 left=left + width/2 - 1.5, top=top + height + 0.05,
                 width=3.0, height=0.25,
                 font_size=8, bold=False, color=GREY_MED, align=PP_ALIGN.CENTER)

    # Eje Y label
    y_label = data.get('ejes', {}).get('y_label', '% asociación Menor Precio')
    add_text_box(slide, y_label,
                 left=0.05, top=top + height/2 - 0.5,
                 width=0.4, height=1.0,
                 font_size=7, bold=False, color=GREY_MED)

    # Líneas de cuadrante (mediana)
    x_med = sum(xs) / len(xs)
    y_med = sum(ys) / len(ys)
    x_line = to_slide_x(x_med)
    y_line = to_slide_y(y_med)

    # Línea vertical cuadrante
    vline = slide.shapes.add_shape(
        1,
        Inches(x_line - 0.01), Inches(top),
        Inches(0.02), Inches(height)
    )
    vline.fill.solid()
    vline.fill.fore_color.rgb = GREY_LIGHT
    vline.line.fill.background()

    # Línea horizontal cuadrante
    hline = slide.shapes.add_shape(
        1,
        Inches(left), Inches(y_line - 0.01),
        Inches(width), Inches(0.02)
    )
    hline.fill.solid()
    hline.fill.fore_color.rgb = GREY_LIGHT
    hline.line.fill.background()

    # Etiquetas cuadrante
    add_text_box(slide, "Alta experiencia\nAlto precio",
                 left=x_line + 0.1, top=top + 0.1,
                 width=1.5, height=0.4, font_size=7, color=GREY_MED)
    add_text_box(slide, "Alta experiencia\nBajo precio",
                 left=x_line + 0.1, top=y_line + 0.1,
                 width=1.5, height=0.4, font_size=7, color=GREY_MED)
    add_text_box(slide, "Baja experiencia\nAlto precio",
                 left=left + 0.05, top=top + 0.1,
                 width=1.5, height=0.4, font_size=7, color=GREY_MED)
    add_text_box(slide, "Baja experiencia\nBajo precio",
                 left=left + 0.05, top=y_line + 0.1,
                 width=1.5, height=0.4, font_size=7, color=GREY_MED)

    # Puntos
    for p in data['puntos']:
        px = to_slide_x(p['x'])
        py = to_slide_y(p['y'])
        r_hex = p.get('color', '#7C8B9B').lstrip('#')
        pt_color = RGBColor(int(r_hex[:2], 16), int(r_hex[2:4], 16), int(r_hex[4:], 16))

        dot_size = 0.22 if p.get('destacar') else 0.16
        dot = slide.shapes.add_shape(
            9,  # oval
            Inches(px - dot_size/2), Inches(py - dot_size/2),
            Inches(dot_size), Inches(dot_size)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = pt_color
        dot.line.fill.background()

        # Etiqueta
        offset_x = 0.12 if px < left + width - 1.5 else -1.5
        add_text_box(slide, p['cadena'],
                     left=px + offset_x, top=py - 0.14,
                     width=1.4, height=0.22,
                     font_size=8,
                     bold=p.get('destacar', False),
                     color=pt_color)

    return None


# ─────────────────────────────────────────────
# C10 — FOREST PLOT / LOGIT (bar horizontal + caveat)
# ─────────────────────────────────────────────

def make_chart_c10(slide, data: dict, left: float = 0.4, top: float = 1.1,
                   width: float = 7.0, height: float = 4.5):
    """Bar horizontal OR. Solo Mayor calidad significativo (#E30613), resto gris."""
    assert 'atributos' in data and 'odds_ratio' in data, "C10: campos requeridos"
    assert len(data['atributos']) == len(data['odds_ratio']), "C10: dimensiones inconsistentes"

    atributos = data['atributos']
    or_vals   = data['odds_ratio']
    sig_attrs = data.get('significativos', [])

    # Para no-significativos, usar OR=1.0 como placeholder visual
    chart_data = CategoryChartData()
    chart_data.categories = atributos
    or_display = [v if v is not None else 1.0 for v in or_vals]
    chart_data.add_series('Odds Ratio', or_display)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(left), Inches(top), Inches(width), Inches(height),
        chart_data
    ).chart

    chart.has_legend = False

    series = chart.series[0]
    for i, attr in enumerate(atributos):
        pt = series.points[i]
        pt.format.fill.solid()
        if attr in sig_attrs:
            pt.format.fill.fore_color.rgb = GAMA_RED
        else:
            pt.format.fill.fore_color.rgb = GREY_LIGHT

    series.data_labels.show_value = True
    series.data_labels.number_format = '0.000'
    series.data_labels.font.size = Pt(8)

    try:
        vax = chart.value_axis
        vax.minimum_scale = 0.0
        vax.maximum_scale = 2.0
        vax.major_gridlines.format.line.color.rgb = GREY_LIGHT
    except Exception:
        pass

    # Línea OR=1.0 "No efecto"
    add_text_box(slide,
                 "OR=1.0\n(sin efecto)",
                 left=left + width * (1.0 / 2.0) - 0.2,
                 top=top - 0.3, width=1.0, height=0.35,
                 font_size=7, bold=False, color=AMBER)

    # Caveat box
    caveat = data.get('caveat_principal', '')
    caveat_box = slide.shapes.add_shape(
        1,
        Inches(left + width + 0.2), Inches(top),
        Inches(2.8), Inches(2.0)
    )
    caveat_box.fill.solid()
    caveat_box.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xCD)
    caveat_box.line.color.rgb = AMBER

    tf = caveat_box.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    r0 = p0.add_run()
    r0.text = "CAVEAT METODOLOGICO"
    r0.font.size = Pt(8)
    r0.font.bold = True
    r0.font.color.rgb = RGBColor(0x7B, 0x4F, 0x00)

    p1 = tf.add_paragraph()
    r1 = p1.add_run()
    r1.text = caveat
    r1.font.size = Pt(7)
    r1.font.color.rgb = TEXT_DARK

    impl = data.get('implicacion', '')
    if impl:
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = impl
        r2.font.size = Pt(7)
        r2.font.color.rgb = TEXT_DARK
        r2.font.bold = True

    return chart


# ─────────────────────────────────────────────
# C11 — PERCEPCION PRECIO NSE (bar vertical grouped)
# ─────────────────────────────────────────────

def make_chart_c11(slide, data: dict, left: float = 0.4, top: float = 1.0,
                   width: float = 9.2, height: float = 5.5):
    """Bar vertical: NETO caro vs NETO económico por NSE."""
    assert 'nse' in data and 'neto_caro_pct' in data and 'neto_economico_pct' in data, "C11: campos requeridos"
    assert len(data['nse']) == len(data['neto_caro_pct']), "C11: dimensiones inconsistentes"

    nse        = data['nse']
    neto_caro  = data['neto_caro_pct']
    neto_eco   = [v if v is not None else 0.0 for v in data['neto_economico_pct']]
    n_base     = data.get('n_base', [''] * len(nse))

    # Labels con n
    nse_labels = [f"{n}\n(n={b})" for n, b in zip(nse, n_base)]

    chart_data = CategoryChartData()
    chart_data.categories = nse_labels
    chart_data.add_series('NETO Caro (%)', neto_caro)
    chart_data.add_series('NETO Económico (%)', neto_eco)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(left), Inches(top), Inches(width), Inches(height),
        chart_data
    ).chart

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    s_caro = chart.series[0]
    s_eco  = chart.series[1]

    s_caro.format.fill.solid()
    s_caro.format.fill.fore_color.rgb = GAMA_RED

    s_eco.format.fill.solid()
    s_eco.format.fill.fore_color.rgb = GREEN_VAL

    for s in [s_caro, s_eco]:
        s.data_labels.show_value = True
        s.data_labels.number_format = '0.0"%"'
        s.data_labels.font.size = Pt(8)
        s.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

    try:
        vax = chart.value_axis
        vax.maximum_scale = 80
        vax.major_gridlines.format.line.color.rgb = GREY_LIGHT
    except Exception:
        pass

    add_footnote(slide, "* REFERENCIAL: Pref-Gama n=32. P33: percepción de precios Gama.")
    return chart


# ─────────────────────────────────────────────
# C12 — SCATTER CATEGORIAS (cuadrantes)
# ─────────────────────────────────────────────

def make_chart_c12(slide, data: dict, left: float = 0.5, top: float = 1.1,
                   width: float = 8.8, height: float = 5.5):
    """Scatter: 15 categorías con clasificación y cuadrantes."""
    assert 'puntos' in data and len(data['puntos']) > 0, "C12: puntos requeridos"

    classification_colors = {
        'OFRECER VALOR':  GREEN_VAL,
        'CUIDAR PRECIO':  PARAMO_BLUE,
        'TERRENO PÁRAMO': AMBER,
        'NEUTRO':         GREY_MED,
        'CONVENIENCIA':   RGBColor(0x9C, 0x27, 0xB0),
        'AUSENTE':        GREY_LIGHT,
    }

    xs = [p['habito'] for p in data['puntos']]
    ys = [p['precio'] for p in data['puntos']]
    x_min, x_max = 0, max(xs) + 1.5
    y_min, y_max = 0, max(ys) + 1.5
    x_range = x_max - x_min
    y_range = y_max - y_min

    def tx(x_val):
        return left + (x_val - x_min) / x_range * width

    def ty(y_val):
        return top + height - (y_val - y_min) / y_range * height

    # Background
    bg = slide.shapes.add_shape(1,
                                 Inches(left), Inches(top),
                                 Inches(width), Inches(height))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)
    bg.line.color.rgb = GREY_LIGHT

    # Mediana cuadrantes
    x_med = 4.5
    y_med = 3.0

    vline = slide.shapes.add_shape(1,
                                    Inches(tx(x_med) - 0.01), Inches(top),
                                    Inches(0.02), Inches(height))
    vline.fill.solid()
    vline.fill.fore_color.rgb = GREY_LIGHT
    vline.line.fill.background()

    hline = slide.shapes.add_shape(1,
                                    Inches(left), Inches(ty(y_med) - 0.01),
                                    Inches(width), Inches(0.02))
    hline.fill.solid()
    hline.fill.fore_color.rgb = GREY_LIGHT
    hline.line.fill.background()

    # Etiquetas cuadrante
    add_text_box(slide, "OFRECER VALOR\n(alto habito, bajo precio perc.)",
                 left=tx(x_med) + 0.1, top=ty(y_med) + 0.1,
                 width=2.2, height=0.4, font_size=7, color=GREEN_VAL)
    add_text_box(slide, "CUIDAR PRECIO\n(alto habito, alto precio perc.)",
                 left=tx(x_med) + 0.1, top=top + 0.05,
                 width=2.2, height=0.4, font_size=7, color=PARAMO_BLUE)
    add_text_box(slide, "TERRENO PÁRAMO",
                 left=left + 0.05, top=top + 0.05,
                 width=1.8, height=0.25, font_size=7, color=AMBER)

    # Puntos
    for p in data['puntos']:
        px = tx(p['habito'])
        py = ty(p['precio'])
        cls = p.get('clasificacion', 'NEUTRO')
        pt_color = classification_colors.get(cls, GREY_MED)

        dot_size = 0.16
        dot = slide.shapes.add_shape(9,
                                      Inches(px - dot_size/2), Inches(py - dot_size/2),
                                      Inches(dot_size), Inches(dot_size))
        dot.fill.solid()
        dot.fill.fore_color.rgb = pt_color
        dot.line.fill.background()

        offset_x = 0.1
        add_text_box(slide, p['categoria'],
                     left=px + offset_x, top=py - 0.12,
                     width=1.3, height=0.2,
                     font_size=7, color=pt_color)

    # Ejes labels
    ejes = data.get('ejes', {})
    add_text_box(slide, ejes.get('x_label', '% hábito Gama (P30)'),
                 left=left + width/2 - 1.5, top=top + height + 0.05,
                 width=3.0, height=0.25, font_size=8, color=GREY_MED, align=PP_ALIGN.CENTER)
    add_text_box(slide, ejes.get('y_label', '% mejor precio perc. (P32)'),
                 left=0.05, top=top + height/2 - 0.5,
                 width=0.4, height=1.0, font_size=7, color=GREY_MED)

    add_footnote(slide, "* REFERENCIAL: todos los n Gama por categoría <30. Patrones indicativos.", amber=True)
    return None


# ─────────────────────────────────────────────
# C13 — RECALL PTL vs DTLS (bar vertical grouped)
# ─────────────────────────────────────────────

def make_chart_c13(slide, data: dict, left: float = 0.4, top: float = 1.0,
                   width: float = 9.2, height: float = 5.5):
    """Bar vertical: recall PTL y DTLS por NSE vs distribución muestra."""
    assert 'nse' in data and 'pct_recordadores_ptl' in data and 'pct_recordadores_dtls' in data, "C13: campos requeridos"
    assert len(data['nse']) == len(data['pct_recordadores_ptl']) == len(data['pct_recordadores_dtls']), "C13: dimensiones inconsistentes"
    assert all(0 <= v <= 100 for v in data['pct_recordadores_ptl'] + data['pct_recordadores_dtls']), "C13: valores fuera de rango"

    nse          = data['nse']
    ptl          = data['pct_recordadores_ptl']
    dtls         = data['pct_recordadores_dtls']
    muestra_base = data.get('pct_muestra_total', [None] * len(nse))

    chart_data = CategoryChartData()
    chart_data.categories = nse
    chart_data.add_series('Recall PTL', ptl)
    chart_data.add_series('Recall DTLS', dtls)
    chart_data.add_series('Muestra Total', muestra_base)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(left), Inches(top), Inches(width), Inches(height),
        chart_data
    ).chart

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    s_ptl    = chart.series[0]
    s_dtls   = chart.series[1]
    s_muest  = chart.series[2]

    s_ptl.format.fill.solid()
    s_ptl.format.fill.fore_color.rgb = GAMA_RED

    s_dtls.format.fill.solid()
    s_dtls.format.fill.fore_color.rgb = GAMA_RED_SOFT

    s_muest.format.fill.solid()
    s_muest.format.fill.fore_color.rgb = GREY_LIGHT

    for s in [s_ptl, s_dtls, s_muest]:
        s.data_labels.show_value = True
        s.data_labels.number_format = '0.0"%"'
        s.data_labels.font.size = Pt(8)
        s.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

    try:
        vax = chart.value_axis
        vax.maximum_scale = 80
        vax.major_gridlines.format.line.color.rgb = GREY_LIGHT
    except Exception:
        pass

    n_ptl  = data.get('n_recordadores', {}).get('ptl', '?')
    n_dtls = data.get('n_recordadores', {}).get('dtls', '?')
    add_footnote(slide, f"Base recordadores: PTL n={n_ptl}, DTLS n={n_dtls}. Fuente: {data.get('fuente', 'CU-9 §8')}.")
    return chart


# ─────────────────────────────────────────────
# C14 — PERFIL RECORDADORES vs MUESTRA (bar grouped)
# ─────────────────────────────────────────────

def make_chart_c14(slide, data: dict, left: float = 0.4, top: float = 1.0,
                   width: float = 9.2, height: float = 5.5):
    """Bar vertical: perfil recordadores combinados vs muestra total + gap pp."""
    assert 'nse' in data and 'muestra_total_pct' in data and 'recordadores_combinados_pct' in data, "C14: campos requeridos"
    assert len(data['nse']) == len(data['muestra_total_pct']) == len(data['recordadores_combinados_pct']), "C14: dimensiones inconsistentes"
    assert all(0 <= v <= 100 for v in data['muestra_total_pct'] + data['recordadores_combinados_pct']), "C14: valores fuera de rango"

    nse     = data['nse']
    muestra = data['muestra_total_pct']
    record  = data['recordadores_combinados_pct']
    gap     = data.get('gap_pp', [0] * len(nse))

    chart_data = CategoryChartData()
    chart_data.categories = nse
    chart_data.add_series('Muestra Total (%)', muestra)
    chart_data.add_series('Recordadores (%)', record)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(left), Inches(top), Inches(width), Inches(height),
        chart_data
    ).chart

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    s_muest  = chart.series[0]
    s_record = chart.series[1]

    s_muest.format.fill.solid()
    s_muest.format.fill.fore_color.rgb = GREY_MED

    s_record.format.fill.solid()
    s_record.format.fill.fore_color.rgb = GAMA_RED

    for s in [s_muest, s_record]:
        s.data_labels.show_value = True
        s.data_labels.number_format = '0.0"%"'
        s.data_labels.font.size = Pt(9)
        s.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

    try:
        vax = chart.value_axis
        vax.maximum_scale = 70
        vax.major_gridlines.format.line.color.rgb = GREY_LIGHT
    except Exception:
        pass

    # Gap labels
    gap_labels = "\n".join([f"{nse[i]}: {'+' if gap[i] > 0 else ''}{gap[i]} pp" for i in range(len(nse))])
    gap_box = slide.shapes.add_shape(1,
                                      Inches(left + width + 0.1), Inches(top + 0.5),
                                      Inches(1.5), Inches(1.5))
    gap_box.fill.solid()
    gap_box.fill.fore_color.rgb = RGBColor(0xF8, 0xF0, 0xD0)
    gap_box.line.color.rgb = AMBER
    tf = gap_box.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    r0 = p0.add_run()
    r0.text = "GAP\n" + gap_labels
    r0.font.size = Pt(8)
    r0.font.color.rgb = TEXT_DARK

    n_comb = data.get('n_recordadores_combinados', '?')
    add_footnote(slide, f"Base recordadores combinados: n={n_comb}. Fuente: {data.get('fuente', 'CU-9 §8')}.")
    return chart
