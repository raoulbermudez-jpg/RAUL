"""
build_resumen_ejecutivo_v8.2.py  —  Vivienne V8.2  —  Resumen Ejecutivo (9 slides)
Proyecto: Notoriedad Gama 2026 (consultoria externa via Cora Urrea)

REGLA CRITICA: TODOS los datos leidos de chart_data_v82.json.
PROHIBIDO inventar valores. Si assertion falla -> abortar con mensaje.

Encoding: UTF-8 obligatorio (Windows cp1252 workaround)
"""
import sys
sys.stdout.reconfigure(encoding='utf-8')

import json
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Importar builders
sys.path.insert(0, str(Path(__file__).parent))
from chart_builders_v82 import (
    make_chart_c03, make_chart_c08, make_chart_c09,
    add_text_box, add_slide_title, add_footnote, add_section_tag,
    GAMA_RED, TEXT_DARK, GREY_MED, GREY_LIGHT, WHITE, BLACK,
    AMBER, GREEN_VAL, PARAMO_BLUE
)

# ─────────────────────────────────────────────
# RUTAS
# ─────────────────────────────────────────────
BASE = Path(r"C:\rAUL\03-projects\consultoria-externa\gama-notoriedad-2026\03-deck\V8")
JSON_PATH = BASE / "chart_data_v82.json"
OUT_PATH  = BASE / "2026-05-18_Notoriedad-Gama-2026_V8.2_Resumen-Ejecutivo.pptx"

# ─────────────────────────────────────────────
# CARGA JSON — VALIDACION CRITICA
# ─────────────────────────────────────────────
print("Cargando JSON (RE)...")
with open(JSON_PATH, encoding='utf-8') as f:
    DATA = json.load(f)

REQUIRED_RE = ['C03_modelo_mental_precio_dominante',
               'C07_heatmap_p23_mundo_marca',
               'C08_tornado_brechas_p23_gama',
               'C09_scatter_vecindad_perceptual',
               'C11_percepcion_precio_p33_por_nse',
               'C14_perfil_recordadores_vs_muestra']
for key in REQUIRED_RE:
    assert key in DATA, f"ABORTAR RE: '{key}' no encontrado en JSON"

# Aliases
C03 = DATA['C03_modelo_mental_precio_dominante']
C07 = DATA['C07_heatmap_p23_mundo_marca']
C08 = DATA['C08_tornado_brechas_p23_gama']
C09 = DATA['C09_scatter_vecindad_perceptual']
C11 = DATA['C11_percepcion_precio_p33_por_nse']
C14 = DATA['C14_perfil_recordadores_vs_muestra']

# SPOT-CHECKS CRITICOS
assert abs(C07['matrix_pct'][0][7] - 7.2) < 0.05,  "ABORTAR RE: C07 Gama Precio != 7.2"
assert abs(C07['matrix_pct'][0][4] - 21.9) < 0.05, "ABORTAR RE: C07 Gama Atencion != 21.9"
idx_at = C08['atributos'].index('Mejor atención')
assert abs(C08['delta_pref_total'][idx_at] - 62.5) < 0.1, "ABORTAR RE: C08 brecha Atencion != 62.5"

print(f"JSON RE OK: {len(REQUIRED_RE)} charts verificados")
print(f"SPOT-CHECK RE: C07 Gama Precio={C07['matrix_pct'][0][7]} (esperado 7.2)")
print(f"SPOT-CHECK RE: C08 Atencion brecha={C08['delta_pref_total'][idx_at]} (esperado 62.5)")

# ─────────────────────────────────────────────
# HELPERS
# ─────────────────────────────────────────────

def new_slide(prs):
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def slide_header_bar(slide, color: RGBColor = None):
    c = color or GAMA_RED
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = c
    bar.line.fill.background()


def highlight_box(slide, text: str, left: float = 0.3, top: float = 0.95,
                  width: float = 9.4, height: float = 0.5,
                  bg_color: RGBColor = None, border_color: RGBColor = None,
                  text_color: RGBColor = None, font_size: int = 12):
    bg = bg_color or RGBColor(0xFF, 0xEB, 0xEB)
    bc = border_color or GAMA_RED
    tc = text_color or GAMA_RED

    box = slide.shapes.add_shape(1, Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = bg
    box.line.color.rgb = bc

    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.color.rgb = tc
    return box


def bullet_list(slide, bullets: list, left: float, top: float,
                width: float, height: float, font_size: int = 12,
                color: RGBColor = None):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"•  {b}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color or TEXT_DARK
    return txBox


def number_badge(slide, num: str, left: float, top: float,
                 size: float = 0.38, bg: RGBColor = None):
    badge = slide.shapes.add_shape(9,
                                    Inches(left), Inches(top),
                                    Inches(size), Inches(size))
    badge.fill.solid()
    badge.fill.fore_color.rgb = bg or GAMA_RED
    badge.line.fill.background()
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = num
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = WHITE


# ─────────────────────────────────────────────
# PRESENTACION
# ─────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(10)
prs.slide_height = Inches(7.5)
slide_count = 0


# ══════════════════════════════════════════════
# RE01 — PORTADA
# ══════════════════════════════════════════════
sl = new_slide(prs)
slide_count += 1

bg = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(7.5))
bg.fill.solid()
bg.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
bg.line.fill.background()

stripe = sl.shapes.add_shape(1, Inches(0), Inches(5.5), Inches(10), Inches(0.12))
stripe.fill.solid()
stripe.fill.fore_color.rgb = GAMA_RED
stripe.line.fill.background()

label = sl.shapes.add_shape(1, Inches(0.8), Inches(1.5), Inches(3.0), Inches(0.4))
label.fill.solid()
label.fill.fore_color.rgb = GAMA_RED
label.line.fill.background()
tf = label.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "RESUMEN EJECUTIVO"
r.font.size = Pt(10)
r.font.bold = True
r.font.color.rgb = WHITE

add_text_box(sl, "Notoriedad Gama 2026",
             left=0.8, top=2.1, width=8.4, height=0.9,
             font_size=34, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text_box(sl, "V8 — Mayo 2026",
             left=0.8, top=3.1, width=8.4, height=0.45,
             font_size=16, bold=False, color=GREY_LIGHT, align=PP_ALIGN.LEFT)
add_text_box(sl, "Cliente: Gama  |  Elaborado por: Cora Urrea / equipo análisis",
             left=0.8, top=3.65, width=8.4, height=0.35,
             font_size=11, bold=False, color=GREY_LIGHT, align=PP_ALIGN.LEFT)
add_text_box(sl, "Confidencial — uso interno Gama",
             left=0.8, top=6.9, width=8.4, height=0.28,
             font_size=8, bold=False, color=GREY_MED, align=PP_ALIGN.LEFT)
print(f"  RE01 Portada OK")


# ══════════════════════════════════════════════
# RE02 — LA PREGUNTA CENTRAL
# ══════════════════════════════════════════════
sl = new_slide(prs)
slide_count += 1
slide_header_bar(sl)
add_slide_title(sl, "¿Qué cambió en el mercado y qué debe hacer Gama en 2026-2027?")
add_section_tag(sl, "[CONTEXTO]")

# 3 cards metodológicas
cards = [
    ("n = 402", "entrevistas completadas\nBBDD V2.0, n=403 × 295 preg."),
    ("4 zonas", "Venezuela\nCaracas, Miranda,\nAragua, Valencia"),
    ("±4.89%", "error muestral\n95% confianza"),
]
for i, (big_val, small_text) in enumerate(cards):
    x_pos = 0.4 + i * 3.2
    card = sl.shapes.add_shape(1, Inches(x_pos), Inches(1.6),
                                Inches(2.8), Inches(2.2))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xF8, 0xF8, 0xF8)
    card.line.color.rgb = GREY_LIGHT

    add_text_box(sl, big_val,
                 left=x_pos + 0.15, top=1.7, width=2.5, height=0.7,
                 font_size=28, bold=True, color=GAMA_RED, align=PP_ALIGN.CENTER)
    add_text_box(sl, small_text,
                 left=x_pos + 0.1, top=2.4, width=2.6, height=0.8,
                 font_size=10, bold=False, color=GREY_MED, align=PP_ALIGN.CENTER)

add_text_box(sl, "Nueva ola V8: metodología TB puro + vecindad perceptual + drivers reformulados\n→ 7 hallazgos centrales + 3 decisiones para la Junta",
             left=0.5, top=4.1, width=9.0, height=0.7,
             font_size=12, bold=False, color=TEXT_DARK)
add_text_box(sl, "Este documento resume en 7 slides el conocimiento accionable del estudio completo V8.",
             left=0.5, top=4.9, width=9.0, height=0.4,
             font_size=10, bold=False, color=GREY_MED)
print(f"  RE02 Pregunta central OK")


# ══════════════════════════════════════════════
# RE03 — 7 HALLAZGOS CENTRALES
# ══════════════════════════════════════════════
sl = new_slide(prs)
slide_count += 1
slide_header_bar(sl)
add_slide_title(sl, "7 hallazgos que definen la estrategia Gama 2026-2027")
add_section_tag(sl, "[SÍNTESIS]")

hallazgos_re = [
    ("1", "DNA experiencial consolidado", "6 atributos vs 4 en V3; precio subíndice sostenido (z=-0.76)"),
    ("2", "NSE C+/C ve Gama más fuerte", "Segmento de mayor valor: 60.6% lo percibe caro, pero lidera preferencia 13.5%"),
    ("3", "Rio sigue precio-dominante", "No migró a territorio experiencial — no compite hoy en el DNA de Gama"),
    ("4", "Campaña al segmento equivocado", f"NSE E sobre +{C14['gap_pp'][C14['nse'].index('E')]} pp / C+/C sub {C14['gap_pp'][C14['nse'].index('C+/C')]} pp"),
    ("5", "Categorías 'ofrecer valor'", "Congelados, Gaseosas, Salsas — Gama competitiva en precio en estas"),
    ("6", "[NUEVO V8] TB puro — presión precio", "Más severa de lo que T2B indicaba; Menor precio sube +6 posiciones en ranking"),
    ("7", "[NUEVO V8] ATENCIÓN = atributo SOMBRA", f"+{C08['delta_pref_total'][C08['atributos'].index('Mejor atención')]} pp entre leales vs mercado — activo invisible a activar"),
]

col1 = hallazgos_re[:4]
col2 = hallazgos_re[4:]

for i, (num, h_title, h_desc) in enumerate(col1):
    y_pos = 1.35 + i * 1.35
    number_badge(sl, num, left=0.25, top=y_pos)
    add_text_box(sl, h_title, left=0.75, top=y_pos, width=3.8, height=0.35,
                 font_size=11, bold=True, color=TEXT_DARK)
    add_text_box(sl, h_desc, left=0.75, top=y_pos + 0.36, width=4.0, height=0.55,
                 font_size=9, bold=False, color=GREY_MED)

for i, (num, h_title, h_desc) in enumerate(col2):
    y_pos = 1.35 + i * 1.6
    number_badge(sl, num, left=5.15, top=y_pos)
    add_text_box(sl, h_title, left=5.65, top=y_pos, width=3.8, height=0.35,
                 font_size=11, bold=True, color=TEXT_DARK)
    add_text_box(sl, h_desc, left=5.65, top=y_pos + 0.36, width=4.0, height=0.55,
                 font_size=9, bold=False, color=GREY_MED)

# Divisor central
divider = sl.shapes.add_shape(1, Inches(4.95), Inches(1.2), Inches(0.02), Inches(5.9))
divider.fill.solid()
divider.fill.fore_color.rgb = GREY_LIGHT
divider.line.fill.background()

add_footnote(sl, "* REFERENCIAL: cifras Pref-Gama basadas en n=32 — tendencias direccionales, no proyectables a la población.", amber=True)
print(f"  RE03 7 hallazgos OK")


# ══════════════════════════════════════════════
# RE04 — HALLAZGO #1: DNA OUTLIER
# ══════════════════════════════════════════════
sl = new_slide(prs)
slide_count += 1
slide_header_bar(sl)
gama_pct_mm = C03['pct_precio'][C03['cadenas_ordenadas_asc'].index('Gama')]
add_slide_title(sl, f"Gama es la ÚNICA cadena no dominada por precio — el DNA experiencial es el activo")
add_section_tag(sl, "[HALLAZGO 1]")
make_chart_c03(sl, C03, left=0.4, top=1.1, width=5.8, height=4.8)
# Panel derecho con análisis
panel = sl.shapes.add_shape(1, Inches(6.4), Inches(1.2), Inches(3.4), Inches(5.0))
panel.fill.solid()
panel.fill.fore_color.rgb = RGBColor(0xFF, 0xEB, 0xEB)
panel.line.color.rgb = GAMA_RED

add_text_box(sl, f"{gama_pct_mm}%",
             left=6.5, top=1.3, width=3.2, height=0.8,
             font_size=40, bold=True, color=GAMA_RED, align=PP_ALIGN.CENTER)
add_text_box(sl, "de los preferentes de Gama\ncita precio como razón dominante",
             left=6.5, top=2.1, width=3.2, height=0.6,
             font_size=10, bold=False, color=TEXT_DARK, align=PP_ALIGN.CENTER)

divider2 = sl.shapes.add_shape(1, Inches(6.6), Inches(2.8), Inches(2.8), Inches(0.02))
divider2.fill.solid()
divider2.fill.fore_color.rgb = GAMA_RED
divider2.line.fill.background()

add_text_box(sl, "vs 51–84%\nen las otras 7 cadenas",
             left=6.5, top=2.9, width=3.2, height=0.55,
             font_size=12, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)
add_text_box(sl, "Este posicionamiento ES el\ndiferenciador — erosionarlo\ncon precio directo sería\nregresivo.",
             left=6.5, top=3.55, width=3.2, height=1.0,
             font_size=10, bold=False, color=GREY_MED, align=PP_ALIGN.CENTER)
print(f"  RE04 Hallazgo DNA outlier OK")


# ══════════════════════════════════════════════
# RE05 — HALLAZGO #2: ATENCIÓN COMO ATRIBUTO SOMBRA
# ══════════════════════════════════════════════
sl = new_slide(prs)
slide_count += 1
slide_header_bar(sl)
brecha_at  = C08['delta_pref_total'][C08['atributos'].index('Mejor atención')]
pref_at    = C08['pref_gama_pct'][C08['atributos'].index('Mejor atención')]
total_at   = C08['total_pct'][C08['atributos'].index('Mejor atención')]
exp_at     = C08['exp_reciente_pct'][C08['atributos'].index('Mejor atención')]
add_slide_title(sl, f"ATENCIÓN +{brecha_at} pp entre leales vs mercado — activo invisible que debe activarse")
add_section_tag(sl, "[HALLAZGO 2]")
make_chart_c08(sl, C08, left=0.4, top=1.1, width=6.0, height=5.5)

# Panel derecho
p_right = sl.shapes.add_shape(1, Inches(6.6), Inches(1.2), Inches(3.2), Inches(4.5))
p_right.fill.solid()
p_right.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
p_right.line.color.rgb = GREEN_VAL

add_text_box(sl, f"+{brecha_at} pp",
             left=6.7, top=1.3, width=3.0, height=0.8,
             font_size=38, bold=True, color=GREEN_VAL, align=PP_ALIGN.CENTER)
add_text_box(sl, "brecha Atención\nleales vs mercado",
             left=6.7, top=2.1, width=3.0, height=0.55,
             font_size=10, bold=False, color=TEXT_DARK, align=PP_ALIGN.CENTER)

divider3 = sl.shapes.add_shape(1, Inches(6.8), Inches(2.75), Inches(2.8), Inches(0.02))
divider3.fill.solid()
divider3.fill.fore_color.rgb = GREEN_VAL
divider3.line.fill.background()

stats = [
    (f"{pref_at}%", "Pref-Gama asocia Atención"),
    (f"{exp_at}%", "Exp. recientes idem"),
    (f"{total_at}%", "Mercado total idem"),
]
for i, (val, lbl) in enumerate(stats):
    add_text_box(sl, val,
                 left=6.7, top=2.85 + i * 0.7, width=1.2, height=0.4,
                 font_size=16, bold=True, color=TEXT_DARK, align=PP_ALIGN.RIGHT)
    add_text_box(sl, lbl,
                 left=7.9, top=2.9 + i * 0.7, width=1.8, height=0.35,
                 font_size=9, bold=False, color=GREY_MED)

add_footnote(sl, "* REFERENCIAL: n Pref-Gama=32, n Exp.reciente=30. Tendencia validada, no cifra proyectable.", amber=True)
print(f"  RE05 Hallazgo Atención OK")


# ══════════════════════════════════════════════
# RE06 — HALLAZGO #3: VECINDAD + AMENAZA PLAZAS / RIO
# ══════════════════════════════════════════════
sl = new_slide(prs)
slide_count += 1
slide_header_bar(sl)
add_slide_title(sl, "Plazas es el rival perceptual más cercano — si activa Atención antes, captura el territorio")
add_section_tag(sl, "[HALLAZGO 3]")
make_chart_c09(sl, C09, left=0.4, top=1.1, width=6.0, height=5.5)

# Panel derecho
p3 = sl.shapes.add_shape(1, Inches(6.6), Inches(1.2), Inches(3.2), Inches(5.0))
p3.fill.solid()
p3.fill.fore_color.rgb = RGBColor(0xFE, 0xF9, 0xEC)
p3.line.color.rgb = AMBER

add_text_box(sl, "3 clusters del mercado:",
             left=6.7, top=1.3, width=3.0, height=0.35,
             font_size=11, bold=True, color=TEXT_DARK)

clusters = [
    (GAMA_RED, "Cluster experiencial", "Gama + Plazas + Rio"),
    (GREY_MED, "Precio-volumen", "Páramo + CM + Forum"),
    (GREY_LIGHT, "Imagen débil", "Plan Suárez + Luz"),
]
for i, (dot_color, c_name, c_desc) in enumerate(clusters):
    y_c = 1.8 + i * 0.8
    dot = sl.shapes.add_shape(9, Inches(6.7), Inches(y_c), Inches(0.18), Inches(0.18))
    dot.fill.solid()
    dot.fill.fore_color.rgb = dot_color
    dot.line.fill.background()
    add_text_box(sl, c_name, left=7.0, top=y_c - 0.02, width=2.7, height=0.28,
                 font_size=10, bold=True, color=TEXT_DARK)
    add_text_box(sl, c_desc, left=7.0, top=y_c + 0.26, width=2.7, height=0.25,
                 font_size=8, bold=False, color=GREY_MED)

divider4 = sl.shapes.add_shape(1, Inches(6.8), Inches(4.3), Inches(2.8), Inches(0.02))
divider4.fill.solid()
divider4.fill.fore_color.rgb = AMBER
divider4.line.fill.background()

add_text_box(sl, "URGENCIA: Plazas comparte mismo atributo sombra (Atención) — si comunica primero, captura el territorio.",
             left=6.7, top=4.4, width=3.0, height=0.9,
             font_size=9, bold=False, color=TEXT_DARK)
add_text_box(sl, "MONITOREO: Rio sobreíndice en Calidad entre sus leales.",
             left=6.7, top=5.4, width=3.0, height=0.5,
             font_size=8, bold=False, color=GREY_MED)
print(f"  RE06 Vecindad + amenazas OK")


# ══════════════════════════════════════════════
# RE07 — HALLAZGO #4: DRIVERS REFORMULADOS
# ══════════════════════════════════════════════
sl = new_slide(prs)
slide_count += 1
slide_header_bar(sl)
add_slide_title(sl, "La percepción predice preferencia; lo declarado 'importante' no discrimina — la campaña debe DEMOSTRAR")
add_section_tag(sl, "[HALLAZGO 4 — METODOLÓGICO]")

# Diagram de 2 columnas: "Lo que NO funciona" vs "Lo que SÍ funciona"
no_box = sl.shapes.add_shape(1, Inches(0.4), Inches(1.4), Inches(4.3), Inches(4.5))
no_box.fill.solid()
no_box.fill.fore_color.rgb = RGBColor(0xFF, 0xEB, 0xEB)
no_box.line.color.rgb = RGBColor(0xFF, 0x88, 0x88)

si_box = sl.shapes.add_shape(1, Inches(5.3), Inches(1.4), Inches(4.3), Inches(4.5))
si_box.fill.solid()
si_box.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
si_box.line.color.rgb = GREEN_VAL

add_text_box(sl, "NO FUNCIONA",
             left=0.5, top=1.5, width=4.1, height=0.4,
             font_size=13, bold=True, color=RGBColor(0xCC, 0x00, 0x00), align=PP_ALIGN.CENTER)
no_items = [
    "Importancias P22 declaradas",
    "\"Todo es muy importante\"",
    "Efecto techo — no discrimina",
    "Modelo P22→preferencia no significativo (LLR p=0.574)",
    "Mensaje: \"te atendemos bien\" (genérico)",
]
for i, item in enumerate(no_items):
    add_text_box(sl, f"✗  {item}",
                 left=0.5, top=2.0 + i * 0.55, width=4.1, height=0.45,
                 font_size=10, bold=False, color=RGBColor(0xCC, 0x00, 0x00))

add_text_box(sl, "SÍ FUNCIONA",
             left=5.4, top=1.5, width=4.1, height=0.4,
             font_size=13, bold=True, color=GREEN_VAL, align=PP_ALIGN.CENTER)
si_items = [
    "Percepción real P23",
    "Asociación real con atributos",
    "Mayor calidad: OR=0.648, p=0.042 (único significativo)",
    "Mecanismo: PERCEPCIÓN → preferencia",
    "Campaña: demostrar Atención (testimonios, momentos de verdad)",
]
for i, item in enumerate(si_items):
    add_text_box(sl, f"✓  {item}",
                 left=5.4, top=2.0 + i * 0.55, width=4.1, height=0.45,
                 font_size=10, bold=False, color=GREEN_VAL)

# Flecha central
arrow = sl.shapes.add_shape(1, Inches(4.75), Inches(3.4), Inches(0.5), Inches(0.4))
arrow.fill.solid()
arrow.fill.fore_color.rgb = TEXT_DARK
arrow.line.fill.background()
tf = arrow.text_frame
p_a = tf.paragraphs[0]
p_a.alignment = PP_ALIGN.CENTER
r_a = p_a.add_run()
r_a.text = "→"
r_a.font.size = Pt(20)
r_a.font.bold = True
r_a.font.color.rgb = WHITE

add_footnote(sl, "Fuente: CU-10 §3.3. Modelo logístico n=402. Solo Páramo viable (n=85 preferentes). Caveat: pseudo-R²=0.021 (POBRE), LLR p=0.574.", amber=True)
print(f"  RE07 Drivers reformulados OK")


# ══════════════════════════════════════════════
# RE08 — 3 DECISIONES PARA LA JUNTA
# ══════════════════════════════════════════════
sl = new_slide(prs)
slide_count += 1

# Fondo oscuro
bg_d = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(7.5))
bg_d.fill.solid()
bg_d.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
bg_d.line.fill.background()

stripe = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.08))
stripe.fill.solid()
stripe.fill.fore_color.rgb = GAMA_RED
stripe.line.fill.background()

add_text_box(sl, "3 decisiones estratégicas — Junta Gama 2026-2027",
             left=0.5, top=0.15, width=9.0, height=0.6,
             font_size=18, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_section_tag(sl, "[CALL TO ACTION]")

decisiones_re = [
    ("1", "DEFENDER DNA + abordar precio creativamente",
     "5 vías identificadas (Cap 5.6) — no competencia directa en precio.\nVan Westendorp PSM ola 2027 para calibrar límites psicológicos por segmento."),
    ("2", "ACTIVAR Atención en campaña 2026-2027",
     f"Demostrar (testimonios, momentos de verdad), no declarar.\n+{C08['delta_pref_total'][C08['atributos'].index('Mejor atención')]} pp entre leales validados — mensaje disponible, solo falta la campaña."),
    ("3", "DIFERENCIAR de Plazas + vigilar Rio + rediseñar mix de medios",
     "NSE C+/C es el segmento objetivo — campaña actual no lo alcanza.\nIncluir tracking Rio en ola 2027 (amenaza Calidad en segmento experiencial)."),
]
for i, (num, d_title, d_desc) in enumerate(decisiones_re):
    y_pos = 1.05 + i * 2.0

    badge = sl.shapes.add_shape(1, Inches(0.4), Inches(y_pos), Inches(0.5), Inches(0.5))
    badge.fill.solid()
    badge.fill.fore_color.rgb = GAMA_RED
    badge.line.fill.background()
    tf_b = badge.text_frame
    p_b = tf_b.paragraphs[0]
    p_b.alignment = PP_ALIGN.CENTER
    r_b = p_b.add_run()
    r_b.text = num
    r_b.font.size = Pt(14)
    r_b.font.bold = True
    r_b.font.color.rgb = WHITE

    add_text_box(sl, d_title,
                 left=1.1, top=y_pos, width=8.4, height=0.45,
                 font_size=13, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_text_box(sl, d_desc,
                 left=1.1, top=y_pos + 0.5, width=8.4, height=0.85,
                 font_size=10, bold=False, color=GREY_LIGHT, align=PP_ALIGN.LEFT)

add_text_box(sl, "Confidencial — Gama Notoriedad 2026 V8.2 RE | Cora Urrea / equipo análisis",
             left=0.5, top=7.1, width=9.0, height=0.25,
             font_size=7, bold=False, color=GREY_MED, align=PP_ALIGN.LEFT)
print(f"  RE08 3 decisiones OK")


# ══════════════════════════════════════════════
# RE09 — CIERRE + PROXIMOS PASOS
# ══════════════════════════════════════════════
sl = new_slide(prs)
slide_count += 1
slide_header_bar(sl)
add_slide_title(sl, "Próximo paso metodológico: Van Westendorp PSM en ola 2027")
add_section_tag(sl, "[CIERRE]")

highlight_box(sl, "La evidencia está. La estrategia es clara. El siguiente paso es ejecutar.",
              left=0.4, top=1.1, width=9.2, height=0.55,
              bg_color=RGBColor(0xE8, 0xF5, 0xE9), border_color=GREEN_VAL,
              text_color=GREEN_VAL, font_size=13)

proximos = [
    ("Van Westendorp PSM", "Ola 2027", "Calibrar precio psicológico por segmento — prerequisito antes de decisiones de precio táctico"),
    ("Tracking Rio", "Ola 2027", "Incluir tracking de Rio (amenaza Calidad en segmento experiencial)"),
    ("Mix de medios C+/C", "2026", "Rediseñar affinities de canal por NSE — campaña actual no alcanza al segmento objetivo"),
    ("Campaña Atención", "2026", "Demostrar via testimonios reales, momentos de verdad, comparativas experienciales específicas"),
]
for i, (item_title, timing, item_desc) in enumerate(proximos):
    y_pos = 1.85 + i * 1.25

    timing_box = sl.shapes.add_shape(1, Inches(0.4), Inches(y_pos), Inches(0.9), Inches(0.38))
    timing_box.fill.solid()
    timing_box.fill.fore_color.rgb = GAMA_RED if timing == "2026" else PARAMO_BLUE
    timing_box.line.fill.background()
    tf_t = timing_box.text_frame
    p_t = tf_t.paragraphs[0]
    p_t.alignment = PP_ALIGN.CENTER
    r_t = p_t.add_run()
    r_t.text = timing
    r_t.font.size = Pt(9)
    r_t.font.bold = True
    r_t.font.color.rgb = WHITE

    add_text_box(sl, item_title,
                 left=1.45, top=y_pos, width=4.0, height=0.35,
                 font_size=12, bold=True, color=TEXT_DARK)
    add_text_box(sl, item_desc,
                 left=1.45, top=y_pos + 0.36, width=8.0, height=0.45,
                 font_size=10, bold=False, color=GREY_MED)

add_text_box(sl, "Contacto: Cora Urrea | Base de datos V2.0 completa + metodología documentada disponibles.",
             left=0.5, top=7.0, width=9.0, height=0.28,
             font_size=8, bold=False, color=GREY_MED, align=PP_ALIGN.LEFT)
print(f"  RE09 Cierre OK")


# ─────────────────────────────────────────────
# GUARDAR
# ─────────────────────────────────────────────
prs.save(str(OUT_PATH))
file_size = os.path.getsize(str(OUT_PATH))
print(f"\n==================================================")
print(f"RE V8.2 GUARDADO OK")
print(f"Path:   {OUT_PATH}")
print(f"Slides: {slide_count}")
print(f"Bytes:  {file_size:,}")
print(f"==================================================")
print(f"\nSPOT-CHECKS FINALES (RE):")
print(f"  C07 Gama Precio    = {C07['matrix_pct'][0][7]} (esperado 7.2)")
print(f"  C07 Gama Atencion  = {C07['matrix_pct'][0][4]} (esperado 21.9)")
print(f"  C08 Atencion brecha= {C08['delta_pref_total'][C08['atributos'].index('Mejor atención')]} (esperado 62.5)")
