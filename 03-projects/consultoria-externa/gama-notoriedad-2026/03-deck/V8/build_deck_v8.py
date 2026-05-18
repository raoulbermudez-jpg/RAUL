import sys
sys.stdout.reconfigure(encoding='utf-8')

import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# PATHS
# ---------------------------------------------------------------------------
BASE_DIR = Path(r"C:\rAUL\03-projects\consultoria-externa\gama-notoriedad-2026\03-deck\V8")
CHARTS_DIR = BASE_DIR / "charts"
OUTPUT_PATH = BASE_DIR / "2026-05-18_Notoriedad-Gama-2026_V8.pptx"

# ---------------------------------------------------------------------------
# BRAND COLORS
# ---------------------------------------------------------------------------
ROJO_GAMA   = RGBColor(0xE3, 0x06, 0x13)
NEGRO_TEXTO = RGBColor(0x1A, 0x1A, 0x1A)
GRIS_MED    = RGBColor(0x6B, 0x6B, 0x6B)
GRIS_CLARO  = RGBColor(0xE5, 0xE5, 0xE5)
AMBAR       = RGBColor(0xF2, 0xA9, 0x00)
VERDE       = RGBColor(0x2D, 0x8F, 0x47)
BLANCO      = RGBColor(0xFF, 0xFF, 0xFF)

# ---------------------------------------------------------------------------
# SLIDE DIMENSIONS (16:9 widescreen)
# ---------------------------------------------------------------------------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ---------------------------------------------------------------------------
# LAYOUT INDEX — use Blank (index 6 in most default masters, fallback to 5)
# ---------------------------------------------------------------------------

def get_blank_layout(prs):
    """Return blank layout from available layouts."""
    for i, layout in enumerate(prs.slide_layouts):
        if layout.name.lower() in ("blank", "en blanco"):
            return layout
    # Fallback: last layout is usually blank
    return prs.slide_layouts[-1]


# ---------------------------------------------------------------------------
# HELPER: add a text box
# ---------------------------------------------------------------------------
def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=None, align=PP_ALIGN.LEFT,
                 font_name="Calibri", wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return txBox


def add_text_box_multiline(slide, lines, left, top, width, height,
                            font_size=14, bold=False, color=None,
                            align=PP_ALIGN.LEFT, font_name="Calibri",
                            bullet_char="•", line_spacing_pt=None):
    """Add a text box with multiple paragraphs (bullet lines)."""
    from pptx.util import Pt as _Pt
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for idx, line in enumerate(lines):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = f"{bullet_char} {line}" if bullet_char else line
        run.font.name = font_name
        run.font.size = _Pt(font_size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color
    return txBox


# ---------------------------------------------------------------------------
# HELPER: fill slide background
# ---------------------------------------------------------------------------
def fill_background(slide, color):
    from pptx.util import Pt as _Pt
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


# ---------------------------------------------------------------------------
# HELPER: add a colored rectangle (callout strip, header bar, etc.)
# ---------------------------------------------------------------------------
def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


# ---------------------------------------------------------------------------
# HELPER: add chart image
# ---------------------------------------------------------------------------
def add_chart_image(slide, chart_path, left, top, width, height):
    if not chart_path.exists():
        print(f"  WARNING: chart not found -> {chart_path}")
        return None
    pic = slide.shapes.add_picture(str(chart_path), left, top, width, height)
    return pic


# ---------------------------------------------------------------------------
# HELPER: add footnote at bottom
# ---------------------------------------------------------------------------
def add_footnote(slide, text):
    add_text_box(
        slide, text,
        left=Inches(0.5), top=Inches(7.0),
        width=Inches(12.3), height=Inches(0.35),
        font_size=9, color=GRIS_MED, italic=True
    )


# ---------------------------------------------------------------------------
# STANDARD SLIDE BUILDERS
# ---------------------------------------------------------------------------

def build_title_bar(slide, title_text, tag_text=None):
    """Red header bar with title. Optional tag badge top-right."""
    # Red bar
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.15), ROJO_GAMA)
    # Title text
    add_text_box(
        slide, title_text,
        left=Inches(0.35), top=Inches(0.12),
        width=Inches(11.5), height=Inches(0.9),
        font_size=22, bold=True, color=BLANCO,
        font_name="Calibri"
    )
    # Tag badge
    if tag_text:
        add_rect(slide, Inches(11.7), Inches(0.18), Inches(1.45), Inches(0.38), NEGRO_TEXTO)
        add_text_box(
            slide, tag_text,
            left=Inches(11.72), top=Inches(0.21),
            width=Inches(1.4), height=Inches(0.32),
            font_size=9, bold=True, color=BLANCO,
            align=PP_ALIGN.CENTER
        )


def slide_portada(prs):
    """S01 — Cover slide."""
    layout = get_blank_layout(prs)
    slide = prs.slides.add_slide(layout)
    fill_background(slide, BLANCO)

    # Top red bar full width
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(2.0), ROJO_GAMA)

    # Main title
    add_text_box(
        slide, "NOTORIEDAD GAMA 2026",
        left=Inches(0.5), top=Inches(0.25),
        width=Inches(12.3), height=Inches(0.85),
        font_size=36, bold=True, color=BLANCO,
        align=PP_ALIGN.CENTER, font_name="Calibri"
    )
    # Subtitle
    add_text_box(
        slide, "Estudio de Brand Health  |  V8 — Mayo 2026",
        left=Inches(0.5), top=Inches(1.1),
        width=Inches(12.3), height=Inches(0.6),
        font_size=20, bold=False, color=BLANCO,
        align=PP_ALIGN.CENTER
    )

    # Body block
    add_text_box(
        slide, "Cliente: Gama",
        left=Inches(1.5), top=Inches(2.5),
        width=Inches(10.3), height=Inches(0.5),
        font_size=16, bold=True, color=NEGRO_TEXTO,
        align=PP_ALIGN.CENTER
    )
    add_text_box(
        slide, "Elaborado por: Cora Urrea / equipo análisis",
        left=Inches(1.5), top=Inches(3.0),
        width=Inches(10.3), height=Inches(0.5),
        font_size=14, color=GRIS_MED, align=PP_ALIGN.CENTER
    )

    # Bottom bar
    add_rect(slide, Inches(0), Inches(6.9), SLIDE_W, Inches(0.6), GRIS_CLARO)
    add_text_box(
        slide, "Confidencial — uso interno Gama",
        left=Inches(0.5), top=Inches(6.95),
        width=Inches(12.3), height=Inches(0.4),
        font_size=10, color=GRIS_MED, align=PP_ALIGN.CENTER
    )


def slide_toc(prs):
    """S02 — Table of contents."""
    layout = get_blank_layout(prs)
    slide = prs.slides.add_slide(layout)
    fill_background(slide, BLANCO)
    build_title_bar(slide, "Ruta narrativa", "[AGENDA]")

    sections = [
        ("Sección 1", "Embudo de marca + posicionamiento espontáneo (Cap. 1)"),
        ("Sección 2", "Top Box puro — prioridades funcionales reales [NUEVO V8]"),
        ("Sección 3", "DNA de marca + Mundo de marca Gama (Cap. 2)"),
        ("Sección 4", "Vecindad perceptual + amenazas competitivas [NUEVO V8]"),
        ("Sección 5", "Drivers reformulados — percepción como predictor (Cap. 4.4)"),
        ("Sección 6", "Precio + estrategia por categorías (Cap. 3 + 5)"),
        ("Sección 7", "Comunicación PTL/DTLS — alcance y perfil (Cap. 4)"),
        ("Sección 8", "Síntesis + 3 decisiones para Junta"),
    ]

    top_start = Inches(1.35)
    row_h = Inches(0.62)
    for i, (sec, desc) in enumerate(sections):
        top = top_start + i * row_h
        # section label
        add_text_box(
            slide, sec,
            left=Inches(0.4), top=top,
            width=Inches(1.6), height=Inches(0.5),
            font_size=11, bold=True, color=ROJO_GAMA
        )
        # description
        add_text_box(
            slide, desc,
            left=Inches(2.1), top=top,
            width=Inches(11.0), height=Inches(0.5),
            font_size=12, color=NEGRO_TEXTO
        )
        # separator line
        add_rect(slide, Inches(0.4), top + Inches(0.52), Inches(12.8), Emu(8000), GRIS_CLARO)


def slide_ficha(prs):
    """S03 — Technical specification."""
    layout = get_blank_layout(prs)
    slide = prs.slides.add_slide(layout)
    fill_background(slide, BLANCO)
    build_title_bar(slide, "Metodología: 402 entrevistas, 4 zonas, ±4.89%", "[METODOLOGÍA]")

    bullets = [
        "n = 402 entrevistas completadas (BBDD V2.0, n=403 × 295 preguntas)",
        "Geografía: 4 zonas Venezuela — Caracas, Miranda, Aragua, Valencia",
        "Fechas de campo: Ola 2026  |  Margen de error: ±4.89% (95% confianza)",
        "Segmentación: NSE A/B, C+, C, D, E  |  Cuota edad × género",
    ]
    add_text_box_multiline(
        slide, bullets,
        left=Inches(0.6), top=Inches(1.4),
        width=Inches(12.1), height=Inches(4.5),
        font_size=16, color=NEGRO_TEXTO
    )


# ---------------------------------------------------------------------------
# GENERIC DATOS SLIDE (chart full)
# ---------------------------------------------------------------------------
def slide_datos(prs, title, chart_file, tag="[DATOS]", footnote=None, speaker_note=None):
    layout = get_blank_layout(prs)
    slide = prs.slides.add_slide(layout)
    fill_background(slide, BLANCO)
    build_title_bar(slide, title, tag)

    chart_path = CHARTS_DIR / chart_file
    # Chart occupies most of slide below header
    add_chart_image(
        slide, chart_path,
        left=Inches(0.3), top=Inches(1.25),
        width=Inches(12.7), height=Inches(5.55)
    )
    if footnote:
        add_footnote(slide, footnote)
    if speaker_note:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = speaker_note


# ---------------------------------------------------------------------------
# GENERIC DATOS + CHART + BULLETS SLIDE
# ---------------------------------------------------------------------------
def slide_datos_bullets(prs, title, chart_file, bullets, tag="[DATOS + ANÁLISIS]",
                         footnote=None, chart_pct=0.52, speaker_note=None):
    """Left chart + right bullets layout."""
    layout = get_blank_layout(prs)
    slide = prs.slides.add_slide(layout)
    fill_background(slide, BLANCO)
    build_title_bar(slide, title, tag)

    chart_w = Inches(13.333 * chart_pct)
    chart_path = CHARTS_DIR / chart_file
    add_chart_image(
        slide, chart_path,
        left=Inches(0.2), top=Inches(1.25),
        width=chart_w, height=Inches(5.5)
    )

    bullets_left = Inches(0.2) + chart_w + Inches(0.2)
    bullets_width = SLIDE_W - bullets_left - Inches(0.2)
    add_text_box_multiline(
        slide, bullets,
        left=bullets_left, top=Inches(1.5),
        width=bullets_width, height=Inches(5.0),
        font_size=13, color=NEGRO_TEXTO
    )
    if footnote:
        add_footnote(slide, footnote)
    if speaker_note:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = speaker_note


# ---------------------------------------------------------------------------
# GENERIC ANALISIS SLIDE (bullets only)
# ---------------------------------------------------------------------------
def slide_analisis(prs, title, bullets, tag="[ANÁLISIS]",
                    footnote=None, accent_color=None, speaker_note=None):
    layout = get_blank_layout(prs)
    slide = prs.slides.add_slide(layout)
    fill_background(slide, BLANCO)
    build_title_bar(slide, title, tag)

    color = accent_color if accent_color else NEGRO_TEXTO
    add_text_box_multiline(
        slide, bullets,
        left=Inches(0.6), top=Inches(1.45),
        width=Inches(12.1), height=Inches(5.2),
        font_size=15, color=color,
        line_spacing_pt=6
    )
    if footnote:
        add_footnote(slide, footnote)
    if speaker_note:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = speaker_note


# ---------------------------------------------------------------------------
# SYNTHESIS SLIDE (S31)
# ---------------------------------------------------------------------------
def slide_sintesis(prs):
    layout = get_blank_layout(prs)
    slide = prs.slides.add_slide(layout)
    fill_background(slide, BLANCO)
    build_title_bar(
        slide,
        "5 hallazgos que definen la estrategia Gama 2026-2027",
        "[SÍNTESIS]"
    )

    hallazgos = [
        "DNA experiencial consolidado: Gama es la ÚNICA cadena no dominada por precio (40.6%)",
        "Atributo sombra = Atención: +62.5 pp entre leales vs mercado amplio — el mercado no lo sabe aún",
        "Vecindad Plazas: rival perceptual más próximo en el mismo atributo sombra — urgencia de diferenciación",
        "Presión precio más severa en TB puro: exigencia del mercado es más rígida de lo que T2B mostraba",
        "Campaña llega al segmento equivocado: NSE E sobre-representado, C+/C sub-representado",
    ]

    top_start = Inches(1.4)
    for i, h in enumerate(hallazgos):
        top = top_start + i * Inches(1.0)
        # Number badge
        add_rect(slide, Inches(0.4), top, Inches(0.38), Inches(0.38), ROJO_GAMA)
        add_text_box(
            slide, str(i + 1),
            left=Inches(0.4), top=top - Inches(0.02),
            width=Inches(0.38), height=Inches(0.38),
            font_size=14, bold=True, color=BLANCO,
            align=PP_ALIGN.CENTER
        )
        add_text_box(
            slide, h,
            left=Inches(0.9), top=top,
            width=Inches(12.0), height=Inches(0.55),
            font_size=13, color=NEGRO_TEXTO
        )


# ---------------------------------------------------------------------------
# CALL TO ACTION SLIDE (S32)
# ---------------------------------------------------------------------------
def slide_cta(prs):
    layout = get_blank_layout(prs)
    slide = prs.slides.add_slide(layout)
    fill_background(slide, BLANCO)
    build_title_bar(
        slide,
        "3 decisiones estratégicas para Junta Gama — 2026-2027",
        "[CALL TO ACTION]"
    )

    decisions = [
        (
            "1. Defender DNA + abordar precio creativamente",
            "5 vías identificadas: promociones diferenciadas, coleccionables, mensaje del anzuelo, lealtad visible, Van Westendorp 2027"
        ),
        (
            "2. Activar Atención en campaña 2026-2027",
            "Demostrar (testimonios, momentos de verdad, comparativas), NO declarar. Anti-mensaje: 'ofrecemos buen servicio'"
        ),
        (
            "3. Diferenciar de Plazas + vigilar Rio en Calidad",
            "Incluir tracking competitivo de Rio + prerrequisito Van Westendorp en diseño ola 2027"
        ),
    ]

    for i, (header, detail) in enumerate(decisions):
        top_base = Inches(1.5) + i * Inches(1.8)
        # Colored strip
        add_rect(slide, Inches(0.3), top_base, Inches(12.7), Inches(0.55), GRIS_CLARO)
        add_text_box(
            slide, header,
            left=Inches(0.5), top=top_base + Inches(0.05),
            width=Inches(12.3), height=Inches(0.45),
            font_size=14, bold=True, color=ROJO_GAMA
        )
        add_text_box(
            slide, detail,
            left=Inches(0.7), top=top_base + Inches(0.6),
            width=Inches(12.0), height=Inches(0.9),
            font_size=13, color=NEGRO_TEXTO
        )


# ---------------------------------------------------------------------------
# BUILD ALL SLIDES
# ---------------------------------------------------------------------------
def build_deck():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    print("Building deck V8...")

    # ---- SECCIÓN 0 ----
    slide_portada(prs)          # S01
    slide_toc(prs)              # S02
    slide_ficha(prs)            # S03

    # ---- SECCIÓN 1 ----
    # S04 [DATOS] Embudo C01
    slide_datos(
        prs,
        title="Embudo de notoriedad espontánea a preferencia — 8 cadenas",
        chart_file="C01_embudo_marca.png",
        tag="[DATOS]",
        speaker_note="[DATOS] Métricas: Notoriedad espontánea, Recordación top of mind, Visita en 30d, Consideración, Preferencia."
    )

    # S05 [ANÁLISIS] Embudo
    slide_analisis(
        prs,
        title="Gama no tiene problema de visibilidad — tiene un problema de conversión",
        bullets=[
            "Notoriedad espontánea alta: Gama posición 1-2 en el mercado",
            "Caída significativa entre consideración y preferencia: el mercado 'conoce pero no elige'",
            "Palanca de crecimiento = mejorar percepción de valor, no construir más awareness",
            "Implicación: inversión en comunicación debe dirigirse a fondo de embudo",
        ],
        tag="[ANÁLISIS]"
    )

    # S06 [DATOS] Preferencia P21 C02
    slide_datos(
        prs,
        title="Preferencia espontánea P21 — ranking total y segmento C+/C",
        chart_file="C02_p21_preferencia.png",
        tag="[DATOS]",
        speaker_note="[DATOS] Pregunta P21 espontánea. NSE C+/C = segmento objetivo premium."
    )

    # S07 [ANÁLISIS] Gama #1 C+/C
    slide_analisis(
        prs,
        title="En C+/C, Gama lidera con 30.8% — Plan Suárez es el #2 con 15.4%",
        bullets=[
            "Gama consolida liderazgo en el segmento de mayor valor por ticket (30.8% C+/C)",
            "Plan Suárez alcanza 15.4% C+/C: rivalidad subestimada en segmento premium",
            "La brecha (≈15 pp) es defendible pero no definitiva — requiere activación sostenida",
            "Hallazgo Cora: diferencia estadística — validar n por celda si n<30",
        ],
        tag="[ANÁLISIS]",
        footnote="REF: n por celda C+/C — validar significancia si n<30"
    )

    # S08 [DATOS+ANÁLISIS] Modelo mental C03
    slide_datos_bullets(
        prs,
        title="Gama es la ÚNICA cadena con modelo mental no dominado por precio (40.6%)",
        chart_file="C03_modelo_mental_precio.png",
        bullets=[
            "Las otras 7 cadenas: 51-84% asociadas a precio como atributo dominante",
            "Gama: 40.6% — outlier estructural, no accidental",
            "Posicionamiento experiencial de Gama ES el diferenciador real",
            "Implicación: no competir en precio directo; proteger este outlier",
        ],
        tag="[DATOS + ANÁLISIS]",
        chart_pct=0.56
    )

    # ---- SECCIÓN 2 ----
    # S09 [DATOS] TB puro reranking C04
    slide_datos(
        prs,
        title="Top Box puro reordena las prioridades del mercado vs T2B",
        chart_file="C04_p22_tb_puro_reranking.png",
        tag="[DATOS]",
        speaker_note="[DATOS] P22: importancia atributos de selección de supermercado. T2B = top 2 boxes; TB = top box únicamente."
    )

    # S10 [ANÁLISIS] TB puro
    slide_analisis(
        prs,
        title="El T2B suavizaba la presión: en TB puro, precio es prioridad más dura",
        bullets=[
            "T2B: precio aparecía importante pero no dominante en el ranking",
            "TB puro: precio escala posición — el mercado tiene exigencia funcional más rígida",
            "La presión de precio que enfrenta Gama es más severa de lo que V3-V7 mostraban",
            "No basta con 'ser experiencial' — hay que abordar precio creativamente con las 5 vías",
        ],
        tag="[ANÁLISIS]",
        accent_color=AMBAR
    )

    # S11 [DATOS] TB puro por marca C05
    slide_datos(
        prs,
        title="TB puro por marca preferida: heatmap de exigencias por cluster de leales",
        chart_file="C05_p22_tb_puro_por_marca.png",
        tag="[DATOS]",
        footnote="REFERENCIAL: n Pref-Gama = 32. Leer con cautela; tendencias, no cifras definitivas.",
        speaker_note="[DATOS] Filas = marcas. Columnas = atributos. Color = % TB puro. Ref: n=32 para Pref-Gama."
    )

    # S12 [ANÁLISIS] Leales Gama bifurcación
    slide_analisis(
        prs,
        title="Los leales de Gama priorizan Rapidez y Atención — Precio cae al noveno lugar",
        bullets=[
            "Pref-Gama (n=32 REF): Rapidez #1 con 84.4%, Atención #2 con 71.9%",
            "Precio ocupa el puesto #9 con solo 43.8% entre sus propios leales",
            "La bifurcación mercado amplio (precio fuerte) vs leales (experiencia) es MÁS extrema en TB puro",
            "El reto: convertir no-leales sin abandonar la propuesta que retiene a los leales",
        ],
        tag="[ANÁLISIS]",
        footnote="REFERENCIAL: n=32. Datos direccionales, no proyectables al universo."
    )

    # ---- SECCIÓN 3 ----
    # S13 [DATOS] DNA z-scores C06
    slide_datos(
        prs,
        title="DNA de marca Gama: z-scores de asociación vs promedio mercado",
        chart_file="C06_dna_zscores.png",
        tag="[DATOS]",
        speaker_note="[DATOS] Z-score >0: sobreíndice vs promedio 8 cadenas. Z-score <0: subíndice."
    )

    # S14 [ANÁLISIS] DNA consolidado
    slide_analisis(
        prs,
        title="Gama sobreíndice en 6 atributos experienciales; precio subíndice sostenido (-0.76 z)",
        bullets=[
            "V8 confirma y extiende V3: 6 atributos experienciales por encima del promedio (vs 4 en V3)",
            "Precio: z = -0.76 — estructura perceptual establecida, no variación aleatoria",
            "DNA experiencial es activo estratégico, no debilidad a corregir",
            "Riesgo: intentar competir en precio directamente puede erosionar el DNA diferenciador",
        ],
        tag="[ANÁLISIS]",
        accent_color=VERDE
    )

    # S15 [DATOS] Mundo de marca C07
    slide_datos(
        prs,
        title="Mapa de imagen P23: 8 cadenas × 10 atributos — posición completa del mercado",
        chart_file="C07_p23_mundo_marca_total.png",
        tag="[DATOS]",
        speaker_note="[DATOS] P23 espontánea de atributos por cadena. Base: total muestra n=402."
    )

    # S16 [ANÁLISIS] Cluster experiencial
    slide_analisis(
        prs,
        title="Gama lidera el cluster experiencial; Limpieza y Seguridad son 'tabla de higiene' del sector",
        bullets=[
            "Gama, Rio y Plazas comparten cluster experiencial — zona de alta competencia perceptual",
            "Otros competidores se anclan en precio/volumen: territorio distinto, menor overlap directo",
            "Limpieza y Seguridad convergen en todas las cadenas: no diferencian, son condición de entrada",
            "Para diferenciarse, Gama debe activar atributos exclusivos dentro del cluster experiencial",
        ],
        tag="[ANÁLISIS]"
    )

    # S17 [DATOS] Brechas P23 tornado C08
    slide_datos(
        prs,
        title="Brechas de imagen Gama — compradores activos vs mercado amplio (tornado P23)",
        chart_file="C08_p23_brechas_pref_vs_total.png",
        tag="[DATOS]",
        speaker_note="[DATOS] Brecha = % Pref-Gama que asocia atributo MENOS % Total que asocia atributo. Positivo = leales ven más."
    )

    # S18 [ANÁLISIS] Atributo sombra
    slide_analisis(
        prs,
        title="ATRIBUTO SOMBRA: Atención está +62.5 pp entre leales — el mercado no lo sabe aún",
        bullets=[
            "Quienes compran en Gama asocian Atención en 62.5 pp más que el mercado amplio",
            "El mercado amplio no percibe la atención de Gama — imagen no llegó a no-clientes",
            "Tarea de comunicación: trasladar imagen validada por leales al mercado amplio",
            "No es crear un atributo nuevo — es amplificar uno que ya existe y está validado internamente",
        ],
        tag="[ANÁLISIS]",
        accent_color=ROJO_GAMA
    )

    # ---- SECCIÓN 4 ----
    # S19 [DATOS] Vecindad perceptual C09
    slide_datos(
        prs,
        title="Vecindad perceptual: mapa de posicionamiento relativo — 8 cadenas",
        chart_file="C09_vecindad_perceptual.png",
        tag="[DATOS]",
        speaker_note="[DATOS] Análisis de correspondencias o MDS sobre P23. Proximidad = similitud perceptual."
    )

    # S20 [ANÁLISIS] Plazas rival más cercano
    slide_analisis(
        prs,
        title="Plazas es el rival perceptual más próximo a Gama — mismo cluster, mismo atributo sombra",
        bullets=[
            "Plazas ocupa el vecindario más cercano a Gama en el mapa perceptual",
            "Comparten cluster experiencial y el atributo sombra es Atención en ambas marcas",
            "Si Plazas activa comunicación de Atención antes que Gama, puede capturar el territorio",
            "Prioridad urgente: diferenciar de Plazas en Atención — calidad, velocidad o trato personalizado",
        ],
        tag="[ANÁLISIS]",
        accent_color=AMBAR
    )

    # S21 [ANÁLISIS] Rio amenaza latente
    slide_analisis(
        prs,
        title="Rio sobreíndice en Calidad (+69 pp entre leales) — si comunica, puede capturar consumidores experienciales",
        bullets=[
            "Rio no compite en el territorio actual de Gama (páramo opuesto en el mapa perceptual)",
            "Atributo sombra de Rio = Calidad: +69 pp entre sus compradores vs mercado amplio",
            "Si Rio activa campaña de calidad, puede atraer al consumidor experiencial que Gama no fidelizó",
            "Monitoreo recomendado: incluir tracking de Rio en diseño ola 2027",
        ],
        tag="[ANÁLISIS]",
        footnote="REFERENCIAL: n Pref-Rio — verificar tamaño de muestra antes de reportar cifra exacta."
    )

    # ---- SECCIÓN 5 ----
    # S22 [DATOS] Regresión logística C10
    slide_datos(
        prs,
        title="Modelo de drivers — regresión logística sobre percepción P23 → preferencia",
        chart_file="C10_drivers_logit_paramo.png",
        tag="[DATOS]",
        footnote="Modelo exploratorio. Validación recomendada con ola 2027 n≥600.",
        speaker_note="[DATOS] Variable dependiente: preferir Gama vs no. Predictores: % asociación por atributo P23. Caveats: n=402, variables binarias, VIF controlado."
    )

    # S23 [ANÁLISIS METODOLÓGICO] Importancias declaradas
    slide_analisis(
        prs,
        title="Lo que la gente DECLARA importante no predice preferencia — la percepción sí",
        bullets=[
            "Importancias P22 (declaradas) sufren 'techo': casi todo es muy importante, no discrimina",
            "Modelo logístico sobre P23 (percepción real) sí predice preferencia Gama con significancia",
            "El mecanismo correcto: demostrar que Gama TIENE el atributo, no que el atributo es importante",
            "Anti-mensaje: 'ofrecemos buen servicio' — genérico, no memorable, no diferencia de Plazas",
        ],
        tag="[ANÁLISIS METODOLÓGICO]"
    )

    # S24 [ANÁLISIS ESTRATÉGICO] Campaña correcta
    slide_analisis(
        prs,
        title="La campaña correcta demuestra Atención superior — no la declara",
        bullets=[
            "Mensaje equivocado: 'en Gama te atendemos bien' — cualquier cadena puede decir lo mismo",
            "Mensaje correcto: testimonios reales, comparativas experienciales, momentos de verdad específicos",
            "Activar EXPERIENCIA + RECONOCIMIENTO como mecanismo creativo central de campaña",
            "Alinear con atributo sombra: trasladar imagen de leales al mercado amplio a través de evidencia",
        ],
        tag="[ANÁLISIS ESTRATÉGICO]",
        accent_color=VERDE
    )

    # ---- SECCIÓN 6 ----
    # S25 [DATOS] Percepción precio NSE C11
    slide_datos(
        prs,
        title="Percepción de precio Gama por nivel socioeconómico (P33)",
        chart_file="C11_p33_percepcion_precio_nse.png",
        tag="[DATOS]",
        speaker_note="[DATOS] P33: 'Los precios de Gama son...' Escala de percepción por NSE. Base: compradores Gama por NSE."
    )

    # S26 [DATOS+ANÁLISIS] Categorías C12
    slide_datos_bullets(
        prs,
        title="Categorías: Congelados y Gaseosas = ofrecer valor; Galletas = cuidar precio; Proteínas = no competir",
        chart_file="C12_p30_p32_categorias_estrategia.png",
        bullets=[
            "Ofrecer valor: Congelados, Gaseosas, Salsas — Gama percibida competitiva en precio",
            "Cuidar precio: Galletas — sensibilidad alta, ajuste táctico necesario",
            "No competir: Proteínas — costo percibido alto, riesgo de esfuerzo sin retorno",
        ],
        tag="[DATOS + ANÁLISIS]",
        chart_pct=0.56
    )

    # S27 [ANÁLISIS] 5 vías creativas precio
    slide_analisis(
        prs,
        title="5 vías para abordar precio sin abandonar el posicionamiento experiencial",
        bullets=[
            "(1) Promociones diferenciadas: beneficio simbólico + precio (no precio como único gancho)",
            "(2) Coleccionables / lealtad visible: monetizar la relación, no el descuento puntual",
            "(3) Mensaje del anzuelo: una categoría competitiva en precio ancla la percepción general",
            "(4) Prerrequisito metodológico: Van Westendorp ola 2027 para calibrar precio psicológico por NSE",
        ],
        tag="[ANÁLISIS]"
    )

    # ---- SECCIÓN 7 ----
    # S28 [DATOS] Recall PTL DTLS C13
    slide_datos(
        prs,
        title="Recall de campaña PTL vs DTLS — desagregado por NSE",
        chart_file="C13_recall_ptl_dtls.png",
        tag="[DATOS]",
        speaker_note="[DATOS] PTL = 'Precios de tu lado'. DTLS = activaciones DTLS. Recall espontáneo + asistido por NSE."
    )

    # S29 [DATOS] Perfil recordadores C14
    slide_datos(
        prs,
        title="¿Quién recuerda la campaña? Perfil recordadores vs muestra total",
        chart_file="C14_perfil_recordadores_vs_muestra.png",
        tag="[DATOS]",
        speaker_note="[DATOS] Índice de sobre/sub-representación por NSE, género, edad. Base: recordadores vs total muestra."
    )

    # S30 [ANÁLISIS] Segmento equivocado
    slide_analisis(
        prs,
        title="La campaña actual sobre-alcanza NSE E y sub-representa C+/C — reorientar",
        bullets=[
            "NSE E sobre-representado entre recordadores: segmento con menor ticket y menor conversión a valor",
            "NSE C+/C sub-representado: segmento objetivo premium no está siendo impactado por la campaña",
            "PTL 'Precios de tu lado' refuerza atributo precio en NSE E — contrario al DNA experiencial de Gama",
            "Acción: rediseñar mix de medios + mensaje para C+/C; validar affinities de canal por NSE",
        ],
        tag="[ANÁLISIS]",
        accent_color=AMBAR
    )

    # ---- SECCIÓN 8 ----
    slide_sintesis(prs)          # S31
    slide_cta(prs)               # S32

    # ---- SAVE ----
    prs.save(str(OUTPUT_PATH))
    size_bytes = OUTPUT_PATH.stat().st_size
    print(f"\nDeck guardado: {OUTPUT_PATH}")
    print(f"Slides totales: {len(prs.slides)}")
    print(f"Tamaño: {size_bytes:,} bytes ({size_bytes/1024:.1f} KB)")
    return len(prs.slides), size_bytes


if __name__ == "__main__":
    n_slides, size = build_deck()
    print(f"\n✅ build_deck_v8.py completado — {n_slides} slides, {size:,} bytes")
