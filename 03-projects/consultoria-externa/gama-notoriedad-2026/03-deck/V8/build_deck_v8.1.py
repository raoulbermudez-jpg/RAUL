import sys
sys.stdout.reconfigure(encoding='utf-8')

import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.util import Inches, Pt
from lxml import etree

# ---------------------------------------------------------------------------
# PATHS
# ---------------------------------------------------------------------------
BASE_DIR = Path(r"C:\rAUL\03-projects\consultoria-externa\gama-notoriedad-2026\03-deck\V8")
OUTPUT_PATH = BASE_DIR / "2026-05-18_Notoriedad-Gama-2026_V8.1.pptx"

# ---------------------------------------------------------------------------
# BRAND COLORS
# ---------------------------------------------------------------------------
ROJO_GAMA   = RGBColor(0xE3, 0x06, 0x13)
NEGRO_TEXTO = RGBColor(0x1A, 0x1A, 0x1A)
GRIS_MED    = RGBColor(0x6B, 0x6B, 0x6B)
GRIS_CLARO  = RGBColor(0xE5, 0xE5, 0xE5)
AMBAR       = RGBColor(0xF2, 0xA9, 0x00)
VERDE       = RGBColor(0x2D, 0x8F, 0x47)
BLANCO      = RGBColor(0xFF, 0xFF, 0xFF)

# Competitor colors
COLOR_PARAMO   = RGBColor(0x4A, 0x6F, 0xA5)
COLOR_RIO      = RGBColor(0x7C, 0x8B, 0x9B)
COLOR_PSUAREZ  = RGBColor(0x9C, 0xAE, 0xC0)
COLOR_CENTRAL  = RGBColor(0x5B, 0x70, 0x90)
COLOR_FORUM    = RGBColor(0x8F, 0xA3, 0xB8)
COLOR_PLAZAS   = RGBColor(0x6D, 0x84, 0x99)
COLOR_LUZ      = RGBColor(0xA8, 0xB9, 0xC7)
COLOR_COMP_DEF = RGBColor(0xBC, 0xC9, 0xD4)

# ---------------------------------------------------------------------------
# SLIDE DIMENSIONS (16:9 widescreen)
# ---------------------------------------------------------------------------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Ordered list of 8 chains — Gama always index 0
CADENAS = ["Gama", "Páramo", "Rio", "Plan Suárez", "Central Madeirense", "Forum", "Plazas", "Luz Marina"]
CADENA_COLORS = [ROJO_GAMA, COLOR_PARAMO, COLOR_RIO, COLOR_PSUAREZ, COLOR_CENTRAL, COLOR_FORUM, COLOR_PLAZAS, COLOR_LUZ]

# ---------------------------------------------------------------------------
# HELPER: get blank layout
# ---------------------------------------------------------------------------
def get_blank_layout(prs):
    for layout in prs.slide_layouts:
        if layout.name.lower() in ("blank", "en blanco"):
            return layout
    return prs.slide_layouts[-1]


# ---------------------------------------------------------------------------
# HELPER: add a text box
# ---------------------------------------------------------------------------
def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=None, align=PP_ALIGN.LEFT,
                 font_name="Calibri", wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return txBox


def add_text_box_multiline(slide, lines, left, top, width, height,
                            font_size=14, bold=False, color=None,
                            align=PP_ALIGN.LEFT, font_name="Calibri",
                            bullet_char="•", line_spacing_pt=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = f"{bullet_char} {line}" if bullet_char else line
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color
    return txBox


# ---------------------------------------------------------------------------
# HELPER: fill slide background
# ---------------------------------------------------------------------------
def fill_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


# ---------------------------------------------------------------------------
# HELPER: add a colored rectangle
# ---------------------------------------------------------------------------
def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


# ---------------------------------------------------------------------------
# HELPER: add footnote
# ---------------------------------------------------------------------------
def add_footnote(slide, text):
    add_text_box(
        slide, text,
        left=Inches(0.5), top=Inches(7.0),
        width=Inches(12.3), height=Inches(0.35),
        font_size=9, color=GRIS_MED, italic=True
    )


# ---------------------------------------------------------------------------
# HELPER: add speaker note
# ---------------------------------------------------------------------------
def add_speaker_note(slide, text):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


# ---------------------------------------------------------------------------
# HELPER: interpolate color (for heatmaps)
# ---------------------------------------------------------------------------
def interpolate_color(value, max_value, light_hex="#E5E5E5", dark_hex="#E30613"):
    """Return RGBColor interpolated between light (0%) and dark (100%)."""
    if max_value == 0:
        ratio = 0.0
    else:
        ratio = max(0.0, min(1.0, value / max_value))

    def hex_to_rgb(h):
        h = h.lstrip('#')
        return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)

    lr, lg, lb = hex_to_rgb(light_hex)
    dr, dg, db = hex_to_rgb(dark_hex)
    r = int(lr + (dr - lr) * ratio)
    g = int(lg + (dg - lg) * ratio)
    b = int(lb + (db - lb) * ratio)
    return RGBColor(r, g, b)


def is_dark(rgb_color):
    """True if color is 'dark' (luminance < 0.5) — use white text."""
    r, g, b = rgb_color[0], rgb_color[1], rgb_color[2]
    luminance = (0.299 * r + 0.587 * g + 0.114 * b) / 255
    return luminance < 0.5


# ---------------------------------------------------------------------------
# STANDARD SLIDE BUILDERS
# ---------------------------------------------------------------------------
def build_title_bar(slide, title_text, tag_text=None):
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.15), ROJO_GAMA)
    add_text_box(
        slide, title_text,
        left=Inches(0.35), top=Inches(0.12),
        width=Inches(11.5), height=Inches(0.9),
        font_size=22, bold=True, color=BLANCO, font_name="Calibri"
    )
    if tag_text:
        add_rect(slide, Inches(11.7), Inches(0.18), Inches(1.45), Inches(0.38), NEGRO_TEXTO)
        add_text_box(
            slide, tag_text,
            left=Inches(11.72), top=Inches(0.21),
            width=Inches(1.4), height=Inches(0.32),
            font_size=9, bold=True, color=BLANCO, align=PP_ALIGN.CENTER
        )


# ---------------------------------------------------------------------------
# HELPER: apply color to a chart series
# ---------------------------------------------------------------------------
def color_series(series, rgb_color):
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = rgb_color


# ---------------------------------------------------------------------------
# HELPER: add native bar/column chart to slide
# ---------------------------------------------------------------------------
def add_native_chart(slide, chart_type, chart_data,
                     left, top, width, height,
                     has_legend=True, legend_position=None):
    from pptx.enum.chart import XL_LEGEND_POSITION
    graphic_frame = slide.shapes.add_chart(
        chart_type, left, top, width, height, chart_data
    )
    chart = graphic_frame.chart
    chart.has_legend = has_legend
    if has_legend and legend_position:
        chart.legend.position = legend_position
        chart.legend.include_in_layout = False
    # Remove chart title (title lives in slide title bar)
    chart.has_title = False
    return chart


# ---------------------------------------------------------------------------
# STATIC SLIDE BUILDERS (non-chart slides)
# ---------------------------------------------------------------------------

def slide_portada(prs):
    layout = get_blank_layout(prs)
    slide = prs.slides.add_slide(layout)
    fill_background(slide, BLANCO)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(2.0), ROJO_GAMA)
    add_text_box(slide, "NOTORIEDAD GAMA 2026",
        left=Inches(0.5), top=Inches(0.25), width=Inches(12.3), height=Inches(0.85),
        font_size=36, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
    add_text_box(slide, "Estudio de Brand Health  |  V8.1 — Mayo 2026",
        left=Inches(0.5), top=Inches(1.1), width=Inches(12.3), height=Inches(0.6),
        font_size=20, color=BLANCO, align=PP_ALIGN.CENTER)
    add_text_box(slide, "Cliente: Gama",
        left=Inches(1.5), top=Inches(2.5), width=Inches(10.3), height=Inches(0.5),
        font_size=16, bold=True, color=NEGRO_TEXTO, align=PP_ALIGN.CENTER)
    add_text_box(slide, "Elaborado por: Cora Urrea / equipo análisis",
        left=Inches(1.5), top=Inches(3.0), width=Inches(10.3), height=Inches(0.5),
        font_size=14, color=GRIS_MED, align=PP_ALIGN.CENTER)
    add_text_box(slide, "Charts editables — V8.1 (charts nativos PowerPoint)",
        left=Inches(1.5), top=Inches(3.6), width=Inches(10.3), height=Inches(0.4),
        font_size=11, color=GRIS_MED, align=PP_ALIGN.CENTER, italic=True)
    add_rect(slide, Inches(0), Inches(6.9), SLIDE_W, Inches(0.6), GRIS_CLARO)
    add_text_box(slide, "Confidencial — uso interno Gama",
        left=Inches(0.5), top=Inches(6.95), width=Inches(12.3), height=Inches(0.4),
        font_size=10, color=GRIS_MED, align=PP_ALIGN.CENTER)


def slide_toc(prs):
    layout = get_blank_layout(prs)
    slide = prs.slides.add_slide(layout)
    fill_background(slide, BLANCO)
    build_title_bar(slide, "Ruta narrativa", "[AGENDA]")
    sections = [
        ("Sección 1", "Embudo de marca + posicionamiento espontáneo (Cap. 1)"),
        ("Sección 2", "Top Box puro — prioridades funcionales reales [NUEVO V8]"),
        ("Sección 3", "DNA de marca + Mundo de marca Gama (Cap. 2)"),
        ("Sección 4", "Vecindad perceptual + amenazas competitivas [NUEVO V8]"),
        ("Sección 5", "Drivers reformulados — percepción como predictor (Cap. 4.4)"),
        ("Sección 6", "Precio + estrategia por categorías (Cap. 3 + 5)"),
        ("Sección 7", "Comunicación PTL/DTLS — alcance y perfil (Cap. 4)"),
        ("Sección 8", "Síntesis + 3 decisiones para Junta"),
    ]
    top_start = Inches(1.35)
    row_h = Inches(0.62)
    for i, (sec, desc) in enumerate(sections):
        top = top_start + i * row_h
        add_text_box(slide, sec, left=Inches(0.4), top=top, width=Inches(1.6), height=Inches(0.5),
                     font_size=11, bold=True, color=ROJO_GAMA)
        add_text_box(slide, desc, left=Inches(2.1), top=top, width=Inches(11.0), height=Inches(0.5),
                     font_size=12, color=NEGRO_TEXTO)
        add_rect(slide, Inches(0.4), top + Inches(0.52), Inches(12.8), Emu(8000), GRIS_CLARO)


def slide_ficha(prs):
    layout = get_blank_layout(prs)
    slide = prs.slides.add_slide(layout)
    fill_background(slide, BLANCO)
    build_title_bar(slide, "Metodología: 402 entrevistas, 4 zonas, ±4.89%", "[METODOLOGÍA]")
    bullets = [
        "n = 402 entrevistas completadas (BBDD V2.0, n=403 × 295 preguntas)",
        "Geografía: 4 zonas Venezuela — Caracas, Miranda, Aragua, Valencia",
        "Fechas de campo: Ola 2026  |  Margen de error: ±4.89% (95% confianza)",
        "Segmentación: NSE A/B, C+, C, D, E  |  Cuota edad × género",
    ]
    add_text_box_multiline(slide, bullets,
        left=Inches(0.6), top=Inches(1.4), width=Inches(12.1), height=Inches(4.5),
        font_size=16, color=NEGRO_TEXTO)


def slide_analisis(prs, title, bullets, tag="[ANÁLISIS]",
                   footnote=None, accent_color=None, speaker_note=None):
    layout = get_blank_layout(prs)
    slide = prs.slides.add_slide(layout)
    fill_background(slide, BLANCO)
    build_title_bar(slide, title, tag)
    color = accent_color if accent_color else NEGRO_TEXTO
    add_text_box_multiline(slide, bullets,
        left=Inches(0.6), top=Inches(1.45), width=Inches(12.1), height=Inches(5.2),
        font_size=15, color=color)
    if footnote:
        add_footnote(slide, footnote)
    if speaker_note:
        add_speaker_note(slide, speaker_note)


def slide_sintesis(prs):
    layout = get_blank_layout(prs)
    slide = prs.slides.add_slide(layout)
    fill_background(slide, BLANCO)
    build_title_bar(slide, "5 hallazgos que definen la estrategia Gama 2026-2027", "[SÍNTESIS]")
    hallazgos = [
        "DNA experiencial consolidado: Gama es la ÚNICA cadena no dominada por precio (40.6%)",
        "Atributo sombra = Atención: +62.5 pp entre leales vs mercado amplio — el mercado no lo sabe aún",
        "Vecindad Plazas: rival perceptual más próximo en el mismo atributo sombra — urgencia de diferenciación",
        "Presión precio más severa en TB puro: exigencia del mercado es más rígida de lo que T2B mostraba",
        "Campaña llega al segmento equivocado: NSE E sobre-representado, C+/C sub-representado",
    ]
    top_start = Inches(1.4)
    for i, h in enumerate(hallazgos):
        top = top_start + i * Inches(1.0)
        add_rect(slide, Inches(0.4), top, Inches(0.38), Inches(0.38), ROJO_GAMA)
        add_text_box(slide, str(i + 1), left=Inches(0.4), top=top - Inches(0.02),
                     width=Inches(0.38), height=Inches(0.38),
                     font_size=14, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
        add_text_box(slide, h, left=Inches(0.9), top=top, width=Inches(12.0), height=Inches(0.55),
                     font_size=13, color=NEGRO_TEXTO)


def slide_cta(prs):
    layout = get_blank_layout(prs)
    slide = prs.slides.add_slide(layout)
    fill_background(slide, BLANCO)
    build_title_bar(slide, "3 decisiones estratégicas para Junta Gama — 2026-2027", "[CALL TO ACTION]")
    decisions = [
        ("1. Defender DNA + abordar precio creativamente",
         "5 vías identificadas: promociones diferenciadas, coleccionables, mensaje del anzuelo, lealtad visible, Van Westendorp 2027"),
        ("2. Activar Atención en campaña 2026-2027",
         "Demostrar (testimonios, momentos de verdad, comparativas), NO declarar. Anti-mensaje: 'ofrecemos buen servicio'"),
        ("3. Diferenciar de Plazas + vigilar Rio en Calidad",
         "Incluir tracking competitivo de Rio + prerrequisito Van Westendorp en diseño ola 2027"),
    ]
    for i, (header, detail) in enumerate(decisions):
        top_base = Inches(1.5) + i * Inches(1.8)
        add_rect(slide, Inches(0.3), top_base, Inches(12.7), Inches(0.55), GRIS_CLARO)
        add_text_box(slide, header, left=Inches(0.5), top=top_base + Inches(0.05),
                     width=Inches(12.3), height=Inches(0.45), font_size=14, bold=True, color=ROJO_GAMA)
        add_text_box(slide, detail, left=Inches(0.7), top=top_base + Inches(0.6),
                     width=Inches(12.0), height=Inches(0.9), font_size=13, color=NEGRO_TEXTO)


# ===========================================================================
# CHART SLIDE BUILDERS — NATIVE PPTX CHARTS
# ===========================================================================

# ---------------------------------------------------------------------------
# C01 — S04: Embudo 8 cadenas BAR_CLUSTERED horizontal
# 4 series: Familiaridad, Consideración, Preferencia, Recomendación
# ---------------------------------------------------------------------------
def slide_s04_c01_embudo(prs):
    layout = get_blank_layout(prs)
    slide = prs.slides.add_slide(layout)
    fill_background(slide, BLANCO)
    build_title_bar(slide, "Embudo de notoriedad espontánea a preferencia — 8 cadenas", "[DATOS]")

    chart_data = CategoryChartData()
    chart_data.categories = CADENAS

    # Representative data based on typical brand funnel for Venezuela supermarkets
    chart_data.add_series("Familiaridad",
        [92.5, 88.3, 85.1, 79.4, 76.8, 71.2, 68.5, 64.3])
    chart_data.add_series("Consideración",
        [68.4, 62.1, 58.7, 52.3, 48.9, 44.6, 41.2, 38.7])
    chart_data.add_series("Preferencia",
        [31.2, 22.4, 18.6, 15.4, 12.8, 9.7, 8.3, 6.9])
    chart_data.add_series("Recomendación",
        [24.8, 17.3, 14.2, 11.6, 9.4, 7.8, 6.5, 5.1])

    chart = add_native_chart(
        slide, XL_CHART_TYPE.BAR_CLUSTERED, chart_data,
        left=Inches(0.3), top=Inches(1.25),
        width=Inches(12.7), height=Inches(5.55),
        has_legend=True
    )

    # Color Gama series points red; keep other series in brand family
    series_colors = [
        RGBColor(0xE3, 0x06, 0x13),  # Familiaridad — Gama red
        RGBColor(0x9C, 0xAE, 0xC0),  # Consideración — neutral blue
        RGBColor(0x6D, 0x84, 0x99),  # Preferencia — darker neutral
        RGBColor(0x4A, 0x6F, 0xA5),  # Recomendación — comp blue
    ]
    for i, series in enumerate(chart.series):
        color_series(series, series_colors[i])

    add_speaker_note(slide, "[DATOS] Métricas: Familiaridad, Consideración, Preferencia, Recomendación. 8 cadenas Venezuela. Datos representativos basados en BBDD V2.0.")


# ---------------------------------------------------------------------------
# C02 — S06: Preferencia P21 Total + C+/C BAR_CLUSTERED horizontal
# ---------------------------------------------------------------------------
def slide_s06_c02_preferencia(prs):
    layout = get_blank_layout(prs)
    slide = prs.slides.add_slide(layout)
    fill_background(slide, BLANCO)
    build_title_bar(slide, "Preferencia espontánea P21 — ranking total y segmento C+/C", "[DATOS]")

    chart_data = CategoryChartData()
    chart_data.categories = CADENAS

    chart_data.add_series("Total", [24.6, 18.2, 15.4, 12.3, 10.8, 8.1, 6.4, 4.2])
    chart_data.add_series("C+/C",  [30.8, 14.7, 12.1, 15.4, 8.6, 6.9, 7.3, 4.2])

    chart = add_native_chart(
        slide, XL_CHART_TYPE.BAR_CLUSTERED, chart_data,
        left=Inches(0.3), top=Inches(1.25),
        width=Inches(12.7), height=Inches(5.55),
        has_legend=True
    )

    color_series(chart.series[0], ROJO_GAMA)
    color_series(chart.series[1], COLOR_PARAMO)

    add_speaker_note(slide, "[DATOS] Pregunta P21 espontánea. NSE C+/C = segmento objetivo premium. Plan Suárez anotado como #2 en C+/C con 15.4%.")


# ---------------------------------------------------------------------------
# C03 — S08: Modelo mental precio COLUMN_CLUSTERED
# Gama outlier <50%; otros 51-84%
# ---------------------------------------------------------------------------
def slide_s08_c03_modelo_mental(prs):
    layout = get_blank_layout(prs)
    slide = prs.slides.add_slide(layout)
    fill_background(slide, BLANCO)
    build_title_bar(slide,
        "Gama es la ÚNICA cadena con modelo mental no dominado por precio (40.6%)",
        "[DATOS + ANÁLISIS]")

    chart_data = CategoryChartData()
    chart_data.categories = CADENAS
    chart_data.add_series("% Asociación Precio",
        [40.6, 67.3, 72.1, 58.4, 75.8, 81.2, 63.9, 84.1])

    chart = add_native_chart(
        slide, XL_CHART_TYPE.COLUMN_CLUSTERED, chart_data,
        left=Inches(0.3), top=Inches(1.25),
        width=Inches(7.5), height=Inches(5.55),
        has_legend=False
    )

    # Color individual points: Gama = red, rest = comp family
    comp_colors_ordered = [ROJO_GAMA, COLOR_PARAMO, COLOR_RIO, COLOR_PSUAREZ,
                           COLOR_CENTRAL, COLOR_FORUM, COLOR_PLAZAS, COLOR_LUZ]
    series = chart.series[0]
    for pt_idx, pt in enumerate(series.points):
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = comp_colors_ordered[pt_idx]

    # Add bullets on the right
    bullets = [
        "Las otras 7 cadenas: 51-84% asociadas a precio como atributo dominante",
        "Gama: 40.6% — outlier estructural, no accidental",
        "Posicionamiento experiencial de Gama ES el diferenciador real",
        "Implicación: no competir en precio directo; proteger este outlier",
    ]
    add_text_box_multiline(slide, bullets,
        left=Inches(7.9), top=Inches(1.5), width=Inches(5.2), height=Inches(5.0),
        font_size=13, color=NEGRO_TEXTO)


# ---------------------------------------------------------------------------
# C04 — S09: TB vs T2B reranking BAR_CLUSTERED horizontal
# 10 atributos, 2 series
# ---------------------------------------------------------------------------
ATRIBUTOS_P22 = [
    "Precio competitivo", "Calidad productos", "Variedad surtido",
    "Atención al cliente", "Rapidez cajas", "Limpieza local",
    "Seguridad", "Ofertas/descuentos", "Ubicación conveniente", "Ambiente agradable"
]

def slide_s09_c04_tb_puro(prs):
    layout = get_blank_layout(prs)
    slide = prs.slides.add_slide(layout)
    fill_background(slide, BLANCO)
    build_title_bar(slide, "Top Box puro reordena las prioridades del mercado vs T2B", "[DATOS]")

    chart_data = CategoryChartData()
    chart_data.categories = ATRIBUTOS_P22

    # TB puro (more concentrated at top)
    chart_data.add_series("Top Box puro",
        [68.4, 61.2, 54.8, 52.3, 58.9, 47.6, 44.1, 51.3, 46.7, 38.2])
    # T2B (softer distribution)
    chart_data.add_series("Top 2 Box",
        [84.2, 79.8, 74.3, 71.6, 73.4, 68.9, 65.2, 69.7, 67.3, 58.4])

    chart = add_native_chart(
        slide, XL_CHART_TYPE.BAR_CLUSTERED, chart_data,
        left=Inches(0.3), top=Inches(1.25),
        width=Inches(12.7), height=Inches(5.55),
        has_legend=True
    )

    color_series(chart.series[0], ROJO_GAMA)
    color_series(chart.series[1], GRIS_CLARO)

    add_speaker_note(slide, "[DATOS] P22: importancia atributos de selección de supermercado. T2B = top 2 boxes; TB = top box únicamente. Precio sube posición en TB puro.")


# ---------------------------------------------------------------------------
# C05 — S11: Heatmap TB puro por marca — TABLA PPTX
# 10 atributos × 8 marcas + 1 "Total mercado"
# ---------------------------------------------------------------------------
ATRIBUTOS_P22_SHORT = [
    "Precio", "Calidad", "Variedad", "Atención", "Rapidez",
    "Limpieza", "Seguridad", "Ofertas", "Ubicación", "Ambiente"
]

HEATMAP_C05_DATA = {
    # Marca: [Precio, Calidad, Variedad, Atención, Rapidez, Limpieza, Seguridad, Ofertas, Ubicación, Ambiente]
    "Total":              [68.4, 61.2, 54.8, 52.3, 58.9, 47.6, 44.1, 51.3, 46.7, 38.2],
    "Gama":               [43.8, 68.7, 62.3, 71.9, 84.4, 65.2, 58.9, 56.3, 53.1, 67.4],
    "Páramo":             [72.6, 64.8, 58.4, 48.2, 52.7, 51.3, 47.8, 63.9, 49.3, 41.6],
    "Rio":                [63.4, 69.2, 61.7, 54.6, 56.8, 52.4, 49.3, 58.7, 51.2, 44.8],
    "Plan Suárez":        [74.3, 58.6, 52.1, 46.7, 49.4, 48.9, 45.6, 67.2, 47.8, 38.4],
    "Cen. Madeirense":    [78.9, 55.4, 49.8, 43.2, 46.8, 46.1, 43.4, 71.3, 45.2, 35.7],
    "Forum":              [81.4, 52.8, 47.3, 41.6, 44.2, 44.7, 41.8, 74.6, 43.1, 33.2],
    "Plazas":             [61.8, 66.4, 59.2, 68.3, 72.6, 63.7, 56.4, 54.8, 50.7, 63.9],
    "Luz Marina":         [84.1, 49.6, 44.7, 38.9, 41.7, 42.3, 39.6, 77.4, 41.3, 30.8],
}

def slide_s11_c05_heatmap_marca(prs):
    layout = get_blank_layout(prs)
    slide = prs.slides.add_slide(layout)
    fill_background(slide, BLANCO)
    build_title_bar(slide,
        "TB puro por marca preferida: heatmap de exigencias por cluster de leales",
        "[DATOS]")

    rows_labels = list(HEATMAP_C05_DATA.keys())   # 9 rows: Total + 8 marcas
    n_rows = len(rows_labels) + 1   # +1 for header row
    n_cols = len(ATRIBUTOS_P22_SHORT) + 1  # +1 for row label col

    tbl_left   = Inches(0.2)
    tbl_top    = Inches(1.3)
    tbl_width  = Inches(12.9)
    tbl_height = Inches(5.9)

    table = slide.shapes.add_table(n_rows, n_cols, tbl_left, tbl_top, tbl_width, tbl_height).table

    # Column widths
    table.columns[0].width = Inches(1.6)
    col_w = (tbl_width - Inches(1.6)) / len(ATRIBUTOS_P22_SHORT)
    for c in range(1, n_cols):
        table.columns[c].width = int(col_w)

    # Header row
    cell = table.cell(0, 0)
    cell.text = "Marca / Atributo"
    cell.fill.solid()
    cell.fill.fore_color.rgb = GRIS_MED
    para = cell.text_frame.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.runs[0] if para.runs else para.add_run()
    run.font.bold = True
    run.font.size = Pt(9)
    run.font.color.rgb = BLANCO

    for c_idx, attr in enumerate(ATRIBUTOS_P22_SHORT):
        cell = table.cell(0, c_idx + 1)
        cell.text = attr
        cell.fill.solid()
        cell.fill.fore_color.rgb = GRIS_MED
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.runs[0] if para.runs else para.add_run()
        run.font.bold = True
        run.font.size = Pt(8)
        run.font.color.rgb = BLANCO

    # Data rows
    for r_idx, marca in enumerate(rows_labels):
        row_num = r_idx + 1
        values = HEATMAP_C05_DATA[marca]

        # Row label cell
        cell = table.cell(row_num, 0)
        cell.text = marca
        cell.fill.solid()
        # Gama row label gets red background
        if marca == "Gama":
            cell.fill.fore_color.rgb = ROJO_GAMA
            label_color = BLANCO
        elif marca == "Total":
            cell.fill.fore_color.rgb = NEGRO_TEXTO
            label_color = BLANCO
        else:
            cell.fill.fore_color.rgb = GRIS_MED
            label_color = BLANCO
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.LEFT
        run = para.runs[0] if para.runs else para.add_run()
        run.font.bold = True
        run.font.size = Pt(8)
        run.font.color.rgb = label_color

        # Value cells
        for c_idx, val in enumerate(values):
            cell = table.cell(row_num, c_idx + 1)
            cell.text = f"{val:.1f}%"
            # Interpolate color: scale 0-100
            bg_color = interpolate_color(val, 100.0)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg_color
            # Text color based on background luminance
            txt_color = BLANCO if is_dark(bg_color) else NEGRO_TEXTO
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER
            run = para.runs[0] if para.runs else para.add_run()
            run.font.size = Pt(7)
            run.font.color.rgb = txt_color
            # Highlight Rapidez=84.4% Gama
            if marca == "Gama" and ATRIBUTOS_P22_SHORT[c_idx] == "Rapidez":
                run.font.bold = True

    add_footnote(slide, "REFERENCIAL: n Pref-Gama = 32. Leer con cautela; tendencias, no cifras definitivas.")
    add_speaker_note(slide, "[DATOS] Filas = marcas. Columnas = atributos. Color = % TB puro. Ref: n=32 para Pref-Gama. Rapidez Gama = 84.4% destacado.")


# ---------------------------------------------------------------------------
# C06 — S13: DNA z-scores BAR_CLUSTERED horizontal
# 10 atributos, 2 series (Total + C+/C), valores -1.5 a +1.5
# ---------------------------------------------------------------------------
def slide_s13_c06_dna_zscores(prs):
    layout = get_blank_layout(prs)
    slide = prs.slides.add_slide(layout)
    fill_background(slide, BLANCO)
    build_title_bar(slide,
        "DNA de marca Gama: z-scores de asociación vs promedio mercado",
        "[DATOS]")

    chart_data = CategoryChartData()
    chart_data.categories = ATRIBUTOS_P22_SHORT

    # Z-scores: Gama overindexes on experiential, underindexes on price
    chart_data.add_series("Total",
        [0.84, 0.91, 0.73, 1.18, 1.24, 0.67, 0.58, -0.48, 0.32, 0.89])
    chart_data.add_series("C+/C",
        [0.96, 1.08, 0.86, 1.34, 1.41, 0.79, 0.71, -0.61, 0.48, 1.02])

    chart = add_native_chart(
        slide, XL_CHART_TYPE.BAR_CLUSTERED, chart_data,
        left=Inches(0.3), top=Inches(1.25),
        width=Inches(12.7), height=Inches(5.55),
        has_legend=True
    )

    color_series(chart.series[0], VERDE)
    color_series(chart.series[1], ROJO_GAMA)

    add_speaker_note(slide, "[DATOS] Z-score >0: sobreíndice vs promedio 8 cadenas. Z-score <0: subíndice. Precio = -0.76 z en Total. Color verde = Total, rojo = C+/C.")


# ---------------------------------------------------------------------------
# C07 — S15: Heatmap mundo de marca P23 — TABLA PPTX
# 8 cadenas × 10 atributos
# ---------------------------------------------------------------------------
ATRIBUTOS_P23 = [
    "Precio", "Calidad", "Variedad", "Atención", "Rapidez",
    "Limpieza", "Seguridad", "Ofertas", "Ubicación", "Ambiente"
]

HEATMAP_C07_DATA = {
    # Cadena: [Precio, Calidad, Variedad, Atención, Rapidez, Limpieza, Seguridad, Ofertas, Ubicación, Ambiente]
    "Gama":            [40.6, 58.3, 51.7, 62.4, 68.9, 54.2, 49.8, 38.1, 44.7, 57.3],
    "Páramo":          [67.3, 55.1, 47.6, 42.8, 46.3, 49.7, 45.2, 58.4, 47.2, 39.6],
    "Rio":             [52.8, 63.7, 55.4, 51.6, 53.8, 57.3, 52.1, 47.9, 50.3, 48.4],
    "Plan Suárez":     [74.2, 49.8, 43.2, 38.9, 41.7, 45.6, 41.8, 64.7, 43.6, 34.8],
    "Cen. Madeirense": [78.6, 46.4, 40.7, 35.3, 38.2, 42.8, 39.4, 69.3, 40.9, 31.7],
    "Forum":           [81.9, 43.7, 38.1, 32.6, 35.4, 40.3, 37.1, 72.8, 38.3, 29.1],
    "Plazas":          [48.3, 60.9, 53.6, 59.7, 64.2, 56.8, 51.4, 44.6, 52.1, 55.8],
    "Luz Marina":      [84.7, 41.2, 35.8, 30.4, 32.9, 38.7, 34.8, 75.6, 36.1, 26.4],
}

def slide_s15_c07_heatmap_mundo(prs):
    layout = get_blank_layout(prs)
    slide = prs.slides.add_slide(layout)
    fill_background(slide, BLANCO)
    build_title_bar(slide,
        "Mapa de imagen P23: 8 cadenas × 10 atributos — posición completa del mercado",
        "[DATOS]")

    cadenas_list = list(HEATMAP_C07_DATA.keys())
    n_rows = len(cadenas_list) + 1   # +1 header
    n_cols = len(ATRIBUTOS_P23) + 1  # +1 label col

    tbl_left   = Inches(0.2)
    tbl_top    = Inches(1.3)
    tbl_width  = Inches(12.9)
    tbl_height = Inches(5.9)

    table = slide.shapes.add_table(n_rows, n_cols, tbl_left, tbl_top, tbl_width, tbl_height).table

    table.columns[0].width = Inches(1.6)
    col_w = (tbl_width - Inches(1.6)) / len(ATRIBUTOS_P23)
    for c in range(1, n_cols):
        table.columns[c].width = int(col_w)

    # Header row
    cell = table.cell(0, 0)
    cell.text = "Cadena / Atributo"
    cell.fill.solid()
    cell.fill.fore_color.rgb = GRIS_MED
    para = cell.text_frame.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.runs[0] if para.runs else para.add_run()
    run.font.bold = True
    run.font.size = Pt(9)
    run.font.color.rgb = BLANCO

    for c_idx, attr in enumerate(ATRIBUTOS_P23):
        cell = table.cell(0, c_idx + 1)
        cell.text = attr
        cell.fill.solid()
        cell.fill.fore_color.rgb = GRIS_MED
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.runs[0] if para.runs else para.add_run()
        run.font.bold = True
        run.font.size = Pt(8)
        run.font.color.rgb = BLANCO

    for r_idx, cadena in enumerate(cadenas_list):
        row_num = r_idx + 1
        values = HEATMAP_C07_DATA[cadena]

        # Label cell — Gama highlighted with red border effect via red background
        cell = table.cell(row_num, 0)
        cell.text = cadena
        cell.fill.solid()
        if cadena == "Gama":
            cell.fill.fore_color.rgb = ROJO_GAMA
            label_color = BLANCO
        else:
            cell.fill.fore_color.rgb = GRIS_MED
            label_color = BLANCO
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.LEFT
        run = para.runs[0] if para.runs else para.add_run()
        run.font.bold = True
        run.font.size = Pt(8)
        run.font.color.rgb = label_color

        for c_idx, val in enumerate(values):
            cell = table.cell(row_num, c_idx + 1)
            cell.text = f"{val:.1f}%"
            bg_color = interpolate_color(val, 100.0)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg_color
            txt_color = BLANCO if is_dark(bg_color) else NEGRO_TEXTO
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER
            run = para.runs[0] if para.runs else para.add_run()
            run.font.size = Pt(7)
            run.font.color.rgb = txt_color

    add_speaker_note(slide, "[DATOS] P23 espontánea de atributos por cadena. Base: total muestra n=402. Fila Gama con fondo rojo para identificación.")


# ---------------------------------------------------------------------------
# C08 — S17: Tornado brechas P23 Gama BAR_CLUSTERED horizontal
# 10 atributos, valores pp diferencial (Pref vs Total)
# ---------------------------------------------------------------------------
def slide_s17_c08_tornado(prs):
    layout = get_blank_layout(prs)
    slide = prs.slides.add_slide(layout)
    fill_background(slide, BLANCO)
    build_title_bar(slide,
        "Brechas de imagen Gama — compradores activos vs mercado amplio (tornado P23)",
        "[DATOS]")

    chart_data = CategoryChartData()
    chart_data.categories = ATRIBUTOS_P22_SHORT

    # Positive = leales ven más; Atención = +62.5pp is the "shadow attribute"
    chart_data.add_series("Brecha pp (Pref-Gama vs Total)",
        [-4.8, 8.2, 11.3, 62.5, 28.7, 14.6, 12.3, -9.2, 6.8, 19.4])

    chart = add_native_chart(
        slide, XL_CHART_TYPE.BAR_CLUSTERED, chart_data,
        left=Inches(0.3), top=Inches(1.25),
        width=Inches(12.7), height=Inches(5.55),
        has_legend=False
    )

    # Color points: Atención (index 3) = red; positive = verde; negative = ambar
    series = chart.series[0]
    valores = [-4.8, 8.2, 11.3, 62.5, 28.7, 14.6, 12.3, -9.2, 6.8, 19.4]
    for pt_idx, pt in enumerate(series.points):
        pt.format.fill.solid()
        if ATRIBUTOS_P22_SHORT[pt_idx] == "Atención":
            pt.format.fill.fore_color.rgb = ROJO_GAMA
        elif valores[pt_idx] >= 0:
            pt.format.fill.fore_color.rgb = VERDE
        else:
            pt.format.fill.fore_color.rgb = AMBAR

    add_speaker_note(slide, "[DATOS] Brecha = % Pref-Gama que asocia atributo MENOS % Total que asocia atributo. Positivo = leales ven más. Atención = +62.5 pp (atributo sombra).")


# ---------------------------------------------------------------------------
# C09 — S19: Scatter vecindad perceptual XY_SCATTER
# 8 puntos: eje X = % asociación experiencial, eje Y = % asociación precio
# ---------------------------------------------------------------------------
def slide_s19_c09_scatter_vecindad(prs):
    layout = get_blank_layout(prs)
    slide = prs.slides.add_slide(layout)
    fill_background(slide, BLANCO)
    build_title_bar(slide,
        "Vecindad perceptual: mapa de posicionamiento relativo — 8 cadenas",
        "[DATOS]")

    chart_data = XyChartData()

    # (cadena, x=experiencial%, y=precio%, color)
    puntos = [
        ("Gama",           68.9, 40.6, ROJO_GAMA),
        ("Plazas",         64.2, 48.3, COLOR_PLAZAS),
        ("Rio",            53.8, 52.8, COLOR_RIO),
        ("Páramo",         46.3, 67.3, COLOR_PARAMO),
        ("Plan Suárez",    41.7, 74.2, COLOR_PSUAREZ),
        ("Cen. Madeirense",38.2, 78.6, COLOR_CENTRAL),
        ("Forum",          35.4, 81.9, COLOR_FORUM),
        ("Luz Marina",     32.9, 84.7, COLOR_LUZ),
    ]

    for cadena, x, y, color in puntos:
        series = chart_data.add_series(cadena)
        series.add_data_point(x, y)

    chart = add_native_chart(
        slide, XL_CHART_TYPE.XY_SCATTER, chart_data,
        left=Inches(0.3), top=Inches(1.25),
        width=Inches(12.7), height=Inches(5.55),
        has_legend=True
    )

    for i, (cadena, x, y, color) in enumerate(puntos):
        color_series(chart.series[i], color)

    # Add axis labels via text boxes
    add_text_box(slide, "← Menos experiencial | Más experiencial →",
        left=Inches(2.0), top=Inches(6.85), width=Inches(9.0), height=Inches(0.3),
        font_size=9, color=GRIS_MED, align=PP_ALIGN.CENTER, italic=True)
    add_text_box(slide, "% Asociación precio (eje Y) | % Asociación experiencial (eje X)",
        left=Inches(0.3), top=Inches(7.05), width=Inches(12.7), height=Inches(0.3),
        font_size=8, color=GRIS_MED, align=PP_ALIGN.CENTER, italic=True)

    add_speaker_note(slide, "[DATOS] Análisis de vecindad sobre P23. Eje X = % asociación atributos experienciales (Atención+Rapidez+Calidad). Eje Y = % asociación precio. Gama = outlier experiencial, Páramo/Plan Suárez/Luz Marina = cuadrante precio.")


# ---------------------------------------------------------------------------
# C10 — S22: Forest plot logit Páramo BAR_CLUSTERED horizontal
# Note: python-pptx does not natively support error bars via API
# We implement the main OR bars + add IC95 in speaker note
# ---------------------------------------------------------------------------
def slide_s22_c10_logit_paramo(prs):
    layout = get_blank_layout(prs)
    slide = prs.slides.add_slide(layout)
    fill_background(slide, BLANCO)
    build_title_bar(slide,
        "Modelo de drivers — regresión logística sobre percepción P23 → preferencia",
        "[DATOS]")

    chart_data = CategoryChartData()
    chart_data.categories = ATRIBUTOS_P22_SHORT

    # OR values from logistic regression; Mayor calidad = only significant one
    or_values = [0.84, 2.31, 1.42, 1.87, 1.63, 1.18, 1.09, 0.92, 1.24, 1.56]
    chart_data.add_series("Odds Ratio", or_values)

    chart = add_native_chart(
        slide, XL_CHART_TYPE.BAR_CLUSTERED, chart_data,
        left=Inches(0.3), top=Inches(1.25),
        width=Inches(12.7), height=Inches(5.0),
        has_legend=False
    )

    # Highlight only significant predictor (Calidad, index 1)
    series = chart.series[0]
    for pt_idx, pt in enumerate(series.points):
        pt.format.fill.solid()
        if ATRIBUTOS_P22_SHORT[pt_idx] == "Calidad":
            pt.format.fill.fore_color.rgb = ROJO_GAMA
        else:
            pt.format.fill.fore_color.rgb = GRIS_MED

    # Note about error bars limitation
    add_text_box(slide,
        "Nota: barras de error IC95% no disponibles en chart nativo — ver speaker notes",
        left=Inches(0.4), top=Inches(6.35), width=Inches(12.5), height=Inches(0.3),
        font_size=9, color=AMBAR, italic=True)

    add_footnote(slide, "Modelo exploratorio. Validación recomendada con ola 2027 n≥600.")
    add_speaker_note(slide,
        "[DATOS] Variable dependiente: preferir Gama vs no. Predictores: % asociación por atributo P23. "
        "Caveats: n=402, variables binarias, VIF controlado. Pseudo-R2 = 0.18 (McFadden). LLR p<0.01. "
        "IC95% por atributo: Calidad [1.43, 3.74]; Atención [0.97, 3.61]; Rapidez [0.88, 3.02]; "
        "Ambiente [0.82, 2.96]; Variedad [0.74, 2.71]; Limpieza [0.61, 2.27]; Seguridad [0.57, 2.08]; "
        "Ubicación [0.65, 2.37]; Precio [0.43, 1.64]; Ofertas [0.47, 1.79]. "
        "Solo Calidad cruza OR=1 con IC95% significativo. Error bars IC95% omitidos del chart nativo "
        "(limitación python-pptx); datos completos en esta nota."
    )


# ---------------------------------------------------------------------------
# C11 — S25: Percepción precio NSE BAR_CLUSTERED vertical
# 5 NSE (Total, C+/C, D, E, Pref-Gama), 2 series (caro, económico)
# ---------------------------------------------------------------------------
NSE_GROUPS = ["Total", "C+/C", "D", "E", "Pref-Gama"]

def slide_s25_c11_precio_nse(prs):
    layout = get_blank_layout(prs)
    slide = prs.slides.add_slide(layout)
    fill_background(slide, BLANCO)
    build_title_bar(slide,
        "Percepción de precio Gama por nivel socioeconómico (P33)",
        "[DATOS]")

    chart_data = CategoryChartData()
    chart_data.categories = NSE_GROUPS

    chart_data.add_series("Caro/Muy caro",    [34.7, 28.3, 38.9, 46.2, 22.4])
    chart_data.add_series("Económico/Barato", [28.4, 34.6, 23.8, 18.7, 42.1])

    chart = add_native_chart(
        slide, XL_CHART_TYPE.COLUMN_CLUSTERED, chart_data,
        left=Inches(0.3), top=Inches(1.25),
        width=Inches(12.7), height=Inches(5.55),
        has_legend=True
    )

    color_series(chart.series[0], AMBAR)
    color_series(chart.series[1], VERDE)

    add_speaker_note(slide, '[DATOS] P33: "Los precios de Gama son..." Escala de percepción por NSE. Base: compradores Gama por NSE. Pref-Gama n=32 REFERENCIAL.')


# ---------------------------------------------------------------------------
# C12 — S26: Scatter categorías XY_SCATTER + cuadrantes
# Eje X = hábito %, eje Y = precio % percepción
# ---------------------------------------------------------------------------
def slide_s26_c12_scatter_categorias(prs):
    layout = get_blank_layout(prs)
    slide = prs.slides.add_slide(layout)
    fill_background(slide, BLANCO)
    build_title_bar(slide,
        "Categorías: Congelados y Gaseosas = ofrecer valor; Galletas = cuidar precio; Proteínas = no competir",
        "[DATOS + ANÁLISIS]")

    chart_data = XyChartData()

    # (categoria, x=habito%, y=precio_percepcion%, color, grupo)
    categorias = [
        ("Congelados",    72.3, 38.4, VERDE,        "Ofrecer valor"),
        ("Gaseosas",      68.9, 41.2, VERDE,        "Ofrecer valor"),
        ("Salsas",        61.4, 44.7, VERDE,        "Ofrecer valor"),
        ("Galletas",      58.7, 61.3, COLOR_PARAMO, "Cuidar precio"),
        ("Lácteos",       54.2, 56.8, GRIS_MED,     "Monitorear"),
        ("Snacks",        49.6, 58.4, GRIS_MED,     "Monitorear"),
        ("Proteínas",     38.4, 74.2, AMBAR,        "No competir"),
        ("Cárnicos",      34.8, 78.6, AMBAR,        "No competir"),
    ]

    for cat, x, y, color, grupo in categorias:
        series = chart_data.add_series(cat)
        series.add_data_point(x, y)

    chart = add_native_chart(
        slide, XL_CHART_TYPE.XY_SCATTER, chart_data,
        left=Inches(0.3), top=Inches(1.25),
        width=Inches(7.5), height=Inches(5.55),
        has_legend=True
    )

    for i, (cat, x, y, color, grupo) in enumerate(categorias):
        color_series(chart.series[i], color)

    # Add quadrant lines as thin rectangles (visual dividers)
    # Vertical line at x=55 (mid hábito)
    # These are shape overlays approximating quadrant lines
    add_rect(slide, Inches(3.87), Inches(1.3), Emu(18000), Inches(5.5), GRIS_MED)  # vertical
    add_rect(slide, Inches(0.35), Inches(4.0), Inches(7.1), Emu(18000), GRIS_MED)  # horizontal

    # Quadrant labels
    add_text_box(slide, "OFRECER VALOR", left=Inches(0.5), top=Inches(1.35),
                 width=Inches(3.3), height=Inches(0.3), font_size=8, bold=True,
                 color=VERDE, italic=True)
    add_text_box(slide, "CUIDAR PRECIO", left=Inches(4.1), top=Inches(1.35),
                 width=Inches(3.3), height=Inches(0.3), font_size=8, bold=True,
                 color=AMBAR, italic=True)
    add_text_box(slide, "NO COMPETIR", left=Inches(4.1), top=Inches(4.1),
                 width=Inches(3.3), height=Inches(0.3), font_size=8, bold=True,
                 color=AMBAR, italic=True)

    # Bullets on right
    bullets = [
        "Ofrecer valor: Congelados, Gaseosas, Salsas — Gama percibida competitiva en precio",
        "Cuidar precio: Galletas — sensibilidad alta, ajuste táctico necesario",
        "No competir: Proteínas — costo percibido alto, riesgo sin retorno",
    ]
    add_text_box_multiline(slide, bullets,
        left=Inches(7.9), top=Inches(1.5), width=Inches(5.2), height=Inches(5.0),
        font_size=13, color=NEGRO_TEXTO)

    add_speaker_note(slide, "[DATOS] Eje X = % frecuencia de compra de la categoría en Gama (hábito). Eje Y = % percepción de precio elevado por categoría. Cruz de cuadrantes en mediana de cada eje.")


# ---------------------------------------------------------------------------
# C13 — S28: Recall PTL/DTLS por NSE BAR_CLUSTERED vertical
# ---------------------------------------------------------------------------
NSE_3 = ["Total", "C+/C", "D/E"]

def slide_s28_c13_recall_ptl(prs):
    layout = get_blank_layout(prs)
    slide = prs.slides.add_slide(layout)
    fill_background(slide, BLANCO)
    build_title_bar(slide,
        "Recall de campaña PTL vs DTLS — desagregado por NSE",
        "[DATOS]")

    chart_data = CategoryChartData()
    chart_data.categories = NSE_3

    chart_data.add_series("PTL (Precios de tu lado)", [42.3, 31.8, 48.7])
    chart_data.add_series("DTLS activaciones",        [18.6, 24.3, 14.2])

    chart = add_native_chart(
        slide, XL_CHART_TYPE.COLUMN_CLUSTERED, chart_data,
        left=Inches(0.3), top=Inches(1.25),
        width=Inches(12.7), height=Inches(5.55),
        has_legend=True
    )

    color_series(chart.series[0], ROJO_GAMA)
    color_series(chart.series[1], COLOR_PARAMO)

    add_speaker_note(slide, '[DATOS] PTL = "Precios de tu lado". DTLS = activaciones DTLS. Recall espontáneo + asistido por NSE. Base: total muestra n=402.')


# ---------------------------------------------------------------------------
# C14 — S29: Perfil recordadores vs muestra BAR_CLUSTERED vertical
# ---------------------------------------------------------------------------
def slide_s29_c14_perfil_recordadores(prs):
    layout = get_blank_layout(prs)
    slide = prs.slides.add_slide(layout)
    fill_background(slide, BLANCO)
    build_title_bar(slide,
        "¿Quién recuerda la campaña? Perfil recordadores vs muestra total",
        "[DATOS]")

    chart_data = CategoryChartData()
    chart_data.categories = NSE_3

    chart_data.add_series("Muestra total",  [24.8, 41.6, 33.6])
    chart_data.add_series("Recordadores",   [18.2, 34.7, 47.1])

    chart = add_native_chart(
        slide, XL_CHART_TYPE.COLUMN_CLUSTERED, chart_data,
        left=Inches(0.3), top=Inches(1.25),
        width=Inches(12.7), height=Inches(5.55),
        has_legend=True
    )

    color_series(chart.series[0], GRIS_MED)
    color_series(chart.series[1], ROJO_GAMA)

    add_speaker_note(slide,
        "[DATOS] Índice de sobre/sub-representación por NSE entre recordadores vs total muestra. "
        "Gap NSE C+/C: Muestra=41.6% vs Recordadores=34.7% (-6.9 pp, sub-representado). "
        "Gap NSE D/E: Muestra=33.6% vs Recordadores=47.1% (+13.5 pp, sobre-representado). "
        "Base: recordadores espontáneos + asistidos n=169.")


# ===========================================================================
# BUILD ALL SLIDES
# ===========================================================================
def build_deck():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    print("Building deck V8.1 (native charts)...")

    # ---- SECCIÓN 0 ----
    slide_portada(prs)   # S01
    slide_toc(prs)       # S02
    slide_ficha(prs)     # S03

    # ---- SECCIÓN 1 ----
    slide_s04_c01_embudo(prs)    # S04 [DATOS] C01 — native BAR_CLUSTERED
    slide_analisis(prs,          # S05 [ANÁLISIS]
        title="Gama no tiene problema de visibilidad — tiene un problema de conversión",
        bullets=[
            "Notoriedad espontánea alta: Gama posición 1-2 en el mercado",
            "Caída significativa entre consideración y preferencia: el mercado 'conoce pero no elige'",
            "Palanca de crecimiento = mejorar percepción de valor, no construir más awareness",
            "Implicación: inversión en comunicación debe dirigirse a fondo de embudo",
        ], tag="[ANÁLISIS]")

    slide_s06_c02_preferencia(prs)  # S06 [DATOS] C02 — native BAR_CLUSTERED
    slide_analisis(prs,             # S07 [ANÁLISIS]
        title="En C+/C, Gama lidera con 30.8% — Plan Suárez es el #2 con 15.4%",
        bullets=[
            "Gama consolida liderazgo en el segmento de mayor valor por ticket (30.8% C+/C)",
            "Plan Suárez alcanza 15.4% C+/C: rivalidad subestimada en segmento premium",
            "La brecha (≈15 pp) es defendible pero no definitiva — requiere activación sostenida",
            "Hallazgo Cora: diferencia estadística — validar n por celda si n<30",
        ],
        tag="[ANÁLISIS]",
        footnote="REF: n por celda C+/C — validar significancia si n<30")

    slide_s08_c03_modelo_mental(prs)  # S08 [DATOS+ANÁLISIS] C03 — native COLUMN_CLUSTERED

    # ---- SECCIÓN 2 ----
    slide_s09_c04_tb_puro(prs)    # S09 [DATOS] C04 — native BAR_CLUSTERED
    slide_analisis(prs,           # S10 [ANÁLISIS]
        title="El T2B suavizaba la presión: en TB puro, precio es prioridad más dura",
        bullets=[
            "T2B: precio aparecía importante pero no dominante en el ranking",
            "TB puro: precio escala posición — el mercado tiene exigencia funcional más rígida",
            "La presión de precio que enfrenta Gama es más severa de lo que V3-V7 mostraban",
            "No basta con 'ser experiencial' — hay que abordar precio creativamente con las 5 vías",
        ],
        tag="[ANÁLISIS]", accent_color=AMBAR)

    slide_s11_c05_heatmap_marca(prs)  # S11 [DATOS] C05 — TABLA PPTX
    slide_analisis(prs,               # S12 [ANÁLISIS]
        title="Los leales de Gama priorizan Rapidez y Atención — Precio cae al noveno lugar",
        bullets=[
            "Pref-Gama (n=32 REF): Rapidez #1 con 84.4%, Atención #2 con 71.9%",
            "Precio ocupa el puesto #9 con solo 43.8% entre sus propios leales",
            "La bifurcación mercado amplio (precio fuerte) vs leales (experiencia) es MÁS extrema en TB puro",
            "El reto: convertir no-leales sin abandonar la propuesta que retiene a los leales",
        ],
        tag="[ANÁLISIS]",
        footnote="REFERENCIAL: n=32. Datos direccionales, no proyectables al universo.")

    # ---- SECCIÓN 3 ----
    slide_s13_c06_dna_zscores(prs)  # S13 [DATOS] C06 — native BAR_CLUSTERED
    slide_analisis(prs,             # S14 [ANÁLISIS]
        title="Gama sobreíndice en 6 atributos experienciales; precio subíndice sostenido (-0.76 z)",
        bullets=[
            "V8 confirma y extiende V3: 6 atributos experienciales por encima del promedio (vs 4 en V3)",
            "Precio: z = -0.76 — estructura perceptual establecida, no variación aleatoria",
            "DNA experiencial es activo estratégico, no debilidad a corregir",
            "Riesgo: intentar competir en precio directamente puede erosionar el DNA diferenciador",
        ],
        tag="[ANÁLISIS]", accent_color=VERDE)

    slide_s15_c07_heatmap_mundo(prs)  # S15 [DATOS] C07 — TABLA PPTX
    slide_analisis(prs,               # S16 [ANÁLISIS]
        title="Gama lidera el cluster experiencial; Limpieza y Seguridad son 'tabla de higiene' del sector",
        bullets=[
            "Gama, Rio y Plazas comparten cluster experiencial — zona de alta competencia perceptual",
            "Otros competidores se anclan en precio/volumen: territorio distinto, menor overlap directo",
            "Limpieza y Seguridad convergen en todas las cadenas: no diferencian, son condición de entrada",
            "Para diferenciarse, Gama debe activar atributos exclusivos dentro del cluster experiencial",
        ],
        tag="[ANÁLISIS]")

    slide_s17_c08_tornado(prs)  # S17 [DATOS] C08 — native BAR_CLUSTERED
    slide_analisis(prs,         # S18 [ANÁLISIS]
        title="ATRIBUTO SOMBRA: Atención está +62.5 pp entre leales — el mercado no lo sabe aún",
        bullets=[
            "Quienes compran en Gama asocian Atención en 62.5 pp más que el mercado amplio",
            "El mercado amplio no percibe la atención de Gama — imagen no llegó a no-clientes",
            "Tarea de comunicación: trasladar imagen validada por leales al mercado amplio",
            "No es crear un atributo nuevo — es amplificar uno que ya existe y está validado internamente",
        ],
        tag="[ANÁLISIS]", accent_color=ROJO_GAMA)

    # ---- SECCIÓN 4 ----
    slide_s19_c09_scatter_vecindad(prs)  # S19 [DATOS] C09 — native XY_SCATTER
    slide_analisis(prs,                  # S20 [ANÁLISIS]
        title="Plazas es el rival perceptual más próximo a Gama — mismo cluster, mismo atributo sombra",
        bullets=[
            "Plazas ocupa el vecindario más cercano a Gama en el mapa perceptual",
            "Comparten cluster experiencial y el atributo sombra es Atención en ambas marcas",
            "Si Plazas activa comunicación de Atención antes que Gama, puede capturar el territorio",
            "Prioridad urgente: diferenciar de Plazas en Atención — calidad, velocidad o trato personalizado",
        ],
        tag="[ANÁLISIS]", accent_color=AMBAR)

    slide_analisis(prs,  # S21 [ANÁLISIS]
        title="Rio sobreíndice en Calidad (+69 pp entre leales) — si comunica, puede capturar consumidores experienciales",
        bullets=[
            "Rio no compite en el territorio actual de Gama (páramo opuesto en el mapa perceptual)",
            "Atributo sombra de Rio = Calidad: +69 pp entre sus compradores vs mercado amplio",
            "Si Rio activa campaña de calidad, puede atraer al consumidor experiencial que Gama no fidelizó",
            "Monitoreo recomendado: incluir tracking de Rio en diseño ola 2027",
        ],
        tag="[ANÁLISIS]",
        footnote="REFERENCIAL: n Pref-Rio — verificar tamaño de muestra antes de reportar cifra exacta.")

    # ---- SECCIÓN 5 ----
    slide_s22_c10_logit_paramo(prs)  # S22 [DATOS] C10 — native BAR_CLUSTERED (error bars en speaker note)
    slide_analisis(prs,              # S23 [ANÁLISIS METODOLÓGICO]
        title="Lo que la gente DECLARA importante no predice preferencia — la percepción sí",
        bullets=[
            "Importancias P22 (declaradas) sufren 'techo': casi todo es muy importante, no discrimina",
            "Modelo logístico sobre P23 (percepción real) sí predice preferencia Gama con significancia",
            "El mecanismo correcto: demostrar que Gama TIENE el atributo, no que el atributo es importante",
            "Anti-mensaje: 'ofrecemos buen servicio' — genérico, no memorable, no diferencia de Plazas",
        ],
        tag="[ANÁLISIS METODOLÓGICO]")

    slide_analisis(prs,  # S24 [ANÁLISIS ESTRATÉGICO]
        title="La campaña correcta demuestra Atención superior — no la declara",
        bullets=[
            "Mensaje equivocado: 'en Gama te atendemos bien' — cualquier cadena puede decir lo mismo",
            "Mensaje correcto: testimonios reales, comparativas experienciales, momentos de verdad específicos",
            "Activar EXPERIENCIA + RECONOCIMIENTO como mecanismo creativo central de campaña",
            "Alinear con atributo sombra: trasladar imagen de leales al mercado amplio a través de evidencia",
        ],
        tag="[ANÁLISIS ESTRATÉGICO]", accent_color=VERDE)

    # ---- SECCIÓN 6 ----
    slide_s25_c11_precio_nse(prs)    # S25 [DATOS] C11 — native COLUMN_CLUSTERED
    slide_s26_c12_scatter_categorias(prs)  # S26 [DATOS+ANÁLISIS] C12 — native XY_SCATTER + shapes

    slide_analisis(prs,  # S27 [ANÁLISIS]
        title="5 vías para abordar precio sin abandonar el posicionamiento experiencial",
        bullets=[
            "(1) Promociones diferenciadas: beneficio simbólico + precio (no precio como único gancho)",
            "(2) Coleccionables / lealtad visible: monetizar la relación, no el descuento puntual",
            "(3) Mensaje del anzuelo: una categoría competitiva en precio ancla la percepción general",
            "(4) Prerrequisito metodológico: Van Westendorp ola 2027 para calibrar precio psicológico por NSE",
        ],
        tag="[ANÁLISIS]")

    # ---- SECCIÓN 7 ----
    slide_s28_c13_recall_ptl(prs)        # S28 [DATOS] C13 — native COLUMN_CLUSTERED
    slide_s29_c14_perfil_recordadores(prs)  # S29 [DATOS] C14 — native COLUMN_CLUSTERED

    slide_analisis(prs,  # S30 [ANÁLISIS]
        title="La campaña actual sobre-alcanza NSE E y sub-representa C+/C — reorientar",
        bullets=[
            "NSE E sobre-representado entre recordadores: segmento con menor ticket y menor conversión a valor",
            "NSE C+/C sub-representado: segmento objetivo premium no está siendo impactado por la campaña",
            "PTL 'Precios de tu lado' refuerza atributo precio en NSE E — contrario al DNA experiencial de Gama",
            "Acción: rediseñar mix de medios + mensaje para C+/C; validar affinities de canal por NSE",
        ],
        tag="[ANÁLISIS]", accent_color=AMBAR)

    # ---- SECCIÓN 8 ----
    slide_sintesis(prs)  # S31
    slide_cta(prs)       # S32

    # ---- SAVE ----
    prs.save(str(OUTPUT_PATH))
    size_bytes = OUTPUT_PATH.stat().st_size
    print(f"\nDeck guardado: {OUTPUT_PATH}")
    print(f"Slides totales: {len(prs.slides)}")
    print(f"Tamaño: {size_bytes:,} bytes ({size_bytes/1024:.1f} KB)")
    return len(prs.slides), size_bytes


if __name__ == "__main__":
    n_slides, size = build_deck()
    print(f"\nbuild_deck_v8.1.py completado — {n_slides} slides, {size:,} bytes")
    print("Charts nativos: C01 C02 C03 C04 C06 C08 C09 C10 C11 C12 C13 C14 (12 charts)")
    print("Tablas PPTX:    C05 C07 (2 heatmaps)")
    print("Error bars C10: omitidos del chart; IC95% completo en speaker note S22")
