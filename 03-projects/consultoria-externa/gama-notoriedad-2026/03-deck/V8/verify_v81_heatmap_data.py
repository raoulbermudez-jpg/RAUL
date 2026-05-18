"""Verifica si los heatmaps V8.1 tienen (a) colores aplicados a celdas y (b) datos correctos.
Comparar con valores reales del reporte CU-10.
"""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.dml.color import RGBColor

V81 = r"C:\rAUL\03-projects\consultoria-externa\gama-notoriedad-2026\03-deck\V8\2026-05-18_Notoriedad-Gama-2026_V8.1.pptx"

# Valores REALES del CU-10 reporte para Gama en P23 mundo de marca (Cap §4.1)
# Cadena | Surtido | Calidad | Precio | Atencion | Promoc | Atractiva | Limpieza | Seguro | Rapidez | V.Dinero
GAMA_P23_REAL = {
    'Surtido': 17.9, 'Calidad': 26.6, 'Precio': 7.2, 'Atencion': 21.9, 'Promoc': 9.0,
    'Atractiva': 28.9, 'Limpieza': 31.1, 'Seguro': 29.6, 'Rapidez': 21.1, 'V.Dinero': 11.4
}

pres = Presentation(V81)

print("="*90)
print("VERIFICACION HEATMAPS V8.1")
print("="*90)

for slide_idx in [10, 14]:  # S11 y S15
    s = pres.slides[slide_idx]
    print(f"\n>>> SLIDE {slide_idx+1} <<<")
    for shape in s.shapes:
        if shape.has_table:
            t = shape.table
            rows = list(t.rows)
            cols = list(t.columns)
            print(f"\nTabla {len(rows)}x{len(cols)} encontrada")

            # Inspeccionar fills de celdas (primeras 3 filas, primeras 3 cols)
            print("\nINSPECCION FILLS (sample 3x3):")
            for r_idx in range(min(3, len(rows))):
                for c_idx in range(min(4, len(cols))):
                    cell = rows[r_idx].cells[c_idx]
                    text = cell.text.strip()[:15]
                    fill_type = "none"
                    fill_color = "none"
                    try:
                        fill = cell.fill
                        if fill.type == 1:  # MSO_FILL.SOLID
                            fill_type = "solid"
                            try:
                                rgb = fill.fore_color.rgb
                                fill_color = f"#{rgb}"
                            except:
                                fill_color = "(theme color)"
                    except Exception as e:
                        fill_type = f"err:{e}"
                    print(f"  [{r_idx},{c_idx}] text='{text}' fill={fill_type} color={fill_color}")

            # Inspeccionar fila Gama (buscar)
            print("\nDATOS FILA GAMA:")
            for r_idx, row in enumerate(rows):
                first_cell_text = row.cells[0].text.strip()
                if 'gama' in first_cell_text.lower():
                    print(f"  Fila {r_idx}: {first_cell_text}")
                    headers = [rows[0].cells[c].text.strip() for c in range(len(cols))]
                    for c_idx in range(1, len(cols)):
                        h = headers[c_idx]
                        v = row.cells[c_idx].text.strip()
                        print(f"    {h}: '{v}'")
                    break

print("\n" + "="*90)
print("REFERENCIA REAL CU-10 §4.1 — Gama P23 imagen total:")
print("="*90)
for k, v in GAMA_P23_REAL.items():
    print(f"  {k}: {v}%")
