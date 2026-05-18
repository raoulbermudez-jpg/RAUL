"""Compara estructura de V8.pptx vs V8.1.pptx slide por slide.
Reporta # imágenes, # charts, # tablas, # textboxes, # shapes.
Identifica qué se perdió y qué quedó.
"""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

V8 = r"C:\rAUL\03-projects\consultoria-externa\gama-notoriedad-2026\03-deck\V8\2026-05-18_Notoriedad-Gama-2026_V8.pptx"
V81 = r"C:\rAUL\03-projects\consultoria-externa\gama-notoriedad-2026\03-deck\V8\2026-05-18_Notoriedad-Gama-2026_V8.1.pptx"

def analyze_slide(slide, slide_num):
    counts = {
        'images': 0,
        'charts': 0,
        'tables': 0,
        'textboxes': 0,
        'shapes_other': 0,
        'placeholders': 0,
    }
    img_sizes = []
    chart_types = []
    table_dims = []
    text_snippets = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                text_snippets.append(text[:80])
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            counts['images'] += 1
            try:
                img_sizes.append((shape.width.inches, shape.height.inches))
            except:
                pass
        elif shape.has_chart:
            counts['charts'] += 1
            try:
                chart_types.append(str(shape.chart.chart_type))
            except:
                chart_types.append('unknown')
        elif shape.has_table:
            counts['tables'] += 1
            try:
                t = shape.table
                table_dims.append((len(list(t.rows)), len(list(t.columns))))
            except:
                pass
        elif shape.has_text_frame:
            counts['textboxes'] += 1
        elif shape.is_placeholder:
            counts['placeholders'] += 1
        else:
            counts['shapes_other'] += 1
    return counts, img_sizes, chart_types, table_dims, text_snippets

def main():
    print("="*90)
    print(f"COMPARACION V8 (PNG charts) vs V8.1 (charts nativos)")
    print("="*90)
    pres_v8 = Presentation(V8)
    pres_v81 = Presentation(V81)
    print(f"\nV8 total slides: {len(pres_v8.slides)}")
    print(f"V8.1 total slides: {len(pres_v81.slides)}")
    print()
    n_compare = min(len(pres_v8.slides), len(pres_v81.slides))
    diffs = []
    for i in range(n_compare):
        s_v8 = pres_v8.slides[i]
        s_v81 = pres_v81.slides[i]
        c_v8, imgs_v8, charts_v8, tables_v8, texts_v8 = analyze_slide(s_v8, i+1)
        c_v81, imgs_v81, charts_v81, tables_v81, texts_v81 = analyze_slide(s_v81, i+1)
        # Title del slide (primer texto largo)
        title = (texts_v8[0] if texts_v8 else '(sin texto)')[:70]
        print(f"\n--- Slide {i+1:02d} ---  {title}")
        print(f"  V8:   img={c_v8['images']}  chart={c_v8['charts']}  table={c_v8['tables']}  txt={c_v8['textboxes']}")
        print(f"  V8.1: img={c_v81['images']}  chart={c_v81['charts']}  table={c_v81['tables']}  txt={c_v81['textboxes']}")
        if imgs_v8:
            sizes = ", ".join([f"{w:.1f}x{h:.1f}in" for w,h in imgs_v8])
            print(f"        V8 imagenes: {sizes}")
        if charts_v81:
            print(f"        V8.1 charts: {', '.join(charts_v81)}")
        if tables_v81:
            dims = ", ".join([f"{r}x{c}" for r,c in tables_v81])
            print(f"        V8.1 tablas: {dims}")
        # Flagging
        if c_v8['images'] != c_v81['images'] or c_v8['charts'] != c_v81['charts'] or c_v8['tables'] != c_v81['tables']:
            diffs.append({
                'slide': i+1,
                'title': title,
                'v8': f"img={c_v8['images']} chart={c_v8['charts']} table={c_v8['tables']}",
                'v81': f"img={c_v81['images']} chart={c_v81['charts']} table={c_v81['tables']}",
            })
    print("\n" + "="*90)
    print("RESUMEN DE DIFERENCIAS ESTRUCTURALES (slides con cambio)")
    print("="*90)
    for d in diffs:
        print(f"\nSlide {d['slide']:02d}: {d['title']}")
        print(f"  V8:   {d['v8']}")
        print(f"  V8.1: {d['v81']}")

    print("\n" + "="*90)
    print("INSPECCION DETALLADA — slides clave de heatmaps (S11, S15)")
    print("="*90)
    for slide_idx in [10, 14]:  # 0-indexed: S11=10, S15=14
        if slide_idx >= len(pres_v8.slides) or slide_idx >= len(pres_v81.slides):
            continue
        print(f"\n>>> Slide {slide_idx+1} <<<")
        s_v8 = pres_v8.slides[slide_idx]
        s_v81 = pres_v81.slides[slide_idx]
        print("\nV8 contenido textual:")
        c_v8, imgs_v8, _, _, texts_v8 = analyze_slide(s_v8, slide_idx+1)
        for t in texts_v8[:8]:
            print(f"  - {t}")
        print(f"\nV8 imagenes: {len(imgs_v8)} (tamano {imgs_v8 if imgs_v8 else 'n/a'})")
        print("\nV8.1 contenido textual:")
        c_v81, _, _, tables_v81, texts_v81 = analyze_slide(s_v81, slide_idx+1)
        for t in texts_v81[:8]:
            print(f"  - {t}")
        if tables_v81:
            print(f"\nV8.1 tablas: {tables_v81} (filas x columnas)")
            # Inspeccionar primera tabla
            for sh in s_v81.shapes:
                if sh.has_table:
                    t = sh.table
                    print(f"  Tabla {len(list(t.rows))}x{len(list(t.columns))}:")
                    # Mostrar primera fila y primera columna
                    rows = list(t.rows)
                    cols = list(t.columns)
                    if rows:
                        header = [c.text.strip()[:15] for c in rows[0].cells]
                        print(f"    Cabecera: {header}")
                    if len(rows) > 1:
                        first_data = [c.text.strip()[:15] for c in rows[1].cells]
                        print(f"    Primera fila datos: {first_data}")
                    break

if __name__ == '__main__':
    main()
