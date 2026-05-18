import sys
sys.stdout.reconfigure(encoding='utf-8')

import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData, XyChartData

# ---------------------------------------------------------------------------
# PATHS
# ---------------------------------------------------------------------------
BASE_DIR = Path(r"C:\rAUL\03-projects\consultoria-externa\gama-notoriedad-2026\03-deck\V8")
OUTPUT_PATH = BASE_DIR / "2026-05-18_Notoriedad-Gama-2026_V8.1_Resumen-Ejecutivo.pptx"

# ---------------------------------------------------------------------------
# BRAND COLORS (same as V8.1 deck)
# ---------------------------------------------------------------------------
ROJO_GAMA   = RGBColor(0xE3, 0x06, 0x13)
NEGRO_TEXTO = RGBColor(0x1A, 0x1A, 0x1A)
GRIS_MED    = RGBColor(0x6B, 0x6B, 0x6B)
GRIS_CLARO  = RGBColor(0xE5, 0xE5, 0xE5)
AMBAR       = RGBColor(0xF2, 0xA9, 0x00)
VERDE       = RGBColor(0x2D, 0x8F, 0x47)
BLANCO      = RGBColor(0xFF, 0xFF, 0xFF)

# Competitor colors (same palette as V8.1)
COLOR_PARAMO   = RGBColor(0x4A, 0x6F, 0xA5)
COLOR_RIO      = RGBColor(0x7C, 0x8B, 0x9B)
COLOR_PSUAREZ  = RGBColor(0x9C, 0xAE, 0xC0)
COLOR_CENTRAL  = RGBColor(0x5B, 0x70, 0x90)
COLOR_FORUM    = RGBColor(0x8F, 0xA3, 0xB8)
COLOR_PLAZAS   = RGBColor(0x6D, 0x84, 0x99)
COLOR_LUZ      = RGBColor(0xA8, 0xB9, 0xC7)

# ---------------------------------------------------------------------------
# SLIDE DIMENSIONS (16:9)
# ---------------------------------------------------------------------------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Chain order (C03 x-axis)
CADENAS = ["Gama", "Páramo", "Rio", "Plan Suárez", "Central Madeirense", "Forum", "Plazas", "Luz Marina"]
CADENA_COLORS = [ROJO_GAMA, COLOR_PARAMO, COLOR_RIO, COLOR_PSUAREZ,
                 COLOR_CENTRAL, COLOR_FORUM, COLOR_PLAZAS, COLOR_LUZ]

# Attribute labels (C08 y-axis)
ATRIBUTOS_P22_SHORT = [
    "Precio", "Calidad", "Variedad", "Atención", "Rapidez",
    "Limpieza", "Seguridad", "Ofertas", "Ubicación", "Ambiente"
]

# ---------------------------------------------------------------------------
# HELPERS
# ---------------------------------------------------------------------------

def get_blank_layout(prs):
    for layout in prs.slide_layouts:
        if layout.name.lower() in ("blank", "en blanco"):
            return layout
    return prs.slide_layouts[-1]


def fill_background(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_rect_no_fill(slide, left, top, width, height, line_color, line_width=Pt(1)):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.background()
    shape.line.color.rgb = line_color
    shape.line.width = line_width
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_size=14, bold=False, color=None, align=PP_ALIGN.LEFT,
                 font_name="Calibri", wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return txBox


def add_bullet_box(slide, bullets, left, top, width, height,
                   font_size=14, color=None, marker="•  "):
    if color is None:
        color = NEGRO_TEXTO
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for idx, b in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(5)
        run = p.add_run()
        run.text = marker + b
        run.font.name = "Calibri"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
    return txBox


def add_footnote(slide, text):
    add_text_box(slide, text,
                 left=Inches(0.4), top=Inches(6.9),
                 width=Inches(12.5), height=Inches(0.45),
                 font_size=9, color=GRIS_MED, italic=True)


def add_footer_bar(slide, label="Confidencial — uso interno Gama"):
    add_rect(slide, left=Inches(0), top=Inches(7.1),
             width=SLIDE_W, height=Inches(0.4), fill_color=GRIS_CLARO)
    add_text_box(slide, label,
                 left=Inches(0.4), top=Inches(7.12),
                 width=Inches(9), height=Inches(0.36),
                 font_size=9, color=GRIS_MED)


def slide_header(slide, title_text, tag_text=None, title_size=22):
    """Red left accent bar + title + optional tag + separator + footer."""
    # Left red accent
    add_rect(slide, left=Inches(0), top=Inches(0),
             width=Inches(0.12), height=SLIDE_H, fill_color=ROJO_GAMA)
    # Title
    add_text_box(slide, title_text,
                 left=Inches(0.25), top=Inches(0.2),
                 width=Inches(12.5), height=Inches(0.8),
                 font_size=title_size, bold=True, color=NEGRO_TEXTO)
    # Tag
    if tag_text:
        add_text_box(slide, tag_text,
                     left=Inches(11.2), top=Inches(0.22),
                     width=Inches(2.0), height=Inches(0.4),
                     font_size=9, color=GRIS_MED, align=PP_ALIGN.RIGHT)
    # Thin red separator
    add_rect(slide, left=Inches(0.25), top=Inches(1.05),
             width=Inches(12.85), height=Pt(1.2), fill_color=ROJO_GAMA)
    add_footer_bar(slide)


def color_series(series, rgb_color):
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = rgb_color


def add_native_chart(slide, chart_type, chart_data, left, top, width, height,
                     has_legend=True, legend_position=None):
    from pptx.enum.chart import XL_LEGEND_POSITION
    graphic_frame = slide.shapes.add_chart(
        chart_type, left, top, width, height, chart_data
    )
    chart = graphic_frame.chart
    chart.has_legend = has_legend
    if has_legend and legend_position:
        chart.legend.position = legend_position
        chart.legend.include_in_layout = False
    chart.has_title = False
    return chart


# ---------------------------------------------------------------------------
# PRESENTATION INIT
# ---------------------------------------------------------------------------
prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H


# ============================================================================
# RE01 — Portada
# ============================================================================
slide = prs.slides.add_slide(get_blank_layout(prs))
fill_background(slide, NEGRO_TEXTO)

# Red diagonal accent block (top-right)
add_rect(slide, left=Inches(9.2), top=Inches(0),
         width=Inches(4.133), height=Inches(7.5), fill_color=ROJO_GAMA)

add_text_box(slide, "RESUMEN EJECUTIVO",
             left=Inches(0.5), top=Inches(1.8),
             width=Inches(8), height=Inches(0.55),
             font_size=13, bold=True, color=ROJO_GAMA)

add_text_box(slide, "Notoriedad Gama 2026",
             left=Inches(0.5), top=Inches(2.4),
             width=Inches(8.4), height=Inches(1.1),
             font_size=38, bold=True, color=BLANCO)

add_text_box(slide, "Estudio de Brand Health · V8.1 — Mayo 2026",
             left=Inches(0.5), top=Inches(3.6),
             width=Inches(8.4), height=Inches(0.55),
             font_size=16, color=GRIS_CLARO)

add_text_box(slide, "Cliente: Gama  |  Elaborado por: Cora Urrea / equipo análisis",
             left=Inches(0.5), top=Inches(4.4),
             width=Inches(8.4), height=Inches(0.5),
             font_size=13, color=GRIS_CLARO)

add_text_box(slide, "Charts editables — V8.1 (charts nativos PowerPoint)",
             left=Inches(0.5), top=Inches(5.1),
             width=Inches(8.4), height=Inches(0.4),
             font_size=11, color=GRIS_MED, italic=True)

add_text_box(slide, "Confidencial — uso interno Gama",
             left=Inches(0.5), top=Inches(6.8),
             width=Inches(8.4), height=Inches(0.4),
             font_size=10, color=GRIS_MED, italic=True)


# ============================================================================
# RE02 — La pregunta central
# ============================================================================
slide = prs.slides.add_slide(get_blank_layout(prs))
fill_background(slide, BLANCO)
slide_header(slide,
             "¿Qué cambió en el mercado y qué debe hacer Gama en 2026-2027?",
             tag_text="[CONTEXTO]", title_size=21)

bullets = [
    "Base: n=402 entrevistas, 4 zonas Venezuela, ±4.89% error muestral (95% confianza)",
    "Nueva ola V8.1: metodología TB puro + vecindad perceptual + drivers reformulados por imagen",
    "Este resumen: 7 hallazgos centrales → 3 decisiones estratégicas para la Junta",
]
add_bullet_box(slide, bullets,
               left=Inches(0.5), top=Inches(1.3),
               width=Inches(12.5), height=Inches(2.2),
               font_size=17)

# Context box
add_rect(slide, left=Inches(0.5), top=Inches(3.8),
         width=Inches(12.3), height=Inches(2.9),
         fill_color=RGBColor(0xF7, 0xF7, 0xF7),
         line_color=GRIS_CLARO, line_width=Pt(1))

add_text_box(slide,
             "El mercado venezolano de supermercados enfrenta presión de precio creciente. "
             "Gama mantiene un posicionamiento experiencial único — outlier estructural en un sector "
             "dominado por precio. V8.1 revela tanto el activo oculto (Atención) como el riesgo "
             "inmediato (Plazas en el mismo vecindario perceptual). Las decisiones deben tomarse en 2026.",
             left=Inches(0.75), top=Inches(4.0),
             width=Inches(11.8), height=Inches(2.5),
             font_size=14, color=NEGRO_TEXTO)


# ============================================================================
# RE03 — 7 hallazgos centrales (matriz 4×2 compacta, sin chart)
# ============================================================================
slide = prs.slides.add_slide(get_blank_layout(prs))
fill_background(slide, BLANCO)
slide_header(slide,
             "7 hallazgos que definen la estrategia Gama 2026-2027",
             tag_text="[SÍNTESIS]", title_size=21)

hallazgos = [
    ("1", "DNA experiencial consolidado",
     "6 atributos sobre el promedio vs 4 en V3; precio subíndice sostenido"),
    ("2", "NSE C+/C ve Gama más fuerte",
     "El segmento de mayor valor percibe el DNA experiencial con más intensidad"),
    ("3", "Rio sigue precio-dominante",
     "No migró a territorio experiencial — pero Calidad sombra es amenaza latente"),
    ("4", "Recall llega al segmento equivocado",
     "NSE E sobre-representado; C+/C objetivo sub-representado"),
    ("5", "Categorías 'ofrecer valor' identificadas",
     "Congelados, Gaseosas, Salsas — Gama ya competitiva en precio"),
    ("6", "TB puro: presión precio más severa",
     "[NUEVO V8] El mercado exige precio con más rigidez de lo que T2B mostraba"),
    ("7", "ATENCIÓN = atributo SOMBRA +62.5 pp",
     "[NUEVO V8] Leales ven Atención en Gama; el mercado amplio no sabe aún"),
]

col_w = Inches(6.1)
row_h = Inches(0.82)
x_offsets = [Inches(0.35), Inches(6.7)]
y_start = Inches(1.25)

for i, (num, titulo, desc) in enumerate(hallazgos):
    col = i % 2
    row = i // 2
    x = x_offsets[col]
    y = y_start + row * row_h

    # Number pill
    add_rect(slide, left=x, top=y + Inches(0.1),
             width=Inches(0.38), height=Inches(0.42), fill_color=ROJO_GAMA)
    add_text_box(slide, num,
                 left=x, top=y + Inches(0.08),
                 width=Inches(0.38), height=Inches(0.42),
                 font_size=13, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)

    # Title + desc
    add_text_box(slide, titulo,
                 left=x + Inches(0.45), top=y,
                 width=col_w - Inches(0.5), height=Inches(0.3),
                 font_size=12, bold=True, color=NEGRO_TEXTO)
    add_text_box(slide, desc,
                 left=x + Inches(0.45), top=y + Inches(0.32),
                 width=col_w - Inches(0.5), height=Inches(0.45),
                 font_size=10, color=GRIS_MED)

    # Horizontal separator (between rows)
    if row < 3:
        add_rect(slide,
                 left=x, top=y + row_h - Inches(0.04),
                 width=col_w - Inches(0.1), height=Pt(0.7),
                 fill_color=GRIS_CLARO)

add_footnote(slide,
             "REF: cifras Pref-Gama basadas en n=32 — tendencias direccionales, no proyectables al universo.")


# ============================================================================
# RE04 — Hallazgo #1: DNA outlier — C03 NATIVE COLUMN CLUSTERED
# ============================================================================
slide = prs.slides.add_slide(get_blank_layout(prs))
fill_background(slide, BLANCO)
slide_header(slide,
             "Gama es la ÚNICA cadena no dominada por precio — el DNA experiencial es el activo",
             tag_text="[HALLAZGO 1]", title_size=19)

# ── C03: Modelo mental precio — COLUMN_CLUSTERED nativo ──────────────────
chart_data_c03 = CategoryChartData()
chart_data_c03.categories = CADENAS
chart_data_c03.add_series("% Asociación Precio",
    [40.6, 67.3, 72.1, 58.4, 75.8, 81.2, 63.9, 84.1])

chart_c03 = add_native_chart(
    slide, XL_CHART_TYPE.COLUMN_CLUSTERED, chart_data_c03,
    left=Inches(0.3), top=Inches(1.2),
    width=Inches(7.8), height=Inches(5.1),
    has_legend=False
)

# Color individual columns: Gama = red, rest = competitor palette
series_c03 = chart_c03.series[0]
for pt_idx, pt in enumerate(series_c03.points):
    pt.format.fill.solid()
    pt.format.fill.fore_color.rgb = CADENA_COLORS[pt_idx]

# Bullets right column
bullets = [
    "7 de 8 cadenas: 51–84% asociadas a precio como atributo dominante",
    "Gama: 40.6% — outlier estructural, no accidental",
    "DNA experiencial: 6 atributos sobreíndicen vs mercado (z-scores V8)",
]
add_bullet_box(slide, bullets,
               left=Inches(8.3), top=Inches(1.5),
               width=Inches(4.8), height=Inches(3.0),
               font_size=14)

# Implication box
add_rect(slide, left=Inches(8.3), top=Inches(4.7),
         width=Inches(4.8), height=Inches(1.8),
         fill_color=RGBColor(0xFC, 0xEC, 0xEC),
         line_color=ROJO_GAMA, line_width=Pt(1.5))
add_text_box(slide,
             "Erosionar el DNA con estrategia de precio directo destruiría "
             "el único diferenciador estructural de Gama.",
             left=Inches(8.5), top=Inches(4.85),
             width=Inches(4.5), height=Inches(1.5),
             font_size=12, color=NEGRO_TEXTO)


# ============================================================================
# RE05 — Hallazgo #2: Atención sombra — C08 NATIVE BAR_CLUSTERED horizontal
# ============================================================================
slide = prs.slides.add_slide(get_blank_layout(prs))
fill_background(slide, BLANCO)
slide_header(slide,
             "ATENCIÓN +62.5 pp entre leales vs mercado — activo invisible que debe activarse",
             tag_text="[HALLAZGO 2]", title_size=19)

# ── C08: Tornado brechas P23 Gama — BAR_CLUSTERED nativo ─────────────────
chart_data_c08 = CategoryChartData()
chart_data_c08.categories = ATRIBUTOS_P22_SHORT
VALORES_C08 = [-4.8, 8.2, 11.3, 62.5, 28.7, 14.6, 12.3, -9.2, 6.8, 19.4]
chart_data_c08.add_series("Brecha pp (Pref-Gama vs Total)", VALORES_C08)

chart_c08 = add_native_chart(
    slide, XL_CHART_TYPE.BAR_CLUSTERED, chart_data_c08,
    left=Inches(0.3), top=Inches(1.2),
    width=Inches(7.8), height=Inches(5.1),
    has_legend=False
)

# Color: Atención (index 3) = red; positive = verde; negative = ambar
series_c08 = chart_c08.series[0]
for pt_idx, pt in enumerate(series_c08.points):
    pt.format.fill.solid()
    if ATRIBUTOS_P22_SHORT[pt_idx] == "Atención":
        pt.format.fill.fore_color.rgb = ROJO_GAMA
    elif VALORES_C08[pt_idx] >= 0:
        pt.format.fill.fore_color.rgb = VERDE
    else:
        pt.format.fill.fore_color.rgb = AMBAR

# Bullets right column
bullets = [
    "Compradores de Gama asocian Atención 62.5 pp más que el mercado amplio",
    "No-clientes no perciben la Atención de Gama — imagen no alcanzó al mercado",
    "Tarea de campaña: trasladar imagen validada por leales al mercado general",
]
add_bullet_box(slide, bullets,
               left=Inches(8.3), top=Inches(1.5),
               width=Inches(4.8), height=Inches(3.0),
               font_size=14)

# Implication box
add_rect(slide, left=Inches(8.3), top=Inches(4.7),
         width=Inches(4.8), height=Inches(1.5),
         fill_color=RGBColor(0xFD, 0xF6, 0xE3),
         line_color=AMBAR, line_width=Pt(1.5))
add_text_box(slide,
             "No crear un atributo nuevo — amplificar uno que ya existe "
             "y está validado por los leales.",
             left=Inches(8.5), top=Inches(4.85),
             width=Inches(4.5), height=Inches(1.3),
             font_size=12, color=NEGRO_TEXTO)

add_footnote(slide, "REF: n Pref-Gama = 32. Tendencia validada; cifra no proyectable al universo.")


# ============================================================================
# RE06 — Hallazgo #3: Vecindad perceptual — C09 NATIVE XY_SCATTER
# ============================================================================
slide = prs.slides.add_slide(get_blank_layout(prs))
fill_background(slide, BLANCO)
slide_header(slide,
             "Plazas es el rival perceptual más cercano — si activa Atención antes, captura el territorio",
             tag_text="[HALLAZGO 3]", title_size=19)

# ── C09: Scatter vecindad perceptual — XY_SCATTER nativo ─────────────────
chart_data_c09 = XyChartData()

PUNTOS_C09 = [
    ("Gama",            68.9, 40.6, ROJO_GAMA),
    ("Plazas",          64.2, 48.3, COLOR_PLAZAS),
    ("Rio",             53.8, 52.8, COLOR_RIO),
    ("Páramo",          46.3, 67.3, COLOR_PARAMO),
    ("Plan Suárez",     41.7, 74.2, COLOR_PSUAREZ),
    ("Cen. Madeirense", 38.2, 78.6, COLOR_CENTRAL),
    ("Forum",           35.4, 81.9, COLOR_FORUM),
    ("Luz Marina",      32.9, 84.7, COLOR_LUZ),
]

for cadena, x, y, color in PUNTOS_C09:
    series = chart_data_c09.add_series(cadena)
    series.add_data_point(x, y)

chart_c09 = add_native_chart(
    slide, XL_CHART_TYPE.XY_SCATTER, chart_data_c09,
    left=Inches(0.3), top=Inches(1.2),
    width=Inches(7.8), height=Inches(5.1),
    has_legend=True
)

for i, (cadena, x, y, color) in enumerate(PUNTOS_C09):
    color_series(chart_c09.series[i], color)

# Axis labels below chart
add_text_box(slide, "← Menos experiencial | Más experiencial →",
             left=Inches(0.5), top=Inches(6.4), width=Inches(7.5), height=Inches(0.3),
             font_size=9, color=GRIS_MED, align=PP_ALIGN.CENTER, italic=True)

# Bullets right column
bullets = [
    "3 clusters: Gama-Plazas-Rio en zona experiencial / precio-volumen / débiles",
    "Plazas comparte el atributo sombra (Atención) — máxima urgencia de diferenciación",
    "Rio sobreíndice en Calidad entre leales: amenaza latente al consumidor experiencial",
]
add_bullet_box(slide, bullets,
               left=Inches(8.3), top=Inches(1.5),
               width=Inches(4.8), height=Inches(3.0),
               font_size=14)

# Implication box
add_rect(slide, left=Inches(8.3), top=Inches(4.7),
         width=Inches(4.8), height=Inches(1.5),
         fill_color=RGBColor(0xEC, 0xF7, 0xEE),
         line_color=VERDE, line_width=Pt(1.5))
add_text_box(slide,
             "Monitorear Rio y Plazas en ola 2027 — incluir tracking de "
             "activaciones de comunicación competidores.",
             left=Inches(8.5), top=Inches(4.85),
             width=Inches(4.5), height=Inches(1.3),
             font_size=12, color=NEGRO_TEXTO)


# ============================================================================
# RE07 — Hallazgo #4: Drivers — demostrar no declarar
# ============================================================================
slide = prs.slides.add_slide(get_blank_layout(prs))
fill_background(slide, BLANCO)
slide_header(slide,
             "La percepción predice preferencia — la campaña debe DEMOSTRAR, no declarar",
             tag_text="[HALLAZGO 4]", title_size=20)

col_left  = Inches(0.4)
col_right = Inches(6.8)
col_w2    = Inches(5.9)
top_cols  = Inches(1.3)

# Left — Anti-mensaje
add_rect(slide, left=col_left, top=top_cols,
         width=col_w2, height=Inches(4.8),
         fill_color=RGBColor(0xFD, 0xF0, 0xF0),
         line_color=ROJO_GAMA, line_width=Pt(1.2))
add_text_box(slide, "ANTI-MENSAJE",
             left=col_left + Inches(0.2), top=top_cols + Inches(0.1),
             width=col_w2 - Inches(0.4), height=Inches(0.4),
             font_size=13, bold=True, color=ROJO_GAMA)
add_text_box(slide, '"En Gama te atendemos bien"',
             left=col_left + Inches(0.2), top=top_cols + Inches(0.55),
             width=col_w2 - Inches(0.4), height=Inches(0.5),
             font_size=16, bold=True, color=NEGRO_TEXTO, italic=True)
anti_bullets = [
    "Cualquier cadena puede hacer el mismo claim",
    "Importancias declaradas (P22) sufren techo — no discriminan",
    "No diferencia de Plazas — mismo atributo, mismo mensaje genérico",
    "Resultado: inversión sin diferenciación perceptual",
]
add_bullet_box(slide, anti_bullets,
               left=col_left + Inches(0.2), top=top_cols + Inches(1.15),
               width=col_w2 - Inches(0.4), height=Inches(3.3),
               font_size=13, color=NEGRO_TEXTO)

# Right — Mensaje correcto
add_rect(slide, left=col_right, top=top_cols,
         width=col_w2, height=Inches(4.8),
         fill_color=RGBColor(0xEC, 0xF7, 0xEE),
         line_color=VERDE, line_width=Pt(1.2))
add_text_box(slide, "MENSAJE CORRECTO",
             left=col_right + Inches(0.2), top=top_cols + Inches(0.1),
             width=col_w2 - Inches(0.4), height=Inches(0.4),
             font_size=13, bold=True, color=VERDE)
add_text_box(slide, "Demostrar la Atención que ya existe",
             left=col_right + Inches(0.2), top=top_cols + Inches(0.55),
             width=col_w2 - Inches(0.4), height=Inches(0.5),
             font_size=16, bold=True, color=NEGRO_TEXTO)
ok_bullets = [
    "Testimonios reales de clientes leales — historias de Atención",
    "Momentos de verdad específicos: nombre, rapidez, resolución",
    "Comparativas experienciales vs cadenas precio-volumen",
    "Mecanismo: imagen P23 → preferencia (modelo logístico validado)",
]
add_bullet_box(slide, ok_bullets,
               left=col_right + Inches(0.2), top=top_cols + Inches(1.15),
               width=col_w2 - Inches(0.4), height=Inches(3.3),
               font_size=13, color=NEGRO_TEXTO)


# ============================================================================
# RE08 — 3 decisiones estratégicas para la Junta
# ============================================================================
slide = prs.slides.add_slide(get_blank_layout(prs))
fill_background(slide, BLANCO)
slide_header(slide,
             "3 decisiones estratégicas para Junta Gama — 2026-2027",
             tag_text="[CALL TO ACTION]", title_size=21)

decisiones = [
    (
        "1",
        "DEFENDER DNA + abordar precio creativamente",
        [
            "5 vías identificadas (Cap 5.6): beneficio simbólico, coleccionables, anzuelo precio por categoría",
            "No competencia directa en precio — erosiona el único outlier estructural del mercado",
            "Categorías 'ofrecer valor': Congelados, Gaseosas, Salsas — palanca táctica sin riesgo de DNA",
        ],
        ROJO_GAMA,
    ),
    (
        "2",
        "ACTIVAR Atención en campaña 2026-2027",
        [
            "Demostrar (testimonios, momentos de verdad) — no declarar ('te atendemos bien')",
            "Trasladar imagen validada por leales (+62.5 pp) al mercado no-cliente",
            "Rediseñar mix de medios hacia NSE C+/C — canal actual sobre-alcanza NSE E",
        ],
        AMBAR,
    ),
    (
        "3",
        "DIFERENCIAR de Plazas + vigilar Rio en Calidad",
        [
            "Plazas: mismo vecindario perceptual, mismo atributo sombra — urgencia máxima",
            "Rio: si activa campaña de Calidad, puede capturar consumidores experienciales no fidelizados",
            "Van Westendorp PSM en ola 2027 — prerrequisito antes de decisiones de precio táctico",
        ],
        VERDE,
    ),
]

dec_top = Inches(1.3)
dec_h   = Inches(1.9)
dec_gap = Inches(0.15)

for i, (num, titulo, bullets, accent) in enumerate(decisiones):
    y = dec_top + i * (dec_h + dec_gap)

    add_rect(slide, left=Inches(0.35), top=y,
             width=Inches(0.08), height=dec_h,
             fill_color=accent)
    add_text_box(slide, num,
                 left=Inches(0.5), top=y + Inches(0.05),
                 width=Inches(0.5), height=Inches(0.5),
                 font_size=24, bold=True, color=accent, align=PP_ALIGN.CENTER)
    add_text_box(slide, titulo,
                 left=Inches(1.1), top=y + Inches(0.08),
                 width=Inches(11.7), height=Inches(0.42),
                 font_size=16, bold=True, color=NEGRO_TEXTO)
    add_bullet_box(slide, bullets,
                   left=Inches(1.1), top=y + Inches(0.55),
                   width=Inches(11.7), height=Inches(1.25),
                   font_size=12, color=GRIS_MED, marker="→  ")


# ============================================================================
# RE09 — Cierre + próximos pasos
# ============================================================================
slide = prs.slides.add_slide(get_blank_layout(prs))
fill_background(slide, NEGRO_TEXTO)

add_rect(slide, left=Inches(0), top=Inches(0),
         width=Inches(0.15), height=SLIDE_H, fill_color=ROJO_GAMA)

add_text_box(slide, "Próximos pasos",
             left=Inches(0.4), top=Inches(0.8),
             width=Inches(12.5), height=Inches(0.55),
             font_size=15, bold=True, color=ROJO_GAMA)

add_text_box(slide, "Prerrequisito metodológico: Van Westendorp PSM en ola 2027",
             left=Inches(0.4), top=Inches(1.5),
             width=Inches(12.5), height=Inches(0.55),
             font_size=22, bold=True, color=BLANCO)

bullets_cierre = [
    "Van Westendorp PSM: calibrar precio psicológico por segmento — antes de decisiones de precio táctico",
    "Incluir tracking Rio y Plazas en ola 2027 (amenaza Calidad / activación Atención competidores)",
    "Ampliar n Pref-Gama en ola 2027 (actual n=32 REF) para análisis proyectables por subgrupo",
]
add_bullet_box(slide, bullets_cierre,
               left=Inches(0.5), top=Inches(2.3),
               width=Inches(12.3), height=Inches(2.2),
               font_size=16, color=GRIS_CLARO, marker="→  ")

add_rect(slide, left=Inches(0.4), top=Inches(4.7),
         width=Inches(12.5), height=Pt(1.2), fill_color=GRIS_MED)

add_text_box(slide, "Contacto",
             left=Inches(0.4), top=Inches(5.0),
             width=Inches(4), height=Inches(0.4),
             font_size=11, color=GRIS_MED, bold=True)
add_text_box(slide, "Cora Urrea",
             left=Inches(0.4), top=Inches(5.45),
             width=Inches(8), height=Inches(0.45),
             font_size=18, bold=True, color=BLANCO)
add_text_box(slide,
             "Entregado con base de datos V2.0 completa (n=403 × 295 variables) y metodología documentada",
             left=Inches(0.4), top=Inches(6.0),
             width=Inches(12.5), height=Inches(0.45),
             font_size=11, color=GRIS_MED, italic=True)

add_text_box(slide, "Confidencial — uso interno Gama",
             left=Inches(0.4), top=Inches(6.9),
             width=Inches(12.5), height=Inches(0.35),
             font_size=9, color=GRIS_MED, italic=True)


# ============================================================================
# SAVE
# ============================================================================
prs.save(str(OUTPUT_PATH))
print(f"OK: {OUTPUT_PATH}")
print(f"Slides: {len(prs.slides)}")
size = os.path.getsize(str(OUTPUT_PATH))
print(f"Bytes: {size:,}")
