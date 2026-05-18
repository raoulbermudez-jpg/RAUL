import sys
sys.stdout.reconfigure(encoding='utf-8')

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Paths ───────────────────────────────────────────────────────────────────
BASE = r"C:\rAUL\03-projects\consultoria-externa\gama-notoriedad-2026\03-deck\V8"
CHARTS = os.path.join(BASE, "charts")
OUTPUT = os.path.join(BASE, "2026-05-18_Notoriedad-Gama-2026_V8_Resumen-Ejecutivo.pptx")

# ── Brand palette ────────────────────────────────────────────────────────────
RED     = RGBColor(0xE3, 0x06, 0x13)   # Rojo Gama
DARK    = RGBColor(0x1A, 0x1A, 0x1A)   # Texto principal
GREY_M  = RGBColor(0x6B, 0x6B, 0x6B)   # Gris medio
GREY_L  = RGBColor(0xE5, 0xE5, 0xE5)   # Gris claro
AMBER   = RGBColor(0xF2, 0xA9, 0x00)   # Ámbar
GREEN   = RGBColor(0x2D, 0x8F, 0x47)   # Verde
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

# ── Slide dimensions (16:9) ──────────────────────────────────────────────────
W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

# ── Helper functions ─────────────────────────────────────────────────────────

def blank_slide(prs):
    """Add a completely blank slide."""
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def fill_bg(slide, color: RGBColor):
    """Fill slide background with a solid color."""
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0)):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_name="Calibri", font_size=Pt(14), bold=False,
             color=DARK, align=PP_ALIGN.LEFT, italic=False,
             word_wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_para(tf, text, font_name="Calibri", font_size=Pt(13), bold=False,
             color=DARK, align=PP_ALIGN.LEFT, space_before=Pt(4)):
    from pptx.util import Pt as Ptt
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    return p


def add_bullet_box(slide, bullets, left, top, width, height,
                   font_size=Pt(14), color=DARK, marker="•  "):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for b in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(5)
        run = p.add_run()
        run.text = marker + b
        run.font.name = "Calibri"
        run.font.size = font_size
        run.font.color.rgb = color
    return txBox


def add_footnote(slide, text):
    add_text(slide, text,
             left=Inches(0.4), top=Inches(6.9),
             width=Inches(12.5), height=Inches(0.45),
             font_size=Pt(9), color=GREY_M, italic=True)


def add_footer_bar(slide, label="Confidencial — uso interno Gama"):
    bar = add_rect(slide,
                   left=Inches(0), top=Inches(7.1),
                   width=W, height=Inches(0.4),
                   fill_color=GREY_L)
    add_text(slide, label,
             left=Inches(0.4), top=Inches(7.12),
             width=Inches(9), height=Inches(0.36),
             font_size=Pt(9), color=GREY_M)


def slide_header(slide, title_text, tag_text=None, title_size=Pt(22)):
    """Red accent bar + title + optional tag pill."""
    # Red left accent
    add_rect(slide, left=Inches(0), top=Inches(0),
             width=Inches(0.12), height=H, fill_color=RED)
    # Title
    add_text(slide, title_text,
             left=Inches(0.25), top=Inches(0.2),
             width=Inches(12.5), height=Inches(0.8),
             font_name="Calibri", font_size=title_size, bold=True, color=DARK)
    # Tag
    if tag_text:
        add_text(slide, tag_text,
                 left=Inches(11.2), top=Inches(0.22),
                 width=Inches(2.0), height=Inches(0.4),
                 font_size=Pt(9), color=GREY_M, align=PP_ALIGN.RIGHT)
    # Thin separator line
    add_rect(slide,
             left=Inches(0.25), top=Inches(1.05),
             width=Inches(12.85), height=Pt(1.2),
             fill_color=RED)
    add_footer_bar(slide)


def chart_path(name):
    return os.path.join(CHARTS, name)


# ════════════════════════════════════════════════════════════════════════════
# RE01 — Portada
# ════════════════════════════════════════════════════════════════════════════
slide = blank_slide(prs)
fill_bg(slide, DARK)

# Red diagonal accent block (top-right)
add_rect(slide,
         left=Inches(9.2), top=Inches(0),
         width=Inches(4.133), height=Inches(7.5),
         fill_color=RED)

# "RESUMEN EJECUTIVO" tag
add_text(slide, "RESUMEN EJECUTIVO",
         left=Inches(0.5), top=Inches(1.8),
         width=Inches(8), height=Inches(0.55),
         font_size=Pt(13), bold=True, color=RED, align=PP_ALIGN.LEFT)

# Main title
add_text(slide, "Notoriedad Gama 2026",
         left=Inches(0.5), top=Inches(2.4),
         width=Inches(8.4), height=Inches(1.1),
         font_name="Calibri", font_size=Pt(38), bold=True, color=WHITE)

# Subtitle
add_text(slide, "Estudio de Brand Health · V8 — Mayo 2026",
         left=Inches(0.5), top=Inches(3.6),
         width=Inches(8.4), height=Inches(0.55),
         font_size=Pt(16), color=GREY_L)

# Client + confidential
add_text(slide, "Cliente: Gama  |  Elaborado por: Cora Urrea / equipo análisis",
         left=Inches(0.5), top=Inches(4.4),
         width=Inches(8.4), height=Inches(0.5),
         font_size=Pt(13), color=GREY_L)

add_text(slide, "Confidencial — uso interno Gama",
         left=Inches(0.5), top=Inches(6.8),
         width=Inches(8.4), height=Inches(0.4),
         font_size=Pt(10), color=GREY_M, italic=True)

# ════════════════════════════════════════════════════════════════════════════
# RE02 — La pregunta central
# ════════════════════════════════════════════════════════════════════════════
slide = blank_slide(prs)
fill_bg(slide, WHITE)
slide_header(slide,
             "¿Qué cambió en el mercado y qué debe hacer Gama en 2026-2027?",
             tag_text="[CONTEXTO]", title_size=Pt(21))

bullets = [
    "Base: n=402 entrevistas, 4 zonas Venezuela, ±4.89% error muestral (95% confianza)",
    "Nueva ola V8: metodología TB puro + vecindad perceptual + drivers reformulados por imagen",
    "Este resumen: 7 hallazgos centrales → 3 decisiones estratégicas para la Junta",
]
add_bullet_box(slide, bullets,
               left=Inches(0.5), top=Inches(1.3),
               width=Inches(12.5), height=Inches(2.2),
               font_size=Pt(17))

# Context box
box = add_rect(slide,
               left=Inches(0.5), top=Inches(3.8),
               width=Inches(12.3), height=Inches(2.9),
               fill_color=RGBColor(0xF7, 0xF7, 0xF7),
               line_color=GREY_L, line_width=Pt(1))

add_text(slide, "El mercado venezolano de supermercados enfrenta presión de precio creciente. "
         "Gama mantiene un posicionamiento experiencial único — outlier estructural en un sector "
         "dominado por precio. V8 revela tanto el activo oculto (Atención) como el riesgo inmediato "
         "(Plazas en el mismo vecindario perceptual). Las decisiones deben tomarse en 2026.",
         left=Inches(0.75), top=Inches(4.0),
         width=Inches(11.8), height=Inches(2.5),
         font_size=Pt(14), color=DARK)

# ════════════════════════════════════════════════════════════════════════════
# RE03 — 7 hallazgos centrales (matriz compacta)
# ════════════════════════════════════════════════════════════════════════════
slide = blank_slide(prs)
fill_bg(slide, WHITE)
slide_header(slide,
             "7 hallazgos que definen la estrategia Gama 2026-2027",
             tag_text="[SÍNTESIS]", title_size=Pt(21))

hallazgos = [
    ("1", "DNA experiencial consolidado", "6 atributos sobre el promedio vs 4 en V3; precio subíndice sostenido"),
    ("2", "NSE C+/C ve Gama más fuerte", "El segmento de mayor valor percibe el DNA experiencial con más intensidad"),
    ("3", "Rio sigue precio-dominante", "No migró a territorio experiencial — pero Calidad sombra es amenaza latente"),
    ("4", "Recall llega al segmento equivocado", "NSE E sobre-representado; C+/C objetivo sub-representado"),
    ("5", "Categorías 'ofrecer valor' identificadas", "Congelados, Gaseosas, Salsas — Gama ya competitiva en precio"),
    ("6", "TB puro: presión precio más severa", "[NUEVO V8] El mercado exige precio con más rigidez de lo que T2B mostraba"),
    ("7", "ATENCIÓN = atributo SOMBRA +62.5 pp", "[NUEVO V8] Leales ven Atención en Gama; el mercado amplio no sabe aún"),
]

col_w = Inches(6.1)
row_h = Inches(0.82)
x_offsets = [Inches(0.35), Inches(6.7)]
y_start = Inches(1.25)

for i, (num, titulo, desc) in enumerate(hallazgos):
    col = i % 2
    row = i // 2
    x = x_offsets[col]
    y = y_start + row * row_h

    # Number pill
    pill = add_rect(slide, left=x, top=y + Inches(0.1),
                    width=Inches(0.38), height=Inches(0.42),
                    fill_color=RED)
    add_text(slide, num,
             left=x, top=y + Inches(0.08),
             width=Inches(0.38), height=Inches(0.42),
             font_size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Titulo + desc
    add_text(slide, titulo,
             left=x + Inches(0.45), top=y,
             width=col_w - Inches(0.5), height=Inches(0.3),
             font_size=Pt(12), bold=True, color=DARK)
    add_text(slide, desc,
             left=x + Inches(0.45), top=y + Inches(0.32),
             width=col_w - Inches(0.5), height=Inches(0.45),
             font_size=Pt(10), color=GREY_M)

    # Separador horizontal
    if row < 3:
        add_rect(slide,
                 left=x, top=y + row_h - Inches(0.04),
                 width=col_w - Inches(0.1), height=Pt(0.7),
                 fill_color=GREY_L)

add_footnote(slide, "REF: cifras Pref-Gama basadas en n=32 — tendencias direccionales, no proyectables al universo.")

# ════════════════════════════════════════════════════════════════════════════
# RE04 — Hallazgo #1: DNA outlier (con chart C03)
# ════════════════════════════════════════════════════════════════════════════
slide = blank_slide(prs)
fill_bg(slide, WHITE)
slide_header(slide,
             "Gama es la ÚNICA cadena no dominada por precio — el DNA experiencial es el activo",
             tag_text="[HALLAZGO 1]", title_size=Pt(19))

# Chart
img_path = chart_path("C03_modelo_mental_precio.png")
slide.shapes.add_picture(img_path,
                         left=Inches(0.3), top=Inches(1.2),
                         width=Inches(7.8), height=Inches(5.1))

# Bullets (right column)
bullets = [
    "7 de 8 cadenas: 51–84% asociadas a precio como atributo dominante",
    "Gama: 40.6% — outlier estructural, no accidental",
    "DNA experiencial: 6 atributos sobreíndicen vs mercado (z-scores V8)",
]
add_bullet_box(slide, bullets,
               left=Inches(8.3), top=Inches(1.5),
               width=Inches(4.8), height=Inches(3.0),
               font_size=Pt(14))

# Implication box
impl = add_rect(slide,
                left=Inches(8.3), top=Inches(4.7),
                width=Inches(4.8), height=Inches(1.8),
                fill_color=RGBColor(0xFC, 0xEC, 0xEC),
                line_color=RED, line_width=Pt(1.5))
add_text(slide,
         "Erosionar el DNA con estrategia de precio directo destruiría el único diferenciador estructural de Gama.",
         left=Inches(8.5), top=Inches(4.85),
         width=Inches(4.5), height=Inches(1.5),
         font_size=Pt(12), bold=False, color=DARK)

# ════════════════════════════════════════════════════════════════════════════
# RE05 — Hallazgo #2: Atención sombra (con chart C08)
# ════════════════════════════════════════════════════════════════════════════
slide = blank_slide(prs)
fill_bg(slide, WHITE)
slide_header(slide,
             "ATENCIÓN +62.5 pp entre leales vs mercado — activo invisible que debe activarse",
             tag_text="[HALLAZGO 2]", title_size=Pt(19))

img_path = chart_path("C08_p23_brechas_pref_vs_total.png")
slide.shapes.add_picture(img_path,
                         left=Inches(0.3), top=Inches(1.2),
                         width=Inches(7.8), height=Inches(5.1))

bullets = [
    "Compradores de Gama asocian Atención 62.5 pp más que el mercado amplio",
    "No-clientes no perciben la Atención de Gama — imagen no alcanzó al mercado",
    "Tarea de campaña: trasladar imagen validada por leales al mercado general",
]
add_bullet_box(slide, bullets,
               left=Inches(8.3), top=Inches(1.5),
               width=Inches(4.8), height=Inches(3.0),
               font_size=Pt(14))

impl = add_rect(slide,
                left=Inches(8.3), top=Inches(4.7),
                width=Inches(4.8), height=Inches(1.5),
                fill_color=RGBColor(0xFD, 0xF6, 0xE3),
                line_color=AMBER, line_width=Pt(1.5))
add_text(slide,
         "No crear un atributo nuevo — amplificar uno que ya existe y está validado por los leales.",
         left=Inches(8.5), top=Inches(4.85),
         width=Inches(4.5), height=Inches(1.3),
         font_size=Pt(12), color=DARK)

add_footnote(slide, "REF: n Pref-Gama = 32. Tendencia validada; cifra no proyectable al universo.")

# ════════════════════════════════════════════════════════════════════════════
# RE06 — Hallazgo #3: Vecindad perceptual + amenaza (con chart C09)
# ════════════════════════════════════════════════════════════════════════════
slide = blank_slide(prs)
fill_bg(slide, WHITE)
slide_header(slide,
             "Plazas es el rival perceptual más cercano — si activa Atención antes, captura el territorio",
             tag_text="[HALLAZGO 3]", title_size=Pt(19))

img_path = chart_path("C09_vecindad_perceptual.png")
slide.shapes.add_picture(img_path,
                         left=Inches(0.3), top=Inches(1.2),
                         width=Inches(7.8), height=Inches(5.1))

bullets = [
    "3 clusters: Gama-Plazas-Rio en zona experiencial / precio-volumen / débiles",
    "Plazas comparte el atributo sombra (Atención) — máxima urgencia de diferenciación",
    "Rio sobreíndice en Calidad entre leales: amenaza latente al consumidor experiencial",
]
add_bullet_box(slide, bullets,
               left=Inches(8.3), top=Inches(1.5),
               width=Inches(4.8), height=Inches(3.0),
               font_size=Pt(14))

impl = add_rect(slide,
                left=Inches(8.3), top=Inches(4.7),
                width=Inches(4.8), height=Inches(1.5),
                fill_color=RGBColor(0xEC, 0xF7, 0xEE),
                line_color=GREEN, line_width=Pt(1.5))
add_text(slide,
         "Monitorear Rio y Plazas en ola 2027 — incluir tracking de activaciones de comunicación competidores.",
         left=Inches(8.5), top=Inches(4.85),
         width=Inches(4.5), height=Inches(1.3),
         font_size=Pt(12), color=DARK)

# ════════════════════════════════════════════════════════════════════════════
# RE07 — Hallazgo #4: Drivers — demostrar no declarar
# ════════════════════════════════════════════════════════════════════════════
slide = blank_slide(prs)
fill_bg(slide, WHITE)
slide_header(slide,
             "La percepción predice preferencia — la campaña debe DEMOSTRAR, no declarar",
             tag_text="[HALLAZGO 4]", title_size=Pt(20))

# Two-column layout: Anti-mensaje vs Mensaje correcto
col_left  = Inches(0.4)
col_right = Inches(6.8)
col_w2    = Inches(5.9)
top_cols  = Inches(1.3)

# Left column — Anti-mensaje
add_rect(slide, left=col_left, top=top_cols,
         width=col_w2, height=Inches(4.8),
         fill_color=RGBColor(0xFD, 0xF0, 0xF0),
         line_color=RED, line_width=Pt(1.2))
add_text(slide, "ANTI-MENSAJE",
         left=col_left + Inches(0.2), top=top_cols + Inches(0.1),
         width=col_w2 - Inches(0.4), height=Inches(0.4),
         font_size=Pt(13), bold=True, color=RED)
add_text(slide, '"En Gama te atendemos bien"',
         left=col_left + Inches(0.2), top=top_cols + Inches(0.55),
         width=col_w2 - Inches(0.4), height=Inches(0.5),
         font_size=Pt(16), bold=True, color=DARK, italic=True)
anti_bullets = [
    "Cualquier cadena puede hacer el mismo claim",
    "Importancias declaradas (P22) sufren techo — no discriminan",
    "No diferencia de Plazas — mismo atributo, mismo mensaje genérico",
    "Resultado: inversión sin diferenciación perceptual",
]
add_bullet_box(slide, anti_bullets,
               left=col_left + Inches(0.2), top=top_cols + Inches(1.15),
               width=col_w2 - Inches(0.4), height=Inches(3.3),
               font_size=Pt(13), color=DARK)

# Right column — Mensaje correcto
add_rect(slide, left=col_right, top=top_cols,
         width=col_w2, height=Inches(4.8),
         fill_color=RGBColor(0xEC, 0xF7, 0xEE),
         line_color=GREEN, line_width=Pt(1.2))
add_text(slide, "MENSAJE CORRECTO",
         left=col_right + Inches(0.2), top=top_cols + Inches(0.1),
         width=col_w2 - Inches(0.4), height=Inches(0.4),
         font_size=Pt(13), bold=True, color=GREEN)
add_text(slide, "Demostrar la Atención que ya existe",
         left=col_right + Inches(0.2), top=top_cols + Inches(0.55),
         width=col_w2 - Inches(0.4), height=Inches(0.5),
         font_size=Pt(16), bold=True, color=DARK)
ok_bullets = [
    "Testimonios reales de clientes leales — historias de Atención",
    "Momentos de verdad específicos: nombre, rapidez, resolución",
    "Comparativas experienciales vs cadenas precio-volumen",
    "Mecanismo: imagen P23 → preferencia (modelo logístico validado)",
]
add_bullet_box(slide, ok_bullets,
               left=col_right + Inches(0.2), top=top_cols + Inches(1.15),
               width=col_w2 - Inches(0.4), height=Inches(3.3),
               font_size=Pt(13), color=DARK)

# ════════════════════════════════════════════════════════════════════════════
# RE08 — 3 decisiones para la Junta
# ════════════════════════════════════════════════════════════════════════════
slide = blank_slide(prs)
fill_bg(slide, WHITE)
slide_header(slide,
             "3 decisiones estratégicas para Junta Gama — 2026-2027",
             tag_text="[CALL TO ACTION]", title_size=Pt(21))

decisiones = [
    (
        "1",
        "DEFENDER DNA + abordar precio creativamente",
        [
            "5 vías identificadas (Cap 5.6): beneficio simbólico, coleccionables, anzuelo precio por categoría",
            "No competencia directa en precio — erosiona el único outlier estructural del mercado",
            "Categorías 'ofrecer valor': Congelados, Gaseosas, Salsas — palanca táctica sin riesgo de DNA",
        ],
        RED,
    ),
    (
        "2",
        "ACTIVAR Atención en campaña 2026-2027",
        [
            "Demostrar (testimonios, momentos de verdad) — no declarar ('te atendemos bien')",
            "Trasladar imagen validada por leales (+62.5 pp) al mercado no-cliente",
            "Rediseñar mix de medios hacia NSE C+/C — canal actual sobre-alcanza NSE E",
        ],
        AMBER,
    ),
    (
        "3",
        "DIFERENCIAR de Plazas + vigilar Rio en Calidad",
        [
            "Plazas: mismo vecindario perceptual, mismo atributo sombra — urgencia máxima",
            "Rio: si activa campaña de Calidad, puede capturar consumidores experienciales no fidelizados",
            "Van Westendorp PSM en ola 2027 — prerrequisito antes de decisiones de precio táctico",
        ],
        GREEN,
    ),
]

dec_top = Inches(1.3)
dec_h   = Inches(1.9)
dec_gap = Inches(0.15)

for i, (num, titulo, bullets, accent) in enumerate(decisiones):
    y = dec_top + i * (dec_h + dec_gap)

    # Accent left bar
    add_rect(slide, left=Inches(0.35), top=y,
             width=Inches(0.08), height=dec_h,
             fill_color=accent)

    # Number
    add_text(slide, num,
             left=Inches(0.5), top=y + Inches(0.05),
             width=Inches(0.5), height=Inches(0.5),
             font_size=Pt(24), bold=True, color=accent, align=PP_ALIGN.CENTER)

    # Title
    add_text(slide, titulo,
             left=Inches(1.1), top=y + Inches(0.08),
             width=Inches(11.7), height=Inches(0.42),
             font_size=Pt(16), bold=True, color=DARK)

    # Bullets
    add_bullet_box(slide, bullets,
                   left=Inches(1.1), top=y + Inches(0.55),
                   width=Inches(11.7), height=Inches(1.25),
                   font_size=Pt(12), color=GREY_M, marker="→  ")

# ════════════════════════════════════════════════════════════════════════════
# RE09 — Cierre
# ════════════════════════════════════════════════════════════════════════════
slide = blank_slide(prs)
fill_bg(slide, DARK)

# Red left accent
add_rect(slide, left=Inches(0), top=Inches(0),
         width=Inches(0.15), height=H, fill_color=RED)

add_text(slide, "Próximos pasos",
         left=Inches(0.4), top=Inches(0.8),
         width=Inches(12.5), height=Inches(0.55),
         font_size=Pt(15), bold=True, color=RED)

add_text(slide, "Prerrequisito metodológico: Van Westendorp PSM en ola 2027",
         left=Inches(0.4), top=Inches(1.5),
         width=Inches(12.5), height=Inches(0.55),
         font_size=Pt(22), bold=True, color=WHITE)

bullets_cierre = [
    "Van Westendorp PSM: calibrar precio psicológico por segmento — antes de decisiones de precio táctico",
    "Incluir tracking Rio y Plazas en ola 2027 (amenaza Calidad / activación Atención competidores)",
    "Ampliar n Pref-Gama en ola 2027 (actual n=32 REF) para análisis proyectables por subgrupo",
]
add_bullet_box(slide, bullets_cierre,
               left=Inches(0.5), top=Inches(2.3),
               width=Inches(12.3), height=Inches(2.2),
               font_size=Pt(16), color=GREY_L, marker="→  ")

# Separator
add_rect(slide,
         left=Inches(0.4), top=Inches(4.7),
         width=Inches(12.5), height=Pt(1.2),
         fill_color=GREY_M)

add_text(slide, "Contacto",
         left=Inches(0.4), top=Inches(5.0),
         width=Inches(4), height=Inches(0.4),
         font_size=Pt(11), color=GREY_M, bold=True)
add_text(slide, "Cora Urrea",
         left=Inches(0.4), top=Inches(5.45),
         width=Inches(8), height=Inches(0.45),
         font_size=Pt(18), bold=True, color=WHITE)
add_text(slide, "Entregado con base de datos V2.0 completa (n=403 × 295 variables) y metodología documentada",
         left=Inches(0.4), top=Inches(6.0),
         width=Inches(12.5), height=Inches(0.45),
         font_size=Pt(11), color=GREY_M, italic=True)

add_text(slide, "Confidencial — uso interno Gama",
         left=Inches(0.4), top=Inches(6.9),
         width=Inches(12.5), height=Inches(0.35),
         font_size=Pt(9), color=GREY_M, italic=True)

# ════════════════════════════════════════════════════════════════════════════
# Save
# ════════════════════════════════════════════════════════════════════════════
prs.save(OUTPUT)
print(f"OK: {OUTPUT}")
print(f"Slides: {len(prs.slides)}")
import os
size = os.path.getsize(OUTPUT)
print(f"Bytes: {size:,}")
