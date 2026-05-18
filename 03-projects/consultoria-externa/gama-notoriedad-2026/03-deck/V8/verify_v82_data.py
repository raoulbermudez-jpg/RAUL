"""Verifica datos en V8.2 heatmaps contra JSON canonico."""
import sys
sys.stdout.reconfigure(encoding='utf-8')
import json
from pptx import Presentation

V82 = r"C:\rAUL\03-projects\consultoria-externa\gama-notoriedad-2026\03-deck\V8\2026-05-18_Notoriedad-Gama-2026_V8.2.pptx"
JSON_PATH = r"C:\rAUL\03-projects\consultoria-externa\gama-notoriedad-2026\03-deck\V8\chart_data_v82.json"

data = json.load(open(JSON_PATH, encoding='utf-8'))
pres = Presentation(V82)

print("="*80)
print("VERIFICACION DATOS V8.2 vs JSON CANONICO")
print("="*80)

# S11 → C05 heatmap TB puro
# Atributo Rapidez caja (idx 2) x Pref-Gama (idx 1) debe ser 84.4
expected_c05 = data['C05_heatmap_tb_puro_por_marca']
expected_atributos = expected_c05['atributos']
expected_segmentos = expected_c05['segmentos']
matrix_c05 = expected_c05['matrix_pct']
rapidez_idx = expected_atributos.index('Rapidez caja')
pref_gama_idx = expected_segmentos.index('Pref-Gama*')
expected_value = matrix_c05[rapidez_idx][pref_gama_idx]
print(f"\nC05 esperado: Rapidez x Pref-Gama = {expected_value}%")

# Buscar tabla en S11 (slide_idx=10)
s11 = pres.slides[10]
for shape in s11.shapes:
    if shape.has_table:
        t = shape.table
        rows = list(t.rows)
        cols = list(t.columns)
        # Encontrar columna Pref-Gama y fila Rapidez
        header_row = [c.text.strip() for c in rows[0].cells]
        rapidez_row_idx = None
        for r_idx, row in enumerate(rows):
            if 'rapidez' in row.cells[0].text.strip().lower():
                rapidez_row_idx = r_idx
                break
        if rapidez_row_idx:
            print(f"  Cabecera S11 tabla: {header_row}")
            print(f"  Fila Rapidez en S11 ({rapidez_row_idx}): {[c.text.strip() for c in rows[rapidez_row_idx].cells]}")
        break

# S15 → C07 heatmap mundo de marca P23
# Gama (idx 0) x Precio (idx 7 atributo "Precio") debe ser 7.2
expected_c07 = data['C07_heatmap_p23_mundo_marca']
cadenas_c07 = expected_c07['cadenas']
atributos_c07 = expected_c07['atributos']
matrix_c07 = expected_c07['matrix_pct']
gama_row = cadenas_c07.index('Gama')
precio_col = atributos_c07.index('Precio')
atencion_col = atributos_c07.index('Atención')
expected_gama_precio = matrix_c07[gama_row][precio_col]
expected_gama_atencion = matrix_c07[gama_row][atencion_col]
print(f"\nC07 esperado: Gama Precio = {expected_gama_precio}% (V8.1 incorrecto: 40.6)")
print(f"C07 esperado: Gama Atención = {expected_gama_atencion}% (V8.1 incorrecto: 62.4)")

s15 = pres.slides[14]
for shape in s15.shapes:
    if shape.has_table:
        t = shape.table
        rows = list(t.rows)
        header_row = [c.text.strip() for c in rows[0].cells]
        gama_row_idx = None
        for r_idx, row in enumerate(rows):
            if row.cells[0].text.strip() == 'Gama':
                gama_row_idx = r_idx
                break
        if gama_row_idx:
            gama_cells = [c.text.strip() for c in rows[gama_row_idx].cells]
            print(f"  Cabecera S15 tabla: {header_row}")
            print(f"  Fila Gama en S15: {gama_cells}")
        break

print("\n" + "="*80)
print("RESULTADO")
print("="*80)
print("Si los valores en las filas mostradas coinciden con los esperados del JSON,")
print("la corrección V8.1 → V8.2 está confirmada.")
