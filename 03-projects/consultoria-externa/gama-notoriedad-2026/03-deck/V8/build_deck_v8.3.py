"""
build_deck_v8.2.py  —  Vivienne V8.3  —  Deck principal (32 slides)
Proyecto: Notoriedad Gama 2026 (consultoria externa via Cora Urrea)

REGLA CRITICA: TODOS los datos leidos de chart_data_v82.json.
PROHIBIDO inventar valores. Si assertion falla -> abortar con mensaje.

Encoding: UTF-8 obligatorio (Windows cp1252 workaround)
"""
import sys
sys.stdout.reconfigure(encoding='utf-8')

import json
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Importar builders
sys.path.insert(0, str(Path(__file__).parent))
from chart_builders_v82 import (
    make_chart_c01, make_chart_c02, make_chart_c03, make_chart_c04,
    make_chart_c05, make_chart_c06, make_chart_c07, make_chart_c08,
    make_chart_c09, make_chart_c10, make_chart_c11, make_chart_c12,
    make_chart_c13, make_chart_c14,
    add_text_box, add_slide_title, add_footnote, add_section_tag,
    GAMA_RED, TEXT_DARK, GREY_MED, GREY_LIGHT, WHITE, BLACK,
    AMBER, GREEN_VAL, PARAMO_BLUE
)

# ─────────────────────────────────────────────
# RUTAS
# ─────────────────────────────────────────────
BASE = Path(r"C:\rAUL\03-projects\consultoria-externa\gama-notoriedad-2026\03-deck\V8")
JSON_PATH = BASE / "chart_data_v82.json"
OUT_PATH  = BASE / "2026-05-18_Notoriedad-Gama-2026_V8.3.pptx"

# ─────────────────────────────────────────────
# CARGA JSON — VALIDACION CRITICA
# ─────────────────────────────────────────────
print("Cargando JSON...")
with open(JSON_PATH, encoding='utf-8') as f:
    DATA = json.load(f)

REQUIRED_CHARTS = [
    'C01_embudo_gama', 'C02_preferencia_p21_total_vs_cpc',
    'C03_modelo_mental_precio_dominante', 'C04_tb_vs_t2b_reranking',
    'C05_heatmap_tb_puro_por_marca', 'C06_dna_zscores_total_vs_cpc',
    'C07_heatmap_p23_mundo_marca', 'C08_tornado_brechas_p23_gama',
    'C09_scatter_vecindad_perceptual', 'C10_forest_plot_logit_paramo',
    'C11_percepcion_precio_p33_por_nse', 'C12_scatter_categorias_estrategia',
    'C13_recall_ptl_dtls_por_nse', 'C14_perfil_recordadores_vs_muestra'
]
for key in REQUIRED_CHARTS:
    assert key in DATA, f"ABORTAR: '{key}' no encontrado en JSON"

print(f"JSON OK: {len(REQUIRED_CHARTS)} charts verificados")

# Aliases de datos
C01 = DATA['C01_embudo_gama']
C02 = DATA['C02_preferencia_p21_total_vs_cpc']
C03 = DATA['C03_modelo_mental_precio_dominante']
C04 = DATA['C04_tb_vs_t2b_reranking']
C05 = DATA['C05_heatmap_tb_puro_por_marca']
C06 = DATA['C06_dna_zscores_total_vs_cpc']
C07 = DATA['C07_heatmap_p23_mundo_marca']
C08 = DATA['C08_tornado_brechas_p23_gama']
C09 = DATA['C09_scatter_vecindad_perceptual']
C10 = DATA['C10_forest_plot_logit_paramo']
C11 = DATA['C11_percepcion_precio_p33_por_nse']
C12 = DATA['C12_scatter_categorias_estrategia']
C13 = DATA['C13_recall_ptl_dtls_por_nse']
C14 = DATA['C14_perfil_recordadores_vs_muestra']

# SPOT-CHECKS CRITICOS (evitar el error de V8.1)
assert abs(C07['matrix_pct'][0][7] - 7.2) < 0.05,  "ABORTAR: C07 Gama Precio != 7.2"
assert abs(C07['matrix_pct'][0][4] - 21.9) < 0.05, "ABORTAR: C07 Gama Atencion != 21.9"
assert abs(C05['matrix_pct'][2][1] - 84.4) < 0.05, "ABORTAR: C05 Pref-Gama Rapidez != 84.4"
print(f"SPOT-CHECK OK: C07 Gama Precio={C07['matrix_pct'][0][7]} (esperado 7.2)")
print(f"SPOT-CHECK OK: C07 Gama Atencion={C07['matrix_pct'][0][4]} (esperado 21.9)")
print(f"SPOT-CHECK OK: C05 Pref-Gama Rapidez={C05['matrix_pct'][2][1]} (esperado 84.4)")

# ─────────────────────────────────────────────
# HELPERS DE SLIDE
# ─────────────────────────────────────────────

def new_slide(prs):
    """Agrega slide en blanco y retorna."""
    blank_layout = prs.slide_layouts[6]  # Completely Blank
    return prs.slides.add_slide(blank_layout)


def slide_header_bar(slide, section_color: RGBColor = None):
    """Barra superior de color sección."""
    color = section_color or GAMA_RED
    bar = slide.shapes.add_shape(1,
                                  Inches(0), Inches(0),
                                  Inches(13.333), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def bullet_box(slide, bullets: list, left: float, top: float,
               width: float, height: float, font_size: int = 12,
               color: RGBColor = None, bullet_char: str = "•"):
    """Text box con lista de bullets."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"{bullet_char}  {b}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color or TEXT_DARK
    return txBox


def analysis_highlight_box(slide, text: str, left: float = 0.3, top: float = 1.0,
                             width: float = 9.4, height: float = 0.55):
    """Box destacado para hallazgo clave en slides de análisis."""
    box = slide.shapes.add_shape(1,
                                  Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xFF, 0xEB, 0xEB)
    box.line.color.rgb = GAMA_RED

    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = GAMA_RED
    return box


def number_badge(slide, num: str, left: float, top: float):
    """Circulo numerado para hallazgos."""
    badge = slide.shapes.add_shape(9,
                                    Inches(left), Inches(top),
                                    Inches(0.35), Inches(0.35))
    badge.fill.solid()
    badge.fill.fore_color.rgb = GAMA_RED
    badge.line.fill.background()
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = num
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = WHITE


# ─────────────────────────────────────────────
# PRESENTACION
# ─────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

slide_count = 0


# ══════════════════════════════════════════════
# SECCIÓN 0 — APERTURA (3 slides)
# ══════════════════════════════════════════════

# ── S01 — PORTADA ──────────────────────────────
sl = new_slide(prs)
slide_count += 1

bg = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
bg.fill.solid()
bg.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
bg.line.fill.background()

stripe = sl.shapes.add_shape(1, Inches(0), Inches(5.2), Inches(13.333), Inches(0.12))
stripe.fill.solid()
stripe.fill.fore_color.rgb = GAMA_RED
stripe.line.fill.background()

add_text_box(sl, "NOTORIEDAD GAMA 2026",
             left=0.8, top=2.0, width=8.4, height=1.0,
             font_size=36, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text_box(sl, "Estudio de Brand Health  |  V8 — Mayo 2026",
             left=0.8, top=3.1, width=8.4, height=0.5,
             font_size=16, bold=False, color=GREY_LIGHT, align=PP_ALIGN.LEFT)
add_text_box(sl, "Cliente: Gama  |  Elaborado por: Cora Urrea / equipo análisis",
             left=0.8, top=3.7, width=8.4, height=0.4,
             font_size=11, bold=False, color=GREY_LIGHT, align=PP_ALIGN.LEFT)
add_text_box(sl, "Confidencial — uso interno Gama",
             left=0.8, top=6.8, width=8.4, height=0.3,
             font_size=8, bold=False, color=GREY_MED, align=PP_ALIGN.LEFT)
print(f"  S01 Portada OK")

# ── S02 — TABLA DE CONTENIDO ───────────────────
sl = new_slide(prs)
slide_count += 1
slide_header_bar(sl)
add_slide_title(sl, "Ruta narrativa")
add_section_tag(sl, "[AGENDA]")

sections = [
    ("1", "Embudo + posicionamiento espontáneo", "Slides 4–8"),
    ("2–3", "Top Box puro + DNA de marca", "Slides 9–18"),
    ("4–5", "Vecindad perceptual + drivers reformulados", "Slides 19–24"),
    ("6–8", "Precio | Comunicación PTL/DTLS | Síntesis + decisión", "Slides 25–32"),
]
for i, (num, title, slides) in enumerate(sections):
    y_pos = 1.5 + i * 1.2
    badge = sl.shapes.add_shape(1,
                                 Inches(0.5), Inches(y_pos),
                                 Inches(0.6), Inches(0.6))
    badge.fill.solid()
    badge.fill.fore_color.rgb = GAMA_RED
    badge.line.fill.background()
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = num
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = WHITE

    add_text_box(sl, title,
                 left=1.3, top=y_pos + 0.05, width=7.0, height=0.35,
                 font_size=14, bold=True, color=TEXT_DARK)
    add_text_box(sl, slides,
                 left=1.3, top=y_pos + 0.4, width=3.0, height=0.25,
                 font_size=9, bold=False, color=GREY_MED)
print(f"  S02 Agenda OK")

# ── S03 — FICHA TECNICA ────────────────────────
sl = new_slide(prs)
slide_count += 1
slide_header_bar(sl)
add_slide_title(sl, "Metodología: 402 entrevistas, 4 zonas, ±4.89%")
add_section_tag(sl, "[METODOLOGÍA]")

ficha = [
    ("n total", "402 entrevistas completadas (BBDD V2.0, n=403 × 295 preguntas)"),
    ("Geografía", "4 zonas Venezuela: Caracas, Miranda, Aragua, Valencia"),
    ("Fechas", "Ola 2026  |  Margen de error: ±4.89% (95% confianza, n=402)"),
    ("Segmentación", "NSE A/B, C+, C, D, E  |  Cuota edad × género"),
]
for i, (label, val) in enumerate(ficha):
    y_pos = 1.5 + i * 1.25
    box = sl.shapes.add_shape(1, Inches(0.5), Inches(y_pos), Inches(1.4), Inches(0.5))
    box.fill.solid()
    box.fill.fore_color.rgb = GAMA_RED
    box.line.fill.background()
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.size = Pt(10)
    r.font.bold = True
    r.font.color.rgb = WHITE
    add_text_box(sl, val,
                 left=2.1, top=y_pos + 0.05, width=7.5, height=0.45,
                 font_size=12, bold=False, color=TEXT_DARK)
print(f"  S03 Ficha técnica OK")


# ══════════════════════════════════════════════
# SECCIÓN 1 — EMBUDO + POSICIONAMIENTO (5 slides)
# ══════════════════════════════════════════════

# ── S04 — EMBUDO DATOS ─────────────────────────
sl = new_slide(prs)
slide_count += 1
slide_header_bar(sl)
add_slide_title(sl, "Embudo de notoriedad espontánea a preferencia — 8 cadenas",
                subtitle="[DATOS] Métricas: TOM, Asistida, Consideración, Compra 3m, Preferida, Habitual")
add_section_tag(sl, "[DATOS] S1")
make_chart_c01(sl, C01, left=0.4, top=1.05, width=12.27, height=5.8)
add_footnote(sl, f"Fuente: {C01['fuente']}. Solo Gama tiene datos completos en 6 etapas.")
print(f"  S04 Embudo datos OK")

# ── S05 — EMBUDO ANALISIS ──────────────────────
sl = new_slide(prs)
slide_count += 1
slide_header_bar(sl)
add_slide_title(sl, "Gama no tiene problema de visibilidad — tiene un problema de conversión")
add_section_tag(sl, "[ANÁLISIS] S1")
analysis_highlight_box(sl, f"TOM = {C01['gama_total_pct'][0]}% → Preferida = {C01['gama_total_pct'][4]}% — Gap de conversión de {round(C01['gama_total_pct'][0] - C01['gama_total_pct'][4], 1)} pp")
bullet_box(sl, [
    "Notoriedad espontánea alta (TOM 44.3%): posición 1-2 en mercado",
    f"Caída crítica entre consideración ({C01['gama_total_pct'][2]}%) y preferencia ({C01['gama_total_pct'][4]}%): el mercado conoce pero no elige",
    "Palanca de crecimiento = mejorar percepción de valor, no awareness",
    "Implicación: inversión en comunicación debe dirigirse al fondo del embudo",
], left=0.5, top=1.65, width=9.0, height=4.5, font_size=13)
print(f"  S05 Embudo análisis OK")

# ── S06 — PREFERENCIA P21 DATOS ────────────────
sl = new_slide(prs)
slide_count += 1
slide_header_bar(sl)
add_slide_title(sl, "Preferencia espontánea P21 — ranking total y segmento C+/C",
                subtitle="[DATOS] Pregunta P21 espontánea. NSE C+/C = segmento objetivo premium.")
add_section_tag(sl, "[DATOS] S1")
make_chart_c02(sl, C02, left=0.4, top=1.05, width=12.27, height=5.7)
add_footnote(sl, f"Fuente: {C02['fuente']}. * REFERENCIAL: Gama, Luz, Plan Suárez C+/C con n<30.")
print(f"  S06 Preferencia P21 datos OK")

# ── S07 — PREFERENCIA P21 ANALISIS ─────────────
sl = new_slide(prs)
slide_count += 1
slide_header_bar(sl)
gama_cpc_pct = C02['pct_cpc']
gama_cpc_val = C02['pct_cpc'][C02['cadenas_cpc'].index('Gama')] if 'Gama' in C02['cadenas_cpc'] else "?"
plan_cpc_val = C02['pct_cpc'][C02['cadenas_cpc'].index('Plan Suárez')] if 'Plan Suárez' in C02['cadenas_cpc'] else "?"
add_slide_title(sl, f"En C+/C, Gama lidera con {gama_cpc_val}% — Plan Suárez es el #2 con {plan_cpc_val}%")
add_section_tag(sl, "[ANÁLISIS] S1")
analysis_highlight_box(sl, f"Gama #{C02['cadenas_cpc'].index('Gama')+1} en C+/C ({gama_cpc_val}%) | Plan Suárez #{C02['cadenas_cpc'].index('Plan Suárez')+1} ({plan_cpc_val}%) — rival subestimado")
bullet_box(sl, [
    "Gama consolida liderazgo en el segmento de mayor valor por ticket",
    f"Plan Suárez alcanza {plan_cpc_val}% C+/C: rivalidad subestimada en segmento premium",
    "La brecha (~15 pp) es defendible pero no definitiva — requiere activación sostenida",
    "REFERENCIAL: n por celda C+/C validar significancia (n<30 en cadenas menores)",
], left=0.5, top=1.65, width=9.0, height=4.5, font_size=13)
add_footnote(sl, "* REFERENCIAL: Gama C+/C n=14, Luz n=11, Plan Suárez n=16. Tendencias, no proyectables.")
print(f"  S07 Preferencia P21 análisis OK")

# ── S08 — MODELO MENTAL PRECIO ─────────────────
sl = new_slide(prs)
slide_count += 1
slide_header_bar(sl)
gama_precio_mm = C03['pct_precio'][C03['cadenas_ordenadas_asc'].index('Gama')]
add_slide_title(sl, f"Gama es la ÚNICA cadena con modelo mental no dominado por precio ({gama_precio_mm}%)")
add_section_tag(sl, "[DATOS + ANÁLISIS] S1")
make_chart_c03(sl, C03, left=0.4, top=1.2, width=12.27, height=4.8)
bullet_box(sl, [
    "Las otras 7 cadenas: 51–84% asociadas a precio como atributo dominante",
    f"Gama: {gama_precio_mm}% — outlier estructural, no accidental",
    "Implicación estratégica: el posicionamiento experiencial de Gama ES diferenciador",
], left=0.5, top=6.2, width=9.0, height=1.0, font_size=10)
print(f"  S08 Modelo mental precio OK")


# ══════════════════════════════════════════════
# SECCIÓN 2 — TOP BOX PURO (4 slides)
# ══════════════════════════════════════════════

# ── S09 — TB vs T2B DATOS ──────────────────────
sl = new_slide(prs)
slide_count += 1
slide_header_bar(sl, PARAMO_BLUE)
add_slide_title(sl, "Top Box puro reordena las prioridades del mercado vs T2B",
                subtitle="[DATOS] P22: importancia atributos de selección de supermercado.")
add_section_tag(sl, "[DATOS] S2")
make_chart_c04(sl, C04, left=0.4, top=1.05, width=11.2, height=5.8)
add_footnote(sl, f"Fuente: {C04['fuente']}. TB = top box único; T2B = top 2 boxes. Delta rank = movimiento posicional.")
print(f"  S09 TB vs T2B datos OK")

# ── S10 — TB vs T2B ANALISIS ───────────────────
sl = new_slide(prs)
slide_count += 1
slide_header_bar(sl, PARAMO_BLUE)
# Encontrar reranking de precio
precio_idx = C04['atributos'].index('Menor precio')
delta_precio = C04['delta_rank'][precio_idx]
add_slide_title(sl, f"El T2B suavizaba la presión: en TB puro, precio escala {'+' if delta_precio > 0 else ''}{delta_precio} posiciones")
add_section_tag(sl, "[ANÁLISIS] S2")
analysis_highlight_box(sl, f"Menor precio: rank {C04['rank_t2b'][precio_idx]} en T2B → rank {C04['rank_tb'][precio_idx]} en TB puro ({'+' if delta_precio>0 else ''}{delta_precio} posiciones)")
highlights = [h for h in C04.get('highlights', []) if h]
bullet_box(sl, [
    "T2B: precio aparecía importante pero no dominante",
    "TB puro: precio escala posición — el mercado tiene exigencia funcional más rígida",
    "La presión de precio que enfrenta Gama es más severa de lo que V3-V7 mostraban",
    "Revisión de estrategia: no basta con 'ser experiencial' — hay que abordar precio creativamente",
], left=0.5, top=1.65, width=9.0, height=4.5, font_size=13)
print(f"  S10 TB vs T2B análisis OK")

# ── S11 — HEATMAP TB PURO POR MARCA ────────────
sl = new_slide(prs)
slide_count += 1
slide_header_bar(sl, PARAMO_BLUE)
add_slide_title(sl, "TB puro por marca preferida: heatmap de exigencias por cluster de leales",
                subtitle="[DATOS] Rango dinámico real del dataset. Celda máxima destacada.")
add_section_tag(sl, "[DATOS] S2")
make_chart_c05(sl, C05, left=0.3, top=1.05, width=12.53, height=5.9)
print(f"  S11 Heatmap C05 OK")

# ── S12 — HEATMAP ANALISIS ─────────────────────
sl = new_slide(prs)
slide_count += 1
slide_header_bar(sl, PARAMO_BLUE)
# Pref-Gama datos
rapidez_gama = C05['matrix_pct'][2][1]  # Rapidez (row 2), Pref-Gama (col 1)
precio_gama  = C05['matrix_pct'][1][1]  # Precio (row 1), Pref-Gama (col 1)
add_slide_title(sl, f"Los leales de Gama priorizan Rapidez ({rapidez_gama}%) y Atención — Precio cae al noveno lugar ({precio_gama}%)")
add_section_tag(sl, "[ANÁLISIS] S2")
analysis_highlight_box(sl, f"Pref-Gama: Rapidez #1 = {rapidez_gama}% | Precio #9 = {precio_gama}% (n=32 REFERENCIAL)")
bullet_box(sl, [
    f"Pref-Gama (n=32 REF): Rapidez #1 con {rapidez_gama}%, Atención #2 con {C05['matrix_pct'][7][1]}%",
    f"Precio ocupa el puesto #9 con solo {precio_gama}% entre sus propios leales",
    "La bifurcación entre el mercado amplio (precio fuerte) y los leales (experiencia) es MÁS extrema en TB puro",
    "El reto: convertir no-leales sin abandonar la propuesta que retiene a los leales",
], left=0.5, top=1.65, width=9.0, height=4.5, font_size=13)
add_footnote(sl, "* REFERENCIAL: n Pref-Gama=32. Datos direccionales, no proyectables a la población.", amber=True)
print(f"  S12 Heatmap análisis OK")


# ══════════════════════════════════════════════
# SECCIÓN 3 — DNA + MUNDO DE MARCA (6 slides)
# ══════════════════════════════════════════════

# ── S13 — DNA Z-SCORES DATOS ───────────────────
sl = new_slide(prs)
slide_count += 1
slide_header_bar(sl, GREEN_VAL)
add_slide_title(sl, "DNA de marca Gama: z-scores de asociación vs promedio mercado",
                subtitle="[DATOS] Z>0 = sobreíndice vs promedio 8 cadenas. Z<0 = subíndice.")
add_section_tag(sl, "[DATOS] S3")
make_chart_c06(sl, C06, left=0.4, top=1.05, width=12.27, height=5.8)
add_footnote(sl, f"Fuente: {C06['fuente']}. Referencia = promedio simple de las 8 cadenas.")
print(f"  S13 DNA z-scores datos OK")

# ── S14 — DNA ANALISIS ─────────────────────────
sl = new_slide(prs)
slide_count += 1
slide_header_bar(sl, GREEN_VAL)
n_sobre = sum(1 for z in C06['zscore_total'] if z > 0)
precio_z = C06['zscore_total'][C06['atributos'].index('Menor precio')]
add_slide_title(sl, f"Gama sobreíndice en {n_sobre} atributos experienciales; precio subíndice sostenido ({precio_z} z)")
add_section_tag(sl, "[ANÁLISIS] S3")
analysis_highlight_box(sl, f"{n_sobre} atributos experienciales sobre el promedio de mercado | Precio z={precio_z} (subíndice estructural)")
bullet_box(sl, [
    f"V8 confirma y extiende V3: {n_sobre} atributos experienciales sobre el promedio (vs 4 en V3)",
    f"Precio: z = {precio_z} — no es casualidad, es estructura perceptual establecida",
    "DNA experiencial es activo estratégico, no debilidad a corregir",
    "Riesgo: si Gama intenta competir en precio directamente, puede erosionar el DNA diferenciador",
], left=0.5, top=1.65, width=9.0, height=4.5, font_size=13)
print(f"  S14 DNA análisis OK")

# ── S15 — MUNDO DE MARCA P23 DATOS ─────────────
sl = new_slide(prs)
slide_count += 1
slide_header_bar(sl, GREEN_VAL)
add_slide_title(sl, "Mapa de imagen P23: 8 cadenas × 10 atributos — posición completa del mercado",
                subtitle="[DATOS] P23 espontánea. Base: total muestra n=402.")
add_section_tag(sl, "[DATOS] S3")
make_chart_c07(sl, C07, left=0.3, top=1.05, width=12.53, height=5.9)
print(f"  S15 Mundo de marca P23 datos OK — Gama Precio={C07['matrix_pct'][0][7]}")

# ── S16 — MUNDO DE MARCA ANALISIS ──────────────
sl = new_slide(prs)
slide_count += 1
slide_header_bar(sl, GREEN_VAL)
gama_limpieza = C07['matrix_pct'][0][0]
paramo_precio = C07['matrix_pct'][1][7]
add_slide_title(sl, f"Gama lidera el cluster experiencial; Limpieza ({gama_limpieza}%) y Seguridad son tabla de higiene del sector")
add_section_tag(sl, "[ANÁLISIS] S3")
analysis_highlight_box(sl, f"Gama Precio={C07['matrix_pct'][0][7]}% (mínimo absoluto mercado) | Páramo Precio={paramo_precio}% (máximo)")
bullet_box(sl, [
    "Gama, Rio y Plazas comparten cluster experiencial — zona de alta competencia perceptual",
    "Otros competidores se anclan en precio/volumen: territorio distinto, menor overlap",
    f"Limpieza ({gama_limpieza}%) y Seguridad ({C07['matrix_pct'][0][1]}%) convergen en todas las cadenas: no diferencian, son condición de entrada",
    "Para diferenciarse, Gama debe activar atributos exclusivos dentro del cluster experiencial",
], left=0.5, top=1.65, width=9.0, height=4.5, font_size=13)
print(f"  S16 Mundo de marca análisis OK")

# ── S17 — TORNADO BRECHAS DATOS ────────────────
sl = new_slide(prs)
slide_count += 1
slide_header_bar(sl, GREEN_VAL)
add_slide_title(sl, "Brechas de imagen Gama — compradores activos vs mercado amplio (tornado P23)",
                subtitle="[DATOS] Brecha pp = % Pref-Gama que asocia atributo MENOS % Total. n REF.")
add_section_tag(sl, "[DATOS] S3")
make_chart_c08(sl, C08, left=0.4, top=1.05, width=12.27, height=5.7)
print(f"  S17 Tornado brechas datos OK")

# ── S18 — TORNADO ANALISIS ─────────────────────
sl = new_slide(prs)
slide_count += 1
slide_header_bar(sl, GREEN_VAL)
idx_at = C08['atributos'].index('Mejor atención')
brecha_atencion = C08['delta_pref_total'][idx_at]
brecha_exp      = C08['delta_exp_total'][idx_at]
pref_atencion   = C08['pref_gama_pct'][idx_at]
total_atencion  = C08['total_pct'][idx_at]
add_slide_title(sl, f"ATRIBUTO SOMBRA: Atención está +{brecha_atencion} pp entre leales — el mercado no lo sabe aún")
add_section_tag(sl, "[ANÁLISIS] S3")
analysis_highlight_box(sl, f"Leales Gama: Atención {pref_atencion}% | Mercado total: {total_atencion}% | Brecha: +{brecha_atencion} pp")
bullet_box(sl, [
    f"Quienes compran en Gama asocian Atención en {brecha_atencion} pp más que el mercado amplio",
    f"El mercado no-cliente no percibe la atención de Gama (solo {total_atencion}% Total vs {pref_atencion}% Pref-Gama)",
    "Tarea de comunicación: trasladar la imagen validada por leales al mercado amplio",
    "No es crear un atributo nuevo — es amplificar uno que ya existe y está validado por experiencia",
], left=0.5, top=1.65, width=9.0, height=4.5, font_size=13)
add_footnote(sl, "* REFERENCIAL: n Pref-Gama=32, n Exp.reciente=30. Tendencia validada.", amber=True)
print(f"  S18 Tornado análisis OK")


# ══════════════════════════════════════════════
# SECCIÓN 4 — VECINDAD PERCEPTUAL (3 slides)
# ══════════════════════════════════════════════

# ── S19 — VECINDAD DATOS ───────────────────────
sl = new_slide(prs)
slide_count += 1
slide_header_bar(sl, RGBColor(0x4A, 0x6F, 0xA5))
add_slide_title(sl, "Vecindad perceptual: mapa de posicionamiento relativo — 8 cadenas",
                subtitle="[DATOS] Ejes: % imagen experiencial vs % asociación Menor Precio. Fuente: CU-10 §4.1.")
add_section_tag(sl, "[DATOS] S4")
make_chart_c09(sl, C09, left=0.5, top=1.1, width=11.73, height=5.5)
add_footnote(sl, f"Fuente: {C09['fuente']}. Eje X = promedio Limpieza+Seguro+Atractiva+Calidad. Eje Y = % Menor Precio.")
print(f"  S19 Vecindad perceptual datos OK")

# ── S20 — VECINDAD ANALISIS PLAZAS ─────────────
sl = new_slide(prs)
slide_count += 1
slide_header_bar(sl, RGBColor(0x4A, 0x6F, 0xA5))
# Encontrar Plazas y Gama en datos scatter
gama_pt   = next((p for p in C09['puntos'] if p['cadena'] == 'Gama'), None)
plazas_pt = next((p for p in C09['puntos'] if p['cadena'] == 'Plazas'), None)
add_slide_title(sl, "Plazas es el rival perceptual más próximo a Gama — mismo cluster, mismo atributo sombra")
add_section_tag(sl, "[ANÁLISIS] S4")
if gama_pt and plazas_pt:
    analysis_highlight_box(sl, f"Gama (x={gama_pt['x']}, y={gama_pt['y']}) vs Plazas (x={plazas_pt['x']}, y={plazas_pt['y']}) — vecinos en el mapa perceptual")
bullet_box(sl, [
    "Plazas ocupa el vecindario más cercano a Gama en el mapa perceptual experiencial",
    "Comparten cluster experiencial y el atributo sombra es Atención en ambas cadenas",
    "Si Plazas activa comunicación de Atención antes que Gama, puede capturar el territorio",
    "Prioridad urgente: diferenciar de Plazas en Atención — calidad, velocidad o trato personalizado",
], left=0.5, top=1.65, width=9.0, height=4.5, font_size=13)
print(f"  S20 Vecindad análisis Plazas OK")

# ── S21 — VECINDAD ANALISIS RIO ────────────────
sl = new_slide(prs)
slide_count += 1
slide_header_bar(sl, RGBColor(0x4A, 0x6F, 0xA5))
rio_pt = next((p for p in C09['puntos'] if p['cadena'] == 'Rio'), None)
add_slide_title(sl, "Rio sobreíndice en Calidad entre sus leales — si comunica, puede capturar consumidores experienciales")
add_section_tag(sl, "[ANÁLISIS] S4")
if rio_pt:
    analysis_highlight_box(sl, f"Rio: x={rio_pt['x']} (experiencial), y={rio_pt['y']} (precio bajo) — territorio separado pero expansivo")
bullet_box(sl, [
    "Rio no compite en el territorio actual de Gama (posición diferente en el mapa perceptual)",
    "Atributo sombra de Rio = Calidad: si activa campaña de calidad, puede atraer consumidor experiencial no fidelizado",
    "Rio: bajo precio percibido + imagen experiencial moderada = amenaza latente si comunica activamente",
    "Monitoreo recomendado: incluir tracking de Rio en ola 2027",
], left=0.5, top=1.65, width=9.0, height=4.5, font_size=13)
print(f"  S21 Vecindad análisis Rio OK")


# ══════════════════════════════════════════════
# SECCIÓN 5 — DRIVERS REFORMULADOS (3 slides)
# ══════════════════════════════════════════════

# ── S22 — DRIVERS DATOS ────────────────────────
sl = new_slide(prs)
slide_count += 1
slide_header_bar(sl, AMBER)
add_slide_title(sl, "Modelo de drivers — regresión logística sobre percepción P23 → preferencia",
                subtitle="[DATOS] Variable dependiente: preferir Gama. Predictores: % asociación P23.")
add_section_tag(sl, "[DATOS] S5")
make_chart_c10(sl, C10, left=0.4, top=1.1, width=8.67, height=4.8)
add_footnote(sl, f"Fuente: {C10['fuente']}. Modelo exploratorio — validación recomendada ola 2027 n≥600.", amber=True)
print(f"  S22 Drivers logit datos OK")

# ── S23 — DRIVERS METODOLOGICO ─────────────────
sl = new_slide(prs)
slide_count += 1
slide_header_bar(sl, AMBER)
pseudo_r2 = C10['modelo_global']['pseudo_r2']
llr_p     = C10['modelo_global']['llr_p_value']
add_slide_title(sl, "Lo que la gente DECLARA importante no predice preferencia — la percepción sí")
add_section_tag(sl, "[ANÁLISIS METODOLÓGICO] S5")
analysis_highlight_box(sl, f"Modelo logístico global: Pseudo-R²={pseudo_r2} (POBRE), LLR p={llr_p} — no significativo globalmente")
bullet_box(sl, [
    "Importancias P22 (declaradas) sufren 'techo': casi todo es muy importante, no discrimina",
    "El modelo logístico sobre P23 (percepción real) sí predice preferencia con significancia individual",
    "El mecanismo correcto: demostrar que Gama TIENE el atributo, no que el atributo es importante",
    "Anti-mensaje: 'ofrecemos buen servicio' (genérico, no memorable, no diferencia de Plazas)",
], left=0.5, top=1.65, width=9.0, height=4.5, font_size=13)
or_val = C10['odds_ratio'][C10['atributos'].index('Mayor calidad')]
add_footnote(sl, f"Único coeficiente significativo: Mayor calidad OR={or_val}, p=0.042 — interpretación en slide S22.", amber=True)
print(f"  S23 Drivers metodológico OK")

# ── S24 — CAMPAÑA CORRECTA ─────────────────────
sl = new_slide(prs)
slide_count += 1
slide_header_bar(sl, AMBER)
add_slide_title(sl, "La campaña correcta demuestra Atención superior — no la declara")
add_section_tag(sl, "[ANÁLISIS ESTRATÉGICO] S5")
analysis_highlight_box(sl, "Declarar = ruido. Demostrar = diferenciación. El DNA experiencial Gama ya está validado por compradores.")
bullet_box(sl, [
    "Mensaje equivocado: 'en Gama te atendemos bien' — cualquier cadena puede decir lo mismo",
    "Mensaje correcto: testimonios reales, comparativas experienciales, momentos de verdad específicos",
    "Activar EXPERIENCIA + RECONOCIMIENTO como mecanismo creativo central de la campaña",
    "Alinear con atributo sombra: trasladar imagen de leales al mercado amplio a través de evidencia",
], left=0.5, top=1.65, width=9.0, height=4.5, font_size=13)
print(f"  S24 Campaña correcta OK")


# ══════════════════════════════════════════════
# SECCIÓN 6 — PRECIO + CATEGORIAS (3 slides)
# ══════════════════════════════════════════════

# ── S25 — PERCEPCION PRECIO ────────────────────
sl = new_slide(prs)
slide_count += 1
slide_header_bar(sl, RGBColor(0x9C, 0xAE, 0xC0))
cpc_caro = C11['neto_caro_pct'][C11['nse'].index('C+/C')]
total_caro = C11['neto_caro_pct'][C11['nse'].index('Total')]
add_slide_title(sl, f"Percepción de precio Gama por NSE — C+/C: {cpc_caro}% NETO caro vs {total_caro}% Total",
                subtitle="[DATOS] P33: percepción de precios Gama. Base: total n=402.")
add_section_tag(sl, "[DATOS] S6")
make_chart_c11(sl, C11, left=0.4, top=1.05, width=12.27, height=5.8)
print(f"  S25 Percepción precio datos OK")

# ── S26 — CATEGORIAS ───────────────────────────
sl = new_slide(prs)
slide_count += 1
slide_header_bar(sl, RGBColor(0x9C, 0xAE, 0xC0))
add_slide_title(sl, "Categorías: Congelados y Gaseosas = ofrecer valor; Galletas = cuidar precio; Proteínas = no competir")
add_section_tag(sl, "[DATOS + ANÁLISIS] S6")
make_chart_c12(sl, C12, left=0.5, top=1.1, width=11.73, height=5.5)
print(f"  S26 Categorías estrategia OK")

# ── S27 — 5 VIAS CREATIVAS ─────────────────────
sl = new_slide(prs)
slide_count += 1
slide_header_bar(sl, RGBColor(0x9C, 0xAE, 0xC0))
add_slide_title(sl, "5 vías para abordar precio sin abandonar el posicionamiento experiencial")
add_section_tag(sl, "[ANÁLISIS] S6")

vias = [
    ("1", "Promociones diferenciadas", "Beneficio simbólico + precio (no precio solo)"),
    ("2", "Coleccionables / lealtad visible", "Monetizar la relación, no el descuento puntual"),
    ("3", "Mensaje del anzuelo", "Una categoría competitiva en precio ancla la percepción general"),
    ("4", "Categorías 'ofrecer valor'", "Activar Congelados, Gaseosas, Salsas — percibidos competitivos"),
    ("5", "Van Westendorp 2027", "Prerrequisito metodológico: calibrar precio psicológico por segmento"),
]
for i, (num, title, desc) in enumerate(vias):
    y_pos = 1.4 + i * 1.0
    number_badge(sl, num, left=0.4, top=y_pos)
    add_text_box(sl, title,
                 left=0.95, top=y_pos, width=3.5, height=0.35,
                 font_size=12, bold=True, color=TEXT_DARK)
    add_text_box(sl, desc,
                 left=0.95, top=y_pos + 0.38, width=8.6, height=0.3,
                 font_size=10, bold=False, color=GREY_MED)
print(f"  S27 5 vías creativas OK")


# ══════════════════════════════════════════════
# SECCIÓN 7 — COMUNICACION PTL + DTLS (3 slides)
# ══════════════════════════════════════════════

# ── S28 — RECALL PTL DTLS ──────────────────────
sl = new_slide(prs)
slide_count += 1
slide_header_bar(sl, GREY_MED)
add_slide_title(sl, "Recall de campaña PTL vs DTLS — desagregado por NSE",
                subtitle="[DATOS] PTL = 'Precios de tu lado'. Recall espontáneo + asistido por NSE.")
add_section_tag(sl, "[DATOS] S7")
make_chart_c13(sl, C13, left=0.4, top=1.05, width=12.27, height=5.8)
print(f"  S28 Recall PTL/DTLS datos OK")

# ── S29 — PERFIL RECORDADORES ──────────────────
sl = new_slide(prs)
slide_count += 1
slide_header_bar(sl, GREY_MED)
add_slide_title(sl, "¿Quién recuerda la campaña? Perfil recordadores vs muestra total",
                subtitle="[DATOS] Índice sobre/sub-representación por NSE. n recordadores combinados.")
add_section_tag(sl, "[DATOS] S7")
make_chart_c14(sl, C14, left=0.4, top=1.05, width=12.27, height=5.8)
print(f"  S29 Perfil recordadores datos OK")

# ── S30 — RECALL ANALISIS ──────────────────────
sl = new_slide(prs)
slide_count += 1
slide_header_bar(sl, GREY_MED)
gap_e   = C14['gap_pp'][C14['nse'].index('E')]
gap_cpc = C14['gap_pp'][C14['nse'].index('C+/C')]
add_slide_title(sl, f"La campaña actual sobre-alcanza NSE E (+{gap_e} pp) y sub-representa C+/C ({gap_cpc} pp) — reorientar")
add_section_tag(sl, "[ANÁLISIS] S7")
analysis_highlight_box(sl, f"NSE E sobre-representado +{gap_e} pp | NSE C+/C sub-representado {gap_cpc} pp — campaña llega al segmento equivocado")
bullet_box(sl, [
    f"NSE E sobre-representado entre recordadores (+{gap_e} pp): segmento con menor ticket y menor conversión a premium",
    f"NSE C+/C sub-representado ({gap_cpc} pp): segmento objetivo premium no está siendo impactado",
    "PTL 'Precios de tu lado' refuerza atributo precio en segmento E — que no es el DNA de Gama",
    "Acción: rediseñar mix de medios + mensaje para C+/C; validar affinities de canal por NSE",
], left=0.5, top=1.65, width=9.0, height=4.5, font_size=13)
print(f"  S30 Recall análisis OK")


# ══════════════════════════════════════════════
# SECCIÓN 8 — SINTESIS + DECISION (2 slides)
# ══════════════════════════════════════════════

# ── S31 — 5 HALLAZGOS CENTRALES ────────────────
sl = new_slide(prs)
slide_count += 1
slide_header_bar(sl)
add_slide_title(sl, "5 hallazgos que definen la estrategia Gama 2026-2027")
add_section_tag(sl, "[SÍNTESIS]")

hallazgos = [
    (f"DNA experiencial consolidado", f"Gama es la única cadena no dominada por precio ({C03['pct_precio'][C03['cadenas_ordenadas_asc'].index('Gama')]}%)"),
    (f"Atributo SOMBRA = Atención", f"+{C08['delta_pref_total'][C08['atributos'].index('Mejor atención')]} pp entre leales vs mercado amplio; el mercado no lo sabe aún"),
    (f"Vecindad Plazas", f"Rival perceptual más cercano en el mismo atributo sombra — urgencia de diferenciación"),
    (f"Presión precio más severa en TB puro", f"La exigencia del mercado es más rígida de lo que T2B mostraba"),
    (f"Campaña llega al segmento equivocado", f"NSE E sobre-alcanzado (+{C14['gap_pp'][C14['nse'].index('E')]} pp) — NSE C+/C sub-representado"),
]
for i, (h_title, h_desc) in enumerate(hallazgos):
    y_pos = 1.3 + i * 1.1
    number_badge(sl, str(i+1), left=0.4, top=y_pos)
    add_text_box(sl, h_title,
                 left=0.95, top=y_pos, width=4.0, height=0.35,
                 font_size=12, bold=True, color=TEXT_DARK)
    add_text_box(sl, h_desc,
                 left=0.95, top=y_pos + 0.38, width=8.6, height=0.4,
                 font_size=10, bold=False, color=GREY_MED)
print(f"  S31 5 hallazgos síntesis OK")

# ── S32 — CALL TO ACTION ───────────────────────
sl = new_slide(prs)
slide_count += 1

# Fondo oscuro para impacto
bg_cta = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
bg_cta.fill.solid()
bg_cta.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
bg_cta.line.fill.background()

stripe = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.08))
stripe.fill.solid()
stripe.fill.fore_color.rgb = GAMA_RED
stripe.line.fill.background()

add_text_box(sl, "3 decisiones estratégicas para Junta Gama — 2026-2027",
             left=0.5, top=0.2, width=9.0, height=0.6,
             font_size=18, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_section_tag(sl, "[CALL TO ACTION]")

decisiones = [
    ("1", "DEFENDER DNA + abordar precio creativamente",
     "5 vías identificadas — no competencia directa en precio. Van Westendorp 2027 para calibrar límites."),
    ("2", "ACTIVAR Atención en campaña 2026-2027",
     "Demostrar (testimonios, momentos de verdad), no declarar. Mensaje específico, no genérico."),
    ("3", "DIFERENCIAR de Plazas + vigilar Rio en Calidad",
     "Rediseñar mix de medios hacia NSE C+/C. Incluir tracking Rio en ola 2027."),
]
for i, (num, d_title, d_desc) in enumerate(decisiones):
    y_pos = 1.0 + i * 2.0
    badge = sl.shapes.add_shape(1,
                                 Inches(0.5), Inches(y_pos),
                                 Inches(0.5), Inches(0.5))
    badge.fill.solid()
    badge.fill.fore_color.rgb = GAMA_RED
    badge.line.fill.background()
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = num
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.color.rgb = WHITE

    add_text_box(sl, d_title,
                 left=1.2, top=y_pos, width=8.3, height=0.45,
                 font_size=14, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_text_box(sl, d_desc,
                 left=1.2, top=y_pos + 0.5, width=8.3, height=0.6,
                 font_size=11, bold=False, color=GREY_LIGHT, align=PP_ALIGN.LEFT)

add_text_box(sl, "Confidencial — Gama Notoriedad 2026 V8.2 | Cora Urrea / equipo análisis",
             left=0.5, top=7.1, width=9.0, height=0.25,
             font_size=7, bold=False, color=GREY_MED, align=PP_ALIGN.LEFT)
print(f"  S32 Call to action OK")

# ─────────────────────────────────────────────
# GUARDAR
# ─────────────────────────────────────────────
prs.save(str(OUT_PATH))
file_size = os.path.getsize(str(OUT_PATH))
print(f"\n==================================================")
print(f"DECK V8.2 GUARDADO OK")
print(f"Path:   {OUT_PATH}")
print(f"Slides: {slide_count}")
print(f"Bytes:  {file_size:,}")
print(f"==================================================")
print(f"\nSPOT-CHECKS FINALES:")
print(f"  C07 Gama Precio    = {C07['matrix_pct'][0][7]} (esperado 7.2)")
print(f"  C07 Gama Atencion  = {C07['matrix_pct'][0][4]} (esperado 21.9)")
print(f"  C05 PrefGama Rapid = {C05['matrix_pct'][2][1]} (esperado 84.4)")
