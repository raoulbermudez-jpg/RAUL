"""
Common renderer for V3 deck slides — Gama Notoriedad 2026.
Renders a slide from a dict-spec using pptx_helpers (V2 visual continuity).
"""
import sys
from pathlib import Path
sys.path.insert(0, r'C:\Raul\03-projects\consultoria-externa\gama-notoriedad-2026\02-analysis')

from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx_helpers import (
    new_presentation, add_blank_slide, set_slide_bg,
    add_title_bar, add_footer, add_caveat_chip, add_notes,
    add_text_block, add_bullets, add_pptx_table,
    GAMA_RED, GAMA_RED_LIGHT, GAMA_RED_DARK,
    GRIS_OSCURO, GRIS_MEDIO, GRIS_CLARO, FONDO_CLARO,
    BLANCO, NEGRO, VERDE_OK, NARANJA_WARN, ROJO_ALERT,
    SLIDE_WIDTH, SLIDE_HEIGHT,
)

sys.stdout.reconfigure(encoding='utf-8') if hasattr(sys.stdout, 'reconfigure') else None

PLOTS_DIR = Path(r"C:\Raul\03-projects\consultoria-externa\gama-notoriedad-2026\02-analysis\cuanti\plots")


def render_portada(prs, slide_n, title, subtitle, meta_lines, version_label='V3'):
    s = add_blank_slide(prs)
    set_slide_bg(s, BLANCO)
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(2.5), SLIDE_HEIGHT)
    band.fill.solid(); band.fill.fore_color.rgb = GAMA_RED
    band.line.fill.background()
    logo = s.shapes.add_textbox(Inches(0.4), Inches(0.5), Inches(2), Inches(1))
    p = logo.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = 'GAMA'
    r.font.size = Pt(36); r.font.bold = True; r.font.color.rgb = BLANCO
    sello = s.shapes.add_textbox(Inches(0.4), Inches(6.3), Inches(2), Inches(0.8))
    p = sello.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = '2026'
    r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = BLANCO
    # version chip
    chip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(5.6), Inches(2), Inches(0.5))
    chip.fill.solid(); chip.fill.fore_color.rgb = BLANCO
    chip.line.color.rgb = BLANCO
    p = chip.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = version_label
    r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = GAMA_RED

    t = s.shapes.add_textbox(Inches(3.0), Inches(2.3), Inches(9.8), Inches(1.5))
    tf = t.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = title
    r.font.size = Pt(32); r.font.bold = True; r.font.color.rgb = GAMA_RED

    sub = s.shapes.add_textbox(Inches(3.0), Inches(3.9), Inches(9.8), Inches(1.0))
    tf = sub.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = subtitle
    r.font.size = Pt(18); r.font.italic = True; r.font.color.rgb = GRIS_OSCURO

    for i, line in enumerate(meta_lines):
        ln = s.shapes.add_textbox(Inches(3.0), Inches(5.1 + i * 0.35), Inches(9.8), Inches(0.35))
        p = ln.text_frame.paragraphs[0]
        r = p.add_run(); r.text = line
        r.font.size = Pt(11); r.font.color.rgb = GRIS_OSCURO
    return s


def render_slide(prs, spec):
    """Render a content slide from a dict spec.

    spec keys:
      n: slide number (int) -> used for bloque_label
      title: str
      bloque: str (e.g. 'Bloque 1' or 'Apéndice')
      takeaway: str (highlighted right under title)
      bullets: list[str|tuple]  (rendered as bullet list)
      table: dict {headers: [...], rows: [[...]]}  (optional)
      image: str path  (optional — embed plot)
      caveat: str  (chip al pie)
      notes: str  (presenter notes, optional)
    """
    s = add_blank_slide(prs)
    bloque_label = spec.get('bloque', '')
    add_title_bar(s, spec['title'], bloque_label=bloque_label)

    cursor_top = 0.95  # Inches below title bar

    # Takeaway (highlighted under title)
    if spec.get('takeaway'):
        tk_box = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.3), Inches(cursor_top),
            Inches(12.7), Inches(0.7),
        )
        tk_box.fill.solid(); tk_box.fill.fore_color.rgb = FONDO_CLARO
        tk_box.line.color.rgb = GAMA_RED; tk_box.line.width = Pt(1)
        tf = tk_box.text_frame; tf.word_wrap = True
        tf.margin_left = Inches(0.15); tf.margin_right = Inches(0.15)
        tf.margin_top = Inches(0.08); tf.margin_bottom = Inches(0.08)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = '★ ' + spec['takeaway']
        r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = GAMA_RED_DARK
        cursor_top += 0.85

    has_image = bool(spec.get('image'))
    has_table = bool(spec.get('table'))

    if has_image or has_table:
        # split: bullets left, visual right
        bullets_width = Inches(7.0)
        visual_left = Inches(7.6)
        visual_width = Inches(5.4)
    else:
        bullets_width = Inches(12.7)
        visual_left = None
        visual_width = None

    if spec.get('bullets'):
        bullets_height = Inches(7.0 - cursor_top - 0.6)
        # Normalize bullets: tuple-of-1 -> bold-only string with ':' suffix
        normalized = []
        for b in spec['bullets']:
            if isinstance(b, tuple) and len(b) == 1:
                normalized.append((b[0], ''))
            else:
                normalized.append(b)
        add_bullets(
            s, normalized,
            Inches(0.3), Inches(cursor_top),
            bullets_width, bullets_height,
            font_size=11,
        )

    if has_image:
        try:
            img_path = spec['image']
            s.shapes.add_picture(
                img_path,
                visual_left, Inches(cursor_top),
                width=visual_width,
            )
        except Exception as e:
            add_text_block(
                s, f'[plot no disponible: {Path(spec["image"]).name}]',
                visual_left, Inches(cursor_top), visual_width, Inches(0.4),
                font_size=10, color=GRIS_MEDIO,
            )

    if has_table:
        tbl = spec['table']
        rows = tbl['rows']
        headers = tbl['headers']
        n_rows = 1 + len(rows)
        table_height = Inches(min(5.0, 0.4 + 0.35 * n_rows))
        add_pptx_table(
            s, headers, rows,
            visual_left, Inches(cursor_top),
            visual_width, table_height,
        )

    # Caveat al pie
    if spec.get('caveat'):
        add_caveat_chip(s, spec['caveat'], top=Inches(6.65))

    add_footer(s, custom=f'V3 · slide {spec.get("n", "")}')
    add_notes(
        s,
        spec.get('notes_preguntas', 'Ver outline VI-1a §slide ' + str(spec.get('n', ''))),
        spec.get('notes_calculo', 'Ver CU-3 v2, CU-4 v2 según corresponda'),
        spec.get('notes_lectura', spec.get('takeaway', 'Ver outline')),
    )
    return s


def render_cover_table_slide(prs, title, bloque, takeaway, headers, rows, caveat=None, n=None):
    """Special: table-centric slide (no bullets)."""
    s = add_blank_slide(prs)
    add_title_bar(s, title, bloque_label=bloque)
    cursor_top = 0.95
    if takeaway:
        tk_box = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.3), Inches(cursor_top),
            Inches(12.7), Inches(0.7),
        )
        tk_box.fill.solid(); tk_box.fill.fore_color.rgb = FONDO_CLARO
        tk_box.line.color.rgb = GAMA_RED; tk_box.line.width = Pt(1)
        tf = tk_box.text_frame; tf.word_wrap = True
        tf.margin_left = Inches(0.15); tf.margin_right = Inches(0.15)
        tf.margin_top = Inches(0.08); tf.margin_bottom = Inches(0.08)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = '★ ' + takeaway
        r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = GAMA_RED_DARK
        cursor_top += 0.85
    n_rows = 1 + len(rows)
    table_h = Inches(min(5.5, 0.4 + 0.35 * n_rows))
    add_pptx_table(s, headers, rows, Inches(0.3), Inches(cursor_top), Inches(12.7), table_h)
    if caveat:
        add_caveat_chip(s, caveat, top=Inches(6.65))
    add_footer(s, custom=f'V3 · slide {n}' if n else 'V3')
    return s
