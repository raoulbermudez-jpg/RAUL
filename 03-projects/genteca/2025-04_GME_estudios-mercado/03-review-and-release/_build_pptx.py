"""
Build: Resultados_Investigacion_Pantallas_GME_VW_para_I+D_v1_2026-05-06.pptx
Audience: GME R&D team (internal)
Purpose: Enable R&D to continue development with actionable recommendations
         on screen selection, UI changes, feature prioritization, and pricing target.
Author: Vivienne (Presentation Designer, RAUL)
Date: 2026-05-06
"""

import sys, os
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Inches, Pt
import copy

# ── Paths ──────────────────────────────────────────────────────────────────────
BASE   = r"C:\RAUL\03-projects\genteca\2025-04_GME_estudios-mercado"
OUT    = os.path.join(BASE, "03-review-and-release",
                      "Resultados_Investigacion_Pantallas_GME_VW_para_I+D_v1_2026-05-06.pptx")
IMG_R  = os.path.join(BASE, "_img_R")
IMG_B  = os.path.join(BASE, "_img_B")
IMG_M  = os.path.join(BASE, "_img_M")
IMG_VW = os.path.join(BASE, "van_westendorp")

# ── Brand palette ──────────────────────────────────────────────────────────────
C_ORANGE  = RGBColor(0xE8, 0x6B, 0x1A)   # Exceline accent orange
C_NAVY    = RGBColor(0x1A, 0x2E, 0x4A)   # Deep navy (primary text / backgrounds)
C_BLUE    = RGBColor(0x1E, 0x5C, 0x9E)   # Corporate blue
C_LGRAY   = RGBColor(0xF2, 0xF4, 0xF7)   # Light gray background
C_MGRAY   = RGBColor(0xB0, 0xBA, 0xC8)   # Mid gray (secondary text)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK   = RGBColor(0x1A, 0x1A, 0x1A)
C_GREEN   = RGBColor(0x21, 0x7A, 0x3E)   # positive / confirm
C_RED     = RGBColor(0xC0, 0x32, 0x2B)   # warning
C_AMBER   = RGBColor(0xD4, 0x8C, 0x0A)   # caution

FONT_BODY  = "Calibri"
FONT_TITLE = "Calibri"

# ── Slide dimensions (16:9 widescreen) ────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK_LAYOUT = prs.slide_layouts[6]   # blank

# ══════════════════════════════════════════════════════════════════════════════
# Helper functions
# ══════════════════════════════════════════════════════════════════════════════

def add_slide():
    return prs.slides.add_slide(BLANK_LAYOUT)

def rect(slide, l, t, w, h, fill_rgb=None, fill_alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE = 1
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.line.fill.background()
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    return shape

def txbox(slide, l, t, w, h):
    return slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))

def styled_text(tf, text, size, bold=False, color=C_BLACK, italic=False, align=PP_ALIGN.LEFT):
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic= italic
    run.font.color.rgb = color
    run.font.name  = FONT_BODY

def add_text(slide, text, l, t, w, h, size=12, bold=False, color=C_BLACK,
             italic=False, align=PP_ALIGN.LEFT, wrap=True):
    box = txbox(slide, l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = wrap
    styled_text(tf, text, size, bold, color, italic, align)
    return box

def add_para(tf, text, size=11, bold=False, color=C_BLACK, italic=False,
             align=PP_ALIGN.LEFT, space_before=0):
    from pptx.util import Pt as _Pt
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = _Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size  = _Pt(size)
    run.font.bold  = bold
    run.font.italic= italic
    run.font.color.rgb = color
    run.font.name  = FONT_BODY
    return p

def img(slide, path, l, t, w, h=None):
    if h:
        return slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))
    else:
        return slide.shapes.add_picture(path, Inches(l), Inches(t), width=Inches(w))

def section_banner(slide, label, color=C_ORANGE):
    """Thin orange top banner with section label."""
    rect(slide, 0, 0, 13.33, 0.18, fill_rgb=color)
    add_text(slide, label, 0.25, 0.01, 6, 0.18,
             size=7.5, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)

def slide_footer(slide, slide_num, total="18"):
    """Consistent footer line."""
    rect(slide, 0, 7.25, 13.33, 0.25, fill_rgb=C_NAVY)
    add_text(slide, "GME — Resultados Investigacion Pantallas + VW | Equipo I+D | Exceline | Confidencial | 2026-05-06",
             0.25, 7.26, 10, 0.22, size=6.5, color=C_MGRAY)
    add_text(slide, f"{slide_num} / {total}", 12.5, 7.26, 0.8, 0.22,
             size=6.5, color=C_MGRAY, align=PP_ALIGN.RIGHT)

def title_block(slide, title, subtitle=None, l=0.5, t=0.25, w=12.33):
    box = txbox(slide, l, t, w, 0.75)
    tf  = box.text_frame
    tf.word_wrap = True
    styled_text(tf, title, 22, bold=True, color=C_NAVY)
    if subtitle:
        add_para(tf, subtitle, size=11, color=C_MGRAY, italic=True, space_before=2)

def callout_box(slide, l, t, w, h, bg=C_LGRAY, border_color=C_ORANGE):
    """Colored callout box."""
    s = rect(slide, l, t, w, h, fill_rgb=bg)
    s.line.color.rgb = border_color
    s.line.width     = Pt(1.5)
    return s

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════════════════════
s1 = add_slide()
rect(s1, 0, 0, 13.33, 7.5, fill_rgb=C_NAVY)                   # full bg navy
rect(s1, 0, 0, 0.35, 7.5, fill_rgb=C_ORANGE)                  # left accent bar

add_text(s1, "EXCELINE", 0.6, 0.7, 5, 0.5,
         size=11, bold=True, color=C_ORANGE, align=PP_ALIGN.LEFT)
add_text(s1, "GME — Protector Monofasico Inteligente", 0.6, 1.15, 11, 0.5,
         size=13, bold=False, color=C_MGRAY, align=PP_ALIGN.LEFT)

# Main title
box_t = txbox(s1, 0.6, 1.75, 11.5, 2.2)
tf = box_t.text_frame
tf.word_wrap = True
styled_text(tf, "Estudio de Pantallas GME:", 30, bold=True, color=C_WHITE)
add_para(tf, "Resultados y Recomendaciones para I+D", 30, bold=False, color=C_ORANGE)

add_text(s1, "Preferencia de pantalla A vs B  |  Voz del tecnico  |  Diferenciacion R/B/M  |  Analisis Van Westendorp",
         0.6, 4.0, 11.5, 0.5, size=11, color=C_MGRAY, align=PP_ALIGN.LEFT)

# Meta row
add_text(s1, "Audiencia: Equipo I+D GME (interno)", 0.6, 4.7, 5.5, 0.35, size=9.5, color=C_MGRAY)
add_text(s1, "Encuesta: Comunidad Exceline Profesional  |  n=29  |  Abr-May 2026",
         0.6, 5.05, 8, 0.35, size=9.5, color=C_MGRAY)
add_text(s1, "Fecha: 2026-05-06  |  Version: v1  |  Confidencial", 0.6, 5.4, 7, 0.35,
         size=9.5, color=C_MGRAY)

# Footer bar
rect(s1, 0, 7.15, 13.33, 0.35, fill_rgb=RGBColor(0x0E, 0x1C, 0x2F))
add_text(s1, "Documento interno. No distribuir fuera del equipo Genteca / Exceline.",
         0.6, 7.18, 12, 0.28, size=8, color=C_MGRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ══════════════════════════════════════════════════════════════════════════════
s2 = add_slide()
rect(s2, 0, 0, 13.33, 7.5, fill_rgb=C_LGRAY)
section_banner(s2, "MARCO")
title_block(s2, "Agenda", "Estructura del deck — 18 slides")
slide_footer(s2, 2)

sections = [
    ("1", "Marco y metodologia",            "Slides 3",     "Encuesta, muestra, caveats metodologicos"),
    ("2", "Hallazgo principal: Pantalla A vs B","Slides 4-7", "Resultado 83% B, comprensibilidad, agrado, recomendacion"),
    ("3", "Voz del tecnico: comentarios abiertos","Slides 8-10","Demandas top, quotes, priorizacion v1/v2 UI"),
    ("4", "Diferenciacion R/B/M",            "Slides 11-12", "Lo que I+D hizo bien — logica real entre modos"),
    ("5", "Pricing: Van Westendorp",         "Slides 13-16", "Metodologia, curvas, segmentos, implicacion para BOM"),
    ("6", "Recomendaciones consolidadas",    "Slides 17-18", "Tabla v1 vs v2, 15 preguntas pendientes a engineering"),
    ("7", "Proximos pasos",                  "Slide 18",     "Engineering, naming, brief Vael"),
]

row_h = 0.68
for i, (num, sec, loc, desc) in enumerate(sections):
    y = 1.15 + i * row_h
    bg = C_WHITE if i % 2 == 0 else RGBColor(0xE8, 0xEC, 0xF3)
    rect(s2, 0.4, y, 12.5, row_h - 0.04, fill_rgb=bg)
    rect(s2, 0.4, y, 0.55, row_h - 0.04, fill_rgb=C_ORANGE)
    add_text(s2, num, 0.4, y + 0.12, 0.55, 0.45, size=15, bold=True,
             color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(s2, sec,  1.05, y + 0.04, 3.8, 0.35, size=11, bold=True, color=C_NAVY)
    add_text(s2, loc,  4.95, y + 0.04, 1.6, 0.35, size=9,  bold=False, color=C_BLUE)
    add_text(s2, desc, 6.6,  y + 0.04, 6.2, 0.35, size=9,  bold=False, color=C_MGRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — METODOLOGIA
# ══════════════════════════════════════════════════════════════════════════════
s3 = add_slide()
rect(s3, 0, 0, 13.33, 7.5, fill_rgb=C_WHITE)
section_banner(s3, "SECCION 1  |  MARCO")
title_block(s3, "Metodologia: encuesta digital a tecnicos Exceline Profesional",
            "Leer los caveats antes de aplicar los resultados — muestra pequena por segmento")
slide_footer(s3, 3)

# 4 method boxes
method = [
    ("Instrumento", C_BLUE,
     "Encuesta online\nComunidad Exceline Profesional\nAbr-May 2026"),
    ("Muestra", C_ORANGE,
     "n = 32 completadas\n29 validas para analisis\n3 descartadas (incompletas)"),
    ("Segmentacion", C_NAVY,
     "Refrigeracion (R): n=15, 52%\nMotores (M): n=8, 28%\nBombas (B): n=6, 20%"),
    ("Precio: Van Westendorp", C_GREEN,
     "4 preguntas PSM por respondiente\nAnalisis de interseccion de curvas\nPuntos PMC / OPP / IPP / PME"),
]

for i, (lbl, col, body) in enumerate(method):
    x = 0.4 + i * 3.2
    rect(s3, x, 1.3, 3.0, 1.6, fill_rgb=col)
    add_text(s3, lbl, x + 0.12, 1.35, 2.76, 0.42,
             size=10.5, bold=True, color=C_WHITE)
    rect(s3, x, 2.72, 3.0, 1.5, fill_rgb=C_LGRAY)
    add_text(s3, body, x + 0.12, 2.78, 2.76, 1.3,
             size=9.5, bold=False, color=C_NAVY)

# Caveat box
callout_box(s3, 0.4, 4.4, 12.5, 2.4, bg=RGBColor(0xFF, 0xF3, 0xE0), border_color=C_AMBER)
box_c = txbox(s3, 0.6, 4.5, 12.1, 2.2)
tf_c  = box_c.text_frame
tf_c.word_wrap = True
styled_text(tf_c, "Caveats metodologicos — leer antes de extrapolar resultados:", 10, bold=True, color=C_AMBER)
caveats = [
    "Muestra pequena por segmento: Bombas (n=6) y Motores (n=8) tienen intervalos de confianza amplios. Resultados son direccionales, no estadisticamente concluyentes para esos segmentos.",
    "Refrigeracion representa 52% del total — el promedio total esta sesgado hacia el segmento mas sensible al precio y con preferencias mas conservadoras.",
    "Van Westendorp mide percepcion de precio, no demanda real ni intencion de compra. No reemplaza analisis de costos BOM ni estrategia comercial.",
    "Encuesta sin ancla de precio competidor: los tecnicos respondieron sin ver precios de referencia — las respuestas VW pueden subestimar disposicion a pagar del segmento de Motores.",
]
for cv in caveats:
    add_para(tf_c, f"•  {cv}", size=8.5, color=C_BLACK, space_before=3)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — PANTALLA A vs B: VISUAL SIDE-BY-SIDE
# ══════════════════════════════════════════════════════════════════════════════
s4 = add_slide()
rect(s4, 0, 0, 13.33, 7.5, fill_rgb=C_WHITE)
section_banner(s4, "SECCION 2  |  HALLAZGO PRINCIPAL: PANTALLA A vs B")
title_block(s4, "Las dos opciones de pantalla principal evaluadas por los tecnicos",
            "Variante R-220 (Refrigeracion) — la misma logica aplica a B-220 y M-220")
slide_footer(s4, 4)

# Option A
rect(s4, 0.4, 1.2, 5.8, 0.45, fill_rgb=C_MGRAY)
add_text(s4, "OPCION A — Pantalla simple", 0.5, 1.22, 5.6, 0.4,
         size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
img_a_path = os.path.join(IMG_R, "image1.png")
img(s4, img_a_path, 0.5, 1.72, 5.5)
add_text(s4, "Muestra estado y valores medidos\nSin setpoints visibles en pantalla principal",
         0.5, 5.6, 5.5, 0.75, size=9, color=C_NAVY)

# Option B
rect(s4, 7.0, 1.2, 5.8, 0.45, fill_rgb=C_ORANGE)
add_text(s4, "OPCION B — Con setpoints visibles", 7.1, 1.22, 5.6, 0.4,
         size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
img_b_path = os.path.join(IMG_R, "image2.png")
img(s4, img_b_path, 7.1, 1.72, 5.5)
add_text(s4, "Muestra estado + valores de ajuste alto/bajo voltaje\ny rangos de sobrecarga/subcarga configurados",
         7.1, 5.6, 5.5, 0.75, size=9, color=C_NAVY)

# VS divider
add_text(s4, "vs", 6.2, 3.2, 0.9, 0.6, size=20, bold=True, color=C_ORANGE,
         align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — RESULTADO: 83% PREFIERE B
# ══════════════════════════════════════════════════════════════════════════════
s5 = add_slide()
rect(s5, 0, 0, 13.33, 7.5, fill_rgb=C_WHITE)
section_banner(s5, "SECCION 2  |  HALLAZGO PRINCIPAL: PANTALLA A vs B")
title_block(s5, "83% de los tecnicos prefiere la Pantalla B — resultado consistente en los tres segmentos",
            "24 de 29 respondientes — preferencia estadisticamente robusta a nivel total")
slide_footer(s5, 5)

# Big number
rect(s5, 0.4, 1.2, 4.0, 2.8, fill_rgb=C_ORANGE)
add_text(s5, "83%", 0.4, 1.4, 4.0, 1.6, size=72, bold=True, color=C_WHITE,
         align=PP_ALIGN.CENTER)
add_text(s5, "Prefieren\nPantalla B", 0.4, 2.95, 4.0, 0.9, size=14, bold=True,
         color=C_WHITE, align=PP_ALIGN.CENTER)
add_text(s5, "(24 de 29 respondientes)", 0.4, 3.88, 4.0, 0.35, size=9,
         color=RGBColor(0xFF, 0xD0, 0xA0), align=PP_ALIGN.CENTER)

# Segment breakdown table
cols = ["Segmento", "n", "Prefiere B", "Prefiere A", "% B"]
rows_data = [
    ("Motores (M)",       "8",  "8",  "0", "100%",  C_GREEN),
    ("Refrigeracion (R)", "15", "12", "3", "80%",   C_BLUE),
    ("Bombas (B)",        "6",  "4",  "2", "67%",   C_AMBER),
    ("TOTAL",             "29", "24", "5", "83%",   C_ORANGE),
]
col_x = [4.8, 6.5, 7.5, 8.7, 9.9]
col_w = [1.6, 0.9, 0.9, 0.9, 0.9]

rect(s5, 4.7, 1.2, 7.8, 0.45, fill_rgb=C_NAVY)
for ci, (ch, cx, cw) in enumerate(zip(cols, col_x, col_w)):
    add_text(s5, ch, cx, 1.22, cw, 0.4, size=9.5, bold=True, color=C_WHITE,
             align=PP_ALIGN.CENTER)

for ri, (seg, n, b, a, pct, col) in enumerate(rows_data):
    y = 1.7 + ri * 0.72
    bg = C_LGRAY if ri % 2 == 0 else C_WHITE
    if ri == 3:  # total row
        bg = RGBColor(0xFF, 0xF0, 0xE0)
    rect(s5, 4.7, y, 7.8, 0.68, fill_rgb=bg)
    rect(s5, 4.7, y, 0.15, 0.68, fill_rgb=col)
    add_text(s5, seg, 4.92, y + 0.14, 1.5, 0.4, size=10, bold=(ri==3), color=C_NAVY)
    for vi, (val, cx, cw) in enumerate(zip([n, b, a, pct], col_x[1:], col_w[1:])):
        add_text(s5, val, cx, y + 0.14, cw, 0.4, size=10.5 if ri==3 else 10,
                 bold=(ri==3 or vi==3), color=col if vi==3 else C_NAVY,
                 align=PP_ALIGN.CENTER)

# Recommendation box
callout_box(s5, 4.7, 4.8, 7.8, 1.5, bg=RGBColor(0xE8, 0xF5, 0xE9), border_color=C_GREEN)
add_text(s5, "RECOMENDACION PARA I+D:",
         4.9, 4.9, 7.4, 0.4, size=10, bold=True, color=C_GREEN)
add_text(s5, "Adoptar Pantalla B como pantalla principal por defecto para los tres modos (R/B/M). "
             "El resultado es unanime en Motores (8/8), mayoritario en Refrigeracion (12/15) "
             "y mayoritario en Bombas (4/6). La ventaja de B sobre A es consistente y robusta "
             "a nivel total.",
         4.9, 5.3, 7.4, 0.9, size=9.5, color=C_NAVY)

# note
add_text(s5, "Nota: Bombas n=6 — resultado direccional, confirmar con muestra mayor en fase 2.",
         4.7, 6.35, 7.8, 0.35, size=8, italic=True, color=C_MGRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — COMPRENSIBILIDAD Y AGRADO
# ══════════════════════════════════════════════════════════════════════════════
s6 = add_slide()
rect(s6, 0, 0, 13.33, 7.5, fill_rgb=C_WHITE)
section_banner(s6, "SECCION 2  |  HALLAZGO PRINCIPAL: PANTALLA A vs B")
title_block(s6, "Pantalla B tambien supera a A en comprensibilidad y agrado",
            "Mas informacion visible no implica mayor confusion — los tecnicos leen la densidad como valor")
slide_footer(s6, 6)

# Comprensibilidad
rect(s6, 0.4, 1.2, 5.9, 0.45, fill_rgb=C_BLUE)
add_text(s6, "COMPRENSIBILIDAD", 0.4, 1.22, 5.9, 0.4,
         size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

comp_data = [
    ("Pantalla A", 86, C_MGRAY),
    ("Pantalla B", 97, C_ORANGE),
]
for i, (lbl, pct, col) in enumerate(comp_data):
    y = 1.75 + i * 1.1
    add_text(s6, lbl, 0.5, y, 2.5, 0.35, size=10, bold=(i==1), color=C_NAVY)
    # bar
    rect(s6, 0.5, y + 0.38, 5.2 * pct / 100, 0.42, fill_rgb=col)
    rect(s6, 0.5, y + 0.38, 5.2, 0.42, fill_rgb=C_LGRAY)
    rect(s6, 0.5, y + 0.38, 5.2 * pct / 100, 0.42, fill_rgb=col)
    add_text(s6, f"{pct}%", 5.8, y + 0.38, 0.5, 0.42, size=11, bold=True, color=col,
             align=PP_ALIGN.LEFT)

add_text(s6, '% que calificaron "comprensible" o "muy comprensible"',
         0.5, 3.95, 5.7, 0.35, size=8, italic=True, color=C_MGRAY)

# Agrado
rect(s6, 7.0, 1.2, 5.9, 0.45, fill_rgb=C_NAVY)
add_text(s6, "AGRADO", 7.0, 1.22, 5.9, 0.4,
         size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

agrado_data = [
    ("Pantalla A", 79, "con 7% negativos", C_MGRAY),
    ("Pantalla B", 90, "sin negativos",    C_ORANGE),
]
for i, (lbl, pct, note, col) in enumerate(agrado_data):
    y = 1.75 + i * 1.1
    add_text(s6, lbl, 7.1, y, 2.5, 0.35, size=10, bold=(i==1), color=C_NAVY)
    rect(s6, 7.1, y + 0.38, 5.2, 0.42, fill_rgb=C_LGRAY)
    rect(s6, 7.1, y + 0.38, 5.2 * pct / 100, 0.42, fill_rgb=col)
    add_text(s6, f"{pct}%", 12.4, y + 0.38, 0.7, 0.42, size=11, bold=True, color=col)
    add_text(s6, note, 7.1, y + 0.82, 5.2, 0.28, size=8, italic=True, color=C_MGRAY)

add_text(s6, '% "alto" o "muy alto" agrado',
         7.1, 3.95, 5.7, 0.35, size=8, italic=True, color=C_MGRAY)

# Key insight callout
callout_box(s6, 0.4, 4.35, 12.5, 1.2, bg=RGBColor(0xE8, 0xF0, 0xFF), border_color=C_BLUE)
add_text(s6, "Lectura clave para I+D:", 0.6, 4.45, 4, 0.35, size=10, bold=True, color=C_BLUE)
add_text(s6, "La Pantalla B muestra mas datos (setpoints, rangos) que la A, pero los tecnicos la comprenden MEJOR "
             "y la encuentran mas agradable. Esto valida la hipotesis de diseno: el tecnico prefiere "
             "transparencia de parametros sobre simplicidad visual. La pantalla no debe ocultar lo que "
             "el equipo hace.",
         0.6, 4.85, 12.1, 0.65, size=9.5, color=C_NAVY)

# Slides of other mockups (remaining 4 screens)
add_text(s6, "Las otras 4 pantallas (Ajustes Perillas, Ajustes Digitales, Reporte Fallas, Configuracion) "
             "son identicas entre variantes A y B — no fueron objeto de evaluacion comparativa en la encuesta.",
         0.4, 5.65, 12.5, 0.5, size=8.5, italic=True, color=C_MGRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — PANTALLAS ADICIONALES (SCREENS 2-5)
# ══════════════════════════════════════════════════════════════════════════════
s7 = add_slide()
rect(s7, 0, 0, 13.33, 7.5, fill_rgb=C_WHITE)
section_banner(s7, "SECCION 2  |  HALLAZGO PRINCIPAL: PANTALLA A vs B")
title_block(s7, "Las 5 pantallas del GME — flujo completo de la interfaz (variante R-220)",
            "Pantallas 2-5 son comunes a todas las variantes — evaluadas cualitativamente via comentarios abiertos")
slide_footer(s7, 7)

screen_labels = [
    ("Pantalla 1B\nPrincipal (adoptada)", "image2.png", IMG_R),
    ("Pantalla 2\nAjustes Fisicos (Perillas)", "image3.png", IMG_R),
    ("Pantalla 3\nAjustes Digitales", "image4.png", IMG_R),
    ("Pantalla 4\nReporte de Fallas", "image5.png", IMG_R),
    ("Pantalla 5\nConfiguracion", "image6.png", IMG_R),
]

for i, (lbl, fname, folder) in enumerate(screen_labels):
    x = 0.3 + i * 2.58
    path = os.path.join(folder, fname)
    img(s7, path, x, 1.35, 2.3)
    add_text(s7, lbl, x, 5.55, 2.3, 0.65, size=8, color=C_NAVY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — VOZ DEL TECNICO: DEMANDAS TOP
# ══════════════════════════════════════════════════════════════════════════════
s8 = add_slide()
rect(s8, 0, 0, 13.33, 7.5, fill_rgb=C_WHITE)
section_banner(s8, "SECCION 3  |  VOZ DEL TECNICO: COMENTARIOS ABIERTOS")
title_block(s8, "Los tecnicos piden 9 features concretas — clasificadas por urgencia y esfuerzo",
            "29/29 respondientes dejaron comentario abierto — señal de engagement alto con el producto")
slide_footer(s8, 8)

demands = [
    # (label, freq_text, category, urgency_color, urgency_label)
    ("Historial cronologico de fallas con timestamp",
     ">=6 menciones", "Firmware + RTC", C_RED, "v1 critico"),
    ("Notificaciones push en tiempo real (falla/reconexion)",
     ">=3 menciones", "Cloud / App", C_AMBER, "v2"),
    ("Multivoltaje 120/220 V — un solo SKU",
     "2 menciones", "Hardware", C_AMBER, "v2"),
    ("Tiempo restante de reconexion visible en pantalla",
     "2 menciones", "UI / Firmware", C_RED, "v1"),
    ("Rango de voltaje hasta 265 V (sobre-voltaje cronico VE)",
     "1 mencion explícita", "Hardware / Firmware", C_AMBER, "v1-evaluar"),
    ("Picos de corriente arranque vs regimen (diferenciacion)",
     "1 mencion", "Firmware (ADC logic)", C_AMBER, "v2"),
    ("Sensibilidad de desconexion <1 segundo",
     "1 mencion (refrigeracion)", "Firmware / ADC", C_RED, "v1 critico"),
    ("Wording: 'Ajustes Fisicos' -> 'Ajustes de Perillas'",
     "1 mencion explícita", "UI (cosmetico)", C_GREEN, "v1 facil"),
    ("NFC como alternativa de pairing / lectura",
     "1 mencion", "Hardware", C_AMBER, "v2"),
]

col_hdrs = ["Demanda del tecnico", "Frecuencia", "Tipo de cambio", "Version target"]
col_x2   = [0.4, 5.8, 7.6, 9.8, 11.4]
col_w2   = [5.3, 1.7, 2.1, 1.5, 1.5]

rect(s8, 0.4, 1.2, 12.5, 0.4, fill_rgb=C_NAVY)
for ci, (ch, cx, cw) in enumerate(zip(col_hdrs, col_x2, col_w2)):
    add_text(s8, ch, cx, 1.22, cw, 0.36, size=8.5, bold=True, color=C_WHITE,
             align=PP_ALIGN.CENTER)

row_h2 = 0.56
for ri, (dem, freq, cat, col, ver) in enumerate(demands):
    y = 1.62 + ri * row_h2
    bg = C_LGRAY if ri % 2 == 0 else C_WHITE
    rect(s8, 0.4, y, 12.5, row_h2 - 0.03, fill_rgb=bg)
    rect(s8, 0.4, y, 0.12, row_h2 - 0.03, fill_rgb=col)
    add_text(s8, dem,  0.55, y + 0.08, 5.1, 0.42, size=8.5, color=C_NAVY)
    add_text(s8, freq, 5.85, y + 0.08, 1.6, 0.42, size=8,   color=C_NAVY, align=PP_ALIGN.CENTER)
    add_text(s8, cat,  7.65, y + 0.08, 2.0, 0.42, size=8,   color=C_NAVY)
    # version badge
    rect(s8, 11.45, y + 0.08, 1.3, 0.34, fill_rgb=col)
    add_text(s8, ver, 11.45, y + 0.09, 1.3, 0.3, size=7.5, bold=True,
             color=C_WHITE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — QUOTES TEXTUALES
# ══════════════════════════════════════════════════════════════════════════════
s9 = add_slide()
rect(s9, 0, 0, 13.33, 7.5, fill_rgb=C_NAVY)
section_banner(s9, "SECCION 3  |  VOZ DEL TECNICO: COMENTARIOS ABIERTOS")
slide_footer(s9, 9)

add_text(s9, "Lo que dijeron los tecnicos — en sus palabras",
         0.5, 0.3, 12, 0.6, size=20, bold=True, color=C_WHITE)

quotes = [
    (
        '"En electronica un ESP32 puede costar $10 siendo programable y con conectividad. '
        'Un producto masivo que incorpore control y monitoreo remoto no tiene por que elevar '
        'excesivamente el precio."',
        "Tecnico de Refrigeracion — respondiente #104356020",
        "Implicacion: la expectativa de precio es acida — I+D debe saber que el tecnico tiene "
        "un benchmark de componente en $10. La defensa del precio lista ($35) requiere "
        "guion de valor tecnico verificable.",
        C_ORANGE
    ),
    (
        '"Seria muy agradable que puedan registrar un historial de las ultimas 10 a 15 fallas. '
        'Este aplicativo tiene mucho potencial. Pueden agregarlo por contacto NFC."',
        "Ing. Asdrúbal García, Tec. Refrigeracion/Climatizacion — respondiente #104362438",
        "Implicacion: el historial cronologico es la demanda funcional #1. NFC es mencionado "
        "como canal alternativo de acceso — requiere antena hardware.",
        C_BLUE
    ),
    (
        ('"La sensibilidad de desconexion es muy importante: una fluctuacion de menos de 1 segundo'
         'deja sin efecto la proteccion. Esto se incrementa en importancia en refrigeracion'
         'y aire acondicionado."'),
        "Tecnico de Refrigeracion — respondiente #104359489",
        "Implicacion: gap tecnico urgente para modo R. Ver Pregunta Q1 al equipo de engineering "
        "(sampling rate del ADC). Esta es una condicion de credibilidad del producto en el segmento "
        "mas grande.",
        C_RED
    ),
]

for i, (q, attr, impl, col) in enumerate(quotes):
    y = 1.05 + i * 2.0
    rect(s9, 0.4, y, 0.08, 1.55, fill_rgb=col)
    add_text(s9, q,    0.6, y,       10.5, 0.9,  size=10, italic=True, color=C_WHITE)
    add_text(s9, f"— {attr}", 0.6, y + 0.92, 7, 0.3, size=8, color=col, italic=True)
    add_text(s9, impl, 8.5, y,       4.4,  1.2,  size=8, color=C_MGRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — PRIORIZACION v1 vs v2 (UI/UX)
# ══════════════════════════════════════════════════════════════════════════════
s10 = add_slide()
rect(s10, 0, 0, 13.33, 7.5, fill_rgb=C_WHITE)
section_banner(s10, "SECCION 3  |  VOZ DEL TECNICO: COMENTARIOS ABIERTOS")
title_block(s10, "Priorizacion de demandas: v1 (octubre 2026) vs v2 (post-lanzamiento)",
            "Criterio: impacto en credibilidad del producto para el segmento dominante (Refrigeracion) + esfuerzo de implementacion")
slide_footer(s10, 10)

# v1 column
rect(s10, 0.4, 1.2, 6.0, 0.45, fill_rgb=C_GREEN)
add_text(s10, "v1 — Lanzamiento octubre 2026", 0.4, 1.22, 6.0, 0.4,
         size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

v1_items = [
    ("Pantalla B como default", "Ninguno — decision de configuracion"),
    ("Fix wording: 'Ajustes de Perillas'", "UI cosmetico — 1 dia de dev"),
    ("Tiempo restante de reconexion visible", "UI / firmware — mostrar countdown en pantalla"),
    ("Validar sensibilidad <1 s en modo R", "Pregunta Q1 a engineering (ADC sampling rate)"),
    ("Log cronologico minimo (si hay RTC)", "Requiere confirmar RTC — ver Pregunta Q5"),
]
for ri, (item, note) in enumerate(v1_items):
    y = 1.72 + ri * 0.75
    bg = C_LGRAY if ri % 2 == 0 else C_WHITE
    rect(s10, 0.4, y, 6.0, 0.71, fill_rgb=bg)
    rect(s10, 0.4, y, 0.12, 0.71, fill_rgb=C_GREEN)
    add_text(s10, item, 0.58, y + 0.04, 5.68, 0.35, size=9.5, bold=True, color=C_NAVY)
    add_text(s10, note, 0.58, y + 0.37, 5.68, 0.28, size=8,   color=C_MGRAY, italic=True)

# v2 column
rect(s10, 7.0, 1.2, 5.9, 0.45, fill_rgb=C_BLUE)
add_text(s10, "v2 — Post-lanzamiento", 7.0, 1.22, 5.9, 0.4,
         size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

v2_items = [
    ("NFC pairing / lectura", "Requiere antena hardware — confirmar con Q13"),
    ("Multivoltaje 120/220 V", "Requiere variante PCB — ver Q10 a engineering"),
    ("Notificaciones push (cloud)", "Requiere arquitectura cloud/relay — fuera de alcance v1"),
    ("Log expandido: 25+ eventos cronologicos", "Requiere RTC confirmado + buffer firmware"),
    ("Rango voltaje hasta 265 V", "Hardware/firmware — evaluar si requiere cambio circuito"),
]
for ri, (item, note) in enumerate(v2_items):
    y = 1.72 + ri * 0.75
    bg = C_LGRAY if ri % 2 == 0 else C_WHITE
    rect(s10, 7.0, y, 5.9, 0.71, fill_rgb=bg)
    rect(s10, 7.0, y, 0.12, 0.71, fill_rgb=C_BLUE)
    add_text(s10, item, 7.18, y + 0.04, 5.58, 0.35, size=9.5, bold=True, color=C_NAVY)
    add_text(s10, note, 7.18, y + 0.37, 5.58, 0.28, size=8,   color=C_MGRAY, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — DIFERENCIACION R/B/M: LO QUE I+D HIZO BIEN
# ══════════════════════════════════════════════════════════════════════════════
s11 = add_slide()
rect(s11, 0, 0, 13.33, 7.5, fill_rgb=C_WHITE)
section_banner(s11, "SECCION 4  |  DIFERENCIACION R/B/M")
title_block(s11, "Los tres modos tienen diferencias reales — no solo cambio de defaults",
            "Hallazgo de validacion tecnica (Vera, 2026-05-06) — sustenta el claim 'tres modos en un protector'")
slide_footer(s11, 11)

# 3 mode columns
modes = [
    ("GME-R220\nRefrigeracion", C_BLUE, [
        ("Tiempo reconexion",    "300 s  (rango 180-600 s)"),
        ("Subcarga %",           "60% I-nom, t=10 s"),
        ("Tiempo reintento",     "240 s"),
        ("Maniobra diaria",      "No aplica"),
        ("Tiempo detec. sobrecarga", "No visible en mockup"),
        ("Alineacion normativa", "Copeland AE4-1365\n(igualacion presiones)"),
        ("Diferenciador clave",  "Tiempos conservadores\npara compresores hermeticos"),
    ]),
    ("GME-B220\nBombas", C_ORANGE, [
        ("Tiempo reconexion",    "120 s  (rango 5-300 s)"),
        ("Subcarga %",           "70% I-nom, t=5 s"),
        ("Tiempo reintento",     "20 min"),
        ("Maniobra diaria",      "ACTIVADA — logica exclusiva de modo B"),
        ("Tiempo detec. sobrecarga", "No visible en mockup"),
        ("Alineacion normativa", "Practica dry-run bombas\n(Franklin SubMonitor ref.)"),
        ("Diferenciador clave",  "Maniobra diaria — unica\nfuncion exclusiva de modo"),
    ]),
    ("GME-M220\nMotores", C_NAVY, [
        ("Tiempo reconexion",    "60 s  (rango 5-300 s)"),
        ("Subcarga %",           "60% I-nom, t=5 s"),
        ("Tiempo reintento",     "240 s"),
        ("Maniobra diaria",      "No aplica"),
        ("Tiempo detec. sobrecarga", "0.1 - 5 s ajustable\n(visible y configurable en modo M)"),
        ("Alineacion normativa", "Logica Class 10 / Class 10A\n(pendiente confirmar)"),
        ("Diferenciador clave",  "Tiempo deteccion sobrecarga\najustable — critico rotor trancado"),
    ]),
]

for mi, (title, col, params) in enumerate(modes):
    x = 0.4 + mi * 4.3
    rect(s11, x, 1.2, 4.1, 0.55, fill_rgb=col)
    add_text(s11, title, x, 1.22, 4.1, 0.5, size=10.5, bold=True,
             color=C_WHITE, align=PP_ALIGN.CENTER)
    for pi, (param, val) in enumerate(params):
        y = 1.8 + pi * 0.71
        bg = RGBColor(0xE8, 0xEC, 0xF5) if pi % 2 == 0 else C_WHITE
        rect(s11, x, y, 4.1, 0.68, fill_rgb=bg)
        add_text(s11, param, x + 0.1, y + 0.03, 1.4, 0.3, size=7.5, bold=True, color=C_MGRAY)
        add_text(s11, val,   x + 0.1, y + 0.33, 3.8, 0.3, size=8.5, color=C_NAVY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — DIFERENCIACION R/B/M: SINTESIS
# ══════════════════════════════════════════════════════════════════════════════
s12 = add_slide()
rect(s12, 0, 0, 13.33, 7.5, fill_rgb=C_LGRAY)
section_banner(s12, "SECCION 4  |  DIFERENCIACION R/B/M")
title_block(s12, "Sintesis: la diferenciacion entre modos tiene sustancia tecnica suficiente para el claim",
            "Con una condicion: engineering debe confirmar dos puntos antes del lanzamiento")
slide_footer(s12, 12)

findings = [
    (C_GREEN, "CONFIRMADO",
     "Valores default distintos en todos los parametros de tiempo y umbral",
     "Esto solo ya justifica 'configurado para cada aplicacion'. Directo de mockups."),
    (C_GREEN, "CONFIRMADO",
     "Maniobra diaria es logica exclusiva del modo Bombas (GME-B220)",
     "No aparece en R ni M. Es la unica funcion de lógica diferente confirmada en mockups."),
    (C_GREEN, "CONFIRMADO",
     "Tiempo de deteccion de sobrecarga configurable solo visible en GME-M220",
     "0.1-5 s ajustable. Critico para proteccion contra rotor trancado. No visible en R/B."),
    (C_AMBER, "CONFIRMAR CON ENGINEERING",
     "¿Es el tiempo de deteccion de sobrecarga un parametro oculto en R y B, o esta fijo?",
     "Si esta fijo en R/B pero configurable en M, el claim 'tres modos' es mas fuerte."),
    (C_AMBER, "CONFIRMAR CON ENGINEERING",
     "¿La maniobra diaria implica logica de firmware diferente o solo un toggle de scheduling?",
     "Si es logica exclusiva (ciclo de test programado), el claim tiene profundidad tecnica real."),
]

for ri, (col, badge, finding, note) in enumerate(findings):
    y = 1.25 + ri * 1.0
    rect(s12, 0.4, y, 12.5, 0.95, fill_rgb=C_WHITE)
    rect(s12, 0.4, y, 0.12, 0.95, fill_rgb=col)
    rect(s12, 0.55, y + 0.08, 1.5, 0.32, fill_rgb=col)
    add_text(s12, badge, 0.55, y + 0.08, 1.5, 0.32, size=7.5, bold=True,
             color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(s12, finding, 2.15, y + 0.05, 10.6, 0.38, size=10, bold=True, color=C_NAVY)
    add_text(s12, note,    2.15, y + 0.45, 10.6, 0.38, size=8.5, color=C_MGRAY, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — VAN WESTENDORP: METODOLOGIA
# ══════════════════════════════════════════════════════════════════════════════
s13 = add_slide()
rect(s13, 0, 0, 13.33, 7.5, fill_rgb=C_WHITE)
section_banner(s13, "SECCION 5  |  PRICING: VAN WESTENDORP")
title_block(s13, "Metodologia Van Westendorp: cuatro preguntas, cuatro curvas, cuatro puntos clave",
            "PSM (Price Sensitivity Meter, 1976) — estandar para estudios de percepcion de precio en nuevos productos")
slide_footer(s13, 13)

# 4 questions
qs = [
    ("Q10", "Demasiado barato",
     '"A que precio le parece TAN BARATO que no confiaria en su calidad?"',
     "Curva descendente — TC", C_RED),
    ("Q11", "Economico / aceptable",
     '"A que precio le parece economico y confiable?"',
     "Curva descendente — C", C_GREEN),
    ("Q12", "Costoso pero aceptable",
     '"A que precio le parece costoso pero aun lo compraria?"',
     "Curva ascendente — E", C_AMBER),
    ("Q13", "Excesivamente costoso",
     '"A que precio le parece excesivamente costoso que NO lo compraria?"',
     "Curva ascendente — TE", C_BLUE),
]
for i, (qn, lbl, q, curve, col) in enumerate(qs):
    x = 0.4 + i * 3.2
    rect(s13, x, 1.2, 3.0, 0.45, fill_rgb=col)
    add_text(s13, f"{qn}  |  {lbl}", x + 0.1, 1.22, 2.8, 0.4,
             size=9.5, bold=True, color=C_WHITE)
    rect(s13, x, 1.68, 3.0, 1.5, fill_rgb=C_LGRAY)
    add_text(s13, q, x + 0.1, 1.72, 2.8, 0.9, size=8.5, italic=True, color=C_NAVY)
    add_text(s13, curve, x + 0.1, 2.65, 2.8, 0.35, size=8, bold=True, color=col)

# 4 key points
points = [
    ("PMC", "Point of Marginal Cheapness",
     "Cota inferior del rango aceptable\nInterseccion: TC ∩ E\n\"Por debajo, dudas de calidad dominan\"",
     C_RED),
    ("OPP", "Optimal Price Point",
     "Precio de menor resistencia psicologica\nInterseccion: TC ∩ TE\n\"Mismo % muy-barato que muy-caro\"",
     C_GREEN),
    ("IPP", "Indifference Price Point",
     "Precio 'justo' para el promedio\nInterseccion: C ∩ E\n\"Mismo % economico que costoso\"",
     C_AMBER),
    ("PME", "Point of Marginal Expensiveness",
     "Cota superior del rango aceptable\nInterseccion: C ∩ TE\n\"Por encima, inasequible domina\"",
     C_BLUE),
]
rect(s13, 0.4, 3.35, 12.5, 0.35, fill_rgb=C_NAVY)
add_text(s13, "Los 4 puntos clave que definen el Rango Aceptable de Precios (RAP)",
         0.5, 3.36, 12, 0.33, size=9.5, bold=True, color=C_WHITE)

for i, (code, name, desc, col) in enumerate(points):
    x = 0.4 + i * 3.2
    rect(s13, x, 3.73, 3.0, 0.55, fill_rgb=col)
    add_text(s13, f"{code} — {name}", x + 0.1, 3.75, 2.8, 0.5,
             size=9, bold=True, color=C_WHITE)
    rect(s13, x, 4.3, 3.0, 1.9, fill_rgb=C_LGRAY)
    add_text(s13, desc, x + 0.1, 4.35, 2.8, 1.75, size=8.5, color=C_NAVY)

# caveat
add_text(s13, "Caveat: VW mide percepcion de precio, no demanda real ni intencion de compra. "
              "Los resultados son input para la decision de precio — no la decision en si misma.",
         0.4, 6.3, 12.5, 0.5, size=8, italic=True, color=C_MGRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — CURVA VW TOTAL
# ══════════════════════════════════════════════════════════════════════════════
s14 = add_slide()
rect(s14, 0, 0, 13.33, 7.5, fill_rgb=C_WHITE)
section_banner(s14, "SECCION 5  |  PRICING: VAN WESTENDORP")
title_block(s14, "Resultado total (n=29): rango aceptable USD 29.90 – 40.33, precio optimo USD 35",
            "OPP e IPP casi coinciden en $35 — mercado con expectativa de precio madura y consistente")
slide_footer(s14, 14)

# VW total chart
vw_total_path = os.path.join(IMG_VW, "vw_total.png")
img(s14, vw_total_path, 0.4, 1.15, 7.8)

# Key numbers panel
kn_data = [
    ("PMC", "$29.90", "Cota inferior\nrango aceptable", C_RED),
    ("OPP", "$35.00", "Precio optimo\n(min. resistencia)", C_GREEN),
    ("IPP", "$35.50", "Precio 'justo'\n(indiferencia)", C_AMBER),
    ("PME", "$40.33", "Techo del\nrango aceptable", C_BLUE),
]
for i, (lbl, val, desc, col) in enumerate(kn_data):
    x = 8.5
    y = 1.2 + i * 1.4
    rect(s14, x, y, 4.4, 1.3, fill_rgb=C_LGRAY)
    rect(s14, x, y, 0.12, 1.3, fill_rgb=col)
    add_text(s14, lbl,  x + 0.2, y + 0.05, 1.2, 0.38, size=10, bold=True, color=col)
    add_text(s14, val,  x + 1.4, y + 0.05, 2.8, 0.5,  size=20, bold=True, color=col)
    add_text(s14, desc, x + 0.2, y + 0.7,  3.9, 0.5,  size=8.5, color=C_NAVY)

# RAP band highlight
callout_box(s14, 8.5, 6.85, 4.4, 0.5, bg=RGBColor(0xE8, 0xF5, 0xE9), border_color=C_GREEN)
add_text(s14, "RAP (Rango Aceptable de Precios): $29.90 – $40.33",
         8.6, 6.88, 4.2, 0.4, size=9.5, bold=True, color=C_GREEN)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — VW POR SEGMENTO
# ══════════════════════════════════════════════════════════════════════════════
s15 = add_slide()
rect(s15, 0, 0, 13.33, 7.5, fill_rgb=C_WHITE)
section_banner(s15, "SECCION 5  |  PRICING: VAN WESTENDORP")
title_block(s15, "Curvas por segmento: Motores tolera mas, Refrigeracion es el constraint critico",
            "El segmento de Motores es el mas premium; Refrigeracion fija el techo practico para un SKU unico")
slide_footer(s15, 15)

seg_charts = [
    ("Motores (n=8)", "vw_motores.png",    C_NAVY,   "PMC $35.3  |  OPP $35.5\nIPP $40.5  |  PME $45.5"),
    ("Bombas (n=6)",  "vw_bombas.png",     C_ORANGE, "PMC $30.0  |  OPP $30.5\nIPP $40.0  |  PME $50.0\n[n=6: solo direccional]"),
    ("Refrigeracion (n=15)", "vw_refrigeracion.png", C_BLUE, "PMC $29.7  |  OPP $30.2\nIPP $30.5  |  PME $39.7"),
]

for i, (title, fname, col, kns) in enumerate(seg_charts):
    x = 0.3 + i * 4.35
    rect(s15, x, 1.2, 4.1, 0.38, fill_rgb=col)
    add_text(s15, title, x, 1.22, 4.1, 0.34, size=10, bold=True,
             color=C_WHITE, align=PP_ALIGN.CENTER)
    path = os.path.join(IMG_VW, fname)
    img(s15, path, x, 1.62, 4.1)
    rect(s15, x, 5.6, 4.1, 1.25, fill_rgb=C_LGRAY)
    add_text(s15, kns, x + 0.15, 5.68, 3.8, 1.1, size=9, color=C_NAVY)

add_text(s15, "Constraint critico para SKU unico: el techo de Refrigeracion ($39.7) fija el precio maximo "
              "antes de perder aceptacion en el 52% de la muestra. Un precio lista de $35 cae dentro del "
              "rango aceptable de los TRES segmentos simultaneamente.",
         0.4, 6.95, 12.5, 0.42, size=8.5, italic=True, color=C_NAVY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — IMPLICACION PARA BOM / I+D
# ══════════════════════════════════════════════════════════════════════════════
s16 = add_slide()
rect(s16, 0, 0, 13.33, 7.5, fill_rgb=C_WHITE)
section_banner(s16, "SECCION 5  |  PRICING: VAN WESTENDORP")
title_block(s16, "Recomendacion de precio: USD 35 lista — y lo que eso implica para el BOM",
            "El analisis VW informa a I+D el espacio de precio disponible — I+D define si el BOM encaja")
slide_footer(s16, 16)

# Pricing recommendation box
rect(s16, 0.4, 1.2, 6.0, 3.0, fill_rgb=C_NAVY)
add_text(s16, "RECOMENDACION DE PRECIO", 0.6, 1.3, 5.6, 0.45,
         size=10, bold=True, color=C_ORANGE)
add_text(s16, "USD 35", 0.6, 1.78, 5.6, 1.1,
         size=52, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text(s16, "Precio lista sugerido", 0.6, 2.88, 5.6, 0.35,
         size=10, color=C_MGRAY, align=PP_ALIGN.CENTER)
add_text(s16, "Floor: USD 30   |   Ceiling: USD 40", 0.6, 3.3, 5.6, 0.35,
         size=9.5, color=C_ORANGE, align=PP_ALIGN.CENTER)
add_text(s16, "Escenario A: SKU unico, misma unidad hardware\nconfigurable por firmware R/B/M",
         0.6, 3.72, 5.6, 0.42, size=8.5, color=C_MGRAY, align=PP_ALIGN.CENTER)

# BOM implications
rect(s16, 7.0, 1.2, 5.9, 0.38, fill_rgb=C_ORANGE)
add_text(s16, "Implicaciones para I+D / BOM", 7.0, 1.22, 5.9, 0.34,
         size=10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

bom_items = [
    ("$35 lista implica costo de fabrica objetivo de ~$15-20",
     "Asumiendo margen distribuidor 30-40% + margen Exceline. "
     "El BOM + mano de obra + certificaciones debe caber en ese rango."),
    ("Cada feature adicional de v1 tiene impacto en BOM",
     "RTC (si se incluye para log cronologico): ~$0.5-2 adicional en BOM. "
     "Antena NFC (si se incluye): ~$1-3. La decision de scope v1 vs v2 tiene costo real."),
    ("El segmento de Motores tolera hasta $45.50",
     "Si el hardware tiene variante con mayor capacidad de corriente para motores industriales, "
     "hay headroom de precio para ese SKU sin perder aceptacion."),
    ("Multivoltaje 120/220 V elevaria precio pero expandiria TAM",
     "Un solo SKU 120/220 elimina la decision de compra al tecnico. "
     "Impacto en BOM a confirmar con engineering (Q10)."),
]
for ri, (hdr, body) in enumerate(bom_items):
    y = 1.65 + ri * 1.2
    rect(s16, 7.0, y, 5.9, 1.15, fill_rgb=C_LGRAY if ri % 2 == 0 else C_WHITE)
    rect(s16, 7.0, y, 0.1, 1.15, fill_rgb=C_ORANGE)
    add_text(s16, hdr,  7.15, y + 0.05, 5.65, 0.38, size=9,   bold=True, color=C_NAVY)
    add_text(s16, body, 7.15, y + 0.45, 5.65, 0.65, size=8,   color=C_MGRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — RECOMENDACIONES CONSOLIDADAS v1 vs v2
# ══════════════════════════════════════════════════════════════════════════════
s17 = add_slide()
rect(s17, 0, 0, 13.33, 7.5, fill_rgb=C_WHITE)
section_banner(s17, "SECCION 6  |  RECOMENDACIONES CONSOLIDADAS PARA I+D")
title_block(s17, "Tabla de decisiones para I+D: lo que va en v1 y lo que va en v2",
            "Basada en datos de encuesta + validacion tecnica Vera + analisis de pricing VW")
slide_footer(s17, 17)

# Column headers
rect(s17, 0.4, 1.2, 5.8, 0.45, fill_rgb=C_GREEN)
add_text(s17, "v1 — Lanzamiento octubre 2026", 0.4, 1.22, 5.8, 0.4,
         size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
rect(s17, 6.5, 1.2, 6.4, 0.45, fill_rgb=C_BLUE)
add_text(s17, "v2 — Post-lanzamiento (roadmap)", 6.5, 1.22, 6.4, 0.4,
         size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

v1_recs = [
    ("Pantalla B", "Adoptar como default los tres modos (R/B/M).\nResultado: 83% total, 100% Motores."),
    ("Fix wording UI", "Cambiar 'Ajustes Fisicos' por 'Ajustes de Perillas'\n(solicitado directamente por tecnico, consistencia con barra inferior)."),
    ("Countdown reconexion", "Mostrar tiempo restante de reconexion en pantalla principal.\nMejorar experiencia en campo — tecnico no llama pensando que fallo."),
    ("Validar sens. <1s modo R", "Confirmar con engineering: ADC sampling rate suficiente\npara detectar fluctuaciones <1 s (Q1). Critico para credibilidad en refrigeracion."),
    ("Log basico si hay RTC", "Si el hardware incluye RTC: implementar log de min. 5 eventos\ncon timestamp en v1. Si no hay RTC, diferir a v2 con RTC externo."),
]

v2_recs = [
    ("NFC pairing", "Lectura de configuracion y diagnostico por NFC.\nRequiere hardware: antena + controlador NFC (Q13)."),
    ("Multivoltaje 120/220 V", "Un solo SKU para ambas tensiones.\nRequiere variante de PCB — impacto BOM a cuantificar (Q10)."),
    ("Notificaciones push", "Alertas en tiempo real via cloud/relay.\nRequiere arquitectura backend — fuera de scope v1."),
    ("Log cronologico 25+ eventos", "Historial completo con timestamp real.\nRequiere RTC confirmado + ampliacion buffer firmware (Q4, Q5)."),
    ("Rango hasta 265 V", "Cubrir zonas VE con sobre-voltaje cronico.\nEvaluar si requiere cambio en circuito de medicion o solo calibracion."),
]

row_h3 = 0.94
for ri in range(5):
    y = 1.7 + ri * row_h3
    bg = C_LGRAY if ri % 2 == 0 else C_WHITE

    # v1
    rect(s17, 0.4, y, 5.8, row_h3 - 0.04, fill_rgb=bg)
    rect(s17, 0.4, y, 0.1, row_h3 - 0.04, fill_rgb=C_GREEN)
    add_text(s17, v1_recs[ri][0], 0.55, y + 0.04, 1.4, 0.36, size=9, bold=True, color=C_GREEN)
    add_text(s17, v1_recs[ri][1], 0.55, y + 0.4,  5.5, 0.46, size=8, color=C_NAVY)

    # v2
    rect(s17, 6.5, y, 6.4, row_h3 - 0.04, fill_rgb=bg)
    rect(s17, 6.5, y, 0.1, row_h3 - 0.04, fill_rgb=C_BLUE)
    add_text(s17, v2_recs[ri][0], 6.65, y + 0.04, 1.6, 0.36, size=9, bold=True, color=C_BLUE)
    add_text(s17, v2_recs[ri][1], 6.65, y + 0.4,  6.1, 0.46, size=8, color=C_NAVY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — 15 PREGUNTAS A ENGINEERING + PROXIMOS PASOS
# ══════════════════════════════════════════════════════════════════════════════
s18 = add_slide()
rect(s18, 0, 0, 13.33, 7.5, fill_rgb=C_WHITE)
section_banner(s18, "SECCION 6+7  |  PREGUNTAS A ENGINEERING Y PROXIMOS PASOS")
title_block(s18, "15 preguntas pendientes a engineering determinan alcance real del producto",
            "Archivo completo: Preguntas_Engineering_GME_15.md — estas son las 5 prioritarias para v1")
slide_footer(s18, 18)

# Priority 5 questions
pqs = [
    ("Q1", "Sensibilidad <1 s",
     "¿A que frecuencia muestrea el ADC? ¿Puede detectar fluctuaciones de voltaje de duracion <1 s?",
     "Credibilidad en Refrigeracion", C_RED),
    ("Q2", "Operacion offline",
     "¿El GME protege activamente si no hay WiFi activo? ¿La logica de proteccion es independiente del webserver?",
     "Gate del claim 'sin cloud'", C_RED),
    ("Q4+Q5", "Log y RTC",
     "¿Existe buffer circular de eventos en firmware? ¿Hay RTC con bateria de respaldo o sincronizacion NTP?",
     "Feature mas solicitada por tecnicos", C_AMBER),
    ("Q9", "Corriente maxima contactor",
     "¿Cuál es la corriente maxima conmutable en AC-3? ¿40 A en mockup es limite real o valor de ejemplo?",
     "Define claims de capacidad del producto", C_AMBER),
    ("Q15", "Logica 3 intentos",
     "¿El firmware reintenta 3 veces antes de bloquear? ¿Se puede configurar? ¿Como se comunica al tecnico?",
     "Riesgo de responsabilidad y confianza", C_AMBER),
]

rect(s18, 0.4, 1.2, 9.0, 0.38, fill_rgb=C_NAVY)
add_text(s18, "5 preguntas prioritarias para v1 (ver archivo completo para las 15)", 0.55, 1.22, 8.7, 0.34,
         size=9.5, bold=True, color=C_WHITE)
add_text(s18, "Urgencia", 9.5, 1.22, 1.4, 0.34, size=9.5, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

for ri, (qn, title, body, implication, col) in enumerate(pqs):
    y = 1.62 + ri * 0.98
    bg = C_LGRAY if ri % 2 == 0 else C_WHITE
    rect(s18, 0.4, y, 9.0, 0.94, fill_rgb=bg)
    rect(s18, 0.4, y, 0.55, 0.94, fill_rgb=col)
    add_text(s18, qn, 0.4, y + 0.26, 0.55, 0.4, size=9.5, bold=True,
             color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(s18, title, 1.02, y + 0.04, 2.2, 0.38, size=10, bold=True, color=C_NAVY)
    add_text(s18, body,  1.02, y + 0.44, 5.5, 0.42, size=8.5, color=C_NAVY)
    add_text(s18, f"Impacto: {implication}", 6.6, y + 0.04, 2.6, 0.84,
             size=8, color=C_MGRAY, italic=True)
    rect(s18, 9.5, y, 1.4, 0.94, fill_rgb=col)
    add_text(s18, "ALTA" if col == C_RED else "MEDIA",
             9.5, y + 0.3, 1.4, 0.38, size=10, bold=True,
             color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(s18, "Bloquea claim" if col == C_RED else "v1 scope",
             9.5, y + 0.62, 1.4, 0.28, size=7, color=C_WHITE, align=PP_ALIGN.CENTER)

# Next steps
rect(s18, 0.4, 6.72, 12.5, 0.45, fill_rgb=C_LGRAY)
add_text(s18, "Proximos pasos:", 0.55, 6.74, 2, 0.4, size=9, bold=True, color=C_NAVY)
add_text(s18,
         "1. Engineering responde 15 preguntas (Preguntas_Engineering_GME_15.md)  |  "
         "2. Focus groups de naming GME (sesiones previas en carpeta investigacion)  |  "
         "3. Brief Vael para arquitectura de mensaje del lanzamiento",
         2.6, 6.74, 10.1, 0.4, size=8.5, color=C_NAVY)

# ══════════════════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════════════════
prs.save(OUT)
print(f"OK: {OUT}")
print(f"Slides: {len(prs.slides)}")
