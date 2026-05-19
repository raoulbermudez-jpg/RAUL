#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Build Word (.docx) + PPTX from deck v3 focused markdown SSOT.

Per `feedback_word_first_pptx_after_policy`:
- Word is the primary deliverable (reviewable, easy to comment).
- PPTX is generated from the same SSOT as visual companion.

Per `feedback_python_windows_encoding`:
- stdout reconfigured to utf-8 to avoid cp1252 errors.

Per Vivienne runtime guardrails (codified in agent runtime 2026-05-18):
- Bullets max 4 per slide (we respect what's in the markdown SSOT).
- Sober brand kit (Gama V0.1 candidate) for lawyer audience.

Token explosion guard: this script does NOT invoke Vivienne LLM; it does
deterministic markdown->docx/pptx transformation. No risk of >32K explosion.
"""
import sys
import re
from pathlib import Path

sys.stdout.reconfigure(encoding="utf-8")

from docx import Document
from docx.shared import Pt, RGBColor, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as P_Inches, Pt as P_Pt
from pptx.dml.color import RGBColor as P_RGBColor

BASE = Path(r"C:\RAUL\03-projects\genteca\portfolio-naming-ip-2026\06-three-pack")
MD = BASE / "03-Deck-lawyer-presentation_v3_focus-thermo-safe-flickerguard.md"
DOCX = BASE / "03-Deck-lawyer-presentation_v3_focus-thermo-safe-flickerguard.docx"
PPTX = BASE / "03-Deck-lawyer-presentation_v3_focus-thermo-safe-flickerguard.pptx"

# ---------- Parse markdown SSOT into slide structure ----------

def parse_slides(md_text: str):
    """
    Parse the deck markdown SSOT.
    Returns: (front_matter_dict, slides_list)
    Each slide is a dict with keys: title, body_raw, speaker_note_raw,
        visual_descriptor_raw
    """
    # Strip frontmatter (between first two --- lines)
    fm = {}
    if md_text.startswith("---"):
        end_fm = md_text.find("\n---\n", 4)
        if end_fm > 0:
            fm_text = md_text[4:end_fm]
            for line in fm_text.split("\n"):
                if ":" in line and not line.lstrip().startswith("-"):
                    k, _, v = line.partition(":")
                    fm[k.strip()] = v.strip().strip('"')
            md_text = md_text[end_fm + 5 :]

    # Split into slides by "## Slide N — " heading
    slide_pattern = re.compile(r"^## Slide (\d+) — (.+)$", re.MULTILINE)
    splits = []
    for m in slide_pattern.finditer(md_text):
        splits.append((m.start(), int(m.group(1)), m.group(2).strip()))

    slides = []
    for i, (start, num, title) in enumerate(splits):
        end = splits[i + 1][0] if i + 1 < len(splits) else len(md_text)
        block = md_text[start:end]

        # Extract Body / Speaker note / Visual descriptor sections
        body = _extract_section(block, "Body:")
        spk = _extract_section(block, "Speaker note:")
        vis = _extract_section(block, "Visual descriptor:")

        slides.append(
            {
                "num": num,
                "title": title,
                "body": body,
                "speaker_note": spk,
                "visual_descriptor": vis,
            }
        )
    return fm, slides


def _extract_section(block: str, marker: str) -> str:
    """Extract text between **marker** and the next ** or --- or end."""
    pattern = re.compile(
        r"\*\*" + re.escape(marker) + r"\*\*\n+(.*?)(?=\n\*\*[A-Z]|\n---|\Z)",
        re.DOTALL,
    )
    m = pattern.search(block)
    if m:
        return m.group(1).strip()
    return ""


# ---------- Word document builder ----------

def build_word(slides, fm):
    doc = Document()

    # Page setup — sober margins for lawyer doc
    for section in doc.sections:
        section.left_margin = Inches(1.0)
        section.right_margin = Inches(1.0)
        section.top_margin = Inches(0.8)
        section.bottom_margin = Inches(0.8)

    # Title page
    title = doc.add_heading(
        "Portfolio de Marcas Tecnológicas Genteca", level=0
    )
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER

    subtitle = doc.add_paragraph()
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = subtitle.add_run("Batch Inicial SAPI Venezuela — Thermo-Safe + FlickerGuard")
    run.font.size = Pt(14)
    run.font.bold = True

    meta = doc.add_paragraph()
    meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
    meta.add_run(
        f"Material de trabajo para Antequera (abogada marcaria) + Liliam Ramirez\n"
        f"Fecha: 2026-05-19\n"
        f"Incorpora las 3 correcciones doctrinales recibidas en la reunión 2026-05-18"
    ).font.size = Pt(11)

    doc.add_page_break()

    # TOC-ish summary (lightweight — Word doesn't auto-generate without manual update)
    doc.add_heading("Contenido", level=1)
    for s in slides:
        p = doc.add_paragraph(style="List Number")
        p.add_run(s["title"])
    doc.add_page_break()

    # Render each slide as a section in Word
    for s in slides:
        # Slide heading
        h = doc.add_heading(f"Slide {s['num']} — {s['title']}", level=1)

        if s["body"]:
            doc.add_heading("Contenido", level=2)
            _render_markdown_block(doc, s["body"])

        if s["speaker_note"]:
            doc.add_heading("Nota de presentación", level=2)
            _render_markdown_block(doc, s["speaker_note"])

        if s["visual_descriptor"]:
            doc.add_heading("Indicación visual", level=2)
            _render_markdown_block(doc, s["visual_descriptor"])

        doc.add_paragraph("")  # separator

    # Footer note
    doc.add_page_break()
    foot = doc.add_paragraph()
    foot.alignment = WD_ALIGN_PARAGRAPH.CENTER
    foot.add_run(
        "Documento producido como entregable Word-first per policy. "
        "PPTX companion derivado del mismo SSOT.\n"
        "Genteca — Portfolio Naming IP 2026 — Versión v3 focused — 2026-05-19"
    ).font.size = Pt(9)

    doc.save(str(DOCX))
    return DOCX


def _render_markdown_block(doc, text: str):
    """
    Render a markdown block into Word.
    Handles: paragraphs, bullets (-), numbered lists, tables (| ... |),
    blockquotes (>), bold (**...**), inline emphasis.
    Simple but adequate for the deck v3 content.
    """
    lines = text.split("\n")
    i = 0
    while i < len(lines):
        ln = lines[i]
        stripped = ln.strip()

        # Blank line
        if not stripped:
            i += 1
            continue

        # Table detection: starts with | and the next line is | --- | --- | ...
        if stripped.startswith("|") and i + 1 < len(lines):
            next_ln = lines[i + 1].strip()
            if re.match(r"^\|[\s\-:\|]+\|$", next_ln):
                # Found a table
                table_lines = [stripped]
                j = i + 1  # separator line
                k = j + 1
                while k < len(lines) and lines[k].strip().startswith("|"):
                    table_lines.append(lines[k].strip())
                    k += 1
                _render_table(doc, [table_lines[0]] + table_lines[2:])  # skip separator
                i = k
                continue

        # Blockquote
        if stripped.startswith(">"):
            quote_lines = []
            while i < len(lines) and (
                lines[i].strip().startswith(">") or lines[i].strip() == ""
            ):
                if lines[i].strip().startswith(">"):
                    quote_lines.append(lines[i].strip().lstrip(">").strip())
                elif quote_lines and lines[i].strip() == "":
                    # blank line inside quote terminates it
                    break
                i += 1
            quote_text = "\n".join(quote_lines)
            p = doc.add_paragraph()
            p.paragraph_format.left_indent = Inches(0.4)
            run = p.add_run(quote_text)
            run.italic = True
            run.font.size = Pt(10)
            continue

        # Bullet (- ...)
        if stripped.startswith("- "):
            p = doc.add_paragraph(style="List Bullet")
            _add_inline_runs(p, stripped[2:])
            i += 1
            continue

        # Numbered list (1. ...)
        if re.match(r"^\d+\.\s", stripped):
            p = doc.add_paragraph(style="List Number")
            content = re.sub(r"^\d+\.\s", "", stripped)
            _add_inline_runs(p, content)
            i += 1
            continue

        # Regular paragraph
        p = doc.add_paragraph()
        _add_inline_runs(p, stripped)
        i += 1


def _render_table(doc, table_lines):
    """Render a markdown table to a Word table."""
    rows = []
    for ln in table_lines:
        cells = [c.strip() for c in ln.strip("|").split("|")]
        rows.append(cells)
    if not rows:
        return
    ncols = max(len(r) for r in rows)
    table = doc.add_table(rows=len(rows), cols=ncols)
    table.style = "Light Grid Accent 1"
    for ri, row in enumerate(rows):
        for ci in range(ncols):
            cell_text = row[ci] if ci < len(row) else ""
            cell = table.rows[ri].cells[ci]
            cell.text = ""
            p = cell.paragraphs[0]
            _add_inline_runs(p, cell_text)
            # Bold header row
            if ri == 0:
                for r in p.runs:
                    r.bold = True
                    r.font.size = Pt(10)
            else:
                for r in p.runs:
                    r.font.size = Pt(10)


def _add_inline_runs(paragraph, text: str):
    """
    Parse **bold** and add runs accordingly. Leaves other markdown as plain.
    """
    parts = re.split(r"(\*\*[^*]+\*\*)", text)
    for part in parts:
        if part.startswith("**") and part.endswith("**"):
            r = paragraph.add_run(part[2:-2])
            r.bold = True
        else:
            paragraph.add_run(part)


# ---------- PowerPoint builder ----------

ACCENT_GREEN = P_RGBColor(0x2E, 0x8B, 0x57)  # NODO-B verde
ACCENT_AMBER = P_RGBColor(0xCC, 0x88, 0x00)  # NODO-C amarillo
ACCENT_RED = P_RGBColor(0xB2, 0x22, 0x22)
TEXT_DARK = P_RGBColor(0x1F, 0x1F, 0x1F)
TEXT_GREY = P_RGBColor(0x55, 0x55, 0x55)


def build_pptx(slides, fm):
    prs = Presentation()
    prs.slide_width = P_Inches(13.333)
    prs.slide_height = P_Inches(7.5)

    blank = prs.slide_layouts[6]

    for s in slides:
        slide = prs.slides.add_slide(blank)

        # Title strip
        tb = slide.shapes.add_textbox(
            P_Inches(0.4), P_Inches(0.3), P_Inches(12.5), P_Inches(0.7)
        )
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"Slide {s['num']} — {s['title']}"
        run.font.size = P_Pt(20)
        run.font.bold = True
        run.font.color.rgb = TEXT_DARK

        # Body box (left 8 inches, leave right for visual descriptor placeholder)
        body_box = slide.shapes.add_textbox(
            P_Inches(0.4), P_Inches(1.2), P_Inches(8.4), P_Inches(5.9)
        )
        body_tf = body_box.text_frame
        body_tf.word_wrap = True
        _fill_textframe_from_markdown(body_tf, s["body"])

        # Visual descriptor as right-side hint (smaller, italic, grey)
        if s["visual_descriptor"]:
            vis_box = slide.shapes.add_textbox(
                P_Inches(9.0), P_Inches(1.2), P_Inches(4.0), P_Inches(5.9)
            )
            vis_tf = vis_box.text_frame
            vis_tf.word_wrap = True
            head = vis_tf.paragraphs[0]
            head.add_run().text = "INDICACIÓN VISUAL"
            head.runs[0].font.bold = True
            head.runs[0].font.size = P_Pt(10)
            head.runs[0].font.color.rgb = TEXT_GREY
            for line in _trim_for_pptx(s["visual_descriptor"]).split("\n"):
                if not line.strip():
                    continue
                p = vis_tf.add_paragraph()
                r = p.add_run()
                r.text = line.strip()
                r.font.size = P_Pt(9)
                r.font.italic = True
                r.font.color.rgb = TEXT_GREY

        # Speaker note
        if s["speaker_note"]:
            notes_tf = slide.notes_slide.notes_text_frame
            notes_tf.text = s["speaker_note"]

        # Footer with slide number
        footer = slide.shapes.add_textbox(
            P_Inches(11.5), P_Inches(7.1), P_Inches(1.7), P_Inches(0.3)
        )
        ftf = footer.text_frame
        ftf.text = f"Genteca IP 2026 — v3 — {s['num']}/{len(slides)}"
        ftf.paragraphs[0].runs[0].font.size = P_Pt(8)
        ftf.paragraphs[0].runs[0].font.color.rgb = TEXT_GREY

    prs.save(str(PPTX))
    return PPTX


def _trim_for_pptx(text: str, max_chars: int = 600) -> str:
    """Visual descriptor in PPTX shouldn't dominate — trim if huge."""
    if len(text) <= max_chars:
        return text
    return text[:max_chars].rstrip() + "..."


def _fill_textframe_from_markdown(tf, text: str):
    """
    Simplified markdown -> pptx text frame.
    Handles bullets, bold inline, tables converted to compact text rows,
    blockquotes shortened.
    """
    if not text:
        return
    lines = text.split("\n")
    first = True
    i = 0
    while i < len(lines):
        ln = lines[i]
        stripped = ln.strip()
        if not stripped:
            i += 1
            continue

        # Table -> compact text block
        if stripped.startswith("|") and i + 1 < len(lines):
            next_ln = lines[i + 1].strip()
            if re.match(r"^\|[\s\-:\|]+\|$", next_ln):
                table_lines = [stripped]
                k = i + 2
                while k < len(lines) and lines[k].strip().startswith("|"):
                    table_lines.append(lines[k].strip())
                    k += 1
                # Render as one paragraph per row, compressed
                for row_line in table_lines:
                    cells = [
                        c.strip() for c in row_line.strip("|").split("|")
                    ]
                    p = _new_para(tf, first)
                    first = False
                    r = p.add_run()
                    r.text = " | ".join(cells)
                    r.font.size = P_Pt(10)
                    r.font.color.rgb = TEXT_DARK
                i = k
                continue

        # Blockquote
        if stripped.startswith(">"):
            quote_lines = []
            while i < len(lines) and lines[i].strip().startswith(">"):
                quote_lines.append(lines[i].strip().lstrip(">").strip())
                i += 1
            quote_text = " ".join(quote_lines)
            # truncate very long quotes for slide; full is in speaker notes
            if len(quote_text) > 380:
                quote_text = quote_text[:380].rstrip() + "..."
            p = _new_para(tf, first)
            first = False
            r = p.add_run()
            r.text = "“" + quote_text + "”"
            r.font.italic = True
            r.font.size = P_Pt(10)
            r.font.color.rgb = TEXT_GREY
            continue

        # Bullet
        if stripped.startswith("- "):
            p = _new_para(tf, first)
            first = False
            _add_runs_pptx(p, "•  " + stripped[2:], font_size=11)
            i += 1
            continue

        # Numbered
        if re.match(r"^\d+\.\s", stripped):
            p = _new_para(tf, first)
            first = False
            _add_runs_pptx(p, stripped, font_size=11)
            i += 1
            continue

        # Paragraph
        p = _new_para(tf, first)
        first = False
        _add_runs_pptx(p, stripped, font_size=11)
        i += 1


def _new_para(tf, first: bool):
    if first:
        return tf.paragraphs[0]
    return tf.add_paragraph()


def _add_runs_pptx(paragraph, text: str, font_size: int = 11):
    parts = re.split(r"(\*\*[^*]+\*\*)", text)
    for part in parts:
        if part.startswith("**") and part.endswith("**"):
            r = paragraph.add_run()
            r.text = part[2:-2]
            r.font.bold = True
            r.font.size = P_Pt(font_size)
            r.font.color.rgb = TEXT_DARK
        else:
            r = paragraph.add_run()
            r.text = part
            r.font.size = P_Pt(font_size)
            r.font.color.rgb = TEXT_DARK


# ---------- Main ----------

def main():
    md_text = MD.read_text(encoding="utf-8")
    fm, slides = parse_slides(md_text)
    print(f"Parsed {len(slides)} slides from {MD.name}")
    word_path = build_word(slides, fm)
    print(f"Word written: {word_path}")
    pptx_path = build_pptx(slides, fm)
    print(f"PPTX written: {pptx_path}")


if __name__ == "__main__":
    main()
