import sys
sys.stdout.reconfigure(encoding='utf-8')

import re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ──────────────────────────────────────────────────────────────
# CONFIG
# ──────────────────────────────────────────────────────────────
MD_PATH = r"C:\Raul\03-projects\genteca\portfolio-naming-ip-2026\06-three-pack\03-Deck-lawyer-presentation_v2.md"
OUT_PATH = r"C:\Raul\03-projects\genteca\portfolio-naming-ip-2026\06-three-pack\03-Deck-lawyer-presentation_v2.pptx"

# Metadata canonica del archivo PPTX (visible en propiedades de Windows / PowerPoint)
META_TITLE   = "Portfolio de Marcas Tecnologicas Genteca - Estrategia de Registro IP"
META_AUTHOR  = "Genteca"
META_COMPANY = "Genteca"
META_SUBJECT = "SAPI Venezuela + IMPI Mexico - Clase 9 Niza"
META_KEYWORDS = "marcas, registro, SAPI, IMPI, Clase 9, propiedad industrial, Genteca"

# Paleta Genteca: fondo blanco, titulo negro, cuerpo gris oscuro, acento azul oscuro
COLOR_BG       = RGBColor(0xFF, 0xFF, 0xFF)  # blanco
COLOR_TITLE    = RGBColor(0x1A, 0x1A, 0x1A)  # negro suave
COLOR_BODY     = RGBColor(0x2C, 0x2C, 0x2C)  # gris oscuro
COLOR_NOTE     = RGBColor(0x55, 0x55, 0x55)  # gris medio (speaker notes)
COLOR_ACCENT   = RGBColor(0x1A, 0x3A, 0x5C)  # azul Genteca
COLOR_BLOCK    = RGBColor(0x08, 0x25, 0x44)  # azul muy oscuro para bloque header
COLOR_BLOCK_TXT= RGBColor(0xFF, 0xFF, 0xFF)  # blanco texto en bloque

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ──────────────────────────────────────────────────────────────
# PARSE MARKDOWN
# ──────────────────────────────────────────────────────────────

def parse_slides(md_path):
    """
    Parse the VI-1 Markdown into a list of slide dicts.
    Each slide has: title, block_label (BLOQUE X — ...), body_lines, speaker_note, visual
    """
    with open(md_path, encoding='utf-8') as f:
        raw = f.read()

    # Split on ## Slide N — or ## BLOQUE headings
    # We collect slides between ## Slide N headings
    slide_pattern = re.compile(r'^## (Slide \d+.*?)$', re.MULTILINE)
    block_pattern = re.compile(r'^## (BLOQUE [A-Z].*?)$', re.MULTILINE)

    # Find all slide positions
    slide_matches = list(slide_pattern.finditer(raw))
    block_matches = list(block_pattern.finditer(raw))

    # Build a map of position -> block label
    block_map = {}
    for m in block_matches:
        block_map[m.start()] = m.group(1)

    def find_block_for(pos):
        last = None
        for bpos, blabel in sorted(block_map.items()):
            if bpos <= pos:
                last = blabel
        return last

    slides = []
    for i, m in enumerate(slide_matches):
        start = m.start()
        end = slide_matches[i+1].start() if i+1 < len(slide_matches) else len(raw)
        section = raw[start:end]
        title_raw = m.group(1)  # "Slide 1 — Portada"
        # Remove "Slide N — " prefix for the title text
        title_clean = re.sub(r'^Slide \d+ — ', '', title_raw).strip()
        block_label = find_block_for(start)

        # Extract body (between **Body:** and next **... header)
        body_lines = []
        note_lines = []
        visual_lines = []
        section_mode = None

        for line in section.split('\n'):
            if line.strip().startswith('**Body:**'):
                section_mode = 'body'
                continue
            if line.strip().startswith('**Speaker note:**'):
                section_mode = 'note'
                continue
            if line.strip().startswith('**Visual descriptor'):
                section_mode = 'visual'
                continue
            if section_mode == 'body':
                body_lines.append(line)
            elif section_mode == 'note':
                note_lines.append(line)
            elif section_mode == 'visual':
                visual_lines.append(line)

        slides.append({
            'title_raw': title_raw,
            'title': title_clean,
            'block': block_label or '',
            'body': '\n'.join(body_lines).strip(),
            'note': '\n'.join(note_lines).strip(),
            'visual': '\n'.join(visual_lines).strip(),
        })

    return slides

# ──────────────────────────────────────────────────────────────
# BUILD PPTX
# ──────────────────────────────────────────────────────────────

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=14,
                bold=False, color=COLOR_BODY, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox

def add_block_header(slide, label, left, top, width, height):
    """Blue block for BLOQUE headings"""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_BLOCK
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = COLOR_BLOCK_TXT

def format_body(text):
    """Clean markdown formatting for body text"""
    # Remove bold markdown
    text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
    # Remove italic
    text = re.sub(r'\*(.+?)\*', r'\1', text)
    # Condense multiple blanks
    text = re.sub(r'\n{3,}', '\n\n', text)
    return text.strip()

def build_pptx(slides, out_path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Set core properties (metadata visible en Windows + PowerPoint)
    cp = prs.core_properties
    cp.title    = META_TITLE
    cp.author   = META_AUTHOR
    cp.last_modified_by = META_AUTHOR
    cp.subject  = META_SUBJECT
    cp.keywords = META_KEYWORDS
    cp.category = "Estrategia IP"
    cp.comments = ""
    cp.revision = 1

    blank_layout = prs.slide_layouts[6]  # completely blank

    current_block = ''

    for s in slides:
        slide = prs.slides.add_slide(blank_layout)
        set_bg(slide, COLOR_BG)

        margin_l = Inches(0.5)
        margin_t = Inches(0.4)
        content_w = Inches(12.3)

        # ── Block header bar (thin, at very top) ──────────────
        block = s.get('block', '')
        if block and block != current_block:
            current_block = block
            add_block_header(slide, block,
                             margin_l, margin_t,
                             content_w, Inches(0.28))
            title_top = margin_t + Inches(0.32)
        else:
            title_top = margin_t

        # ── Title ─────────────────────────────────────────────
        add_textbox(slide,
                    margin_l, title_top,
                    content_w, Inches(0.65),
                    s['title'],
                    font_size=22, bold=True, color=COLOR_ACCENT,
                    align=PP_ALIGN.LEFT)

        # ── Separator line (thin rectangle) ───────────────────
        sep_top = title_top + Inches(0.68)
        sep = slide.shapes.add_shape(1, margin_l, sep_top, content_w, Emu(18000))
        sep.fill.solid()
        sep.fill.fore_color.rgb = COLOR_ACCENT
        sep.line.fill.background()

        # ── Body ──────────────────────────────────────────────
        body_top = sep_top + Inches(0.12)
        body_h   = Inches(4.6)
        body_txt = format_body(s['body']) if s['body'] else ''
        if body_txt:
            add_textbox(slide,
                        margin_l, body_top,
                        content_w, body_h,
                        body_txt,
                        font_size=13, color=COLOR_BODY)

        # ── Footer: visual descriptor hint (small) ─────────────
        if s.get('visual'):
            vis_short = s['visual'][:120] + ('...' if len(s['visual']) > 120 else '')
            vis_txt = '[Visual: ' + format_body(vis_short) + ']'
            add_textbox(slide,
                        margin_l, Inches(6.9),
                        content_w, Inches(0.45),
                        vis_txt,
                        font_size=9, color=RGBColor(0x99, 0x99, 0x99),
                        align=PP_ALIGN.LEFT)

        # ── Speaker notes ─────────────────────────────────────
        if s.get('note'):
            note_tf = slide.notes_slide.notes_text_frame
            note_tf.text = format_body(s['note'])

    prs.save(out_path)
    print(f'PPTX guardado en: {out_path}')
    print(f'Total slides: {len(slides)}')


# ──────────────────────────────────────────────────────────────
# MAIN
# ──────────────────────────────────────────────────────────────
if __name__ == '__main__':
    slides = parse_slides(MD_PATH)
    print(f'Slides parseados: {len(slides)}')
    for i, s in enumerate(slides, 1):
        print(f'  {i:02d}. {s["title_raw"]}')
    build_pptx(slides, OUT_PATH)
