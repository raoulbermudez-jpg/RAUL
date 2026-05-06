"""
build_v6.py — Junta_GSM_empaque_v6.pptx
Vivienne — Presentation Designer — Sistema /RAUL/
2026-05-04
"""

import sys, os
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from pptx.enum.dml import MSO_THEME_COLOR
import copy

# ── Paths ──────────────────────────────────────────────────────────────────────
OUT_DIR  = r'C:\Raul\03-projects\genteca\2026-04_GSM-MB-RB-RF_empaque\03-review-and-release\Junta_PPTX_v2'
PNG_DIR  = OUT_DIR
OUT_FILE = os.path.join(OUT_DIR, 'Junta_GSM_empaque_v6.pptx')

MOCK_A = os.path.join(PNG_DIR, 'Atlas_mockup_frente_A.png')
MOCK_B = os.path.join(PNG_DIR, 'Atlas_mockup_frente_B.png')
MOCK_C = os.path.join(PNG_DIR, 'Atlas_mockup_frente_C.png')
MOCK_D = os.path.join(PNG_DIR, 'Atlas_mockup_frente_D.png')

# ── Palette ────────────────────────────────────────────────────────────────────
ORANGE   = RGBColor(0xFF, 0x82, 0x00)   # #FF8200 Exceline
DARK_BG  = RGBColor(0x1A, 0x1A, 0x1A)   # almost-black background
NEAR_BLK = RGBColor(0x12, 0x12, 0x12)   # stripe
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LGREY    = RGBColor(0xCC, 0xCC, 0xCC)   # body copy
MGREY    = RGBColor(0x88, 0x88, 0x88)   # secondary
DGREY    = RGBColor(0x33, 0x33, 0x33)   # table row alt
MIDGREY  = RGBColor(0x55, 0x55, 0x55)   # table border

# ── Slide dimensions (widescreen 16:9) ─────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.50)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # blank layout

# ── Helpers ────────────────────────────────────────────────────────────────────

def new_slide():
    slide = prs.slides.add_slide(BLANK)
    # dark background rectangle
    bg = slide.shapes.add_shape(1, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()
    # orange top stripe 6px
    top = slide.shapes.add_shape(1, 0, 0, W, Pt(6))
    top.fill.solid(); top.fill.fore_color.rgb = ORANGE; top.line.fill.background()
    # orange bottom stripe 6px
    bot = slide.shapes.add_shape(1, 0, H - Pt(6), W, Pt(6))
    bot.fill.solid(); bot.fill.fore_color.rgb = ORANGE; bot.line.fill.background()
    return slide


def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    txb.word_wrap = wrap
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic    = italic
    # font name with Calibri fallback
    run.font.name = 'Calibri'
    return txb


def add_label(slide, text, l, t, w=Inches(3), color=ORANGE, size=9):
    """Small all-caps orange label above a title."""
    add_text(slide, text.upper(), l, t, w, Pt(18),
             size=size, bold=True, color=color, align=PP_ALIGN.LEFT)


def orange_bar(slide, l, t, w=Inches(0.07), h=Inches(0.9)):
    """Vertical orange accent bar to left of title."""
    bar = slide.shapes.add_shape(1, l, t, w, h)
    bar.fill.solid(); bar.fill.fore_color.rgb = ORANGE
    bar.line.fill.background()


def add_divider(slide, t, color=MGREY):
    line = slide.shapes.add_shape(1, Inches(0.6), t, W - Inches(1.2), Pt(1))
    line.fill.solid(); line.fill.fore_color.rgb = color
    line.line.fill.background()


def add_mockup(slide, png_path, l, t, w, h):
    if os.path.exists(png_path):
        slide.shapes.add_picture(png_path, l, t, w, h)
    else:
        placeholder = slide.shapes.add_shape(1, l, t, w, h)
        placeholder.fill.solid(); placeholder.fill.fore_color.rgb = DGREY
        placeholder.line.color.rgb = ORANGE
        add_text(slide, '[MOCKUP]', l, t + h/2 - Pt(12), w, Pt(24),
                 size=14, color=MGREY, align=PP_ALIGN.CENTER)


def _set_cell_fill(cell, rgb_color):
    tc   = cell._tc
    tcPr = tc.get_or_add_tcPr()
    sf   = tcPr.get_or_change_to_solidFill()
    srgb = sf.get_or_change_to_srgbClr()
    srgb.set('val', str(rgb_color))

def add_table_simple(slide, rows_data, l, t, w, h, header_color=ORANGE):
    """rows_data: list of lists of strings. First row = header."""
    nr = len(rows_data)
    nc = len(rows_data[0])
    tbl = slide.shapes.add_table(nr, nc, l, t, w, h).table
    col_w = w // nc
    for ci in range(nc):
        tbl.columns[ci].width = col_w
    for ri, row in enumerate(rows_data):
        for ci, cell_text in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = cell_text
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size  = Pt(11)
                    run.font.color.rgb = WHITE if ri == 0 else LGREY
                    run.font.bold  = (ri == 0)
                    run.font.name  = 'Calibri'
            # cell fill (pptx v1 API)
            if ri == 0:
                _set_cell_fill(cell, header_color)
            elif ri % 2 == 0:
                _set_cell_fill(cell, DGREY)
            else:
                _set_cell_fill(cell, NEAR_BLK)
    return tbl


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — PORTADA
# ─────────────────────────────────────────────────────────────────────────────
slide = new_slide()

# Center accent block
accent = slide.shapes.add_shape(1, Inches(0), Inches(2.6), W, Inches(2.3))
accent.fill.solid(); accent.fill.fore_color.rgb = RGBColor(0x0F, 0x0F, 0x0F)
accent.line.fill.background()

# Left orange stripe
left_bar = slide.shapes.add_shape(1, 0, Inches(2.6), Inches(0.12), Inches(2.3))
left_bar.fill.solid(); left_bar.fill.fore_color.rgb = ORANGE
left_bar.line.fill.background()

add_text(slide, 'GENTECA  |  EXCELINE PROFESIONAL',
         Inches(0.7), Inches(0.45), Inches(12), Inches(0.5),
         size=11, color=MGREY, bold=False)

add_text(slide, 'Empaque Linea Exceline GSM',
         Inches(0.7), Inches(2.65), Inches(11), Inches(0.85),
         size=36, bold=True, color=WHITE)

add_text(slide, 'Decision de Junta Directiva   |   Cuatro Alternativas A / B / C / D   |   QR   |   Venta Cruzada',
         Inches(0.7), Inches(3.45), Inches(11), Inches(0.45),
         size=15, color=LGREY)

add_text(slide, 'v2  |  2026-05-04',
         Inches(0.7), Inches(6.6), Inches(4), Inches(0.4),
         size=11, color=MGREY)

add_text(slide, 'CONFIDENCIAL  |  Solo para uso interno Genteca',
         Inches(7.5), Inches(6.6), Inches(5.5), Inches(0.4),
         size=10, color=MGREY, align=PP_ALIGN.RIGHT)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — PUNTO DE PARTIDA
# ─────────────────────────────────────────────────────────────────────────────
slide = new_slide()

add_label(slide, 'CONTEXTO', Inches(0.7), Inches(0.45))
orange_bar(slide, Inches(0.55), Inches(0.85))
add_text(slide, 'La Junta ya decidio comunicar las dos innovaciones. Hoy elige como.',
         Inches(0.75), Inches(0.85), Inches(11.8), Inches(0.85),
         size=26, bold=True, color=WHITE)

add_divider(slide, Inches(1.9))

# Three columns
cols = [
    ('DECISION PREVIA - CERRADA',
     'En sesion anterior la Junta aprobó comunicar las dos innovaciones del GSM en el empaque: el sensor NTC y el tiempo de respuesta ante parpadeos < 30 ms.'),
    ('LO QUE HOY SE DECIDE',
     'La Junta elige cual de las cuatro arquitecturas de comunicacion del tiro (frente del empaque) representa mejor la propuesta de valor del nuevo GSM.'),
    ('POR QUE IMPORTA HOY',
     'Esta decision desbloquea la instruccion a Oz para el arte final y el inicio de produccion. Sin decision, el empaque no avanza.'),
]
for i, (hdr, body) in enumerate(cols):
    lx = Inches(0.6 + i * 4.2)
    add_text(slide, hdr, lx, Inches(2.1), Inches(3.9), Inches(0.4),
             size=10, bold=True, color=ORANGE)
    add_text(slide, body, lx, Inches(2.55), Inches(3.9), Inches(3.5),
             size=14, color=LGREY)

# Four alternatives summary box
box = slide.shapes.add_shape(1, Inches(0.6), Inches(5.85), Inches(12.1), Inches(1.2))
box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x22)
box.line.color.rgb = ORANGE

add_text(slide, 'CUATRO ALTERNATIVAS EN MESA:  A - Completitud  |  B - Dos diferenciadores unicos  |  C - Una frase narrativa  |  D - QR como eje (nueva)',
         Inches(0.8), Inches(5.95), Inches(11.7), Inches(1.0),
         size=13, bold=True, color=WHITE)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — LAS DOS INNOVACIONES
# ─────────────────────────────────────────────────────────────────────────────
slide = new_slide()

add_label(slide, 'LAS INNOVACIONES', Inches(0.7), Inches(0.45))
orange_bar(slide, Inches(0.55), Inches(0.85))
add_text(slide, 'Dos diferenciadores que ningun competidor venezolano puede igualar hoy',
         Inches(0.75), Inches(0.85), Inches(11.8), Inches(0.75),
         size=26, bold=True, color=WHITE)

add_divider(slide, Inches(1.75))

# Innovation 1
box1 = slide.shapes.add_shape(1, Inches(0.6), Inches(1.9), Inches(5.9), Inches(4.9))
box1.fill.solid(); box1.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x1E)
box1.line.color.rgb = ORANGE

add_text(slide, 'INNOVACION 1', Inches(0.8), Inches(2.0), Inches(5.5), Inches(0.35),
         size=9, bold=True, color=ORANGE)
add_text(slide, 'Tiempo de respuesta < 30 ms ante parpadeos',
         Inches(0.8), Inches(2.35), Inches(5.5), Inches(0.7),
         size=19, bold=True, color=WHITE)
add_text(slide, '< 0,03 s',
         Inches(0.8), Inches(3.1), Inches(5.5), Inches(0.9),
         size=48, bold=True, color=ORANGE)
add_text(slide,
    'El GSM detecta y desconecta ante parpadeos (fluctuaciones rapidas del voltaje) '
    'en menos de 30 milisegundos. La generacion anterior tardaba 150 ms: una mejora '
    'de 5x en el propio producto. Ningun competidor venezolano publica un dato comparable.',
    Inches(0.8), Inches(4.05), Inches(5.5), Inches(1.5),
    size=13, color=LGREY)

# Innovation 2
box2 = slide.shapes.add_shape(1, Inches(6.8), Inches(1.9), Inches(5.9), Inches(4.9))
box2.fill.solid(); box2.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x1E)
box2.line.color.rgb = MGREY

add_text(slide, 'INNOVACION 2', Inches(7.0), Inches(2.0), Inches(5.5), Inches(0.35),
         size=9, bold=True, color=ORANGE)
add_text(slide, 'Sensor NTC incorporado',
         Inches(7.0), Inches(2.35), Inches(5.5), Inches(0.7),
         size=19, bold=True, color=WHITE)
add_text(slide, 'Autoproteccion termica activa',
         Inches(7.0), Inches(3.1), Inches(5.5), Inches(0.55),
         size=22, bold=False, color=LGREY)
add_text(slide,
    'El sensor NTC detecta calentamiento excesivo causado por sobrecorrientes '
    'o conexiones deficientes, y desconecta la carga antes de que cables o '
    'bornes se danen. Protege el protector mismo y la instalacion. '
    'Ningun competidor venezolano comunica NTC o autoproteccion termica en empaque.',
    Inches(7.0), Inches(3.75), Inches(5.5), Inches(1.8),
    size=13, color=LGREY)

add_text(slide, 'Lenguaje en empaque: "El mas rapido ante parpadeos (< 0,03 s)" + "Sensor NTC incorporado*"',
         Inches(0.6), Inches(7.0), Inches(12.1), Inches(0.38),
         size=10, color=MGREY, italic=True)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — TABLA COMPETITIVA CORREGIDA
# ─────────────────────────────────────────────────────────────────────────────
slide = new_slide()

add_label(slide, 'INTELIGENCIA DE MERCADO', Inches(0.7), Inches(0.45))
orange_bar(slide, Inches(0.55), Inches(0.85))
add_text(slide, 'El mercado: factor 2x contra Breakermatic, hasta 10x contra el resto',
         Inches(0.75), Inches(0.85), Inches(11.8), Inches(0.75),
         size=26, bold=True, color=WHITE)

add_divider(slide, Inches(1.75))

rows = [
    ['Marca', 'Tiempo de respuesta', 'Fuente / confianza', 'Factor vs Genteca (tipico 25 ms)'],
    ['Genteca Exceline GSM (nueva)', '20-30 ms  [< 0,03 s]', 'Laboratorio I&D Genteca — confirmado Liliam', '—  REFERENCIA'],
    ['Breakermatic', '32-64 ms  (50 ms tipico)', 'Hoja tecnica + lab Genteca — confianza ALTA', '~2x'],
    ['WellSpec', '200-300 ms', 'Claim marketing distribuidor — confianza MEDIA', '~10x'],
    ['Powest Airmatic', '1.000 ms', 'Datasheet oficial Powest — confianza ALTA', '~40x'],
    ['TQ / Avtek / JVTRONIC / Protektor', 'No publican dato', 'Multiples fuentes — ausencia confirmada', 'No calculable'],
]

add_table_simple(slide, rows, Inches(0.6), Inches(1.85), Inches(12.1), Inches(3.8))

add_text(slide, 'COMO LEER ESTA TABLA',
         Inches(0.6), Inches(5.85), Inches(12.1), Inches(0.3),
         size=9, bold=True, color=ORANGE)

bullets = [
    '"El mas rapido del mercado" sigue siendo correcto en todos los escenarios — Genteca supera a Breakermatic incluso en el caso mas desfavorable (30 ms Genteca vs 32 ms Breakermatic).',
    '"Hasta 10x mas rapido" es real y honesto solo si se refiere al grueso del mercado (WellSpec, Powest). Contra Breakermatic la ventaja tipica es 2x. El cuantificador "hasta" es critico.',
    'El dato de Breakermatic fue verificado en laboratorio Genteca (Liliam I&D, mayo 2026) — reemplaza el dato anterior de distribuidores Colombia que reportaba 1.000-1.500 ms.',
]
for i, b in enumerate(bullets):
    add_text(slide, chr(0x2022) + '  ' + b,
             Inches(0.6), Inches(6.22) + Pt(i * 42), Inches(12.1), Pt(38),
             size=11, color=LGREY)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — FLICKERS VS PICOS (nueva en v2)
# ─────────────────────────────────────────────────────────────────────────────
slide = new_slide()

add_label(slide, 'FUNDAMENTO TECNICO — NUEVA EN V2', Inches(0.7), Inches(0.45))
orange_bar(slide, Inches(0.55), Inches(0.85))
add_text(slide, 'Por que el GSM protege auténticamente al equipo inverter',
         Inches(0.75), Inches(0.85), Inches(11.8), Inches(0.75),
         size=26, bold=True, color=WHITE)

add_divider(slide, Inches(1.75))

# Three columns: the distinction
cols3 = [
    ('EL EQUIPO INVERTER YA TIENE PROTECCION...',
     'Equipos inverter modernos incluyen MOV (varistores) de fabrica.\n\n'
     'El MOV protege contra picos de alta energia (rayos, transitorios impulsivos).\n\n'
     'Certificacion tipica: IEC 61000-4-5 nivel 3-4 (surges hasta 2-4 kV).\n\n'
     'Esa proteccion cumple su funcion. No esta en discusion.',
     MGREY, '...PERO PARA PICOS, NO PARA PARPADEOS'),
    ('EL MECANISMO DE DANO EN FLICKERS',
     'El bus DC del inverter tiene un capacitor de hold-up con autonomia tipica de 8-25 ms.\n\n'
     'Un flicker (caida y recuperacion rapida) agota ese capacitor si el protector no actua a tiempo.\n\n'
     'Consecuencia: overcurrent en los IGBT durante la recuperacion + estres termico ciclico '
     'del capacitor + posibles resets del microcontrolador de control.\n\n'
     'El dano es acumulativo, no instantaneo.',
     ORANGE, 'POR QUE LA VELOCIDAD DE RESPUESTA IMPORTA'),
    ('< 30 ms: LLEGA A TIEMPO. > 100 ms: LLEGA TARDE.',
     'Un protector con respuesta < 30 ms interrumpe el suministro dentro del hold-up time tipico '
     'del bus DC (8-25 ms), protegiendo la electronica antes de que ocurra el dano.\n\n'
     'Un protector con 100-200 ms llega cuando el evento ya ocurrio.\n\n'
     'Los MOV del equipo no "ven" caidas de voltaje — operan solo ante sobretension. '
     'IEC 61000-4-11 es la norma de inmunidad a voltage sags — distinta de IEC 61000-4-5.',
     MGREY, 'FUENTE: PAXS_INVERTER_FLICKERS_V1 | NIVEL CONFIANZA: ALTO'),
]
for i, (hdr, body, col, footer) in enumerate(cols3):
    lx = Inches(0.55 + i * 4.25)
    box = slide.shapes.add_shape(1, lx, Inches(1.9), Inches(4.1), Inches(4.55))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0x1C, 0x1C, 0x1C)
    box.line.color.rgb = col
    add_text(slide, hdr, lx + Inches(0.12), Inches(1.98), Inches(3.85), Inches(0.42),
             size=9, bold=True, color=col)
    add_text(slide, body, lx + Inches(0.12), Inches(2.45), Inches(3.85), Inches(3.4),
             size=11, color=LGREY)
    add_text(slide, footer, lx + Inches(0.12), Inches(6.25), Inches(3.85), Inches(0.25),
             size=8, color=MGREY, italic=True)

add_text(slide,
    'NOTA: El argumento de "proteccion autentica inverter" vive detras del QR (en D) y en el argumentario ESC — no en el frente del empaque (Bruna §7.4).',
    Inches(0.6), Inches(6.9), Inches(12.1), Inches(0.38),
    size=10, color=MGREY, italic=True)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — LO QUE LA CADENA GATEO
# ─────────────────────────────────────────────────────────────────────────────
slide = new_slide()

add_label(slide, 'GATE DE CLAIMS', Inches(0.7), Inches(0.45))
orange_bar(slide, Inches(0.55), Inches(0.85))
add_text(slide, 'Todo lo que aparece en el empaque tiene sello de Bruna. Lo que no, se sabe por que.',
         Inches(0.75), Inches(0.85), Inches(11.8), Inches(0.75),
         size=24, bold=True, color=WHITE)

add_divider(slide, Inches(1.75))

rows_g = [
    ['Claim', 'Estado', 'Ruta Bruna', 'Aplica a'],
    ['"El mas rapido ante parpadeos (< 0,03 s)"', 'APROBADO con caveat', 'BR-2 §2 Claim A + §7.3', 'A / B / C / D'],
    ['"Protege tecnologia Inverter"', 'APROBADO con caveat', 'BR-2 §2 Claim C', 'A / B / C / D'],
    ['"Sensor NTC incorporado*"', 'APROBADO (asterisco obligatorio)', 'BR-2 §2 Claim D + §7.1', 'A / B / C'],
    ['"Hasta 10 veces mas rapido"', 'APROBADO con caveat §7.2 en retiro', 'BR-2 §7.2 — EXCLUSIVO D', 'Solo D'],
    ['"Sensor NTC: protege contra calentamiento*"', 'Gate confirmatorio pendiente Bruna', 'Adaptacion Claim D para arq. causal D', 'Solo D'],
    ['Lengueta "AVERIGUALO"', 'Gate confirmatorio pendiente Bruna', 'Formulacion nueva — bajo riesgo segun Solenne', 'Solo D'],
    ['"Autenticamente protege tecnologia Inverter"', 'RECHAZADO para frente empaque', 'Bruna §7.4 — va en QR y ESC', 'Ninguna'],
]
add_table_simple(slide, rows_g, Inches(0.6), Inches(1.85), Inches(12.1), Inches(4.3))

add_text(slide,
    'Para A/B/C: cero gates pendientes. Para D: dos confirmaciones de bajo riesgo antes del arte final. '
    'El asterisco NTC fue REESCRITO en v2 (sin temperaturas numericas) — la formulacion de v1 queda obsoleta.',
    Inches(0.6), Inches(6.35), Inches(12.1), Inches(0.45),
    size=11, color=LGREY, italic=True)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — LA PREGUNTA DE HOY: CUATRO ALTERNATIVAS
# ─────────────────────────────────────────────────────────────────────────────
slide = new_slide()

add_label(slide, 'LA DECISION DE HOY', Inches(0.7), Inches(0.45))
orange_bar(slide, Inches(0.55), Inches(0.85))
add_text(slide, 'Cuatro arquitecturas de comunicacion en mesa',
         Inches(0.75), Inches(0.85), Inches(11.8), Inches(0.75),
         size=28, bold=True, color=WHITE)

add_divider(slide, Inches(1.75))

alts = [
    ('A', 'Completitud de claims',
     'Tres innovaciones en el tiro + lengueta "La Proteccion mas completa".\nRatificada por Owner y Canudas.',
     ORANGE),
    ('B', 'Dos diferenciadores unicos',
     'Solo velocidad y NTC en el tiro. Inverter se desarrolla en el retiro.\nRecomendada por Aurelio y Atlas (primaria).',
     ORANGE),
    ('C', 'Una frase narrativa',
     '"Actua en < 0,03 s antes de que la fluctuacion llegue a tu equipo".\nLa mas accesible para consumidor residencial.',
     LGREY),
    ('D', 'QR como eje — NUEVA',
     '"Hasta 10 veces mas rapido" + flecha causal + QR en zona de titulos.\nLengueta "NUEVO / AVERIGUALO". Statement de innovacion.',
     ORANGE),
]
for i, (letter, title, body, col) in enumerate(alts):
    lx = Inches(0.55 + i * 3.2)
    box = slide.shapes.add_shape(1, lx, Inches(1.95), Inches(3.05), Inches(4.8))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x1E)
    box.line.color.rgb = col if letter in ('A', 'B', 'D') else MGREY
    add_text(slide, letter, lx + Inches(0.1), Inches(2.05), Inches(0.6), Inches(0.65),
             size=40, bold=True, color=ORANGE)
    add_text(slide, title, lx + Inches(0.1), Inches(2.7), Inches(2.8), Inches(0.4),
             size=12, bold=True, color=WHITE)
    add_text(slide, body, lx + Inches(0.1), Inches(3.15), Inches(2.8), Inches(2.4),
             size=11, color=LGREY)

# Recomendacion note
rec_box = slide.shapes.add_shape(1, Inches(0.55), Inches(6.9), Inches(12.1), Inches(0.42))
rec_box.fill.solid(); rec_box.fill.fore_color.rgb = RGBColor(0x28, 0x1A, 0x00)
rec_box.line.color.rgb = ORANGE
add_text(slide,
    'RECOMENDACION AURELIO v2:  B como primaria para canal tecnico/ferreteria de mostrador  |  '
    'D como alternativa explícita si la Junta prioriza statement de innovacion',
    Inches(0.75), Inches(6.92), Inches(11.9), Inches(0.4),
    size=11, bold=True, color=ORANGE)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — ALTERNATIVA A
# ─────────────────────────────────────────────────────────────────────────────
slide = new_slide()

add_label(slide, 'ALTERNATIVA A', Inches(0.7), Inches(0.45))
orange_bar(slide, Inches(0.55), Inches(0.85))
add_text(slide, 'A: Completitud de tres innovaciones — el instalador tecnico ve todo en el tiro',
         Inches(0.75), Inches(0.85), Inches(11.8), Inches(0.75),
         size=23, bold=True, color=WHITE)

add_divider(slide, Inches(1.75))

# Mockup
add_mockup(slide, MOCK_A, Inches(0.55), Inches(1.85), Inches(3.7), Inches(5.3))

# Content
add_text(slide, 'LENGUETA / CINTILLA', Inches(4.5), Inches(1.9), Inches(8.6), Inches(0.3),
         size=9, bold=True, color=ORANGE)
add_text(slide, 'NUEVO  |  LA PROTECCION MAS COMPLETA',
         Inches(4.5), Inches(2.2), Inches(8.6), Inches(0.45),
         size=16, bold=True, color=WHITE)

add_text(slide, 'TRES CLAIMS EN EL TIRO', Inches(4.5), Inches(2.75), Inches(8.6), Inches(0.3),
         size=9, bold=True, color=ORANGE)
claims_a = [
    'El mas rapido ante parpadeos (< 0,03 s)',
    'Protege tecnologia Inverter',
    'Sensor NTC incorporado*',
]
for i, c in enumerate(claims_a):
    add_text(slide, chr(0x25B6) + '  ' + c,
             Inches(4.5), Inches(3.1) + Pt(i * 32), Inches(8.6), Pt(28),
             size=13, color=LGREY)

add_text(slide, 'APUESTA PRINCIPAL', Inches(4.5), Inches(4.15), Inches(8.6), Inches(0.3),
         size=9, bold=True, color=ORANGE)
add_text(slide,
    'El comprador tecnico encuentra los tres argumentos que valida para recomendar el producto.',
    Inches(4.5), Inches(4.45), Inches(8.6), Inches(0.45),
    size=13, color=LGREY)

add_text(slide, 'A QUIEN FAVORECE MAS', Inches(4.5), Inches(5.05), Inches(8.6), Inches(0.3),
         size=9, bold=True, color=ORANGE)
add_text(slide, 'Instalador tecnico. Consumidor residencial solo si Atlas resuelve jerarquia visual muy clara.',
         Inches(4.5), Inches(5.35), Inches(8.6), Inches(0.45),
         size=13, color=LGREY)

add_text(slide, 'RIESGO PRINCIPAL', Inches(4.5), Inches(5.9), Inches(8.6), Inches(0.3),
         size=9, bold=True, color=ORANGE)
add_text(slide,
    'Densidad: cuatro elementos en blister pequeno. Conexion causal velocidad→inverter '
    'no es explicita en el texto — depende del diseno de Atlas.',
    Inches(4.5), Inches(6.2), Inches(8.6), Inches(0.55),
    size=13, color=LGREY)

add_text(slide, 'GATE: Cerrado (cero pendientes) | DISTANCIA DE PUNTO DE PARTIDA OWNER: Cero',
         Inches(0.55), Inches(7.08), Inches(12.5), Inches(0.32),
         size=9, color=MGREY)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — ALTERNATIVA B
# ─────────────────────────────────────────────────────────────────────────────
slide = new_slide()

add_label(slide, 'ALTERNATIVA B  —  PRIMARIA AURELIO + ATLAS', Inches(0.7), Inches(0.45))
orange_bar(slide, Inches(0.55), Inches(0.85))
add_text(slide, 'B: Los dos diferenciadores unicos dominan el tiro — inverter se desarrolla en el retiro',
         Inches(0.75), Inches(0.85), Inches(11.8), Inches(0.75),
         size=23, bold=True, color=WHITE)

add_divider(slide, Inches(1.75))

add_mockup(slide, MOCK_B, Inches(0.55), Inches(1.85), Inches(3.7), Inches(5.3))

add_text(slide, 'LENGUETA / CINTILLA', Inches(4.5), Inches(1.9), Inches(8.6), Inches(0.3),
         size=9, bold=True, color=ORANGE)
add_text(slide, 'NUEVO  |  LA PROTECCION MAS COMPLETA',
         Inches(4.5), Inches(2.2), Inches(8.6), Inches(0.45),
         size=16, bold=True, color=WHITE)
add_text(slide, '(Lengueta confirmada por Bruna §6 Opcion 3 — duda editorial v1 resuelta)',
         Inches(4.5), Inches(2.65), Inches(8.6), Inches(0.3),
         size=9, color=MGREY, italic=True)

add_text(slide, 'DOS CLAIMS EN EL TIRO', Inches(4.5), Inches(3.0), Inches(8.6), Inches(0.3),
         size=9, bold=True, color=ORANGE)
for i, c in enumerate(['El mas rapido ante parpadeos (< 0,03 s)  —  DATO DOMINANTE EN TIPOGRAFIA MAYOR',
                        'Sensor NTC incorporado*']):
    add_text(slide, chr(0x25B6) + '  ' + c,
             Inches(4.5), Inches(3.3) + Pt(i * 32), Inches(8.6), Pt(28),
             size=13, color=LGREY if i == 1 else WHITE)

add_text(slide, '"Protege tecnologia Inverter" migra al RETIRO con argumento causal completo',
         Inches(4.5), Inches(4.05), Inches(8.6), Inches(0.4),
         size=12, color=ORANGE, italic=True)

add_text(slide, 'APUESTA PRINCIPAL', Inches(4.5), Inches(4.55), Inches(8.6), Inches(0.3),
         size=9, bold=True, color=ORANGE)
add_text(slide,
    'El dato "< 0,03 s" tiene espacio para ser visualmente dominante. '
    'El instalador tecnico que lee el retiro encuentra el argumento de inverter con desarrollo causal completo.',
    Inches(4.5), Inches(4.85), Inches(8.6), Inches(0.55),
    size=13, color=LGREY)

add_text(slide, 'A QUIEN FAVORECE MAS', Inches(4.5), Inches(5.5), Inches(8.6), Inches(0.3),
         size=9, bold=True, color=ORANGE)
add_text(slide, 'Instalador tecnico (tiro + retiro). Menos efectiva para consumidor que solo ve el tiro.',
         Inches(4.5), Inches(5.8), Inches(8.6), Inches(0.4),
         size=13, color=LGREY)

add_text(slide, 'RIESGO PRINCIPAL', Inches(4.5), Inches(6.3), Inches(8.6), Inches(0.3),
         size=9, bold=True, color=ORANGE)
add_text(slide,
    '"Protege tecnologia Inverter" invisible para el comprador que no voltea el empaque. '
    'Condicion B-3-3: Atlas incluye elemento visual tiro→retiro.',
    Inches(4.5), Inches(6.6), Inches(8.6), Inches(0.45),
    size=13, color=LGREY)

add_text(slide, 'GATE: Cerrado | NOVEDAD V2: Duda editorial lengueta resuelta por Bruna §6 Opcion 3',
         Inches(0.55), Inches(7.08), Inches(12.5), Inches(0.32),
         size=9, color=MGREY)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — ALTERNATIVA C
# ─────────────────────────────────────────────────────────────────────────────
slide = new_slide()

add_label(slide, 'ALTERNATIVA C', Inches(0.7), Inches(0.45))
orange_bar(slide, Inches(0.55), Inches(0.85))
add_text(slide, 'C: Una frase que traduce la innovacion tecnica a consecuencia cotidiana',
         Inches(0.75), Inches(0.85), Inches(11.8), Inches(0.75),
         size=23, bold=True, color=WHITE)

add_divider(slide, Inches(1.75))

add_mockup(slide, MOCK_C, Inches(0.55), Inches(1.85), Inches(3.7), Inches(5.3))

add_text(slide, 'LENGUETA / CINTILLA', Inches(4.5), Inches(1.9), Inches(8.6), Inches(0.3),
         size=9, bold=True, color=ORANGE)
add_text(slide, 'NUEVO  |  EL MAS RAPIDO DE LA CATEGORIA',
         Inches(4.5), Inches(2.2), Inches(8.6), Inches(0.45),
         size=16, bold=True, color=WHITE)

add_text(slide, 'FRASE DOMINANTE', Inches(4.5), Inches(2.75), Inches(8.6), Inches(0.3),
         size=9, bold=True, color=ORANGE)
add_text(slide, '"Actua en < 0,03 s antes de que la fluctuacion llegue a tu equipo"',
         Inches(4.5), Inches(3.05), Inches(8.6), Inches(0.55),
         size=15, bold=True, color=WHITE)

add_text(slide, 'SUB-TEXTO COMPACTO', Inches(4.5), Inches(3.7), Inches(8.6), Inches(0.3),
         size=9, bold=True, color=ORANGE)
add_text(slide, 'Sensor NTC* + Proteccion Inverter incluidos',
         Inches(4.5), Inches(4.0), Inches(8.6), Inches(0.4),
         size=13, color=LGREY)

add_text(slide, 'APUESTA PRINCIPAL', Inches(4.5), Inches(4.5), Inches(8.6), Inches(0.3),
         size=9, bold=True, color=ORANGE)
add_text(slide,
    'La mas accesible para el consumidor final residencial. El vinculo < 0,03 s ↔ inverter '
    'esta en el texto, no en el diseno. Atlas confirmo factibilidad de espacio (2 lineas, tipografia 24px).',
    Inches(4.5), Inches(4.8), Inches(8.6), Inches(0.65),
    size=13, color=LGREY)

add_text(slide, 'A QUIEN FAVORECE MAS', Inches(4.5), Inches(5.55), Inches(8.6), Inches(0.3),
         size=9, bold=True, color=ORANGE)
add_text(slide, 'Consumidor final residencial. Valida para instalador tecnico con menor peso del dato cuantitativo.',
         Inches(4.5), Inches(5.85), Inches(8.6), Inches(0.4),
         size=13, color=LGREY)

add_text(slide, 'RIESGO PRINCIPAL', Inches(4.5), Inches(6.35), Inches(8.6), Inches(0.3),
         size=9, bold=True, color=ORANGE)
add_text(slide,
    'Menos tecnica — el instalador puede sentirla como publicidad. '
    'NTC e inverter en sub-texto pueden no leerse si el ojo se detiene en la frase dominante.',
    Inches(4.5), Inches(6.65), Inches(8.6), Inches(0.45),
    size=13, color=LGREY)

add_text(slide, 'GATE: Cerrado | DISTANCIA OWNER: Alta (arquitectura diferente)',
         Inches(0.55), Inches(7.08), Inches(12.5), Inches(0.32),
         size=9, color=MGREY)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11 — ALTERNATIVA D (nueva en v2)
# ─────────────────────────────────────────────────────────────────────────────
slide = new_slide()

# Special orange-tinted header for new slide
new_hdr = slide.shapes.add_shape(1, 0, 0, W, Inches(1.75))
new_hdr.fill.solid(); new_hdr.fill.fore_color.rgb = RGBColor(0x22, 0x0C, 0x00)
new_hdr.line.fill.background()

add_label(slide, 'ALTERNATIVA D  —  NUEVA EN V2  |  ALTERNATIVA EXPLÍCITA AURELIO', Inches(0.7), Inches(0.45))
orange_bar(slide, Inches(0.55), Inches(0.85))
add_text(slide, 'D: QR como eje + "Hasta 10 veces mas rapido" + flecha causal',
         Inches(0.75), Inches(0.85), Inches(11.8), Inches(0.75),
         size=23, bold=True, color=WHITE)

add_divider(slide, Inches(1.75))

add_mockup(slide, MOCK_D, Inches(0.55), Inches(1.85), Inches(3.7), Inches(5.3))

add_text(slide, 'LENGUETA / CINTILLA', Inches(4.5), Inches(1.9), Inches(8.6), Inches(0.3),
         size=9, bold=True, color=ORANGE)
add_text(slide, 'NUEVO  |  AVERIGUALO',
         Inches(4.5), Inches(2.2), Inches(8.6), Inches(0.45),
         size=16, bold=True, color=WHITE)
add_text(slide, '(Gate Bruna pendiente — bajo riesgo)',
         Inches(4.5), Inches(2.65), Inches(8.6), Inches(0.3),
         size=9, color=MGREY, italic=True)

add_text(slide, 'BLOQUE VISUAL UNITARIO EN ZONA DE TITULOS', Inches(4.5), Inches(3.0), Inches(8.6), Inches(0.3),
         size=9, bold=True, color=ORANGE)

claims_d = [
    ('Hasta 10 veces mas rapido  [QR]', WHITE, True),
    ('< 0,03 s ante parpadeos', LGREY, False),
    ('   ↓  Protege tecnologia Inverter', LGREY, False),
    ('   ↓  Sensor NTC: protege contra calentamiento*', LGREY, False),
]
for i, (txt, col, bld) in enumerate(claims_d):
    add_text(slide, txt,
             Inches(4.5), Inches(3.3) + Pt(i * 30), Inches(8.6), Pt(26),
             size=13, color=col, bold=bld)

add_text(slide, 'APUESTA PRINCIPAL', Inches(4.5), Inches(4.65), Inches(8.6), Inches(0.3),
         size=9, bold=True, color=ORANGE)
add_text(slide,
    '"Hasta 10 veces mas rapido" comunica la brecha competitiva de forma intuitiva para cualquier comprador. '
    'La flecha conectora hace visible la cadena causal sin texto adicional. '
    'El QR en zona de titulos es pionero en la categoria venezolana.',
    Inches(4.5), Inches(4.95), Inches(8.6), Inches(0.7),
    size=12, color=LGREY)

add_text(slide, 'CONDICION SINE QUA NON', Inches(4.5), Inches(5.75), Inches(8.6), Inches(0.3),
         size=9, bold=True, color=ORANGE)
add_text(slide,
    'D es viable SOLO SI los contenidos del QR (landing tecnica + video) estan producidos '
    'y publicados antes de imprimir. Sin ese contenido, D es una promesa vacia. '
    'Si la Junta elige D, comprometer workstream QR en ESTA SESION.',
    Inches(4.5), Inches(6.05), Inches(8.6), Inches(0.65),
    size=12, color=ORANGE)

add_text(slide, 'GATE: 2 elementos pendientes (bajo riesgo) | DISTANCIA OWNER: Alta | QR: bloqueante para produccion',
         Inches(0.55), Inches(7.08), Inches(12.5), Inches(0.32),
         size=9, color=MGREY)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 12 — COMPARACION LADO A LADO
# ─────────────────────────────────────────────────────────────────────────────
slide = new_slide()

add_label(slide, 'COMPARACION', Inches(0.7), Inches(0.45))
orange_bar(slide, Inches(0.55), Inches(0.85))
add_text(slide, 'B maximiza el diferenciador mas unico; D cuantifica la brecha; C maximiza accesibilidad',
         Inches(0.75), Inches(0.85), Inches(11.8), Inches(0.75),
         size=22, bold=True, color=WHITE)

add_divider(slide, Inches(1.75))

# Four mockups side by side
for i, (png, lbl) in enumerate([
        (MOCK_A, 'A'), (MOCK_B, 'B'), (MOCK_C, 'C'), (MOCK_D, 'D')]):
    lx = Inches(0.45 + i * 3.25)
    add_text(slide, lbl, lx, Inches(1.82), Inches(0.5), Inches(0.4),
             size=18, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    add_mockup(slide, png, lx, Inches(2.2), Inches(3.05), Inches(2.45))

# Compact comparison table
rows_c = [
    ['Dimension', 'A', 'B', 'C', 'D'],
    ['Claims en el tiro', '3 badges', '2 badges', '1 frase dominante', '3 claims + flecha + QR'],
    ['Dato protagonista visual', '< 0,03 s (72px)', '< 0,03 s (96px)', 'Frase narrativa', '"Hasta 10x" (52px)'],
    ['Claim comparativo cuantificado', 'No', 'No', 'No', 'Si ("Hasta 10x")'],
    ['Audiencia primaria', 'Instalador tecnico', 'Instalador tecnico', 'Consumidor final', 'Mixta'],
    ['Gate pendiente', 'Ninguno', 'Ninguno', 'Ninguno', '2 elementos (bajo riesgo)'],
    ['Dependencia QR', 'Baja', 'Baja-media', 'Baja', 'ALTA (bloqueante)'],
    ['Atlas score visual (de 40)', '25', '30', '36', '29'],
]
add_table_simple(slide, rows_c, Inches(0.45), Inches(4.8), Inches(12.4), Inches(2.5))


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 13 — CINCO TENSIONES
# ─────────────────────────────────────────────────────────────────────────────
slide = new_slide()

add_label(slide, 'TENSIONES A RESOLVER', Inches(0.7), Inches(0.45))
orange_bar(slide, Inches(0.55), Inches(0.85))
add_text(slide, 'Cinco preguntas que solo la Junta puede responder',
         Inches(0.75), Inches(0.85), Inches(11.8), Inches(0.75),
         size=26, bold=True, color=WHITE)

add_divider(slide, Inches(1.75))

tensions = [
    ('1', 'Densidad vs. respiro visual',
     'En el punto de venta donde se compra el GSM, el comprador tipico dedica 3-5 seg al frente?  '
     'Ferreteria de mostrador (densidad de A sostenible) o libre servicio (B o C mas efectivas; D apuesta por el escaneo)?',
     'A / B', 'C / D'),
    ('2', 'Superlativo cualitativo vs. cuantitativo vs. comparativo',
     '"La Proteccion mas completa" (A/B) vs. "El mas rapido de la categoria" (C) vs. "Hasta 10x mas rapido" (D). '
     'Los tres son defendibles y aprobados por Bruna. El tercero impacta mas pero requiere caveat.',
     'A / B', 'C / D'),
    ('3', 'Continuidad de marca vs. innovacion visible',
     'A es la variacion mas cercana al empaque conocido. D es la mas diferente. '
     'El tecnico debe ver evolucion (A/B) o un salto que capture compradores nuevos (C/D)?',
     'A / B', 'C / D'),
    ('4', 'Audiencia tecnica vs. audiencia residencial',
     'El empaque es uno solo para todos los modelos con copy tecnico (A/B) '
     'o se optimiza para perfil mixto (C/D)?',
     'A / B', 'C / D'),
    ('5 — NUEVA V2', 'Superlativo no cuantificado vs. claim con dato comparativo (10x)',
     '"El mas rapido" (A/B/C) dice que Genteca gana sin decir por cuanto. '
     '"Hasta 10x mas rapido" (D) dice que la brecha es enorme — mas persuasivo para el instalador que ya sabe que la velocidad importa. '
     'Riesgo: Breakermatic es 2x, no 10x; el comprador que lo conoce puede sentir que el claim esta exagerado.',
     'A / B / C', 'D'),
]

for i, (num, title, body, fav_left, fav_right) in enumerate(tensions):
    ty = Inches(1.9) + Pt(i * 76)
    add_text(slide, num, Inches(0.55), ty, Inches(0.5), Pt(32),
             size=16, bold=True, color=ORANGE)
    add_text(slide, title, Inches(1.05), ty, Inches(7.5), Pt(28),
             size=13, bold=True, color=WHITE)
    add_text(slide, body, Inches(1.05), ty + Pt(28), Inches(7.5), Pt(42),
             size=10, color=LGREY)
    add_text(slide, 'Apunta a: ' + fav_left,
             Inches(8.65), ty, Inches(2.1), Pt(28),
             size=9, color=ORANGE, align=PP_ALIGN.RIGHT)
    add_text(slide, 'o: ' + fav_right,
             Inches(10.85), ty, Inches(1.9), Pt(28),
             size=9, color=LGREY, align=PP_ALIGN.RIGHT)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 14 — RECOMENDACION DE LA CADENA
# ─────────────────────────────────────────────────────────────────────────────
slide = new_slide()

add_label(slide, 'RECOMENDACION', Inches(0.7), Inches(0.45))
orange_bar(slide, Inches(0.55), Inches(0.85))
add_text(slide, 'Aurelio y Atlas recomiendan B — y presentan D como alternativa explicita',
         Inches(0.75), Inches(0.85), Inches(11.8), Inches(0.75),
         size=24, bold=True, color=WHITE)

add_divider(slide, Inches(1.75))

# Aurelio box
box_a = slide.shapes.add_shape(1, Inches(0.55), Inches(1.9), Inches(6.0), Inches(4.85))
box_a.fill.solid(); box_a.fill.fore_color.rgb = RGBColor(0x1C, 0x1C, 0x1C)
box_a.line.color.rgb = ORANGE

add_text(slide, 'AURELIO — Content Strategist', Inches(0.75), Inches(2.0), Inches(5.6), Inches(0.32),
         size=9, bold=True, color=ORANGE)
add_text(slide, 'Primaria: Alternativa B',
         Inches(0.75), Inches(2.32), Inches(5.6), Inches(0.45),
         size=18, bold=True, color=WHITE)
add_text(slide,
    'Los dos diferenciadores unicos en el tiro con el dato dominante. '
    'El argumento de inverter se desarrolla en el retiro con argumento causal completo. '
    'Arquitectura mas limpia que D. Sin gates pendientes. Sin dependencia de contenido digital. '
    'Con el dato Breakermatic corregido, "el mas rapido ante parpadeos" sigue siendo correcto y '
    'tiene respaldo documental verificado en laboratorio Genteca.',
    Inches(0.75), Inches(2.8), Inches(5.6), Inches(2.0),
    size=12, color=LGREY)

add_text(slide, 'Alternativa explicita: D — si la Junta prioriza statement de innovacion',
         Inches(0.75), Inches(4.85), Inches(5.6), Inches(0.4),
         size=13, bold=True, color=ORANGE)
add_text(slide,
    'D es viable si — y solo si — los contenidos del QR se comprometen en esta misma sesion. '
    'Sin ese compromiso, D no es presentable en punto de venta.',
    Inches(0.75), Inches(5.28), Inches(5.6), Inches(0.65),
    size=12, color=LGREY)

# Atlas box
box_b = slide.shapes.add_shape(1, Inches(6.8), Inches(1.9), Inches(6.0), Inches(4.85))
box_b.fill.solid(); box_b.fill.fore_color.rgb = RGBColor(0x1C, 0x1C, 0x1C)
box_b.line.color.rgb = MGREY

add_text(slide, 'ATLAS — Static Visual Production Lead', Inches(7.0), Inches(2.0), Inches(5.6), Inches(0.32),
         size=9, bold=True, color=ORANGE)
add_text(slide, 'Primaria: Alternativa B-3',
         Inches(7.0), Inches(2.32), Inches(5.6), Inches(0.45),
         size=18, bold=True, color=WHITE)
add_text(slide,
    'B-3 sigue siendo la apuesta visual mas solida para canal ferreteria tecnica. '
    '"< 0,03 s" a 96px es el elemento con mayor poder de detencion del ojo en el punto de venta. '
    'No exige que el comprador escanee nada para captar el mensaje central. '
    'Duda editorial sobre la lengueta quedo resuelta por Bruna §6 Opcion 3.',
    Inches(7.0), Inches(2.8), Inches(5.6), Inches(2.0),
    size=12, color=LGREY)

add_text(slide, 'Matiz nuevo sobre D:',
         Inches(7.0), Inches(4.85), Inches(5.6), Inches(0.32),
         size=12, bold=True, color=ORANGE)
add_text(slide,
    'D plantea un caso visual genuinamente fuerte: convierte el empaque en una invitacion activa, '
    'no en un catalogo de especificaciones. "Hasta 10x mas rapido" es mas intuitivo que "< 0,03 s" '
    'para el comprador sin referencia de milisegundos. Atlas considera B y D equivalentes en fuerza '
    'visual si las condiciones de QR estan resueltas.',
    Inches(7.0), Inches(5.18), Inches(5.6), Inches(1.4),
    size=12, color=LGREY)

add_text(slide, 'Si la Junta pide una tercera via entre B y D: Alternativa A sigue siendo correcta y ratificada por el Owner.',
         Inches(0.55), Inches(7.0), Inches(12.1), Inches(0.42),
         size=11, color=MGREY, italic=True)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 15 — QR COMO VEHICULO DE COMUNICACION PROFUNDA
# ─────────────────────────────────────────────────────────────────────────────
slide = new_slide()

add_label(slide, 'DECISION ADICIONAL — QR', Inches(0.7), Inches(0.45))
orange_bar(slide, Inches(0.55), Inches(0.85))
add_text(slide, 'El QR: complemento en A/B/C — argumento central en D',
         Inches(0.75), Inches(0.85), Inches(11.8), Inches(0.75),
         size=26, bold=True, color=WHITE)

add_divider(slide, Inches(1.75))

# Left: A/B/C QR
box_l = slide.shapes.add_shape(1, Inches(0.55), Inches(1.9), Inches(5.9), Inches(4.5))
box_l.fill.solid(); box_l.fill.fore_color.rgb = RGBColor(0x1C, 0x1C, 0x1C)
box_l.line.color.rgb = MGREY

add_text(slide, 'EN A / B / C — QR COMO COMPLEMENTO', Inches(0.75), Inches(2.0), Inches(5.5), Inches(0.32),
         size=9, bold=True, color=ORANGE)
add_text(slide, 'No bloqueante para produccion del empaque.',
         Inches(0.75), Inches(2.32), Inches(5.5), Inches(0.32),
         size=12, bold=True, color=LGREY)

qr_abc = [
    'Frases candidatas (Aurelio recomienda Candidata 1):',
    '  C1 (tecnica): "Por que < 0,03 s? Escanea aqui" — alta conversion con instalador',
    '  C2 (beneficio): "Descubre como protege tu equipo"',
    '  C3 (enigma): "Tecnologia NTC explicada aqui"',
    '',
    'Contenidos defendibles segun Paxs (pueden ir al QR):',
    '  + Equipos inverter traen MOV — protegen picos, no parpadeos',
    '  + Brecha reconocida en IEC 61000-4-11',
    '  + Protector < 30 ms actua dentro del hold-up time del bus DC',
    '  + Sensor NTC detecta fallas silenciosas del protector',
    '  + Fabricantes como Mitsubishi recomiendan estabilizador en red inestable',
]
for i, line in enumerate(qr_abc):
    add_text(slide, line, Inches(0.75), Inches(2.7) + Pt(i * 22), Inches(5.5), Pt(20),
             size=10, color=LGREY if not line.startswith('  +') else LGREY,
             bold=(not line.startswith(' ') and ':' in line))

# Right: D QR
box_r = slide.shapes.add_shape(1, Inches(6.7), Inches(1.9), Inches(6.1), Inches(4.5))
box_r.fill.solid(); box_r.fill.fore_color.rgb = RGBColor(0x22, 0x0C, 0x00)
box_r.line.color.rgb = ORANGE

add_text(slide, 'EN D — QR COMO EJE CENTRAL', Inches(6.9), Inches(2.0), Inches(5.7), Inches(0.32),
         size=9, bold=True, color=ORANGE)
add_text(slide, 'La lengueta "AVERIGUALO" y el QR son el sistema de comunicacion de D.',
         Inches(6.9), Inches(2.32), Inches(5.7), Inches(0.45),
         size=12, bold=True, color=WHITE)

qr_d = [
    'Contenidos que debe desarrollar el QR de D:',
    '  + Narrativa flickers vs. picos completa (Paxs §2-§4)',
    '  + Mecanismo de dano en tarjeta inverter: bus DC capacitor + IGBT stress',
    '  + Comparativa tecnica detallada vs. Breakermatic, WellSpec, Powest',
    '  + Argumento "proteccion autentica inverter" (rechazado para el frente)',
    '  + Video explicativo 60 seg para consumidor final',
    '  + Version extendida 2-3 min para instalador tecnico',
    '',
    'Argumentos que NO se sostienen bajo cuestionamiento (anti-mensajes):',
    '  X "Los competidores con > 30 ms NO protegen los equipos inverter"',
    '  X "Solo < 30 ms es proteccion autentica"',
    '  X "Los equipos inverter se danan porque no tienen proteccion interna"',
]
for i, line in enumerate(qr_d):
    color = ORANGE if line.startswith('  X') else LGREY
    add_text(slide, line, Inches(6.9), Inches(2.85) + Pt(i * 22), Inches(5.7), Pt(20),
             size=10, color=color,
             bold=(':' in line and not line.startswith(' ')))

add_text(slide,
    'DECISION PENDIENTE: que cubre hoy el QR existente del empaque? (manual, garantia, web generica) — desbloquea decision de opcion de QR para cualquier alternativa.',
    Inches(0.55), Inches(6.55), Inches(12.1), Inches(0.55),
    size=10, color=ORANGE, italic=True)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 16 — VENTA CRUZADA EXCELINE PROFESIONAL
# ─────────────────────────────────────────────────────────────────────────────
slide = new_slide()

add_label(slide, 'DECISION ADICIONAL — VENTA CRUZADA', Inches(0.7), Inches(0.45))
orange_bar(slide, Inches(0.55), Inches(0.85))
add_text(slide, 'El disclaimer del empaque es una palanca comercial hacia Exceline Profesional',
         Inches(0.75), Inches(0.85), Inches(11.8), Inches(0.75),
         size=24, bold=True, color=WHITE)

add_divider(slide, Inches(1.75))

# Left column
add_text(slide, 'EL DISCLAIMER EN EL EMPAQUE', Inches(0.55), Inches(1.9), Inches(6.0), Inches(0.32),
         size=9, bold=True, color=ORANGE)

disc_box = slide.shapes.add_shape(1, Inches(0.55), Inches(2.25), Inches(6.0), Inches(0.75))
disc_box.fill.solid(); disc_box.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x22)
disc_box.line.color.rgb = ORANGE
add_text(slide, '"No reemplaza los breakers termomagneticos de la instalacion electrica"',
         Inches(0.7), Inches(2.32), Inches(5.7), Inches(0.6),
         size=13, bold=True, color=WHITE, italic=True)

add_text(slide,
    'Este disclaimer — preservado desde v3 del empaque y presente en las cuatro alternativas — '
    'reconoce que el GSM es protector de voltaje, no protector de sobrecarga. '
    'Genteca Exceline Profesional tiene breakers termomagneticos en catalogo.',
    Inches(0.55), Inches(3.1), Inches(6.0), Inches(1.1),
    size=13, color=LGREY)

add_text(slide, 'LA PALANCA COMERCIAL', Inches(0.55), Inches(4.3), Inches(6.0), Inches(0.32),
         size=9, bold=True, color=ORANGE)
add_text(slide,
    'El disclaimer naturaliza la necesidad del breaker termomagnetico. '
    'Si el GSM incluye en el QR y en el argumentario un modulo de "instalacion completa con Exceline" '
    '— GSM + breakers Exceline Profesional — el disclaimer se convierte en un llamado a la accion '
    'comercial en lugar de una advertencia pasiva.',
    Inches(0.55), Inches(4.65), Inches(6.0), Inches(1.55),
    size=13, color=LGREY)

# Right column
add_text(slide, 'DECISION QUE LA JUNTA DEBE REGISTRAR', Inches(6.8), Inches(1.9), Inches(6.0), Inches(0.32),
         size=9, bold=True, color=ORANGE)

dec_box = slide.shapes.add_shape(1, Inches(6.8), Inches(2.25), Inches(6.0), Inches(2.0))
dec_box.fill.solid(); dec_box.fill.fore_color.rgb = RGBColor(0x22, 0x0C, 0x00)
dec_box.line.color.rgb = ORANGE
add_text(slide,
    'Autoriza la Junta el modulo de "instalacion completa con Exceline" en los materiales de lanzamiento digitales?\n\n'
    'Canal: QR (en cualquier alternativa) + argumentario ESC + landing de innovaciones.',
    Inches(6.95), Inches(2.35), Inches(5.7), Inches(1.8),
    size=13, color=WHITE)

add_text(slide, 'LO QUE INCLUIRIA EL MODULO', Inches(6.8), Inches(4.35), Inches(6.0), Inches(0.32),
         size=9, bold=True, color=ORANGE)
module_items = [
    'Por que el GSM no reemplaza el breaker (explicacion honesta del disclaimer)',
    'Que tipo de instalacion necesita GSM + breaker complementario',
    'Exceline Profesional: referencias de breakers termomagneticos compatibles',
    'Llamado a la accion: "Consulta a tu electricista sobre tu instalacion completa"',
]
for i, item in enumerate(module_items):
    add_text(slide, chr(0x2022) + '  ' + item,
             Inches(6.8), Inches(4.72) + Pt(i * 30), Inches(6.0), Pt(28),
             size=12, color=LGREY)

add_text(slide, 'Este modulo no bloquea ninguna alternativa — es decision independiente de la eleccion A/B/C/D.',
         Inches(0.55), Inches(7.0), Inches(12.1), Inches(0.42),
         size=11, color=MGREY, italic=True)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 17 — PENDIENTES OPERATIVOS
# ─────────────────────────────────────────────────────────────────────────────
slide = new_slide()

add_label(slide, 'PENDIENTES OPERATIVOS', Inches(0.7), Inches(0.45))
orange_bar(slide, Inches(0.55), Inches(0.85))
add_text(slide, 'El datasheet de I&D bloquea imprenta — pero no bloquea la decision de hoy',
         Inches(0.75), Inches(0.85), Inches(11.8), Inches(0.75),
         size=24, bold=True, color=WHITE)

add_divider(slide, Inches(1.75))

rows_p = [
    ['Pendiente', 'Responsable', 'Impacto', 'Aplica a'],
    ['P-1: Datasheet I&D actualizado con < 30 ms', 'Owner + I&D Genteca', 'BLOQUEANTE imprenta — sin excepcion', 'A / B / C / D'],
    ['P-2: Bruna refresh "AVERIGUALO" + NTC D', 'Bruna + Solenne', 'Bloqueante arte final — bajo riesgo', 'Solo D'],
    ['P-3: Contenidos QR producidos y publicados', 'Owner (asignar Nerea + Luma/Vela)', 'BLOQUEANTE imprenta de D', 'Solo D'],
    ['P-4: Caveat "Hasta 10x" en retiro — texto literal Bruna §7.2 a Oz', 'Oz (instruccion Owner)', 'Condicion operativa de D', 'Solo D'],
    ['P-5: Certificacion externa < 30 ms (FONDONORMA/UL/IEC)', 'Raoul + Liliam I&D + Vera', 'No bloquea — fortalece argumento', 'A / B / C / D'],
    ['P-6: Decision QR existente del empaque (funcion actual)', 'Raoul', 'Desbloquea diseno QR en cualquier alternativa', 'A / B / C / D'],
    ['P-7: Atlas confirma condicion diseno B-3-3 (hilo tiro→retiro)', 'Atlas + Oz', 'Condicion de diseno menor', 'Solo B'],
]
add_table_simple(slide, rows_p, Inches(0.55), Inches(1.9), Inches(12.2), Inches(4.3))

add_text(slide,
    'Ninguno de estos pendientes bloquea la decision de Junta de hoy. '
    'P-1, P-3 y P-6 deben quedar asignados antes de salir de esta sesion.',
    Inches(0.55), Inches(6.35), Inches(12.2), Inches(0.55),
    size=12, color=LGREY, italic=True)

add_text(slide, 'INSTRUCCION A OZ solo se emite tras: (1) decision Junta + (2) datasheet I&D confirmado.',
         Inches(0.55), Inches(7.0), Inches(12.2), Inches(0.42),
         size=11, bold=True, color=ORANGE)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 18 — PROXIMOS PASOS
# ─────────────────────────────────────────────────────────────────────────────
slide = new_slide()

add_label(slide, 'PROXIMOS PASOS', Inches(0.7), Inches(0.45))
orange_bar(slide, Inches(0.55), Inches(0.85))
add_text(slide, 'La decision de Junta hoy desbloquea el Paso 9 del pipeline CSC',
         Inches(0.75), Inches(0.85), Inches(11.8), Inches(0.75),
         size=26, bold=True, color=WHITE)

add_divider(slide, Inches(1.75))

decisions = [
    ('ELIGE A',
     'Solenne: delta v5 copy definitivo A.\nAtlas: mockup final A con jerarquia visual clara.\nOz: instrucciones + datasheet I&D (P-1).',
     ORANGE),
    ('ELIGE B',
     'Solenne: delta v5.\nAtlas: confirma condicion diseno B-3-3.\nOz: instrucciones + datasheet I&D (P-1).',
     ORANGE),
    ('ELIGE C',
     'Atlas: confirma espacio tipografico frase dominante.\nSolenne: delta v5 si factible.\nOz: instrucciones + datasheet I&D (P-1).',
     LGREY),
    ('ELIGE D',
     'Owner: compromete QR (quien + plazo + QR existente).\nBruna: confirma gates lengueta + NTC.\nSolenne: delta v5 con caveats D.\nOz: instrucciones SOLO tras QR listo.',
     ORANGE),
    ('PIDE AJUSTES',
     'Solenne: SO-1 v3 con ajustes especificos.\nGate Bruna si los ajustes afectan claims.\nNueva ronda antes de instruccion Oz.',
     LGREY),
]

for i, (dec, steps, col) in enumerate(decisions):
    lx = Inches(0.55 + i * 2.56)
    box = slide.shapes.add_shape(1, lx, Inches(1.95), Inches(2.45), Inches(4.8))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x1E)
    box.line.color.rgb = col
    add_text(slide, dec, lx + Inches(0.1), Inches(2.05), Inches(2.25), Inches(0.45),
             size=15, bold=True, color=col)
    add_text(slide, steps, lx + Inches(0.1), Inches(2.55), Inches(2.25), Inches(3.8),
             size=11, color=LGREY)

# Common actions
rec_box = slide.shapes.add_shape(1, Inches(0.55), Inches(6.88), Inches(12.2), Inches(0.5))
rec_box.fill.solid(); rec_box.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x22)
rec_box.line.color.rgb = ORANGE
add_text(slide,
    'ACCIONES COMUNES A CUALQUIER DECISION:  '
    '(1) Confirmar funcion QR existente (P-6)  |  '
    '(2) Asignar P-1 datasheet I&D  |  '
    '(3) Registrar autorizacion venta cruzada Exceline Profesional',
    Inches(0.75), Inches(6.92), Inches(12.0), Inches(0.44),
    size=10, bold=True, color=WHITE)


# ─────────────────────────────────────────────────────────────────────────────
# SAVE
# ─────────────────────────────────────────────────────────────────────────────
prs.save(OUT_FILE)
print(f'Saved: {OUT_FILE}')
print(f'Slides: {len(prs.slides)}')
