import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from lxml import etree
import copy
import os

# ── Palette ────────────────────────────────────────────────────────────────────
NAVY       = RGBColor(0x0D, 0x1F, 0x3C)   # dark navy — background slides
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE     = RGBColor(0xFF, 0x8C, 0x00)   # Exceline/Genteca accent
LIGHT_GREY = RGBColor(0xF2, 0xF4, 0xF7)   # content slide background
MID_GREY   = RGBColor(0x4A, 0x5A, 0x6B)   # body text on light bg
DARK_TEXT  = RGBColor(0x0D, 0x1F, 0x3C)   # headings on light bg
DIVIDER    = RGBColor(0xD0, 0xD5, 0xDD)   # subtle divider line

# ── Slide dimensions (16:9 widescreen) ─────────────────────────────────────────
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK_LAYOUT = prs.slide_layouts[6]   # completely blank

IMG_PATH = r"C:\Raul\03-projects\genteca\2026-04_GSM-MB-RB-RF_empaque\03-review-and-release\Junta_PPTX_v2\Atlas_mockup_frente_B.png"
OUTPUT   = r"C:\Raul\03-projects\genteca\2026-04_GSM-MB-RB-RF_empaque\03-review-and-release\Junta_GSM_empaque_B-sin-NTC_propuesta.pptx"

# ── Helper functions ────────────────────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(BLANK_LAYOUT)

def fill_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill_color=None, line_color=None, line_width=Pt(0)):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        l, t, w, h
    )
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width if line_width else Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, l, t, w, h, text, font_name="Calibri", font_size=Pt(14),
                bold=False, color=WHITE, align=PP_ALIGN.LEFT,
                wrap=True, italic=False, para_space_before=Pt(0), line_spacing=None):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_before = para_space_before
    if line_spacing:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox

def add_multi_para_textbox(slide, l, t, w, h, paras, font_name="Calibri",
                            default_size=Pt(13), default_color=MID_GREY,
                            default_align=PP_ALIGN.LEFT, wrap=True):
    """
    paras: list of dicts with keys:
      text, size, bold, color, align, italic, space_before, line_spacing, bullet (bool)
    """
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    first = True
    for pd in paras:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = pd.get('align', default_align)
        if pd.get('space_before'):
            p.space_before = pd['space_before']
        if pd.get('line_spacing'):
            p.line_spacing = pd['line_spacing']
        run = p.add_run()
        run.text = pd.get('text', '')
        run.font.name = font_name
        run.font.size = pd.get('size', default_size)
        run.font.bold = pd.get('bold', False)
        run.font.color.rgb = pd.get('color', default_color)
        run.font.italic = pd.get('italic', False)
    return txBox

def accent_bar(slide, top=Inches(0.12)):
    """Thin orange top bar across full width."""
    add_rect(slide, Inches(0), top, SLIDE_W, Inches(0.06), fill_color=ORANGE)

def slide_number_label(slide, num, total=10):
    add_textbox(slide, Inches(12.3), Inches(7.1), Inches(0.9), Inches(0.28),
                f"{num} / {total}", font_size=Pt(9), color=RGBColor(0xAA, 0xB0, 0xBB),
                align=PP_ALIGN.RIGHT)

def section_header_bg(slide):
    """Light grey background for content slides."""
    fill_bg(slide, LIGHT_GREY)

def dark_slide_bg(slide):
    fill_bg(slide, NAVY)

def add_divider(slide, top, color=DIVIDER):
    add_rect(slide, Inches(0.6), top, Inches(12.13), Pt(1.2), fill_color=color)

# ── SLIDE 1 — Cover ─────────────────────────────────────────────────────────────
s1 = add_slide()
dark_slide_bg(s1)
accent_bar(s1, top=Inches(0))

# Full-width orange block at top
add_rect(s1, Inches(0), Inches(0), SLIDE_W, Inches(0.55), fill_color=ORANGE)

# Logo area placeholder (top-left)
add_textbox(s1, Inches(0.45), Inches(0.1), Inches(3.5), Inches(0.38),
            "GENTECA  |  EXCELINE GSM",
            font_size=Pt(12), bold=True, color=NAVY, align=PP_ALIGN.LEFT)

# Main title
add_textbox(s1, Inches(1.2), Inches(2.1), Inches(10.9), Inches(1.3),
            "Empaque Línea Exceline GSM",
            font_size=Pt(42), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_textbox(s1, Inches(1.2), Inches(3.35), Inches(10.9), Inches(0.85),
            "Propuesta Refinada — Alternativa B",
            font_size=Pt(26), bold=False, color=ORANGE, align=PP_ALIGN.CENTER)

add_textbox(s1, Inches(1.2), Inches(4.15), Inches(10.9), Inches(0.5),
            "Modelos GSM-MB · GSM-RB · GSM-RF · GSM-RE",
            font_size=Pt(16), bold=False, color=RGBColor(0xCC, 0xD3, 0xDE), align=PP_ALIGN.CENTER)

# Date and session label
add_textbox(s1, Inches(1.2), Inches(5.3), Inches(10.9), Inches(0.38),
            "Junta Directiva — 6 de mayo de 2026",
            font_size=Pt(13), bold=False, color=RGBColor(0xAA, 0xB4, 0xC4), align=PP_ALIGN.CENTER)

# Bottom accent strip
add_rect(s1, Inches(0), Inches(7.3), SLIDE_W, Inches(0.2), fill_color=ORANGE)
slide_number_label(s1, 1)

# ── SLIDE 2 — Contexto ──────────────────────────────────────────────────────────
s2 = add_slide()
section_header_bg(s2)
accent_bar(s2)

add_textbox(s2, Inches(0.6), Inches(0.25), Inches(11), Inches(0.55),
            "Contexto — De dónde venimos",
            font_size=Pt(22), bold=True, color=DARK_TEXT, align=PP_ALIGN.LEFT)
add_divider(s2, Inches(0.85))

bullets_s2 = [
    {"text": "En la sesión anterior la Junta evaluó cuatro alternativas de empaque para la línea Exceline GSM (modelos MB, RB, RF, RE) y seleccionó la Alternativa B.",
     "size": Pt(14.5), "color": MID_GREY, "space_before": Pt(14)},
    {"text": "La Alternativa B posiciona la velocidad de respuesta ante parpadeos (< 0,03 s) como argumento diferenciador primario y añade un segundo claim de función térmica en el frente del empaque.",
     "size": Pt(14.5), "color": MID_GREY, "space_before": Pt(10)},
    {"text": "Un miembro de la Junta señaló que la formulación original del segundo claim — \"Sensor NTC incorporado*\" — comunica el componente técnico interno en lugar de la función que recibe el comprador, y expone información de propiedad industrial sin agregar valor perceptible al consumidor.",
     "size": Pt(14.5), "color": DARK_TEXT, "bold": True, "space_before": Pt(10)},
    {"text": "Esa observación fue trabajada en detalle por el equipo de marketing y diseño. Esta presentación expone el resultado: una versión refinada de la Alternativa B con un único cambio quirúrgico en el frente.",
     "size": Pt(14.5), "color": MID_GREY, "space_before": Pt(10)},
]
add_multi_para_textbox(s2, Inches(0.6), Inches(1.0), Inches(12.13), Inches(5.6), bullets_s2)

add_textbox(s2, Inches(0.6), Inches(6.8), Inches(12), Inches(0.4),
            "Todo lo demás de la Alternativa B permanece intacto — este es un cambio de una sola frase.",
            font_size=Pt(12), color=ORANGE, bold=True, align=PP_ALIGN.LEFT, italic=True)
slide_number_label(s2, 2)

# ── SLIDE 3 — El cambio en una frase ───────────────────────────────────────────
s3 = add_slide()
dark_slide_bg(s3)
accent_bar(s3, top=Inches(0))
add_rect(s3, Inches(0), Inches(0), SLIDE_W, Inches(0.55), fill_color=ORANGE)
add_textbox(s3, Inches(0.45), Inches(0.1), Inches(5), Inches(0.38),
            "GENTECA  |  EXCELINE GSM", font_size=Pt(12), bold=True, color=NAVY)

add_textbox(s3, Inches(0.6), Inches(0.7), Inches(12), Inches(0.55),
            "El cambio en una frase",
            font_size=Pt(22), bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_divider(s3, Inches(1.3), color=RGBColor(0x2A, 0x3F, 0x60))

# Before / After boxes
# ANTES
add_rect(s3, Inches(0.6), Inches(1.65), Inches(5.7), Inches(2.6),
         fill_color=RGBColor(0x1C, 0x2E, 0x4A), line_color=RGBColor(0x55, 0x66, 0x80), line_width=Pt(1))
add_textbox(s3, Inches(0.75), Inches(1.75), Inches(2.5), Inches(0.38),
            "ANTES", font_size=Pt(10), bold=True, color=RGBColor(0xAA, 0xB4, 0xC4),
            align=PP_ALIGN.LEFT)
add_textbox(s3, Inches(0.75), Inches(2.2), Inches(5.3), Inches(0.8),
            "Sensor NTC incorporado*",
            font_size=Pt(22), bold=True, color=RGBColor(0xCC, 0xCC, 0xCC),
            align=PP_ALIGN.LEFT)
add_textbox(s3, Inches(0.75), Inches(3.05), Inches(5.3), Inches(1.0),
            "Comunica el componente interno. Expone información de propiedad industrial sin beneficio adicional para el comprador.",
            font_size=Pt(12), color=RGBColor(0x99, 0xA8, 0xBB), align=PP_ALIGN.LEFT)

# Arrow
add_textbox(s3, Inches(6.5), Inches(2.6), Inches(0.9), Inches(0.6),
            "→", font_size=Pt(36), bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

# AHORA
add_rect(s3, Inches(7.0), Inches(1.65), Inches(5.7), Inches(2.6),
         fill_color=RGBColor(0x0D, 0x2A, 0x1A), line_color=ORANGE, line_width=Pt(1.5))
add_textbox(s3, Inches(7.15), Inches(1.75), Inches(2.5), Inches(0.38),
            "AHORA", font_size=Pt(10), bold=True, color=ORANGE, align=PP_ALIGN.LEFT)
add_textbox(s3, Inches(7.15), Inches(2.2), Inches(5.3), Inches(0.8),
            "Autoprotección térmica activa*",
            font_size=Pt(22), bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_textbox(s3, Inches(7.15), Inches(3.05), Inches(5.3), Inches(1.0),
            "Comunica la función. El comprador entiende el beneficio. El competidor no recibe información sobre el componente.",
            font_size=Pt(12), color=RGBColor(0xBB, 0xCC, 0xBB), align=PP_ALIGN.LEFT)

# Razon central
add_rect(s3, Inches(0.6), Inches(4.55), Inches(12.13), Inches(1.6),
         fill_color=RGBColor(0x15, 0x28, 0x45))
add_textbox(s3, Inches(0.85), Inches(4.65), Inches(11.7), Inches(0.4),
            "PRINCIPIO", font_size=Pt(10), bold=True, color=ORANGE, align=PP_ALIGN.LEFT)
add_textbox(s3, Inches(0.85), Inches(5.05), Inches(11.7), Inches(0.9),
            "Un claim de empaque debe comunicar lo que el producto hace por el comprador, no cómo está construido internamente. La reformulación preserva toda la protección comercial y elimina la exposición de ingeniería.",
            font_size=Pt(13), color=WHITE, align=PP_ALIGN.LEFT)

add_rect(s3, Inches(0), Inches(7.3), SLIDE_W, Inches(0.2), fill_color=ORANGE)
slide_number_label(s3, 3)

# ── SLIDE 4 — Por qué esta formulación gana ─────────────────────────────────────
s4 = add_slide()
section_header_bg(s4)
accent_bar(s4)

add_textbox(s4, Inches(0.6), Inches(0.25), Inches(12), Inches(0.55),
            "Por qué \"Autoprotección térmica activa*\" gana sobre las alternativas",
            font_size=Pt(20), bold=True, color=DARK_TEXT, align=PP_ALIGN.LEFT)
add_divider(s4, Inches(0.85))

reasons = [
    ("Cero revelación técnica",
     "La formulación no menciona componente alguno. Un competidor que quiera copiar la frase tendrá que implementar la función para sostener el claim — lo que eleva el estándar de la categoría en lugar de regalar un secreto industrial."),
    ("Lenguaje del comprador, no del ingeniero",
     "\"Autoprotección\" es vocabulario que el instalador técnico y el consumidor residencial entienden sin formación especializada. \"Sensor NTC\" requiere conocimiento electrónico que la mayoría del mercado no tiene."),
    ("Territorio competitivo libre en el mercado venezolano",
     "Ningún competidor venezolano usa esta formulación. El espacio está vacante. La marca puede apropiarse del territorio de forma defendible."),
    ("Prefijo \"auto\" corrige el malentendido más frecuente",
     "Cuatro de cada diez compradores interpretan \"protección térmica\" como protección del equipo conectado. El prefijo \"auto\" aclara que la protección es del propio protector y la instalación, antes de que el asterisco sea leído."),
]

top = Inches(1.1)
for i, (title, body) in enumerate(reasons):
    # Number circle placeholder
    add_rect(s4, Inches(0.6), top + Inches(i * 1.4), Inches(0.38), Inches(0.38),
             fill_color=ORANGE)
    add_textbox(s4, Inches(0.62), top + Inches(i * 1.4) - Pt(2), Inches(0.35), Inches(0.4),
                str(i + 1), font_size=Pt(13), bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_textbox(s4, Inches(1.15), top + Inches(i * 1.4) - Pt(3), Inches(11.2), Inches(0.35),
                title, font_size=Pt(14), bold=True, color=DARK_TEXT, align=PP_ALIGN.LEFT)
    add_textbox(s4, Inches(1.15), top + Inches(i * 1.4) + Inches(0.32), Inches(11.2), Inches(0.75),
                body, font_size=Pt(12.5), color=MID_GREY, align=PP_ALIGN.LEFT)

slide_number_label(s4, 4)

# ── SLIDE 5 — Frente actualizado (mockup) ───────────────────────────────────────
s5 = add_slide()
dark_slide_bg(s5)
accent_bar(s5, top=Inches(0))
add_rect(s5, Inches(0), Inches(0), SLIDE_W, Inches(0.55), fill_color=ORANGE)
add_textbox(s5, Inches(0.45), Inches(0.1), Inches(5), Inches(0.38),
            "GENTECA  |  EXCELINE GSM", font_size=Pt(12), bold=True, color=NAVY)

add_textbox(s5, Inches(0.6), Inches(0.65), Inches(12), Inches(0.5),
            "El frente actualizado — Alternativa B refinada",
            font_size=Pt(20), bold=True, color=WHITE, align=PP_ALIGN.LEFT)

# Insert mockup image centered
img_w = Inches(3.8)
img_h = Inches(5.8)
img_l = (SLIDE_W - img_w) / 2
img_t = Inches(1.2)
s5.shapes.add_picture(IMG_PATH, img_l, img_t, img_w, img_h)

# Caption below image
add_textbox(s5, Inches(0.6), Inches(7.08), Inches(12.13), Inches(0.3),
            "Mockup frente Alternativa B — claim térmico actualizado · arte final sujeto a validación técnica de I+D",
            font_size=Pt(10), color=RGBColor(0xAA, 0xB4, 0xC4), align=PP_ALIGN.CENTER, italic=True)

add_rect(s5, Inches(0), Inches(7.3), SLIDE_W, Inches(0.2), fill_color=ORANGE)
slide_number_label(s5, 5)

# ── SLIDE 6 — Texto del frente palabra por palabra ──────────────────────────────
s6 = add_slide()
section_header_bg(s6)
accent_bar(s6)

add_textbox(s6, Inches(0.6), Inches(0.25), Inches(12), Inches(0.55),
            "Texto del frente — palabra por palabra",
            font_size=Pt(22), bold=True, color=DARK_TEXT, align=PP_ALIGN.LEFT)
add_divider(s6, Inches(0.85))

# Three labeled boxes
boxes = [
    ("LENGÜETA / CINTILLA",
     "NUEVO\nLA PROTECCIÓN MÁS COMPLETA",
     RGBColor(0x1A, 0x2E, 0x50), WHITE, Inches(0.6)),
    ("CLAIM 1 — Elemento dominante (mayor jerarquía tipográfica)",
     "El más rápido ante parpadeos (< 0,03 s)",
     RGBColor(0xFF, 0x8C, 0x00), NAVY, Inches(3.55)),
    ("CLAIM 2 — Segundo elemento (jerarquía claramente menor)",
     "Autoprotección térmica activa*\n\n* Ver información de autoprotección térmica al reverso del empaque.",
     RGBColor(0xF2, 0xF4, 0xF7), DARK_TEXT, Inches(5.9)),
]

for label, content, bg, fg, top_pos in boxes:
    box_h = Inches(1.35)
    add_rect(s6, Inches(0.6), top_pos, Inches(12.13), box_h, fill_color=bg,
             line_color=DIVIDER, line_width=Pt(0.5))
    add_textbox(s6, Inches(0.8), top_pos + Pt(8), Inches(11.7), Inches(0.3),
                label, font_size=Pt(9), bold=True,
                color=ORANGE if bg != RGBColor(0xFF, 0x8C, 0x00) else NAVY,
                align=PP_ALIGN.LEFT)
    add_textbox(s6, Inches(0.8), top_pos + Inches(0.35), Inches(11.7), Inches(0.9),
                content, font_size=Pt(16), bold=True, color=fg, align=PP_ALIGN.LEFT)

slide_number_label(s6, 6)

# ── SLIDE 7 — Texto del reverso parte 1 (características) ──────────────────────
s7 = add_slide()
section_header_bg(s7)
accent_bar(s7)

add_textbox(s7, Inches(0.6), Inches(0.25), Inches(12), Inches(0.55),
            "Texto del reverso — Características (parte 1 de 2)",
            font_size=Pt(22), bold=True, color=DARK_TEXT, align=PP_ALIGN.LEFT)
add_divider(s7, Inches(0.85))

add_textbox(s7, Inches(0.6), Inches(0.9), Inches(3.5), Inches(0.38),
            "SECCIÓN: CARACTERÍSTICAS",
            font_size=Pt(10), bold=True, color=ORANGE, align=PP_ALIGN.LEFT)

caract_bullets = [
    "•  El más rápido ante parpadeos: tiempo de respuesta < 30 ms (< 0,03 s), especialmente adecuado para equipos con tecnología inverter y cargas de arranque corto.",
    "•  Protege tecnología Inverter: la velocidad de respuesta de < 0,03 s minimiza la exposición de la electrónica de control inverter a condiciones de inestabilidad de red.",
    "•  Autoprotección térmica activa*: autoprotección del protector y del cableado de la instalación ante corrientes excesivas o conexiones deficientes.",
    "•  Curva de disparo de tiempo inverso: respuesta más rápida ante mayor sobrecarga — configuración técnica complementaria al breaker termomagnético del circuito.",
    "•  Protege contra sobre voltaje, bajo voltaje, parpadeos e inestabilidad de la red eléctrica.",
    "•  Supresor de picos incorporado.",
    "•  Temporizado inteligente de reconexión.",
]

paras_s7 = []
for b in caract_bullets:
    paras_s7.append({"text": b, "size": Pt(13), "color": MID_GREY, "space_before": Pt(7)})

add_multi_para_textbox(s7, Inches(0.6), Inches(1.35), Inches(12.13), Inches(5.6), paras_s7)

add_textbox(s7, Inches(0.6), Inches(6.85), Inches(12), Inches(0.38),
            "INFORMACIÓN DE SEGURIDAD:  No reemplaza los breakers termomagnéticos de la instalación eléctrica.",
            font_size=Pt(11), bold=True, color=DARK_TEXT, align=PP_ALIGN.LEFT)

slide_number_label(s7, 7)

# ── SLIDE 8 — Texto del reverso parte 2 (caveats) ──────────────────────────────
s8 = add_slide()
section_header_bg(s8)
accent_bar(s8)

add_textbox(s8, Inches(0.6), Inches(0.25), Inches(12), Inches(0.55),
            "Texto del reverso — Caveats y asteriscos (parte 2 de 2)",
            font_size=Pt(22), bold=True, color=DARK_TEXT, align=PP_ALIGN.LEFT)
add_divider(s8, Inches(0.85))

caveat_blocks = [
    ("* VELOCIDAD — texto obligatorio, sin modificar",
     "El tiempo de desconexión de menos de 30 milisegundos (< 0,03 s) aplica ante parpadeos (fluctuaciones rápidas del voltaje de la red eléctrica) e inestabilidad de la red. No aplica a la desconexión ante sobre voltaje o bajo voltaje pronunciados, cuyo tiempo de desconexión es de 0,4 a 3 segundos según la intensidad de la falla. Genteca es el protector enchufable monofásico con el menor tiempo de respuesta verificado ante parpadeos en el mercado venezolano. Según pruebas de laboratorio I+D Genteca y verificación comparativa de competidores."),
    ("* INVERTER — texto obligatorio, sin modificar",
     "La protección ante parpadeos (fluctuaciones rápidas del voltaje de la red) que ofrece este protector es especialmente beneficiosa para equipos con tecnología inverter, cuya electrónica de control es sensible a variaciones rápidas del voltaje. Este protector no reemplaza la protección contra transientes de alta energía (descargas atmosféricas, conmutación inductiva) presente en el equipo inverter de fábrica. Ambas protecciones son complementarias."),
    ("* AUTOPROTECCIÓN TÉRMICA ACTIVA — texto obligatorio, sin modificar",
     "Autoprotección térmica activa: sensor de temperatura ubicado junto al relé de potencia. Detecta calentamiento excesivo causado por sobrecorrientes o por conexiones deficientes (bornes flojos o falsos contactos) y desconecta la carga antes de que el cableado o las conexiones se dañen. Protege al protector mismo y a la instalación eléctrica. Para cargas de baja demanda de corriente, protege al protector y al cableado, pero no actúa como protección de sobrecarga directa de la carga conectada. Funciona como capa adicional de protección térmica; no reemplaza al interruptor termomagnético de la instalación."),
]

top = Inches(0.95)
for label, body in caveat_blocks:
    add_textbox(s8, Inches(0.6), top, Inches(12.13), Inches(0.3),
                label, font_size=Pt(9.5), bold=True, color=ORANGE, align=PP_ALIGN.LEFT)
    add_textbox(s8, Inches(0.6), top + Inches(0.3), Inches(12.13), Inches(1.05),
                body, font_size=Pt(11), color=MID_GREY, align=PP_ALIGN.LEFT)
    top += Inches(1.55)

slide_number_label(s8, 8)

# ── SLIDE 9 — Lo que NO cambia ──────────────────────────────────────────────────
s9 = add_slide()
dark_slide_bg(s9)
accent_bar(s9, top=Inches(0))
add_rect(s9, Inches(0), Inches(0), SLIDE_W, Inches(0.55), fill_color=ORANGE)
add_textbox(s9, Inches(0.45), Inches(0.1), Inches(5), Inches(0.38),
            "GENTECA  |  EXCELINE GSM", font_size=Pt(12), bold=True, color=NAVY)

add_textbox(s9, Inches(0.6), Inches(0.65), Inches(12), Inches(0.5),
            "Lo que NO cambia — Alternativa B permanece intacta en todo lo demás",
            font_size=Pt(20), bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_divider(s9, Inches(1.25), color=RGBColor(0x2A, 0x3F, 0x60))

unchanged = [
    ("Claim de velocidad en el frente",
     "\"El más rápido ante parpadeos (< 0,03 s)\" — intacto. Es el argumento diferenciador primario de la Alternativa B. Verificado en el mercado venezolano por I+D Genteca."),
    ("Caveat de velocidad en el reverso",
     "Texto literal obligatorio completo — sin modificar. Incluye la comparación verificada con competidores y la especificación de los escenarios de aplicación."),
    ("Argumento Inverter en el reverso",
     "\"Protege tecnología Inverter\" con argumento causal completo en bullets de Características. Caveat literal completo en bloque de asteriscos. Sin cambio."),
    ("Disclaimer breaker termomagnético",
     "\"No reemplaza los breakers termomagnéticos de la instalación eléctrica\" — intacto. Presente en INFORMACIÓN DE SEGURIDAD del reverso."),
    ("Lengüeta del frente",
     "\"NUEVO / LA PROTECCIÓN MÁS COMPLETA\" — intacta. Posicionamiento de línea sin cambio."),
    ("Curva de disparo de tiempo inverso",
     "Bullet en Características — intacto. Sin cambio."),
]

top = Inches(1.45)
col_w = Inches(5.7)
for i, (title, body) in enumerate(unchanged):
    col = i % 2
    row = i // 2
    l = Inches(0.6) + col * Inches(6.5)
    t = top + row * Inches(1.6)
    add_rect(s9, l, t, col_w, Inches(1.45), fill_color=RGBColor(0x15, 0x28, 0x45))
    add_textbox(s9, l + Inches(0.15), t + Pt(8), col_w - Inches(0.3), Inches(0.35),
                title, font_size=Pt(11.5), bold=True, color=ORANGE, align=PP_ALIGN.LEFT)
    add_textbox(s9, l + Inches(0.15), t + Inches(0.4), col_w - Inches(0.3), Inches(0.9),
                body, font_size=Pt(11), color=RGBColor(0xCC, 0xD8, 0xE8), align=PP_ALIGN.LEFT)

add_rect(s9, Inches(0), Inches(7.3), SLIDE_W, Inches(0.2), fill_color=ORANGE)
slide_number_label(s9, 9)

# ── SLIDE 10 — Decisión y próximos pasos ────────────────────────────────────────
s10 = add_slide()
section_header_bg(s10)
accent_bar(s10)

add_textbox(s10, Inches(0.6), Inches(0.25), Inches(12), Inches(0.55),
            "Decisión solicitada y próximos pasos",
            font_size=Pt(22), bold=True, color=DARK_TEXT, align=PP_ALIGN.LEFT)
add_divider(s10, Inches(0.85))

# Decision box
add_rect(s10, Inches(0.6), Inches(1.0), Inches(12.13), Inches(1.85),
         fill_color=NAVY, line_color=ORANGE, line_width=Pt(1.5))
add_textbox(s10, Inches(0.8), Inches(1.1), Inches(11.7), Inches(0.35),
            "DECISIÓN SOLICITADA A LA JUNTA HOY",
            font_size=Pt(10), bold=True, color=ORANGE, align=PP_ALIGN.LEFT)
add_textbox(s10, Inches(0.8), Inches(1.45), Inches(11.7), Inches(1.2),
            "Confirmar \"Autoprotección térmica activa*\" como segundo claim del frente de la Alternativa B, y el texto literal del caveat correspondiente en el reverso (presentado en el slide anterior), como versión de producción del empaque Exceline GSM.",
            font_size=Pt(14), bold=False, color=WHITE, align=PP_ALIGN.LEFT)

# Next steps
add_textbox(s10, Inches(0.6), Inches(3.05), Inches(12), Inches(0.38),
            "PRÓXIMOS PASOS — tras confirmación de la Junta",
            font_size=Pt(11), bold=True, color=DARK_TEXT, align=PP_ALIGN.LEFT)

steps = [
    ("1", "Validación técnica con I+D",
     "El equipo de I+D confirma el mecanismo del componente térmico y valida que el adjetivo \"activa\" es técnicamente sostenible. Si el resultado es negativo, el claim pasa a \"Autoprotección térmica*\" (formulación de respaldo ya definida, sin impacto en el caveat del reverso)."),
    ("2", "Arte final por parte del equipo de diseño",
     "El equipo de diseño actualiza el arte final del tiro con el badge aprobado y remite el archivo para revisión. El proceso incluye la verificación del datasheet I+D con el valor de tiempo de respuesta verificado (< 30 ms)."),
    ("3", "Entrega al proveedor de empaque",
     "Una vez cerrado el arte final y validada la especificación técnica, el arte aprobado se comparte con el proveedor de empaque para producción. Los pasos 1 y 2 deben estar cerrados antes de esta instrucción."),
]

top = Inches(3.55)
for num, title, body in steps:
    add_rect(s10, Inches(0.6), top, Inches(0.42), Inches(0.42), fill_color=ORANGE)
    add_textbox(s10, Inches(0.63), top - Pt(1), Inches(0.37), Inches(0.44),
                num, font_size=Pt(14), bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_textbox(s10, Inches(1.18), top - Pt(2), Inches(11.1), Inches(0.35),
                title, font_size=Pt(13), bold=True, color=DARK_TEXT, align=PP_ALIGN.LEFT)
    add_textbox(s10, Inches(1.18), top + Inches(0.3), Inches(11.1), Inches(0.75),
                body, font_size=Pt(12), color=MID_GREY, align=PP_ALIGN.LEFT)
    top += Inches(1.15)

slide_number_label(s10, 10)

# ── Document properties — no internal metadata ──────────────────────────────────
core = prs.core_properties
core.author   = "Genteca"
core.title    = "Empaque Linea Exceline GSM - Propuesta Refinada"
core.subject  = "Presentacion Junta Directiva"
core.keywords = ""
core.comments = ""
core.last_modified_by = "Genteca"

# ── Save ────────────────────────────────────────────────────────────────────────
prs.save(OUTPUT)
print(f"SAVED: {OUTPUT}")
print(f"SLIDES: {len(prs.slides)}")
