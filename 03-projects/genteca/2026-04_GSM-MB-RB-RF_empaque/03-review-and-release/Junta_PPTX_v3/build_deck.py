"""Build short junta deck for 2026-05-11 — Thermo-Safe + inverter package."""
import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Brand palette (extracted from VRB 2026-05-06 deck)
BG_NAVY = RGBColor(0x1A, 0x2E, 0x4A)
CARD_NAVY = RGBColor(0x0D, 0x1E, 0x33)
CARD_LIGHTER = RGBColor(0x24, 0x3D, 0x5C)
ACCENT_ORANGE = RGBColor(0xF5, 0xA6, 0x23)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SUBTLE_BLUE = RGBColor(0xD6, 0xE4, 0xF0)
GRAY_LIGHT = RGBColor(0xC8, 0xC8, 0xC8)
GRAY_MED = RGBColor(0x80, 0x80, 0x80)
GRAY_DARK = RGBColor(0x5A, 0x64, 0x70)
RED_CONF = RGBColor(0xFF, 0x60, 0x60)
DARK_INK = RGBColor(0x1A, 0x2E, 0x4A)
PALE_CARD = RGBColor(0xF5, 0xF5, 0xF0)
GREEN_OK = RGBColor(0x4C, 0xAF, 0x50)
RED_NOK = RGBColor(0xE5, 0x73, 0x73)
YELLOW_MED = RGBColor(0xF5, 0xD0, 0x6A)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
    shape.shadow.inherit = False
    return shape


def add_text(slide, x, y, w, h, text, *, font="Calibri", size=12, bold=False,
             color=WHITE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    if isinstance(text, str):
        text = [text]
    for i, line in enumerate(text):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_multi_text(slide, x, y, w, h, runs, *, align=PP_ALIGN.LEFT,
                   anchor=MSO_ANCHOR.TOP, line_spacing=None):
    """runs: list of dicts with text, font, size, bold, color, optional newline_after."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    first = True
    for r in runs:
        if r.get("paragraph") and not first:
            p = tf.add_paragraph()
            p.alignment = r.get("align", align)
            if line_spacing:
                p.line_spacing = line_spacing
        run = p.add_run()
        run.text = r["text"]
        run.font.name = r.get("font", "Calibri")
        run.font.size = Pt(r.get("size", 12))
        run.font.bold = r.get("bold", False)
        run.font.color.rgb = r.get("color", WHITE)
        first = False
    return tb


def add_brand_header(slide, page_label=None):
    # Orange stripe
    add_rect(slide, Emu(0), Inches(0.55), Inches(13.333), Pt(4), ACCENT_ORANGE)
    # GENTECA tag
    add_text(slide, Inches(0.45), Inches(0.18), Inches(2), Inches(0.35),
             "GENTECA", font="Calibri Light", size=14, bold=True, color=ACCENT_ORANGE)
    # CONFIDENCIAL stamp top right
    add_text(slide, Inches(11.0), Inches(0.18), Inches(2.3), Inches(0.35),
             "CONFIDENCIAL", font="Calibri", size=10, bold=True, color=RED_CONF,
             align=PP_ALIGN.RIGHT)
    if page_label:
        add_text(slide, Inches(0.45), Inches(7.05), Inches(12.5), Inches(0.3),
                 page_label, font="Calibri", size=8.5, color=GRAY_LIGHT,
                 align=PP_ALIGN.LEFT)


FOOTER = "Genteca  |  Junta Directiva  |  2026-05-11  |  Confidencial"


# ============================================================================
# Build presentation
# ============================================================================
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

# ---------------------------------------------------------------------------
# SLIDE 1 — COVER
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG_NAVY)
# Center accent block
add_rect(s, Inches(0), Inches(2.6), Inches(13.333), Inches(0.06), ACCENT_ORANGE)
# GENTECA top
add_text(s, Inches(0.6), Inches(0.6), Inches(4), Inches(0.6),
         "GENTECA", font="Calibri Light", size=20, bold=True, color=ACCENT_ORANGE)
# Title
add_text(s, Inches(0.6), Inches(2.95), Inches(12), Inches(1.1),
         "Avance empaque GSM-MB / RB / RF / RE",
         font="Calibri Light", size=38, bold=True, color=WHITE)
# Subtitle
add_text(s, Inches(0.6), Inches(4.05), Inches(12), Inches(0.6),
         "Recomendación final por audiencia — decisión solicitada hoy",
         font="Calibri Light", size=20, bold=False, color=SUBTLE_BLUE)
# Authorship sello
add_text(s, Inches(0.6), Inches(6.25), Inches(8), Inches(0.4),
         "Elaborada entre Alberto Betancourt y Raoul Bermúdez",
         font="Calibri", size=13, bold=True, color=ACCENT_ORANGE)
# Junta sello bottom right
add_multi_text(s, Inches(8.5), Inches(6.25), Inches(4.4), Inches(0.8), [
    {"text": "Junta Directiva Genteca", "size": 12, "color": GRAY_LIGHT, "bold": False},
    {"text": "2026-05-11", "size": 12, "color": GRAY_LIGHT, "bold": False, "paragraph": True},
    {"text": "CONFIDENCIAL", "size": 11, "color": RED_CONF, "bold": True, "paragraph": True},
], align=PP_ALIGN.RIGHT)

# ---------------------------------------------------------------------------
# SLIDE 2 — Lo firme y lo que se reabre
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG_NAVY)
add_brand_header(s, FOOTER)
add_text(s, Inches(0.45), Inches(0.85), Inches(12), Inches(0.7),
         "Punto de partida y qué cambia hoy",
         font="Calibri Light", size=26, bold=True, color=WHITE)
add_text(s, Inches(0.45), Inches(1.5), Inches(12), Inches(0.5),
         "El punto de partida construido hasta hoy + las decisiones que pedimos resolver en esta sesión.",
         font="Calibri", size=13, color=ACCENT_ORANGE)

# Two-column layout
col_w = Inches(6.05)
col_h = Inches(4.7)
col_y = Inches(2.2)

# Left col — FIRME
add_rect(s, Inches(0.45), col_y, col_w, col_h, CARD_NAVY)
add_rect(s, Inches(0.45), col_y, col_w, Inches(0.55), CARD_LIGHTER)
add_text(s, Inches(0.7), Inches(2.27), col_w, Inches(0.4),
         "PUNTO DE PARTIDA", font="Calibri", size=12, bold=True, color=ACCENT_ORANGE)
items_firme = [
    "Alternativa B aprobada",
    "Velocidad (<30 ms) como Claim 1 dominante",
    "QR dinámico en lengüeta",
    "Caveats al retiro",
]
for i, item in enumerate(items_firme):
    add_multi_text(s, Inches(0.7), Inches(2.95) + Inches(0.55) * i, Inches(5.6), Inches(0.5), [
        {"text": "•  ", "size": 14, "color": ACCENT_ORANGE, "bold": True},
        {"text": item, "size": 14, "color": WHITE},
    ])

# Right col — SE DECIDE HOY
add_rect(s, Inches(6.83), col_y, col_w, col_h, CARD_NAVY)
add_rect(s, Inches(6.83), col_y, col_w, Inches(0.55), CARD_LIGHTER)
add_text(s, Inches(7.08), Inches(2.27), col_w, Inches(0.4),
         "SE DECIDE HOY", font="Calibri", size=12, bold=True, color=ACCENT_ORANGE)
items_hoy = [
    "Nombre del badge térmico para el frente",
    "Cómo se trata NTC en el empaque (al lado del nombre fantasía)",
    "Inverter al frente — sub-decisión: revisar la posición previa",
    "Diferenciación del lenguaje técnico vs. masivo",
]
for i, item in enumerate(items_hoy):
    add_multi_text(s, Inches(7.08), Inches(2.95) + Inches(0.55) * i, Inches(5.6), Inches(0.6), [
        {"text": "•  ", "size": 14, "color": ACCENT_ORANGE, "bold": True},
        {"text": item, "size": 14, "color": WHITE},
    ])

# ---------------------------------------------------------------------------
# SLIDE 3 — Dos públicos, dos estrategias
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG_NAVY)
add_brand_header(s, FOOTER)
add_text(s, Inches(0.45), Inches(0.85), Inches(12), Inches(0.7),
         "Dos públicos, dos estrategias",
         font="Calibri Light", size=26, bold=True, color=WHITE)
add_text(s, Inches(0.45), Inches(1.5), Inches(12), Inches(0.5),
         "El frente del empaque debe hablar el idioma del comprador, no el del producto.",
         font="Calibri", size=13, color=ACCENT_ORANGE)

# Table headers
table_y = Inches(2.25)
col1_x = Inches(0.45)
col2_x = Inches(3.85)
col3_x = Inches(8.65)
c1_w = Inches(3.30)
c23_w = Inches(4.7)
header_h = Inches(0.7)
row_h = Inches(0.95)

# Header row
add_rect(s, col2_x, table_y, c23_w, header_h, CARD_LIGHTER)
add_rect(s, col3_x, table_y, c23_w, header_h, CARD_LIGHTER)
add_text(s, col2_x, table_y, c23_w, header_h,
         "GSM-MB / RB / RF — bornera", font="Calibri", size=15, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, col3_x, table_y, c23_w, header_h,
         "GSM-RE — enchufable", font="Calibri", size=15, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

rows = [
    ("Quién compra", "Técnico instalador", "Consumidor masivo"),
    ("Cómo decide", "Credenciales técnicas — busca el dato", "Busca el beneficio inmediato"),
    ("Lenguaje del frente", "Funcional + nombre registrable", "Beneficio del beneficio"),
    ("Badge propuesto", "Tecnología Thermo-Safe™ NTC", "Protege tu instalación ante corrientes extremas"),
]
for i, (label, c2, c3) in enumerate(rows):
    y = table_y + header_h + row_h * i
    bg = CARD_NAVY if i % 2 == 0 else BG_NAVY
    add_rect(s, col1_x, y, c1_w, row_h, CARD_LIGHTER if i % 2 == 0 else CARD_NAVY)
    add_rect(s, col2_x, y, c23_w, row_h, bg)
    add_rect(s, col3_x, y, c23_w, row_h, bg)
    add_text(s, col1_x + Inches(0.15), y, c1_w - Inches(0.3), row_h, label,
             font="Calibri", size=12, bold=True, color=ACCENT_ORANGE,
             anchor=MSO_ANCHOR.MIDDLE)
    badge_size = 16 if "Badge" in label else 13
    badge_bold = True if "Badge" in label else False
    add_text(s, col2_x + Inches(0.2), y, c23_w - Inches(0.4), row_h, c2,
             font="Calibri", size=badge_size, bold=badge_bold, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, col3_x + Inches(0.2), y, c23_w - Inches(0.4), row_h, c3,
             font="Calibri", size=badge_size, bold=badge_bold, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)

# ---------------------------------------------------------------------------
# Helper: build a "frente del empaque" mock for slides 4 and 5
# ---------------------------------------------------------------------------
def add_packaging_mock(slide, x, y, w, h, *, lengueta_text, claim1_main, claim1_sub,
                       claim2_label, claim2_text):
    # Card background
    add_rect(slide, x, y, w, h, PALE_CARD)
    # Top "tab" (lengüeta)
    tab_h = Inches(0.85)
    add_rect(slide, x, y, w, tab_h, ACCENT_ORANGE)
    add_text(slide, x, y, w, tab_h,
             lengueta_text, font="Calibri", size=22, bold=True,
             color=DARK_INK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Section: CLAIM 1
    sec1_y = y + tab_h + Inches(0.25)
    add_text(slide, x + Inches(0.3), sec1_y, w - Inches(0.6), Inches(0.3),
             "CLAIM 1 — ELEMENTO DOMINANTE", font="Calibri", size=10, bold=True,
             color=GRAY_DARK)
    add_text(slide, x + Inches(0.3), sec1_y + Inches(0.32), w - Inches(0.6), Inches(0.6),
             claim1_main, font="Calibri", size=22, bold=True,
             color=DARK_INK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_multi_text(slide, x + Inches(0.3), sec1_y + Inches(0.95),
                   w - Inches(0.6), Inches(0.6), [
        {"text": "El mejor para:  ", "size": 13, "color": GRAY_DARK, "bold": False},
        {"text": "INVERTER", "size": 26, "color": DARK_INK, "bold": True},
    ], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Divider
    sec2_y = sec1_y + Inches(1.7)
    add_rect(slide, x + Inches(0.5), sec2_y, w - Inches(1.0), Pt(1), GRAY_DARK)

    # Section: CLAIM 2 (badge térmico)
    sec2_y = sec2_y + Inches(0.18)
    add_text(slide, x + Inches(0.3), sec2_y, w - Inches(0.6), Inches(0.3),
             claim2_label, font="Calibri", size=10, bold=True, color=GRAY_DARK)
    add_text(slide, x + Inches(0.3), sec2_y + Inches(0.35), w - Inches(0.6), Inches(0.7),
             claim2_text, font="Calibri", size=20, bold=True,
             color=DARK_INK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Asterisk note bottom
    add_text(slide, x + Inches(0.3), y + h - Inches(0.4), w - Inches(0.6), Inches(0.3),
             "*  Asterisco — caveat literal en el retiro",
             font="Calibri", size=9, color=GRAY_MED)

# ---------------------------------------------------------------------------
# SLIDE 4 — Recomendación tiro: MB / RB / RF
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG_NAVY)
add_brand_header(s, FOOTER)
add_text(s, Inches(0.45), Inches(0.85), Inches(12), Inches(0.7),
         "Recomendación tiro — MB / RB / RF (bornera, técnico)",
         font="Calibri Light", size=24, bold=True, color=WHITE)
add_text(s, Inches(0.45), Inches(1.5), Inches(12), Inches(0.5),
         "Texto completo del frente del empaque.",
         font="Calibri", size=13, color=ACCENT_ORANGE)

# Packaging mock (left)
add_packaging_mock(s, Inches(0.6), Inches(2.2), Inches(6.0), Inches(4.5),
                   lengueta_text="NUEVO  ◻QR  PROTECCIÓN MÁS COMPLETA",
                   claim1_main="El más rápido ante parpadeos (< 0,03 s)",
                   claim1_sub="",
                   claim2_label="CLAIM 2 — BADGE TÉRMICO (jerarquía menor)",
                   claim2_text="Tecnología Thermo-Safe™ NTC  *")

# Razón al lado derecho
right_x = Inches(7.1)
add_rect(s, right_x, Inches(2.2), Inches(5.8), Inches(4.5), CARD_NAVY)
add_rect(s, right_x, Inches(2.2), Inches(5.8), Inches(0.55), CARD_LIGHTER)
add_text(s, right_x + Inches(0.25), Inches(2.27), Inches(5.5), Inches(0.4),
         "POR QUÉ ESTA VARIANTE",
         font="Calibri", size=12, bold=True, color=ACCENT_ORANGE)

reasons_mb = [
    ("Credencial técnica visible.",
     "El técnico reconoce NTC como dato relevante; el badge no le hace perder esa pista."),
    ("Anti-copia activada.",
     "Thermo-Safe™ es nombre fantasía registrable. Una vez registrado, ningún competidor puede usarlo."),
    ("Inverter al frente.",
     "Cierra la deuda con la pregunta más recurrente del canal técnico, anclada en velocidad."),
]
ry = Inches(2.95)
for title, body in reasons_mb:
    add_multi_text(s, right_x + Inches(0.3), ry, Inches(5.3), Inches(1.2), [
        {"text": title + " ", "size": 14, "bold": True, "color": WHITE},
        {"text": body, "size": 13, "color": SUBTLE_BLUE},
    ])
    ry += Inches(1.2)

# ---------------------------------------------------------------------------
# SLIDE 5 — Recomendación tiro: RE
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG_NAVY)
add_brand_header(s, FOOTER)
add_text(s, Inches(0.45), Inches(0.85), Inches(12), Inches(0.7),
         "Recomendación tiro — RE (enchufable, masivo)",
         font="Calibri Light", size=24, bold=True, color=WHITE)
add_text(s, Inches(0.45), Inches(1.5), Inches(12), Inches(0.5),
         "Texto completo del frente del empaque para consumo masivo.",
         font="Calibri", size=13, color=ACCENT_ORANGE)

add_packaging_mock(s, Inches(0.6), Inches(2.2), Inches(6.0), Inches(4.5),
                   lengueta_text="NUEVO  ◻QR  PROTECCIÓN MÁS COMPLETA",
                   claim1_main="El más rápido ante parpadeos (< 0,03 s)",
                   claim1_sub="",
                   claim2_label="CLAIM 2 — BENEFICIO PARA EL CONSUMIDOR",
                   claim2_text="Protege tu instalación ante corrientes extremas  *")

right_x = Inches(7.1)
add_rect(s, right_x, Inches(2.2), Inches(5.8), Inches(4.5), CARD_NAVY)
add_rect(s, right_x, Inches(2.2), Inches(5.8), Inches(0.55), CARD_LIGHTER)
add_text(s, right_x + Inches(0.25), Inches(2.27), Inches(5.5), Inches(0.4),
         "POR QUÉ ESTA VARIANTE",
         font="Calibri", size=12, bold=True, color=ACCENT_ORANGE)

reasons_re = [
    ("Beneficio inmediato y concreto.",
     "El comprador masivo no escanea QR ni lee retiro al decidir; el frente debe entregar el resultado, no la tecnología."),
    ("Inverter al frente — decisivo.",
     "22% del mercado 12k BTU es inverter hoy; proyección >50% en 3 años. Si no lo decimos, perdemos en el anaquel."),
    ("Coherencia con la propuesta de Jesús María.",
     "El \"beneficio del beneficio\" para el masivo. La función técnica queda en el retiro y en piezas digitales."),
]
ry = Inches(2.95)
for title, body in reasons_re:
    add_multi_text(s, right_x + Inches(0.3), ry, Inches(5.3), Inches(1.2), [
        {"text": title + " ", "size": 14, "bold": True, "color": WHITE},
        {"text": body, "size": 13, "color": SUBTLE_BLUE},
    ])
    ry += Inches(1.2)

# ---------------------------------------------------------------------------
# SLIDE 6 — Por qué Thermo-Safe™ y no Escudo Térmico solo
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG_NAVY)
add_brand_header(s, FOOTER)
add_text(s, Inches(0.45), Inches(0.85), Inches(12), Inches(0.7),
         "El nombre que se registra protege; el descriptor solo informa",
         font="Calibri Light", size=24, bold=True, color=WHITE)
add_text(s, Inches(0.45), Inches(1.5), Inches(12), Inches(0.5),
         "Tres opciones evaluadas — la combinación gana.",
         font="Calibri", size=13, color=ACCENT_ORANGE)

# 3-row table
ty = Inches(2.2)
ch = Inches(0.55)
rh = Inches(1.3)
cx0 = Inches(0.45)
cx1 = Inches(3.15)
cx2 = Inches(6.55)
cx3 = Inches(9.95)
cw0 = Inches(2.65)
cw1 = Inches(3.35)
cw2 = Inches(3.35)
cw3 = Inches(2.95)

add_rect(s, cx0, ty, cw0, ch, CARD_LIGHTER)
add_rect(s, cx1, ty, cw1, ch, CARD_LIGHTER)
add_rect(s, cx2, ty, cw2, ch, CARD_LIGHTER)
add_rect(s, cx3, ty, cw3, ch, CARD_LIGHTER)
for cx, cw, lbl in [(cx0, cw0, "Opción"), (cx1, cw1, "Qué dice"),
                    (cx2, cw2, "Qué pasa con SAPI"), (cx3, cw3, "Qué ganamos")]:
    add_text(s, cx + Inches(0.15), ty, cw, ch, lbl,
             font="Calibri", size=11, bold=True, color=ACCENT_ORANGE,
             anchor=MSO_ANCHOR.MIDDLE)

opts = [
    ('Solo "Escudo Térmico"', "Función técnica",
     "Riesgo medio — puede leerse descriptivo", "Cero protección anti-copia", RED_NOK),
    ('Solo "Tecnología Thermo-Safe™"', "Nombre fantasía puro",
     "Mejor perfil de registro", "Anti-copia, pero el técnico pierde la pista de NTC", YELLOW_MED),
    ('"Tecnología Thermo-Safe™ NTC"  ✓', "Nombre fantasía + credencial",
     "Thermo-Safe™ se registra; NTC va al lado sin pretender registro",
     "Anti-copia + credencial técnica visible", GREEN_OK),
]
for i, (a, b, c, d, marker) in enumerate(opts):
    y = ty + ch + rh * i
    bg = CARD_NAVY if i != 2 else CARD_LIGHTER
    add_rect(s, cx0, y, cw0, rh, bg)
    add_rect(s, cx1, y, cw1, rh, bg)
    add_rect(s, cx2, y, cw2, rh, bg)
    add_rect(s, cx3, y, cw3, rh, bg)
    # Marker stripe at left
    add_rect(s, cx0, y, Inches(0.08), rh, marker)
    bold_a = i == 2
    add_text(s, cx0 + Inches(0.2), y, cw0 - Inches(0.3), rh, a,
             font="Calibri", size=12, bold=bold_a, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, cx1 + Inches(0.15), y, cw1 - Inches(0.3), rh, b,
             font="Calibri", size=12, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, cx2 + Inches(0.15), y, cw2 - Inches(0.3), rh, c,
             font="Calibri", size=12, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, cx3 + Inches(0.15), y, cw3 - Inches(0.3), rh, d,
             font="Calibri", size=12, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

# Bottom block — submisión paralela
sb_y = ty + ch + rh * 3 + Inches(0.2)
add_rect(s, Inches(0.45), sb_y, Inches(12.45), Inches(1.15), CARD_NAVY)
add_multi_text(s, Inches(0.65), sb_y + Inches(0.12), Inches(12.0), Inches(0.45), [
    {"text": "Estrategia de submisión:  ", "size": 13, "bold": True, "color": ACCENT_ORANGE},
    {"text": "se introducen varios candidatos en paralelo (Escudo Térmico / Thermo-Shield / Thermo-Safe). El que pase primero ancla la marca; los demás quedan como respaldo.",
     "size": 13, "color": WHITE},
])
add_multi_text(s, Inches(0.65), sb_y + Inches(0.6), Inches(12.0), Inches(0.45), [
    {"text": "Nota sobre Thermo-Shield:  ", "size": 12, "bold": True, "color": ACCENT_ORANGE},
    {"text": "marca internacional activa en categorías afines (revestimientos térmicos, clase distinta). Coexistencia legalmente posible por especialidad, pero refuerza la prioridad de ",
     "size": 12, "color": WHITE},
    {"text": "Thermo-Safe™ ", "size": 12, "bold": True, "color": WHITE},
    {"text": "como candidato principal.", "size": 12, "color": WHITE},
])

# ---------------------------------------------------------------------------
# SLIDE 7 — Inverter al frente (sub-decisión)
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG_NAVY)
add_brand_header(s, FOOTER)
add_text(s, Inches(0.45), Inches(0.85), Inches(12), Inches(0.7),
         "Inverter al frente — sub-decisión a reabrir",
         font="Calibri Light", size=24, bold=True, color=WHITE)
add_text(s, Inches(0.45), Inches(1.5), Inches(12.5), Inches(0.5),
         "Posición previa: inverter al retiro.   Posición que recomendamos hoy: inverter al frente.",
         font="Calibri", size=13, color=ACCENT_ORANGE)

# Three vertical blocks
bx = [Inches(0.45), Inches(4.83), Inches(9.21)]
bw = Inches(4.10)
by = Inches(2.25)
bh = Inches(4.6)
titles = [
    "El mercado se mueve hacia inverter",
    "Lo que dice el mercado",
    "Lo que ya tenemos para decirlo",
]
bodies = [
    "Lo que era novedad hace algunos años hoy es estándar de facto, y va camino a ser expectativa por defecto en el corto plazo.\n\nEl empaque debe acompañar esa transición. Colocar inverter al frente no es repetir lo viejo — es alinear el producto con la dirección del mercado.",
    "22% del mercado 12k BTU es inverter hoy.\n\n>50% en proyección a 3 años.\n\nSi el competidor lo dice y nosotros no, perdemos la decisión en el anaquel — no por producto, por información ausente.",
    "Velocidad <30 ms.\n\nSomos los más rápidos. La velocidad es la razón técnica para proteger inverter mejor que nadie.\n\nInverter al frente cierra el círculo: damos la prueba (velocidad) y el dato que el comprador busca (inverter) en el mismo plano.",
]
for i in range(3):
    add_rect(s, bx[i], by, bw, bh, CARD_NAVY)
    add_rect(s, bx[i], by, bw, Inches(0.6), CARD_LIGHTER)
    # Badge number
    add_rect(s, bx[i] + Inches(0.2), by + Inches(0.13), Inches(0.35), Inches(0.35),
             ACCENT_ORANGE)
    add_text(s, bx[i] + Inches(0.2), by + Inches(0.13), Inches(0.35), Inches(0.35),
             str(i + 1), font="Calibri", size=14, bold=True, color=DARK_INK,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, bx[i] + Inches(0.7), by + Inches(0.1), bw - Inches(0.9), Inches(0.5),
             titles[i], font="Calibri", size=13, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, bx[i] + Inches(0.3), by + Inches(0.85), bw - Inches(0.6), bh - Inches(1.0),
             bodies[i], font="Calibri", size=13, color=SUBTLE_BLUE)

# Recommendation strip below
rs_y = by + bh + Inches(0.05)
add_rect(s, Inches(0.45), rs_y, Inches(12.45), Inches(0.0), ACCENT_ORANGE)

# ---------------------------------------------------------------------------
# SLIDE 8 — Mecánica del registro: ™ → ®
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG_NAVY)
add_brand_header(s, FOOTER)
add_text(s, Inches(0.45), Inches(0.85), Inches(12), Inches(0.7),
         "Salimos al mercado con el ™. Volvemos con ® cuando llegue.",
         font="Calibri Light", size=24, bold=True, color=WHITE)
add_text(s, Inches(0.45), Inches(1.5), Inches(12), Inches(0.5),
         "Mecánica del registro aplicada al caso Thermo-Safe.",
         font="Calibri", size=13, color=ACCENT_ORANGE)

# Timeline three boxes
tx = [Inches(0.5), Inches(4.78), Inches(9.06)]
tw = Inches(4.0)
ty_t = Inches(2.4)
th = Inches(2.3)
labels = ["HOY", "D + 30 días", "D + 6 a 12 meses"]
events = [
    "Empaque al mercado con\n«Tecnología Thermo-Safe™ NTC»",
    "Solicitud SAPI introducida\n(presupuesto IP autorizado)",
    "Certificado emitido\n«Thermo-Safe®»  +  «Exclusivo de Genteca»",
]
for i in range(3):
    add_rect(s, tx[i], ty_t, tw, th, CARD_NAVY)
    add_rect(s, tx[i], ty_t, tw, Inches(0.55), ACCENT_ORANGE)
    add_text(s, tx[i], ty_t, tw, Inches(0.55), labels[i],
             font="Calibri", size=14, bold=True, color=DARK_INK,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, tx[i] + Inches(0.3), ty_t + Inches(0.7), tw - Inches(0.6),
             th - Inches(0.85), events[i],
             font="Calibri", size=14, color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    if i < 2:
        # Arrow between
        ax = tx[i] + tw + Inches(0.05)
        add_rect(s, ax, ty_t + th / 2 - Inches(0.03), Inches(0.18), Inches(0.06),
                 ACCENT_ORANGE)

# Why-it-matters block
wy = ty_t + th + Inches(0.25)
wh_box = Inches(2.0)
add_rect(s, Inches(0.45), wy, Inches(12.45), wh_box, CARD_NAVY)
add_rect(s, Inches(0.45), wy, Inches(12.45), Inches(0.55), CARD_LIGHTER)
add_text(s, Inches(0.65), wy + Inches(0.07), Inches(12), Inches(0.4),
         "POR QUÉ HACERLO ASÍ", font="Calibri", size=12, bold=True, color=ACCENT_ORANGE)

points = [
    "Adoptar el ™ inmediato es legal y comercialmente correcto.",
    "30 días es el plazo desde el primer uso público para introducir el papeleo SAPI.",
    "Cuando llegue el ®, se libera la palabra «exclusivo» — segunda ola de comunicación gratis.",
    "Submisión paralela de varios candidatos: el abogado marcario decide la combinación óptima.",
]
for i, pt in enumerate(points):
    add_multi_text(s, Inches(0.7), wy + Inches(0.65) + Inches(0.32) * i,
                   Inches(12), Inches(0.32), [
        {"text": "•  ", "size": 13, "bold": True, "color": ACCENT_ORANGE},
        {"text": pt, "size": 13, "color": WHITE},
    ])

# ---------------------------------------------------------------------------
# SLIDE 9 — Argumento de fondo: protección contra copias
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG_NAVY)
add_brand_header(s, FOOTER)
add_text(s, Inches(0.45), Inches(0.85), Inches(12.5), Inches(0.7),
         "El problema de 20 años con una solución estructural",
         font="Calibri Light", size=26, bold=True, color=WHITE)
add_text(s, Inches(0.45), Inches(1.5), Inches(12.5), Inches(0.5),
         "Las copias chinas siempre van a aparecer. Lo que cambia es qué pueden copiar.",
         font="Calibri", size=13, color=ACCENT_ORANGE)

# Big quote-style narrative block
nb_y = Inches(2.3)
nb_h = Inches(3.6)
add_rect(s, Inches(0.8), nb_y, Inches(11.7), nb_h, CARD_NAVY)
add_rect(s, Inches(0.8), nb_y, Inches(0.08), nb_h, ACCENT_ORANGE)

add_multi_text(s, Inches(1.2), nb_y + Inches(0.35), Inches(11.0), nb_h - Inches(0.7), [
    {"text": "Durante dos décadas, las copias chinas han imitado diseño de caja, software interno, posicionamiento. Cambiamos la caja, cambian la suya. Si el nombre de la tecnología no está registrado, el competidor puede registrarlo primero — y entonces es Genteca quien recibe la notificación de infracción por una innovación que es nuestra.",
     "size": 16, "color": WHITE},
    {"text": "", "size": 8, "paragraph": True, "color": WHITE},
    {"text": "Registrar el nombre cambia la asimetría. Pueden copiar la tecnología; no pueden usar la palabra. Y la palabra es lo que el técnico pide en el anaquel:",
     "size": 16, "color": WHITE, "paragraph": True},
    {"text": "", "size": 6, "paragraph": True, "color": WHITE},
    {"text": "«el que tiene Thermo-Safe».",
     "size": 18, "color": ACCENT_ORANGE, "bold": True, "paragraph": True},
])

# Closing strip
cs_y = nb_y + nb_h + Inches(0.25)
add_rect(s, Inches(0.45), cs_y, Inches(12.45), Inches(0.75), CARD_LIGHTER)
add_multi_text(s, Inches(0.7), cs_y, Inches(12), Inches(0.75), [
    {"text": "El costo del registro está acotado:  ", "size": 14, "bold": True, "color": ACCENT_ORANGE},
    {"text": "$880–$1,640 por marca con abogado marcario.   ", "size": 14, "color": WHITE},
    {"text": "El costo de no registrar no tiene techo.", "size": 14, "bold": True, "color": WHITE},
], anchor=MSO_ANCHOR.MIDDLE)

# ---------------------------------------------------------------------------
# SLIDE 10 — Decisión solicitada
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG_NAVY)
add_brand_header(s, FOOTER)
add_text(s, Inches(0.45), Inches(0.85), Inches(12), Inches(0.7),
         "Lo que pedimos aprobar hoy",
         font="Calibri Light", size=26, bold=True, color=WHITE)
add_text(s, Inches(0.45), Inches(1.5), Inches(12.5), Inches(0.5),
         "Una decisión principal con tres sub-decisiones empacadas + una autorización presupuestaria.",
         font="Calibri", size=13, color=ACCENT_ORANGE)

decisions = [
    ("1", "Variante recomendada por audiencia",
     "MB / RB / RF (bornera): «Tecnología Thermo-Safe™ NTC» en el frente.\n"
     "RE (enchufable): «Protege tu instalación ante corrientes extremas» en el frente."),
    ("2", "Inverter al frente — sub-decisión",
     "Reabrir y aprobar la inclusión de «El mejor para: INVERTER» en el frente de los cuatro empaques."),
    ("3", "Diferenciación técnico / masivo",
     "Aprobar tratamientos distintos para empaque de bornera (técnico) vs. empaque enchufable (masivo)."),
    ("4", "Presupuesto IP — registro SAPI",
     "Autorizar inicio de proceso con abogado marcario. Submisión paralela de candidatos (Thermo-Safe / Thermo-Shield / Escudo Térmico). Costo estimado: $880–$1,640 por marca."),
]
dy = Inches(2.25)
dh = Inches(1.15)
for i, (num, title, body) in enumerate(decisions):
    y = dy + dh * i
    add_rect(s, Inches(0.45), y, Inches(12.45), dh - Inches(0.1), CARD_NAVY)
    # Number circle
    add_rect(s, Inches(0.65), y + Inches(0.22), Inches(0.6), Inches(0.6), ACCENT_ORANGE)
    add_text(s, Inches(0.65), y + Inches(0.22), Inches(0.6), Inches(0.6),
             num, font="Calibri", size=22, bold=True, color=DARK_INK,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(1.4), y + Inches(0.13), Inches(11), Inches(0.4),
             title, font="Calibri", size=15, bold=True, color=WHITE)
    add_text(s, Inches(1.4), y + Inches(0.5), Inches(11), Inches(0.6),
             body, font="Calibri", size=12, color=SUBTLE_BLUE)

# ---------------------------------------------------------------------------
# SLIDE 11 — Próximos pasos
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG_NAVY)
add_brand_header(s, FOOTER)
add_text(s, Inches(0.45), Inches(0.85), Inches(12), Inches(0.7),
         "De la decisión a la imprenta",
         font="Calibri Light", size=26, bold=True, color=WHITE)
add_text(s, Inches(0.45), Inches(1.5), Inches(12.5), Inches(0.5),
         "Próximos pasos si la decisión es positiva.",
         font="Calibri", size=13, color=ACCENT_ORANGE)

# Three columns
phases = [
    ("Esta semana", [
        "Mockup final del frente con la variante aprobada.",
        "Validación con vendedores para captar fricción residual.",
    ]),
    ("Próximas 2 semanas", [
        "Submisión SAPI iniciada con abogado marcario.",
        "Primer lote de candidatos sometido en paralelo.",
        "Ajustes de arte según las primeras lecturas.",
    ]),
    ("Mes 1–2 post-aprobación", [
        "Producción de planchas e impresión.",
        "Salida al mercado con ™.",
        "Estrategia de comunicación al canal preparada para la transición a ®.",
    ]),
]
px = [Inches(0.45), Inches(4.83), Inches(9.21)]
pw = Inches(4.10)
py = Inches(2.3)
ph = Inches(4.5)
for i, (label, items) in enumerate(phases):
    add_rect(s, px[i], py, pw, ph, CARD_NAVY)
    add_rect(s, px[i], py, pw, Inches(0.6), CARD_LIGHTER)
    add_text(s, px[i], py, pw, Inches(0.6), label,
             font="Calibri", size=14, bold=True, color=ACCENT_ORANGE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    iy = py + Inches(0.85)
    for it in items:
        add_multi_text(s, px[i] + Inches(0.3), iy, pw - Inches(0.5), Inches(1.3), [
            {"text": "•  ", "size": 14, "bold": True, "color": ACCENT_ORANGE},
            {"text": it, "size": 13, "color": WHITE},
        ])
        iy += Inches(1.05)

# ---------------------------------------------------------------------------
# SLIDE 12 — Anexo opcional: estrategia ampliada IP
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG_NAVY)
add_brand_header(s, FOOTER)
add_text(s, Inches(0.45), Inches(0.85), Inches(12.5), Inches(0.7),
         "Anexo — una decisión que abre una conversación más amplia",
         font="Calibri Light", size=24, bold=True, color=WHITE)
add_text(s, Inches(0.45), Inches(1.5), Inches(12.5), Inches(0.5),
         "Solo si surge la pregunta — el método aplicado a Thermo-Safe™ se replica al resto del portafolio.",
         font="Calibri", size=13, color=ACCENT_ORANGE)

# Table
ty = Inches(2.2)
ch = Inches(0.5)
rh = Inches(0.55)
cx0 = Inches(0.45)
cx1 = Inches(1.15)
cx2 = Inches(3.85)
cx3 = Inches(8.25)
cw0 = Inches(0.7)
cw1 = Inches(2.7)
cw2 = Inches(4.4)
cw3 = Inches(4.65)

add_rect(s, cx0, ty, cw0, ch, CARD_LIGHTER)
add_rect(s, cx1, ty, cw1, ch, CARD_LIGHTER)
add_rect(s, cx2, ty, cw2, ch, CARD_LIGHTER)
add_rect(s, cx3, ty, cw3, ch, CARD_LIGHTER)
for cx, cw, lbl in [(cx0, cw0, "#"), (cx1, cw1, "Marca candidata"),
                    (cx2, cw2, "Tecnología que nombra"),
                    (cx3, cw3, "Línea Genteca")]:
    add_text(s, cx + Inches(0.15), ty, cw, ch, lbl,
             font="Calibri", size=11, bold=True, color=ACCENT_ORANGE,
             anchor=MSO_ANCHOR.MIDDLE)

candidates = [
    ("1", "ForensLog™", "Histórico de fallas (100 / 20)", "GIII+, GME"),
    ("2", "FlickerGuard™", "Detección de parpadeos <30 ms", "GSM-L y nuevos GSM"),
    ("3", "Thermo-Safe™", "Función térmica integrada", "GSM-MB / RB / RF / RE  (decisión hoy)"),
    ("4", "TripleLock™", "Bloqueo a la tercera falla", "GIII+, GUCT, GOC"),
    ("5", "TaskMemory™", "Memoria de estado", "GRN-MV"),
    ("6", "ThermalCurve™", "Modelo curva fría / caliente", "GIII+, GUCT, GOC"),
    ("7", "StaggerStart™", "Arranque escalonado aleatorio", "Múltiples líneas"),
]
for i, (n, marca, tec, linea) in enumerate(candidates):
    y = ty + ch + rh * i
    bg = CARD_NAVY if i % 2 == 0 else BG_NAVY
    add_rect(s, cx0, y, cw0, rh, bg)
    add_rect(s, cx1, y, cw1, rh, bg)
    add_rect(s, cx2, y, cw2, rh, bg)
    add_rect(s, cx3, y, cw3, rh, bg)
    is_today = (i == 2)
    name_color = ACCENT_ORANGE if is_today else WHITE
    add_text(s, cx0 + Inches(0.15), y, cw0, rh, n,
             font="Calibri", size=12, bold=True, color=ACCENT_ORANGE,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, cx1 + Inches(0.15), y, cw1 - Inches(0.3), rh, marca,
             font="Calibri", size=13, bold=True, color=name_color,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, cx2 + Inches(0.15), y, cw2 - Inches(0.3), rh, tec,
             font="Calibri", size=12, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, cx3 + Inches(0.15), y, cw3 - Inches(0.3), rh, linea,
             font="Calibri", size=12, color=SUBTLE_BLUE, anchor=MSO_ANCHOR.MIDDLE)

# Footnote
fy = ty + ch + rh * 7 + Inches(0.25)
add_rect(s, Inches(0.45), fy, Inches(12.45), Inches(0.75), CARD_NAVY)
add_text(s, Inches(0.65), fy, Inches(12), Inches(0.75),
         "Decisión separada — no se solicita aprobación hoy. El método validado con Thermo-Safe™ se aplica al resto cuando la junta lo decida.",
         font="Calibri", size=13, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

# ---------------------------------------------------------------------------
# Save
# ---------------------------------------------------------------------------
out_path = (
    "03-projects/genteca/2026-04_GSM-MB-RB-RF_empaque/03-review-and-release/"
    "Junta_PPTX_v3/Junta_GSM_empaque_v3.2_2026-05-11.pptx"
)
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Total slides: {len(prs.slides)}")
