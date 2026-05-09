"""
build_pptx.py
Generador de presentacion ejecutiva — Junta Directiva Genteca
Proyecto: 2026-04_GSM-MB-RB-RF_empaque
Fecha: 2026-05-06
Fuente principal: AU-1 v2.2 (dual-variantes) + AU-1 v2.1 (caveats invariables)
v3 cambios: paleta V1 verde -> gris neutro; caveats divididos en 1 slide compartida
            + 2 slides termico especificas; anexo IP defensiva (A1-A5).
v4 cambios: incorpora Variante 3 (V3) propuesta Junta Directiva (Jesus Maria).
            3 sub-opciones V3a/V3b/V3c. Comparacion 3 columnas. Slides nuevas
            tiro V3, retiro caract V3, caveat termico V3, argumentos/riesgos/
            posicion/pasos/cierre actualizados. Paleta V3: siena/crema arena.
"""

import sys
sys.stdout.reconfigure(encoding="utf-8")

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
import os

# ---------------------------------------------------------------------------
# PALETA CORPORATIVA — Genteca / Exceline
# ---------------------------------------------------------------------------
C_AZUL_OSCURO  = RGBColor(0x1A, 0x2E, 0x4A)
C_AZUL_MEDIO   = RGBColor(0x1F, 0x4E, 0x79)
C_AZUL_CLARO   = RGBColor(0xD6, 0xE4, 0xF0)
C_AMBAR        = RGBColor(0xF5, 0xA6, 0x23)
C_GRIS_OSCURO  = RGBColor(0x40, 0x40, 0x40)
C_GRIS_MEDIO   = RGBColor(0x80, 0x80, 0x80)
C_GRIS_LINEA   = RGBColor(0xC8, 0xC8, 0xC8)
C_BLANCO       = RGBColor(0xFF, 0xFF, 0xFF)

# V1: gris carbon neutro
C_GRIS_V1      = RGBColor(0x5A, 0x64, 0x70)
C_GRIS_V1_CLR  = RGBColor(0xE5, 0xE7, 0xEA)

# V2: azul medio
C_AZUL_V2      = RGBColor(0x1F, 0x4E, 0x79)
C_AZUL_V2_CLR  = RGBColor(0xD6, 0xE4, 0xF0)

# V3: siena/crema arena — neutro, sin connotacion de "preferido"
# Diferente de gris (V1), azul (V2) y dorado opaco del anexo
C_SIENA_V3     = RGBColor(0x8B, 0x5E, 0x3C)   # #8B5E3C — siena oscuro
C_SIENA_V3_CLR = RGBColor(0xF5, 0xED, 0xDE)   # #F5EDDE — crema arena claro

# Paleta ANEXO — neutra, diferenciada del cuerpo principal
C_ANEXO_HDR    = RGBColor(0x2E, 0x2E, 0x2E)
C_ANEXO_ACENTO = RGBColor(0x8A, 0x6E, 0x3C)
C_ANEXO_FONDO  = RGBColor(0xF2, 0xF0, 0xEB)

# Dimensiones slide 16:9
W = Inches(13.333)
H = Inches(7.5)

# ---------------------------------------------------------------------------
# TEXTOS LITERALES DEL EMPAQUE
# ---------------------------------------------------------------------------

TIRO_LENGUETA = "NUEVO\nLA PROTECCION MAS COMPLETA"
TIRO_CLAIM1   = "El mas rapido ante parpadeos (< 0,03 s)"

TIRO_CLAIM2_V1 = "Autoproteccion termica activa*"
TIRO_CLAIM2_V2 = "Escudo Termico NTC*"

# V3 sub-opciones — textos literales (no modificar)
TIRO_CLAIM2_V3A = "Respaldo termico ante el breaker*"
TIRO_CLAIM2_V3B = "Respaldo termico ante fallas del termomagnetico*"
TIRO_CLAIM2_V3C = "Ultima linea de defensa electrica*"
TIRO_SUBTEXTO_V3C = "De tu instalacion electrica"

CARACT_COMUNES = [
    "- El mas rapido ante parpadeos: tiempo de respuesta < 30 ms (< 0,03 s),\n"
    "  especialmente adecuado para equipos con tecnologia inverter y cargas de arranque corto.",
    "- Protege tecnologia Inverter: la velocidad de respuesta de < 0,03 s minimiza la exposicion\n"
    "  de la electronica de control inverter a condiciones de inestabilidad de red.",
    "- Curva de disparo de tiempo inverso: respuesta mas rapida ante mayor sobrecarga --\n"
    "  configuracion tecnica complementaria al breaker termomagnetico del circuito.",
    "- Protege contra sobre voltaje, bajo voltaje, parpadeos e inestabilidad de la red electrica.",
    "- Supresor de picos incorporado.",
    "- Temporizado inteligente de reconexion.",
]

CARACT_TERMICO_V1 = (
    "- Autoproteccion termica activa*: autoproteccion del protector y del cableado de la\n"
    "  instalacion ante corrientes excesivas o conexiones deficientes."
)
CARACT_TERMICO_V2 = (
    "- Escudo Termico NTC*: autoproteccion del protector y del cableado de la\n"
    "  instalacion ante corrientes excesivas o conexiones deficientes."
)
CARACT_TERMICO_V3A = (
    "- Respaldo termico ante el breaker*: autoproteccion del protector y del cableado de la\n"
    "  instalacion ante corrientes excesivas o conexiones deficientes."
)
CARACT_TERMICO_V3B = (
    "- Respaldo termico ante fallas del termomagnetico*: autoproteccion del protector y del cableado\n"
    "  de la instalacion ante corrientes excesivas o conexiones deficientes."
)
CARACT_TERMICO_V3C = (
    "- Ultima linea de defensa electrica*: autoproteccion del protector y del cableado de la\n"
    "  instalacion ante corrientes excesivas o conexiones deficientes."
)

INFO_SEGURIDAD = "No reemplaza los breakers termomagneticos de la instalacion electrica."

CAVEAT_VELOCIDAD = (
    "*El tiempo de desconexion de menos de 30 milisegundos (< 0,03 s) aplica ante parpadeos "
    "(fluctuaciones rapidas del voltaje de la red electrica) e inestabilidad de la red. "
    "No aplica a la desconexion ante sobre voltaje o bajo voltaje pronunciados, cuyo tiempo "
    "de desconexion es de 0,4 a 3 segundos segun la intensidad de la falla. Genteca es el "
    "protector enchufable monofasico con el menor tiempo de respuesta verificado ante "
    "parpadeos en el mercado venezolano. Segun pruebas de laboratorio I&D Genteca y "
    "verificacion comparativa de competidores."
)

CAVEAT_INVERTER = (
    "*La proteccion ante parpadeos (fluctuaciones rapidas del voltaje de la red) que ofrece "
    "este protector es especialmente beneficiosa para equipos con tecnologia inverter, cuya "
    "electronica de control es sensible a variaciones rapidas del voltaje. Este protector no "
    "reemplaza la proteccion contra transientes de alta energia (descargas atmosfericas, "
    "conmutacion inductiva) presente en el equipo inverter de fabrica. Ambas protecciones "
    "son complementarias."
)

CUERPO_TERMICO = (
    "sensor de temperatura ubicado junto al rele de potencia. "
    "Detecta calentamiento excesivo causado por sobrecorrientes o por conexiones deficientes "
    "(bornes flojos o falsos contactos) y desconecta la carga antes de que el cableado o las "
    "conexiones se danen. Protege al protector mismo y a la instalacion electrica. "
    "Para cargas de baja demanda de corriente, protege al protector y al cableado, pero no "
    "actua como proteccion de sobrecarga directa de la carga conectada. "
    "Funciona como capa adicional de proteccion termica; no reemplaza al interruptor "
    "termomagnetico de la instalacion."
)

CAVEAT_TERMICO_V1 = "*Autoproteccion termica activa: " + CUERPO_TERMICO
CAVEAT_TERMICO_V2 = "*Escudo Termico NTC: " + CUERPO_TERMICO
CAVEAT_TERMICO_V3A = "*Respaldo termico ante el breaker: " + CUERPO_TERMICO
CAVEAT_TERMICO_V3B = "*Respaldo termico ante fallas del termomagnetico: " + CUERPO_TERMICO
CAVEAT_TERMICO_V3C = "*Ultima linea de defensa electrica: " + CUERPO_TERMICO

# ---------------------------------------------------------------------------
# HELPERS
# ---------------------------------------------------------------------------

def add_slide(prs, layout_idx=6):
    layout = prs.slide_layouts[layout_idx]
    return prs.slides.add_slide(layout)


def rgb(r, g, b):
    return RGBColor(r, g, b)


def add_textbox(slide, left, top, width, height,
                text="", font_size=14, bold=False, italic=False,
                color=None, bg_color=None, align=PP_ALIGN.LEFT,
                border_color=None, word_wrap=True, font_name="Calibri"):
    from pptx.util import Pt
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap

    if bg_color:
        fill = txBox.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

    if border_color:
        ln = txBox.line
        ln.color.rgb = border_color
        ln.width = Pt(0.75)

    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_name
    if color:
        run.font.color.rgb = color
    return txBox


def add_multiline_textbox(slide, left, top, width, height,
                          lines, font_size=11, bold_first=False,
                          color=None, bg_color=None, border_color=None,
                          line_spacing=None, font_name="Calibri"):
    from pptx.oxml.ns import qn
    from lxml import etree
    from pptx.util import Pt

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    if bg_color:
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = bg_color
    if border_color:
        txBox.line.color.rgb = border_color
        txBox.line.width = Pt(0.75)

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.name = font_name
        run.font.bold = (bold_first and i == 0)
        if color:
            run.font.color.rgb = color
    return txBox


def add_rect_bg(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(
        1,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    return shape


def title_bar(slide, text, sub=None, bg=C_AZUL_OSCURO, fg=C_BLANCO, sub_fg=C_AMBAR):
    add_rect_bg(slide, Inches(0), Inches(0), W, Inches(1.1), bg)
    add_textbox(slide, Inches(0.3), Inches(0.1), Inches(12.7), Inches(0.6),
                text=text, font_size=22, bold=True, color=fg, font_name="Calibri Light")
    if sub:
        add_textbox(slide, Inches(0.3), Inches(0.68), Inches(12.7), Inches(0.38),
                    text=sub, font_size=13, bold=False, color=sub_fg, font_name="Calibri")


def title_bar_anexo(slide, text, sub=None):
    add_rect_bg(slide, Inches(0), Inches(0), W, Inches(1.1), C_ANEXO_HDR)
    add_rect_bg(slide, Inches(0), Inches(0), Inches(1.2), Inches(1.1), C_ANEXO_ACENTO)
    add_textbox(slide, Inches(0.05), Inches(0.35), Inches(1.12), Inches(0.42),
                text="ANEXO", font_size=11, bold=True, color=C_BLANCO,
                align=PP_ALIGN.CENTER, font_name="Calibri")
    add_textbox(slide, Inches(1.3), Inches(0.1), Inches(11.7), Inches(0.6),
                text=text, font_size=20, bold=True, color=C_BLANCO, font_name="Calibri Light")
    if sub:
        add_textbox(slide, Inches(1.3), Inches(0.68), Inches(11.7), Inches(0.38),
                    text=sub, font_size=12, bold=False, color=C_ANEXO_ACENTO, font_name="Calibri")


def footer_bar(slide, text="Genteca | Junta Directiva | 2026-05-06 | Confidencial"):
    add_rect_bg(slide, Inches(0), Inches(7.1), W, Inches(0.4), C_AZUL_OSCURO)
    add_textbox(slide, Inches(0.3), Inches(7.12), Inches(12.7), Inches(0.35),
                text=text, font_size=9, color=C_GRIS_LINEA, font_name="Calibri")


def footer_bar_anexo(slide, num_label=""):
    add_rect_bg(slide, Inches(0), Inches(7.1), W, Inches(0.4), C_ANEXO_HDR)
    label = f"Genteca | Junta Directiva | 2026-05-06 | Confidencial | {num_label}" if num_label else \
            "Genteca | Junta Directiva | 2026-05-06 | Confidencial"
    add_textbox(slide, Inches(0.3), Inches(7.12), Inches(12.7), Inches(0.35),
                text=label, font_size=9, color=C_GRIS_LINEA, font_name="Calibri")


def label_chip(slide, left, top, text, bg, fg=C_BLANCO):
    add_rect_bg(slide, left, top, Inches(2.2), Inches(0.32), bg)
    add_textbox(slide, left + Inches(0.08), top + Inches(0.03),
                Inches(2.1), Inches(0.28),
                text=text, font_size=10, bold=True, color=fg, font_name="Calibri")


# ---------------------------------------------------------------------------
# SLIDES — CUERPO PRINCIPAL
# ---------------------------------------------------------------------------

def slide_portada(prs):
    slide = add_slide(prs)
    add_rect_bg(slide, Inches(0), Inches(0), W, H, C_AZUL_OSCURO)
    add_rect_bg(slide, Inches(0), Inches(3.35), W, Inches(0.06), C_AMBAR)
    add_textbox(slide, Inches(0.5), Inches(0.45), Inches(5), Inches(0.5),
                text="GENTECA", font_size=16, bold=True, color=C_AMBAR,
                font_name="Calibri Light")
    add_textbox(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(1.6),
                text="Avance empaque GSM-MB / RB / RF / RE",
                font_size=34, bold=True, color=C_BLANCO, font_name="Calibri Light")
    add_textbox(slide, Inches(0.5), Inches(2.85), Inches(12.3), Inches(0.6),
                text="Tres variantes del badge termico: B-sin-NTC, B-con-NTC y enfoque beneficio del beneficio",
                font_size=20, bold=False, color=C_AZUL_CLARO, font_name="Calibri Light")
    add_rect_bg(slide, Inches(0.5), Inches(3.6), Inches(12.3), Inches(0.7),
                RGBColor(0x0D, 0x1E, 0x33))
    add_textbox(slide, Inches(0.65), Inches(3.68), Inches(12.0), Inches(0.55),
                text="Decision solicitada: cual variante (y sub-opcion si V3) avanza a arte final",
                font_size=16, bold=True, color=C_AMBAR, font_name="Calibri")
    add_textbox(slide, Inches(0.5), Inches(4.55), Inches(6), Inches(0.38),
                text="Junta Directiva Genteca", font_size=13, color=C_GRIS_LINEA,
                font_name="Calibri")
    add_textbox(slide, Inches(0.5), Inches(4.9), Inches(6), Inches(0.38),
                text="2026-05-06", font_size=13, color=C_GRIS_LINEA,
                font_name="Calibri")
    add_textbox(slide, Inches(0.5), Inches(5.25), Inches(6), Inches(0.38),
                text="CONFIDENCIAL", font_size=11, bold=True,
                color=RGBColor(0xFF, 0x60, 0x60), font_name="Calibri")


def slide_contexto(prs):
    slide = add_slide(prs)
    title_bar(slide,
              "Contexto y pregunta a la Junta",
              "La Alternativa B esta aprobada. La decision de hoy es el nombre del badge termico (3 variantes).")
    footer_bar(slide)

    bloques = [
        {
            "num": "1",
            "titulo": "La Junta eligio Alternativa B",
            "texto": (
                "En sesion anterior, la Junta selecciono la Alternativa B para el empaque GSM-MB / RB / RF / RE. "
                "Esa decision esta cerrada: dos claims en el tiro (velocidad + funcion termica), "
                "mas la tecnologia Inverter desarrollada en el retiro."
            )
        },
        {
            "num": "2",
            "titulo": "El claim de velocidad no cambia",
            "texto": (
                "\"El mas rapido ante parpadeos (< 0,03 s)\" -- es el elemento dominante del tiro "
                "en las tres variantes. Identico, aprobado, invariable."
            )
        },
        {
            "num": "3",
            "titulo": "Como se nombra la funcion termica: tres ejes posibles",
            "texto": (
                "Variante 1 (B-sin-NTC): nombre funcional puro, sin revelar componente. "
                "Variante 2 (B-con-NTC): NTC visible en empaque, tesis anti-copia activada. "
                "Variante 3 (beneficio del beneficio, Jesus Maria): el claim comunica la consecuencia ultima -- "
                "proteccion de la instalacion si el breaker falla. Tres sub-opciones de tono (V3a/V3b/V3c). "
                "Jose Miguel Canudas elige cual eje avanza y, si es V3, cual sub-opcion."
            )
        },
    ]

    tops = [Inches(1.3), Inches(3.0), Inches(4.55)]
    for i, bloque in enumerate(bloques):
        top = tops[i]
        add_rect_bg(slide, Inches(0.35), top, Inches(0.45), Inches(0.45), C_AMBAR)
        add_textbox(slide, Inches(0.35), top + Inches(0.02), Inches(0.45), Inches(0.42),
                    text=bloque["num"], font_size=18, bold=True, color=C_BLANCO,
                    align=PP_ALIGN.CENTER, font_name="Calibri")
        add_textbox(slide, Inches(0.95), top - Inches(0.02), Inches(11.8), Inches(0.4),
                    text=bloque["titulo"], font_size=14, bold=True, color=C_AZUL_OSCURO,
                    font_name="Calibri")
        add_textbox(slide, Inches(0.95), top + Inches(0.36), Inches(11.8), Inches(0.75),
                    text=bloque["texto"], font_size=11, color=C_GRIS_OSCURO,
                    font_name="Calibri")

    add_rect_bg(slide, Inches(0.35), Inches(6.4), Inches(12.6), Inches(0.55),
                C_AZUL_OSCURO)
    add_textbox(slide, Inches(0.55), Inches(6.45), Inches(12.3), Inches(0.45),
                text="Nivel 1: ¿V1, V2 o V3? -- Nivel 2 (solo si V3): ¿V3a, V3b o V3c?",
                font_size=13, bold=True, color=C_AMBAR, font_name="Calibri")


def slide_comparacion(prs):
    """Slide 3 — Comparacion lado a lado: 3 columnas V1 / V2 / V3."""
    slide = add_slide(prs)
    title_bar(slide,
              "Comparacion -- Variante 1 vs. Variante 2 vs. Variante 3",
              "La diferencia entre variantes es exclusivamente el badge del segundo claim del tiro.")
    footer_bar(slide)

    # 3 columnas: ancho reducido para caber las tres
    col_w = Inches(4.2)
    col1_x = Inches(0.22)
    col2_x = Inches(4.57)
    col3_x = Inches(8.92)
    row_start = Inches(1.2)
    row_h = Inches(0.52)

    # Headers de columna
    add_rect_bg(slide, col1_x, row_start, col_w, Inches(0.55), C_GRIS_V1)
    add_textbox(slide, col1_x + Inches(0.1), row_start + Inches(0.08),
                col_w - Inches(0.15), Inches(0.42),
                text="VARIANTE 1 -- B-sin-NTC", font_size=12, bold=True,
                color=C_BLANCO, font_name="Calibri")

    add_rect_bg(slide, col2_x, row_start, col_w, Inches(0.55), C_AZUL_V2)
    add_textbox(slide, col2_x + Inches(0.1), row_start + Inches(0.08),
                col_w - Inches(0.15), Inches(0.42),
                text="VARIANTE 2 -- B-con-NTC", font_size=12, bold=True,
                color=C_BLANCO, font_name="Calibri")

    add_rect_bg(slide, col3_x, row_start, col_w, Inches(0.55), C_SIENA_V3)
    add_textbox(slide, col3_x + Inches(0.1), row_start + Inches(0.08),
                col_w - Inches(0.15), Inches(0.42),
                text="VARIANTE 3 -- Beneficio del beneficio", font_size=12, bold=True,
                color=C_BLANCO, font_name="Calibri")

    filas = [
        ("Badge en el tiro",
         "Autoproteccion termica activa*",
         "Escudo Termico NTC*",
         "V3a: Respaldo termico ante el breaker*\nV3b: Respaldo termico ante fallas del termomagnetico*\nV3c: Ultima linea de defensa electrica*"),
        ("Eje de comunicacion",
         "Funcion tecnica del protector.",
         "Componente NTC visible como diferenciador.",
         "Consecuencia ultima: respaldo si el breaker falla o esta mal seleccionado."),
        ("Requiere verificacion SAPI VE antes de imprenta",
         "No.",
         "Si. Una semana adicional.",
         "No para V3a/V3b/V3c (claim funcional). Confirmar con abogado si se desea TM."),
        ("Condicion tecnica Vera P-5",
         "Si (si NTC es fusible unico, badge de respaldo aprobado).",
         "No aplica.",
         "No aplica. El claim no alude al mecanismo NTC."),
        ("Riesgo de malentendido",
         "\"Activa\" puede leerse como que el protector actua sobre el equipo conectado.",
         "NTC puede percibirse como marca del componente, no del protector.",
         "V3c (\"ultima linea\") requiere subtexto obligatorio en tiro. V3a coloquialismo confirmar."),
        ("Tiempo estimado de cierre",
         "Pocos dias.",
         "Aprox. una semana (SAPI VE).",
         "Pocos dias (V3a/V3b). V3c requiere ajuste de layout tiro por subtexto."),
    ]

    y = row_start + Inches(0.6)
    for idx, (dim, v1, v2, v3) in enumerate(filas):
        # Altura variable segun contenido de V3
        extra = 0.38 if "\n" in v3 else 0
        rh = Inches(0.52 + extra)
        bg = C_BLANCO if idx % 2 == 0 else RGBColor(0xF5, 0xF8, 0xFC)

        for cx in [col1_x, col2_x, col3_x]:
            add_rect_bg(slide, cx, y, col_w, rh, bg, border_color=C_GRIS_LINEA)

        # Dimension label (span col1)
        add_textbox(slide, col1_x + Inches(0.06), y + Inches(0.03),
                    col_w - Inches(0.1), Inches(0.2),
                    text=dim, font_size=8, bold=True, color=C_GRIS_MEDIO,
                    font_name="Calibri")

        for cx, val in [(col1_x, v1), (col2_x, v2), (col3_x, v3)]:
            add_textbox(slide, cx + Inches(0.06), y + Inches(0.2),
                        col_w - Inches(0.12), rh - Inches(0.22),
                        text=val, font_size=9, color=C_GRIS_OSCURO,
                        font_name="Calibri")

        y += rh


def slide_tiro(prs, variante, badge_claim2, color_header, color_header_light):
    """Tiro completo para V1 o V2."""
    num = "1" if variante == "V1" else "2"
    nombre = "B-sin-NTC" if variante == "V1" else "B-con-NTC"
    nota_tm = "" if variante == "V1" else "  (sin TM hasta verificacion SAPI VE)"

    slide = add_slide(prs)
    title_bar(slide,
              f"Variante {num} -- {nombre}: tiro completo",
              f"Textos literales aprobados del frente del empaque{nota_tm}",
              bg=color_header)
    footer_bar(slide)

    caja_left = Inches(1.8)
    caja_top = Inches(1.25)
    caja_w = Inches(9.7)
    caja_h = Inches(5.5)

    add_rect_bg(slide, caja_left, caja_top, caja_w, caja_h,
                color_header_light, border_color=color_header)
    add_rect_bg(slide, caja_left, caja_top, caja_w, Inches(0.38), color_header)
    add_textbox(slide, caja_left + Inches(0.15), caja_top + Inches(0.05),
                caja_w - Inches(0.3), Inches(0.32),
                text="TIRO / FRENTE DEL EMPAQUE -- TEXTOS APROBADOS",
                font_size=10, bold=True, color=C_BLANCO, font_name="Calibri")

    add_textbox(slide, caja_left + Inches(0.3), caja_top + Inches(0.55),
                caja_w - Inches(0.6), Inches(0.28),
                text="LENGUETA", font_size=9, bold=True,
                color=color_header, font_name="Calibri")
    add_rect_bg(slide, caja_left + Inches(0.3), caja_top + Inches(0.82),
                caja_w - Inches(0.6), Inches(0.75),
                C_BLANCO, border_color=color_header)
    add_textbox(slide, caja_left + Inches(0.45), caja_top + Inches(0.88),
                caja_w - Inches(0.8), Inches(0.65),
                text=TIRO_LENGUETA, font_size=13, bold=True,
                color=C_AZUL_OSCURO, font_name="Calibri")

    add_textbox(slide, caja_left + Inches(0.3), caja_top + Inches(1.72),
                caja_w - Inches(0.6), Inches(0.28),
                text="CLAIM 1 -- ELEMENTO DOMINANTE", font_size=9, bold=True,
                color=color_header, font_name="Calibri")
    add_rect_bg(slide, caja_left + Inches(0.3), caja_top + Inches(1.98),
                caja_w - Inches(0.6), Inches(0.72),
                C_BLANCO, border_color=color_header)
    add_textbox(slide, caja_left + Inches(0.45), caja_top + Inches(2.04),
                caja_w - Inches(0.8), Inches(0.62),
                text=TIRO_CLAIM1, font_size=18, bold=True,
                color=C_AZUL_OSCURO, font_name="Calibri")

    add_textbox(slide, caja_left + Inches(0.3), caja_top + Inches(2.85),
                caja_w - Inches(0.6), Inches(0.28),
                text="CLAIM 2 -- BADGE TERMICO (jerarquia menor)",
                font_size=9, bold=True, color=color_header, font_name="Calibri")
    add_rect_bg(slide, caja_left + Inches(0.3), caja_top + Inches(3.11),
                caja_w - Inches(0.6), Inches(0.72),
                color_header_light, border_color=color_header)
    add_rect_bg(slide, caja_left + Inches(0.28), caja_top + Inches(3.09),
                Inches(0.04), Inches(0.76), color_header)
    add_textbox(slide, caja_left + Inches(0.45), caja_top + Inches(3.17),
                caja_w - Inches(0.8), Inches(0.62),
                text=badge_claim2, font_size=16, bold=True,
                color=color_header, font_name="Calibri")

    add_textbox(slide, caja_left + Inches(0.3), caja_top + Inches(4.05),
                caja_w - Inches(0.6), Inches(0.35),
                text="* Asterisco obligatorio -- remite al caveat literal del retiro (ver slides siguientes)",
                font_size=9, italic=True, color=C_GRIS_MEDIO, font_name="Calibri")


def slide_tiro_v3(prs):
    """Tiro V3: tres sub-opciones en bloques paralelos."""
    slide = add_slide(prs)
    title_bar(slide,
              "Variante 3 -- Beneficio del beneficio: tiro completo (3 sub-opciones)",
              "Tres registros tonales para el mismo eje conceptual -- la Junta elige sub-opcion si V3 avanza",
              bg=C_SIENA_V3)
    footer_bar(slide)

    # Encabezado de contexto
    add_rect_bg(slide, Inches(0.3), Inches(1.18), Inches(12.7), Inches(0.34),
                C_SIENA_V3_CLR, border_color=C_SIENA_V3)
    add_textbox(slide, Inches(0.45), Inches(1.22), Inches(12.4), Inches(0.28),
                text="Lengeta y Claim 1 son identicos a V1/V2. Solo el badge termico (Claim 2) difiere entre sub-opciones.",
                font_size=9.5, italic=True, color=C_SIENA_V3, font_name="Calibri")

    # Tres bloques paralelos
    sub_opciones = [
        {
            "label": "V3a -- Tecnica-coloquial",
            "badge": TIRO_CLAIM2_V3A,
            "subtexto": None,
            "nota": "Coloquialismo \"breaker\": confirmar registro de audiencia objetivo antes de imprenta.",
        },
        {
            "label": "V3b -- Tecnica-formal (menor riesgo combinado)",
            "badge": TIRO_CLAIM2_V3B,
            "subtexto": None,
            "nota": "Formulacion mas formal. Verificar espacio disponible en layout de diseno.",
        },
        {
            "label": "V3c -- Emocional / loss aversion",
            "badge": TIRO_CLAIM2_V3C,
            "subtexto": TIRO_SUBTEXTO_V3C,
            "nota": "Subtexto obligatorio adyacente en tiro fisico. Requiere ajuste de layout con disenador.",
        },
    ]

    col_w = Inches(4.1)
    col_xs = [Inches(0.3), Inches(4.62), Inches(8.94)]
    bloque_top = Inches(1.62)
    bloque_h = Inches(5.28)

    for col_x, sub in zip(col_xs, sub_opciones):
        # Marco del bloque
        add_rect_bg(slide, col_x, bloque_top, col_w, bloque_h,
                    C_SIENA_V3_CLR, border_color=C_SIENA_V3)

        # Header sub-opcion
        add_rect_bg(slide, col_x, bloque_top, col_w, Inches(0.36), C_SIENA_V3)
        add_textbox(slide, col_x + Inches(0.1), bloque_top + Inches(0.06),
                    col_w - Inches(0.15), Inches(0.28),
                    text=sub["label"], font_size=9.5, bold=True,
                    color=C_BLANCO, font_name="Calibri")

        y = bloque_top + Inches(0.44)

        # Lengueta
        add_textbox(slide, col_x + Inches(0.15), y, col_w - Inches(0.25), Inches(0.2),
                    text="LENGUETA", font_size=8, bold=True, color=C_SIENA_V3, font_name="Calibri")
        y += Inches(0.2)
        add_rect_bg(slide, col_x + Inches(0.15), y, col_w - Inches(0.25), Inches(0.58),
                    C_BLANCO, border_color=C_SIENA_V3)
        add_textbox(slide, col_x + Inches(0.25), y + Inches(0.06),
                    col_w - Inches(0.42), Inches(0.48),
                    text=TIRO_LENGUETA, font_size=9, bold=True,
                    color=C_AZUL_OSCURO, font_name="Calibri")
        y += Inches(0.65)

        # Claim 1
        add_textbox(slide, col_x + Inches(0.15), y, col_w - Inches(0.25), Inches(0.2),
                    text="CLAIM 1", font_size=8, bold=True, color=C_SIENA_V3, font_name="Calibri")
        y += Inches(0.2)
        add_rect_bg(slide, col_x + Inches(0.15), y, col_w - Inches(0.25), Inches(0.58),
                    C_BLANCO, border_color=C_SIENA_V3)
        add_textbox(slide, col_x + Inches(0.25), y + Inches(0.06),
                    col_w - Inches(0.42), Inches(0.48),
                    text=TIRO_CLAIM1, font_size=10, bold=True,
                    color=C_AZUL_OSCURO, font_name="Calibri")
        y += Inches(0.65)

        # Claim 2 — badge diferenciado
        add_textbox(slide, col_x + Inches(0.15), y, col_w - Inches(0.25), Inches(0.2),
                    text="CLAIM 2 -- BADGE TERMICO", font_size=8, bold=True,
                    color=C_SIENA_V3, font_name="Calibri")
        y += Inches(0.2)

        badge_h = Inches(0.72) if not sub["subtexto"] else Inches(1.1)
        add_rect_bg(slide, col_x + Inches(0.15), y, col_w - Inches(0.25), badge_h,
                    C_SIENA_V3_CLR, border_color=C_SIENA_V3)
        add_rect_bg(slide, col_x + Inches(0.13), y, Inches(0.04), badge_h, C_SIENA_V3)
        add_textbox(slide, col_x + Inches(0.28), y + Inches(0.08),
                    col_w - Inches(0.42), Inches(0.45),
                    text=sub["badge"], font_size=11, bold=True,
                    color=C_SIENA_V3, font_name="Calibri")
        if sub["subtexto"]:
            add_textbox(slide, col_x + Inches(0.28), y + Inches(0.54),
                        col_w - Inches(0.42), Inches(0.35),
                        text=sub["subtexto"], font_size=9.5, italic=True,
                        color=C_SIENA_V3, font_name="Calibri")
        y += badge_h + Inches(0.18)

        # Nota condicion
        add_textbox(slide, col_x + Inches(0.15), y,
                    col_w - Inches(0.25), Inches(0.7),
                    text=sub["nota"], font_size=8.5, italic=True,
                    color=C_GRIS_MEDIO, font_name="Calibri")


def slide_retiro_caract(prs, variante, bullet_termico, color_header, color_header_light):
    """Retiro CARACTERISTICAS para V1 o V2."""
    num = "1" if variante == "V1" else "2"
    nombre = "B-sin-NTC" if variante == "V1" else "B-con-NTC"

    slide = add_slide(prs)
    title_bar(slide,
              f"Variante {num} -- {nombre}: retiro -- CARACTERISTICAS",
              "Textos literales aprobados del reverso del empaque (bullets de caracteristicas)",
              bg=color_header)
    footer_bar(slide)

    bullets_full = CARACT_COMUNES[:2] + [bullet_termico] + CARACT_COMUNES[2:]

    caja_left = Inches(0.4)
    caja_top = Inches(1.25)
    caja_w = Inches(12.5)

    add_rect_bg(slide, caja_left, caja_top, caja_w, Inches(0.38), color_header)
    add_textbox(slide, caja_left + Inches(0.15), caja_top + Inches(0.05),
                caja_w - Inches(0.3), Inches(0.32),
                text="RETIRO / REVERSO DEL EMPAQUE -- SECCION CARACTERISTICAS",
                font_size=10, bold=True, color=C_BLANCO, font_name="Calibri")

    y = caja_top + Inches(0.45)
    for i, bullet in enumerate(bullets_full):
        bg = color_header_light if i % 2 == 0 else C_BLANCO
        num_lines = bullet.count("\n") + 1 + (len(bullet) // 110)
        bh = Inches(0.38) * max(1, num_lines)

        add_rect_bg(slide, caja_left, y, caja_w, bh, bg, border_color=C_GRIS_LINEA)
        add_textbox(slide, caja_left + Inches(0.15), y + Inches(0.04),
                    caja_w - Inches(0.3), bh - Inches(0.06),
                    text=bullet, font_size=10.5, color=C_GRIS_OSCURO,
                    font_name="Calibri")
        y += bh

    y += Inches(0.12)
    add_rect_bg(slide, caja_left, y, caja_w, Inches(0.42),
                RGBColor(0xFF, 0xF3, 0xCD), border_color=RGBColor(0xD4, 0xA0, 0x17))
    add_textbox(slide, caja_left + Inches(0.15), y + Inches(0.06),
                caja_w - Inches(0.3), Inches(0.34),
                text="INFORMACION DE SEGURIDAD:  " + INFO_SEGURIDAD,
                font_size=10.5, bold=True,
                color=RGBColor(0x7B, 0x4F, 0x00), font_name="Calibri")


def slide_retiro_caract_v3(prs):
    """Retiro CARACTERISTICAS V3 — muestra las 3 variantes del bullet termico."""
    slide = add_slide(prs)
    title_bar(slide,
              "Variante 3 -- retiro: CARACTERISTICAS (3 sub-opciones del bullet termico)",
              "El resto de los bullets es invariable. Solo el bullet termico difiere segun sub-opcion elegida.",
              bg=C_SIENA_V3)
    footer_bar(slide)

    caja_left = Inches(0.4)
    caja_top = Inches(1.2)
    caja_w = Inches(12.5)

    # Nota aclaratoria
    add_rect_bg(slide, caja_left, caja_top, caja_w, Inches(0.34),
                C_SIENA_V3_CLR, border_color=C_SIENA_V3)
    add_textbox(slide, caja_left + Inches(0.15), caja_top + Inches(0.06),
                caja_w - Inches(0.25), Inches(0.26),
                text="Bullets en fondo crema: identicos en V3a, V3b y V3c. Bullet termico: se muestra en las 3 variantes de header.",
                font_size=9, italic=True, color=C_SIENA_V3, font_name="Calibri")

    y = caja_top + Inches(0.42)

    # Bullets comunes 1 y 2
    for i, bullet in enumerate(CARACT_COMUNES[:2]):
        bg = C_SIENA_V3_CLR if i % 2 == 0 else C_BLANCO
        num_lines = bullet.count("\n") + 1 + (len(bullet) // 110)
        bh = Inches(0.38) * max(1, num_lines)
        add_rect_bg(slide, caja_left, y, caja_w, bh, bg, border_color=C_GRIS_LINEA)
        add_textbox(slide, caja_left + Inches(0.15), y + Inches(0.04),
                    caja_w - Inches(0.3), bh - Inches(0.06),
                    text=bullet, font_size=10, color=C_GRIS_OSCURO, font_name="Calibri")
        y += bh

    # Bloque bullet termico — 3 sub-opciones
    add_rect_bg(slide, caja_left, y, caja_w, Inches(0.32), C_SIENA_V3)
    add_textbox(slide, caja_left + Inches(0.15), y + Inches(0.05),
                caja_w - Inches(0.25), Inches(0.25),
                text="BULLET TERMICO -- TRES VARIANTES DE HEADER (elegir segun sub-opcion V3a/V3b/V3c)",
                font_size=8.5, bold=True, color=C_BLANCO, font_name="Calibri")
    y += Inches(0.32)

    bullets_v3 = [
        ("V3a", CARACT_TERMICO_V3A),
        ("V3b", CARACT_TERMICO_V3B),
        ("V3c", CARACT_TERMICO_V3C),
    ]
    for etiqueta, bullet in bullets_v3:
        num_lines = bullet.count("\n") + 1 + (len(bullet) // 105)
        bh = Inches(0.38) * max(1, num_lines)
        add_rect_bg(slide, caja_left, y, caja_w, bh,
                    C_SIENA_V3_CLR, border_color=C_SIENA_V3)
        add_rect_bg(slide, caja_left, y, Inches(0.52), bh, C_SIENA_V3)
        add_textbox(slide, caja_left + Inches(0.04), y + Inches(0.08),
                    Inches(0.46), Inches(0.28),
                    text=etiqueta, font_size=8.5, bold=True, color=C_BLANCO,
                    align=PP_ALIGN.CENTER, font_name="Calibri")
        add_textbox(slide, caja_left + Inches(0.6), y + Inches(0.04),
                    caja_w - Inches(0.7), bh - Inches(0.06),
                    text=bullet, font_size=10, color=C_GRIS_OSCURO, font_name="Calibri")
        y += bh

    # Bullets comunes restantes
    for i, bullet in enumerate(CARACT_COMUNES[2:]):
        bg = C_BLANCO if i % 2 == 0 else C_SIENA_V3_CLR
        num_lines = bullet.count("\n") + 1 + (len(bullet) // 110)
        bh = Inches(0.38) * max(1, num_lines)
        add_rect_bg(slide, caja_left, y, caja_w, bh, bg, border_color=C_GRIS_LINEA)
        add_textbox(slide, caja_left + Inches(0.15), y + Inches(0.04),
                    caja_w - Inches(0.3), bh - Inches(0.06),
                    text=bullet, font_size=10, color=C_GRIS_OSCURO, font_name="Calibri")
        y += bh

    y += Inches(0.1)
    add_rect_bg(slide, caja_left, y, caja_w, Inches(0.42),
                RGBColor(0xFF, 0xF3, 0xCD), border_color=RGBColor(0xD4, 0xA0, 0x17))
    add_textbox(slide, caja_left + Inches(0.15), y + Inches(0.06),
                caja_w - Inches(0.3), Inches(0.34),
                text="INFORMACION DE SEGURIDAD:  " + INFO_SEGURIDAD,
                font_size=10.5, bold=True,
                color=RGBColor(0x7B, 0x4F, 0x00), font_name="Calibri")


def slide_caveats_compartidos(prs):
    """Caveats heredados compartidos: velocidad + inverter — identicos en V1, V2 y V3."""
    slide = add_slide(prs)
    title_bar(
        slide,
        "Retiro -- CAVEATS compartidos (las tres variantes)",
        "Caveats de velocidad e inverter: texto literal aprobado -- identico en V1, V2 y V3",
        bg=C_AZUL_OSCURO
    )
    footer_bar(slide)

    add_rect_bg(slide, Inches(0.4), Inches(1.18), Inches(12.5), Inches(0.34),
                RGBColor(0xF0, 0xF4, 0xF8), border_color=C_GRIS_LINEA)
    add_textbox(slide, Inches(0.55), Inches(1.22), Inches(12.2), Inches(0.28),
                text="Estos dos caveats son invariables. Aplican palabra por palabra a las tres variantes sin modificacion.",
                font_size=9.5, italic=True, color=C_GRIS_MEDIO, font_name="Calibri")

    caja_left = Inches(0.4)
    caja_w = Inches(12.5)

    caveats_compartidos = [
        {
            "header": "CAVEAT 1 -- VELOCIDAD (Bruna §7.3 -- invariable en V1, V2 y V3)",
            "texto": CAVEAT_VELOCIDAD,
        },
        {
            "header": "CAVEAT 2 -- INVERTER (Bruna §2 Claim C -- invariable en V1, V2 y V3)",
            "texto": CAVEAT_INVERTER,
        },
    ]

    y = Inches(1.62)
    for i, cav in enumerate(caveats_compartidos):
        hdr_h = Inches(0.36)
        add_rect_bg(slide, caja_left, y, caja_w, hdr_h, C_AZUL_OSCURO)
        add_textbox(slide, caja_left + Inches(0.12), y + Inches(0.05),
                    caja_w - Inches(0.25), hdr_h - Inches(0.08),
                    text=cav["header"], font_size=9.5, bold=True,
                    color=C_BLANCO, font_name="Calibri")
        y += hdr_h

        char_est = len(cav["texto"])
        txt_h = Inches(0.5) + Inches(char_est / 480)

        add_rect_bg(slide, caja_left, y, caja_w, txt_h,
                    RGBColor(0xF5, 0xF8, 0xFC), border_color=C_GRIS_LINEA)
        add_textbox(slide, caja_left + Inches(0.15), y + Inches(0.08),
                    caja_w - Inches(0.3), txt_h - Inches(0.1),
                    text=cav["texto"], font_size=10, color=C_GRIS_OSCURO,
                    font_name="Calibri")
        y += txt_h + Inches(0.2)


def slide_caveat_termico(prs, variante, caveat_termico, color_header, color_header_light):
    """Caveat termico especifico para V1 o V2."""
    num = "1" if variante == "V1" else "2"
    nombre = "B-sin-NTC" if variante == "V1" else "B-con-NTC"
    ref_bruna = "Bruna §8.4" if variante == "V1" else "Bruna §10"

    slide = add_slide(prs)
    title_bar(
        slide,
        f"Variante {num} -- {nombre}: retiro -- CAVEAT TERMICO",
        f"Unico caveat que difiere entre variantes -- texto literal aprobado ({ref_bruna})",
        bg=color_header
    )
    footer_bar(slide)

    add_rect_bg(slide, Inches(0.4), Inches(1.18), Inches(12.5), Inches(0.34),
                color_header_light, border_color=color_header)
    add_textbox(slide, Inches(0.55), Inches(1.22), Inches(12.2), Inches(0.28),
                text="Este es el caveat termico especifico de esta variante. "
                     "Los caveats de velocidad e inverter son identicos en V1, V2 y V3 (ver slide anterior).",
                font_size=9.5, italic=True, color=color_header, font_name="Calibri")

    caja_left = Inches(0.4)
    caja_w = Inches(12.5)
    y = Inches(1.62)

    hdr_label = (
        f"CAVEAT 3 -- TERMICO ({ref_bruna}) -- "
        f"{'VARIANTE 1: Autoproteccion termica activa' if variante == 'V1' else 'VARIANTE 2: Escudo Termico NTC'}"
    )

    hdr_h = Inches(0.38)
    add_rect_bg(slide, caja_left, y, caja_w, hdr_h, color_header)
    add_textbox(slide, caja_left + Inches(0.12), y + Inches(0.05),
                caja_w - Inches(0.25), hdr_h - Inches(0.08),
                text=hdr_label, font_size=9.5, bold=True,
                color=C_BLANCO, font_name="Calibri")
    y += hdr_h

    char_est = len(caveat_termico)
    txt_h = Inches(0.6) + Inches(char_est / 420)

    add_rect_bg(slide, caja_left, y, caja_w, txt_h,
                color_header_light, border_color=color_header)
    add_rect_bg(slide, caja_left, y, Inches(0.06), txt_h, color_header)
    add_textbox(slide, caja_left + Inches(0.22), y + Inches(0.1),
                caja_w - Inches(0.35), txt_h - Inches(0.15),
                text=caveat_termico, font_size=10.5, color=C_GRIS_OSCURO,
                font_name="Calibri")
    y += txt_h + Inches(0.25)

    add_rect_bg(slide, caja_left, y, caja_w, Inches(0.38),
                RGBColor(0xF5, 0xF8, 0xFC), border_color=C_GRIS_LINEA)
    add_textbox(slide, caja_left + Inches(0.15), y + Inches(0.07),
                caja_w - Inches(0.3), Inches(0.28),
                text="Caveats 1 y 2 (velocidad e inverter): identicos en V1, V2 y V3. Ver slide anterior.",
                font_size=9.5, italic=True, color=C_GRIS_MEDIO, font_name="Calibri")


def slide_caveat_termico_v3(prs):
    """Caveat termico V3 — 3 headers + cuerpo invariable compartido."""
    slide = add_slide(prs)
    title_bar(
        slide,
        "Variante 3 -- retiro: CAVEAT TERMICO (3 sub-opciones de header)",
        "Cuerpo del caveat identico para V3a, V3b y V3c -- solo el header cambia segun sub-opcion",
        bg=C_SIENA_V3
    )
    footer_bar(slide)

    add_rect_bg(slide, Inches(0.4), Inches(1.18), Inches(12.5), Inches(0.34),
                C_SIENA_V3_CLR, border_color=C_SIENA_V3)
    add_textbox(slide, Inches(0.55), Inches(1.22), Inches(12.2), Inches(0.28),
                text="Los caveats de velocidad e inverter son identicos en V1, V2 y V3 (ver slide anterior). "
                     "El cuerpo §8.3 se muestra una sola vez al pie.",
                font_size=9.5, italic=True, color=C_SIENA_V3, font_name="Calibri")

    caja_left = Inches(0.4)
    caja_w = Inches(12.5)
    y = Inches(1.62)

    headers_v3 = [
        ("V3a", "Respaldo termico ante el breaker:", CAVEAT_TERMICO_V3A),
        ("V3b", "Respaldo termico ante fallas del termomagnetico:", CAVEAT_TERMICO_V3B),
        ("V3c", "Ultima linea de defensa electrica:", CAVEAT_TERMICO_V3C),
    ]

    for etiqueta, header_text, _ in headers_v3:
        hdr_h = Inches(0.38)
        add_rect_bg(slide, caja_left, y, caja_w, hdr_h, C_SIENA_V3)
        add_rect_bg(slide, caja_left, y, Inches(0.52), hdr_h, RGBColor(0x5C, 0x3A, 0x1E))
        add_textbox(slide, caja_left + Inches(0.04), y + Inches(0.08),
                    Inches(0.46), Inches(0.24),
                    text=etiqueta, font_size=8.5, bold=True, color=C_BLANCO,
                    align=PP_ALIGN.CENTER, font_name="Calibri")
        add_textbox(slide, caja_left + Inches(0.6), y + Inches(0.07),
                    caja_w - Inches(0.7), Inches(0.28),
                    text="CAVEAT 3 -- TERMICO -- " + header_text,
                    font_size=9.5, bold=True, color=C_BLANCO, font_name="Calibri")
        y += hdr_h

    # Cuerpo invariable — mostrar una sola vez
    y += Inches(0.15)
    add_rect_bg(slide, caja_left, y, caja_w, Inches(0.3), C_AZUL_OSCURO)
    add_textbox(slide, caja_left + Inches(0.15), y + Inches(0.05),
                caja_w - Inches(0.25), Inches(0.23),
                text="CUERPO §8.3 -- IDENTICO PARA LAS 3 SUB-OPCIONES V3a / V3b / V3c",
                font_size=8.5, bold=True, color=C_BLANCO, font_name="Calibri")
    y += Inches(0.3)

    char_est = len(CUERPO_TERMICO)
    txt_h = Inches(0.5) + Inches(char_est / 450)
    add_rect_bg(slide, caja_left, y, caja_w, txt_h,
                C_SIENA_V3_CLR, border_color=C_SIENA_V3)
    add_rect_bg(slide, caja_left, y, Inches(0.06), txt_h, C_SIENA_V3)
    add_textbox(slide, caja_left + Inches(0.22), y + Inches(0.1),
                caja_w - Inches(0.35), txt_h - Inches(0.15),
                text=CUERPO_TERMICO, font_size=10.5, color=C_GRIS_OSCURO,
                font_name="Calibri")
    y += txt_h + Inches(0.2)

    add_rect_bg(slide, caja_left, y, caja_w, Inches(0.38),
                RGBColor(0xF5, 0xF8, 0xFC), border_color=C_GRIS_LINEA)
    add_textbox(slide, caja_left + Inches(0.15), y + Inches(0.07),
                caja_w - Inches(0.3), Inches(0.28),
                text="Caveats 1 y 2 (velocidad e inverter): identicos en V1, V2 y V3. Ver slide anterior.",
                font_size=9.5, italic=True, color=C_GRIS_MEDIO, font_name="Calibri")


def slide_argumentos(prs):
    """Argumentos por variante — 3 columnas V1 / V2 / V3."""
    slide = add_slide(prs)
    title_bar(slide,
              "Argumentos por variante",
              "¿Por que elegir cada una? -- V3 propuesto por Jesus Maria (Junta Directiva)")
    footer_bar(slide)

    col_w = Inches(4.1)
    col1_x = Inches(0.22)
    col2_x = Inches(4.57)
    col3_x = Inches(8.92)
    top_hdr = Inches(1.2)

    add_rect_bg(slide, col1_x, top_hdr, col_w, Inches(0.48), C_GRIS_V1)
    add_textbox(slide, col1_x + Inches(0.1), top_hdr + Inches(0.08),
                col_w - Inches(0.15), Inches(0.35),
                text="VARIANTE 1 -- B-sin-NTC: a favor",
                font_size=11, bold=True, color=C_BLANCO, font_name="Calibri")

    add_rect_bg(slide, col2_x, top_hdr, col_w, Inches(0.48), C_AZUL_V2)
    add_textbox(slide, col2_x + Inches(0.1), top_hdr + Inches(0.08),
                col_w - Inches(0.15), Inches(0.35),
                text="VARIANTE 2 -- B-con-NTC: a favor",
                font_size=11, bold=True, color=C_BLANCO, font_name="Calibri")

    add_rect_bg(slide, col3_x, top_hdr, col_w, Inches(0.48), C_SIENA_V3)
    add_textbox(slide, col3_x + Inches(0.1), top_hdr + Inches(0.08),
                col_w - Inches(0.15), Inches(0.35),
                text="VARIANTE 3 -- Beneficio del beneficio: a favor",
                font_size=11, bold=True, color=C_BLANCO, font_name="Calibri")

    args_v1 = [
        ("Preserva la posicion original de Canudas",
         "Comunica la funcion sin exponer NTC. Territorio vacante confirmado; ningun competidor venezolano usa esta formulacion."),
        ("Sin condicion marcaria nueva",
         "Unico pendiente: Vera P-5 -- ya abierto desde v2.1. Sin semana adicional de abogado marcario."),
        ("Cierre mas rapido",
         "Arte directo contra textos vigentes. Pocos dias a instruccion de diseno grafico."),
        ("Corrige el malentendido del 40%",
         "\"Auto\" comunica que el protector se protege a si mismo -- el prefijo actua en el tiro antes de que el comprador lea el retiro."),
    ]

    args_v2 = [
        ("Responde al pedido literal de Canudas",
         "\"Hagan un ejemplo con las siglas NTC.\" Esta variante es ese ejemplo. Presentar solo V1 seria no responder el brief."),
        ("Activa la tesis anti-copia",
         "NTC visible permite bautizar la funcion con nombre propio. Competidor que replique la placa NTC no puede imprimir \"Escudo Termico NTC\" sin dispute marcario."),
        ("Mayor impacto estrategico",
         "Score 6.45/10 en evaluacion del universo de candidatos -- el mas alto entre opciones con NTC."),
        ("Aislada de la condicion Vera P-5",
         "No usa el adjetivo \"activa\" -- si P-5 resulta desfavorable para V1, V2 no necesita ninguna modificacion."),
    ]

    args_v3 = [
        ("Conexion emocional mas profunda (Jesus Maria)",
         "No comunica la funcion tecnica -- comunica la consecuencia ultima: proteccion si el breaker falla o esta mal seleccionado. Loss aversion activado."),
        ("Reputacion del instalador como vector",
         "El tecnico instalador se identifica con el claim: prescribir un protector con este badge refuerza su reputacion ante el cliente final."),
        ("Tres registros tonales disponibles",
         "V3a (coloquial), V3b (formal) y V3c (emocional) permiten ajustar el tono sin cambiar el eje conceptual. La Junta elige registro."),
        ("Sin riesgo marcario nuevo",
         "Ningun sub-claim de V3 requiere verificacion SAPI VE para uso como claim funcional. Sin semana adicional."),
    ]

    row_h = Inches(1.35)
    y = top_hdr + Inches(0.55)
    for v1_arg, v2_arg, v3_arg in zip(args_v1, args_v2, args_v3):
        for col_x, arg, bg_clr, brd_clr, txt_clr in [
            (col1_x, v1_arg, C_GRIS_V1_CLR, C_GRIS_V1, C_GRIS_V1),
            (col2_x, v2_arg, C_AZUL_V2_CLR, C_AZUL_V2, C_AZUL_V2),
            (col3_x, v3_arg, C_SIENA_V3_CLR, C_SIENA_V3, C_SIENA_V3),
        ]:
            add_rect_bg(slide, col_x, y, col_w, row_h, bg_clr, border_color=brd_clr)
            add_textbox(slide, col_x + Inches(0.1), y + Inches(0.06),
                        col_w - Inches(0.18), Inches(0.34),
                        text=arg[0], font_size=9.5, bold=True,
                        color=txt_clr, font_name="Calibri")
            add_textbox(slide, col_x + Inches(0.1), y + Inches(0.38),
                        col_w - Inches(0.18), Inches(0.9),
                        text=arg[1], font_size=9, color=C_GRIS_OSCURO,
                        font_name="Calibri")
        y += row_h + Inches(0.04)


def slide_riesgos(prs):
    """Riesgos y condiciones — 3 columnas V1 / V2 / V3."""
    slide = add_slide(prs)
    title_bar(slide,
              "Riesgos y condiciones pre-imprenta",
              "Que debe resolverse antes de que el arte entre a produccion")
    footer_bar(slide)

    col_w = Inches(4.1)
    col1_x = Inches(0.22)
    col2_x = Inches(4.57)
    col3_x = Inches(8.92)
    top_hdr = Inches(1.2)

    for col_x, label, bg in [
        (col1_x, "VARIANTE 1 -- B-sin-NTC: riesgos y condiciones", C_GRIS_V1),
        (col2_x, "VARIANTE 2 -- B-con-NTC: riesgos y condiciones", C_AZUL_V2),
        (col3_x, "VARIANTE 3 -- Beneficio del beneficio: riesgos y condiciones", C_SIENA_V3),
    ]:
        add_rect_bg(slide, col_x, top_hdr, col_w, Inches(0.45), bg)
        add_textbox(slide, col_x + Inches(0.1), top_hdr + Inches(0.07),
                    col_w - Inches(0.15), Inches(0.35),
                    text=label, font_size=10, bold=True, color=C_BLANCO, font_name="Calibri")

    riesgos_v1 = [
        (Inches(0.62), "Sin riesgo marcario nuevo",
         "Ningun riesgo nuevo respecto a la version anterior. La formulacion no requiere verificacion SAPI VE."),
        (Inches(1.08), "Condicion tecnica Vera P-5",
         "I&D Genteca debe confirmar si el NTC opera con curva continua o como fusible de corte unico. "
         "Si es fusible unico: badge pasa a \"Autoproteccion termica*\" (formulacion de respaldo aprobada). "
         "No bloquea la Junta; si bloquea imprenta."),
        (Inches(0.62), "Datasheet I&D < 30 ms",
         "Condicion de produccion vigente desde v2 -- identica para las tres variantes. Sin cambio."),
    ]

    riesgos_v2 = [
        (Inches(0.82), "Verificacion SAPI VE obligatoria antes de imprenta",
         "\"Escudo Termico\" -- riesgo medio-alto de objecion por descriptividad. Requiere abogado marcario. Estimado: una semana adicional."),
        (Inches(0.78), "Simbolo TM suspendido",
         "No imprimir TM hasta SAPI VE favorable. El termino puede usarse como claim funcional; "
         "sin garantia de proteccion marcaria hasta que el proceso concluya."),
        (Inches(0.78), "Postura ante riesgo SAPI: decision post-Junta",
         "Si SAPI objeta: funciona como claim funcional sin TM. Si SAPI aprueba: TM puede incorporarse antes de imprenta."),
        (Inches(0.62), "Datasheet I&D < 30 ms",
         "Condicion de produccion vigente desde v2 -- identica para las tres variantes. Sin cambio."),
    ]

    riesgos_v3 = [
        (Inches(0.68), "V3a: confirmar coloquialismo con audiencia",
         "\"Breaker\" es reconocido en VE pero es anglicismo informal. Confirmar con muestra de instaladores tecnicos antes de imprenta."),
        (Inches(0.68), "V3b: espacio de diseno",
         "\"Respaldo termico ante fallas del termomagnetico\" es la formulacion mas larga de las tres. Verificar espacio disponible en layout con disenador grafico."),
        (Inches(0.88), "V3c: subtexto obligatorio en tiro fisico",
         "\"Ultima linea de defensa electrica\" requiere subtexto adyacente \"De tu instalacion electrica\" en el tiro. "
         "Sin el subtexto, el claim puede leerse como defensa general, no de la instalacion especifica."),
        (Inches(0.62), "Datasheet I&D < 30 ms",
         "Condicion de produccion vigente desde v2 -- identica para las tres variantes. Sin cambio."),
    ]

    def render_col(col_x, riesgos, bg_clr, brd_clr, txt_clr):
        y = top_hdr + Inches(0.52)
        for rh, titulo, texto in riesgos:
            add_rect_bg(slide, col_x, y, col_w, rh, bg_clr, border_color=brd_clr)
            add_textbox(slide, col_x + Inches(0.1), y + Inches(0.05),
                        col_w - Inches(0.18), Inches(0.26),
                        text=titulo, font_size=9.5, bold=True,
                        color=txt_clr, font_name="Calibri")
            add_textbox(slide, col_x + Inches(0.1), y + Inches(0.3),
                        col_w - Inches(0.18), rh - Inches(0.35),
                        text=texto, font_size=9, color=C_GRIS_OSCURO,
                        font_name="Calibri")
            y += rh + Inches(0.05)

    render_col(col1_x, riesgos_v1, C_GRIS_V1_CLR, C_GRIS_V1, C_GRIS_V1)
    render_col(col2_x, riesgos_v2, C_AZUL_V2_CLR, C_AZUL_V2, C_AZUL_V2)
    render_col(col3_x, riesgos_v3, C_SIENA_V3_CLR, C_SIENA_V3, C_SIENA_V3)


def slide_recomendacion(prs):
    """Posicion del equipo — 3 ejes evaluados."""
    slide = add_slide(prs)
    title_bar(slide,
              "Posicion del equipo (sin imponer variante)",
              "Las tres variantes son ejecutables. El equipo presenta las tres para que la Junta decida.")
    footer_bar(slide)

    voces = [
        {
            "rol": "Analisis de riesgo",
            "posicion": (
                "V1: menor carga de condiciones nuevas -- Vera P-5 es el unico pendiente, ya abierto. "
                "V2: anade SAPI VE (una semana) y TM suspendido. "
                "V3: sin condicion marcaria nueva; V3a requiere confirmacion de audiencia, V3b verificacion de espacio, V3c subtexto en tiro. "
                "-> V1 o V3 si el criterio es menor friccion de produccion; V2 si es la tesis anti-copia."
            ),
            "acento": C_AMBAR,
        },
        {
            "rol": "Estrategia de marca y mensaje",
            "posicion": (
                "V2: mayor impacto estrategico si Genteca quiere establecer NTC como termino de categoria. "
                "V3: mayor conexion emocional con el consumidor y reputacional con el instalador -- "
                "eje no explotado por ningun competidor venezolano. "
                "V1: correcto y diferenciador, pero no abre ni el camino NTC ni el eje emocional."
            ),
            "acento": C_AZUL_V2,
        },
        {
            "rol": "Verificacion competitiva",
            "posicion": (
                "Vacancia territorial confirmada para las tres variantes en el mercado venezolano. "
                "Riesgo de descriptividad de \"Escudo Termico\" (V2): medio-alto en SAPI VE. "
                "Los claims de V3 son funcionales y no descriptivos del componente -- perfil marcario mas limpio que V2. "
                "-> SAPI VE recomendada antes de imprenta si gana V2."
            ),
            "acento": C_GRIS_V1,
        },
        {
            "rol": "Produccion de contenido y diseno",
            "posicion": (
                "Las tres variantes tienen copy completo listo para instruccion. "
                "V3c requiere ajuste de layout de tiro por subtexto obligatorio -- delta menor con el disenador grafico. "
                "Sin preferencia editorial entre las tres."
            ),
            "acento": C_SIENA_V3,
        },
    ]

    y = Inches(1.25)
    for v in voces:
        rh = Inches(1.32)
        add_rect_bg(slide, Inches(0.3), y, Inches(0.06), rh, v["acento"])
        add_rect_bg(slide, Inches(0.4), y, Inches(12.5), rh,
                    RGBColor(0xF8, 0xF9, 0xFA), border_color=C_GRIS_LINEA)
        add_textbox(slide, Inches(0.55), y + Inches(0.07),
                    Inches(12.1), Inches(0.28),
                    text=v["rol"].upper(), font_size=9.5, bold=True,
                    color=v["acento"], font_name="Calibri")
        add_textbox(slide, Inches(0.55), y + Inches(0.34),
                    Inches(12.1), Inches(0.92),
                    text=v["posicion"], font_size=10, color=C_GRIS_OSCURO,
                    font_name="Calibri")
        y += rh + Inches(0.05)


def slide_proximos_pasos(prs):
    """Proximos pasos — 3 escenarios principales + sub-escenario V3."""
    slide = add_slide(prs)
    title_bar(slide,
              "Proximos pasos segun la decision",
              "La instruccion al diseno grafico es inmediata en los tres escenarios principales.")
    footer_bar(slide)

    col_w = Inches(4.1)
    col1_x = Inches(0.22)
    col2_x = Inches(4.57)
    col3_x = Inches(8.92)
    top_hdr = Inches(1.2)

    for col_x, label, bg in [
        (col1_x, "Si gana Variante 1 -- B-sin-NTC", C_GRIS_V1),
        (col2_x, "Si gana Variante 2 -- B-con-NTC", C_AZUL_V2),
        (col3_x, "Si gana Variante 3 -- Beneficio del beneficio", C_SIENA_V3),
    ]:
        add_rect_bg(slide, col_x, top_hdr, col_w, Inches(0.55), bg)
        add_textbox(slide, col_x + Inches(0.1), top_hdr + Inches(0.1),
                    col_w - Inches(0.15), Inches(0.42),
                    text=label, font_size=11, bold=True, color=C_BLANCO, font_name="Calibri")

    pasos_v1 = [
        ("1", "Instruccion inmediata al diseno grafico",
         "Contra los textos aprobados vigentes. Arte puede comenzar."),
        ("2", "Vera P-5 en paralelo",
         "I&D confirma mecanismo NTC. Si es fusible unico: badge pasa a \"Autoproteccion termica*\". El cuerpo del caveat no cambia."),
        ("3", "Datasheet I&D < 30 ms",
         "Condicion de imprenta vigente. Sin cambio."),
        ("CIERRE", "Estimado: pocos dias",
         "Arte final -> imprenta, sujeto a Vera P-5 + datasheet."),
    ]

    pasos_v2 = [
        ("1", "Instruccion al diseno grafico con textos V2",
         "Arte puede comenzar sin TM."),
        ("2", "Verificacion SAPI VE en paralelo (obligatoria antes de imprenta)",
         "Abogado marcario -- \"Escudo Termico\" clase 9. Estimado: una semana."),
        ("3", "Datasheet I&D < 30 ms",
         "Condicion de imprenta vigente. Sin cambio."),
        ("CIERRE", "Estimado: aprox. una semana",
         "Arte final -> imprenta, sujeto a SAPI VE + datasheet."),
    ]

    pasos_v3 = [
        ("1", "Sub-decision: V3a, V3b o V3c",
         "La Junta elige el registro tonal. Instruccion al disenador con la sub-opcion elegida."),
        ("2", "Condicion especifica de la sub-opcion",
         "V3a: confirmar coloquialismo con instaladores. V3b: verificar espacio de layout. V3c: ajuste de tiro con subtexto obligatorio."),
        ("3", "Datasheet I&D < 30 ms",
         "Condicion de imprenta vigente. Sin cambio."),
        ("CIERRE", "Estimado: pocos dias",
         "Arte final -> imprenta, sujeto a condicion de sub-opcion + datasheet."),
    ]

    row_heights = [Inches(0.78), Inches(1.05), Inches(0.68), Inches(0.72)]

    def render_pasos(col_x, pasos, color_var):
        y = top_hdr + Inches(0.62)
        for i, (num_label, titulo, texto) in enumerate(pasos):
            rh = row_heights[i]
            is_cierre = (i == 3)
            bg = C_AZUL_OSCURO if is_cierre else RGBColor(0xF8, 0xF9, 0xFA)
            num_bg = C_AMBAR if is_cierre else color_var
            fg_tit = C_AMBAR if is_cierre else color_var
            fg_txt = C_GRIS_LINEA if is_cierre else C_GRIS_OSCURO

            add_rect_bg(slide, col_x, y, col_w, rh, bg, border_color=C_GRIS_LINEA)
            add_rect_bg(slide, col_x, y, Inches(0.38), rh, num_bg)
            add_textbox(slide, col_x + Inches(0.02), y + Inches(0.08),
                        Inches(0.35), Inches(0.36),
                        text=num_label, font_size=9, bold=True, color=C_BLANCO,
                        align=PP_ALIGN.CENTER, font_name="Calibri")
            add_textbox(slide, col_x + Inches(0.44), y + Inches(0.06),
                        col_w - Inches(0.52), Inches(0.28),
                        text=titulo, font_size=9.5, bold=True, color=fg_tit,
                        font_name="Calibri")
            add_textbox(slide, col_x + Inches(0.44), y + Inches(0.32),
                        col_w - Inches(0.52), rh - Inches(0.37),
                        text=texto, font_size=9, color=fg_txt, font_name="Calibri")
            y += rh + Inches(0.04)

    render_pasos(col1_x, pasos_v1, C_GRIS_V1)
    render_pasos(col2_x, pasos_v2, C_AZUL_V2)
    render_pasos(col3_x, pasos_v3, C_SIENA_V3)


def slide_cierre(prs):
    """Cierre — dos niveles de pregunta a la Junta."""
    slide = add_slide(prs)

    add_rect_bg(slide, Inches(0), Inches(0), W, H, C_AZUL_OSCURO)
    add_rect_bg(slide, Inches(0), Inches(3.5), W, Inches(0.06), C_AMBAR)

    add_textbox(slide, Inches(0.5), Inches(0.5), Inches(12.3), Inches(0.55),
                text="GENTECA | JUNTA DIRECTIVA | 2026-05-06",
                font_size=12, color=C_GRIS_LINEA, font_name="Calibri Light")

    add_textbox(slide, Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.9),
                text="La decision de hoy",
                font_size=30, bold=True, color=C_AMBAR, font_name="Calibri Light")

    # Pregunta nivel 1
    add_rect_bg(slide, Inches(0.5), Inches(1.98), Inches(12.3), Inches(0.68),
                RGBColor(0x0D, 0x1E, 0x33))
    add_textbox(slide, Inches(0.7), Inches(2.06), Inches(11.9), Inches(0.55),
                text="Nivel 1 -- ¿Que variante avanza a arte final?  V1 (B-sin-NTC)  |  V2 (B-con-NTC)  |  V3 (Beneficio del beneficio)",
                font_size=16, bold=True, color=C_BLANCO, font_name="Calibri")

    # Pregunta nivel 2
    add_rect_bg(slide, Inches(0.5), Inches(2.74), Inches(12.3), Inches(0.6),
                RGBColor(0x1A, 0x1A, 0x2E))
    add_textbox(slide, Inches(0.7), Inches(2.82), Inches(11.9), Inches(0.48),
                text="Nivel 2 (solo si V3) -- ¿Que sub-opcion?  V3a (tecnica-coloquial)  |  V3b (tecnica-formal)  |  V3c (emocional)",
                font_size=13, bold=True, color=C_SIENA_V3_CLR, font_name="Calibri")

    add_textbox(slide, Inches(0.5), Inches(3.65), Inches(12.3), Inches(0.6),
                text="¿Hay algun angulo o sensibilidad no contemplado que deba incorporarse\n"
                     "antes de la instruccion al diseno grafico?",
                font_size=14, color=C_AZUL_CLARO, font_name="Calibri Light")

    add_textbox(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.35),
                text="Lo que NO decide la Junta hoy:",
                font_size=11, bold=True, color=C_GRIS_MEDIO, font_name="Calibri")
    no_decide = [
        "Vera P-5 (mecanismo NTC): tarea tecnica paralela de I&D Genteca",
        "Verificacion SAPI VE: aplica solo si gana V2; no requerida para V1 ni V3 como claim funcional",
        "Datasheet I&D < 30 ms: condicion de imprenta vigente, sin cambio",
        "Sub-condicion V3c (subtexto en tiro): ajuste de layout con disenador grafico post-decision",
    ]
    for i, item in enumerate(no_decide):
        add_textbox(slide, Inches(0.7), Inches(4.88 + i * 0.36), Inches(12.0), Inches(0.34),
                    text="-- " + item, font_size=10, color=C_GRIS_LINEA,
                    font_name="Calibri")

    footer_bar(slide, "Genteca | Junta Directiva | 2026-05-06 | Confidencial")


# ---------------------------------------------------------------------------
# SLIDES — ANEXO (A1-A5) — sin modificacion respecto a version anterior
# ---------------------------------------------------------------------------

def slide_anexo_apertura(prs):
    slide = add_slide(prs)
    add_rect_bg(slide, Inches(0), Inches(0), W, H, C_ANEXO_FONDO)
    title_bar_anexo(
        slide,
        "Bautizado en anglicismos como estrategia de IP defensiva",
        "Analisis exploratorio para discusion posterior"
    )
    footer_bar_anexo(slide, "A1 de 5")
    add_rect_bg(slide, Inches(0), Inches(3.2), W, Inches(0.05), C_ANEXO_ACENTO)
    add_rect_bg(slide, Inches(0.5), Inches(1.35), Inches(12.3), Inches(1.65),
                C_BLANCO, border_color=C_ANEXO_ACENTO)
    add_textbox(slide, Inches(0.7), Inches(1.5), Inches(11.9), Inches(1.4),
                text=(
                    "Las opciones presentadas en el cuerpo principal no incluyen este enfoque. "
                    "Esta seccion es un abreboca a discusiones futuras sobre como blindar legalmente "
                    "cada innovacion tecnica de Genteca."
                ),
                font_size=15, color=C_GRIS_OSCURO, font_name="Calibri Light")
    props = [
        ("Registro mas probable",
         "SAPI VE registra con mayor probabilidad palabras inventadas o anglicismos no descriptivos "
         "que descriptores comunes en espanol."),
        ("Barrera competitiva real",
         "Si Genteca registra la marca en anglicismo, la competencia que copie la placa nunca podra "
         "usar ese nombre en su empaque -- aunque tenga el componente."),
        ("Posicionamiento premium",
         "Los anglicismos transmiten estandar internacional. Refuerzan el posicionamiento "
         "Exceline Profesional en el canal ferretero VE."),
    ]
    col_w = Inches(3.8)
    col_xs = [Inches(0.5), Inches(4.65), Inches(8.8)]
    for col_x, prop in zip(col_xs, props):
        add_rect_bg(slide, col_x, Inches(3.45), col_w, Inches(2.9),
                    C_BLANCO, border_color=C_GRIS_LINEA)
        add_rect_bg(slide, col_x, Inches(3.45), col_w, Inches(0.06), C_ANEXO_ACENTO)
        add_textbox(slide, col_x + Inches(0.15), Inches(3.6), col_w - Inches(0.3), Inches(0.42),
                    text=prop[0], font_size=12, bold=True, color=C_ANEXO_ACENTO,
                    font_name="Calibri")
        add_textbox(slide, col_x + Inches(0.15), Inches(4.08), col_w - Inches(0.3), Inches(2.1),
                    text=prop[1], font_size=10.5, color=C_GRIS_OSCURO, font_name="Calibri")


def slide_anexo_por_que(prs):
    slide = add_slide(prs)
    add_rect_bg(slide, Inches(0), Inches(0), W, H, C_ANEXO_FONDO)
    title_bar_anexo(
        slide,
        "Por que registrar en anglicismos",
        "Tres razones estructurales -- no es estetica, es estrategia de IP"
    )
    footer_bar_anexo(slide, "A2 de 5")
    razones = [
        {
            "num": "1",
            "titulo": "Mayor probabilidad de registro exitoso en SAPI VE",
            "texto": (
                "El SAPI Venezuela registra con mayor facilidad palabras inventadas o anglicismos no descriptivos "
                "que descriptores comunes en espanol. \"Escudo Termico\" es descriptivo puro -- riesgo medio-alto "
                "de objecion. \"Thermo-Safe\" o \"Thermal Shield\" no describen el mecanismo del producto: "
                "perfil mas limpio para registro formal."
            ),
        },
        {
            "num": "2",
            "titulo": "La marca registrada en anglicismo blinda a Genteca frente a copias",
            "texto": (
                "Si un competidor replica la placa NTC -- situacion hipotetica pero plausible dado el costo bajo "
                "del componente --, jamas podra usar el nombre registrado de Genteca en su empaque. "
                "El producto puede ser fisicamente identico; el nombre no puede copiarse sin exposicion legal. "
                "Esto convierte el bautizado en una barrera competitiva que el hardware solo no puede proveer."
            ),
        },
        {
            "num": "3",
            "titulo": "Sofisticacion percibida: anglicismos = estandar internacional",
            "texto": (
                "En el canal ferretero VE, los instaladores tecnicos -- prescriptores del 80% de la decision "
                "de compra -- asocian anglicismos en empaque electrico con calidad tecnica superior. "
                "\"Inverter\", \"digital\", \"smart\" ya circulan sin friccion. \"Thermo-Safe\" sigue ese patron. "
                "Refuerza el posicionamiento Exceline Profesional sin costo adicional."
            ),
        },
    ]
    y = Inches(1.28)
    colors_num = [C_ANEXO_ACENTO, C_AZUL_OSCURO, C_GRIS_OSCURO]
    for i, r in enumerate(razones):
        rh = Inches(1.68)
        add_rect_bg(slide, Inches(0.35), y, Inches(12.6), rh, C_BLANCO, border_color=C_GRIS_LINEA)
        add_rect_bg(slide, Inches(0.35), y, Inches(0.55), rh, colors_num[i])
        add_textbox(slide, Inches(0.35), y + Inches(0.55), Inches(0.55), Inches(0.5),
                    text=r["num"], font_size=20, bold=True, color=C_BLANCO,
                    align=PP_ALIGN.CENTER, font_name="Calibri")
        add_textbox(slide, Inches(1.02), y + Inches(0.1), Inches(11.7), Inches(0.38),
                    text=r["titulo"], font_size=12, bold=True, color=C_ANEXO_HDR, font_name="Calibri")
        add_textbox(slide, Inches(1.02), y + Inches(0.48), Inches(11.7), Inches(1.1),
                    text=r["texto"], font_size=10.5, color=C_GRIS_OSCURO, font_name="Calibri")
        y += rh + Inches(0.07)


def slide_anexo_candidatos(prs):
    slide = add_slide(prs)
    add_rect_bg(slide, Inches(0), Inches(0), W, H, C_ANEXO_FONDO)
    title_bar_anexo(
        slide,
        "Candidatos en anglicismos -- proteccion termica",
        "Territorio competitivo VE -- snapshot Orlan OL-1 §Refresh 2026-05-06"
    )
    footer_bar_anexo(slide, "A3 de 5")
    cols = [Inches(2.8), Inches(3.2), Inches(3.2), Inches(2.8)]
    headers = ["Candidato", "Vacancia VE (clase 9)", "Conflicto detectado", "Perfil de riesgo IP"]
    col_xs = [Inches(0.4), Inches(3.25), Inches(6.5), Inches(9.75)]
    y_hdr = Inches(1.22)
    for i, (hdr, cw, cx) in enumerate(zip(headers, cols, col_xs)):
        add_rect_bg(slide, cx, y_hdr, cw, Inches(0.44), C_ANEXO_HDR)
        add_textbox(slide, cx + Inches(0.1), y_hdr + Inches(0.08), cw - Inches(0.2), Inches(0.32),
                    text=hdr, font_size=10, bold=True, color=C_BLANCO, font_name="Calibri")
    candidatos = [
        ("Thermo-Safe",
         "Vacante -- ningun uso en protectores electricos VE",
         "Sin conflicto detectado (preliminar)",
         "Bajo -- perfil mas limpio del grupo"),
        ("Thermal Shield",
         "Vacante -- ningun uso en protectores electricos VE",
         "Sin conflicto detectado en clase 9 (preliminar)",
         "Bajo -- uso en otros rubros (embalajes, tejados), no electrico"),
        ("Thermo-Shield",
         "Vacante -- ningun uso en protectores electricos VE",
         "Sin conflicto detectado en clase 9 (preliminar)",
         "Bajo -- situacion similar a Thermal Shield"),
        ("Escudo Termico",
         "Vacante en protectores electricos VE",
         "Sin conflicto en clase 9 -- pero uso extenso en espanol como generico",
         "Medio-alto -- descriptor puro en espanol, riesgo de objecion SAPI por descriptividad"),
    ]
    y = y_hdr + Inches(0.44)
    for j, (nombre, vacancia, conflicto, riesgo) in enumerate(candidatos):
        rh = Inches(0.92)
        bg = C_BLANCO if j % 2 == 0 else RGBColor(0xF7, 0xF5, 0xF0)
        for i, (texto, cw, cx) in enumerate(zip([nombre, vacancia, conflicto, riesgo], cols, col_xs)):
            add_rect_bg(slide, cx, y, cw, rh, bg, border_color=C_GRIS_LINEA)
            is_name = (i == 0)
            add_textbox(slide, cx + Inches(0.1), y + Inches(0.1),
                        cw - Inches(0.2), rh - Inches(0.15),
                        text=texto, font_size=10, bold=is_name,
                        color=C_ANEXO_ACENTO if is_name else C_GRIS_OSCURO,
                        font_name="Calibri")
        y += rh
    y += Inches(0.12)
    add_rect_bg(slide, Inches(0.4), y, Inches(12.5), Inches(0.5),
                RGBColor(0xFF, 0xF8, 0xE8), border_color=C_ANEXO_ACENTO)
    add_textbox(slide, Inches(0.55), y + Inches(0.08), Inches(12.2), Inches(0.38),
                text=(
                    "Nota obligatoria: Cualquiera de estas opciones requiere verificacion SAPI VE con abogado marcario "
                    "antes de imprenta. Los datos de Orlan son orientativos, no sustituto de consulta registral formal."
                ),
                font_size=9.5, italic=True, color=RGBColor(0x6B, 0x4E, 0x10), font_name="Calibri")


def slide_anexo_expansivo(prs):
    slide = add_slide(prs)
    add_rect_bg(slide, Inches(0), Inches(0), W, H, C_ANEXO_FONDO)
    title_bar_anexo(
        slide,
        "Aplicacion expansiva: cada innovacion es un activo comercial",
        "El patron de bautizado puede aplicarse a toda la cartera de IP de Genteca"
    )
    footer_bar_anexo(slide, "A4 de 5")
    add_rect_bg(slide, Inches(0.4), Inches(1.25), Inches(12.5), Inches(0.72),
                C_BLANCO, border_color=C_ANEXO_ACENTO)
    add_rect_bg(slide, Inches(0.4), Inches(1.25), Inches(0.08), Inches(0.72), C_ANEXO_ACENTO)
    add_textbox(slide, Inches(0.6), Inches(1.35), Inches(12.1), Inches(0.55),
                text=(
                    "Si Genteca adopta el habito de bautizar cada innovacion tecnica con nombre de marca propio, "
                    "cada avance de I&D se convierte en un activo comercial blindado -- no solo una caracteristica del producto."
                ),
                font_size=12, color=C_GRIS_OSCURO, font_name="Calibri")
    ejemplos = [
        {
            "innovacion": "Proteccion termica (NTC)",
            "aplicacion": "Thermo-Safe / Thermal Shield / Thermo-Shield",
            "efecto": "Competidor con NTC no puede usar el nombre registrado en su empaque.",
            "estado": "En evaluacion hoy",
        },
        {
            "innovacion": "Curva inversa de voltaje",
            "aplicacion": "Nombre de marca para el algoritmo de disparo (ej. VoltCurve)",
            "efecto": "Diferencia a Genteca como empresa de software de proteccion, no solo hardware.",
            "estado": "Exploratorio",
        },
        {
            "innovacion": "Algoritmo del microcontrolador",
            "aplicacion": "Nombre de firmware / motor de decision (ej. ProGuard Engine)",
            "efecto": "El 'cerebro' del protector se vuelve activo con nombre y barrera de marca.",
            "estado": "Exploratorio",
        },
        {
            "innovacion": "Software de control / UI",
            "aplicacion": "Nombre comercial propietario para la plataforma de gestion",
            "efecto": "Cada version del software puede tener numero de version y nombre de marca.",
            "estado": "Exploratorio",
        },
    ]
    y = Inches(2.12)
    col_xs = [Inches(0.4), Inches(3.5), Inches(7.3), Inches(10.5)]
    col_ws = [Inches(3.05), Inches(3.75), Inches(3.15), Inches(2.4)]
    col_hdrs = ["Innovacion", "Aplicacion de bautizado", "Efecto de barrera", "Estado"]
    for cx, cw, hdr in zip(col_xs, col_ws, col_hdrs):
        add_rect_bg(slide, cx, y, cw, Inches(0.36), C_ANEXO_HDR)
        add_textbox(slide, cx + Inches(0.08), y + Inches(0.06), cw - Inches(0.15), Inches(0.26),
                    text=hdr, font_size=9.5, bold=True, color=C_BLANCO, font_name="Calibri")
    y += Inches(0.36)
    for k, ej in enumerate(ejemplos):
        rh = Inches(0.96)
        bg = C_BLANCO if k % 2 == 0 else RGBColor(0xF7, 0xF5, 0xF0)
        valores = [ej["innovacion"], ej["aplicacion"], ej["efecto"], ej["estado"]]
        for cx, cw, val in zip(col_xs, col_ws, valores):
            add_rect_bg(slide, cx, y, cw, rh, bg, border_color=C_GRIS_LINEA)
            is_estado = (val == ej["estado"])
            fc = C_ANEXO_ACENTO if ej["estado"] == "En evaluacion hoy" and is_estado else C_GRIS_OSCURO
            add_textbox(slide, cx + Inches(0.08), y + Inches(0.08),
                        cw - Inches(0.15), rh - Inches(0.12),
                        text=val, font_size=9.5,
                        bold=(is_estado and ej["estado"] == "En evaluacion hoy"),
                        color=fc, font_name="Calibri")
        y += rh
    y += Inches(0.1)
    add_textbox(slide, Inches(0.4), y, Inches(12.5), Inches(0.32),
                text="Nota: los nombres de ejemplo (VoltCurve, ProGuard Engine) son ilustrativos. "
                     "Este anexo no propone nombres especificos -- establece el patron.",
                font_size=9, italic=True, color=C_GRIS_MEDIO, font_name="Calibri")


def slide_anexo_pregunta(prs):
    slide = add_slide(prs)
    add_rect_bg(slide, Inches(0), Inches(0), W, H, C_ANEXO_HDR)
    add_rect_bg(slide, Inches(0), Inches(3.6), W, Inches(0.06), C_ANEXO_ACENTO)
    footer_bar_anexo(slide, "A5 de 5")
    add_rect_bg(slide, Inches(0.5), Inches(0.5), Inches(1.6), Inches(0.38), C_ANEXO_ACENTO)
    add_textbox(slide, Inches(0.55), Inches(0.55), Inches(1.5), Inches(0.3),
                text="ANEXO", font_size=11, bold=True, color=C_BLANCO,
                align=PP_ALIGN.CENTER, font_name="Calibri")
    add_textbox(slide, Inches(0.5), Inches(1.08), Inches(12.3), Inches(0.75),
                text="Preguntas abiertas para la Junta",
                font_size=28, bold=True, color=C_ANEXO_ACENTO, font_name="Calibri Light")
    add_textbox(slide, Inches(0.5), Inches(1.82), Inches(12.3), Inches(0.48),
                text="Estas preguntas no demandan decision hoy. Son semilla para conversacion posterior.",
                font_size=13, italic=True, color=C_GRIS_LINEA, font_name="Calibri Light")
    preguntas = [
        ("¿Vale la pena explorar este enfoque sistematico de bautizado de innovaciones?",
         "El cuerpo principal de hoy ya abre esa puerta con el badge termico. "
         "La pregunta es si Genteca lo adopta como habito sistematico para toda la cartera de IP."),
        ("¿Para que innovaciones de Genteca tendria mas sentido empezar?",
         "La curva inversa de voltaje, el algoritmo del microcontrolador, el software de control: "
         "cada uno es candidato a bautizado. ¿Cual tiene mayor prioridad comercial?"),
        ("¿Que presupuesto comercial-legal considerar para implementarlo?",
         "Cada bautizado formal requiere consulta con abogado marcario y registro SAPI VE. "
         "La economia de escala favorece registrar varios en un mismo proceso."),
    ]
    y = Inches(2.5)
    for i, (pregunta, contexto) in enumerate(preguntas):
        rh = Inches(1.42)
        add_rect_bg(slide, Inches(0.4), y, Inches(0.06), rh, C_ANEXO_ACENTO)
        add_rect_bg(slide, Inches(0.5), y, Inches(12.4), rh,
                    RGBColor(0x1A, 0x1A, 0x1A), border_color=RGBColor(0x3A, 0x3A, 0x3A))
        add_textbox(slide, Inches(0.68), y + Inches(0.1), Inches(12.1), Inches(0.42),
                    text=pregunta, font_size=13, bold=True, color=C_BLANCO, font_name="Calibri")
        add_textbox(slide, Inches(0.68), y + Inches(0.54), Inches(12.1), Inches(0.82),
                    text=contexto, font_size=10.5, color=C_GRIS_LINEA, font_name="Calibri")
        y += rh + Inches(0.06)


# ---------------------------------------------------------------------------
# MAIN
# ---------------------------------------------------------------------------

def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    # --- CUERPO PRINCIPAL ---
    # Slide 1 — Portada
    slide_portada(prs)

    # Slide 2 — Contexto y pregunta
    slide_contexto(prs)

    # Slide 3 — Comparacion lado a lado (3 columnas)
    slide_comparacion(prs)

    # Slide 4 — Variante 1: tiro completo
    slide_tiro(prs, "V1", TIRO_CLAIM2_V1, C_GRIS_V1, C_GRIS_V1_CLR)

    # Slide 5 — Variante 2: tiro completo
    slide_tiro(prs, "V2", TIRO_CLAIM2_V2, C_AZUL_V2, C_AZUL_V2_CLR)

    # Slide 6 — Variante 3: tiro completo (3 sub-opciones) — NUEVA
    slide_tiro_v3(prs)

    # Slide 7 — Variante 1: retiro CARACTERISTICAS
    slide_retiro_caract(prs, "V1", CARACT_TERMICO_V1, C_GRIS_V1, C_GRIS_V1_CLR)

    # Slide 8 — Variante 2: retiro CARACTERISTICAS
    slide_retiro_caract(prs, "V2", CARACT_TERMICO_V2, C_AZUL_V2, C_AZUL_V2_CLR)

    # Slide 9 — Variante 3: retiro CARACTERISTICAS (3 sub-opciones) — NUEVA
    slide_retiro_caract_v3(prs)

    # Slide 10 — CAVEATS COMPARTIDOS (velocidad + inverter — V1, V2 y V3)
    slide_caveats_compartidos(prs)

    # Slide 11 — Variante 1: caveat termico especifico
    slide_caveat_termico(prs, "V1", CAVEAT_TERMICO_V1, C_GRIS_V1, C_GRIS_V1_CLR)

    # Slide 12 — Variante 2: caveat termico especifico
    slide_caveat_termico(prs, "V2", CAVEAT_TERMICO_V2, C_AZUL_V2, C_AZUL_V2_CLR)

    # Slide 13 — Variante 3: caveat termico (3 headers + cuerpo invariable) — NUEVA
    slide_caveat_termico_v3(prs)

    # Slide 14 — Argumentos por variante (actualizado con V3)
    slide_argumentos(prs)

    # Slide 15 — Riesgos y condiciones (actualizado con V3)
    slide_riesgos(prs)

    # Slide 16 — Posicion del equipo (actualizado con 3 ejes)
    slide_recomendacion(prs)

    # Slide 17 — Proximos pasos (actualizado con 3 escenarios)
    slide_proximos_pasos(prs)

    # Slide 18 — Cierre (actualizado con 2 niveles de pregunta)
    slide_cierre(prs)

    # --- ANEXO ---
    slide_anexo_apertura(prs)   # A1
    slide_anexo_por_que(prs)    # A2
    slide_anexo_candidatos(prs) # A3
    slide_anexo_expansivo(prs)  # A4
    slide_anexo_pregunta(prs)   # A5

    out_path = (
        "C:/RAUL/03-projects/genteca/2026-04_GSM-MB-RB-RF_empaque"
        "/03-review-and-release/Junta_PPTX_v3"
        "/Avance_dual-variantes_GSM-empaque_2026-05-06.pptx"
    )
    prs.save(out_path)
    print(f"Archivo generado: {out_path}")
    return out_path


if __name__ == "__main__":
    build()
