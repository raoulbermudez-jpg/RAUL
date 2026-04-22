"""
gen_gst_r_pptx_v2.py
Generates GST-R_etiquetas_brief_v2.pptx from the source Markdown copy.
python-pptx 1.0.2  |  Python 3.x
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy, os

# ── Output paths ────────────────────────────────────────────────────────────
PATHS = [
    r"C:/WorkspaceIA/PROJECTS/Genteca/Work In Progress/GST-R_etiquetas_brief_v2.pptx",
    r"C:/WorkspaceIA/Owner Inbox/2026-04-20-gst-r-etiquetas-brief-ozwaldo-v2.pptx",
]

# ── Palette ──────────────────────────────────────────────────────────────────
C_DARK_BG    = RGBColor(0x1E, 0x1E, 0x1E)
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_GRAY = RGBColor(0x80, 0x80, 0x80)
C_FOOTER_TXT = RGBColor(0xA0, 0xA0, 0xA0)
C_BLACK_SLIDE= RGBColor(0xFF, 0xFF, 0xFF)  # body text on white slides
C_SECTION_LBL= RGBColor(0x44, 0x44, 0x44)
C_ORANGE     = RGBColor(0xFF, 0x8C, 0x00)
C_GOLD       = RGBColor(0xB8, 0x86, 0x0B)

# Product color map
PROD_COLORS = {
    "RT": RGBColor(0x2E, 0x7D, 0x32),   # verde
    "RD": RGBColor(0x1A, 0x1A, 0x1A),   # negro
    "RM": RGBColor(0x54, 0x6E, 0x7A),   # gris
    "RR": RGBColor(0x15, 0x65, 0xC0),   # azul
}

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ── Helpers ──────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_layout(prs):
    return prs.slide_layouts[6]  # completely blank


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=11,
                bold=False, color=None, align=PP_ALIGN.LEFT,
                font_name="Calibri", wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.name = font_name
    if color:
        run.font.color.rgb = color
    return txBox


def add_footer(slide, text="Brief Etiquetas GST-R v2 | Confidencial Genteca"):
    add_textbox(
        slide,
        Inches(0.2), Inches(7.1),
        Inches(12.9), Inches(0.3),
        text,
        font_size=8, color=C_FOOTER_TXT,
        align=PP_ALIGN.CENTER
    )


def add_multiline_textbox(slide, left, top, width, height, lines,
                          base_font_size=11, base_color=None):
    """lines = list of (text, font_size, bold, color)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for (text, fsize, fbold, fcolor) in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = text
        run.font.size = Pt(fsize)
        run.font.bold = fbold
        run.font.name = "Calibri"
        run.font.color.rgb = fcolor if fcolor else (base_color or C_BLACK_SLIDE)
    return txBox


# ── SLIDE 1 — Cover ───────────────────────────────────────────────────────────

def build_cover(prs):
    slide = prs.slides.add_slide(blank_layout(prs))

    # Dark background
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, C_DARK_BG)

    # Accent bar (top)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.06), RGBColor(0x55,0x55,0x55))

    # Main title
    add_textbox(
        slide,
        Inches(0.7), Inches(0.8),
        Inches(12), Inches(1.4),
        "Brief de Diseño — Etiquetas GST-R Nueva Línea",
        font_size=32, bold=True, color=C_WHITE,
        align=PP_ALIGN.LEFT
    )

    # Subtitle
    add_textbox(
        slide,
        Inches(0.7), Inches(2.1),
        Inches(10), Inches(0.5),
        "Exceline Profesional  |  Genteca",
        font_size=16, bold=False, color=RGBColor(0xCC,0xCC,0xCC),
        align=PP_ALIGN.LEFT
    )

    # Date
    add_textbox(
        slide,
        Inches(0.7), Inches(2.65),
        Inches(6), Inches(0.35),
        "2026-04-20",
        font_size=12, color=RGBColor(0x99,0x99,0x99),
        align=PP_ALIGN.LEFT
    )

    # Font note
    add_textbox(
        slide,
        Inches(0.7), Inches(3.05),
        Inches(11), Inches(0.3),
        "Nota: fuente de producción es Montserrat. Este archivo usa Calibri Bold como sustituto.",
        font_size=9, color=RGBColor(0x88,0x88,0x88), align=PP_ALIGN.LEFT
    )

    # Product cards  (4 colored strips)
    products = [
        ("ProTransfer",  "GST-RT", PROD_COLORS["RT"]),
        ("ProDigital",   "GST-RD", PROD_COLORS["RD"]),
        ("ProMotor",     "GST-RM", PROD_COLORS["RM"]),
        ("ProFrio",      "GST-RR", PROD_COLORS["RR"]),
    ]
    card_w = Inches(2.8)
    card_h = Inches(1.1)
    gap    = Inches(0.22)
    start_x= Inches(0.7)
    card_y = Inches(3.55)

    for i, (name, code, color) in enumerate(products):
        cx = start_x + i * (card_w + gap)
        add_rect(slide, cx, card_y, card_w, card_h, color)
        add_textbox(
            slide, cx + Inches(0.12), card_y + Inches(0.1),
            card_w - Inches(0.2), Inches(0.5),
            name, font_size=15, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT
        )
        add_textbox(
            slide, cx + Inches(0.12), card_y + Inches(0.58),
            card_w - Inches(0.2), Inches(0.35),
            code, font_size=11, bold=False, color=RGBColor(0xDD,0xDD,0xDD),
            align=PP_ALIGN.LEFT
        )

    # Footer
    add_textbox(
        slide,
        Inches(0), Inches(7.05),
        SLIDE_W, Inches(0.35),
        "Confidencial — Uso interno Genteca",
        font_size=9, color=C_FOOTER_TXT, align=PP_ALIGN.CENTER
    )


# ── Reusable: product header band ────────────────────────────────────────────

def add_product_header(slide, product_key, fantasy_name, code_label, zone_text):
    color = PROD_COLORS[product_key]
    # Header band
    band_h = Inches(1.05)
    add_rect(slide, 0, 0, SLIDE_W, band_h, color)

    # Gold accent line for RD
    if product_key == "RD":
        add_rect(slide, 0, band_h - Inches(0.05), SLIDE_W, Inches(0.05), C_GOLD)

    # Fantasy name
    add_textbox(
        slide,
        Inches(0.4), Inches(0.06),
        Inches(7), Inches(0.55),
        fantasy_name,
        font_size=24, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT
    )
    # Code
    add_textbox(
        slide,
        Inches(0.4), Inches(0.58),
        Inches(9), Inches(0.3),
        code_label,
        font_size=14, bold=False, color=RGBColor(0xEE,0xEE,0xEE), align=PP_ALIGN.LEFT
    )
    # Zone
    add_textbox(
        slide,
        Inches(0.4), Inches(0.83),
        Inches(9), Inches(0.22),
        f"Zona: {zone_text}",
        font_size=12, bold=False, color=RGBColor(0xDD,0xDD,0xDD), align=PP_ALIGN.LEFT
    )


# ── Section label helper ──────────────────────────────────────────────────────

def section_label(slide, text, top_y, prod_key):
    color = PROD_COLORS[prod_key]
    add_rect(slide, Inches(0.3), top_y, Inches(3.5), Inches(0.28), color)
    add_textbox(
        slide, Inches(0.35), top_y + Inches(0.02),
        Inches(3.4), Inches(0.26),
        text, font_size=9, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT
    )


def bullet_line(slide, left, top, width, height, text, bullet_color, font_size=11):
    """Single bullet line with colored square bullet."""
    # Bullet square
    sq = Inches(0.09)
    add_rect(slide, left, top + Inches(0.06), sq, sq, bullet_color)
    add_textbox(
        slide,
        left + sq + Inches(0.1), top,
        width - sq - Inches(0.1), height,
        text,
        font_size=font_size, color=RGBColor(0x22,0x22,0x22), align=PP_ALIGN.LEFT
    )


# ══════════════════════════════════════════════════════════════════════════════
#  PRODUCT SLIDES — generic builder
# ══════════════════════════════════════════════════════════════════════════════

def build_product_slides(prs, prod_key, fantasy_name, code_full, zone,
                         applications, features, tech_data, models,
                         sec_b_items, sec_c_resolved, sec_c_pending):
    """Generates 2 slides (A and B+C) for one product."""

    bcolor = PROD_COLORS[prod_key]

    # ── Slide A ───────────────────────────────────────────────────────────────
    slide_a = prs.slides.add_slide(blank_layout(prs))
    add_rect(slide_a, 0, 0, SLIDE_W, SLIDE_H, C_WHITE)
    add_product_header(slide_a, prod_key, fantasy_name, code_full, zone)
    add_footer(slide_a)

    y = Inches(1.2)
    col1_x = Inches(0.35)
    col1_w = Inches(5.9)
    col2_x = Inches(6.7)
    col2_w = Inches(6.4)

    # -- LEFT COLUMN --

    # Aplicaciones
    section_label(slide_a, "SECCIÓN A — APLICACIONES", y, prod_key)
    y += Inches(0.33)
    for app in applications:
        bullet_line(slide_a, col1_x, y, col1_w, Inches(0.28), app, bcolor)
        y += Inches(0.3)

    y += Inches(0.12)
    section_label(slide_a, "CARACTERÍSTICAS CLAVE", y, prod_key)
    y += Inches(0.33)
    for feat in features:
        bullet_line(slide_a, col1_x, y, col1_w, Inches(0.28), feat, bcolor, font_size=10)
        y += Inches(0.3)

    # -- RIGHT COLUMN --
    ry = Inches(1.2)

    section_label(slide_a, "DATOS TÉCNICOS", ry, prod_key)
    ry += Inches(0.33)
    for td_line in tech_data:
        add_textbox(
            slide_a, col2_x, ry, col2_w, Inches(0.28),
            td_line, font_size=10, color=RGBColor(0x22,0x22,0x22)
        )
        ry += Inches(0.3)

    ry += Inches(0.12)
    section_label(slide_a, "MODELOS DISPONIBLES", ry, prod_key)
    ry += Inches(0.33)
    for mod in models:
        bullet_line(slide_a, col2_x, ry, col2_w, Inches(0.28), mod, bcolor)
        ry += Inches(0.3)

    # ENCABEZADO COPY block (right side, lower)
    ry += Inches(0.2)
    section_label(slide_a, "ENCABEZADO (COPY PARA OZWALDO)", ry, prod_key)
    ry += Inches(0.33)
    header_lines = [
        (f"[Fondo {_product_color_name(prod_key)}]", 9, False, RGBColor(0x66,0x66,0x66)),
        (fantasy_name, 13, True, bcolor),
        (code_full, 10, False, RGBColor(0x33,0x33,0x33)),
    ]
    add_multiline_textbox(slide_a, col2_x, ry, col2_w, Inches(1.0), header_lines)

    # ── Slide B ───────────────────────────────────────────────────────────────
    slide_b = prs.slides.add_slide(blank_layout(prs))
    add_rect(slide_b, 0, 0, SLIDE_W, SLIDE_H, C_WHITE)
    add_product_header(slide_b, prod_key, fantasy_name, code_full, zone)
    add_footer(slide_b)

    # Two-column layout
    bx1 = Inches(0.35)
    bx2 = Inches(6.7)
    bw  = Inches(6.1)
    by  = Inches(1.2)

    # SEC B
    section_label(slide_b, "SECCIÓN B — NOTAS DE DISEÑO (para Ozwaldo)", by, prod_key)
    by += Inches(0.35)

    b_lines = []
    for item in sec_b_items:
        if item.startswith("**") and item.endswith("**"):
            # bold heading
            b_lines.append((item.strip("*"), 10, True, RGBColor(0x22,0x22,0x22)))
        else:
            b_lines.append((item, 10, False, RGBColor(0x33,0x33,0x33)))

    add_multiline_textbox(slide_b, bx1, by, bw, Inches(5.3), b_lines)

    # SEC C
    cy = Inches(1.2)
    section_label(slide_b, "SECCIÓN C — DECISIONES", cy, prod_key)
    cy += Inches(0.35)

    c_lines = []
    c_lines.append(("Resueltas en esta sesión:", 10, True, RGBColor(0x22,0x22,0x22)))
    for r in sec_c_resolved:
        c_lines.append(("• " + r, 10, False, RGBColor(0x22,0x22,0x22)))
    c_lines.append(("", 6, False, RGBColor(0x22,0x22,0x22)))
    c_lines.append(("Pendientes:", 10, True, RGBColor(0xCC,0x44,0x00)))
    for idx, p_item in enumerate(sec_c_pending, 1):
        c_lines.append((f"{idx}. {p_item}", 10, False, RGBColor(0x55,0x22,0x00)))

    add_multiline_textbox(slide_b, bx2, cy, bw, Inches(5.3), c_lines)


def _product_color_name(prod_key):
    return {"RT": "verde", "RD": "negro + bordes dorados", "RM": "gris oscuro", "RR": "azul"}[prod_key]


# ── ANNEXE — Summary table slide ──────────────────────────────────────────────

def build_annex_table(prs):
    """Single slide with the summary comparison table."""
    slide = prs.slides.add_slide(blank_layout(prs))
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, C_WHITE)

    # Header band
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.55), C_DARK_BG)
    add_textbox(
        slide,
        Inches(0.4), Inches(0.08),
        Inches(12), Inches(0.4),
        "ANEXO — Tabla Resumen de Copy por Producto (v2)",
        font_size=16, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT
    )
    add_footer(slide)

    # Table data
    headers = ["Elemento", "GST-RT ProTransfer", "GST-RD ProDigital", "GST-RM ProMotor", "GST-RR ProFrio"]
    rows = [
        ["Color header",       "Verde",              "Negro + Dorado",     "Gris",                "Azul"],
        ["Bullet 1",           "TD AJUSTABLE 0,5–10 s","PANTALLA LCD DESMONTABLE","CURVA INVERSA","PROTECCIÓN CICLADO CORTO"],
        ["Bullet 2",           "TC MÍNIMO 5 s",      "HISTORIAL 20 FALLAS","TD 0,5–3 s AUTO",    "CURVA INVERSA"],
        ["Bullet 3",           "IDEAL ATS",          "REARME MANUAL/AUTO", "RECONEXIÓN DESDE 5 s","V-ALTO Y V-BAJO AJ."],
        ["TD en etiqueta",     "0,5–10 s (aj.)",     "1–30 s (digital)",   "0,5–3 s (curva)",    "0,5–3 s (curva)"],
        ["TC en etiqueta",     "5–600 s",            "0–600 s",            "5–300 s",             "180–600 s (mín 3 min)"],
        ["Diales",             "2 (V-bajo, TC)",     "Ninguno",            "2 (V-bajo, TC)",      "3 (V-bajo, V-alto, TC)"],
        ["Badge diferenciador","ATS/TRANSFERENCIAS", "LCD + HISTORIAL",    "CURVA INVERSA",       "CICLADO CORTO 3 MIN"],
        ["Dimensiones etiq.",  "70×90 mm",           "Por definir (Oz)",   "70×90 mm",            "70×90 mm"],
        ["Modelos",            "220 / 440",          "120 / 220 / 440",    "220 / 440",           "220 / 440"],
        ["Normas",             "COVENIN + IEC",      "IEC + UL 508 + CE",  "COVENIN + IEC",       "COVENIN + IEC"],
    ]

    col_widths = [Inches(1.8), Inches(2.6), Inches(2.6), Inches(2.6), Inches(2.6)]
    row_h = Inches(0.42)
    t_left = Inches(0.25)
    t_top  = Inches(0.65)

    prod_header_colors = [C_DARK_BG, PROD_COLORS["RT"], PROD_COLORS["RD"], PROD_COLORS["RM"], PROD_COLORS["RR"]]

    # Header row
    cx = t_left
    for ci, (hdr, cw, hc) in enumerate(zip(headers, col_widths, prod_header_colors)):
        add_rect(slide, cx, t_top, cw, row_h, hc)
        add_textbox(
            slide, cx + Inches(0.05), t_top + Inches(0.06),
            cw - Inches(0.1), row_h - Inches(0.08),
            hdr, font_size=9, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT
        )
        cx += cw

    # Data rows
    for ri, row in enumerate(rows):
        ry = t_top + row_h * (ri + 1)
        cx = t_left
        row_bg = RGBColor(0xF5,0xF5,0xF5) if ri % 2 == 0 else C_WHITE
        for ci, (cell, cw) in enumerate(zip(row, col_widths)):
            add_rect(slide, cx, ry, cw, row_h, row_bg)
            fc = RGBColor(0x22,0x22,0x22)
            add_textbox(
                slide, cx + Inches(0.06), ry + Inches(0.06),
                cw - Inches(0.1), row_h - Inches(0.08),
                cell, font_size=9, bold=(ci == 0), color=fc, align=PP_ALIGN.LEFT
            )
            cx += cw


# ══════════════════════════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════════════════════════

def main():
    prs = new_prs()

    # ── Slide 1 — Cover ───────────────────────────────────────────────────────
    build_cover(prs)

    # ── Slides 2–3 — GST-RT ProTransfer ──────────────────────────────────────
    build_product_slides(
        prs,
        prod_key="RT",
        fantasy_name="ProTransfer",
        code_full="GST-RT — Supervisor Trifásico de Voltaje para Tableros y Transferencias",
        zone="Tableros con ATS, generadores, acometidas",
        applications=[
            "Transferencias automáticas (ATS)",
            "Tableros generales de distribución",
            "Generadores de respaldo",
        ],
        features=[
            "TD AJUSTABLE: 0,5–10 s",
            "TC MÍNIMO GARANTIZADO: 5 s",
            "IDEAL PARA ATS / TRANSFERENCIAS AUTOMÁTICAS",
        ],
        tech_data=[
            "TD: 0,5–10 s",
            "TC: 5–600 s",
            "Voltaje: 208/220 V~ (mod. 220)  ·  440/480 V~ (mod. 440)",
        ],
        models=[
            "GST-RT220: 208/220 V~",
            "GST-RT440: 440/480 V~",
        ],
        sec_b_items=[
            "**Dimensiones recomendadas:**",
            "70 × 90 mm (punto de partida). Casing: 80 × 100 × 38 mm.",
            "",
            "**Color dominante:**",
            "Verde — mismos Pantone de la línea monofásica Exceline Profesional.",
            "Texto sobre fondo verde: blanco. Acento para badges: naranja.",
            "",
            "**Tipografía:** Montserrat (fuente oficial Exceline Profesional).",
            "",
            "**Elemento visual prioritario:**",
            "ProTransfer: elemento más grande del encabezado.",
            "Badge «IDEAL PARA ATS / TRANSFERENCIAS»: segundo elemento de mayor peso.",
            "",
            "**Diales gráficos (2 diales):**",
            "Dial 1 — BAJO VOLTAJE (ajustable): 165–200 V~, zona naranja ~185–195 V~",
            "Dial 2 — RECONEXIÓN TC (ajustable): 5–600 s, zona naranja ~30–60 s",
            "Sobre voltaje fijo (264 V~): mostrar en tabla técnica, NO como dial.",
            "",
            "**Terminales:**",
            "Izquierdo — ENTRADA: L1, L2, L3",
            "Derecho — SALIDA: 95 (COM), 96 (NC), 98 (NA)",
            "",
            "**Íconos sugeridos:**",
            "Tablero eléctrico / rack de distribución",
            "Símbolo de transferencia (flecha bidireccional entre dos fuentes)",
            "NO incluir símbolo de curva inversa — este producto NO la tiene.",
        ],
        sec_c_resolved=[
            "Dimensiones casing: 80×100×38 mm → etiqueta 70×90 mm sugerida.",
            "Sin frases comparativas ni referencias a competidores.",
            "Colores: mismos Pantone de la línea monofásica (Ozwaldo los maneja).",
            "Fuente confirmada: Montserrat.",
        ],
        sec_c_pending=[
            "¿El sobre voltaje del GST-RT es fijo (264 V~) o ajustable? Si I&D lo hace ajustable, añadir tercer dial.",
            "Confirmación del Pantone verde exacto con Keiddys.",
            "¿La frase de argumento va impresa en etiqueta o solo en materiales PDV?",
        ],
    )

    # ── Slides 4–5 — GST-RD ProDigital ───────────────────────────────────────
    build_product_slides(
        prs,
        prod_key="RD",
        fantasy_name="ProDigital",
        code_full="GST-RD — Supervisor Trifásico Digital de Voltaje",
        zone="Acometidas críticas, subtableros, variadores, PLC",
        applications=[
            "Variadores de frecuencia y arrancadores suaves",
            "Subtableros industriales con equipos sensibles",
            "Controladores lógicos (PLC) y salas de control",
        ],
        features=[
            "PANTALLA LCD DESMONTABLE — se monta en la puerta del tablero sin abrir el tablero",
            "HISTORIAL 20 FALLAS — tipo, valor, fase, fecha y hora",
            "REARME MANUAL / AUTO SELECCIONABLE",
        ],
        tech_data=[
            "TD: 1–30 s",
            "TC: 0–600 s",
            "Voltaje: 120 V~  ·  208/220 V~  ·  440/480 V~",
            "Desbalance: 2–10% ajustable  |  IP20",
        ],
        models=[
            "GST-RD120: 120 V~",
            "GST-RD220: 208/220 V~",
            "GST-RD440: 440/480 V~",
        ],
        sec_b_items=[
            "**Dimensiones:**",
            "A definir por Ozwaldo. Casing diferente a los otros tres modelos.",
            "Punto de partida sugerido: 85 × 110 mm (casing físico 105 × 90 × 68 mm).",
            "",
            "**Tratamiento visual diferenciado:**",
            "Etiqueta NEGRA con bordes dorados oscuros elegantes — identidad exclusiva ProDigital.",
            "Negro + dorado oscuro: premium, digital, industrial de alta gama.",
            "Texto blanco sobre fondo negro.",
            "",
            "**Tipografía:** Montserrat (fuente oficial Exceline Profesional).",
            "",
            "**Elemento visual prioritario:**",
            "Pantalla LCD desmontable: diferenciador #1. Incluir representación gráfica mostrando",
            "valores de voltaje de las 3 fases (ej: VL1-L2: 218V / VL2-L3: 221V / VL3-L1: 219V).",
            "Concepto «desmontable / se monta en la puerta del tablero» debe quedar visualmente claro.",
            "",
            "**Sin diales gráficos.** El GST-RD ajusta por botones digitales, no por perillas.",
            "No incluir diales analógicos — confundirían al instalador.",
            "",
            "**Terminales:**",
            "Izquierdo — ENTRADA: L1, L2, L3",
            "Derecho — SALIDA: 95 (COM), 96 (NC), 98 (NA)",
            "Puerto GIO para GIO-Link (opcional) — si hay espacio en lateral.",
            "",
            "**Nota GIO-Link / Modbus:**",
            "El módulo GIO-Link (Modbus RTU / RS485) se vende POR SEPARADO como accesorio.",
            "No incluir en bullets de la etiqueta base.",
        ],
        sec_c_resolved=[
            "Modbus eliminado del copy base — GIO-Link se vende por separado.",
            "Tres modelos confirmados: GST-RD120 / GST-RD220 / GST-RD440.",
            "Diferenciador visual principal: pantalla LCD desmontable.",
            "Sin referencias comparativas.",
            "Tratamiento visual: negro con bordes dorados oscuros elegantes.",
        ],
        sec_c_pending=[
            "Confirmación del código de pedido GST-RD con I&D (nomenclatura final).",
            "¿Se añade mención al GIO-Link como accesorio opcional en el pie? Decisión de Keiddys.",
            "Dimensión exacta de la etiqueta: Ozwaldo debe ajustar según cara real del producto.",
        ],
    )

    # ── Slides 6–7 — GST-RM ProMotor ─────────────────────────────────────────
    build_product_slides(
        prs,
        prod_key="RM",
        fantasy_name="ProMotor",
        code_full="GST-RM — Supervisor Trifásico de Voltaje para Motores y Bombas",
        zone="Motores, bombas, ventiladores industriales",
        applications=[
            "Motores industriales de inducción",
            "Bombas centrífugas y sumergibles",
            "Sistemas hidroneumáticos",
        ],
        features=[
            "CURVA INVERSA DE VOLTAJE — reacción proporcional a la falla",
            "TD 0,5–3 s AUTOMÁTICO según severidad del voltaje",
            "RECONEXIÓN DESDE 5 s — rearranque rápido de motores",
        ],
        tech_data=[
            "TD: 0,5–3 s (curva inv. auto)",
            "TC: 5–300 s",
            "Voltaje: 208/220 V~ (mod. 220)  ·  440/480 V~ (mod. 440)",
        ],
        models=[
            "GST-RM220: 208/220 V~",
            "GST-RM440: 440/480 V~",
        ],
        sec_b_items=[
            "**Dimensiones recomendadas:**",
            "70 × 90 mm. Casing: 80 × 100 × 38 mm (igual que GST-RT y GST-RR).",
            "",
            "**Color dominante:**",
            "Gris oscuro — mismos Pantone de la línea monofásica Exceline Profesional.",
            "Texto sobre fondo gris oscuro: blanco. Acento para badges: naranja.",
            "",
            "**Tipografía:** Montserrat (fuente oficial Exceline Profesional).",
            "",
            "**Elemento visual prioritario:**",
            "Badge «CURVA INVERSA DE VOLTAJE»: diferenciador absoluto. Badge rectangular naranja",
            "con texto blanco, posición inmediatamente debajo del nombre ProMotor.",
            "",
            "**Símbolo de curva inversa:**",
            "Línea curva descendente (tiempo en eje Y, desviación de voltaje en eje X).",
            "",
            "**Diales gráficos (2 diales):**",
            "Dial 1 — BAJO VOLTAJE (ajustable): 165–200 V~ (mod. 220), zona naranja",
            "Dial 2 — RECONEXIÓN TC (ajustable): 5–300 s",
            "NO hay dial de TD — la curva inversa es automática.",
            "",
            "**Terminales:**",
            "Izquierdo — ENTRADA: L1, L2, L3",
            "Derecho — SALIDA: 95 (COM), 96 (NC), 98 (NA)",
            "",
            "**Nota crítica:** Este producto NO incluye protección de ciclado corto.",
            "Evitar iconografía de refrigeración (ese es el GST-RR ProFrio).",
        ],
        sec_c_resolved=[
            "Nomenclatura confirmada: GST-RM (no GST-RG). Modelos: GST-RM220 / GST-RM440.",
            "Sin referencias comparativas.",
            "Dimensiones casing: 80×100×38 mm → etiqueta 70×90 mm.",
            "Colores: mismos Pantone línea monofásica (Ozwaldo los maneja).",
        ],
        sec_c_pending=[
            "TD exacto curva inversa: ¿0,5–3 s es publicable como especificación oficial? I&D debe confirmar.",
            "¿El sobre voltaje del GST-RM es fijo o ajustable? Impacta número de diales.",
            "¿Hay perilla de TD en la versión con curva inversa? Confirmar con I&D.",
        ],
    )

    # ── Slides 8–9 — GST-RR ProFrio ──────────────────────────────────────────
    build_product_slides(
        prs,
        prod_key="RR",
        fantasy_name="ProFrio",
        code_full="GST-RR — Supervisor Trifásico de Voltaje para Refrigeración y Aire Acondicionado",
        zone="Compresores de refrigeración y aire acondicionado",
        applications=[
            "Compresores de A/A central y sistemas VRF",
            "Cuartos fríos y cámaras frigoríficas",
            "Chillers y equipos industriales de refrigeración",
        ],
        features=[
            "PROTECCIÓN DE CICLADO CORTO — ciclo de espera 3 min obligatorio",
            "CURVA INVERSA DE VOLTAJE — reacción proporcional a la falla",
            "VOLTAJE ALTO Y BAJO AJUSTABLES — control independiente por perilla",
        ],
        tech_data=[
            "TD: 0,5–3 s (curva inv. auto)",
            "TC: 180–600 s (mín. 3 min)",
            "Voltaje: 208/220 V~ (mod. 220)  ·  440/480 V~ (mod. 440)",
        ],
        models=[
            "GST-RR220: 208/220 V~",
            "GST-RR440: 440/480 V~",
        ],
        sec_b_items=[
            "**Dimensiones recomendadas:**",
            "70 × 90 mm (o 70 × 95 mm si los 3 diales + doble badge requieren más espacio).",
            "Casing: 80 × 100 × 38 mm.",
            "",
            "**Color dominante:**",
            "Azul — mismos Pantone línea monofásica Exceline Profesional (línea de refrigeración).",
            "Coherencia visual GSM-R120B/R150B ↔ GST-RR.",
            "Texto blanco. Acento para badges críticos: naranja.",
            "",
            "**Tipografía:** Montserrat (fuente oficial Exceline Profesional).",
            "",
            "**Elemento visual prioritario:**",
            "Badge «PROTECCIÓN DE CICLADO CORTO» con valor «3 MIN» en número grande.",
            "Es el diferenciador exclusivo de toda la línea GST-R.",
            "",
            "**Diales gráficos (3 diales — único producto de la línea con 3 diales):**",
            "Dial 1 — BAJO VOLTAJE: 165–200 V~ (mod. 220), zona naranja",
            "Dial 2 — ALTO VOLTAJE (ajustable): 230–270 V~ — diferenciador vs. GST-RT y GST-RM.",
            "Dial 3 — RECONEXIÓN TC: 180–600 s — marcar mínimo 180 s (3 min).",
            "El piso de 180 s es restricción de hardware — técnico NO puede configurar por debajo.",
            "Sugerencia: zona prohibida o tope visual en escala del dial antes de 180 s.",
            "",
            "**Terminales:**",
            "Izquierdo — ENTRADA: L1, L2, L3",
            "Derecho — SALIDA: 95 (COM), 96 (NC), 98 (NA)",
        ],
        sec_c_resolved=[
            "Sin referencias comparativas.",
            "Dimensiones casing: 80×100×38 mm → etiqueta 70×90 mm (o 70×95 mm).",
            "Colores: mismos Pantone línea monofásica de refrigeración.",
            "TC mínimo 180 s (3 min) confirmado como restricción de hardware.",
        ],
        sec_c_pending=[
            "TD exacto curva inversa: ¿0,5–3 s es publicable? I&D debe confirmar.",
            "¿El Pantone azul se unifica con GSM-R120B o se diferencia para la línea trifásica? Keiddys decide.",
            "¿El sobre voltaje del GST-RR sigue siendo ajustable con la adición de curva inversa? I&D valida.",
        ],
    )

    # ── Slide 10 — Annex table ────────────────────────────────────────────────
    build_annex_table(prs)

    # ── Save ──────────────────────────────────────────────────────────────────
    for path in PATHS:
        os.makedirs(os.path.dirname(path), exist_ok=True)
        prs.save(path)
        size_kb = os.path.getsize(path) // 1024
        print(f"Saved: {path}  ({size_kb} KB)")

    print(f"\nTotal slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
