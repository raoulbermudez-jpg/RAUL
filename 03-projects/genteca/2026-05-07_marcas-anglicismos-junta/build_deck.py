# -*- coding: utf-8 -*-
"""
build_deck.py
Genera 2026-05-07_Junta_marcas-anglicismos_v1.pptx
Vivienne — Presentation Designer · Genteca · Mayo 2026
"""

import sys
import os

sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ─────────────────────────────────────────────
# PALETA GENTECA
# ─────────────────────────────────────────────
NARANJA   = RGBColor(0xFF, 0x82, 0x00)   # Pantone 151 C
GRIS_COOL = RGBColor(0x76, 0x76, 0x76)   # Cool Gray 9 C
BLANCO    = RGBColor(0xFF, 0xFF, 0xFF)
NEGRO     = RGBColor(0x1A, 0x1A, 0x1A)
GRIS_CLARO= RGBColor(0xF2, 0xF2, 0xF2)
GRIS_MED  = RGBColor(0xCC, 0xCC, 0xCC)
VERDE_OK  = RGBColor(0x4C, 0xAF, 0x50)
AMARILLO  = RGBColor(0xFF, 0xC1, 0x07)
ROJO_SUAVE= RGBColor(0xFF, 0xEB, 0xEE)
ROJO_TEXTO= RGBColor(0xC6, 0x28, 0x28)
AZUL_MARINO=RGBColor(0x1A, 0x23, 0x7E)
AZUL_CLARO= RGBColor(0xE8, 0xEA, 0xF6)

# ─────────────────────────────────────────────
# HELPERS
# ─────────────────────────────────────────────

def rgb(r, g, b):
    return RGBColor(r, g, b)

def add_textbox(slide, text, left, top, width, height,
                font_name="Futura", font_size=18, bold=False,
                italic=False, color=NEGRO, align=PP_ALIGN.LEFT,
                bg_color=None, wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    if bg_color:
        fill = txBox.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=None):
    from pptx.util import Inches, Pt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_footer(slide, text="CONFIDENCIAL — Uso exclusivo Junta Directiva",
               extra_left=None):
    """Pie de slide con clasificación + opcionalmente logo text."""
    add_textbox(slide, text,
                left=0.2, top=7.05, width=12.8, height=0.35,
                font_name="Futura", font_size=7, bold=False,
                color=GRIS_COOL, align=PP_ALIGN.CENTER)

def add_speaker_notes(slide, notes_text):
    tf = slide.notes_slide.notes_text_frame
    tf.text = notes_text

def add_title_bar(slide, title, subtitle=None,
                  bg=NEGRO, title_color=BLANCO,
                  sub_color=NARANJA, title_size=28, sub_size=14):
    """Barra de título estándar en la parte superior del slide."""
    add_rect(slide, 0, 0, 13.33, 1.5 if subtitle else 1.1, bg)
    add_textbox(slide, title,
                left=0.35, top=0.12, width=12.6, height=0.7,
                font_name="Futura", font_size=title_size, bold=True,
                color=title_color, align=PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(slide, subtitle,
                    left=0.35, top=0.8, width=12.6, height=0.45,
                    font_name="Futura", font_size=sub_size, bold=False,
                    color=sub_color, align=PP_ALIGN.LEFT)

def bullet_block(slide, items, left, top, width, height,
                 font_name="Futura", font_size=12, color=NEGRO,
                 bullet_char="•", line_spacing=1.15):
    """
    Añade un bloque de bullets. items = list of str.
    Retorna el textbox.
    """
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = f"{bullet_char} {item}"
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
    return txBox

def add_table(slide, data, col_widths, left, top, row_height=0.38,
              header_bg=NEGRO, header_fg=BLANCO,
              row_bg_alt=GRIS_CLARO, font_size=9,
              font_name="Gill Sans MT", col_aligns=None):
    """
    data: list of rows (first row = header).
    col_widths: list of floats in inches.
    Returns table shape.
    """
    from pptx.util import Inches, Pt
    rows = len(data)
    cols = len(data[0])
    total_w = sum(col_widths)
    tbl = slide.shapes.add_table(
        rows, cols,
        Inches(left), Inches(top),
        Inches(total_w), Inches(row_height * rows)
    ).table

    # set col widths
    for ci, w in enumerate(col_widths):
        tbl.columns[ci].width = Inches(w)

    for ri, row in enumerate(data):
        for ci, cell_text in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = str(cell_text)
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            align = PP_ALIGN.LEFT
            if col_aligns and ci < len(col_aligns):
                align = col_aligns[ci]
            p.alignment = align
            run = p.runs[0] if p.runs else p.add_run()
            run.font.name = font_name
            run.font.size = Pt(font_size)
            run.font.bold = (ri == 0)
            if ri == 0:
                run.font.color.rgb = header_fg
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_bg
            else:
                run.font.color.rgb = NEGRO
                bg = BLANCO if ri % 2 == 1 else row_bg_alt
                cell.fill.solid()
                cell.fill.fore_color.rgb = bg
    return tbl

def color_table_cell(tbl, row, col, bg_color, text_color=None):
    cell = tbl.cell(row, col)
    cell.fill.solid()
    cell.fill.fore_color.rgb = bg_color
    if text_color:
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.color.rgb = text_color

# ─────────────────────────────────────────────
# PRESENTACIÓN
# ─────────────────────────────────────────────

OUTPUT = r"C:\RAUL\03-projects\genteca\2026-05-07_marcas-anglicismos-junta\03-review-and-release\2026-05-07_Junta_marcas-anglicismos_v1.pptx"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]  # blank

# ═══════════════════════════════════════════
# S-1 PORTADA
# ═══════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, NEGRO)

# Franja naranja superior sutil
add_rect(s, 0, 0, 13.333, 0.08, NARANJA)

# Logo placeholder (texto estilizado)
add_textbox(s, "GENTECA",
            left=5.2, top=0.55, width=3.0, height=0.55,
            font_name="Futura", font_size=28, bold=True,
            color=NARANJA, align=PP_ALIGN.CENTER)
add_textbox(s, "———————————",
            left=4.5, top=1.0, width=4.3, height=0.3,
            font_name="Futura", font_size=10,
            color=NARANJA, align=PP_ALIGN.CENTER)

# Titular 2 líneas
add_textbox(s, "Exceline™ Registra sus Algoritmos.",
            left=1.0, top=1.8, width=11.3, height=1.1,
            font_name="Futura", font_size=40, bold=True,
            color=BLANCO, align=PP_ALIGN.CENTER)
add_textbox(s, "Antes de que la Competencia los Identifique.",
            left=1.0, top=2.8, width=11.3, height=1.0,
            font_name="Futura", font_size=32, bold=True,
            color=BLANCO, align=PP_ALIGN.CENTER)

# Subtítulo naranja
add_textbox(s,
    "Estrategia de Propiedad Intelectual — Marcas Tecnológicas Clase 9 SAPI Venezuela",
    left=1.5, top=3.85, width=10.3, height=0.5,
    font_name="Futura", font_size=14, bold=False,
    color=NARANJA, align=PP_ALIGN.CENTER)

# Línea separadora
add_rect(s, 2.5, 4.45, 8.33, 0.04, GRIS_COOL)

# Mayo 2026 y presentado por
add_textbox(s, "Mayo 2026",
            left=1.0, top=4.6, width=5.5, height=0.35,
            font_name="Futura", font_size=12,
            color=GRIS_MED, align=PP_ALIGN.CENTER)
add_textbox(s, "Presentado por: Raoul Bermúdez",
            left=6.8, top=4.6, width=5.5, height=0.35,
            font_name="Futura", font_size=12,
            color=GRIS_MED, align=PP_ALIGN.CENTER)

# Clasificación pie
add_textbox(s, "CONFIDENCIAL — Uso exclusivo Junta Directiva",
            left=0.2, top=7.05, width=12.8, height=0.35,
            font_name="Futura", font_size=7,
            color=GRIS_COOL, align=PP_ALIGN.CENTER)

add_speaker_notes(s,
    "Buenos días / tardes. Esta presentación tiene un solo objetivo: que la Junta tome hoy una decisión sobre si Genteca adopta la estrategia de registrar los nombres de sus tecnologías propietarias como marcas en SAPI Venezuela.\n\n"
    "No es una presentación de antecedentes. El trabajo de investigación ya está hecho. Lo que presento es una recomendación clara, con la ruta de implementación lista.\n\n"
    "Voy a comenzar con la pregunta que sé que la Junta va a formular — porque es la más razonable. Y voy a responderla antes de que la hagan."
)

# ═══════════════════════════════════════════
# S-2 LA PARADOJA ANTICIPADA
# ═══════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, BLANCO)

add_title_bar(s,
    "“Exceline” es una Palabra Inglesa. Llevan Décadas Usándola.",
    subtitle="La paradoja que la Junta va a formular — y su respuesta.",
    bg=NEGRO, title_color=BLANCO, sub_color=NARANJA,
    title_size=22, sub_size=12)

# Columna izquierda — gris
add_rect(s, 0.25, 1.65, 6.1, 5.1, GRIS_CLARO)
add_textbox(s, "LO QUE LA JUNTA VA A PENSAR",
            left=0.35, top=1.75, width=5.9, height=0.45,
            font_name="Futura", font_size=11, bold=True,
            color=GRIS_COOL, align=PP_ALIGN.CENTER)

preguntas = [
    "¿Cómo vamos a registrar nombres en inglés\nsi los venezolanos usamos español?",
    "¿No rechaza SAPI marcas en idioma extranjero?",
    "¿No es contradictorio que Genteca, empresa venezolana,\n use anglicismos en sus marcas?",
    "¿Exceline misma no es inglés?",
]
for i, q in enumerate(preguntas):
    add_textbox(s, q,
                left=0.4, top=2.25 + i * 1.0, width=5.8, height=0.85,
                font_name="Futura", font_size=10, italic=True,
                color=NEGRO, align=PP_ALIGN.LEFT)

# Columna derecha — naranja suave
add_rect(s, 6.55, 1.65, 6.5, 5.1, rgb(0xFF, 0xF3, 0xE0))
add_textbox(s, "LA RESPUESTA",
            left=6.65, top=1.75, width=6.3, height=0.45,
            font_name="Futura", font_size=11, bold=True,
            color=NARANJA, align=PP_ALIGN.CENTER)

respuestas = [
    "SAPI no exige español. Prohíbe descriptores literales\n(Art. 34 LPI). El idioma no es la variable relevante.",
    "No. SAPI rechaza signos que describan el producto.\nLos signos de fantasía en cualquier idioma son registrables.",
    "No. Exceline Profesional ya usa anglicismos técnicos\n(inverter, relay, encoder) sin fricción en el canal VE.",
    "Exactamente. Y Exceline es registrable precisamente\nporque no describe ningún protector. Identifica uno.",
]
for i, r in enumerate(respuestas):
    add_textbox(s, r,
                left=6.6, top=2.25 + i * 1.0, width=6.3, height=0.85,
                font_name="Futura", font_size=10,
                color=NEGRO, align=PP_ALIGN.LEFT)

add_footer(s)
add_speaker_notes(s,
    "Anticipo la objeción más frecuente: ¿cómo vamos a registrar nombres en inglés?\n\n"
    "La respuesta corta: SAPI no exige español. Lo que SAPI rechaza son los descriptores literales, no el idioma extranjero.\n\n"
    "La respuesta larga está en esta diapositiva. Cada pregunta que la Junta tiene tiene una respuesta técnica y legal precisa. "
    "Y la paradoja más bella de esta estrategia es que Exceline — la marca madre — ya es un anglicismo. Llevan décadas usándola. El patrón ya está instalado."
)

# ═══════════════════════════════════════════
# S-3 DESCRIPCIÓN VS NOMBRE FANTASÍA
# ═══════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, BLANCO)

add_title_bar(s,
    "La Distinción Crítica: Descripción vs. Nombre Fantasía.",
    subtitle="Lo que SAPI rechaza no es el inglés. Es el descriptor literal.",
    bg=NEGRO, title_color=BLANCO, sub_color=NARANJA,
    title_size=22, sub_size=12)

data = [
    ["Tipo", "RECHAZADO por SAPI", "APROBADO por SAPI"],
    ["Anglicismo descriptivo", '"Thermal Shield"\n(escudo térmico)', '"Thermo-Safe™"'],
    ["Por qué", "Describe literalmente la función", "No describe — identifica"],
    ["Analogía Venezuela", '"Inverter" para AA (genérico)', '"Exceline™" para protectores (nombre propio)'],
    ["Artículo LPI", "Art. 34 — descriptividad", "Signo de fantasía — registrable"],
]

tbl = add_table(s, data,
                col_widths=[2.8, 4.5, 4.5],
                left=0.35, top=1.75,
                row_height=0.75,
                header_bg=NEGRO, header_fg=BLANCO,
                font_size=10, font_name="Gill Sans MT")

# Resaltar columna RECHAZADO en rojo suave
for ri in range(1, 5):
    color_table_cell(tbl, ri, 1, rgb(0xFF, 0xEB, 0xEE), ROJO_TEXTO)

# Resaltar columna APROBADO en verde suave
for ri in range(1, 5):
    color_table_cell(tbl, ri, 2, rgb(0xE8, 0xF5, 0xE9), rgb(0x1B, 0x5E, 0x20))

# Línea cierre
add_rect(s, 0.35, 5.6, 12.6, 0.55, NEGRO)
add_textbox(s,
    "La clave no está en el idioma. Está en si el nombre describe o identifica.",
    left=0.45, top=5.65, width=12.4, height=0.4,
    font_name="Futura", font_size=13, bold=True,
    color=BLANCO, align=PP_ALIGN.CENTER)

add_footer(s)
add_speaker_notes(s,
    "Esta es la distinción técnica que determina si SAPI otorga o niega el registro.\n\n"
    "El Art. 34 de la LPI venezolana declara inadmisibles los signos que describan directamente el producto o su función. No menciona el idioma.\n\n"
    "La diferencia entre 'Thermal Shield' y 'Thermo-Safe' no es estética. Es jurídica. El primero describe. El segundo identifica.\n\n"
    "Y el canal venezolano ya vive esta distinción: el técnico dice 'inverter' para cualquier compresor de velocidad variable — es genérico. "
    "Pero dice 'Exceline' para un protector específico — eso es una marca. Lo mismo que proponemos con StaggerStart, ForensLog y los demás."
)

# ═══════════════════════════════════════════
# S-4 EL PATRÓN EN EL PORTAFOLIO PROPIO
# ═══════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, BLANCO)

add_title_bar(s,
    "Nuestros Mejores Nombres No Describen. Nombran.",
    subtitle="El patrón ganador ya está en nuestro portafolio.",
    bg=NEGRO, title_color=BLANCO, sub_color=NARANJA,
    title_size=24, sub_size=12)

# Bloque 1 — gris
add_rect(s, 0.25, 1.75, 4.0, 4.9, GRIS_CLARO)
add_textbox(s, "EL NOMBRE DE LA MARCA",
            left=0.35, top=1.85, width=3.8, height=0.45,
            font_name="Futura", font_size=10, bold=True,
            color=GRIS_COOL, align=PP_ALIGN.CENTER)
add_textbox(s, "Exceline™\nGenius™",
            left=0.35, top=2.4, width=3.8, height=0.9,
            font_name="Futura", font_size=22, bold=True,
            color=NEGRO, align=PP_ALIGN.CENTER)
add_textbox(s, "No describen.\nIdentifican.",
            left=0.35, top=3.4, width=3.8, height=0.6,
            font_name="Futura", font_size=11, italic=True,
            color=GRIS_COOL, align=PP_ALIGN.CENTER)

# Bloque 2 — naranja
add_rect(s, 4.55, 1.75, 4.3, 4.9, rgb(0xFF, 0xF3, 0xE0))
add_rect(s, 4.55, 1.75, 4.3, 0.45, NARANJA)
add_textbox(s, "LOS NOMBRES QUE PROPONEMOS",
            left=4.65, top=1.78, width=4.1, height=0.38,
            font_name="Futura", font_size=9, bold=True,
            color=BLANCO, align=PP_ALIGN.CENTER)
for i, nm in enumerate(["StaggerStart™", "ForensLog™", "FlickerGuard™"]):
    add_textbox(s, nm,
                left=4.65, top=2.35 + i * 0.75, width=4.1, height=0.6,
                font_name="Futura", font_size=18, bold=True,
                color=NARANJA, align=PP_ALIGN.CENTER)
add_textbox(s, "Nombran el algoritmo.",
            left=4.65, top=4.7, width=4.1, height=0.5,
            font_name="Futura", font_size=11, italic=True,
            color=NEGRO, align=PP_ALIGN.CENTER)

# Bloque 3 — blanco / cursiva (no registrar)
add_rect(s, 9.05, 1.75, 4.0, 4.9, BLANCO)
add_rect(s, 9.05, 1.75, 4.0, 0.04, GRIS_MED)
add_textbox(s, "LO QUE NO VAMOS A REGISTRAR",
            left=9.15, top=1.85, width=3.8, height=0.45,
            font_name="Futura", font_size=9, bold=True,
            color=GRIS_COOL, align=PP_ALIGN.CENTER)
descriptores = [
    "Arranque inteligente\nescalonado",
    "Historial de 100 fallas",
    "Detección de\nparpadeos 20ms",
]
for i, d in enumerate(descriptores):
    add_textbox(s, d,
                left=9.15, top=2.4 + i * 1.1, width=3.8, height=0.85,
                font_name="Futura", font_size=12, italic=True,
                color=GRIS_COOL, align=PP_ALIGN.CENTER)
add_textbox(s, "(descriptor, no nombre)",
            left=9.15, top=5.7, width=3.8, height=0.4,
            font_name="Futura", font_size=9, italic=True,
            color=GRIS_MED, align=PP_ALIGN.CENTER)

# Anclaje pie
add_rect(s, 0.25, 6.72, 12.85, 0.42, rgb(0xF5, 0xF5, 0xF5))
add_textbox(s,
    "Los descriptores van al empaque y al datasheet.  Los nombres van a SAPI.",
    left=0.35, top=6.75, width=12.6, height=0.35,
    font_name="Futura", font_size=11, bold=True,
    color=NARANJA, align=PP_ALIGN.CENTER)

add_footer(s)
add_speaker_notes(s,
    "El patrón ganador ya existe en el portafolio de Genteca. Exceline no describe ningún protector eléctrico. Genius no describe ningún relé. "
    "Pero cualquier técnico en Venezuela los identifica inmediatamente con el producto correcto.\n\n"
    "Eso es exactamente lo que proponemos hacer con los algoritmos propietarios: bautizarlos con nombres que identifiquen, no que describan.\n\n"
    "Los descriptores — arranque escalonado, historial de 100 fallas — pertenecen al empaque, al datasheet, a la hoja de especificaciones. "
    "No pertenecen al registro de SAPI. Lo que va a SAPI es el nombre de fantasía que los identifica."
)

# ═══════════════════════════════════════════
# S-5 LA BARRERA DEL PRIMER REGISTRO
# ═══════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, BLANCO)

add_title_bar(s,
    "Registrar Primero: La Barrera que el Follower No Puede Comprar.",
    subtitle="Por qué el registro en SAPI es una ventaja competitiva real, no un trámite.",
    bg=NEGRO, title_color=BLANCO, sub_color=NARANJA,
    title_size=21, sub_size=12)

# Diagrama 4 pasos — lado izquierdo
add_textbox(s, "¿Qué pasa si Genteca NO registra?",
            left=0.3, top=1.7, width=6.0, height=0.4,
            font_name="Futura", font_size=12, bold=True,
            color=ROJO_TEXTO, align=PP_ALIGN.LEFT)

pasos = [
    ("1", "Genteca desarrolla y comunica la innovación"),
    ("2", "Competidor identifica el nombre y registra primero en SAPI"),
    ("3", "SAPI otorga la marca al competidor"),
    ("4", "Genteca recibe notificación de infracción"),
]
colors_p = [GRIS_CLARO, rgb(0xFF, 0xF3, 0xE0), rgb(0xFF, 0xEB, 0xEE), rgb(0xFF, 0xCB, 0xD1)]
for i, (num, texto) in enumerate(pasos):
    add_rect(s, 0.3, 2.15 + i * 0.88, 5.9, 0.78, colors_p[i])
    add_textbox(s, num,
                left=0.4, top=2.2 + i * 0.88, width=0.5, height=0.65,
                font_name="Futura", font_size=22, bold=True,
                color=NARANJA, align=PP_ALIGN.CENTER)
    add_textbox(s, texto,
                left=0.95, top=2.2 + i * 0.88, width=5.1, height=0.65,
                font_name="Futura", font_size=11,
                color=NEGRO, align=PP_ALIGN.LEFT)

# Lado derecho — lo que el registro garantiza
add_textbox(s, "Lo que el registro garantiza:",
            left=6.8, top=1.7, width=6.0, height=0.4,
            font_name="Futura", font_size=12, bold=True,
            color=VERDE_OK, align=PP_ALIGN.LEFT)

garantias = [
    "Exclusividad 10 años renovables sobre el nombre",
    "Derecho de oposición frente a registros similares",
    "Barrera competitiva sin precio de copia",
    "Activo de propiedad intelectual en balance",
]
bullet_block(s, garantias,
             left=6.8, top=2.15, width=6.0, height=2.8,
             font_name="Futura", font_size=12, color=NEGRO,
             bullet_char="✓")

# Cierre
add_rect(s, 0.3, 6.3, 12.75, 0.55, NEGRO)
add_textbox(s,
    "El costo de registrar: $880–$1,640 por marca.  |  El costo de no registrar: ningún límite.",
    left=0.4, top=6.35, width=12.55, height=0.42,
    font_name="Futura", font_size=12, bold=True,
    color=NARANJA, align=PP_ALIGN.CENTER)

add_footer(s)
add_speaker_notes(s,
    "El registro de marca en SAPI no es un trámite de cumplimiento. Es una barrera competitiva.\n\n"
    "El escenario sin registro es simple: Genteca desarrolla, comunica, y cualquier actor que preste atención puede registrar el nombre antes. "
    "SAPI opera por orden de llegada. El primero que registra tiene la marca. El segundo que llega — aunque haya desarrollado la tecnología — recibe una notificación de infracción.\n\n"
    "Con registro: exclusividad por 10 años renovables. Derecho a oponerse a cualquier registro similar. Una barrera que ningún competidor puede comprar porque ya está ocupada.\n\n"
    "El costo del registro es conocido: entre $880 y $1,640 por marca con abogado marcario. El costo de no registrar no tiene techo."
)

# ═══════════════════════════════════════════
# S-6 PATRÓN COMPETIDORES GLOBALES
# ═══════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, BLANCO)

add_title_bar(s,
    "Industria Global ya Probó este Patrón. En Venezuela, la Pista Está Vacía.",
    subtitle="Los grandes registraron nombres fantasía — no descriptores técnicos.",
    bg=NEGRO, title_color=BLANCO, sub_color=NARANJA,
    title_size=20, sub_size=12)

data6 = [
    ["Marca", "Empresa", "Lo que NO dice", "Lo que SÍ nombra"],
    ["SIRIUS", "Siemens", "arrancador directo", "línea de arranque"],
    ["TeSys", "Schneider", "sistema de control de motor", "ecosistema"],
    ["Easergy", "Schneider", "protección de red", "familia de relés"],
    ["ArmorStart", "Rockwell", "motor protegido", "concepto de armor"],
    ["Relion", "ABB", "relé de protección", "plataforma"],
    ["xStart", "Eaton", "arrancador controlado", "arranque flexible"],
]

tbl6 = add_table(s, data6,
                 col_widths=[2.2, 2.4, 4.4, 3.4],
                 left=0.35, top=1.75,
                 row_height=0.55,
                 header_bg=NEGRO, header_fg=BLANCO,
                 font_size=10, font_name="Gill Sans MT")

# Columna "Lo que NO dice" en rojo suave
for ri in range(1, 7):
    color_table_cell(tbl6, ri, 2, rgb(0xFF, 0xEB, 0xEE), ROJO_TEXTO)

# Columna "Lo que SÍ nombra" en verde suave
for ri in range(1, 7):
    color_table_cell(tbl6, ri, 3, rgb(0xE8, 0xF5, 0xE9), rgb(0x1B, 0x5E, 0x20))

# Bullet cierre
add_rect(s, 0.35, 5.65, 12.6, 0.9, rgb(0xFF, 0xF3, 0xE0))
add_textbox(s,
    "En Venezuela: ningún competidor local ha registrado nombres de algoritmos propietarios en Clase 9.\n"
    "La pista está libre. Genteca puede ser el primero.",
    left=0.5, top=5.7, width=12.3, height=0.8,
    font_name="Futura", font_size=12, bold=True,
    color=NARANJA, align=PP_ALIGN.LEFT)

add_footer(s)
add_speaker_notes(s,
    "La industria global de protección eléctrica lleva décadas aplicando exactamente este patrón. "
    "Siemens no registró 'arrancador directo'. Registró SIRIUS. Schneider no registró 'sistema de control de motor'. Registró TeSys.\n\n"
    "Ninguno de esos nombres describe la función. Todos identifican una plataforma, una línea, un concepto. Ese es el estándar de la industria global.\n\n"
    "Ahora la pregunta es: ¿qué pasa en Venezuela? La respuesta es que la pista está vacía. "
    "Ningún fabricante local ha registrado nombres de algoritmos propietarios en Clase 9 SAPI. "
    "Genteca puede ser el primer fabricante venezolano en hacer lo que Siemens, Schneider, ABB y Rockwell ya hicieron hace décadas."
)

# ═══════════════════════════════════════════
# S-7 LOS 7 CANDIDATOS: ESTADO ACTUAL
# ═══════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, BLANCO)

add_title_bar(s,
    "Los 7 Candidatos de Genteca: Estado Actual de Cada Uno.",
    subtitle="Top 7 con HDE confirmada — diferentes niveles de madurez para el filing.",
    bg=NEGRO, title_color=BLANCO, sub_color=NARANJA,
    title_size=22, sub_size=12)

data7 = [
    ["#", "Marca", "Tecnología ancla", "Producto", "Estado"],
    ["1", "StaggerStart™", "Staggered Intelligent Reconnect", "GSM-MB/RB/RE", "✅ Listo para filing"],
    ["2", "ForensLog™", "100-Fault Forensic History", "GIII+MV", "✅ Listo para filing"],
    ["3", "FlickerGuard™", "20ms Flicker Detection", "GSM-L", "⚠ Ver Advertencia A"],
    ["4", "Thermo-Safe™", "Función térmica integrada", "GSM-MB/RB/RE", "⚠⚠ Ver Advertencia B"],
    ["5", "TripleLock™", "Triple-Fault Lockout", "GIII+/GUCT/GOC", "⚠ Ver Advertencia C"],
    ["6", "TaskMemory™", "Smart Task Memory", "GRN-MV", "⚠ Ver Advertencia D"],
    ["7", "ThermalCurve™", "Cold/Hot Thermal Curve Model", "GIII+/GUCT/GOC", "⚠⚠ Ver Advertencia E"],
]

tbl7 = add_table(s, data7,
                 col_widths=[0.4, 1.8, 3.5, 2.0, 2.3],
                 left=0.2, top=1.65,
                 row_height=0.48,
                 header_bg=NEGRO, header_fg=BLANCO,
                 font_size=9, font_name="Gill Sans MT")

# Color por estado
color_table_cell(tbl7, 1, 4, rgb(0xE8, 0xF5, 0xE9), rgb(0x1B, 0x5E, 0x20))
color_table_cell(tbl7, 2, 4, rgb(0xE8, 0xF5, 0xE9), rgb(0x1B, 0x5E, 0x20))
color_table_cell(tbl7, 3, 4, rgb(0xFF, 0xF9, 0xC4), rgb(0x7B, 0x61, 0x00))
color_table_cell(tbl7, 4, 4, rgb(0xFF, 0xF3, 0xE0), rgb(0xE6, 0x51, 0x00))
color_table_cell(tbl7, 5, 4, rgb(0xFF, 0xF9, 0xC4), rgb(0x7B, 0x61, 0x00))
color_table_cell(tbl7, 6, 4, rgb(0xFF, 0xF9, 0xC4), rgb(0x7B, 0x61, 0x00))
color_table_cell(tbl7, 7, 4, rgb(0xFF, 0xF3, 0xE0), rgb(0xE6, 0x51, 0x00))

# Advertencias en pie — texto pequeño
advertencias = (
    "Advertencia A — FlickerGuard™: Candidato a registro en Clase 9 — SAPI Venezuela. "
    "Verificación preliminar completada. Requiere búsqueda OMPI ampliada antes del filing.\n"
    "Advertencia B — Thermo-Safe™: Candidato a registro en Clase 9 — SAPI Venezuela. "
    "Verificación preliminar completada. Requiere búsqueda OMPI ampliada Y documentación técnica interna "
    "(HDE para uso público de la función NTC) antes del filing.\n"
    "Advertencia C — TripleLock™: Candidato a registro en Clase 9 — SAPI Venezuela. "
    "Verificación preliminar completada. Requiere búsqueda sectorial ampliada (potencial colisión con marcas "
    "del sector seguridad física por el término \"Lock\").\n"
    "Advertencia D — TaskMemory™: Candidato a registro en Clase 9 — SAPI Venezuela. "
    "Verificación preliminar completada. Riesgo de colisión con marcas de software — requiere búsqueda urgente antes del filing.\n"
    "Advertencia E — ThermalCurve™: Candidato a registro — pendiente de documentación técnica interna. "
    "No incluido en primer batch SAPI."
)

add_rect(s, 0.2, 5.5, 12.93, 1.68, rgb(0xF9, 0xF9, 0xF9))
add_textbox(s, advertencias,
            left=0.3, top=5.52, width=12.73, height=1.6,
            font_name="Gill Sans MT", font_size=6.5,
            color=GRIS_COOL, align=PP_ALIGN.LEFT)

add_footer(s)
add_speaker_notes(s,
    "Tenemos 7 candidatos identificados, todos con HDE — Hoja de Datos de Especificaciones — confirmada por el equipo de ingeniería.\n\n"
    "Dos están listos para filing inmediato: StaggerStart y ForensLog. El análisis preliminar OMPI no muestra colisiones significativas y la documentación técnica está completa.\n\n"
    "Los otros cinco tienen diferentes niveles de trabajo previo requerido. Las advertencias al pie son precisas: "
    "no son señales de alarma sobre la registrabilidad del nombre, sino indicaciones de pasos adicionales de debida diligencia antes de presentar la solicitud FM-02.\n\n"
    "La Junta no necesita resolver estas advertencias hoy. Solo necesita decidir si adopta la estrategia. "
    "El equipo legal y el abogado marcario gestionan el detalle de cada candidato."
)

# ═══════════════════════════════════════════
# S-8 THERMO-SAFE VS THERMOSHIELD NTC
# ═══════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, BLANCO)

add_title_bar(s,
    "Thermo-Safe™ No es ThermoShield NTC. El Pivote Importa.",
    subtitle="Por qué cambiamos el nombre candidato — y qué ganamos con eso.",
    bg=NEGRO, title_color=BLANCO, sub_color=NARANJA,
    title_size=22, sub_size=12)

# ANTES
add_rect(s, 0.5, 1.8, 5.5, 2.2, rgb(0xFF, 0xEB, 0xEE))
add_rect(s, 0.5, 1.8, 5.5, 0.45, rgb(0xC6, 0x28, 0x28))
add_textbox(s, "ANTES",
            left=0.6, top=1.83, width=5.3, height=0.38,
            font_name="Futura", font_size=12, bold=True,
            color=BLANCO, align=PP_ALIGN.CENTER)
add_textbox(s, "ThermoShield NTC",
            left=0.6, top=2.35, width=5.3, height=0.9,
            font_name="Futura", font_size=26, bold=True,
            color=rgb(0xC6, 0x28, 0x28), align=PP_ALIGN.CENTER)
add_textbox(s, "✗",
            left=5.6, top=2.5, width=0.5, height=0.7,
            font_name="Futura", font_size=36, bold=True,
            color=rgb(0xC6, 0x28, 0x28), align=PP_ALIGN.CENTER)

# Flecha
add_textbox(s, "→",
            left=6.1, top=2.4, width=1.0, height=0.8,
            font_name="Futura", font_size=40, bold=True,
            color=NARANJA, align=PP_ALIGN.CENTER)

# DESPUÉS
add_rect(s, 7.3, 1.8, 5.5, 2.2, rgb(0xE8, 0xF5, 0xE9))
add_rect(s, 7.3, 1.8, 5.5, 0.45, VERDE_OK)
add_textbox(s, "DESPUÉS",
            left=7.4, top=1.83, width=5.3, height=0.38,
            font_name="Futura", font_size=12, bold=True,
            color=BLANCO, align=PP_ALIGN.CENTER)
add_textbox(s, "Thermo-Safe™",
            left=7.4, top=2.35, width=5.3, height=0.9,
            font_name="Futura", font_size=26, bold=True,
            color=VERDE_OK, align=PP_ALIGN.CENTER)
add_textbox(s, "✓",
            left=12.6, top=2.5, width=0.45, height=0.7,
            font_name="Futura", font_size=36, bold=True,
            color=VERDE_OK, align=PP_ALIGN.CENTER)

# 3 razones
razones = [
    ("\U0001f512", "NTC rompe el patrón ganador",
     "Incorporar el componente en el nombre revela el secreto industrial "
     "y crea riesgo de descriptividad por referencia técnica directa."),
    ("⇄", "Coherencia con decisión de Junta en empaques",
     "La Junta ya decidió no mencionar NTC en empaques GSM. "
     "El nombre de marca debe ser consistente con esa decisión."),
    ("\U0001f4dc", "Mejor perfil de registrabilidad SAPI",
     "Thermo-Safe es signo de fantasía puro. ThermoShield NTC "
     "presenta riesgo de descriptividad combinada (Art. 34)."),
]
for i, (icono, titulo, desc) in enumerate(razones):
    add_rect(s, 0.3 + i * 4.35, 4.3, 4.1, 2.0, GRIS_CLARO)
    add_textbox(s, titulo,
                left=0.4 + i * 4.35, top=4.38, width=3.9, height=0.45,
                font_name="Futura", font_size=11, bold=True,
                color=NEGRO, align=PP_ALIGN.CENTER)
    add_textbox(s, desc,
                left=0.4 + i * 4.35, top=4.88, width=3.9, height=1.35,
                font_name="Futura", font_size=9.5,
                color=GRIS_COOL, align=PP_ALIGN.CENTER)

add_footer(s)
add_speaker_notes(s,
    "Este slide responde a una pregunta específica: ¿por qué cambiamos de ThermoShield NTC a Thermo-Safe?\n\n"
    "La razón más importante es estratégica: incluir 'NTC' en el nombre de marca hace pública la existencia del componente NTC en el producto. "
    "Genteca ha decidido — correctamente — que el sensor NTC es un secreto industrial. El nombre de la marca no puede contradecir esa decisión.\n\n"
    "La segunda razón es legal: ThermoShield NTC presenta mayor riesgo de descriptividad combinada ante SAPI. "
    "Thermo-Safe no hace referencia a ningún componente ni tecnología específica — identifica una función de seguridad, que es exactamente el estándar del signo de fantasía.\n\n"
    "La tercera razón es coherencia: la Junta ya tomó la decisión sobre los empaques GSM. El nombre de marca en SAPI debe ser consistente con lo que aparece en el canal."
)

# ═══════════════════════════════════════════
# S-9 CONVERGENCIA IP + EMPAQUES GSM
# ═══════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, BLANCO)

add_title_bar(s,
    "Thermo-Safe™ Ancla el Deck de Empaques GSM. No lo Contradice.",
    subtitle="La misma decisión estratégica vista desde IP y desde canal.",
    bg=NEGRO, title_color=BLANCO, sub_color=NARANJA,
    title_size=21, sub_size=12)

# Línea 1 — Empaques GSM
add_rect(s, 0.3, 1.75, 12.7, 0.4, GRIS_CLARO)
add_textbox(s, "LÍNEA 1 — Decisión empaques GSM",
            left=0.4, top=1.78, width=12.5, height=0.35,
            font_name="Futura", font_size=10, bold=True,
            color=GRIS_COOL, align=PP_ALIGN.LEFT)

bloques_l1 = [
    ("Alternativa B\nJunta", rgb(0xE8, 0xF5, 0xE9), rgb(0x1B, 0x5E, 0x20)),
    ("Badge:\n\"Autoprotección\ntérmica activa*\"", rgb(0xFF, 0xF9, 0xC4), rgb(0x7B, 0x61, 0x00)),
    ("Consumidor ve\nbeneficio,\nno componente", GRIS_CLARO, GRIS_COOL),
]
for i, (txt, bg, fg) in enumerate(bloques_l1):
    add_rect(s, 0.3 + i * 4.3, 2.22, 4.0, 1.3, bg)
    add_textbox(s, txt,
                left=0.4 + i * 4.3, top=2.28, width=3.8, height=1.1,
                font_name="Futura", font_size=11,
                color=fg, align=PP_ALIGN.CENTER)

# Línea 2 — Decisión IP
add_rect(s, 0.3, 3.7, 12.7, 0.4, rgb(0xFF, 0xF3, 0xE0))
add_textbox(s, "LÍNEA 2 — Decisión IP propuesta",
            left=0.4, top=3.73, width=12.5, height=0.35,
            font_name="Futura", font_size=10, bold=True,
            color=NARANJA, align=PP_ALIGN.LEFT)

bloques_l2 = [
    ("Nombre candidato\nSAPI:\nThermo-Safe™", rgb(0xFF, 0xF3, 0xE0), NARANJA),
    ("Nombre identifica\nfunción sin revelar\ncomponente", rgb(0xFF, 0xF9, 0xC4), rgb(0x7B, 0x61, 0x00)),
    ("Genteca registra\npropiedad de\nfunción, no de\ncomponente", GRIS_CLARO, GRIS_COOL),
]
for i, (txt, bg, fg) in enumerate(bloques_l2):
    add_rect(s, 0.3 + i * 4.3, 4.18, 4.0, 1.45, bg)
    add_textbox(s, txt,
                left=0.4 + i * 4.3, top=4.24, width=3.8, height=1.3,
                font_name="Futura", font_size=11,
                color=fg, align=PP_ALIGN.CENTER)

# Punto convergencia
add_rect(s, 0.3, 5.75, 12.7, 0.7, NARANJA)
add_textbox(s,
    "Una sola estrategia coherente: el componente NTC es secreto industrial. "
    "Lo que se comunica — en el empaque y en la marca — es el beneficio y la función.",
    left=0.4, top=5.78, width=12.5, height=0.62,
    font_name="Futura", font_size=11, bold=True,
    color=BLANCO, align=PP_ALIGN.CENTER)

# Advertencia B pie
add_textbox(s,
    "Advertencia B — Thermo-Safe™: Requiere búsqueda OMPI ampliada + HDE pública NTC antes del filing.",
    left=0.3, top=6.5, width=12.7, height=0.3,
    font_name="Gill Sans MT", font_size=7,
    color=GRIS_COOL, align=PP_ALIGN.LEFT)

add_footer(s)
add_speaker_notes(s,
    "Este slide conecta dos decisiones que la Junta ya tomó o está tomando en paralelo.\n\n"
    "Por el lado del empaque GSM: la Junta eligió la Alternativa B — badge de beneficio 'Autoprotección térmica activa', "
    "sin mencionar el componente NTC. El consumidor ve el resultado, no la tecnología interna.\n\n"
    "Por el lado de IP: el nombre Thermo-Safe en SAPI hace exactamente lo mismo. Identifica la función de autoprotección térmica "
    "sin revelar que la tecnología usa un sensor NTC específico.\n\n"
    "No son dos decisiones. Es una sola estrategia aplicada en dos frentes: canal y registro. El NTC es y seguirá siendo secreto industrial. "
    "Lo que Genteca comunica — y lo que Genteca registra — es la función y el beneficio."
)

# ═══════════════════════════════════════════
# S-10 EL MÉTODO EN 3 PASOS
# ═══════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, BLANCO)

add_title_bar(s,
    "El Caso Ancla Demuestra el Método. Ahora Escala.",
    subtitle="Thermo-Safe™ es el prototipo — el portafolio completo sigue la misma lógica.",
    bg=NEGRO, title_color=BLANCO, sub_color=NARANJA,
    title_size=24, sub_size=12)

pasos10 = [
    ("Paso 1\nIDENTIFICAR", "\U0001f50d",
     "Función térmica integrada confirmada por ingeniería.\n"
     "La HDE existe. El algoritmo está documentado."),
    ("Paso 2\nBAUTIZAR", "\U0001f3f7",
     "No 'Sensor NTC' sino 'Thermo-Safe™'.\n"
     "Signo de fantasía registrable (Art. 34 LPI)."),
    ("Paso 3\nREGISTRAR", "\U0001f4dc",
     "Máximo 30 días desde primer uso público.\n"
     "Proceso SAPI: 6–12 meses."),
]
bg_pasos = [GRIS_CLARO, rgb(0xFF, 0xF3, 0xE0), rgb(0xE8, 0xF5, 0xE9)]
for i, (titulo, icono, desc) in enumerate(pasos10):
    add_rect(s, 0.3 + i * 4.35, 1.75, 4.1, 4.1, bg_pasos[i])
    add_textbox(s, titulo,
                left=0.4 + i * 4.35, top=1.83, width=3.9, height=0.75,
                font_name="Futura", font_size=16, bold=True,
                color=NARANJA, align=PP_ALIGN.CENTER)
    add_textbox(s, icono,
                left=0.4 + i * 4.35, top=2.65, width=3.9, height=0.8,
                font_name="Segoe UI Emoji", font_size=36,
                color=NEGRO, align=PP_ALIGN.CENTER)
    add_textbox(s, desc,
                left=0.4 + i * 4.35, top=3.55, width=3.9, height=2.2,
                font_name="Futura", font_size=11,
                color=NEGRO, align=PP_ALIGN.LEFT)
    # Flecha entre bloques
    if i < 2:
        add_textbox(s, "→",
                    left=4.35 + i * 4.35, top=3.4, width=0.4, height=0.6,
                    font_name="Futura", font_size=28, bold=True,
                    color=NARANJA, align=PP_ALIGN.CENTER)

# Pie
add_rect(s, 0.3, 6.1, 12.7, 0.55, GRIS_CLARO)
add_textbox(s,
    "El mismo método aplicado a todos los candidatos: "
    "StaggerStart™ · ForensLog™ · FlickerGuard™ · TripleLock™ · TaskMemory™ · ThermalCurve™",
    left=0.4, top=6.13, width=12.5, height=0.45,
    font_name="Futura", font_size=10, bold=True,
    color=NARANJA, align=PP_ALIGN.CENTER)

add_footer(s)
add_speaker_notes(s,
    "Thermo-Safe no es solo un nombre candidato. Es la demostración del método completo.\n\n"
    "Paso 1: Ingeniería confirma que la función existe, tiene una HDE documentada y es diferenciadora. "
    "Sin esa confirmación técnica, no hay marca que registrar.\n\n"
    "Paso 2: En lugar de usar el nombre técnico del componente o de la función, se construye un signo de fantasía "
    "que identifica la función sin describirla. Ese es el trabajo de bautismo.\n\n"
    "Paso 3: El registro debe iniciarse en un plazo máximo de 30 días desde que el nombre se usa públicamente por primera vez. "
    "Después de ese punto, cualquier tercero que lo vea puede registrarlo antes que Genteca.\n\n"
    "Este mismo método — identificar, bautizar, registrar — se aplica a los otros seis candidatos. Thermo-Safe es el prototipo. Ahora escala."
)

# ═══════════════════════════════════════════
# S-11 PORTAFOLIO ACTUAL + PIPELINE FUTURO
# ═══════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, BLANCO)

add_title_bar(s,
    "Cada Algoritmo Propio es un Activo en Espera de Nombre.",
    subtitle="La estrategia no termina en estos 7 — comienza con ellos.",
    bg=NEGRO, title_color=BLANCO, sub_color=NARANJA,
    title_size=24, sub_size=12)

# Zona verde — Batch 1
add_rect(s, 0.3, 1.75, 12.7, 0.4, VERDE_OK)
add_textbox(s, "BATCH 1 — LISTO PARA FILING",
            left=0.4, top=1.78, width=12.5, height=0.35,
            font_name="Futura", font_size=11, bold=True,
            color=BLANCO, align=PP_ALIGN.CENTER)

lineas_b1 = [
    ("Línea GSM", "StaggerStart™  /  Thermo-Safe™  /  FlickerGuard™"),
    ("Línea Genius", "ForensLog™  /  TripleLock™  /  ThermalCurve™"),
    ("Línea GRN", "TaskMemory™"),
]
for i, (linea, marcas) in enumerate(lineas_b1):
    add_rect(s, 0.3, 2.22 + i * 0.78, 12.7, 0.72, rgb(0xE8, 0xF5, 0xE9))
    add_textbox(s, linea,
                left=0.4, top=2.27 + i * 0.78, width=2.2, height=0.58,
                font_name="Futura", font_size=10, bold=True,
                color=VERDE_OK, align=PP_ALIGN.LEFT)
    add_textbox(s, marcas,
                left=2.7, top=2.27 + i * 0.78, width=10.0, height=0.58,
                font_name="Futura", font_size=12, bold=True,
                color=NEGRO, align=PP_ALIGN.LEFT)

# Zona amarilla — Batch 2
add_rect(s, 0.3, 4.6, 12.7, 0.4, AMARILLO)
add_textbox(s, "BATCH 2 — DOCUMENTACIÓN PENDIENTE",
            left=0.4, top=4.63, width=12.5, height=0.35,
            font_name="Futura", font_size=11, bold=True,
            color=NEGRO, align=PP_ALIGN.CENTER)

pendientes = [
    "NTC en GSM-MB/RB/RE  (HDE uso público pendiente)",
    "Curva inversa t-v en GST-R  (documentación interna pendiente)",
    "<30ms en GST-R  (HDE pendiente)",
]
for i, p in enumerate(pendientes):
    add_rect(s, 0.3, 5.07 + i * 0.42, 12.7, 0.38, rgb(0xFF, 0xFD, 0xE7))
    add_textbox(s, p,
                left=0.5, top=5.1 + i * 0.42, width=12.3, height=0.32,
                font_name="Futura", font_size=10,
                color=NEGRO, align=PP_ALIGN.LEFT)

# Pie
add_textbox(s,
    "Advertencia: Candidato a registro — pendiente de documentación técnica interna. No incluido en primer batch SAPI.",
    left=0.3, top=6.42, width=12.7, height=0.3,
    font_name="Gill Sans MT", font_size=7, italic=True,
    color=GRIS_COOL, align=PP_ALIGN.LEFT)

add_footer(s)
add_speaker_notes(s,
    "Esta diapositiva muestra el estado real del pipeline.\n\n"
    "El Batch 1 está materialmente listo para iniciar el proceso con el abogado marcario. "
    "StaggerStart y ForensLog tienen la documentación completa y el análisis preliminar sin colisiones.\n\n"
    "FlickerGuard está en Batch 1 con una advertencia: requiere búsqueda OMPI ampliada antes de presentar la solicitud. "
    "Esto es trabajo de semanas, no meses.\n\n"
    "El Batch 2 depende de que el equipo de ingeniería complete la documentación técnica interna — las HDE para uso público. "
    "Esto no bloquea el Batch 1. Los dos pueden correr en paralelo.\n\n"
    "Y más allá de estos 7: cada algoritmo propio que Exceline Profesional desarrolle desde hoy entra automáticamente al pipeline de bautismo y registro. "
    "La estrategia no es un proyecto. Es un proceso institucional permanente."
)

# ═══════════════════════════════════════════
# S-12 ARQUITECTURA 3 NIVELES
# ═══════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, BLANCO)

add_title_bar(s,
    "Arquitectura de Tres Niveles: Sin Crear Nueva Marca Paraguas.",
    subtitle="Exceline Profesional como contenedor — los nombres bautizados como propiedades.",
    bg=NEGRO, title_color=BLANCO, sub_color=NARANJA,
    title_size=22, sub_size=12)

# Nivel 1 — gris arriba (más ancho)
add_rect(s, 1.5, 1.75, 10.3, 1.4, GRIS_CLARO)
add_textbox(s, "NIVEL 1 — Marca contenedor (existente)",
            left=1.6, top=1.8, width=10.1, height=0.38,
            font_name="Futura", font_size=9, bold=True,
            color=GRIS_COOL, align=PP_ALIGN.CENTER)
add_textbox(s, "Exceline Profesional™",
            left=1.6, top=2.2, width=10.1, height=0.8,
            font_name="Futura", font_size=28, bold=True,
            color=NEGRO, align=PP_ALIGN.CENTER)

# Nivel 2 — naranja (medio)
add_rect(s, 2.5, 3.25, 8.3, 1.65, rgb(0xFF, 0xF3, 0xE0))
add_rect(s, 2.5, 3.25, 8.3, 0.38, NARANJA)
add_textbox(s, "NIVEL 2 — Nombres bautizados de tecnología",
            left=2.6, top=3.27, width=8.1, height=0.32,
            font_name="Futura", font_size=9, bold=True,
            color=BLANCO, align=PP_ALIGN.CENTER)
add_textbox(s,
    "StaggerStart™  ·  ForensLog™  ·  FlickerGuard™\n"
    "Thermo-Safe™  ·  TripleLock™  ·  TaskMemory™  ·  ThermalCurve™",
    left=2.6, top=3.68, width=8.1, height=1.15,
    font_name="Futura", font_size=14, bold=True,
    color=NARANJA, align=PP_ALIGN.CENTER)

# Nivel 3 — gris claro abajo (más estrecho)
add_rect(s, 3.7, 5.0, 5.9, 1.35, rgb(0xFA, 0xFA, 0xFA))
add_textbox(s, "NIVEL 3 — Descriptores técnicos (empaque / datasheet)",
            left=3.8, top=5.05, width=5.7, height=0.38,
            font_name="Futura", font_size=9, bold=True,
            color=GRIS_COOL, align=PP_ALIGN.CENTER)
add_textbox(s, "<30ms  ·  100 fallas  ·  Triple bloqueo  ·  Curva inversa t-v  ·  <1% subcarga",
            left=3.8, top=5.48, width=5.7, height=0.75,
            font_name="Futura", font_size=10, italic=True,
            color=GRIS_COOL, align=PP_ALIGN.CENTER)

# Cita pie
add_rect(s, 0.3, 6.48, 12.7, 0.45, GRIS_CLARO)
add_textbox(s, '"Los descriptores explican.  Los nombres son la propiedad."',
            left=0.4, top=6.51, width=12.5, height=0.38,
            font_name="Futura", font_size=12, bold=True, italic=True,
            color=NARANJA, align=PP_ALIGN.CENTER)

add_footer(s)
add_speaker_notes(s,
    "Esta arquitectura responde a una pregunta que la Junta puede tener: ¿estamos creando una nueva marca paraguas? No.\n\n"
    "Exceline Profesional ya existe como contenedor. No tocamos esa decisión de marca.\n\n"
    "Lo que registramos en el Nivel 2 son los nombres de las tecnologías propietarias dentro de ese contenedor. "
    "Son propiedades individuales — cada una con su propio certificado SAPI, su propia vigencia de 10 años, su propia oponibilidad.\n\n"
    "El Nivel 3 — los descriptores técnicos — no se registra. Va al empaque, al datasheet, al manual. "
    "Ahí es exactamente donde pertenece: comunicando la especificación, no la marca.\n\n"
    "La regla es simple: los descriptores explican lo que el producto hace. Los nombres son lo que Genteca posee."
)

# ═══════════════════════════════════════════
# S-13 CANAL VE Y ANGLICISMO TÉCNICO
# ═══════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, BLANCO)

add_title_bar(s,
    "Canal Venezolano: el Anglicismo Técnico Funciona Cuando Identifica.",
    subtitle="El mercado local ya acepta esta convención — sin fricción.",
    bg=NEGRO, title_color=BLANCO, sub_color=NARANJA,
    title_size=22, sub_size=12)

# Columna izquierda
add_rect(s, 0.3, 1.75, 5.9, 4.55, GRIS_CLARO)
add_textbox(s, "El instalador técnico venezolano dice...",
            left=0.4, top=1.82, width=5.7, height=0.42,
            font_name="Futura", font_size=11, bold=True,
            color=GRIS_COOL, align=PP_ALIGN.CENTER)

terminos = ["inverter\n(no \"compresor de velocidad variable\")",
            "breaker", "relay", "encoder"]
for i, t in enumerate(terminos):
    add_rect(s, 0.4, 2.35 + i * 0.92, 5.7, 0.78, BLANCO)
    add_textbox(s, t,
                left=0.5, top=2.4 + i * 0.92, width=5.5, height=0.65,
                font_name="Futura", font_size=14, bold=True,
                color=NEGRO, align=PP_ALIGN.CENTER)

# Flecha central
add_textbox(s, "≡",
            left=6.25, top=3.5, width=1.0, height=0.8,
            font_name="Futura", font_size=40, bold=True,
            color=NARANJA, align=PP_ALIGN.CENTER)
add_textbox(s, "equivalencia\nconvencional",
            left=6.1, top=4.3, width=1.3, height=0.55,
            font_name="Futura", font_size=7, italic=True,
            color=GRIS_COOL, align=PP_ALIGN.CENTER)

# Columna derecha
add_rect(s, 7.4, 1.75, 5.6, 4.55, rgb(0xFF, 0xF3, 0xE0))
add_textbox(s, "Los candidatos Genteca",
            left=7.5, top=1.82, width=5.4, height=0.42,
            font_name="Futura", font_size=11, bold=True,
            color=NARANJA, align=PP_ALIGN.CENTER)

candidatos13 = ["StaggerStart™", "FlickerGuard™", "Thermo-Safe™", "TripleLock™"]
for i, c in enumerate(candidatos13):
    add_rect(s, 7.5, 2.35 + i * 0.92, 5.4, 0.78, rgb(0xFF, 0xFB, 0xF0))
    add_textbox(s, c,
                left=7.6, top=2.4 + i * 0.92, width=5.2, height=0.65,
                font_name="Futura", font_size=16, bold=True,
                color=NARANJA, align=PP_ALIGN.CENTER)

# Cita pie
add_rect(s, 0.3, 6.45, 12.7, 0.7, NARANJA)
add_textbox(s,
    '"Un técnico que recomienda \'el que tiene StaggerStart\' está haciendo lo mismo que el que recomienda \'el que tiene SIRIUS de Siemens\'."',
    left=0.4, top=6.48, width=12.5, height=0.62,
    font_name="Futura", font_size=10, italic=True,
    color=BLANCO, align=PP_ALIGN.CENTER)

add_footer(s)
add_speaker_notes(s,
    "La objeción de fondo sobre el anglicismo técnico en Venezuela ya tiene respuesta empírica — y el mercado la da cada día.\n\n"
    "El instalador técnico venezolano no dice 'compresor de velocidad variable'. Dice 'inverter'. No dice 'interruptor automático'. Dice 'breaker'. "
    "No dice 'rele'. Dice 'relay'. Son anglicismos técnicos adoptados por convención, sin fricción, por todo el canal profesional.\n\n"
    "Los candidatos de Genteca siguen exactamente ese mismo patrón. StaggerStart, FlickerGuard, Thermo-Safe — son exactamente el tipo de denominación "
    "que un técnico venezolano adoptará como referencia de producto, de la misma manera que adoptó SIRIUS como referencia de arranque Siemens.\n\n"
    "El mercado ya habla este idioma. Lo que proponemos es que Genteca lo registre antes de que alguien más lo haga."
)

# ═══════════════════════════════════════════
# S-14 LA DECISIÓN BINARIA
# ═══════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, BLANCO)

add_title_bar(s,
    "La Pregunta Binaria: ¿Adopta Genteca Esta Estrategia?",
    subtitle="Una sola decisión de Junta desbloquea el pipeline completo.",
    bg=NEGRO, title_color=BLANCO, sub_color=NARANJA,
    title_size=24, sub_size=12)

# Pregunta central
add_rect(s, 1.0, 1.75, 11.3, 1.55, GRIS_CLARO)
add_textbox(s,
    "¿Adopta Genteca la estrategia sistemática de bautizar cada innovación técnica propia\n"
    "con un nombre de marca registrable en SAPI — comenzando con los candidatos identificados?",
    left=1.1, top=1.82, width=11.1, height=1.4,
    font_name="Futura", font_size=13, bold=True,
    color=NEGRO, align=PP_ALIGN.CENTER)

# Bifurcación — Rama SÍ
add_rect(s, 0.3, 3.55, 5.9, 2.6, rgb(0xE8, 0xF5, 0xE9))
add_rect(s, 0.3, 3.55, 5.9, 0.48, VERDE_OK)
add_textbox(s, "SI",
            left=0.4, top=3.58, width=5.7, height=0.38,
            font_name="Futura", font_size=16, bold=True,
            color=BLANCO, align=PP_ALIGN.CENTER)
si_items = [
    "Pipeline de registro se activa",
    "Equipo legal instruye abogado marcario",
    "Genteca es primero en Clase 9 VE",
    "Activos IP en balance desde hoy",
]
bullet_block(s, si_items,
             left=0.4, top=4.1, width=5.7, height=1.9,
             font_name="Futura", font_size=11,
             color=rgb(0x1B, 0x5E, 0x20), bullet_char="✓")

# Rama NO
add_rect(s, 7.1, 3.55, 5.9, 2.6, rgb(0xFF, 0xEB, 0xEE))
add_rect(s, 7.1, 3.55, 5.9, 0.48, rgb(0xC6, 0x28, 0x28))
add_textbox(s, "NO",
            left=7.2, top=3.58, width=5.7, height=0.38,
            font_name="Futura", font_size=16, bold=True,
            color=BLANCO, align=PP_ALIGN.CENTER)
no_items = [
    "Nombres permanecen no registrados",
    "Riesgo de apropiación por terceros abierto",
    "Pipeline no se activa",
    "Ventana de primer registro se cierra",
]
bullet_block(s, no_items,
             left=7.2, top=4.1, width=5.7, height=1.9,
             font_name="Futura", font_size=11,
             color=rgb(0xC6, 0x28, 0x28), bullet_char="✗")

# Línea cierre naranja
add_rect(s, 0.3, 6.35, 12.7, 0.55, NARANJA)
add_textbox(s,
    "La estrategia no cuesta más si se ejecuta.  Solo cuesta más si se deja para después.",
    left=0.4, top=6.38, width=12.5, height=0.42,
    font_name="Futura", font_size=13, bold=True,
    color=BLANCO, align=PP_ALIGN.CENTER)

add_footer(s)
add_speaker_notes(s,
    "Esta es la diapositiva de decisión. La pregunta es binaria — Sí o No — y cada respuesta tiene consecuencias claras.\n\n"
    "Si la respuesta es Sí: el equipo legal recibe instrucciones hoy, el abogado marcario es seleccionado en los próximos 10 días, "
    "y Genteca inicia el proceso de ser el primer fabricante venezolano con marcas de algoritmos propietarios en Clase 9.\n\n"
    "Si la respuesta es No: los nombres desarrollados en este proceso permanecen sin protección. "
    "Cualquier tercero que los vea — en una feria, en un catálogo, en una conversación — puede registrarlos antes que Genteca. "
    "Y la ventana de primer registro se cierra de forma irrecuperable.\n\n"
    "El costo del registro es fijo y conocido. El costo de no registrar no tiene techo — y se descuenta del futuro."
)

# ═══════════════════════════════════════════
# S-15 PRIORIZACIÓN SI SÍ
# ═══════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, BLANCO)

add_title_bar(s,
    "Si la Respuesta es Sí: Tres Decisiones de Priorización.",
    subtitle="La Junta decide modalidad de avance — no los abogados.",
    bg=NEGRO, title_color=BLANCO, sub_color=NARANJA,
    title_size=24, sub_size=12)

# Decisión A
add_textbox(s, "Decisión A — ¿Modalidad de filing?",
            left=0.3, top=1.65, width=8.0, height=0.38,
            font_name="Futura", font_size=11, bold=True,
            color=NEGRO, align=PP_ALIGN.LEFT)

data15a = [
    ["Opción", "Marcas incluidas", "Inversión estimada", "Ventaja", "Consideración"],
    ["Opción 1\n7 simultáneas",
     "StaggerStart · ForensLog · FlickerGuard\nThermo-Safe · TripleLock · TaskMemory · ThermalCurve",
     "$6,160–$11,480",
     "Cobertura total inmediata",
     "Thermo-Safe y otros ⚠ aún tienen trabajo previo"],
    ["Opción 2\nPrimer batch",
     "StaggerStart · ForensLog · FlickerGuard",
     "$2,640–$4,920",
     "Menor inversión, más sólidos primero",
     "Restantes expuestos mientras se completa documentación"],
]

tbl15a = add_table(s, data15a,
                   col_widths=[1.5, 3.5, 1.8, 2.2, 2.6],
                   left=0.3, top=2.08,
                   row_height=0.62,
                   header_bg=NEGRO, header_fg=BLANCO,
                   font_size=8.5, font_name="Gill Sans MT")

# Decisión B
add_textbox(s, "Decisión B — ¿Quién gestiona?",
            left=0.3, top=3.45, width=8.0, height=0.38,
            font_name="Futura", font_size=11, bold=True,
            color=NEGRO, align=PP_ALIGN.LEFT)

add_rect(s, 0.3, 3.88, 12.7, 0.75, GRIS_CLARO)
gestiona = ("Abogado marcario VE — especialista Clase 9.  "
            "Genteca aporta: HDE + nombre candidato + lista de productos.  "
            "Plazo SAPI: 6–12 meses desde presentación FM-02.")
add_textbox(s, gestiona,
            left=0.4, top=3.92, width=12.5, height=0.65,
            font_name="Futura", font_size=10,
            color=NEGRO, align=PP_ALIGN.LEFT)

# Advertencia crítica — escalación
add_rect(s, 0.3, 4.72, 12.7, 1.55, rgb(0xFF, 0xEB, 0xEE))
add_rect(s, 0.3, 4.72, 12.7, 0.38, rgb(0xC6, 0x28, 0x28))
add_textbox(s, "ADVERTENCIA CRÍTICA — ESCALACIÓN OWNER #1 (URGENTE)",
            left=0.4, top=4.74, width=12.5, height=0.32,
            font_name="Futura", font_size=9, bold=True,
            color=BLANCO, align=PP_ALIGN.CENTER)
add_textbox(s,
    "Ningún material público debe afirmar \"Exceline® registrada\" sin certificado verificado. "
    "Verificar estado SAPI de \"Exceline\" como paso previo urgente antes de cualquier comunicación externa. "
    "Esta verificación debe instruirse al abogado marcario en la semana 1.",
    left=0.4, top=5.14, width=12.5, height=1.08,
    font_name="Futura", font_size=10,
    color=rgb(0xC6, 0x28, 0x28), align=PP_ALIGN.LEFT)

add_footer(s)
add_speaker_notes(s,
    "Asumiendo que la respuesta a la decisión del slide anterior es Sí, hay tres decisiones operativas que la Junta debe tomar hoy.\n\n"
    "Decisión A — Modalidad: La Opción 1 (7 marcas simultáneas) da cobertura total pero tiene un costo más alto y algunos candidatos aún requieren trabajo previo. "
    "La Opción 2 (primer batch) es más conservadora en inversión y solo incluye los candidatos con documentación completa.\n\n"
    "Decisión B — Gestión: El abogado marcario especialista en Clase 9 es quien presenta. Genteca aporta la documentación técnica. "
    "El plazo de SAPI es 6 a 12 meses — es independiente de nuestra velocidad de presentación.\n\n"
    "La advertencia al pie es urgente y no es negociable: antes de distribuir cualquier material externo que mencione Exceline con símbolo registrado, "
    "debemos verificar con el abogado que SAPI efectivamente tiene un certificado activo. "
    "Si no existe certificado, usamos exclusivamente el símbolo ™. El ® sin certificado es una afirmación falsa con consecuencias legales."
)

# ═══════════════════════════════════════════
# S-16 PRÓXIMOS PASOS + TIMELINE
# ═══════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, BLANCO)

add_title_bar(s,
    "De la Decisión de Hoy a los Primeros Certificados: 6–12 Meses.",
    subtitle="Ruta crítica post-aprobación de Junta.",
    bg=NEGRO, title_color=BLANCO, sub_color=NARANJA,
    title_size=24, sub_size=12)

hitos = [
    ("SEM 1-2\nURGENTE", NARANJA, BLANCO,
     "• Verificar SAPI Exceline\n• Seleccionar abogado marcario\n• Instruir addendum Vera\n• Definir modalidad filing"),
    ("SEM 3-4", GRIS_COOL, BLANCO,
     "• Búsqueda fonética batch 1\n• Completar HDE batch 1"),
    ("MES 2", GRIS_CLARO, NEGRO,
     "• Solicitud FM-02 batch 1\n• Búsqueda OMPI ampliada\n• Publicación oposición 30 días"),
    ("MES 2-3", GRIS_CLARO, NEGRO,
     "• Completar HDE batch 2\n• Filing batch 2 (si aprobado)"),
    ("MES 6-12", VERDE_OK, BLANCO,
     "• Certificados batch 1\n• Filing batch 2\n• Pipeline permanente activo"),
]

hito_w = 2.45
for i, (periodo, bg, fg, contenido) in enumerate(hitos):
    add_rect(s, 0.2 + i * 2.6, 1.72, hito_w, 0.55, bg)
    add_textbox(s, periodo,
                left=0.2 + i * 2.6, top=1.74, width=hito_w, height=0.48,
                font_name="Futura", font_size=10, bold=True,
                color=fg, align=PP_ALIGN.CENTER)
    add_rect(s, 0.2 + i * 2.6, 2.32, hito_w, 3.3, rgb(0xF9, 0xF9, 0xF9) if bg == GRIS_CLARO else rgb(0xFF, 0xF3, 0xE0) if bg == NARANJA else rgb(0xE8, 0xF5, 0xE9) if bg == VERDE_OK else rgb(0xF5, 0xF5, 0xF5))
    add_textbox(s, contenido,
                left=0.25 + i * 2.6, top=2.37, width=hito_w - 0.1, height=3.15,
                font_name="Futura", font_size=9,
                color=NEGRO, align=PP_ALIGN.LEFT)

# Línea conectora
add_rect(s, 0.2, 2.57, 12.9, 0.06, GRIS_MED)

# Recuadro pie
add_rect(s, 0.3, 5.8, 12.7, 0.75, NEGRO)
add_textbox(s, "™ ahora — ® solo con certificado SAPI emitido y verificado",
            left=0.4, top=5.85, width=12.5, height=0.6,
            font_name="Futura", font_size=16, bold=True,
            color=NARANJA, align=PP_ALIGN.CENTER)

add_footer(s)
add_speaker_notes(s,
    "Esta es la ruta crítica desde hoy hasta los primeros certificados.\n\n"
    "Las semanas 1 y 2 son de acción urgente: verificación del estado de Exceline en SAPI, selección del abogado, instrucción al equipo técnico.\n\n"
    "Las semanas 3 y 4 son de trabajo técnico-legal: búsqueda fonética y finalización de documentación.\n\n"
    "El Mes 2 es cuando se presenta la solicitud FM-02 ante SAPI. A partir de ahí, el proceso está en manos de SAPI.\n\n"
    "El rango de 6 a 12 meses para los certificados es el tiempo oficial de procesamiento SAPI. No es una estimación pesimista — es el estándar histórico.\n\n"
    "La regla de cierre es la más importante: ™ se puede usar desde hoy, en cualquier material, sin certificado. "
    "® se usa exclusivamente después de que el abogado confirme que SAPI emitió el certificado. Sin excepción."
)

# ═══════════════════════════════════════════
# S-A1 APÉNDICE: PROCESO SAPI
# ═══════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, BLANCO)

add_rect(s, 0, 0, 13.333, 1.4, GRIS_CLARO)
add_textbox(s, "APÉNDICE",
            left=0.35, top=0.08, width=2.5, height=0.38,
            font_name="Futura", font_size=11, bold=True,
            color=NARANJA, align=PP_ALIGN.LEFT)
add_textbox(s, "El Proceso SAPI: 4 Etapas, 6–12 Meses, Clase 9 Confirmada.",
            left=0.35, top=0.12, width=12.6, height=0.62,
            font_name="Futura", font_size=20, bold=True,
            color=NEGRO, align=PP_ALIGN.LEFT)
add_textbox(s, "Base de proceso para el equipo legal y el abogado marcario.",
            left=0.35, top=0.82, width=12.6, height=0.42,
            font_name="Futura", font_size=12,
            color=NARANJA, align=PP_ALIGN.LEFT)

data_a1 = [
    ["Etapa", "Qué ocurre", "Quién actúa", "Duración"],
    ["1. Búsqueda fonética",
     "SAPI verifica si existe marca igual o confusamente similar registrada",
     "Abogado + SAPI", "2–4 semanas"],
    ["2. Solicitud FM-02",
     "Presentación formal de la solicitud de registro con documentación técnica",
     "Abogado", "1–2 semanas"],
    ["3. Publicación + oposición",
     "SAPI publica en boletín oficial; terceros tienen 30 días para oponerse",
     "SAPI / terceros", "30 días"],
    ["4. Emisión de certificado",
     "SAPI emite el certificado de registro — vigencia 10 años renovables",
     "SAPI", "Variable"],
]

tbl_a1 = add_table(s, data_a1,
                   col_widths=[2.5, 5.3, 2.5, 1.9],
                   left=0.35, top=1.55,
                   row_height=0.72,
                   header_bg=NEGRO, header_fg=BLANCO,
                   font_size=10, font_name="Gill Sans MT")

# Clase 9
add_rect(s, 0.35, 4.6, 12.6, 1.35, rgb(0xFF, 0xF3, 0xE0))
add_textbox(s, "Clase 9 Niza — confirmada para los 7 candidatos",
            left=0.45, top=4.65, width=12.4, height=0.4,
            font_name="Futura", font_size=11, bold=True,
            color=NARANJA, align=PP_ALIGN.LEFT)
add_textbox(s,
    "Clase 9: Aparatos e instrumentos científicos, náuticos, geodésicos, fotográficos, cinematográficos, "
    "ópticos, de pesaje, de medida, de señalización, de control (inspección), de salvamento y de enseñanza; "
    "aparatos e instrumentos para la conducción, distribución, transformación, acumulación, regulación o "
    "control de la electricidad. — Clasificador Internacional de Niza.",
    left=0.45, top=5.1, width=12.4, height=0.75,
    font_name="Gill Sans MT", font_size=8.5, italic=True,
    color=NEGRO, align=PP_ALIGN.LEFT)

add_textbox(s,
    "Costo por marca con abogado: $880–$1,640.  Vigencia: 10 años renovables "
    "(verificar interpretación LPI 1955 con abogado marcario).",
    left=0.35, top=6.05, width=12.6, height=0.35,
    font_name="Gill Sans MT", font_size=8,
    color=GRIS_COOL, align=PP_ALIGN.LEFT)

add_footer(s)
add_speaker_notes(s,
    "Este apéndice documenta el proceso SAPI para referencia del equipo legal y del abogado marcario.\n\n"
    "Las 4 etapas son secuenciales. El cuello de botella principal es la etapa de publicación y oposición — son 30 días que no se pueden acortar.\n\n"
    "Clase 9 de Niza es la clasificación correcta para todos los candidatos: cubre instrumentos de control, regulación y distribución de electricidad. "
    "Un protector eléctrico con algoritmos propietarios cae claramente dentro de esta clase.\n\n"
    "El rango de costo — $880 a $1,640 por marca — incluye honorarios del abogado y tasas SAPI. "
    "La variabilidad depende del abogado y de la complejidad de la búsqueda fonética."
)

# ═══════════════════════════════════════════
# S-A2 APÉNDICE: MATRIZ IP COMPLETA
# ═══════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, BLANCO)

add_rect(s, 0, 0, 13.333, 1.4, GRIS_CLARO)
add_textbox(s, "APÉNDICE",
            left=0.35, top=0.08, width=2.5, height=0.38,
            font_name="Futura", font_size=11, bold=True,
            color=NARANJA, align=PP_ALIGN.LEFT)
add_textbox(s, "Matriz IP Completa: Los 7 Candidatos en Detalle.",
            left=0.35, top=0.12, width=12.6, height=0.62,
            font_name="Futura", font_size=20, bold=True,
            color=NEGRO, align=PP_ALIGN.LEFT)
add_textbox(s, "Estado, riesgo, acción requerida y batch asignado por candidato.",
            left=0.35, top=0.82, width=12.6, height=0.42,
            font_name="Futura", font_size=12,
            color=NARANJA, align=PP_ALIGN.LEFT)

data_a2 = [
    ["#", "Marca", "Tecnología ancla", "Producto", "Estado OMPI", "Riesgo", "Acción previa filing", "Batch"],
    ["1", "StaggerStart™", "Staggered Intelligent Reconnect", "GSM-MB/RB/RE", "Preliminar OK", "Bajo", "HDE lista", "Batch 1"],
    ["2", "ForensLog™", "100-Fault Forensic History", "GIII+MV", "Preliminar OK", "Bajo", "HDE lista", "Batch 1"],
    ["3", "FlickerGuard™", "20ms Flicker Detection", "GSM-L", "Búsqueda ampliada pendiente", "Medio", "Búsqueda OMPI", "Batch 1 c/adv."],
    ["4", "Thermo-Safe™", "Función térmica integrada", "GSM-MB/RB/RE", "Búsqueda ampliada pendiente", "Medio-Alto", "OMPI + HDE NTC", "Batch 2"],
    ["5", "TripleLock™", "Triple-Fault Lockout", "GIII+/GUCT/GOC", "Búsqueda sectorial pendiente", "Medio", "Búsqueda sectorial", "Batch 2"],
    ["6", "TaskMemory™", "Smart Task Memory", "GRN-MV", "Búsqueda software urgente", "Alto", "OMPI urgente", "Batch 2"],
    ["7", "ThermalCurve™", "Cold/Hot Thermal Curve Model", "GIII+/GUCT/GOC", "No iniciada", "Alto", "HDE + evaluación desc.", "Pendiente"],
]

tbl_a2 = add_table(s, data_a2,
                   col_widths=[0.35, 1.55, 2.4, 1.7, 1.9, 1.0, 1.8, 1.2],
                   left=0.1, top=1.55,
                   row_height=0.47,
                   header_bg=NEGRO, header_fg=BLANCO,
                   font_size=7.5, font_name="Gill Sans MT")

# Color Riesgo
riesgo_colors = {
    "Bajo": (rgb(0xE8, 0xF5, 0xE9), rgb(0x1B, 0x5E, 0x20)),
    "Medio": (rgb(0xFF, 0xF9, 0xC4), rgb(0x7B, 0x61, 0x00)),
    "Medio-Alto": (rgb(0xFF, 0xF3, 0xE0), rgb(0xE6, 0x51, 0x00)),
    "Alto": (rgb(0xFF, 0xEB, 0xEE), rgb(0xC6, 0x28, 0x28)),
}
riesgo_vals = ["Bajo", "Bajo", "Medio", "Medio-Alto", "Medio", "Alto", "Alto"]
for ri, rv in enumerate(riesgo_vals):
    bg, fg = riesgo_colors[rv]
    color_table_cell(tbl_a2, ri + 1, 5, bg, fg)

# Advertencias pie
adv_a2 = (
    "Advertencia A — FlickerGuard™: Candidato Clase 9 SAPI Venezuela. Verificación preliminar completada. Requiere búsqueda OMPI ampliada antes del filing.\n"
    "Advertencia B — Thermo-Safe™: Candidato Clase 9 SAPI Venezuela. Requiere búsqueda OMPI ampliada Y documentación HDE pública NTC antes del filing.\n"
    "Advertencia C — TripleLock™: Candidato Clase 9 SAPI Venezuela. Requiere búsqueda sectorial ampliada (potencial colisión sector seguridad física — término \"Lock\").\n"
    "Advertencia D — TaskMemory™: Candidato Clase 9 SAPI Venezuela. Riesgo colisión marcas software — búsqueda urgente antes del filing.\n"
    "Advertencia E — ThermalCurve™: Candidato a registro — pendiente documentación técnica interna. No incluido en primer batch SAPI."
)
add_textbox(s, adv_a2,
            left=0.1, top=5.02, width=13.1, height=1.2,
            font_name="Gill Sans MT", font_size=6,
            color=GRIS_COOL, align=PP_ALIGN.LEFT)

add_footer(s)
add_speaker_notes(s,
    "Esta es la matriz de referencia completa para el equipo legal y el abogado marcario.\n\n"
    "Los campos clave son Estado OMPI, Riesgo y Acción previa al filing. Estas tres columnas determinan el orden y la modalidad de presentación.\n\n"
    "Batch 1 son los candidatos que pueden moverse a presentación FM-02 con trabajo mínimo de preparación.\n\n"
    "Batch 2 son los candidatos que requieren trabajo técnico adicional — búsquedas más amplias o documentación interna — antes de poder presentar.\n\n"
    "ThermalCurve está marcado como Pendiente porque el riesgo de descriptividad no está resuelto. Necesita evaluación legal específica antes de determinar si es registrable en su forma actual."
)

# ═══════════════════════════════════════════
# S-A3 APÉNDICE: MARCO LEGAL
# ═══════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, BLANCO)

add_rect(s, 0, 0, 13.333, 1.4, GRIS_CLARO)
add_textbox(s, "APÉNDICE",
            left=0.35, top=0.08, width=2.5, height=0.38,
            font_name="Futura", font_size=11, bold=True,
            color=NARANJA, align=PP_ALIGN.LEFT)
add_textbox(s, "Marco Legal Aplicado: LPI 1955 y Doctrina de Equivalencia Perceptual.",
            left=0.35, top=0.12, width=12.6, height=0.62,
            font_name="Futura", font_size=18, bold=True,
            color=NEGRO, align=PP_ALIGN.LEFT)
add_textbox(s, "Base jurídica para el abogado marcario y la Junta.",
            left=0.35, top=0.82, width=12.6, height=0.42,
            font_name="Futura", font_size=12,
            color=NARANJA, align=PP_ALIGN.LEFT)

# Columna izquierda — articulado
add_rect(s, 0.3, 1.55, 6.1, 4.7, GRIS_CLARO)
add_textbox(s, "ARTICULADO LEGAL",
            left=0.4, top=1.62, width=5.9, height=0.38,
            font_name="Futura", font_size=10, bold=True,
            color=GRIS_COOL, align=PP_ALIGN.CENTER)

legal_text = (
    "Art. 34 LPI — Causales de inadmisibilidad:\n"
    "• Signos descriptivos del producto o sus características\n"
    "• Signos genéricos del producto\n"
    "• Signos que induzcan error sobre el origen\n\n"
    "Art. 34 NO prohíbe:\n"
    "• Palabras extranjeras no descriptoras literales\n"
    "• Signos de fantasía en cualquier idioma\n"
    "• Neologismos o combinaciones originales\n\n"
    "Doctrina de equivalencia perceptual:\n"
    "Un signo es inadmisible si el público lo percibe\n"
    "como descriptivo — independientemente del idioma.\n"
    "Un signo de fantasía es admisible — si no describe\n"
    "directamente el producto, aunque sea en inglés."
)
add_textbox(s, legal_text,
            left=0.4, top=2.08, width=5.9, height=4.0,
            font_name="Gill Sans MT", font_size=9.5,
            color=NEGRO, align=PP_ALIGN.LEFT)

# Columna derecha — ejemplos
add_rect(s, 6.65, 1.55, 6.3, 4.7, BLANCO)
add_textbox(s, "EJEMPLOS DE VALIDACIÓN",
            left=6.75, top=1.62, width=6.1, height=0.38,
            font_name="Futura", font_size=10, bold=True,
            color=NARANJA, align=PP_ALIGN.CENTER)

ej = [
    ('"Escudo Térmico NTC"', rgb(0xFF, 0xEB, 0xEE), ROJO_TEXTO,
     "Riesgo Art. 34: describe literalmente función + componente."),
    ('"Thermo-Safe™"', rgb(0xE8, 0xF5, 0xE9), rgb(0x1B, 0x5E, 0x20),
     "Distintivo: identifica función de seguridad sin describir tecnología."),
    ('"ThermoShield™"', rgb(0xFF, 0xF9, 0xC4), rgb(0x7B, 0x61, 0x00),
     "Zona gris: Shield puede percibirse como descriptor. Requiere análisis."),
]
for i, (nm, bg, fg, desc) in enumerate(ej):
    add_rect(s, 6.75, 2.12 + i * 1.3, 6.1, 1.15, bg)
    add_textbox(s, nm,
                left=6.85, top=2.17 + i * 1.3, width=6.0, height=0.48,
                font_name="Futura", font_size=14, bold=True,
                color=fg, align=PP_ALIGN.CENTER)
    add_textbox(s, desc,
                left=6.85, top=2.65 + i * 1.3, width=6.0, height=0.52,
                font_name="Gill Sans MT", font_size=9,
                color=NEGRO, align=PP_ALIGN.LEFT)

# Recuadro naranja
add_rect(s, 0.3, 6.38, 12.7, 0.62, rgb(0xFF, 0xF3, 0xE0))
add_rect(s, 0.3, 6.38, 0.08, 0.62, NARANJA)
add_textbox(s,
    "Regla de símbolo:  ™ = uso en comercio sin certificado SAPI.  "
    "® = uso exclusivo después de emitido y verificado el certificado SAPI.",
    left=0.45, top=6.42, width=12.5, height=0.52,
    font_name="Futura", font_size=10, bold=True,
    color=NARANJA, align=PP_ALIGN.LEFT)

add_footer(s)
add_speaker_notes(s,
    "Este apéndice es para el abogado marcario y para cualquier miembro de la Junta que quiera entender la base jurídica de la estrategia.\n\n"
    "El Art. 34 de la LPI de 1955 establece las causales de inadmisibilidad. Ninguna de ellas menciona el idioma extranjero. "
    "La causal relevante es la descriptividad — y la doctrina de equivalencia perceptual determina que la prueba es cómo el público percibe el signo, no en qué idioma está escrito.\n\n"
    "Los tres ejemplos ilustran el espectro: un signo que describe claramente es inadmisible, un signo de fantasía es admisible, y una zona gris requiere análisis más profundo.\n\n"
    "La regla de símbolo al pie es operativa para toda la organización: ™ desde hoy, ® solo con papel en mano."
)

# ═══════════════════════════════════════════
# S-A4 APÉNDICE: LAS 3 ESCALACIONES
# ═══════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, BLANCO)

add_rect(s, 0, 0, 13.333, 1.4, GRIS_CLARO)
add_textbox(s, "APÉNDICE",
            left=0.35, top=0.08, width=2.5, height=0.38,
            font_name="Futura", font_size=11, bold=True,
            color=NARANJA, align=PP_ALIGN.LEFT)
add_textbox(s, "Las Tres Escalaciones que el Owner Debe Resolver Antes del Filing.",
            left=0.35, top=0.12, width=12.6, height=0.62,
            font_name="Futura", font_size=19, bold=True,
            color=NEGRO, align=PP_ALIGN.LEFT)
add_textbox(s, "Decisiones que no puede tomar el abogado — las toma la Junta.",
            left=0.35, top=0.82, width=12.6, height=0.42,
            font_name="Futura", font_size=12,
            color=NARANJA, align=PP_ALIGN.LEFT)

# Bloque 1 — URGENTE rojo
add_rect(s, 0.3, 1.55, 12.7, 1.55, rgb(0xFF, 0xEB, 0xEE))
add_rect(s, 0.3, 1.55, 12.7, 0.42, rgb(0xC6, 0x28, 0x28))
add_textbox(s, "ESCALACIÓN 1 — URGENTE: Verificación estado SAPI de \"Exceline\"",
            left=0.4, top=1.57, width=12.5, height=0.36,
            font_name="Futura", font_size=10, bold=True,
            color=BLANCO, align=PP_ALIGN.LEFT)
add_textbox(s,
    "Acción: Owner instruye abogado marcario para verificar si \"Exceline\" tiene certificado SAPI activo y vigente.\n"
    "Plazo: Antes de cualquier distribución externa de materiales con símbolo ®.\n"
    "Responsable: Owner + abogado marcario.",
    left=0.4, top=2.0, width=12.5, height=1.02,
    font_name="Futura", font_size=9.5,
    color=ROJO_TEXTO, align=PP_ALIGN.LEFT)

# Bloque 2 — amarillo
add_rect(s, 0.3, 3.2, 12.7, 1.55, rgb(0xFF, 0xFD, 0xE7))
add_rect(s, 0.3, 3.2, 12.7, 0.42, AMARILLO)
add_textbox(s, "ESCALACIÓN 2: Addendum HDE Vera — NTC GSM, curva inversa t-v GST-R, <30ms GST-R",
            left=0.4, top=3.22, width=12.5, height=0.36,
            font_name="Futura", font_size=10, bold=True,
            color=NEGRO, align=PP_ALIGN.LEFT)
add_textbox(s,
    "Acción: Owner instruye a Vera completar addendum HDE para las 3 funciones pendientes de documentación pública.\n"
    "Plazo: Antes del filing del Batch 2.\n"
    "Responsable: Owner → Vera → I+D.",
    left=0.4, top=3.65, width=12.5, height=1.02,
    font_name="Futura", font_size=9.5,
    color=NEGRO, align=PP_ALIGN.LEFT)

# Bloque 3 — azul marino
add_rect(s, 0.3, 4.85, 12.7, 1.65, AZUL_CLARO)
add_rect(s, 0.3, 4.85, 12.7, 0.42, AZUL_MARINO)
add_textbox(s, "ESCALACIÓN 3: Decisión de presupuesto IP — bloquea instrucción al abogado",
            left=0.4, top=4.87, width=12.5, height=0.36,
            font_name="Futura", font_size=10, bold=True,
            color=BLANCO, align=PP_ALIGN.LEFT)
add_textbox(s,
    "Opción A: 7 marcas simultáneas — inversión estimada $6,160–$11,480.\n"
    "Opción B: Primer batch (StaggerStart · ForensLog · FlickerGuard) — inversión estimada $2,640–$4,920.\n"
    "Esta decisión no puede delegarse al abogado — es una decisión de Owner y Junta.\n"
    "Responsable: Owner / Junta — hoy.",
    left=0.4, top=5.3, width=12.5, height=1.12,
    font_name="Futura", font_size=9.5,
    color=AZUL_MARINO, align=PP_ALIGN.LEFT)

add_footer(s)
add_speaker_notes(s,
    "Este apéndice documenta las tres escalaciones que requieren decisión del Owner y la Junta — no del equipo técnico ni del abogado.\n\n"
    "Escalación 1 es urgente e inmediata: verificar el estado real de Exceline en SAPI. Si no hay certificado activo, no puede usarse ® en ningún material. "
    "Esta verificación no puede esperar al inicio del proceso de registro de los 7 candidatos.\n\n"
    "Escalación 2 es la instrucción a Vera para completar la documentación técnica de las funciones que aún no tienen HDE para uso público. "
    "Sin esa documentación, el Batch 2 no puede presentarse.\n\n"
    "Escalación 3 es la decisión de presupuesto que bloquea la instrucción al abogado. "
    "El abogado no puede empezar sin saber cuántas marcas debe gestionar. "
    "Esta es la decisión más operativa — y es la que más directamente afecta la velocidad de implementación."
)

# ─────────────────────────────────────────────
# GUARDAR
# ─────────────────────────────────────────────
prs.save(OUTPUT)
print(f"OK — archivo guardado: {OUTPUT}")
print(f"Total slides: {len(prs.slides)}")
